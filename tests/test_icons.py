"""Tests for icons.py -- KONE pictograms as native PowerPoint shapes.

The path conversion is tested on hand-written path data so the rules
are pinned regardless of which sprite revision is installed; the tests
that read the real sprite guard against a revision that stops parsing,
which would silently yield zero icons rather than fail.
"""

import pytest

from deckguard import icons


def _ops(data):
    return icons.path_to_drawingml(data)


# --------------------------------------------------------------------------
# path conversion
# --------------------------------------------------------------------------


def test_absolute_move_and_line():
    assert _ops("M10 20 L30 40") == [("moveTo", [(10, 20)]), ("lnTo", [(30, 40)])]


def test_relative_commands_accumulate_from_the_current_point():
    """Everything after the opening M in this sprite is relative, so
    getting this wrong would misplace all 609 icons rather than fail."""
    assert _ops("M10 10 l5 5 l5 5") == [
        ("moveTo", [(10, 10)]), ("lnTo", [(15, 15)]), ("lnTo", [(20, 20)])
    ]


def test_horizontal_and_vertical_shorthands_become_full_points():
    """DrawingML has no h/v -- they carry the unchanged axis across."""
    assert _ops("M10 10 h20 v30 H5 V0") == [
        ("moveTo", [(10, 10)]), ("lnTo", [(30, 10)]), ("lnTo", [(30, 40)]),
        ("lnTo", [(5, 40)]), ("lnTo", [(5, 0)]),
    ]


def test_a_cubic_curve_keeps_both_control_points():
    ops = _ops("M0 0 C10 0 20 10 20 20")
    assert ops[1] == ("cubicBezTo", [(10, 0), (20, 10), (20, 20)])


def test_a_smooth_curve_reflects_the_previous_control_point():
    """`s` gives only the second control point; the first is the mirror
    of the previous curve's through the current point. 269 of the 609
    icons use this, so a wrong reflection dents a lot of outlines."""
    ops = _ops("M0 0 C10 0 20 10 20 20 S40 30 40 40")
    # previous control was (20, 10), current point (20, 20) -> (20, 30)
    assert ops[2] == ("cubicBezTo", [(20, 30), (40, 30), (40, 40)])


def test_a_smooth_curve_with_no_preceding_curve_starts_at_the_point():
    ops = _ops("M10 10 S30 20 30 30")
    assert ops[1] == ("cubicBezTo", [(10, 10), (30, 20), (30, 30)])


def test_close_returns_to_the_start_of_the_subpath():
    ops = _ops("M10 10 L50 10 Z l5 5")
    assert ops[2] == ("close", [])
    assert ops[3] == ("lnTo", [(15, 15)])  # relative to (10,10), not (50,10)


def test_extra_pairs_after_a_moveto_are_implicit_linetos():
    assert _ops("M0 0 10 10 20 20") == [
        ("moveTo", [(0, 0)]), ("lnTo", [(10, 10)]), ("lnTo", [(20, 20)])
    ]


def test_several_subpaths_are_kept_in_one_path_so_holes_stay_holes():
    """A counter is a subpath wound the other way. Split across separate
    `<a:path>` elements the hole fills in solid."""
    geom = icons.custgeom_xml(_ops("M0 0 h10 v10 h-10 Z M2 2 h6 v6 h-6 Z"))
    paths = geom.findall(f".//{{{icons.A_NS}}}path")
    assert len(paths) == 1
    assert len(paths[0].findall(f"{{{icons.A_NS}}}moveTo")) == 2


def test_custgeom_declares_the_icon_coordinate_space():
    geom = icons.custgeom_xml(_ops("M0 0 h10"))
    path = geom.find(f".//{{{icons.A_NS}}}path")
    assert path.get("w") == str(icons.VIEWBOX) == path.get("h")


# --------------------------------------------------------------------------
# the shipped sprite
# --------------------------------------------------------------------------

needs_sprite = pytest.mark.skipif(not icons.load_icons(), reason="icon sprite not installed")


@needs_sprite
def test_the_sprite_yields_the_expected_catalogue():
    names = icons.icon_names()
    assert len(names) > 500
    assert "elevator" in names and "escalator" in names
    assert not any(n.startswith("i-") for n in names), "the sprite prefix is stripped"


@needs_sprite
def test_every_icon_converts_to_a_well_formed_path():
    """The failure this guards is silent: a sprite revision that stops
    matching produces empty geometry, and an empty custGeom draws
    nothing at all rather than raising."""
    for name, data in icons.load_icons().items():
        ops = icons.path_to_drawingml(data)
        assert ops, f"{name} produced no geometry"
        assert ops[0][0] == "moveTo", f"{name} does not open with a moveTo"


@needs_sprite
def test_no_icon_needs_an_elliptical_arc():
    """The whole approach rests on this. DrawingML has no arc segment,
    so an icon using one would need decomposing into beziers -- and
    would silently lose that segment here."""
    import re

    offenders = [n for n, d in icons.load_icons().items() if re.search(r"[Aa]", d)]
    assert offenders == []


@needs_sprite
def test_find_icons_prefers_exact_then_prefix():
    assert icons.find_icons("elevator")[0] == "elevator"
    assert all("cloud" in n for n in icons.find_icons("cloud"))
    assert icons.find_icons("") == []


@needs_sprite
def test_an_icon_lands_as_an_editable_shape_not_a_picture():
    """The point of the exercise: a designer can select it and edit the
    points, it recolours like any shape, and the file carries no image
    part for it at all."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = icons.add_icon(slide, "elevator", (100, 100, 96, 96))

    assert shape is not None
    assert shape.shape_type.name == "FREEFORM"
    assert shape._element.spPr.find(f"{{{icons.A_NS}}}custGeom") is not None
    assert str(shape.fill.fore_color.rgb) == icons.KONE_BLUE
    assert not [s for s in slide.shapes if s.shape_type == 13], "no pictures"


@needs_sprite
def test_an_icon_carries_no_theme_drop_shadow():
    """`add_shape` attaches a `<p:style>` with `effectRef idx=2`, which
    renders as a drop shadow -- visible on every pictogram until the
    style block was removed. A pictogram with a shadow is off-brand."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = icons.add_icon(slide, "elevator", (10, 10, 96, 96))
    assert "<p:style>" not in shape._element.xml


@needs_sprite
def test_an_icon_is_centred_in_a_wider_box_and_hugs_what_it_draws():
    """The shape is sized to the geometry, not to the nominal 1024
    frame, so it is not square unless the glyph is -- but the frame is
    still what sets the SCALE, which is what keeps a grid of icons
    looking the same size."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = icons.add_icon(slide, "elevator", (0, 0, 200, 100))

    # inside the box on the short axis, and centred on the long one
    assert 0 <= shape.top and shape.top + shape.height <= 200 * 9525
    assert 50 * 9525 <= shape.left <= 150 * 9525
    # and no bigger than the frame it was scaled against
    assert shape.width <= 100 * 9525 * 1.3


@needs_sprite
def test_two_icons_in_one_grid_share_a_scale():
    """Sizing each shape to its own geometry must not resize the icon:
    a wide glyph and a tall one still have to read as the same size."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    wide = icons.add_icon(slide, "wind", (0, 0, 100, 100))
    tall = icons.add_icon(slide, "elevator", (200, 0, 100, 100))
    biggest = max(wide.width, wide.height, tall.width, tall.height)
    smallest = min(max(wide.width, wide.height), max(tall.width, tall.height))
    assert biggest / smallest < 1.3


@needs_sprite
def test_an_unknown_icon_returns_none_rather_than_drawing_a_wrong_one():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert icons.add_icon(slide, "not-a-real-icon", (0, 0, 96, 96)) is None
    assert len(slide.shapes) == 0


def test_no_sprite_means_no_icons_rather_than_a_crash(monkeypatch):
    monkeypatch.setattr(icons, "sprite_path", lambda: None)
    icons.load_icons.cache_clear()
    try:
        assert icons.load_icons() == {}
        assert icons.icon_names() == []
    finally:
        icons.load_icons.cache_clear()


@needs_sprite
def test_most_of_the_set_escapes_its_own_viewbox():
    """Not a defect to fix in the assets -- a fact to design around.
    452 of the 609 overflow by 5-15%, which is ordinary glyph overshoot;
    six overflow by more than half and are genuinely broken."""
    overflows = [icons.viewbox_overflow(n) for n in icons.icon_names()]
    assert sum(1 for o in overflows if o > 0) > 400
    badly = [n for n in icons.icon_names() if icons.viewbox_overflow(n) > 0.5]
    assert "onbattery" in badly and len(badly) < 20


def test_clamping_holds_geometry_inside_the_viewbox():
    ops = [("moveTo", [(-40.0, 500.0)]), ("lnTo", [(2302.0, 1500.0)])]
    assert icons.clamp_to_viewbox(ops) == [
        ("moveTo", [(0.0, 500.0)]), ("lnTo", [(float(icons.VIEWBOX), float(icons.VIEWBOX))])
    ]


def test_clamping_leaves_a_well_behaved_icon_untouched():
    ops = [("moveTo", [(10.0, 20.0)]), ("cubicBezTo", [(1.0, 2.0), (3.0, 4.0), (5.0, 6.0)])]
    assert icons.clamp_to_viewbox(ops) == ops


@needs_sprite
def test_a_badly_overflowing_icon_is_contained_by_default():
    """`onbattery` reaches y=2302 against a 1024 viewBox, and unlike the
    ordinary 5-15% overshoot that is a junk tail rather than part of the
    drawing -- it spilled down four rows of an icon grid. Only the six
    in BROKEN_ICONS are clamped; clamping the rest chopped the base off
    a bucket and an elevator car."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = icons.add_icon(slide, "onbattery", (100, 100, 96, 96))

    path = shape._element.spPr.find(f"{{{icons.A_NS}}}custGeom").find(
        f".//{{{icons.A_NS}}}path")
    coords = [int(p.get("x")) for p in path.iter(f"{{{icons.A_NS}}}pt")]
    coords += [int(p.get("y")) for p in path.iter(f"{{{icons.A_NS}}}pt")]
    assert min(coords) >= 0 and max(coords) <= icons.VIEWBOX


@needs_sprite
def test_an_ordinary_icon_keeps_its_overshoot():
    """Every symbol declares `overflow="visible"`, so the 5-15% most
    icons extend past the nominal frame is intentional and drawn.
    Clamping it -- which this did by default for a while -- flattened
    the bottom off a mop and bucket, reported as icons looking chopped."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = icons.add_icon(slide, "cleaning", (0, 0, 96, 96))
    path = shape._element.spPr.find(f"{{{icons.A_NS}}}custGeom").find(
        f".//{{{icons.A_NS}}}path")

    ops = icons.path_to_drawingml(icons.load_icons()["cleaning"])
    def span(o):
        ys = [y for _, pts in o for _, y in pts]
        return max(ys) - min(ys)
    assert icons.viewbox_overflow("cleaning") > 0, "this icon overshoots its frame"
    assert int(path.get("h")) == round(span(ops)), "the overshoot must survive"
    assert round(span(ops)) > round(span(icons.clamp_to_viewbox(ops)))


@needs_sprite
def test_the_sprite_declares_its_overflow_visible():
    """The fact the clamping bug turned on. Checked against the asset
    rather than against a harness of my own making -- wrapping a path in
    a fresh <svg> without this attribute clips it, and that is how I
    convinced myself browsers clip these."""
    import re

    text = icons.sprite_path().read_text(errors="replace")
    symbols = re.findall(r"<symbol[^>]*>", text)
    assert symbols and all('overflow="visible"' in s for s in symbols)
