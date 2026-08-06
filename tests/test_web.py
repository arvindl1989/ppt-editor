"""Tests for the web app's single Transform capability (plan -> review
-> execute -> audit) plus read-only Audit, auth, and download safety.
No real network calls: AI-dependent paths run keyless (deterministic
degradation) or with the planner monkeypatched."""

import importlib
import re

from fastapi.testclient import TestClient

from tests.helpers import add_rectangle, add_slide, body_run, new_deck, set_run, title_run


def _write_violating_deck(path):
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Calibri", color_hex="005EB8")
    prs.save(str(path))


def _write_three_slide_deck(path):
    prs = new_deck()
    c = add_slide(prs)
    title_run(c).text = "Annual Review"
    m = add_slide(prs)
    title_run(m).text = "Resolution rate"
    body_run(m).text = "91.2% resolved"
    e = add_slide(prs)
    title_run(e).text = "Thank you"
    prs.save(str(path))


def _client(tmp_path, monkeypatch, password=None):
    monkeypatch.setenv("DECKGUARD_WEB_STORAGE", str(tmp_path / "storage"))
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    if password:
        monkeypatch.setenv("DECKGUARD_WEB_PASSWORD", password)
    else:
        monkeypatch.delenv("DECKGUARD_WEB_PASSWORD", raising=False)
    from deckguard import web as web_mod

    importlib.reload(web_mod)  # pick up the env vars set above
    return TestClient(web_mod.app), web_mod


def _post_plan(client, deck_path, reference_path=None):
    files = {"file": (deck_path.name, deck_path.open("rb"), "application/octet-stream")}
    if reference_path is not None:
        files["reference"] = (reference_path.name, reference_path.open("rb"), "application/octet-stream")
    return client.post("/plan", files=files)


# --------------------------------------------------------------------------
# home / audit
# --------------------------------------------------------------------------


def test_index_renders_the_single_transform_tool(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    resp = client.get("/")
    assert resp.status_code == 200
    assert 'action="/plan"' in resp.text
    assert 'formaction="/audit"' in resp.text
    # the four old tabs are gone for real
    for retired in ('action="/fix"', 'action="/learn"', 'action="/create"', 'action="/redesign"'):
        assert retired not in resp.text


def test_audit_flow(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)

    with deck.open("rb") as f:
        resp = client.post("/audit", files={"file": ("d.pptx", f, "application/octet-stream")})
    assert resp.status_code == 200
    assert "critical" in resp.text
    assert "/download/" in resp.text
    assert "Transform this deck" in resp.text  # audit points at the one tool that can fix it


def test_audit_requires_a_file(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    resp = client.post("/audit", data={})
    assert resp.status_code == 400
    assert "Audit needs an uploaded" in resp.text


def test_corrupt_file_gives_clean_error_not_500(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not really a pptx")
    with bad.open("rb") as f:
        resp = client.post("/audit", files={"file": ("bad.pptx", f, "application/octet-stream")})
    assert resp.status_code == 400
    assert "valid .pptx" in resp.text
    assert "Traceback" not in resp.text


# --------------------------------------------------------------------------
# plan -> review
# --------------------------------------------------------------------------


def test_plan_renders_review_page_with_per_slide_actions(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    deck = tmp_path / "d.pptx"
    _write_three_slide_deck(deck)

    resp = _post_plan(client, deck)

    assert resp.status_code == 200
    assert "Review plan" in resp.text
    # one action radio set per slide, defaulting to rebuild; keep is always offered
    for idx in (1, 2, 3):
        assert f'name="action_{idx}"' in resp.text
    assert 'value="rebuild" checked' in resp.text
    assert 'value="keep"' in resp.text
    # previews present: current wireframes + proposed layout boxes
    assert "aspect-ratio:1280/720" in resp.text
    assert "Cover B" in resp.text  # the cover swap proposal is visible
    m = re.search(r'action="/transform/([a-f0-9]+)"', resp.text)
    assert m, "review form must post to /transform/{token}"


def test_plan_requires_deck_or_brief(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    resp = client.post("/plan", data={})
    assert resp.status_code == 400
    assert "Upload a deck" in resp.text


def test_plan_brief_only_is_gated_on_server_api_key(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)  # no ANTHROPIC_API_KEY
    resp = client.post("/plan", data={"brief": "a deck about elevators"})
    assert resp.status_code == 400
    assert "ANTHROPIC_API_KEY" in resp.text


def test_plan_brief_only_renders_archetype_review(tmp_path, monkeypatch):
    monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key-never-used")
    client, web_mod = _client(tmp_path, monkeypatch)
    monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key-never-used")  # reload cleared it

    from deckguard.transform import SlidePlan, TransformPlan

    def fake_plan(brief, **kwargs):
        return TransformPlan(
            slides=[SlidePlan(index=1, default_action="new", archetype={
                "archetype": "hero_stat", "eyebrow": "KPI", "value": "91%", "caption": "c", "support": "s",
            })],
            ai_suggestions_ran=True, deck_title="Planned deck",
        )

    monkeypatch.setattr(web_mod, "plan_transform_from_brief", fake_plan)

    resp = client.post("/plan", data={"brief": "a deck about elevators"})

    assert resp.status_code == 200
    assert "Planned deck" in resp.text
    assert 'name="include_1"' in resp.text
    assert "hero_stat" in resp.text


# --------------------------------------------------------------------------
# transform (execute)
# --------------------------------------------------------------------------


def test_transform_executes_choices_and_downloads(tmp_path, monkeypatch):
    import io

    from pptx import Presentation

    client, _ = _client(tmp_path, monkeypatch)
    deck = tmp_path / "d.pptx"
    _write_three_slide_deck(deck)

    token = re.search(r'action="/transform/([a-f0-9]+)"', _post_plan(client, deck).text).group(1)

    resp = client.post(f"/transform/{token}", data={"action_1": "rebuild", "action_2": "keep", "action_3": "rebuild"})

    assert resp.status_code == 200
    assert "Transformed" in resp.text
    assert "Remaining findings" in resp.text

    m = re.search(r'/download/([a-f0-9]+)/transformed\.pptx', resp.text)
    assert m
    dl = client.get(m.group(0))
    assert dl.status_code == 200
    prs = Presentation(io.BytesIO(dl.content))
    assert len(prs.slides) == 3
    assert prs.slides[0].slide_layout.name == "Cover B"  # slide 1 rebuilt (cover swap)
    assert prs.slides[1].shapes.title.text_frame.text == "Resolution rate"  # slide 2 kept

    report = client.get(m.group(0).replace("transformed.pptx", "transform.json"))
    assert report.status_code == 200


def test_transform_with_reference_shows_similarity_report(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    deck = tmp_path / "d.pptx"
    _write_three_slide_deck(deck)
    reference = tmp_path / "ref.pptx"
    _write_three_slide_deck(reference)

    token = re.search(r'action="/transform/([a-f0-9]+)"', _post_plan(client, deck, reference).text).group(1)
    resp = client.post(f"/transform/{token}", data={"action_1": "keep", "action_2": "keep", "action_3": "keep"})

    assert resp.status_code == 200
    assert "Vs. reference deck" in resp.text
    assert "layouts match" in resp.text


def test_transform_expired_token_gives_clean_message(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    resp = client.post("/transform/deadbeef", data={})
    assert resp.status_code == 404
    assert "expired" in resp.text


# --------------------------------------------------------------------------
# download safety / auth
# --------------------------------------------------------------------------


def test_download_rejects_path_traversal_and_unknown_token(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch)
    assert client.get("/download/deadbeef/nothing.pptx").status_code == 404
    assert client.get("/download/deadbeef/..%2F..%2Fetc%2Fpasswd").status_code == 404
    assert client.get("/download/not-a-token!/x.pptx").status_code == 404


def test_password_gate(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch, password="s3cret")
    assert client.get("/").status_code == 401
    assert client.get("/", auth=("anyone", "s3cret")).status_code == 200
    assert client.get("/", auth=("anyone", "wrong")).status_code == 401


def test_result_page_names_slides_that_need_a_manual_redraw(tmp_path, monkeypatch):
    """The tool already detects slides the reference deck redrew from
    scratch; the result page has to actually say so, or the user finds
    out in the meeting."""
    from deckguard import webtemplates as tpl

    audit = {"summary": {"critical": 0, "major": 0, "minor": 0}, "violations": [],
             "suppressed_archetype_findings": 0}
    outcome = {"rebuilt": [1], "archetype_swapped": [], "reference_carryover": [2, 3],
               "kept": [], "layouts_used": {}, "needs_manual_redraw": [3, 4, 5],
               "duplicate_logos_removed": 1}
    html = tpl.transform_result_page("d.pptx", outcome, audit, None, {"pptx": "/a", "json": "/b"})

    assert "Slides 3, 4, 5 need a manual redraw" in html
    assert "duplicate logo(s) removed" in html

    quiet = tpl.transform_result_page(
        "d.pptx", {**outcome, "needs_manual_redraw": [], "duplicate_logos_removed": 0},
        audit, None, {"pptx": "/a", "json": "/b"},
    )
    assert "manual redraw" not in quiet


def test_a_keyless_server_still_offers_every_readable_slide_an_archetype(tmp_path, monkeypatch):
    """The complaint this closes: with no ANTHROPIC_API_KEY the route
    switched the archetype step off entirely, so dense slides showed
    "Structure kept" and nothing else. Structural matching runs
    regardless, and the card states what the mapping would cost."""
    client, _ = _client(tmp_path, monkeypatch)  # no ANTHROPIC_API_KEY
    deck = tmp_path / "d.pptx"
    _write_three_slide_deck(deck)

    resp = _post_plan(client, deck)

    assert resp.status_code == 200
    assert 'value="archetype"' in resp.text
    assert "Archetype suggestions are switched off" in resp.text  # honest about why


def test_the_plan_page_can_set_every_slide_at_once(tmp_path, monkeypatch):
    """A twenty-slide deck meant twenty individual decisions before the
    reviewer could press Transform, and most of them are the same one."""
    client, _ = _client(tmp_path, monkeypatch)
    deck = tmp_path / "d.pptx"
    _write_three_slide_deck(deck)

    resp = _post_plan(client, deck)

    assert resp.status_code == 200
    for action in ("archetype", "rebuild", "keep"):
        assert f'data-bulk="{action}"' in resp.text
    # the bar has to be INSIDE the form, or it scrolls with nothing to act on
    form = resp.text.index('action="/transform/')
    assert form < resp.text.index('class="card bulkbar"') < resp.text.index("</form>")
    # and the buttons must not submit the form on the way past
    bar = resp.text[resp.text.index('class="card bulkbar"'):]
    assert 'type="submit"' not in bar[:bar.index("</div>")]
    assert 'id="bulk-note"' in resp.text


def test_the_brief_plan_page_can_include_or_exclude_every_slide_at_once(tmp_path, monkeypatch):
    client, _ = _client(tmp_path, monkeypatch, password=None)
    monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key-not-used")
    import deckguard.web as web

    importlib.reload(web)
    from deckguard import webtemplates as templates

    entries = [{"index": i, "archetype_name": f"a{i}", "proposed_html": "<i>x</i>"}
               for i in (1, 2, 3)]
    html = templates.transform_review_page("b", "tok", entries, "brief", True)

    assert 'data-bulk-check="1"' in html and 'data-bulk-check="0"' in html
    assert html.count('name="include_') == 3
    assert html.index('class="card bulkbar"') < html.index('name="include_1"')


def test_a_slide_locked_to_one_action_is_not_reported_as_refusing_it():
    """A slide no template layout fits carries a hidden `keep` rather
    than radios. Counting only radios made "Keep as-is" report it as a
    slide that does not offer keeping -- the one thing it does offer."""
    from deckguard import webtemplates as templates

    entries = [
        {"index": 1, "default_action": "rebuild", "layout_name": "Two Content",
         "current_html": "<i>a</i>", "proposed_html": "<i>b</i>"},
        {"index": 2, "default_action": "keep", "reason": "dense copy",
         "current_html": "<i>a</i>"},
    ]
    html = templates.transform_review_page("d", "tok", entries, "deck", True)

    # slide 2 is locked: a hidden input, no radios
    assert '<input type="hidden" name="action_2" value="keep">' in html
    assert 'name="action_2" value="rebuild"' not in html
    # the script counts a hidden input already carrying the wanted value
    assert "i.type === 'hidden'" in html
