"""Tests for deckguard's outline-driven deck composition (compose):
generating a new, on-brand .pptx (or appending on-brand slides to an
existing one) from a structured content outline, using the org master
template's own sanctioned layouts.
"""

import textwrap

import pytest
from pptx import Presentation

from deckguard.compose import (
    ComposeError,
    Outline,
    SlideSpec,
    _chrome_idxs,
    _normalize_block,
    _select_layout,
    build_deck,
    load_outline,
)
from deckguard.config import default_config_path, load_config
from deckguard.inventory import build_inventory
from deckguard.rules_engine import audit_deck
from deckguard.slide_import import default_template_path
from tests.helpers import make_pattern_png

TEMPLATE_PATH = default_template_path()
pytestmark = pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="bundled template asset not present")


def _write_outline(tmp_path, yaml_text):
    path = tmp_path / "outline.yaml"
    path.write_text(textwrap.dedent(yaml_text), encoding="utf-8")
    return path


# --------------------------------------------------------------------------
# load_outline
# --------------------------------------------------------------------------


def test_load_outline_parses_slides(tmp_path):
    path = _write_outline(
        tmp_path,
        """
        slides:
          - kind: cover
            title: "Hello"
          - kind: content
            title: "Body"
            bullets: ["one", "two"]
        """,
    )
    outline = load_outline(path)
    assert len(outline.slides) == 2
    assert outline.slides[0].kind == "cover"
    assert outline.slides[1].bullets == ["one", "two"]


def test_load_outline_requires_slides_key(tmp_path):
    path = _write_outline(tmp_path, "title: not a deck\n")
    with pytest.raises(ComposeError, match="slides"):
        load_outline(path)


def test_load_outline_rejects_empty_slides_list(tmp_path):
    path = _write_outline(tmp_path, "slides: []\n")
    with pytest.raises(ComposeError, match="non-empty"):
        load_outline(path)


def test_load_outline_rejects_unknown_kind(tmp_path):
    path = _write_outline(tmp_path, "slides:\n  - kind: interpretive-dance\n")
    with pytest.raises(ComposeError, match="unknown kind"):
        load_outline(path)


def test_load_outline_requires_kind_field(tmp_path):
    path = _write_outline(tmp_path, "slides:\n  - title: no kind here\n")
    with pytest.raises(ComposeError, match="needs a 'kind'"):
        load_outline(path)


def test_load_outline_missing_file_raises_compose_error(tmp_path):
    with pytest.raises(ComposeError, match="not found"):
        load_outline(tmp_path / "nope.yaml")


# --------------------------------------------------------------------------
# _normalize_block
# --------------------------------------------------------------------------


def test_normalize_block_accepts_plain_strings():
    assert _normalize_block(["a", "b"]) == [(0, "a"), (0, "b")]


def test_normalize_block_accepts_level_dicts():
    block = _normalize_block(["top", {"level": 2, "text": "nested"}])
    assert block == [(0, "top"), (2, "nested")]


def test_normalize_block_rejects_bad_item():
    with pytest.raises(ComposeError):
        _normalize_block([123])


# --------------------------------------------------------------------------
# layout selection
# --------------------------------------------------------------------------


def _layouts_by_name():
    prs = Presentation(str(TEMPLATE_PATH))
    return {layout.name: layout for master in prs.slide_masters for layout in master.slide_layouts}


def test_select_layout_cover_default_and_variant():
    layouts = _layouts_by_name()
    cache = {}
    assert _select_layout(SlideSpec(kind="cover"), layouts, cache) == "Cover B"
    assert _select_layout(SlideSpec(kind="cover", variant="e"), layouts, cache) == "Cover E"


def test_select_layout_cover_rejects_bad_variant():
    with pytest.raises(ComposeError, match="A-F"):
        _select_layout(SlideSpec(kind="cover", variant="Z"), _layouts_by_name(), {})


def test_select_layout_quote_variant():
    layouts = _layouts_by_name()
    assert _select_layout(SlideSpec(kind="quote", variant="c"), layouts, {}) == "Quote C"


def test_select_layout_section_variant():
    layouts = _layouts_by_name()
    assert _select_layout(SlideSpec(kind="section"), layouts, {}) == "Section divider (just title)"
    assert _select_layout(SlideSpec(kind="section", variant="numbered"), layouts, {}) == "Section divider (numbering)"


def test_select_layout_content_grows_with_column_count():
    layouts = _layouts_by_name()
    cache = {}
    one = _select_layout(SlideSpec(kind="content", bullets=["a"]), layouts, cache)
    three = _select_layout(SlideSpec(kind="content", columns=[["a"], ["b"], ["c"]]), layouts, cache)
    assert one == "Title and content A"
    assert three.startswith("Three content")


def test_select_layout_explicit_override_wins():
    layouts = _layouts_by_name()
    spec = SlideSpec(kind="content", bullets=["a", "b", "c"], layout="Blank")
    assert _select_layout(spec, layouts, {}) == "Blank"


def test_select_layout_explicit_override_must_exist():
    with pytest.raises(ComposeError, match="does not exist"):
        _select_layout(SlideSpec(kind="content", layout="Not A Real Layout"), _layouts_by_name(), {})


def test_select_layout_stat_falls_back_when_more_blocks_than_any_candidate_holds():
    layouts = _layouts_by_name()
    huge = SlideSpec(kind="stat", stats=[{"number": str(i)} for i in range(50)])
    with pytest.raises(ComposeError, match="no layout fits"):
        _select_layout(huge, layouts, {})


def test_select_layout_content_with_images_picks_a_picture_capable_layout():
    """A plain content slide (no images) never needs a picture placeholder,
    so it should keep picking the same text-only layout as always -- adding
    picture-carrying candidates to CONTENT_LAYOUT_CANDIDATES must not change
    that. Once images are present, though, a layout with enough PICTURE
    placeholders needs to actually be selected (this used to be impossible:
    no candidate had any picture placeholder at all, so any content-kind
    slide with images raised "no layout fits")."""
    layouts = _layouts_by_name()
    no_images = _select_layout(SlideSpec(kind="content", title="T", bullets=["a", "b"]), layouts, {})
    assert no_images == "Title and content A"

    with_images = _select_layout(
        SlideSpec(kind="content", title="T", bullets=["a", "b"], images=[b"x", b"y"]), layouts, {}
    )
    assert with_images in layouts
    assert layouts[with_images].placeholders
    pic_count = sum(
        1 for ph in layouts[with_images].placeholders
        if ph.placeholder_format.type is not None and ph.placeholder_format.type.name == "PICTURE"
    )
    assert pic_count >= 2


# --------------------------------------------------------------------------
# _chrome_idxs — logo/tagline placeholders are identified, not populated
# --------------------------------------------------------------------------


def test_chrome_idxs_flags_logo_placeholder_by_name():
    layouts = _layouts_by_name()
    cover_a = layouts["Cover A"]
    logo_idx = next(ph.placeholder_format.idx for ph in cover_a.placeholders if "logo" in ph.name.lower())
    assert logo_idx in _chrome_idxs(cover_a)


def test_chrome_idxs_empty_for_layout_with_no_logo_placeholder():
    layouts = _layouts_by_name()
    # Cover B's logo is a static (non-placeholder) shape on the layout itself.
    assert _chrome_idxs(layouts["Cover B"]) == set()


# --------------------------------------------------------------------------
# build_deck — fresh create
# --------------------------------------------------------------------------


def _outline(*specs) -> Outline:
    return Outline(slides=list(specs))


def test_build_deck_creates_requested_slide_count_and_layouts(tmp_path):
    outline = _outline(
        SlideSpec(kind="cover", title="Cover"),
        SlideSpec(kind="agenda", title="Agenda", bullets=["One", "Two"]),
        SlideSpec(kind="end", title="Thank you"),
    )
    out_path = tmp_path / "deck.pptx"
    result = build_deck(outline, str(out_path))

    assert result.slide_count == 3
    assert result.layouts_used == ["Cover B", "Agenda A", "Outro"]
    prs = Presentation(str(out_path))
    assert len(prs.slides) == 3
    assert [s.slide_layout.name for s in prs.slides] == result.layouts_used


def test_build_deck_populates_title_and_bullets(tmp_path):
    outline = _outline(
        SlideSpec(kind="content", title="Three priorities", columns=[["Fast"], ["Safe"], ["Simple"]]),
    )
    out_path = tmp_path / "deck.pptx"
    build_deck(outline, str(out_path))

    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    texts = [shp.text_frame.text for shp in slide.shapes if shp.has_text_frame and shp.text_frame.text.strip()]
    assert "Three priorities" in texts
    assert "Fast" in texts and "Safe" in texts and "Simple" in texts


def test_build_deck_places_raw_image_bytes_not_just_file_paths(tmp_path):
    """SlideSpec.images has always accepted file paths (what a human-written
    YAML outline gives it); redesign.py's image carryover instead has raw
    image bytes on hand (straight from the source deck's shapes, no temp
    file involved) -- _fill_images needs to accept both."""
    png_bytes = make_pattern_png(tmp_path / "src.png", seed=7).read_bytes()
    outline = _outline(
        SlideSpec(kind="content", title="Has pictures", bullets=["a"], images=[png_bytes]),
    )
    out_path = tmp_path / "deck.pptx"
    build_deck(outline, str(out_path))

    prs = Presentation(str(out_path))
    blobs = []
    for shp in prs.slides[0].shapes:
        try:
            blobs.append(shp.image.blob)
        except (AttributeError, ValueError):
            continue
    assert png_bytes in blobs


def test_build_deck_quote_slide_places_quote_and_author(tmp_path):
    outline = _outline(
        SlideSpec(kind="quote", title="Voice of customer", quote_text="It just works.", quote_author="A happy user"),
    )
    out_path = tmp_path / "deck.pptx"
    build_deck(outline, str(out_path))

    prs = Presentation(str(out_path))
    texts = {
        shp.text_frame.text for shp in prs.slides[0].shapes
        if shp.has_text_frame and shp.text_frame.text.strip()
    }
    assert "It just works." in texts
    assert "A happy user" in texts


def test_build_deck_leaves_logo_placeholder_empty_so_it_inherits(tmp_path):
    """A logo/tagline placeholder is never written to -- its <p:spPr>
    stays empty so PowerPoint falls back to the layout's own baked-in
    image. Confirm compose never touches it."""
    outline = _outline(SlideSpec(kind="cover", title="Cover", variant="A"))
    out_path = tmp_path / "deck.pptx"
    build_deck(outline, str(out_path))

    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    layout = slide.slide_layout
    logo_idx = next(ph.placeholder_format.idx for ph in layout.placeholders if "logo" in ph.name.lower())
    logo_shape = next(ph for ph in slide.placeholders if ph.placeholder_format.idx == logo_idx)
    assert logo_shape.text_frame.text == ""


def test_build_deck_result_is_clean_by_construction(tmp_path):
    """Text with no all-caps-looking content should audit clean --
    zero critical or major violations -- since fix_deck's pass resolves
    the inherited-color findings before saving."""
    outline = _outline(
        SlideSpec(kind="cover", title="Quarterly review", subtitle="People flow, reimagined"),
        SlideSpec(kind="content", title="What matters", bullets=["Reliability", "Speed", "Trust"]),
    )
    out_path = tmp_path / "deck.pptx"
    result = build_deck(outline, str(out_path))
    assert result.manual_review == []

    prs = Presentation(str(out_path))
    violations = audit_deck(build_inventory(prs), load_config(default_config_path()))
    assert not any(v.severity in ("critical", "major") for v in violations)


# --------------------------------------------------------------------------
# build_deck — append to an existing deck
# --------------------------------------------------------------------------


def test_build_deck_append_preserves_existing_slides_untouched(tmp_path):
    legacy_path = tmp_path / "legacy.pptx"
    legacy = Presentation()
    legacy.slides.add_slide(legacy.slide_layouts[0]).shapes.title.text_frame.text = "Existing content"
    legacy.save(str(legacy_path))
    original_bytes = legacy_path.read_bytes()

    outline = _outline(SlideSpec(kind="section", title="New section"))
    out_path = tmp_path / "appended.pptx"
    result = build_deck(outline, str(out_path), existing_deck_path=str(legacy_path))

    # The source file itself was only ever read, never opened for writing.
    assert legacy_path.read_bytes() == original_bytes

    prs = Presentation(str(out_path))
    assert len(prs.slides) == 2
    assert result.slide_count == 1
    first_slide_text = [
        shp.text_frame.text for shp in prs.slides[0].shapes if shp.has_text_frame and shp.text_frame.text.strip()
    ]
    assert first_slide_text == ["Existing content"]
    new_slide_text = [
        shp.text_frame.text for shp in prs.slides[1].shapes if shp.has_text_frame and shp.text_frame.text.strip()
    ]
    assert "New section" in new_slide_text


def test_build_deck_append_scopes_manual_review_to_new_slides_only(tmp_path):
    legacy_path = tmp_path / "legacy.pptx"
    legacy = Presentation()
    legacy.slides.add_slide(legacy.slide_layouts[0]).shapes.title.text_frame.text = "OLD ALL CAPS TITLE"
    legacy.save(str(legacy_path))

    outline = _outline(SlideSpec(kind="section", title="A perfectly normal title"))
    out_path = tmp_path / "appended.pptx"
    result = build_deck(outline, str(out_path), existing_deck_path=str(legacy_path))

    # The pre-existing slide's own all-caps title is real, but out of
    # scope for THIS operation -- only new slides are ever reported on.
    assert result.manual_review == []
