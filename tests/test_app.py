"""The one-page tool: give it something, get a deck, edit it.

There is no mode to choose and no approval step, so these tests are
mostly about the four ways input arrives and what each produces.
"""

import importlib
import re

import pytest
from fastapi.testclient import TestClient


@pytest.fixture
def client():
    from deckguard import web as web_mod

    return TestClient(importlib.reload(web_mod).app)


def _token(html: str) -> str:
    found = re.search(r"/download/([a-f0-9]+)/deck\.pptx", html)
    assert found, "no deck was offered for download"
    return found.group(1)


def test_the_landing_page_asks_for_input_not_for_a_mode(client):
    r = client.get("/")
    assert r.status_code == 200
    assert "Build a KONE deck" in r.text
    # one button, not a choice of products
    assert r.text.count('type="submit"') == 1
    for gone in ("Audit only", "Plan transform", "/audit", "/plan"):
        assert gone not in r.text, gone


def test_picked_slides_are_honoured_exactly_and_in_order(client):
    r = client.post("/generate", data={
        "audience": "internal", "title": "Probe",
        "pick": ["internal:8", "internal:16"],
    })
    assert r.status_code == 200
    assert "numbered_icon_row_6" in r.text and "kone_numbers" in r.text
    body = r.text
    assert body.index("numbered_icon_row_6") < body.index("kone_numbers")


def test_giving_it_nothing_says_so_rather_than_building_an_empty_deck(client):
    r = client.post("/generate", data={"audience": "internal"})
    assert r.status_code == 400
    assert "something to work from" in r.text


def test_a_deck_is_downloadable_the_moment_it_is_built(client):
    """No approval step: generating produces the file."""
    r = client.post("/generate", data={"audience": "internal", "pick": ["internal:8"]})
    deck = client.get(f"/download/{_token(r.text)}/deck.pptx")
    assert deck.status_code == 200
    assert deck.content[:2] == b"PK"


def test_edits_reorder_drop_and_duplicate_then_rebuild(client):
    r = client.post("/generate", data={
        "audience": "internal", "pick": ["internal:8", "internal:16", "internal:20"]})
    token = _token(r.text)
    edited = client.post(f"/rebuild/{token}", data={
        "o:0": "30", "o:1": "10", "o:2": "20", "dup:2": "on",
    })
    assert edited.status_code == 200

    from io import BytesIO

    from pptx import Presentation

    deck = client.get(f"/download/{token}/deck.pptx")
    slides = Presentation(BytesIO(deck.content)).slides
    # three picked, one duplicated, plus the master's cover and outro
    assert len(slides) == 6


def test_an_expired_token_says_so(client):
    r = client.post("/rebuild/" + "0" * 32, data={})
    assert r.status_code == 404


# --------------------------------------------------------------------------
# using someone else's deck
# --------------------------------------------------------------------------


def _a_deck(tmp_path):
    from deckguard import assemble

    out = tmp_path / "source.pptx"
    plan = assemble.plan(audience="internal", title="Source",
                         picks=["internal:8", "internal:16", "internal:20"])
    assemble.build(plan, str(out))
    return out


def test_an_uploaded_deck_becomes_templates_that_actually_render(client, tmp_path):
    """The trap this keeps falling into: a template that is registered
    but has no content, or content but no role styles, builds a deck of
    blank pages and downloads happily."""
    source = _a_deck(tmp_path)
    with source.open("rb") as handle:
        r = client.post("/generate", files={"deck": ("mine.pptx", handle.read())},
                        data={"audience": "internal", "title": "From my deck"})
    assert r.status_code == 200
    assert "templates read from your deck" in r.text

    from io import BytesIO

    from pptx import Presentation

    deck = client.get(f"/download/{_token(r.text)}/deck.pptx")
    slides = list(Presentation(BytesIO(deck.content)).slides)
    body = slides[1:-1]          # between the master's retained cover and outro
    assert body, "nothing was built from the mined templates"
    with_text = [
        s for s in body
        if any(sh.text_frame.text.strip() for sh in s.shapes
               if getattr(sh, "has_text_frame", False))
    ]
    assert len(with_text) >= len(body) // 2, "mined slides came out blank"


def test_a_deck_that_cannot_be_read_is_reported_not_crashed(client):
    r = client.post("/generate", files={"deck": ("broken.pptx", b"not a pptx at all")},
                    data={"audience": "internal", "brief": "", "pick": ["internal:8"]})
    assert r.status_code == 200
    assert "Could not read that deck" in r.text


def test_only_pptx_is_accepted(client):
    r = client.post("/generate", files={"deck": ("notes.txt", b"hello")},
                    data={"audience": "internal"})
    assert r.status_code == 400
    assert ".pptx" in r.text


# --------------------------------------------------------------------------
# preflight
# --------------------------------------------------------------------------


def test_preflight_runs_on_every_build_and_says_when_it_is_clean(client):
    r = client.post("/generate", data={"audience": "internal", "pick": ["internal:8"]})
    assert "Preflight" in r.text


def test_preflight_ignores_the_masters_zero_sized_placeholders(tmp_path):
    """The master's latent DATE / FOOTER / SLIDE_NUMBER placeholders come
    back as 0x0 boxes parked at y=720. Flagging them made every deck
    report findings nobody could act on, which teaches people to ignore
    preflight entirely."""
    from deckguard import assemble

    out = tmp_path / "clean.pptx"
    plan = assemble.plan(audience="internal", title="T", picks=["internal:8", "internal:16"])
    checks = assemble.build(plan, str(out))
    floor_hits = [m for _n, m in checks["findings"] if "past the floor" in m]
    assert not floor_hits, floor_hits


def test_preflight_still_catches_a_real_violation(tmp_path):
    from pptx import Presentation
    from pptx.util import Emu, Pt

    from deckguard import assemble

    out = tmp_path / "dirty.pptx"
    plan = assemble.plan(audience="internal", title="T", picks=["internal:8"])
    assemble.build(plan, str(out))

    prs = Presentation(str(out))
    box = prs.slides[1].shapes.add_textbox(Emu(400000), Emu(400000), Emu(2000000), Emu(300000))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "- a dash pretending to be a bullet"
    run.font.size = Pt(12)
    prs.save(str(out))

    findings = assemble.preflight(str(out))["findings"]
    assert any("dash standing in for a bullet" in m for _n, m in findings), findings


def test_preflight_catches_blue_inter_and_a_bold_flag_but_not_a_blue_figure(tmp_path):
    """Two rules from a review of real output, and the exception that
    makes them safe to assert.

    "Inter is never blue" is how the rule gets quoted, and taken
    literally it is wrong: BRAND_MODE states it under *headlines and
    body*, then sets `stat_value` in `#1450F5` on purpose so a figure
    reads as the number rather than the pair. A blanket assertion would
    report every correct `kone_numbers` slide as a defect, which is the
    fastest way to get a check switched off.

    Weight is the other one: it comes from the SemiBold family, never
    from a bold flag on top of it.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt

    from deckguard import assemble

    out = tmp_path / "type.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def run(top, text, face, colour, bold=False):
        box = slide.shapes.add_textbox(Emu(400000), Emu(top), Emu(4000000), Emu(400000))
        piece = box.text_frame.paragraphs[0].add_run()
        piece.text = text
        piece.font.name, piece.font.bold = face, bold
        piece.font.size = Pt(14)
        piece.font.color.rgb = RGBColor.from_string(colour)

    run(400000, "Ajith Kumar", "Inter", "1450F5")                       # fault
    run(1000000, "Workstream heading", "Inter SemiBold", "141414", True)  # fault
    run(1600000, "70%", "Inter", "1450F5")                              # correct
    run(2200000, "SCOPE", "KONE Information", "1450F5")                 # correct
    prs.save(str(out))

    findings = [m for _n, m in assemble.preflight(str(out))["findings"]]
    assert any("Inter set in KONE Blue" in m and "Ajith" in m for m in findings), findings
    assert any("bold flag" in m for m in findings), findings
    assert not any("70%" in m for m in findings), "a blue stat figure is the brand's own"
    assert not any("SCOPE" in m for m in findings), "KONE Information may be blue"


# --------------------------------------------------------------------------
# the brief path
# --------------------------------------------------------------------------


class _Stream:
    def __init__(self, message):
        self._message = message

    def __enter__(self):
        return self

    def __exit__(self, *exc):
        return False

    def get_final_message(self):
        return self._message


class _Message:
    def __init__(self, payload):
        self.content = [type("Block", (), {"type": "text", "text": payload})()]
        self.stop_reason = "end_turn"
        self.usage = type("U", (), {"input_tokens": 10, "output_tokens": 20})()


class _FakeClient:
    """Stands in for `anthropic.Anthropic`, so the whole planning path
    runs without a key. Nothing else exercised it, which is how a
    missing `import anthropic` reached production: the module imported
    fine and the call was never made."""

    def __init__(self, payload):
        self.calls = []
        outer = self

        class _Messages:
            def stream(self, **kwargs):
                outer.calls.append(kwargs)
                return _Stream(_Message(payload))

        self.messages = _Messages()


_SPEC = """{"slides": [
  {"archetype": "title_content", "title": "Survey findings",
   "bullets": ["Twelve units past economic repair"]},
  {"archetype": "kone_numbers", "title": "At scale",
   "stats": [{"value": "12", "label": "Frontlines"}]}
]}"""


def test_a_brief_is_planned_into_slides(monkeypatch):
    from deckguard import planner

    client = _FakeClient(_SPEC)
    spec, usage = planner.call_claude_for_kone_spec(
        "We migrated Request Management to ServiceNow in six weeks.",
        api_key="test-key", client=client,
    )
    assert [s["archetype"] for s in spec["slides"]] == ["title_content", "kone_numbers"]
    assert usage.output_tokens == 20
    assert client.calls, "the planner never called the model"


def test_the_brief_path_builds_a_real_deck(tmp_path, monkeypatch):
    """End to end with the model faked: brief in, .pptx out."""
    from deckguard import assemble, planner

    monkeypatch.setattr(
        planner, "call_claude_for_kone_spec",
        lambda brief, **kw: (__import__("json").loads(_SPEC), None),
    )
    plan = assemble.plan(brief="Subject: ServiceNow migration complete",
                         audience="external", title="")
    assert plan["title"] == "ServiceNow migration complete", "title read off the subject line"
    out = tmp_path / "brief.pptx"
    checks = assemble.build(plan, str(out))
    assert out.is_file() and checks["slides"] >= 2


def test_a_planning_failure_is_reported_not_swallowed(client, monkeypatch):
    from deckguard import planner

    def _boom(*_a, **_kw):
        raise planner.PlanningError("ANTHROPIC_API_KEY is not set")

    monkeypatch.setattr(planner, "call_claude_for_kone_spec", _boom)
    r = client.post("/generate", data={"brief": "Build me a deck", "audience": "internal"})
    assert r.status_code == 400
    assert "ANTHROPIC_API_KEY" in r.text


def test_planner_translates_a_rate_limit_into_something_actionable():
    import anthropic

    from deckguard import planner

    import httpx

    response = httpx.Response(429, request=httpx.Request("POST", "https://api.anthropic.com"))

    class _Boom:
        class messages:
            @staticmethod
            def stream(**_kw):
                raise anthropic.APIStatusError("busy", response=response, body=None)

    with pytest.raises(planner.PlanningError, match="rate-limited or overloaded"):
        planner._stream_final_message(_Boom(), model="m", messages=[])


# --------------------------------------------------------------------------
# choosing slides, not marching down the list
# --------------------------------------------------------------------------


def _notes_for(audience="internal"):
    """The instruction the planner actually receives, captured without
    calling the model."""
    from deckguard import assemble, planner

    seen = {}

    def _capture(brief, **kw):
        seen["notes"] = kw.get("notes")
        return {"slides": [{"archetype": "title_content", "title": "x"}]}, None

    original = planner.call_claude_for_kone_spec
    planner.call_claude_for_kone_spec = _capture
    try:
        assemble.plan(brief="An email.", audience=audience)
    finally:
        planner.call_claude_for_kone_spec = original
    return seen["notes"]


def test_the_planner_is_given_a_menu_not_a_running_order():
    """"Keep them in the set's order" read as "march down this list from
    the top", and with bare names there was no other signal to go on. A
    brief came back using the same first handful every time."""
    notes = _notes_for("internal")
    assert "keep them in the set's order" not in notes.lower()
    assert "MENU, not a running order" in notes
    assert "Reach across the whole menu" in notes


def test_every_archetype_offered_comes_with_its_job():
    notes = _notes_for("external")
    from deckguard import brandmode as B

    for slide in B.slides_in("external"):
        assert f"{slide['archetype']} — " in notes, slide["archetype"]


def test_the_planner_is_told_not_to_reuse_a_layout_or_emit_an_outro():
    notes = _notes_for("internal")
    assert "ONCE unless the content genuinely repeats" in notes
    # the master's Thank you is retained, so a planned outro doubles it
    assert "Do NOT emit an outro" in notes


def test_variety_reports_repeats_but_forgives_dividers():
    """A divider once per section is correct; the same content layout
    five times is a planner that fell back on position."""
    from deckguard import assemble

    mix = assemble.variety({"slides": [
        {"archetype": "divider_numbering"}, {"archetype": "title_content"},
        {"archetype": "divider_numbering"}, {"archetype": "title_content"},
        {"archetype": "title_content"}, {"archetype": "hero_stat"},
    ]})
    assert mix["total"] == 6 and mix["distinct"] == 3
    assert mix["repeats"] == [("title_content", 3)], mix["repeats"]


def test_the_result_page_shows_how_many_distinct_layouts(client):
    r = client.post("/generate", data={
        "audience": "internal", "title": "Probe",
        "pick": ["internal:8", "internal:8", "internal:16"]})
    assert "Distinct layouts" in r.text
    assert "Reusing" in r.text, "a repeated content layout should be called out"


def test_a_failed_archetype_install_is_recorded_not_only_swallowed():
    """`layouts.install` is wrapped in a bare except so a missing spec
    file cannot stop the engine's own archetypes working. That same
    catch hid a NameError which silently disabled every derived
    archetype -- the registry just came back smaller and nothing said
    why."""
    from deckguard import registry

    registry._load_archetypes()
    assert registry.INSTALL_ERRORS == [], registry.INSTALL_ERRORS


def test_the_contact_line_clears_the_footer(tmp_path):
    """INTERNAL_25 slide 24 puts it at top:566. It shipped at 654 and
    landed exactly on the footer date at 658 -- and preflight missed the
    overlap by three pixels, because the floor check had been loosened
    to 680 to avoid flagging the footer itself."""
    from pptx import Presentation

    from deckguard import assemble, brandmode as bm

    out = tmp_path / "rl.pptx"
    plan = {"title": "T", "date": "1 March 2026", "slides": [{
        "archetype": "resource_links", "title": "Out of scope",
        "contact": "EUR and Nordics own these directly.",
        "tiles": [{"icon": "target", "label": "One"}],
    }]}
    checks = assemble.build(plan, str(out))
    assert not [m for _n, m in checks["findings"] if "floor" in m], checks["findings"]

    prs = Presentation(str(out))
    px = prs.slide_width / 1280
    contact = next(
        sh for sh in list(prs.slides)[1].shapes
        if getattr(sh, "has_text_frame", False) and "EUR and Nordics" in sh.text_frame.text
    )
    assert (contact.top + contact.height) / px <= bm.FLOOR + 1
    assert not contact.name.startswith("Chrome"), "content must not be named as chrome"


def test_chrome_is_named_so_content_can_be_held_to_the_real_floor(tmp_path):
    from pptx import Presentation

    from deckguard import assemble

    out = tmp_path / "c.pptx"
    assemble.build({"title": "T", "date": "1 March 2026", "slides": [
        {"archetype": "title_content", "title": "X", "bullets": ["One"]}]}, str(out))
    named = [sh.name for slide in Presentation(str(out)).slides for sh in slide.shapes
             if (sh.name or "").startswith("Chrome")]
    assert named, "the footer and page number should be named as chrome"


# --------------------------------------------------------------------------
# what the deck should cover
# --------------------------------------------------------------------------


def test_the_sections_a_user_ticked_reach_the_planner():
    from deckguard import assemble, planner

    seen = {}

    def _capture(brief, **kw):
        seen["notes"] = kw.get("notes")
        return {"slides": [{"archetype": "title_content", "title": "x"}]}, None

    original = planner.call_claude_for_kone_spec
    planner.call_claude_for_kone_spec = _capture
    try:
        assemble.plan(brief="An email.", audience="internal",
                      sections=["numbers", "timeline"])
    finally:
        planner.call_claude_for_kone_spec = original

    notes = seen["notes"]
    assert "MUST cover these" in notes
    from deckguard import brandmode as B

    assert B.DECK_SECTIONS["numbers"]["label"] in notes
    assert B.DECK_SECTIONS["timeline"]["label"] in notes
    # and it must name layouts that can carry them, or the instruction is
    # a wish rather than a choice
    assert any(a in notes for a in B.DECK_SECTIONS["timeline"]["internal"])


def test_no_sections_ticked_leaves_the_instruction_alone():
    notes = _notes_for("internal")
    assert "MUST cover these" not in notes


def test_every_section_offers_a_built_archetype_to_both_audiences():
    """A section whose shortlist names nothing the renderer can build is
    an instruction the planner cannot follow."""
    from deckguard import brandmode as B
    from deckguard.registry import _load_archetypes

    built = set(_load_archetypes().ARCHETYPES)
    for key, entry in B.DECK_SECTIONS.items():
        for audience in ("internal", "external"):
            names = entry[audience]
            assert names, f"{key}/{audience} has no shortlist"
            assert [n for n in names if n in built], f"{key}/{audience}: {names}"


def test_the_home_page_offers_the_sections_as_chips():
    from html import escape

    from deckguard import brandmode as B, screens

    html = screens.home()
    assert html.count('name="section"') == len(B.DECK_SECTIONS)
    for entry in B.DECK_SECTIONS.values():
        assert escape(entry["label"]) in html


# --------------------------------------------------------------------------
# a divider is a pause, so it reads centred
# --------------------------------------------------------------------------


def test_the_divider_sits_in_the_middle_of_the_page(tmp_path):
    """It shipped with the numeral and title jammed under the top edge.
    KONE's own dividers centre them."""
    from pptx import Presentation

    from deckguard import assemble

    out = tmp_path / "d.pptx"
    assemble.build({"title": "T", "date": "1 March 2026", "slides": [
        {"archetype": "divider_numbering", "number": "03",
         "eyebrow": "Marketing Hub", "title": "Brand guidelines"}]}, str(out))

    prs = Presentation(str(out))
    px = prs.slide_width / 1280
    shapes = {sh.text_frame.text.strip(): sh for sh in list(prs.slides)[1].shapes
              if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()}
    numeral = shapes["03"]
    mid = (numeral.top + numeral.height / 2) / px
    assert 300 <= mid <= 420, f"numeral centred at {mid:.0f}, want near 360"
    assert shapes["Brand guidelines"].top / px > 250


# --------------------------------------------------------------------------
# previews that look like the slide
# --------------------------------------------------------------------------


def test_the_set_archetypes_have_rendered_thumbnails():
    """Every archetype the picker offers AND can build should have a
    real picture of itself. A missing one is not fatal -- the wireframe
    still draws -- but it is a regeneration someone forgot."""
    from deckguard import thumbs
    from deckguard.registry import _load_archetypes

    built = set(_load_archetypes().ARCHETYPES)
    missing = [n for n in thumbs.set_archetypes()
               if n in built and thumbs.path_for(n) is None]
    assert not missing, f"run `python -m deckguard.thumbs`: {missing}"


def test_a_tile_shows_the_render_and_falls_back_to_the_wireframe():
    from deckguard import screens, thumbs

    name = next(n for n in thumbs.set_archetypes() if thumbs.path_for(n))
    assert f'src="/preview/{name}.png"' in screens._thumb(name)
    # nothing rendered for a mined design, so it must still draw something
    assert 'data-dg-frame="1"' in screens._thumb("no_such_archetype_at_all")


def test_the_preview_route_serves_a_png_and_refuses_a_traversal():
    from fastapi.testclient import TestClient

    from deckguard import thumbs
    from deckguard.web import app

    client = TestClient(app)
    name = next(n for n in thumbs.set_archetypes() if thumbs.path_for(n))
    ok = client.get(f"/preview/{name}.png")
    assert ok.status_code == 200 and ok.content[:4] == b"\x89PNG"
    assert client.get("/preview/..%2f..%2fweb.py.png").status_code == 404
    assert client.get("/preview/nope.png").status_code == 404


def test_grouped_sample_content_varies_between_cells():
    """Four identical cells preview as a rendering fault rather than as
    a row of four things."""
    from deckguard.preview import sample_content

    content = sample_content("kone_numbers")
    for value in content.values():
        if isinstance(value, list) and len(value) > 1 and isinstance(value[0], dict):
            rendered = ["|".join(str(v) for v in item.values()) for item in value]
            assert len(set(rendered)) > 1, rendered


def test_a_picture_slot_naming_a_missing_file_is_refilled():
    """The skill's own sample for `image_section_divider` pointed at a
    path that no longer exists. The engine drew its sand fallback under
    the scrim, so the slide came back as a dark smear with the headline
    invisible on it -- worse than an empty slot, and it read as a
    rendering bug rather than as a missing file."""
    import os

    from deckguard.registry import fill_empty_photo_slots

    spec = {"slides": [{"archetype": "image_section_divider",
                        "image": "/no/such/photo.jpg", "title": "Part two"}]}
    assert fill_empty_photo_slots(spec) == 1
    assert os.path.isfile(spec["slides"][0]["image"])


def test_a_photo_that_is_there_is_left_alone():
    from deckguard.registry import _photo_library, fill_empty_photo_slots

    mine = _photo_library()[0]
    spec = {"slides": [{"archetype": "image_section_divider", "image": mine}]}
    assert fill_empty_photo_slots(spec) == 0
    assert spec["slides"][0]["image"] == mine


def test_the_quote_slide_sets_its_quote_in_quote_type():
    """Ported from the master's boxes, every region came through as 16px
    body: the quotation sat in a 349px pink panel at footnote size and
    the attribution under it looked exactly the same."""
    from deckguard.registry import _load_archetypes

    from deckguard import brandmode as bm

    regions = {r["content"]: r for r in _load_archetypes().ARCHETYPES["quote_b"]["regions"]}
    # The role is now SETTLED at install rather than left for a draw-time
    # lookup that never passed a width -- so this is `quote_lg`, not the
    # unsized `quote`. Asserted through the brand rather than as a
    # literal, because the name is a consequence of the panel's width:
    # 657px is over the 600px threshold, so the quote takes 30px.
    quote = bm.resolve(regions["quote"]["role"])
    assert quote and quote["px"] == 30, regions["quote"]["role"]
    assert regions["quote"]["role"].startswith("quote")
    assert regions["attribution"]["role"] == "attribution"
    assert "body2" not in regions and "body3" not in regions


def test_a_cover_title_is_never_cut_mid_word():
    """A review of real output found the cover of a finished deck
    reading "...ONE Week MOD deployment with you reg" -- the brief's
    first line hard-cut at 70 characters and set at 76px.

    A character count is the wrong instrument for a headline. What
    replaces it is still a heuristic, and only runs when neither the
    user nor the planner supplied a title, but it never breaks a word.
    """
    from deckguard.assemble import TITLE_MAX, _title_from

    lines = [
        "I would like to share our plan of ONE Week MOD deployment with you "
        "regarding the September launch",
        "Subject: ONE Week MOD deployment",
        "Q3 results, EMEA and APAC, with commentary on the service business",
        "Short one",
        "A" * 120 + " tail",
        "",
    ]
    for line in lines:
        title = _title_from(line)
        if not title:
            continue
        # Every word in the result is a word that was in the source: a
        # mid-word cut would leave a fragment that is not.
        source = set(line.split())
        assert all(word in source for word in title.split()), (line, title)
        assert title == title.strip()

    long = _title_from(lines[0])
    assert len(long) <= TITLE_MAX
    assert not long.endswith((" with", " of", " the", " to"))
