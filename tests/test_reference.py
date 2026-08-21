"""What the tool learned from a real KONE deck.

"Life, upgraded in ONE week" -- a finished deck, not a spec. Its cover
and its card grid were measured rather than described, so these hold the
tool to the measurements: a cut cover whose panes are where that deck
puts them, and a card grid that reads like slide 8 without repeating its
two accessibility mistakes.
"""

import re

from deckguard import brandmode as bm
from deckguard import layouts as L


# --------------------------------------------------------------------------
# the cover
# --------------------------------------------------------------------------


def test_the_cut_is_measured_not_derived():
    """It used to be four equal panes on the content column, marked
    `derived` in the source. The reference bakes its cut into a
    transparent PNG; reading that alpha channel back gives panes 289px
    wide on a 330px pitch, full bleed, with staggered depths.

    The stagger is the thing: four equal panes read as a filmstrip."""
    cut = L._CUT_COVERS["cover_a_cut4"]
    assert cut["band"] == [0, 0, 1280, 422], "full bleed, not the content column"
    assert "derived" not in cut

    xs = [p[0] for p in cut["panes"]]
    widths = {p[1] for p in cut["panes"]}
    depths = [p[2] for p in cut["panes"]]
    assert xs == [0, 330, 660, 990]
    assert widths <= {289, 290}
    assert len(set(depths)) == 4, "every pane a different depth"
    assert xs[-1] + cut["panes"][-1][1] == 1280, "the last pane runs to the edge"


def test_every_deck_opens_on_the_cut(tmp_path):
    from pptx import Presentation

    from deckguard import assemble

    out = tmp_path / "d.pptx"
    assemble.build({"title": "T", "date": "1 March 2026", "slides": [
        {"archetype": "title_content", "title": "X", "bullets": ["One"]}]}, str(out))

    slides = list(Presentation(str(out)).slides)
    cover = slides[0]
    names = [sh.name for sh in cover.shapes]
    # Not by mask NAME: which renderer draws the cut depends on whether
    # the gallery is installed, and the two name their rectangles
    # differently. One banner with at least four shapes over it is the
    # cut either way.
    banners = [sh for sh in cover.shapes
               if sh.shape_type and sh.shape_type.name == "PICTURE"
               and sh.width > 3000000]
    assert len(banners) == 1, names
    assert len([sh for sh in cover.shapes
                if sh.shape_type and sh.shape_type.name == "AUTO_SHAPE"]) >= 4, names
    assert "Logo" in names and "Tagline" in names, names


def test_a_deck_that_already_has_a_cover_does_not_get_a_second(tmp_path):
    from pptx import Presentation

    from deckguard import assemble

    out = tmp_path / "d.pptx"
    assemble.build({"title": "T", "date": "1 March 2026", "slides": [
        {"archetype": "cover_f_fullbleed", "title": "X"},
        {"archetype": "title_content", "title": "Y", "bullets": ["One"]}]}, str(out))
    # cover, content, retained Thank you
    assert len(list(Presentation(str(out)).slides)) == 3


def test_a_cover_headline_is_set_at_cover_size():
    """Bound from the master it resolved to 30px, so a cover read like a
    content slide. The reference sets 76."""
    from deckguard.registry import _load_archetypes

    for name in ("cover_a_cut4", "cover_b_cut3", "cover_f_fullbleed"):
        spec = _load_archetypes().ARCHETYPES[name]
        title = next(r for r in spec["regions"] if r.get("content") == "title")
        assert title["role"].startswith("cover_title"), (name, title["role"])


def test_the_engine_knows_every_role_the_brand_defines():
    """The engine shipped 28 roles against the brand's 42, and asking a
    cover for `cover_title` killed the build on a KeyError mid-draw."""
    from deckguard.registry import _load_archetypes

    styles = _load_archetypes().E.ROLE_STYLE
    assert not [r for r in bm.TYPE_SCALE if r not in styles]


# --------------------------------------------------------------------------
# the card grid
# --------------------------------------------------------------------------


def test_the_grid_matches_the_reference_pitch():
    grid = L.card_grid(12)["cards"]
    assert len(grid) == 12
    assert sorted({c["box"][0] for c in grid}) == list(bm.CARD_COL_X)
    assert {c["box"][2] for c in grid} == {bm.CARD_W}


def test_the_grid_never_runs_past_the_floor():
    """The reference's own third row ends at y=681, past the floor. A
    twelve-cell grid shortens its cards rather than hanging them over
    the footer."""
    for cells in (1, 4, 8, 12):
        grid = L.card_grid(cells)["cards"]
        assert max(c["box"][1] + c["box"][3] for c in grid) <= bm.FLOOR


def test_a_short_grid_is_centred_rather_than_hung_from_the_top():
    """Eight cards in a three-row band left the bottom third empty."""
    eight = L.card_grid(8)["cards"]
    top = min(c["box"][1] for c in eight)
    bottom = max(c["box"][1] + c["box"][3] for c in eight)
    assert top > bm.TIGHT_CONTENT_Y, "the block is centred, not pinned"
    assert abs((top - bm.TIGHT_CONTENT_Y) - (bm.FLOOR - bottom)) <= 2


def test_the_accent_is_on_the_rule_and_never_on_the_type(tmp_path):
    """The reference sets its caps labels IN the accent -- mint on white,
    pale blue on white -- and those are the only two cards on the slide
    you cannot read. The rule carries the coding instead.

    Preflight is what enforces it, and it caught the arrow glyph being
    coloured the same way after the labels had been fixed."""
    from deckguard import assemble

    checks = assemble.build({"title": "T", "date": "1 March 2026", "audience": "internal",
        "slides": [{"archetype": "card_grid", "title": "Grid", "cards": [
            {"label": f"Card {n}", "text": "Something worth reading."}
            for n in range(1, 9)]}]}, str(tmp_path / "c.pptx"))
    off_brand = [m for _n, m in checks["findings"] if "not black, white or KONE Blue" in m]
    assert not off_brand, off_brand

    # ...and the rules DO carry the secondary palette
    accents = {c["accent"] for c in L.card_grid(8)["cards"]}
    assert bm.MINT in accents and bm.PINK in accents


def test_a_card_is_rounded_and_carries_a_shadow(tmp_path):
    """The one place the square-corner rule is relaxed, deliberately."""
    from pptx import Presentation

    from deckguard import assemble

    out = tmp_path / "c.pptx"
    assemble.build({"title": "T", "date": "1 March 2026", "slides": [
        {"archetype": "card_grid", "title": "Grid",
         "cards": [{"label": "One", "text": "x"}]}]}, str(out))

    cards = [sh for sh in list(Presentation(str(out)).slides)[1].shapes
             if (sh.name or "") == "Card"]
    assert cards, "no card drawn"
    assert "roundRect" in cards[0]._element.xml
    assert "outerShdw" in cards[0]._element.xml


# --------------------------------------------------------------------------
# the tighter band
# --------------------------------------------------------------------------


def test_a_slide_with_no_eyebrow_moves_up_to_the_reference_band():
    from deckguard.registry import _load_archetypes

    spec = _load_archetypes().ARCHETYPES["agenda_c_split"]
    title = next(r for r in spec["regions"] if r.get("content") == "title")
    assert title["box"][1] == bm.TIGHT_TITLE_Y


def test_a_slide_with_an_eyebrow_keeps_its_band():
    """Shifting the block by the title's delta drove the eyebrow off the
    top edge -- it sits ABOVE the title. The reference agrees: the slide
    that carries an eyebrow puts its title at y=95, which is where these
    already were."""
    from deckguard.registry import _load_archetypes

    spec = _load_archetypes().ARCHETYPES["title_content"]
    eyebrow = next(r for r in spec["regions"] if r.get("content") == "eyebrow")
    title = next(r for r in spec["regions"] if r.get("content") == "title")
    assert eyebrow["box"][1] > 10, "the eyebrow must stay on the slide"
    assert eyebrow["box"][1] < title["box"][1]


def test_nothing_on_the_tight_band_lost_its_internal_spacing():
    """The block moves; its parts do not move relative to each other."""
    from deckguard.registry import _load_archetypes

    spec = _load_archetypes().ARCHETYPES["agenda_c_split"]
    items = next(g for g in spec["groups"] if g["content"] == "items")
    pitch = {b[1] - a[1] for a, b in zip(items["origins"], items["origins"][1:])}
    assert len(pitch) == 1, f"row pitch drifted: {pitch}"


# --------------------------------------------------------------------------
# one sand
# --------------------------------------------------------------------------


def test_there_is_exactly_one_sand():
    """`F3EEE6` and `F3EEEA` were both in circulation and both shipped,
    on adjacent slides of the same deck. The real deck measures EA."""
    import pathlib

    assert bm.SAND == "F3EEEA"
    root = pathlib.Path(__file__).resolve().parent.parent / "src" / "deckguard"
    # The quoted literal, not the word: `brandmode` mentions the old
    # value in the comment that records why it changed.
    stale = [p.name for p in root.glob("*.py") if '"F3EEE6"' in p.read_text()]
    assert not stale, stale


# --------------------------------------------------------------------------
# the deviation meter
# --------------------------------------------------------------------------


def test_the_meter_is_read_from_the_designers_file_not_hard_coded():
    from deckguard import meter

    assert meter.METER_FILE.name == "meter.json"
    assert len(meter.stops()) == 4
    assert [s["n"] for s in meter.stops()] == [1, 2, 3, 4]


def test_pools_are_cumulative():
    """Stop 4 contains stop 1. A layout does not become ineligible by
    moving the control to the right."""
    from deckguard import meter

    pools = [meter.pool_for_stop(n) for n in (1, 2, 3, 4)]
    for smaller, larger in zip(pools, pools[1:]):
        assert smaller < larger, "each stop must contain the one before it"


def test_audience_is_inferred_from_the_stop_and_there_is_no_second_switch():
    from deckguard import meter, screens

    assert meter.audience_for_stop(1) == "external"
    assert meter.audience_for_stop(2) == "external"
    assert meter.audience_for_stop(3) == "internal"
    assert meter.audience_for_stop(4) == "internal"

    page = screens.home()
    assert 'name="stop"' in page
    assert 'name="audience"' not in page, "the meter IS the audience control"


def test_stops_one_and_two_offer_no_secondary_colour_field():
    """Tiers 1-2 are external and the external field policy is white
    plus blue. An archetype whose declared field is a secondary colour
    cannot sit there -- which is why `agenda_c_split` is tier 3 despite
    its geometry being a modest deviation."""
    from deckguard import meter
    from deckguard.registry import _load_archetypes

    built = set(_load_archetypes().ARCHETYPES)
    allowed = {"white", "photo", "blue", ""}
    for stop in (1, 2):
        audience = meter.audience_for_stop(stop)
        for name in meter.pool_for_stop(stop) & built:
            field = next((s["field"] for s in bm.slides_in(audience)
                          if s["archetype"] == name), "")
            assert field in allowed, f"stop {stop}: {name} declares {field}"


def test_the_menu_the_planner_sees_is_filtered_to_the_stop():
    """The filter IS the enforcement -- a model cannot choose a layout it
    was never shown, so nothing downstream re-validates the choice."""
    from deckguard import assemble, meter

    for stop in (1, 4):
        menu = assemble._menu(meter.audience_for_stop(stop), stop)
        names = [line.strip().split(" ")[0] for line in menu.split("\n")
                 if line.startswith("  ") and not line.startswith("      ")]
        assert names, stop
        assert set(names) <= meter.pool_for_stop(stop), stop

    wide = assemble._menu("internal", 4)
    narrow = assemble._menu("external", 1)
    assert len(wide.split("\n")) > len(narrow.split("\n"))


def test_the_picker_shows_only_the_current_stop():
    """Not greyed-out tiles for the stops above: showing someone a
    layout they cannot use invites them to argue with the control they
    just set."""
    from deckguard import meter, screens

    for stop in (1, 4):
        tiles = screens._slide_tiles(stop)
        offered = set(re.findall(r'<span>([a-z_0-9]+)</span>', tiles))
        assert offered, stop
        assert offered <= meter.pool_for_stop(stop), stop


def test_a_stop_out_of_range_is_clamped_rather_than_crashing():
    from deckguard import meter

    assert meter.stop(0)["n"] == 1
    assert meter.stop(99)["n"] == 4
    assert meter.audience_for_stop(0) == "external"


def test_the_summary_says_what_the_stop_costs():
    from deckguard import meter

    assert "customer-safe" in meter.summary(1)
    assert "internal only" in meter.summary(4)
    # stop 1 is short three layouts until the unbuilt five are drawn
    assert "not built yet" in meter.summary(1)


def test_the_page_posts_a_stop_and_the_app_reads_it():
    from fastapi.testclient import TestClient

    from deckguard.web import app

    client = TestClient(app)
    r = client.post("/generate", data={"stop": "1", "pick": ["external:7"]})
    assert r.status_code == 200
    # stop 1 is external, so an external slide was accepted
    assert "title_content" in r.text


def test_the_page_stopped_explaining_itself():
    """Four paragraphs of hint copy nobody reads twice. The meter says
    its own consequence in one line instead."""
    import re as _re

    from deckguard import screens

    page = screens.home()
    prose = " ".join(_re.sub(r"<[^>]+>", " ", b) for b in
                     _re.findall(r'<p class="(?:hint|lede)"[^>]*>(.*?)</p>', page, _re.S))
    assert len(_re.findall(r"[A-Za-z][A-Za-z',.-]*", prose)) < 40


def test_the_meter_ships_inside_the_package():
    """It lived in `docs/` first, which is outside the package -- a pip
    install would not have shipped it and the meter would have come up
    with no stops on the deploy while working perfectly here."""
    from deckguard import meter

    assert "site-packages" in str(meter.METER_FILE) or "src/deckguard" in str(meter.METER_FILE)
    assert meter.METER_FILE.is_file()
    assert meter.stops(), "no stops means no meter"


def test_the_packaged_meter_matches_the_one_design_sent():
    """Two copies drift. The packaged one is read; the one in `docs/` is
    the handoff as received, and this says when they disagree."""
    import json
    import pathlib

    from deckguard import meter

    received = (pathlib.Path(__file__).resolve().parent.parent
                / "docs" / "design-handoff" / "meter.json")
    if not received.is_file():
        return
    assert json.loads(received.read_text()) == json.loads(
        meter.METER_FILE.read_text()), "copy docs/design-handoff/meter.json into assets"
