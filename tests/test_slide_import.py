"""Tests for cross-presentation slide import (deckguard migrate).

Deliberately uses the real bundled template asset rather than a
synthetic fixture -- this feature is fundamentally about raw OOXML
package surgery (copying a master/layouts/slides/media across files and
remapping relationship IDs), which a mocked-out fixture wouldn't
meaningfully exercise. The target deck side still uses the existing
in-memory `new_deck()` helper, so nothing here depends on any file
outside this repository.
"""

import zipfile

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckguard.slide_import import (
    _next_master_or_layout_id,
    _next_sld_id,
    default_template_path,
    import_cover_and_outro_parts,
    replace_intro_outro,
)
from tests.helpers import add_slide, new_deck, title_run

TEMPLATE_PATH = default_template_path()
pytestmark = pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="bundled template asset not present")


def test_next_master_or_layout_id_starts_at_schema_floor_on_empty_target():
    # ST_SlideMasterId/ST_SlideLayoutId minInclusive per pml.xsd.
    target = {"ppt/presentation.xml": b"<p:presentation></p:presentation>"}
    assert _next_master_or_layout_id(target) == 2147483648


def test_next_master_or_layout_id_respects_existing_slide_master_id():
    target = {
        "ppt/presentation.xml": (
            b'<p:presentation><p:sldMasterIdLst>'
            b'<p:sldMasterId id="2147483700" r:id="rId1"/>'
            b"</p:sldMasterIdLst></p:presentation>"
        )
    }
    assert _next_master_or_layout_id(target) == 2147483701


def test_next_master_or_layout_id_respects_existing_slide_layout_id_in_any_master():
    """Regression test for the real corruption bug: sldLayoutId and
    sldMasterId share one ID namespace, so a layout ID buried in a
    *second* slideMaster's XML must still be found and avoided -- this
    is exactly the collision that made a real deck's output invalid."""
    target = {
        "ppt/presentation.xml": b"<p:presentation></p:presentation>",
        "ppt/slideMasters/slideMaster1.xml": (
            b'<p:sldMaster><p:sldLayoutIdLst>'
            b'<p:sldLayoutId id="2147483730" r:id="rId1"/>'
            b"</p:sldLayoutIdLst></p:sldMaster>"
        ),
        "ppt/slideMasters/slideMaster2.xml": (
            b'<p:sldMaster><p:sldLayoutIdLst>'
            b'<p:sldLayoutId id="2147483700" r:id="rId1"/>'
            b"</p:sldLayoutIdLst></p:sldMaster>"
        ),
    }
    # Must be strictly greater than every id found across *both* masters,
    # not just the one a naive implementation might check first.
    assert _next_master_or_layout_id(target) == 2147483731


def test_next_sld_id_starts_at_schema_floor_on_empty_target():
    # ST_SlideId minInclusive per pml.xsd.
    target = {"ppt/presentation.xml": b"<p:presentation></p:presentation>"}
    assert _next_sld_id(target) == 256


def test_next_sld_id_respects_existing_slide_ids_and_stays_below_master_range():
    target = {
        "ppt/presentation.xml": (
            b'<p:presentation><p:sldIdLst>'
            b'<p:sldId id="256" r:id="rId2"/>'
            b'<p:sldId id="300" r:id="rId3"/>'
            b"</p:sldIdLst></p:presentation>"
        )
    }
    next_id = _next_sld_id(target)
    assert next_id == 301
    # ST_SlideId maxExclusive: must never stray into the master/layout range.
    assert next_id < 2147483648


def _no_duplicate_zip_entries(path) -> bool:
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
    return len(names) == len(set(names))


def _no_malformed_xml(path) -> list:
    bad = []
    with zipfile.ZipFile(path) as z:
        for name in z.namelist():
            if name.endswith((".xml", ".rels")):
                try:
                    etree.fromstring(z.read(name))
                except Exception as exc:  # noqa: BLE001
                    bad.append((name, str(exc)))
    return bad


def test_import_cover_and_outro_parts_produces_a_valid_package(tmp_path):
    target = new_deck()
    add_slide(target, layout_idx=1)
    target_path = tmp_path / "target.pptx"
    target.save(str(target_path))

    out_path = tmp_path / "out.pptx"
    cover_part, outro_part = import_cover_and_outro_parts(
        str(target_path), str(TEMPLATE_PATH), str(out_path), cover_slide_num=3, outro_slide_num=11
    )

    assert _no_duplicate_zip_entries(out_path)
    assert _no_malformed_xml(out_path) == []

    prs = Presentation(str(out_path))
    assert len(prs.slides) == 3  # original 1 + imported cover + imported outro
    layout_names = {s.slide_layout.name for s in prs.slides}
    assert "Cover B" in layout_names
    assert "Outro" in layout_names


def test_imported_layout_images_are_readable(tmp_path):
    target = new_deck()
    add_slide(target, layout_idx=1)
    target_path = tmp_path / "target.pptx"
    target.save(str(target_path))

    out_path = tmp_path / "out.pptx"
    import_cover_and_outro_parts(str(target_path), str(TEMPLATE_PATH), str(out_path), 3, 11)

    prs = Presentation(str(out_path))
    found_picture = False
    for slide in prs.slides:
        for shape in slide.slide_layout.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                found_picture = True
                img = shape.image  # raises if the relationship/media chain is broken
                assert len(img.blob) > 0
    assert found_picture


def test_replace_intro_outro_replaces_non_standard_cover_and_outro(tmp_path):
    target = new_deck()
    slide = add_slide(target, layout_idx=1)  # stock "Title and Content" -- not a KONE Cover layout
    title_run(slide).text = "My Deck Title"
    target_path = tmp_path / "target.pptx"
    target.save(str(target_path))

    out_path = tmp_path / "out.pptx"
    result = replace_intro_outro(str(target_path), str(out_path), template_path=str(TEMPLATE_PATH))

    assert result["cover_replaced"] is True
    assert result["outro_replaced"] is True
    assert result["cover_title"] == "My Deck Title"

    assert _no_duplicate_zip_entries(out_path)
    assert _no_malformed_xml(out_path) == []

    prs = Presentation(str(out_path))
    assert len(prs.slides) == 1  # net unchanged: 1 slide in, 1 slide out (cover==outro==only slide, both replaced)
    assert prs.slides[0].slide_layout.name in ("Cover B", "Outro")


def test_replace_intro_outro_carries_title_into_new_cover_placeholder(tmp_path):
    target = new_deck()
    slide = add_slide(target, layout_idx=1)
    title_run(slide).text = "Carried-over title"
    add_slide(target, layout_idx=1)  # second slide so cover and outro are distinct
    target_path = tmp_path / "target.pptx"
    target.save(str(target_path))

    out_path = tmp_path / "out.pptx"
    replace_intro_outro(str(target_path), str(out_path), template_path=str(TEMPLATE_PATH))

    prs = Presentation(str(out_path))
    assert prs.slides[0].slide_layout.name == "Cover B"
    assert prs.slides[0].shapes.title.text_frame.text == "Carried-over title"
    assert prs.slides[-1].slide_layout.name == "Outro"


def test_replace_intro_outro_no_op_when_already_on_template(tmp_path):
    target = new_deck()
    slide = add_slide(target, layout_idx=1)
    title_run(slide).text = "Title"
    add_slide(target, layout_idx=1)
    target_path = tmp_path / "target.pptx"
    target.save(str(target_path))

    once_path = tmp_path / "once.pptx"
    replace_intro_outro(str(target_path), str(once_path), template_path=str(TEMPLATE_PATH))

    twice_path = tmp_path / "twice.pptx"
    result = replace_intro_outro(str(once_path), str(twice_path), template_path=str(TEMPLATE_PATH))

    assert result == {"cover_replaced": False, "outro_replaced": False}
    prs = Presentation(str(twice_path))
    assert len(prs.slides) == 2  # unchanged from the first migration


def test_replace_intro_outro_survives_full_deckguard_fix_afterward(tmp_path):
    """Regression test for a real bug: deleting the old cover/outro slide
    left a dangling relationship in presentation.xml.rels (sldId removed,
    but the relationship itself wasn't), which caused python-pptx to
    instantiate a second, uncached Part object for that partname the
    next time anything elsewhere in the deck lazily touched an rPr --
    producing a genuine 'Duplicate name' zip write, not just clutter.
    Reproduced via a full audit/fix pass, since that's what actually
    touches every run in the deck and is what surfaced this originally."""
    import warnings

    from deckguard.config import load_config, default_config_path
    from deckguard.fixer import fix_deck

    target = new_deck()
    slide = add_slide(target, layout_idx=1)
    title_run(slide).text = "Title"
    add_slide(target, layout_idx=1)
    target_path = tmp_path / "target.pptx"
    target.save(str(target_path))

    migrated_path = tmp_path / "migrated.pptx"
    replace_intro_outro(str(target_path), str(migrated_path), template_path=str(TEMPLATE_PATH))

    config = load_config(default_config_path())
    prs = Presentation(str(migrated_path))
    fixed_path = tmp_path / "fixed.pptx"

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        fix_deck(prs, config, source_path=str(migrated_path), output_path=str(fixed_path), dry_run=False)
        dupes = [str(w.message) for w in caught if "Duplicate name" in str(w.message)]

    assert dupes == []
    assert _no_duplicate_zip_entries(fixed_path)
