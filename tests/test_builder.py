"""The manual slide builder: pick from a set, fill it, build.

The flow the tool leads with -- a brief is a way to pre-fill these same
slides, not a separate path.
"""

import re

import pytest
from fastapi.testclient import TestClient



@pytest.fixture
def client():
    """Built from a freshly imported module.

    `test_web` reloads `deckguard.web` to pick up env vars, which
    rebinds the module object and its STORAGE_ROOT. A client holding
    the app imported at collection time then writes decks one place and
    reads them from another, so downloads 404 depending on test order.
    """
    import importlib

    from deckguard import web as web_mod

    return TestClient(importlib.reload(web_mod).app)


def _compose(client, picks):
    r = client.post("/build/compose", data={"pick": picks})
    assert r.status_code == 200, r.text[:300]
    return re.search(r"/build/generate/([a-f0-9]+)", r.text).group(1), r.text


def test_the_picker_offers_both_sets_and_every_slide_in_them(client):
    """All 50 appear; the ones whose archetype is not drawn yet appear
    as unpickable tiles rather than being hidden, so the set still reads
    as 25 and the gap is visible."""
    from deckguard import brandmode as B

    r = client.get("/build")
    assert r.status_code == 200
    assert "Internal 25" in r.text and "External 25" in r.text
    for audience in B.set_names():
        for slide in B.slides_in(audience):
            assert slide["archetype"] in r.text


def test_composing_offers_a_slot_for_every_key_the_renderer_reads(client):
    _token, html = _compose(client, ["external:9"])
    # two_content reads title + items[2]{label, bullets}
    assert 'name="v:external-9:title"' in html
    assert 'name="v:external-9:items"' in html
    # and says how to fill a repeating group rather than leaving it guessable
    assert "One per line" in html


def test_a_picture_slot_is_offered_but_optional(client):
    _token, html = _compose(client, ["external:12"])   # text_picture_a
    assert 'name="p:external-12"' in html
    assert "automatically" in html


def test_order_reorders_duplicate_duplicates_and_drop_drops(client):
    token, _ = _compose(client, ["external:6", "external:9", "external:23"])
    r = client.post(f"/build/generate/{token}", data={
        "v:external-6:title": "The turn",
        "v:external-9:title": "Two columns",
        "v:external-23:title": "Phasing",
        "o:external-6": "30", "o:external-9": "10", "o:external-23": "20",
        "x:external-23": "on",
    })
    assert r.status_code == 200
    deck = client.get(f"/download/{token}/deck.pptx")
    assert deck.status_code == 200

    from io import BytesIO

    from pptx import Presentation

    slides = Presentation(BytesIO(deck.content)).slides
    text = [" ".join(sh.text_frame.text for sh in s.shapes
                     if getattr(sh, "has_text_frame", False)) for s in slides]
    body = [t for t in text if any(k in t for k in ("turn", "Two columns", "Phasing"))]
    assert "Two columns" in body[0], body      # order 10 first
    assert "Phasing" in body[1] and "Phasing" in body[2], body  # duplicated
    assert "The turn" in body[3], body         # order 30 last


def test_dropping_every_slide_says_so_rather_than_building_nothing(client):
    token, _ = _compose(client, ["external:6"])
    r = client.post(f"/build/generate/{token}",
                    data={"d:external-6": "on"}, follow_redirects=False)
    assert r.status_code == 400
    assert "nothing to build" in r.text


def test_an_empty_slot_is_left_out_rather_than_filled(client):
    """A slide the author did not finish comes out short. Placeholder
    text that reaches a customer is worse than a gap the author sees."""
    from deckguard.web import _spec_from_form

    plan = {"audience": "external", "slides": [
        {"id": "x", "archetype": "two_content",
         "slots": [("title", "text"), ("items", "list of up to 2 × {label, bullets}")]}]}
    spec = _spec_from_form(plan, {"v:x:title": "Kept", "v:x:items": "   "})
    assert spec["slides"][0] == {"archetype": "two_content", "title": "Kept"}


def test_a_repeating_group_is_one_item_per_line_split_on_pipes(client):
    from deckguard.web import _parse_slot

    parsed = _parse_slot("Q1 | Survey\nQ2 | Build", "list of up to 4 × {period, text}")
    assert parsed == [{"period": "Q1", "text": "Survey"},
                      {"period": "Q2", "text": "Build"}]
    # a short line leaves the trailing fields empty rather than shifting them
    assert _parse_slot("Q1", "list of up to 4 × {period, text}") == [
        {"period": "Q1", "text": ""}]
    # a plain string list stays a list of strings
    assert _parse_slot("one\ntwo", "list of strings") == ["one", "two"]
    # a scalar slot is untouched
    assert _parse_slot("just text", "text") == "just text"


def test_a_stale_build_token_says_so(client):
    r = client.post("/build/generate/" + "0" * 32, data={})
    assert r.status_code == 404
    assert "expired" in r.text


def test_a_build_with_no_slides_picked_asks_for_one(client):
    r = client.post("/build/compose", data={})
    assert r.status_code == 400
    assert "Pick at least one slide" in r.text


# --------------------------------------------------------------------------
# previews
# --------------------------------------------------------------------------


def test_every_pickable_slide_shows_a_wireframe_preview(client):
    """The preview is drawn from the same region and group data the pptx
    renderer consumes, so what you pick is the shape you get."""
    r = client.get("/build")
    pickable = r.text.count('name="pick"')
    assert r.text.count("data-dg-frame") == pickable
    assert pickable >= 40


def test_slides_whose_archetype_is_not_built_are_shown_but_not_pickable(client):
    """Five of the 50 are named by the sets and not yet drawn. Picking
    one built a blank slide with no warning, which is worse than not
    offering it."""
    from deckguard.registry import _load_archetypes

    from deckguard import brandmode as B

    built = set(_load_archetypes().ARCHETYPES)
    unbuilt = {a for a in B.canonical_archetypes() if a not in built}
    assert unbuilt, "if every archetype is built, this guard can go"

    r = client.get("/build")
    assert "NOT BUILT YET" in r.text.upper()
    for name in unbuilt:
        # named on the page, but never inside a pickable label
        assert name in r.text
    # and the count of checkboxes is short by exactly the unbuilt slides
    total = sum(len(B.slides_in(a)) for a in B.set_names())
    unbuilt_slides = sum(1 for a in B.set_names() for s in B.slides_in(a)
                         if s["archetype"] in unbuilt)
    assert r.text.count('name="pick"') == total - unbuilt_slides


def test_a_forged_pick_of_an_unbuilt_archetype_is_refused(client):
    """The checkbox is absent, so this only happens to a crafted post --
    but a blank slide reaching a customer deck is worth the guard."""
    from deckguard.registry import _load_archetypes

    from deckguard import brandmode as B

    built = set(_load_archetypes().ARCHETYPES)
    victim = next(s for s in B.slides_in("external") if s["archetype"] not in built)
    r = client.post("/build/compose", data={"pick": [f"external:{victim['n']}"]})
    assert r.status_code == 400
    assert "not built yet" in r.text.lower()


def test_a_list_value_never_reaches_a_preview_as_python_source():
    """`['Standardise intake', ...]` was rendering verbatim on every
    archetype whose bullet slot is not role `bullets`."""
    from deckguard.preview import archetype_preview_html, sample_content

    for name in ("title_content", "title_subtitle_content_a", "two_content"):
        html = archetype_preview_html(name, sample_content(name))
        assert "['" not in html and "', '" not in html, name


def test_sample_content_is_preview_only_and_never_invents_deck_copy():
    """The builder leaves an unfilled slot empty; these words exist so a
    thumbnail is recognisable, and must not leak into a build."""
    from deckguard.preview import sample_content

    from deckguard.web import _spec_from_form

    assert sample_content("two_content")            # the preview has words
    plan = {"audience": "external", "slides": [
        {"id": "x", "archetype": "two_content", "slots": [("title", "text")]}]}
    built = _spec_from_form(plan, {})
    assert built["slides"][0] == {"archetype": "two_content"}
