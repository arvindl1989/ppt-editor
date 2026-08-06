"""Tests for the unified Transform pipeline (transform.py): plan ->
human review -> execute -> audit. No real network calls -- the AI
suggestion step injects the same fake Anthropic-shaped client the rest
of the suite uses, and everything else is deterministic."""

import json

import pytest

from deckguard.skill_bridge import _skill_dir
from deckguard.slide_import import default_template_path
from deckguard.transform import (
    audit_transform_result,
    execute_transform,
    execute_transform_from_brief,
    plan_transform,
    plan_transform_from_brief,
    reference_similarity,
)
from tests.helpers import add_slide, body_run, make_solid_png, new_deck, title_run
from tests.test_redesign import _FakeClient, _FakeResponse, _kone_slide, _kone_spec_json

TEMPLATE_PATH = default_template_path()
pytestmark = pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="bundled template asset not present")

_skill_installed = (_skill_dir() / "kone_deck_creator.py").is_file()
needs_skill = pytest.mark.skipif(not _skill_installed, reason="kone-deck-generator skill not installed")


def _three_slide_deck(tmp_path):
    prs = new_deck()
    c = add_slide(prs)
    title_run(c).text = "Annual Review"
    m = add_slide(prs)
    title_run(m).text = "Resolution rate"
    body_run(m).text = "91.2% resolved"
    e = add_slide(prs)
    title_run(e).text = "Thank you"
    path = tmp_path / "src.pptx"
    prs.save(str(path))
    return path


def _hero_stat_override(index):
    return json.dumps({"overrides": [
        {"outline_index": index, "archetype": "hero_stat", "eyebrow": "Resolution rate", "value": "91.2%",
         "caption": "of requests cleared", "support": "91.2% resolved"},
    ]})


@needs_skill
def test_plan_transform_merges_deterministic_proposals_with_ai_suggestions(tmp_path):
    src = _three_slide_deck(tmp_path)
    plan = plan_transform(str(src), client=_FakeClient(_FakeResponse(_hero_stat_override(2))))

    by_index = {s.index: s for s in plan.slides}
    assert by_index[1].default_action == "rebuild" and by_index[1].layout_name == "Cover B"
    assert by_index[2].default_action == "archetype" and by_index[2].archetype["archetype"] == "hero_stat"
    assert by_index[3].default_action == "rebuild" and by_index[3].layout_name == "Outro"
    assert plan.ai_suggestions_ran


def test_plan_transform_degrades_to_deterministic_with_no_api_key(tmp_path, monkeypatch):
    """No key, no client -> the plan still succeeds AND still offers
    archetypes, matched structurally rather than by a model. The AI step
    must never be load-bearing -- and "not load-bearing" has to mean
    degrading to a deterministic answer, not to no answer at all."""
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    src = _three_slide_deck(tmp_path)

    plan = plan_transform(str(src))

    assert plan.ai_suggestions_ran is False
    assert {s.default_action for s in plan.slides} == {"rebuild"}
    body = [s for s in plan.slides if s.index == 2][0]
    assert body.archetype is not None
    assert body.archetype_source == "structural"


@needs_skill
def test_execute_transform_honors_per_slide_choices(tmp_path):
    """The human decision point: keep one slide untouched, accept the
    archetype for another, rebuild the rest -- exactly and only those
    choices execute."""
    from pptx import Presentation

    src = _three_slide_deck(tmp_path)
    plan = plan_transform(str(src), client=_FakeClient(_FakeResponse(_hero_stat_override(2))))
    out = tmp_path / "out.pptx"

    outcome = execute_transform(str(src), str(out), plan, actions={3: "keep"})

    assert outcome.rebuilt == [1]
    assert outcome.archetype_swapped == [2]
    assert outcome.kept == [3]
    assert outcome.layouts_used == {1: "Cover B", 2: "hero_stat"}

    prs = Presentation(str(out))
    assert len(prs.slides) == 3
    assert any(
        "91.2%" in s.text_frame.text for s in prs.slides[1].shapes if getattr(s, "has_text_frame", False)
    )
    assert prs.slides[2].shapes.title.text_frame.text == "Thank you"  # structure kept (brand patches still apply deck-wide)


@needs_skill
def test_execute_transform_downgrades_archetype_choice_with_nothing_to_render(tmp_path):
    """A stray "archetype" action for a slide that has NO archetype at
    all in the plan degrades to a rebuild -- never invents content,
    never fails. (A slide the structural matcher did find an archetype
    for is a different case: that one renders, see below.)"""
    src = _three_slide_deck(tmp_path)
    plan = plan_transform(str(src), suggest_archetypes=False)
    out = tmp_path / "out.pptx"

    outcome = execute_transform(str(src), str(out), plan, actions={2: "archetype"})

    assert 2 in outcome.rebuilt
    assert outcome.archetype_swapped == []


@needs_skill
def test_a_structurally_matched_archetype_actually_renders(tmp_path):
    """The point of matching without a model: choosing that offered
    archetype has to execute, not quietly fall back to a rebuild."""
    src = _three_slide_deck(tmp_path)
    plan = plan_transform(str(src), suggest_archetypes=True, client=None, api_key=None)
    body = [s for s in plan.slides if s.index == 2][0]
    assert body.archetype_source == "structural"

    out = tmp_path / "out.pptx"
    outcome = execute_transform(str(src), str(out), plan, actions={2: "archetype"})

    assert outcome.archetype_swapped == [2]


@needs_skill
def test_audit_transform_result_excludes_archetype_slides_from_the_report(tmp_path):
    """The trap this closes: rules_engine flagging the archetype
    engine's own styling as a text violation. The transform audit
    excludes archetype slides from the report and says how many findings
    it suppressed, instead of reporting correct slides as defects.

    The muted caption grey was the original trigger and no longer is --
    `layouts._correct_grey_ink` takes it off the brand's type roles at
    install, because `#727272` was never in the palette. The exclusion
    still has to hold for anything else the engine draws deliberately."""
    src = _three_slide_deck(tmp_path)
    plan = plan_transform(str(src), client=_FakeClient(_FakeResponse(_hero_stat_override(2))))
    out = tmp_path / "out.pptx"
    outcome = execute_transform(str(src), str(out), plan)

    audit = audit_transform_result(str(out), archetype_indices=set(outcome.archetype_swapped))

    assert "suppressed_archetype_findings" in audit
    assert all(v.slide_index != 2 for v in audit["violations"])

    # and the count is the number actually withheld, not a constant
    unfiltered = audit_transform_result(str(out), archetype_indices=set())
    withheld = [v for v in unfiltered["violations"] if v.slide_index == 2]
    assert audit["suppressed_archetype_findings"] == len(withheld)


def test_reference_similarity_reports_layout_matches_and_divergence(tmp_path):
    src = _three_slide_deck(tmp_path)

    sim = reference_similarity(str(src), str(src))

    assert sim["slides_compared"] == 3
    assert sim["layout_matches"] == 3
    assert sim["colors_not_in_reference"] == []
    assert sim["fonts_not_in_reference"] == []


@needs_skill
def test_plan_and_execute_from_brief_builds_only_approved_slides(tmp_path):
    spec_json = _kone_spec_json(
        "Planned deck",
        _kone_slide("agenda_contents", title="Agenda", items=[{"number": "01", "item": "One"}]),
        _kone_slide("hero_stat", eyebrow="KPI", value="91%", caption="resolved", support="s"),
    )
    plan = plan_transform_from_brief("A brief.", client=_FakeClient(_FakeResponse(spec_json)))

    assert plan.deck_title == "Planned deck"
    assert [s.archetype["archetype"] for s in plan.slides] == ["agenda_contents", "hero_stat"]

    from pptx import Presentation

    out = tmp_path / "out.pptx"
    outcome = execute_transform_from_brief(str(out), plan, approved_indices={2})

    assert outcome.kept == [1]
    prs = Presentation(str(out))
    assert len(prs.slides) == 3  # retained cover + 1 approved body + retained outro
    assert any(
        "91%" in s.text_frame.text for s in prs.slides[1].shapes if getattr(s, "has_text_frame", False)
    )


def test_execute_transform_all_keep_still_applies_brand_patches(tmp_path):
    """"keep" spares a slide's structure, not the brand rules -- an
    all-keep transform must still run the deck-wide fix pass (the
    apply_rebrand fast path it sits on returns a plain copy with no fix
    at all, which would make Transform's baseline a silent no-op)."""
    from pptx import Presentation

    from tests.helpers import set_run

    prs = new_deck()
    s = add_slide(prs)
    set_run(title_run(s), text="Title", font="Calibri", color_hex="005EB8")  # legacy blue + off-brand font
    src = tmp_path / "src.pptx"
    prs.save(str(src))

    plan = plan_transform(str(src), suggest_archetypes=False)
    out = tmp_path / "out.pptx"
    outcome = execute_transform(str(src), str(out), plan, actions={1: "keep"})

    assert outcome.kept == [1]
    run = Presentation(str(out)).slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == "141414"  # heading forced to brand black despite "keep" (heading_always_dark)


def _overloaded_deck(tmp_path, n_boxes=8):
    """A slide with more text boxes than ANY org-template layout can
    hold (MAX_TEXT_BLOCKS is 3) -- the single most common shape in a
    real deck, and the case that used to be written off as "can't
    happen" with no alternative offered."""
    from pptx.util import Inches

    prs = new_deck()
    s = add_slide(prs)
    title_run(s).text = "Dense slide"
    for i in range(n_boxes):
        box = s.shapes.add_textbox(Inches(1), Inches(0.4 * i + 1), Inches(4), Inches(0.35))
        box.text_frame.text = f"Point {i}: real content"
    path = tmp_path / "dense.pptx"
    prs.save(str(path))
    return path


@needs_skill
def test_slide_too_big_for_any_org_layout_is_still_offered_an_archetype(tmp_path):
    """Regression for the real complaint: on a production deck, 129 of
    168 slides were refused with "more body text blocks than any
    template layout can hold" and shown with NO options. Those slides
    are perfectly readable -- only the ORG TEMPLATE can't hold them --
    so an archetype must still be offered."""
    src = _overloaded_deck(tmp_path)

    plan = plan_transform(str(src), client=_FakeClient(_FakeResponse(_hero_stat_override(1))))

    slide = plan.slides[0]
    assert slide.archetype is not None, "an archetype must still be offered for it"
    assert slide.archetype["archetype"] == "hero_stat"
    # And it becomes the DEFAULT -- the archetype is the only route that
    # can hold this slide, so it shouldn't need hunting for.
    assert slide.default_action == "archetype"

    # Without a suggestion it stays keep-only and says why -- honest, but
    # that was the ONLY outcome before this fix.
    plain = plan_transform(str(src), suggest_archetypes=False)
    assert plain.slides[0].default_action == "keep"
    assert "text blocks" in (plain.slides[0].reason or "")


@needs_skill
def test_archetype_is_honored_on_a_slide_no_org_layout_can_hold(tmp_path):
    """The other half: choosing that offered archetype must actually
    execute, not silently degrade back to keep."""
    from pptx import Presentation

    src = _overloaded_deck(tmp_path)
    plan = plan_transform(str(src), client=_FakeClient(_FakeResponse(_hero_stat_override(1))))
    out = tmp_path / "out.pptx"

    outcome = execute_transform(str(src), str(out), plan, actions={1: "archetype"})

    assert outcome.archetype_swapped == [1]
    assert outcome.kept == []
    assert any(
        "91.2%" in s.text_frame.text
        for s in Presentation(str(out)).slides[0].shapes
        if getattr(s, "has_text_frame", False)
    )


def test_reference_deck_actually_drives_styling(tmp_path):
    """Regression for "it's still editing decks based on the back end
    yaml": with a reference uploaded, the reference's OWN answer for a
    shape must win over what brand_rules.yaml would do on its own.

    Decisive setup -- the old panel is legacy KONE blue #005EB8, which
    the bundled yaml remaps to #1450F5. The reference deck says that
    same panel should be sand #F3EEE6. If the output is sand, the
    uploaded reference drove the result; if it's #1450F5, only the
    backend yaml did (the bug being fixed)."""
    from pptx import Presentation

    from tests.helpers import add_rectangle

    def _deck(panel_hex):
        prs = new_deck()
        c = add_slide(prs)
        title_run(c).text = "Cover"
        mid = add_slide(prs)
        title_run(mid).text = "Panel slide"
        add_rectangle(mid, name="Panel", fill_hex=panel_hex, left_in=1, top_in=2, width_in=3, height_in=2)
        e = add_slide(prs)
        title_run(e).text = "Thank you"
        return prs

    src = tmp_path / "src.pptx"
    _deck("005EB8").save(str(src))  # legacy blue -> yaml alone would make this #1450F5
    reference = tmp_path / "ref.pptx"
    _deck("F3EEE6").save(str(reference))  # the reference's own answer: sand

    plan = plan_transform(str(src), reference_path=str(reference), suggest_archetypes=False)
    out = tmp_path / "out.pptx"
    outcome = execute_transform(
        str(src), str(out), plan, actions={2: "keep"}, reference_path=str(reference),
    )

    # The reference measurably influenced the run (either mechanism:
    # learned deck-wide rules, or per-shape transplant).
    assert outcome.learned_colors + outcome.transplanted_shapes > 0

    panel = next(sh for sh in Presentation(str(out)).slides[1].shapes if sh.name == "Panel")
    assert str(panel.fill.fore_color.rgb) == "F3EEE6", "the reference's answer must beat the yaml default"


def _put_picture_on_master(prs, image_path, left, top, width, height, name):
    """python-pptx can't add a picture to a slide master directly, so
    add it to a slide and move the <p:pic> into the master's spTree --
    which is structurally what layout import does when it drags the
    reference's master across."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.util import Emu

    holder = prs.slides[0]
    pic = holder.shapes.add_picture(str(image_path), Emu(left), Emu(top), Emu(width), Emu(height))
    pic.name = name
    image_part = holder.part.related_part(pic._element.blip_rId)
    el = pic._element
    el.getparent().remove(el)
    master = prs.slide_masters[0]
    # re-point the blip at a relationship the MASTER part owns, or the
    # picture's bytes become unreadable from its new home
    el.blipFill.blip.rEmbed = master.part.relate_to(image_part, RT.IMAGE)
    master.shapes._spTree.append(el)


def test_reference_sourced_master_keeps_one_logo_not_two(tmp_path):
    """Regression for a real report: "THE KONE LOGO IS APPEARING TWICE".

    On a transform with a reference deck, the reference's master comes
    across carrying its own logo at 11.76in -- just outside the 11.9in
    scan region -- so the logo fixer doesn't see it and stamps
    deckguard's bundled mark beside it. Both render, overlapping, on
    every carried-over slide. The reference's own mark must win."""
    from pptx import Presentation

    from deckguard.transform import _dedupe_reference_master_logos

    ref_logo = tmp_path / "ref_logo.png"
    bundled = tmp_path / "bundled_logo.png"
    make_solid_png(ref_logo, (20, 60, 245))
    make_solid_png(bundled, (255, 255, 255))

    reference = tmp_path / "ref.pptx"
    ref_prs = new_deck()
    add_slide(ref_prs)
    _put_picture_on_master(ref_prs, ref_logo, 10753104, 182880, 914400, 274320, "Picture 6")
    ref_prs.save(str(reference))

    out = tmp_path / "out.pptx"
    out_prs = new_deck()
    add_slide(out_prs)
    _put_picture_on_master(out_prs, ref_logo, 10753104, 182880, 914400, 274320, "Picture 6")
    _put_picture_on_master(out_prs, bundled, 10984613, 182880, 914400, 274320, "Picture 8")
    out_prs.save(str(out))

    assert _dedupe_reference_master_logos(str(out), str(reference)) == 1

    names = [s.name for s in Presentation(str(out)).slide_masters[0].shapes if s.shape_type.name == "PICTURE"]
    assert names == ["Picture 6"], "the reference's own mark survives, the stamped duplicate goes"


def test_non_overlapping_master_pictures_are_left_alone(tmp_path):
    """The dedupe is targeted at the overlapping-logo case only -- a
    master legitimately carrying two separate images (logo + a corner
    device) must survive untouched."""
    from pptx import Presentation

    from deckguard.transform import _dedupe_reference_master_logos

    ref_logo = tmp_path / "ref_logo.png"
    other = tmp_path / "other.png"
    make_solid_png(ref_logo, (20, 60, 245))
    make_solid_png(other, (255, 225, 65))

    reference = tmp_path / "ref.pptx"
    ref_prs = new_deck()
    add_slide(ref_prs)
    _put_picture_on_master(ref_prs, ref_logo, 10753104, 182880, 914400, 274320, "Logo")
    ref_prs.save(str(reference))

    out = tmp_path / "out.pptx"
    out_prs = new_deck()
    add_slide(out_prs)
    _put_picture_on_master(out_prs, ref_logo, 10753104, 182880, 914400, 274320, "Logo")
    _put_picture_on_master(out_prs, other, 182880, 6000000, 914400, 274320, "Corner device")
    out_prs.save(str(out))

    assert _dedupe_reference_master_logos(str(out), str(reference)) == 0
    assert len([s for s in Presentation(str(out)).slide_masters[0].shapes if s.shape_type.name == "PICTURE"]) == 2


def test_slides_the_reference_redrew_are_reported_not_shipped_silently(tmp_path):
    """A reference deck whose diagram slides were REDRAWN by hand can't
    be matched by restyling -- exact_transplant already detects that and
    flags those slides, but transform.py used to throw the list away, so
    the tool knew the slide didn't match and never said so. It must
    surface as `needs_manual_redraw`."""
    from pptx.util import Inches

    from tests.helpers import add_rectangle

    def _deck(shape_names):
        prs = new_deck()
        c = add_slide(prs)
        title_run(c).text = "Cover"
        mid = add_slide(prs)
        title_run(mid).text = "Process"
        for i, n in enumerate(shape_names):
            add_rectangle(mid, name=n, fill_hex="005EB8", left_in=1 + i, top_in=2, width_in=0.8, height_in=1)
        e = add_slide(prs)
        title_run(e).text = "Thank you"
        return prs

    src = tmp_path / "src.pptx"
    _deck([f"Step {i}" for i in range(6)]).save(str(src))
    reference = tmp_path / "ref.pptx"
    # Same slide, redrawn: different shape identities in different places.
    ref_prs = _deck([f"Node {i}" for i in range(6)])
    for i, sh in enumerate(s for s in ref_prs.slides[1].shapes if s.name.startswith("Node")):
        sh.left, sh.top = Inches(0.5 + i * 1.7), Inches(4)
    ref_prs.save(str(reference))

    plan = plan_transform(str(src), reference_path=str(reference), suggest_archetypes=False)
    outcome = execute_transform(
        str(src), str(tmp_path / "out.pptx"), plan, actions={2: "keep"}, reference_path=str(reference),
    )

    assert 2 in outcome.needs_manual_redraw, "the redrawn slide must be reported to the user"

    # ...and a reference that DOES line up must not cry wolf.
    matching = tmp_path / "matching.pptx"
    _deck([f"Step {i}" for i in range(6)]).save(str(matching))
    quiet = execute_transform(
        str(src), str(tmp_path / "out2.pptx"),
        plan_transform(str(src), reference_path=str(matching), suggest_archetypes=False),
        actions={2: "keep"}, reference_path=str(matching),
    )
    assert quiet.needs_manual_redraw == []


def test_plan_says_suggestions_are_off_when_the_server_has_no_key(tmp_path, monkeypatch):
    """Reported from a real run: a keyless server reported archetype
    suggestions as having RUN, then offered nothing but "keep" on ten of
    twelve slides. The flag has to distinguish "couldn't ask" from
    "asked and got nothing", or the review page can't explain itself."""
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    src = _three_slide_deck(tmp_path)

    plan = plan_transform(str(src))
    assert plan.ai_suggestions_ran is False
    # ...but archetypes are still offered, matched structurally
    assert any(s.archetype_source == "structural" for s in plan.slides)

    with_key = plan_transform(str(src), client=_FakeClient(_FakeResponse(_hero_stat_override(2))))
    assert with_key.ai_suggestions_ran is True


def test_an_empty_slide_is_kept_not_given_an_archetype(tmp_path):
    """Found on a real campaign deck: a "Slogan" slide carrying only its
    date, footer and page number was handed an archetype, which rendered
    an empty slide. A slide with nothing to say gets nothing proposed."""
    from tests.helpers import add_slide, new_deck, title_run

    prs = new_deck()
    c = add_slide(prs)
    title_run(c).text = "Cover"
    blank = add_slide(prs)  # a slide with no text and no images at all
    del blank
    e = add_slide(prs)
    title_run(e).text = "Thank you"
    src = tmp_path / "src.pptx"
    prs.save(str(src))

    plan = plan_transform(str(src))

    empty = [s for s in plan.slides if s.index == 2][0]
    assert empty.archetype is None
    assert empty.default_action == "keep"


def test_rebuilt_slides_keep_their_date_footer_and_page_number(tmp_path):
    """python-pptx does not clone a layout's date, footer and
    slide-number placeholders -- OOXML calls them "latent". So every
    REBUILT slide lost its footer while kept slides held theirs: a real
    13-slide deck came back with page numbers on three of them."""
    from pptx import Presentation

    from deckguard.transform import restore_footer_chrome
    from tests.helpers import add_slide, new_deck, title_run

    prs = new_deck()
    for n in range(3):
        s = add_slide(prs)
        title_run(s).text = f"Slide {n}"
    src = tmp_path / "src.pptx"
    prs.save(str(src))

    plan = plan_transform(str(src), suggest_archetypes=False)
    out = tmp_path / "out.pptx"
    execute_transform(str(src), str(out), plan)

    # idempotent: running it again adds nothing further
    again = restore_footer_chrome(str(out), str(src))
    numbers = []
    for slide in Presentation(str(out)).slides:
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type.name == "SLIDE_NUMBER":
                    numbers.append(shape.text_frame.text.strip())
            except Exception:  # noqa: BLE001
                continue
    assert numbers == ["1", "2", "3"], numbers
    assert again == 0, "already-present chrome must not be duplicated"
