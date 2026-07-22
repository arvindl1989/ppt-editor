import pytest
from pptx.enum.dml import MSO_THEME_COLOR

from deckguard import colors as colors_mod
from tests.helpers import new_deck, set_theme_slot


def test_normalize_hex():
    assert colors_mod.normalize_hex("#1450f5") == "1450F5"
    assert colors_mod.normalize_hex("1450F5") == "1450F5"
    with pytest.raises(ValueError):
        colors_mod.normalize_hex("#ZZZZZZ")
    with pytest.raises(ValueError):
        colors_mod.normalize_hex("#12345")


def test_rgb_distance_and_apply_brightness():
    assert colors_mod.rgb_distance((0, 0, 0), (0, 0, 0)) == 0
    assert colors_mod.rgb_distance((0, 0, 0), (255, 0, 0)) == 255

    lightened = colors_mod.apply_brightness((0, 94, 184), 0.5)
    assert all(lightened[i] > (0, 94, 184)[i] for i in (0, 1, 2) if (0, 94, 184)[i] < 255)

    darkened = colors_mod.apply_brightness((0, 94, 184), -0.5)
    assert darkened == (0, 47, 92)


def test_remap_theme_colors_rewrites_srgb_slot():
    prs = new_deck()
    set_theme_slot(prs, "accent1", "005EB8")

    changes = colors_mod.remap_theme_colors(prs, {"005EB8": "1450F5"})

    assert len(changes) == 1
    assert changes[0]["slot"] == "accent1"
    master = prs.slide_masters[0]
    scheme = colors_mod.get_theme_scheme(master)
    assert scheme.slots["accent1"] == "1450F5"


def test_remap_theme_colors_replaces_syscolor_reference():
    prs = new_deck()
    # dk1 defaults to a <a:sysClr val="windowText" lastClr="000000"/>
    changes = colors_mod.remap_theme_colors(prs, {"000000": "141414"})
    assert changes[0]["slot"] == "dk1"
    scheme = colors_mod.get_theme_scheme(prs.slide_masters[0])
    assert scheme.slots["dk1"] == "141414"


def test_effective_rgb_resolves_theme_tint():
    prs = new_deck()
    set_theme_slot(prs, "accent1", "1450F5")
    scheme = colors_mod.get_theme_scheme(prs.slide_masters[0])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    run = slide.shapes.title.text_frame.paragraphs[0].add_run()
    run.text = "x"
    run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    run.font.color.brightness = 0.6  # a 60%-tint of accent1, like the brand's tint scale

    eff = colors_mod.effective_rgb(run.font.color, scheme)
    expected = colors_mod.apply_brightness(colors_mod.hex_to_rgb("1450F5"), 0.6)
    assert eff == expected


def test_effective_rgb_none_for_unset_color():
    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    run = slide.shapes.title.text_frame.paragraphs[0].add_run()
    run.text = "x"
    assert colors_mod.effective_rgb(run.font.color, None) is None


def test_relative_luminance_black_and_white_are_the_extremes():
    assert colors_mod.relative_luminance((0, 0, 0)) == 0.0
    assert colors_mod.relative_luminance((255, 255, 255)) == 1.0


def test_contrast_ratio_black_vs_white_is_maximal():
    assert colors_mod.contrast_ratio((0, 0, 0), (255, 255, 255)) == pytest.approx(21.0, abs=0.01)
    assert colors_mod.contrast_ratio((10, 10, 10), (10, 10, 10)) == pytest.approx(1.0, abs=0.01)


def test_expected_contrast_text_hex_matches_general_branding_examples():
    """Regression test pinned to General_Branding.docx's explicit
    examples: white text on KONE Blue, black text on a pale secondary color."""
    assert colors_mod.expected_contrast_text_hex("1450F5", "141414", "FFFFFF") == "FFFFFF"
    assert colors_mod.expected_contrast_text_hex("F3EEE6", "141414", "FFFFFF") == "141414"
    assert colors_mod.expected_contrast_text_hex("FFCDD7", "141414", "FFFFFF") == "141414"
    assert colors_mod.expected_contrast_text_hex("141414", "141414", "FFFFFF") == "FFFFFF"
    assert colors_mod.expected_contrast_text_hex("FFFFFF", "141414", "FFFFFF") == "141414"


def test_is_greyish_true_for_neutral_tones():
    assert colors_mod.is_greyish((217, 217, 217)) is True  # #D9D9D9, a real panel color found in production
    assert colors_mod.is_greyish((0, 0, 0)) is True
    assert colors_mod.is_greyish((255, 255, 255)) is True
    assert colors_mod.is_greyish((59, 66, 69)) is True  # #3B4245, close-to-neutral dark charcoal


def test_is_greyish_false_for_saturated_brand_colors():
    assert colors_mod.is_greyish((20, 80, 245)) is False  # KONE Blue
    assert colors_mod.is_greyish((255, 95, 40)) is False  # KONE Red
    assert colors_mod.is_greyish((0, 113, 185)) is False  # #0071B9
