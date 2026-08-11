"""Tests for the schematic preview renderers (preview.py). These
assert on structure and content (geometry percentages, colors, text,
fail-soft behavior) -- visual quality was verified by rendering the
output in a real browser."""

import pytest

from deckguard.inventory import build_inventory
from deckguard.preview import (
    archetype_preview_html,
    org_layout_preview_html,
    slide_preview_html,
)
from deckguard.legacy.skill_bridge import _skill_dir
from deckguard.legacy.slide_import import default_template_path
from tests.helpers import add_rectangle, add_slide, body_run, new_deck, set_run, title_run

needs_skill = pytest.mark.skipif(
    not (_skill_dir() / "kone_deck_creator.py").is_file(), reason="kone-deck-generator skill not installed"
)


def _deck_with_panel():
    prs = new_deck()
    s = add_slide(prs)
    title_run(s).text = "Resolution rate held at 91%"
    body_run(s).text = "Requests doubled"
    box = add_rectangle(s, name="Panel", fill_hex="1450F5", left_in=7, top_in=4.5, width_in=3, height_in=1.5)
    set_run(box.text_frame.paragraphs[0].add_run(), text="Panel text", color_hex="FFFFFF")
    return prs


def test_slide_preview_shows_real_text_and_fill_colors():
    prs = _deck_with_panel()
    inv = build_inventory(prs)

    html = slide_preview_html(inv.slides[0], prs.slide_width / 914400, prs.slide_height / 914400)

    assert "Resolution rate held at 91%" in html
    assert "background:#1450F5" in html  # the panel's real fill
    assert "aspect-ratio:1280/720" in html
    assert "<script" not in html


def test_slide_preview_escapes_hostile_slide_text():
    prs = new_deck()
    s = add_slide(prs)
    title_run(s).text = '<img src=x onerror=alert(1)>'
    inv = build_inventory(prs)

    html = slide_preview_html(inv.slides[0])

    assert "<img" not in html
    assert "&lt;img" in html


def test_slide_preview_is_fail_soft():
    html = slide_preview_html(object())  # nothing SlideRecord-shaped at all
    assert "aspect-ratio:1280/720" in html  # the fallback card, not an exception


@needs_skill
def test_archetype_preview_draws_from_the_skills_own_region_data():
    html = archetype_preview_html("hero_stat", {
        "eyebrow": "Resolution rate", "value": "91.2%", "caption": "cleared", "support": "674 of 739",
    })

    assert "91.2%" in html
    assert "text-transform:uppercase" in html  # eyebrow role styled per kone_engine.ROLE_STYLE


@needs_skill
def test_archetype_preview_shows_picture_slots_as_placeholders():
    html = archetype_preview_html("three_picture_cards", {
        "title": "Three offerings",
        "cards": [{"heading": "Care", "bullets": ["24/7"]}],
    })

    assert "PHOTO" in html
    assert "Care" in html


def test_archetype_preview_is_fail_soft_for_unknown_archetype():
    html = archetype_preview_html("not_a_real_archetype", {})
    assert "not_a_real_archetype" in html  # fallback card names it


def test_org_layout_preview_places_verbatim_content_in_real_placeholders():
    from pptx import Presentation

    tmpl = Presentation(str(default_template_path()))
    layout = next(l for m in tmpl.slide_masters for l in m.slide_layouts if l.name == "Title and content A")

    html = org_layout_preview_html(
        layout, tmpl.slide_width, tmpl.slide_height,
        "My title", [["Point one", "Point two"]], image_count=0,
    )

    assert "My title" in html
    assert "Point one" in html
    assert "border:1px dashed" in html  # placeholder outlines


def _text_shape(width_in, height_in, text, size_pt):
    """A slide record carrying exactly one text box, for sizing tests."""
    from tests.helpers import add_slide, new_deck, set_run, title_run
    from pptx.util import Inches

    prs = new_deck()
    s = add_slide(prs)
    box = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(width_in), Inches(height_in))
    box.text_frame.text = text
    set_run(box.text_frame.paragraphs[0].runs[0], text=text, size_pt=size_pt)
    title_run(s).text = ""
    return prs, box


def test_slide_preview_sizes_text_in_the_inline_axis_not_cqh():
    """Regression for previews whose type had nothing to do with the
    preview box: `container-type:inline-size` establishes only an INLINE
    query container, so `cqh` had no container to resolve against and
    silently fell back to the viewport height -- a 60pt heading rendered
    at 100px inside a 372px-tall frame. Type must be sized in `cqw`."""
    import re

    from pptx import Presentation

    from deckguard.inventory import build_inventory
    from deckguard.preview import slide_preview_html

    import tempfile, pathlib

    with tempfile.TemporaryDirectory() as td:
        prs, _ = _text_shape(6.0, 3.0, "Ascenseur", 60.0)
        path = pathlib.Path(td) / "d.pptx"
        prs.save(str(path))
        inv = build_inventory(Presentation(str(path)))

    html = slide_preview_html(inv.slides[0])

    assert "cqh" not in html, "block-axis container units don't resolve here"
    sizes = [float(m) for m in re.findall(r"font-size:([\d.]+)cqw", html)]
    assert sizes, "the run must carry an explicit size"
    # 60pt on a 960pt-wide slide is 6.25% of the inline axis.
    assert any(abs(s - 6.25) < 0.05 for s in sizes)


def test_overfull_text_box_shrinks_instead_of_clipping_mid_word():
    """A 60pt word in a 2in column can't be drawn at 60pt without the
    frame slicing it in half. PowerPoint autofits; so must the preview."""
    from deckguard.preview import _fit_scale

    class _R:
        def __init__(self, text, size_pt):
            self.text, self.size_pt = text, size_pt

    class _P:
        def __init__(self, runs):
            self.runs = runs

    roomy = [_P([_R("Ascenseur", 24.0)])]
    assert _fit_scale(roomy, 6.0, 3.0) == 1.0, "text that fits is never shrunk"

    # one word wider than its column
    narrow = [_P([_R("hydraulique", 60.0)])]
    assert _fit_scale(narrow, 2.0, 3.0) < 1.0

    # many lines taller than the box
    tall = [_P([_R("word " * 40, 28.0)])]
    assert _fit_scale(tall, 4.0, 0.6) < 1.0

    # and it never shrinks into illegibility
    assert _fit_scale([_P([_R("x" * 400, 80.0)])], 1.0, 0.3) >= 0.45

    # unknown geometry is left alone
    assert _fit_scale(narrow, None, 3.0) == 1.0
