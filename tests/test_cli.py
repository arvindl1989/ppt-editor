import json

import pytest
from click.testing import CliRunner

from deckguard.cli import main
from tests.helpers import add_slide, new_deck, set_run, title_run


def _write_violating_deck(path):
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Calibri", color_hex="005EB8")
    prs.save(str(path))


def _write_clean_deck(path):
    from pptx.enum.text import PP_ALIGN

    from tests.helpers import body_run, set_theme_font, set_theme_slot

    prs = new_deck()
    set_theme_slot(prs, "dk1", "141414")
    set_theme_font(prs, "majorFont", "Inter")
    set_theme_font(prs, "minorFont", "Inter")
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Inter Semi Bold", color_hex="141414")
    slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    set_run(body_run(slide), text="Body copy", font="Inter", color_hex="141414")
    slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    prs.save(str(path))


def test_validate_rules_default_ok():
    runner = CliRunner()
    result = runner.invoke(main, ["validate-rules"])
    assert result.exit_code == 0
    assert "OK" in result.output


def test_validate_rules_bad_file(tmp_path):
    bad = tmp_path / "bad.yaml"
    bad.write_text("colors: {approved: ['#GGGGGG']}", encoding="utf-8")
    runner = CliRunner()
    result = runner.invoke(main, ["validate-rules", str(bad)])
    assert result.exit_code == 1


def test_fix_rebuild_layouts_flag_rebuilds_non_standard_layout(tmp_path):
    """`deckguard fix --rebuild-layouts`: rebuilds a non_standard_layout
    slide onto an approved layout, verbatim, before patching colors/
    fonts -- opt-in, since default `fix` behavior must stay
    patch-only (see test_fix_default_does_not_rebuild_layouts below)."""
    from pptx import Presentation

    from deckguard.slide_import import default_template_path

    deck = tmp_path / "d.pptx"
    prs = new_deck()
    add_slide(prs)
    middle = add_slide(prs)
    title_run(middle).text = "Middle slide"
    add_slide(prs)
    prs.save(str(deck))

    runner = CliRunner()
    result = runner.invoke(main, ["fix", str(deck), "--rebuild-layouts", "--out", str(tmp_path)])

    assert result.exit_code == 0, result.output
    assert "rebuilt onto an approved layout" in result.output

    out_prs = Presentation(str(tmp_path / "d_fixed.pptx"))
    approved_names = {
        layout.name for master in Presentation(str(default_template_path())).slide_masters for layout in master.slide_layouts
    }
    assert out_prs.slides[1].slide_layout.name in approved_names
    assert out_prs.slides[1].shapes.title.text_frame.text == "Middle slide"


def test_fix_default_does_not_rebuild_layouts(tmp_path):
    deck = tmp_path / "d.pptx"
    prs = new_deck()
    add_slide(prs)
    middle = add_slide(prs)
    title_run(middle).text = "Middle slide"
    original_layout_name = middle.slide_layout.name
    add_slide(prs)
    prs.save(str(deck))

    runner = CliRunner()
    result = runner.invoke(main, ["fix", str(deck), "--out", str(tmp_path)])

    assert result.exit_code == 0, result.output
    assert "rebuilt onto an approved layout" not in result.output

    from pptx import Presentation

    out_prs = Presentation(str(tmp_path / "d_fixed.pptx"))
    assert out_prs.slides[1].slide_layout.name == original_layout_name


def test_inspect_json_is_valid_json(tmp_path):
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)
    runner = CliRunner()
    result = runner.invoke(main, ["inspect", str(deck), "--format", "json"])
    assert result.exit_code == 0
    data = json.loads(result.stdout)
    assert data["slides"][0]["shapes"]


def test_audit_exits_nonzero_when_critical_violation_present(tmp_path):
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)
    runner = CliRunner()
    result = runner.invoke(main, ["audit", str(deck), "--format", "json"])
    assert result.exit_code == 1
    data = json.loads(result.stdout)
    assert data["summary"]["critical"] >= 1


def test_audit_exits_zero_for_clean_deck(tmp_path):
    deck = tmp_path / "clean.pptx"
    _write_clean_deck(deck)
    runner = CliRunner()
    result = runner.invoke(main, ["audit", str(deck), "--format", "json"])
    assert result.exit_code == 0


def test_audit_folder_mode_writes_summary_csv(tmp_path):
    _write_violating_deck(tmp_path / "a.pptx")
    _write_clean_deck(tmp_path / "b.pptx")
    out_dir = tmp_path / "reports"
    runner = CliRunner()
    result = runner.invoke(main, ["audit", str(tmp_path), "--out", str(out_dir)])
    assert (out_dir / "summary.csv").exists()
    content = (out_dir / "summary.csv").read_text()
    assert "a.pptx" in content and "b.pptx" in content


def test_fix_dry_run_leaves_source_untouched(tmp_path):
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)
    original = deck.read_bytes()

    runner = CliRunner()
    result = runner.invoke(main, ["fix", str(deck), "--dry-run", "--out", str(tmp_path / "out")])
    assert result.exit_code == 0
    assert deck.read_bytes() == original
    assert not (tmp_path / "out" / "d_fixed.pptx").exists()
    assert (tmp_path / "out" / "d_changelog.json").exists()
    assert (tmp_path / "out" / "d_changelog.md").exists()


def test_fix_writes_fixed_file(tmp_path):
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)

    runner = CliRunner()
    result = runner.invoke(main, ["fix", str(deck), "--out", str(tmp_path / "out")])
    assert result.exit_code == 0
    assert (tmp_path / "out" / "d_fixed.pptx").exists()


def test_hash_logo_prints_hash(tmp_path):
    from tests.helpers import make_pattern_png

    img = make_pattern_png(tmp_path / "logo.png", seed=9)
    runner = CliRunner()
    result = runner.invoke(main, ["hash-logo", str(img)])
    assert result.exit_code == 0
    assert len(result.stdout.strip()) == 16


def test_inspect_on_corrupt_file_gives_clean_error(tmp_path):
    bad = tmp_path / "not_a_deck.pptx"
    bad.write_bytes(b"not a real pptx file")
    runner = CliRunner()
    result = runner.invoke(main, ["inspect", str(bad)])
    assert result.exit_code == 1
    assert "error" in result.output.lower()
    assert "Traceback" not in result.output


def test_migrate_replaces_non_standard_cover_and_outro(tmp_path):
    from deckguard.slide_import import default_template_path

    if not default_template_path().exists():
        pytest.skip("bundled template asset not present")

    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)  # a single, non-KONE-layout slide
    runner = CliRunner()
    result = runner.invoke(main, ["migrate", str(deck), "--out", str(tmp_path)])
    assert result.exit_code == 0, result.output
    assert "cover replaced" in result.output
    assert "outro replaced" in result.output
    assert (tmp_path / "d_migrated.pptx").exists()


def test_redesign_mode_brand_needs_no_api_key(tmp_path, monkeypatch):
    """--mode brand is fully deterministic -- unlike --mode rewrite
    (default), it must run with no ANTHROPIC_API_KEY set at all."""
    from deckguard.slide_import import default_template_path

    if not default_template_path().exists():
        pytest.skip("bundled template asset not present")

    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)
    out_path = tmp_path / "rebranded.pptx"

    runner = CliRunner()
    result = runner.invoke(main, ["redesign", str(deck), "--out", str(out_path), "--mode", "brand"])

    assert result.exit_code == 0, result.output
    assert "mode: brand" in result.output
    assert out_path.exists()


def test_redesign_mode_rewrite_still_requires_api_key(tmp_path, monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)
    out_path = tmp_path / "out.pptx"

    runner = CliRunner()
    result = runner.invoke(main, ["redesign", str(deck), "--out", str(out_path)])

    assert result.exit_code == 1
    assert "ANTHROPIC_API_KEY is not set" in result.output


def test_redesign_mode_brand_without_deck_errors_cleanly(tmp_path, monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    runner = CliRunner()
    result = runner.invoke(main, ["redesign", "--out", str(tmp_path / "out.pptx"), "--mode", "brand"])

    assert result.exit_code == 1
    assert "needs a DECK" in result.output


def test_redesign_mode_brand_with_review_requires_api_key(tmp_path, monkeypatch):
    """Unlike plain --mode brand, --review makes one small API call and
    needs a key even though the rest of brand mode doesn't."""
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)
    runner = CliRunner()
    result = runner.invoke(
        main, ["redesign", str(deck), "--out", str(tmp_path / "out.pptx"), "--mode", "brand", "--review"]
    )

    assert result.exit_code == 1
    assert "ANTHROPIC_API_KEY is not set" in result.output


def test_redesign_review_rejects_rewrite_mode(tmp_path, monkeypatch):
    monkeypatch.setenv("ANTHROPIC_API_KEY", "test-key")
    deck = tmp_path / "d.pptx"
    _write_violating_deck(deck)
    runner = CliRunner()
    result = runner.invoke(
        main, ["redesign", str(deck), "--out", str(tmp_path / "out.pptx"), "--mode", "rewrite", "--review"]
    )

    assert result.exit_code == 1
    assert "--review only applies to --mode brand" in result.output


def _write_deck_pair_for_learn(old_path, new_path):
    from tests.helpers import body_run

    old_prs = new_deck()
    slide = add_slide(old_prs)
    set_run(body_run(slide), text="Body copy", font="Arial", color_hex="AABBCC")
    old_prs.save(str(old_path))

    new_prs = new_deck()
    slide2 = add_slide(new_prs)
    set_run(body_run(slide2), text="Body copy", font="Inter", color_hex="1450F5")
    new_prs.save(str(new_path))


def test_learn_transform_writes_a_rebuilt_deck(tmp_path, monkeypatch):
    """--transform rebuilds OLD_DECK onto the org template's own approved
    layouts using the just-learned colors/fonts -- the same engine
    `redesign --mode brand` uses -- not just a proposals report."""
    from deckguard.slide_import import default_template_path

    if not default_template_path().exists():
        pytest.skip("bundled template asset not present")

    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    old_path = tmp_path / "old.pptx"
    new_path = tmp_path / "new.pptx"
    _write_deck_pair_for_learn(old_path, new_path)
    transform_path = tmp_path / "transformed.pptx"

    runner = CliRunner()
    result = runner.invoke(
        main, ["learn", str(old_path), str(new_path), "--transform", str(transform_path)]
    )

    assert result.exit_code == 0, result.output
    assert transform_path.exists()
    assert "wrote:" in result.output


def test_learn_review_requires_transform(tmp_path):
    old_path = tmp_path / "old.pptx"
    new_path = tmp_path / "new.pptx"
    _write_deck_pair_for_learn(old_path, new_path)

    runner = CliRunner()
    result = runner.invoke(main, ["learn", str(old_path), str(new_path), "--review"])

    assert result.exit_code == 1
    assert "--review only applies together with --transform" in result.output


def test_learn_review_requires_api_key(tmp_path, monkeypatch):
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    old_path = tmp_path / "old.pptx"
    new_path = tmp_path / "new.pptx"
    _write_deck_pair_for_learn(old_path, new_path)

    runner = CliRunner()
    result = runner.invoke(
        main, ["learn", str(old_path), str(new_path), "--transform", str(tmp_path / "out.pptx"), "--review"]
    )

    assert result.exit_code == 1
    assert "ANTHROPIC_API_KEY is not set" in result.output


# --------------------------------------------------------------------------
# create-archetype -- no AI, builds from a hand-written archetype spec
# --------------------------------------------------------------------------


def _kone_skill_available() -> bool:
    from deckguard.skill_bridge import _skill_dir

    return (_skill_dir() / "kone_deck_creator.py").is_file()


needs_kone_skill = pytest.mark.skipif(not _kone_skill_available(), reason="kone-deck-generator skill not installed")


@needs_kone_skill
def test_create_archetype_builds_a_deck_from_a_hand_written_spec(tmp_path):
    spec_path = tmp_path / "spec.json"
    spec_path.write_text(json.dumps({
        "title": "Hand-authored deck",
        "slides": [
            {"archetype": "agenda_contents", "title": "Agenda",
             "items": [{"number": "01", "item": "Point one"}]},
        ],
    }), encoding="utf-8")
    out_path = tmp_path / "out.pptx"

    runner = CliRunner()
    result = runner.invoke(main, ["create-archetype", str(spec_path), "--out", str(out_path)])

    assert result.exit_code == 0, result.output
    assert out_path.exists()
    assert "agenda_contents" in result.output
    from pptx import Presentation
    assert len(Presentation(str(out_path)).slides) == 3  # Cover F + 1 body + Outro


@needs_kone_skill
def test_create_archetype_rejects_an_unknown_archetype(tmp_path):
    spec_path = tmp_path / "spec.json"
    spec_path.write_text(json.dumps({
        "title": "T", "slides": [{"archetype": "not_a_real_archetype", "title": "x"}],
    }), encoding="utf-8")

    runner = CliRunner()
    result = runner.invoke(main, ["create-archetype", str(spec_path), "--out", str(tmp_path / "out.pptx")])

    assert result.exit_code == 1
    assert "unknown archetype" in result.output


def test_create_archetype_rejects_invalid_json(tmp_path):
    spec_path = tmp_path / "spec.json"
    spec_path.write_text("not json", encoding="utf-8")

    runner = CliRunner()
    result = runner.invoke(main, ["create-archetype", str(spec_path), "--out", str(tmp_path / "out.pptx")])

    assert result.exit_code == 1
    assert "not valid JSON" in result.output


def test_sync_skill_copies_changed_files_and_removes_retired_ones(tmp_path, monkeypatch):
    """`deckguard sync-skill`: after a skill update, one copy step (plus a
    commit) is all a deploy needs -- verify changed files copy, retired
    files are removed from the vendored copy, and unchanged files are
    left alone."""
    import deckguard.skill_bridge as skill_bridge
    from deckguard.cli import main as cli_main

    # A fake "updated skill" source and a fake vendored destination.
    src = tmp_path / "skill"
    (src / "fonts").mkdir(parents=True)
    (src / "kone_deck_creator.py").write_text("print('v2')", encoding="utf-8")
    (src / "catalog.json").write_text("{}", encoding="utf-8")
    dest = tmp_path / "vendored"
    dest.mkdir()
    (dest / "kone_deck_creator.py").write_text("print('v1')", encoding="utf-8")
    (dest / "retired_module.py").write_text("gone in v2", encoding="utf-8")
    monkeypatch.setattr(skill_bridge, "_VENDORED_SKILL_DIR", dest)

    runner = CliRunner()
    result = runner.invoke(cli_main, ["sync-skill", "--source", str(src)])

    assert result.exit_code == 0, result.output
    assert (dest / "kone_deck_creator.py").read_text(encoding="utf-8") == "print('v2')"
    assert (dest / "catalog.json").is_file()
    assert not (dest / "retired_module.py").exists()
    assert "1 updated, 1 added, 1 removed" in result.output


def test_sync_skill_refuses_to_sync_the_vendored_copy_onto_itself(tmp_path, monkeypatch):
    import deckguard.skill_bridge as skill_bridge
    from deckguard.cli import main as cli_main

    dest = tmp_path / "vendored"
    dest.mkdir()
    (dest / "kone_deck_creator.py").write_text("x", encoding="utf-8")
    monkeypatch.setattr(skill_bridge, "_VENDORED_SKILL_DIR", dest)

    runner = CliRunner()
    result = runner.invoke(cli_main, ["sync-skill", "--source", str(dest)])

    assert result.exit_code == 1
    assert "IS the vendored copy" in result.output


def test_visual_check_reports_what_a_deck_lays_out_to(tmp_path):
    """`visual-check` measures rendered layout rather than reading XML.
    Without a browser it must say so and exit non-zero, not crash."""
    from deckguard.visual import playwright_available

    prs = new_deck()
    s = add_slide(prs)
    set_run(title_run(s), text="Readable title", font="Inter", color_hex="141414")
    deck = tmp_path / "d.pptx"
    prs.save(str(deck))

    result = CliRunner().invoke(main, ["visual-check", str(deck), "--format", "json"])

    if not playwright_available():
        assert result.exit_code == 1
        assert "Playwright" in result.output
        return

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload["ran"] is True
    assert payload["frames_measured"] == 1
    assert payload["summary"]["major"] == 0
