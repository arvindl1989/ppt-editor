"""Tests for skill_bridge.py -- deckguard's bridge to the (out-of-repo,
per-machine) kone-deck-generator skill, used only for redesign_deck's
"no source deck, just a brief" path. No real network calls -- these
inject the same fake Anthropic-shaped client the rest of the redesign
test suite uses.
"""

import json
from pathlib import Path

import pytest

from deckguard.compose import Outline, SlideSpec
from deckguard.redesign import RedesignError
from deckguard.skill_bridge import (
    _load_archetypes,
    _skill_dir,
    _validate_kone_spec,
    apply_archetype_overrides_to_deck,
    build_deck_via_skill,
    build_deck_with_archetypes,
    call_claude_for_kone_spec,
    select_archetype_overrides,
    select_archetype_overrides_for_rebrand,
)
from tests.test_redesign import _FakeClient, _FakeResponse, _kone_slide, _kone_spec_json

_skill_installed = (_skill_dir() / "kone_deck_creator.py").is_file()
needs_skill = pytest.mark.skipif(not _skill_installed, reason="kone-deck-generator skill not installed")


def _known_archetypes():
    return set(_load_archetypes().ARCHETYPES.keys())


def test_build_deck_via_skill_raises_a_clean_error_when_the_skill_is_not_installed(tmp_path, monkeypatch):
    """Doesn't need the real skill on disk -- exercises the opposite case:
    an actionable RedesignError, never a raw ImportError/FileNotFoundError,
    when KONE_DECK_GENERATOR_DIR points nowhere real."""
    import deckguard.skill_bridge as skill_bridge

    monkeypatch.setattr(skill_bridge, "_creator_module", None)
    monkeypatch.setattr(skill_bridge, "_archetypes_module", None)
    monkeypatch.setattr(skill_bridge, "_catalog_cache", None)
    monkeypatch.setenv("KONE_DECK_GENERATOR_DIR", str(tmp_path / "does-not-exist"))

    with pytest.raises(RedesignError, match="isn't installed"):
        build_deck_via_skill("A brief.", str(tmp_path / "out.pptx"), client=_FakeClient(_FakeResponse("{}")))


@needs_skill
def test_validate_kone_spec_accepts_a_well_formed_spec():
    spec = json.loads(
        _kone_spec_json(
            "Deck title",
            _kone_slide("image_section_divider", eyebrow="Overview", title="A new era"),
            _kone_slide("agenda_contents", title="Agenda", items=[{"number": "01", "item": "Kickoff"}]),
            _kone_slide(
                "three_stats", title="By the numbers",
                stats=[
                    {"label": "Volume", "value": "2x", "desc": "growth"},
                    {"label": "Resolution", "value": "91%", "desc": "resolved"},
                    {"label": "Phases", "value": "3", "desc": "planned"},
                ],
            ),
            _kone_slide("quote_context", quote="Ship it.", attribution="A KONE engineer"),
        )
    )
    _validate_kone_spec(spec, _known_archetypes())  # must not raise


@needs_skill
def test_validate_kone_spec_rejects_unknown_archetype():
    spec = {"title": "T", "slides": [_kone_slide("carousel", title="x")]}
    with pytest.raises(RedesignError, match="unknown archetype"):
        _validate_kone_spec(spec, _known_archetypes())


@needs_skill
def test_validate_kone_spec_rejects_a_non_object_slide():
    spec = {"title": "T", "slides": ["not a dict"]}
    with pytest.raises(RedesignError, match="not a JSON object"):
        _validate_kone_spec(spec, _known_archetypes())


def test_validate_kone_spec_rejects_missing_title():
    with pytest.raises(RedesignError, match="missing deck 'title'"):
        _validate_kone_spec({"slides": [{"archetype": "three_stats"}]}, {"three_stats"})


def test_validate_kone_spec_rejects_zero_slides():
    with pytest.raises(RedesignError, match="zero slides"):
        _validate_kone_spec({"title": "T", "slides": []}, {"three_stats"})


@needs_skill
def test_validate_kone_spec_accepts_text_over_typical_length_guidance():
    """Character limits are advisory (the renderer's shrink-to-fit
    handles overflow) -- validation must not reject a merely-long value."""
    spec = {
        "title": "T",
        "slides": [_kone_slide("agenda_contents", title="x" * 200, items=[{"number": "01", "item": "y" * 300}])],
    }
    _validate_kone_spec(spec, _known_archetypes())  # must not raise


@needs_skill
def test_call_claude_for_kone_spec_sends_the_brief_and_returns_usage():
    response_json = _kone_spec_json("T", _kone_slide("agenda_contents", title="X", items=[]))
    client = _FakeClient(_FakeResponse(response_json, input_tokens=1200, output_tokens=300))

    spec, usage = call_claude_for_kone_spec("Build a deck about X.", client=client)

    assert spec["title"] == "T"
    assert usage.input_tokens == 1200 and usage.output_tokens == 300
    sent = client.messages.calls[0]
    assert "Build a deck about X." in sent["messages"][0]["content"]
    # no rigid structured-output schema for this path -- see skill_bridge's
    # own docstring for why (content shape varies too much per archetype)
    assert "output_config" not in sent


@needs_skill
def test_call_claude_for_kone_spec_system_prompt_includes_the_real_archetype_catalog():
    response_json = _kone_spec_json("T", _kone_slide("agenda_contents", title="X", items=[]))
    client = _FakeClient(_FakeResponse(response_json))

    call_claude_for_kone_spec("Build a deck about X.", client=client)

    system = client.messages.calls[0]["system"]
    archetypes = _load_archetypes()
    for name in archetypes.ARCHETYPES:
        assert name in system


@needs_skill
def test_call_claude_for_kone_spec_strips_markdown_code_fences():
    response_json = "```json\n" + _kone_spec_json("T", _kone_slide("agenda_contents", title="X", items=[])) + "\n```"
    client = _FakeClient(_FakeResponse(response_json))

    spec, _usage = call_claude_for_kone_spec("Build a deck about X.", client=client)

    assert spec["title"] == "T"


@needs_skill
def test_build_deck_via_skill_renders_a_real_pptx(tmp_path):
    response_json = _kone_spec_json(
        "My deck",
        _kone_slide("agenda_contents", title="Agenda", items=[{"number": "01", "item": "Point one"}]),
    )
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "out.pptx"
    compose_result, redesign_result = build_deck_via_skill("A brief.", str(out_path), client=client)

    assert out_path.exists()
    assert compose_result.slide_count == 3  # Cover F + 1 body + Outro
    assert compose_result.layouts_used == ["Cover F", "agenda_contents", "Outro"]
    assert redesign_result.outline is None
    assert redesign_result.usage.input_tokens > 0


@needs_skill
def test_build_deck_via_skill_renders_content_over_the_old_hard_limit(tmp_path):
    """Regression guard: a value that would have tripped the retired
    hard character-limit validation must still render successfully."""
    response_json = _kone_spec_json(
        "My deck",
        _kone_slide("quote_context", label="Voice of the business, in their own words", quote="Ship it.", attribution="A KONE engineer"),
    )
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "out.pptx"
    compose_result, _redesign_result = build_deck_via_skill("A brief.", str(out_path), client=client)

    assert out_path.exists()
    assert compose_result.slide_count == 3


@needs_skill
def test_build_deck_via_skill_raises_a_clean_error_for_an_invalid_spec(tmp_path):
    response_json = _kone_spec_json("My deck", _kone_slide("not_a_real_archetype", title="Compare"))
    client = _FakeClient(_FakeResponse(response_json))

    with pytest.raises(RedesignError, match="doesn't fit the kone-deck-generator skill's archetypes"):
        build_deck_via_skill("A brief.", str(tmp_path / "out.pptx"), client=client)


# --------------------------------------------------------------------------
# select_archetype_overrides -- additive, fail-closed pass over an
# already-planned compose.py outline
# --------------------------------------------------------------------------


def _outline_item(kind="content", **overrides):
    item = {
        "source_slide_index": 1, "kind": kind, "title": "A title", "subtitle": None,
        "bullets": ["Point one"], "columns": [], "quote_text": None, "quote_author": None,
        "quote_label": None, "stats": [], "milestones": [], "variant": None,
    }
    item.update(overrides)
    return item


@needs_skill
def test_select_archetype_overrides_returns_a_valid_override():
    items = [_outline_item(kind="stat", title="Resolution", bullets=[])]
    response = json.dumps({"overrides": [
        {"outline_index": 0, "archetype": "hero_stat", "eyebrow": "Resolution", "value": "91%",
         "caption": "resolved", "support": "91% resolved"},
    ]})
    client = _FakeClient(_FakeResponse(response))

    overrides = select_archetype_overrides(items, client=client)

    assert overrides == {0: {
        "archetype": "hero_stat", "eyebrow": "Resolution", "value": "91%",
        "caption": "resolved", "support": "91% resolved",
    }}


def test_select_archetype_overrides_skips_the_call_entirely_with_no_candidate_kinds():
    items = [_outline_item(kind="cover"), _outline_item(kind="quote"), _outline_item(kind="end")]
    client = _FakeClient(_FakeResponse("{}"))

    overrides = select_archetype_overrides(items, client=client)

    assert overrides == {}
    assert client.messages.calls == []


@needs_skill
def test_select_archetype_overrides_rejects_an_unknown_archetype_name():
    items = [_outline_item(kind="content")]
    response = json.dumps({"overrides": [{"outline_index": 0, "archetype": "not_a_real_archetype", "title": "x"}]})
    client = _FakeClient(_FakeResponse(response))

    overrides = select_archetype_overrides(items, client=client)

    assert overrides == {}


@needs_skill
def test_select_archetype_overrides_rejects_an_out_of_range_outline_index():
    items = [_outline_item(kind="content")]
    response = json.dumps({"overrides": [{"outline_index": 5, "archetype": "hero_stat", "value": "x"}]})
    client = _FakeClient(_FakeResponse(response))

    overrides = select_archetype_overrides(items, client=client)

    assert overrides == {}


def test_select_archetype_overrides_fails_closed_on_malformed_json():
    items = [_outline_item(kind="content")]
    client = _FakeClient(_FakeResponse("not json"))

    overrides = select_archetype_overrides(items, client=client)

    assert overrides == {}


def test_select_archetype_overrides_fails_closed_on_refusal():
    items = [_outline_item(kind="content")]
    client = _FakeClient(_FakeResponse('{"overrides": []}', stop_reason="refusal"))

    overrides = select_archetype_overrides(items, client=client)

    assert overrides == {}


# --------------------------------------------------------------------------
# build_deck_with_archetypes
# --------------------------------------------------------------------------


def test_build_deck_with_archetypes_delegates_to_compose_build_deck_with_no_overrides(tmp_path):
    outline = Outline(slides=[SlideSpec(kind="content", title="Highlights", bullets=["Grew nicely"])])

    from deckguard.compose import build_deck as compose_build_deck

    expected_out = tmp_path / "expected.pptx"
    expected = compose_build_deck(outline, str(expected_out))

    got_out = tmp_path / "got.pptx"
    got = build_deck_with_archetypes(outline, str(got_out), overrides={})

    assert got.slide_count == expected.slide_count
    assert got.layouts_used == expected.layouts_used


@needs_skill
def test_build_deck_with_archetypes_renders_an_override_and_skips_fix_deck_for_it(tmp_path):
    outline = Outline(slides=[
        SlideSpec(kind="content", title="Highlights", bullets=["Grew nicely"]),
        SlideSpec(kind="stat", title="Resolution"),  # content irrelevant -- overridden below
    ])
    overrides = {1: {
        "archetype": "hero_stat", "eyebrow": "Resolution", "value": "91.2%",
        "caption": "of all requests cleared within the focus period",
        "support": "674 of 739 tickets resolved.",
    }}

    out_path = tmp_path / "out.pptx"
    result = build_deck_with_archetypes(outline, str(out_path), overrides=overrides)

    assert result.slide_count == 2
    assert result.layouts_used[1] == "hero_stat"

    from pptx import Presentation
    from pptx.dml.color import RGBColor

    prs = Presentation(str(out_path))
    archetype_slide = prs.slides[1]
    texts = [
        s.text_frame.text for s in archetype_slide.shapes
        if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()
    ]
    assert any("91.2%" in t for t in texts)

    # The muted-grey caption/support text (#727272, kone_engine.GREY) is
    # NOT in brand_rules.yaml's approved-colors list -- if fix_deck had
    # run over this slide it would have flagged or remapped it. Confirm
    # it survives untouched.
    grey_runs = [
        run for s in archetype_slide.shapes if getattr(s, "has_text_frame", False)
        for p in s.text_frame.paragraphs for run in p.runs
        if run.font.color.type is not None and str(run.font.color.rgb) == "727272"
    ]
    assert grey_runs, "expected at least one run still using kone_engine's own grey #727272"


@needs_skill
def test_build_deck_with_archetypes_preserves_slide_order_for_a_middle_override(tmp_path):
    """The trickier reorder case: archetype slides are physically added
    to the presentation LAST (after fix_deck runs over everything
    else), so a middle-position override specifically exercises the
    _sldIdLst reassembly -- not just "append one archetype slide at
    the end", which wouldn't catch a reordering bug."""
    outline = Outline(slides=[
        SlideSpec(kind="content", title="First", bullets=["A"]),
        SlideSpec(kind="stat", title="Middle (overridden)"),
        SlideSpec(kind="content", title="Third", bullets=["C"]),
    ])
    overrides = {1: {
        "archetype": "hero_stat", "eyebrow": "Middle", "value": "42%",
        "caption": "in the middle", "support": "middle slide check",
    }}

    out_path = tmp_path / "out.pptx"
    result = build_deck_with_archetypes(outline, str(out_path), overrides=overrides)

    assert result.layouts_used == ["Title and content A", "hero_stat", "Title and content A"]

    from pptx import Presentation

    prs = Presentation(str(out_path))
    titles = [
        next(s.text_frame.text for s in slide.shapes if getattr(s, "has_text_frame", False) and s.text_frame.text.strip())
        for slide in prs.slides
    ]
    assert titles[0] == "First"
    assert "42%" in titles[1] or titles[1] == "MIDDLE"
    assert titles[2] == "Third"


# --------------------------------------------------------------------------
# apply_rebrand's (mode='brand', review=True only) archetype coexistence:
# select_archetype_overrides_for_rebrand + apply_archetype_overrides_to_deck
# --------------------------------------------------------------------------


def _build_simple_deck(out_path, n=3):
    from deckguard.compose import build_deck as compose_build_deck

    outline = Outline(slides=[SlideSpec(kind="content", title=f"Slide {i}", bullets=[f"Point {i}"]) for i in range(1, n + 1)])
    compose_build_deck(outline, str(out_path))


@needs_skill
def test_select_archetype_overrides_for_rebrand_returns_a_valid_override():
    from deckguard.retemplate import SlideProfile

    profiles = {2: SlideProfile(title="Resolution", text_blocks=[[(0, "91% resolved")]], images=[], eligible=True)}
    response = json.dumps({"overrides": [
        {"outline_index": 2, "archetype": "hero_stat", "eyebrow": "Resolution", "value": "91%",
         "caption": "c", "support": "s"},
    ]})
    client = _FakeClient(_FakeResponse(response))

    overrides = select_archetype_overrides_for_rebrand(profiles, client=client)

    assert overrides == {2: {
        "archetype": "hero_stat", "eyebrow": "Resolution", "value": "91%", "caption": "c", "support": "s",
    }}


def test_select_archetype_overrides_for_rebrand_skips_the_call_with_no_candidates():
    client = _FakeClient(_FakeResponse("{}"))

    overrides = select_archetype_overrides_for_rebrand({}, client=client)

    assert overrides == {}
    assert client.messages.calls == []


@needs_skill
def test_apply_archetype_overrides_to_deck_swaps_a_slide_in_place(tmp_path):
    deck_path = tmp_path / "deck.pptx"
    _build_simple_deck(deck_path, n=3)

    overrides = {2: {
        "archetype": "hero_stat", "eyebrow": "Middle", "value": "42%", "caption": "c", "support": "s",
    }}
    layout_by_index = apply_archetype_overrides_to_deck(str(deck_path), overrides)

    assert layout_by_index == {2: "hero_stat"}

    from pptx import Presentation

    prs = Presentation(str(deck_path))
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text_frame.text == "Slide 1"
    assert prs.slides[2].shapes.title.text_frame.text == "Slide 3"
    texts = [
        s.text_frame.text for s in prs.slides[1].shapes
        if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()
    ]
    assert any("42%" in t for t in texts)


def test_apply_archetype_overrides_to_deck_is_a_noop_with_no_overrides(tmp_path):
    deck_path = tmp_path / "deck.pptx"
    _build_simple_deck(deck_path, n=1)

    result = apply_archetype_overrides_to_deck(str(deck_path), {})

    assert result == {}


# --------------------------------------------------------------------------
# Source-image carryover into an archetype's picture slots
# --------------------------------------------------------------------------


def _blobs(tmp_path, n):
    """n visually distinct image blobs, as raw bytes -- the same shape
    SlideProfile.images / _attach_source_images carry."""
    from tests.helpers import make_pattern_png

    return [Path(make_pattern_png(tmp_path / f"p{i}.png", seed=i)).read_bytes() for i in range(n)]


def _pictures(slide):
    return [s for s in slide.shapes if s.shape_type is not None and s.shape_type.name == "PICTURE"]


@needs_skill
def test_archetype_override_carries_a_source_image_into_a_single_picture_slot(tmp_path):
    """The gap this closes: an image-bearing slide used to render its
    archetype's picture slot as an empty sand placeholder, because
    kone_engine._image() wants a path while everything upstream carries
    raw bytes."""
    from pptx import Presentation

    outline = Outline(slides=[SlideSpec(kind="content", title="Recap", bullets=["a"], images=_blobs(tmp_path, 1))])
    overrides = {0: {
        "archetype": "numbered_summary_picture", "title": "Recap",
        "items": [{"number": "01", "text": "First point"}],
    }}

    out_path = tmp_path / "out.pptx"
    build_deck_with_archetypes(outline, str(out_path), overrides=overrides)

    pics = _pictures(Presentation(str(out_path)).slides[0])
    assert len(pics) == 1
    assert len(pics[0].image.blob) > 0


@needs_skill
def test_archetype_override_carries_one_image_per_group_item(tmp_path):
    from pptx import Presentation

    outline = Outline(slides=[SlideSpec(kind="content", title="Options", bullets=["a"], images=_blobs(tmp_path, 3))])
    overrides = {0: {
        "archetype": "three_picture_cards", "title": "Options",
        "cards": [{"heading": f"Card {i}", "bullets": ["x"]} for i in range(3)],
    }}

    out_path = tmp_path / "out.pptx"
    build_deck_with_archetypes(outline, str(out_path), overrides=overrides)

    pics = _pictures(Presentation(str(out_path)).slides[0])
    assert len(pics) == 3
    assert len({p.image.blob for p in pics}) == 3  # each card got its OWN image, not the same one repeated


@needs_skill
def test_archetype_override_still_renders_with_no_source_images(tmp_path):
    """A picture archetype chosen for a slide that has no images must
    still render (kone_engine falls back to its sand placeholder), never
    raise -- image injection is strictly additive."""
    outline = Outline(slides=[SlideSpec(kind="content", title="Options", bullets=["a"], images=[])])
    overrides = {0: {
        "archetype": "three_picture_cards", "title": "Options",
        "cards": [{"heading": "Card", "bullets": ["x"]}],
    }}

    result = build_deck_with_archetypes(outline, str(tmp_path / "out.pptx"), overrides=overrides)

    assert result.layouts_used == ["three_picture_cards"]


@needs_skill
def test_apply_archetype_overrides_to_deck_carries_images(tmp_path):
    """The brand-mode (--review) path takes its images via an explicit
    images_by_index arg, since it works from a finished deck rather than
    a compose.py outline."""
    from pptx import Presentation

    deck_path = tmp_path / "deck.pptx"
    _build_simple_deck(deck_path, n=2)

    overrides = {2: {"archetype": "numbered_summary_picture", "title": "Recap",
                     "items": [{"number": "01", "text": "Point"}]}}
    apply_archetype_overrides_to_deck(str(deck_path), overrides, images_by_index={2: _blobs(tmp_path, 1)})

    pics = _pictures(Presentation(str(deck_path)).slides[1])
    assert len(pics) == 1 and len(pics[0].image.blob) > 0


@needs_skill
def test_planning_prompt_reports_image_count_and_never_leaks_a_photo_path():
    """The skill's own SAMPLES carry absolute local photo paths (they
    render a standalone gallery); echoing those into the prompt would
    invite the model to emit or invent a path, when picture slots are
    filled automatically."""
    from deckguard.skill_bridge import _kone_archetype_guide

    guide = _kone_archetype_guide()
    assert "/assets/photos/" not in guide
    assert "filled automatically" in guide

    items = [_outline_item(kind="content", images=[b"\x89PNG-fake"])]
    client = _FakeClient(_FakeResponse(json.dumps({"overrides": []})))
    select_archetype_overrides(items, client=client)

    sent = client.messages.calls[0]["messages"][0]["content"]
    assert '"image_count": 1' in sent
    assert "PNG-fake" not in sent  # raw bytes never serialized


@needs_skill
def test_image_slots_are_derived_from_the_skills_own_archetype_data():
    """The picture-slot map is no longer hand-maintained -- it's derived
    from the skill's ARCHETYPES regions at runtime, so a skill update
    that adds/renames a picture archetype flows through with no code
    change here. Spot-check the three slot shapes plus the figure
    exclusion (a chart slot must never be treated as a photo slot)."""
    from deckguard.skill_bridge import _archetype_image_slots, archetype_image_capacity

    slots = _archetype_image_slots()
    assert slots["numbered_summary_picture"] == ("single", "image")
    assert slots["three_picture_cards"] == ("group", "cards", "image")
    assert "chart_commentary" not in slots  # figure-role: render() overwrites it with bundled art
    assert "segment_breakdown" not in slots
    assert archetype_image_capacity("chart_commentary") == 0
    assert archetype_image_capacity("numbered_summary_picture") == 1


def test_photo_library_is_reachable_without_the_skills_installed(monkeypatch):
    """The deploy has no ~/.claude/skills, so the photos have to be
    vendored or picture slots stay empty there forever."""
    from pathlib import Path

    from deckguard.skill_bridge import _photo_library

    monkeypatch.delenv("KONE_DESIGN_DIR", raising=False)
    monkeypatch.setattr(Path, "home", classmethod(lambda cls: Path("/nonexistent-home")))

    photos = _photo_library()
    assert photos, "a vendored photo library must exist"
    assert all(p.endswith(".jpg") for p in photos)


def test_empty_photo_slots_are_filled_for_a_from_scratch_deck():
    """Reported on a real build: the review previews drew PHOTO slots,
    the delivered deck had blank sand blocks where they were, because
    `kone_engine._image()` draws a sand rectangle when handed no path
    and nothing carries images into a build from a brief."""
    from deckguard.skill_bridge import _archetype_image_slots, fill_empty_photo_slots

    slots = _archetype_image_slots()
    single = next((n for n, s in slots.items() if s[0] == "single"), None)
    group = next((n for n, s in slots.items() if s[0] == "group"), None)
    assert single and group, "the skill must expose both slot shapes"

    spec = {"title": "t", "slides": [
        {"archetype": single, "title": "The growth of marketing"},
        {"archetype": group, "title": "Compare", "items": [{"label": "A"}, {"label": "B"}]},
    ]}
    assert fill_empty_photo_slots(spec) == 3

    key = slots[single][1]
    assert spec["slides"][0][key].endswith(".jpg")
    item_key = slots[group][2]
    photos = {item[item_key] for item in spec["slides"][1]["items"]}
    assert len(photos) == 2, "the library is walked, not repeated down the deck"


def test_filling_never_overwrites_an_image_the_deck_already_supplied():
    """A transform carrying the source deck's own images always wins --
    this fallback is only for slots nothing else can fill."""
    from deckguard.skill_bridge import _archetype_image_slots, fill_empty_photo_slots

    slots = _archetype_image_slots()
    single = next(n for n, s in slots.items() if s[0] == "single")
    key = slots[single][1]

    spec = {"title": "t", "slides": [{"archetype": single, "title": "x", key: "/from/the/deck.png"}]}
    assert fill_empty_photo_slots(spec) == 0
    assert spec["slides"][0][key] == "/from/the/deck.png"


def test_filling_is_deterministic_for_the_same_brief():
    from deckguard.skill_bridge import _archetype_image_slots, fill_empty_photo_slots

    slots = _archetype_image_slots()
    single = next(n for n, s in slots.items() if s[0] == "single")
    key = slots[single][1]

    def _built():
        spec = {"title": "t", "slides": [{"archetype": single, "title": "Same brief"}]}
        fill_empty_photo_slots(spec)
        return spec["slides"][0][key]

    assert _built() == _built()


def test_archetype_suggestions_availability_reports_a_keyless_server(monkeypatch):
    """The review page told a user suggestions had run on a server with
    no API key, then offered nothing but "keep" on ten of twelve slides
    with no explanation."""
    from deckguard.skill_bridge import archetype_suggestions_available

    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    assert archetype_suggestions_available() is False
    assert archetype_suggestions_available(api_key="k") is True
    assert archetype_suggestions_available(client=object()) is True

    monkeypatch.setenv("ANTHROPIC_API_KEY", "k")
    assert archetype_suggestions_available() is True


def test_every_archetype_exposes_a_machine_readable_shape():
    """The matcher reads the skill's own region/group data, so a skill
    update changes what can be matched with no code change here."""
    from deckguard.skill_bridge import archetype_signatures

    sigs = archetype_signatures()
    assert len(sigs) >= 20
    assert all(s["name"] and isinstance(s["capacity"], int) for s in sigs)
    # only an archetype needing content nothing can synthesise (a real
    # table) is excluded from matching
    assert sum(1 for s in sigs if s["unfillable"]) <= 2


def test_a_slide_is_matched_to_an_archetype_without_any_model():
    """The capability the whole thing turns on: the model was the ONLY
    route to an archetype, so a keyless server offered dense slides
    nothing but "keep"."""
    from deckguard.skill_bridge import match_archetypes

    blocks = [[f"Point {i}: something the slide says"] for i in range(8)]
    candidates = match_archetypes("Subject: get ready for your next project", blocks, image_count=4)

    assert candidates, "a readable slide must always have somewhere to go"
    best = candidates[0]
    assert best["content"]["archetype"] == best["archetype"]
    # 8 chunks into a 6-slot archetype: 6 kept, 2 honestly reported lost
    assert best["capacity"] >= 5
    assert best["dropped"] == len(blocks) - best["capacity"]


def test_matching_prefers_the_archetype_whose_capacity_fits():
    from deckguard.skill_bridge import match_archetypes

    three = match_archetypes("A title", [["A", "detail"], ["B", "detail"], ["C", "detail"]])
    assert three[0]["capacity"] == 3 and three[0]["dropped"] == 0

    two = match_archetypes("A title", [["Left", "why"], ["Right", "why"]])
    assert two[0]["dropped"] == 0


def test_a_stat_slide_matches_a_stat_archetype():
    """Archetypes with a value slot are only offered when the slide
    actually has a number to put in it."""
    from deckguard.skill_bridge import archetype_signatures, match_archetypes

    value_archetypes = {s["name"] for s in archetype_signatures() if s["needs_value"]}
    with_stat = match_archetypes("Resolution rate", [["91.2%", "of requests cleared"]])
    assert any(c["archetype"] in value_archetypes for c in with_stat)

    without = match_archetypes("Some words", [["No numbers here", "none at all"]])
    assert all(c["archetype"] not in value_archetypes for c in without)


def test_matched_content_uses_the_slides_own_words():
    from deckguard.skill_bridge import match_archetypes

    best = match_archetypes("Quarter in review", [["Requests in", "739 across three teams"]])[0]
    flat = json.dumps(best["content"])
    assert "Quarter in review" in flat
    assert "Requests in" in flat


def test_gallery_only_archetype_names_are_reported_not_silently_substituted():
    """Reported from a real build: a brief asking for COVER_A_CUT4 /
    DIVIDER_D / END_LOGO / TITLE_TEXT_SPLIT got a deck built from
    entirely different archetypes with nothing said about it. Those four
    are in kone-design's 56-name gallery, not in the 23 the engine
    renders -- only 17 names exist in both vocabularies."""
    from deckguard.skill_bridge import check_brief_archetypes

    brief = (
        "Create a marketing hub review deck, use COVER_A_CUT4 for the title slide, "
        "DIVIDER_D for dividers, END_LOGO for the outro, TITLE_TEXT_SPLIT for slide 2, "
        "HOW_IT_WORKS_3STEP for slide 3 and 2 slides of THREE_PICTURE_CARDS."
    )
    result = check_brief_archetypes(brief)

    # All six now resolve: gallery 1 plus TITLE_TEXT_SPLIT are parsed out
    # of the gallery's own markup and merged into the engine registry, so
    # the names people write briefs against are the names it can build.
    assert {r["requested"] for r in result["exact"]} == {
        "HOW_IT_WORKS_3STEP", "THREE_PICTURE_CARDS",
        "COVER_A_CUT4", "DIVIDER_D", "END_LOGO", "TITLE_TEXT_SPLIT",
    }
    assert result["unknown"] == []
    assert result["master_slide"] == []


def test_ordinary_prose_is_not_mistaken_for_an_archetype_request():
    from deckguard.skill_bridge import check_brief_archetypes

    result = check_brief_archetypes(
        "A deck about our year_end results and the go_to_market plan for next year."
    )
    assert not any(result[k] for k in result)


def test_archetype_names_resolve_case_insensitively():
    from deckguard.skill_bridge import resolve_archetype_name

    assert resolve_archetype_name("HERO_STAT") == ("hero_stat", "exact")
    assert resolve_archetype_name("hero_stat") == ("hero_stat", "exact")
    assert resolve_archetype_name("Divider_D") == ("divider_d", "exact")
    assert resolve_archetype_name("END_LOGO") == ("end_logo", "exact")
    # a gallery-2 name with no engine equivalent still reports honestly
    assert resolve_archetype_name("REPORT_8CELL")[1] == "unknown"
    assert resolve_archetype_name("NOT_A_REAL_ONE") == (None, "unknown")
