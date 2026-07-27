import hashlib

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from deckguard import logo as logo_mod
from deckguard.config import load_config, default_config_path
from deckguard.fixer import fix_deck
from deckguard.inventory import build_inventory
from deckguard.rules_engine import audit_deck
from tests.helpers import (
    add_picture_to_container,
    add_rectangle,
    add_shadow_effect,
    add_slide,
    body_run,
    make_pattern_png,
    new_deck,
    set_background_image,
    set_run,
    set_theme_font,
    set_theme_slot,
    title_run,
)
from pptx.oxml.ns import qn

CONFIG = load_config(default_config_path())


def _violating_deck():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Calibri", color_hex="005EB8")
    run = set_run(body_run(slide), text="body copy", font="Inter", color_hex="141414")
    add_shadow_effect(run)
    slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_theme_slot(prs, "accent1", "005EB8")
    set_theme_font(prs, "majorFont", "Calibri")
    set_theme_font(prs, "minorFont", "Calibri")
    return prs


def test_fix_applies_changes_and_reduces_violations(tmp_path):
    prs = _violating_deck()
    before_count = len(audit_deck(build_inventory(prs), CONFIG))

    out_path = tmp_path / "out.pptx"
    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=str(out_path), dry_run=False)

    assert len(report.changes) > 0
    assert out_path.exists()

    saved = Presentation(str(out_path))
    after_count = len(audit_deck(build_inventory(saved), CONFIG))
    assert after_count < before_count


def test_fix_dry_run_writes_no_output_file(tmp_path):
    prs = _violating_deck()
    out_path = tmp_path / "out.pptx"

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=str(out_path), dry_run=True)

    assert report.dry_run is True
    assert report.output_path is None
    assert not out_path.exists()
    # dry-run still computes what *would* change, for the preview/changelog
    assert len(report.changes) > 0


def test_fix_never_touches_the_source_file(tmp_path):
    src_path = tmp_path / "source.pptx"
    _violating_deck().save(str(src_path))
    original_bytes = src_path.read_bytes()
    original_hash = hashlib.sha256(original_bytes).hexdigest()

    prs = Presentation(str(src_path))
    out_path = tmp_path / "fixed.pptx"
    fix_deck(prs, CONFIG, source_path=str(src_path), output_path=str(out_path), dry_run=False)

    assert src_path.read_bytes() == original_bytes
    assert hashlib.sha256(src_path.read_bytes()).hexdigest() == original_hash
    assert out_path.exists()
    assert out_path.read_bytes() != original_bytes


def test_fix_is_idempotent_on_a_clean_deck():
    prs = new_deck()
    # python-pptx's default template theme is itself legacy (black
    # dk1, Calibri major/minor font) — bring it up to brand first so this
    # test reflects a deck that's *already* on-brand, not the stock template.
    set_theme_slot(prs, "dk1", "141414")
    set_theme_font(prs, "majorFont", "Inter")
    set_theme_font(prs, "minorFont", "Inter")

    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Inter Semi Bold", color_hex="141414")
    set_run(body_run(slide), text="Body copy", font="Inter", color_hex="141414")
    slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # The stock template's own DATE/SLIDE_NUMBER layout placeholders don't
    # match the org template's position either -- bring those up to brand
    # too, same reasoning as the theme/font lines above.
    from deckguard.fixer import _reference_placeholder_geometry
    from deckguard.slide_import import default_template_path

    layout = slide.slide_layout
    containers = [layout.placeholders, layout.slide_master.shapes]
    for container in containers:
        for ph in container:
            if not getattr(ph, "is_placeholder", False):
                continue
            ph_type = ph.placeholder_format.type
            if ph_type is not None and ph_type.name in ("DATE", "SLIDE_NUMBER"):
                geom = _reference_placeholder_geometry(default_template_path(), ph_type.name)
                if geom:
                    ph.left, ph.top, ph.width, ph.height = geom

    report = fix_deck(prs, CONFIG, source_path="clean.pptx", output_path=None, dry_run=True)
    # Nothing to fix except the (unfixable) slide-size mismatch and
    # non-KONE layout names inherent to the default python-pptx template.
    assert report.changes == []
    assert all(v.rule in ("slide_size", "non_standard_layout") for v in report.manual_review)


def test_fix_report_summary_counts_match():
    prs = _violating_deck()
    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)
    summary = report.summary
    assert summary["changes_applied"] == len(report.changes)
    assert summary["manual_review_required"] == len(report.manual_review)


def _panel_config(min_area_sq_in=8.0):
    """Minimal config isolating just the layout-panel-fill remap, so these
    tests aren't coupled to the real brand_rules.yaml's exact thresholds."""
    return {
        "colors": {
            "approved": ["#1450F5", "#EDEFF0"],
            "remap": {},
            "layout_panel_remap": {"#EDEFF0": "#1450F5"},
            "layout_panel_min_area_sq_in": min_area_sq_in,
        },
        "fonts": {"approved": ["Inter"], "remap": {}},
        "typography_rules": {},
        "logo": {},
        "layout": {},
        "audit": {"fail_on": []},
    }


def _add_rectangle_to_layout(layout, name, fill_hex, left_in, top_in, width_in, height_in):
    """python-pptx's LayoutShapes has no add_shape() -- layout shapes are
    normally inherited from a template, not built via the API. Build the
    rectangle on a scratch slide (where add_shape works) and move its XML
    element into the layout's shape tree instead."""
    from copy import deepcopy

    scratch = Presentation()
    scratch_slide = scratch.slides.add_slide(scratch.slide_layouts[6])
    shape = add_rectangle(
        scratch_slide, name=name, fill_hex=fill_hex,
        left_in=left_in, top_in=top_in, width_in=width_in, height_in=height_in,
    )
    elem = deepcopy(shape._element)
    layout.shapes._spTree.append(elem)
    return layout.shapes[-1]


def _fill_hex(shape):
    spPr = shape._element.find(qn("p:spPr"))
    solid_fill = spPr.find(qn("a:solidFill"))
    srgb = solid_fill.find(qn("a:srgbClr"))
    return srgb.get("val")


def test_fix_remaps_large_background_panel_defined_on_a_layout():
    """Regression test for a real bug: a shape's fill can live entirely on
    the slide LAYOUT it inherits from (e.g. a full-panel background behind
    a picture/content placeholder), never on any slide -- so ordinary
    slide-level color remap never sees or fixes it. This was the actual
    root cause behind a real deck's "grey box that should be blue" report:
    the legacy deck's layout still had the panel as the old gray, while
    the reference deck's equivalent layout had it corrected to KONE Blue."""
    prs = new_deck()
    layout = prs.slide_layouts[1]
    panel = _add_rectangle_to_layout(layout, "Big Panel", "EDEFF0", left_in=1, top_in=1, width_in=4, height_in=3)
    add_slide(prs, layout_idx=1)

    report = fix_deck(prs, _panel_config(), source_path="in.pptx", output_path=None, dry_run=True)

    assert _fill_hex(panel) == "1450F5"
    assert any(c.scope == "layout" and c.shape_name == "Big Panel" for c in report.changes)


def test_fix_leaves_small_layout_shapes_untouched():
    """The exact same literal color is legitimately reused on a layout for
    small accent/placeholder shapes (e.g. unselected tab backgrounds)
    whose color must stay static -- only large background panels should
    ever be treated as the "brand panel" and remapped."""
    prs = new_deck()
    layout = prs.slide_layouts[1]
    small = _add_rectangle_to_layout(layout, "Tab", "EDEFF0", left_in=1, top_in=1, width_in=1, height_in=0.5)
    add_slide(prs, layout_idx=1)

    fix_deck(prs, _panel_config(min_area_sq_in=8.0), source_path="in.pptx", output_path=None, dry_run=True)

    assert _fill_hex(small) == "EDEFF0"


def test_fix_layout_panel_remap_is_a_no_op_when_not_configured():
    prs = new_deck()
    layout = prs.slide_layouts[1]
    panel = _add_rectangle_to_layout(layout, "Big Panel", "EDEFF0", left_in=1, top_in=1, width_in=4, height_in=3)
    add_slide(prs, layout_idx=1)

    config = _panel_config()
    config["colors"]["layout_panel_remap"] = {}

    fix_deck(prs, config, source_path="in.pptx", output_path=None, dry_run=True)

    assert _fill_hex(panel) == "EDEFF0"


def test_fix_corrects_text_color_for_legibility_on_kone_blue():
    """End-to-end: fix_deck must rewrite black-on-blue text to white, not
    leave it for manual review -- this is auto-fixable by construction."""
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Blue panel", fill_hex="1450F5", left_in=1, top_in=1, width_in=3, height_in=1)
    box.text_frame.text = "Label"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    from pptx.dml.color import RGBColor

    run.font.color.rgb = RGBColor.from_string("141414")

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "text_contrast" for c in report.changes)
    assert str(run.font.color.rgb) == "FFFFFF"


def test_fix_sets_explicit_color_for_legibility_when_run_has_no_color_at_all():
    """End-to-end regression test for the real bug report: text with no
    explicit color on a KONE Blue panel must come out of fix() explicitly
    white, not stay silently inherited (and therefore black)."""
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Blue panel", fill_hex="1450F5", left_in=1, top_in=1, width_in=3, height_in=1)
    box.text_frame.text = "Label"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    # no run.font.color set at all -- inherits

    fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert str(run.font.color.rgb) == "FFFFFF"


def test_fix_converts_a_grey_scheme_tint_panel_fill_to_literal_f3eee6():
    """Regression test for a real bug report: a deck's big content panel
    used <a:schemeClr val="bg1"><a:lumMod val="85000"/></a:schemeClr> --
    a theme-relative tint, not a literal color -- which resolved to grey
    (#D9D9D9) but was never fixed, because legacy_color's auto-fix used
    to unconditionally skip any non-literal-RGB color to avoid guessing.
    That skip doesn't apply here: colors.approved keeps white (bg1)
    correct everywhere else it's used un-tinted, so the theme can't be
    globally remapped -- this specific tinted shape needs its own fix,
    and unlisted_panel_fallback's target is fully deterministic, not a guess."""
    from lxml import etree

    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Panel", left_in=1, top_in=1, width_in=3, height_in=2)

    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    spPr = box._element.spPr
    solid_fill = etree.SubElement(spPr, f"{{{A_NS}}}solidFill")
    scheme_clr = etree.SubElement(solid_fill, f"{{{A_NS}}}schemeClr")
    scheme_clr.set("val", "bg1")
    etree.SubElement(scheme_clr, f"{{{A_NS}}}lumMod").set("val", "85000")

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "legacy_color" and c.shape_name == "Panel" for c in report.changes)
    assert str(box.fill.fore_color.rgb) == "F3EEE6"


def _logo_config(old_hash, new_logo_path):
    """Minimal config isolating just old-logo replacement, so these tests
    aren't coupled to the real brand_rules.yaml's own (currently empty,
    per a separate finding) old_logo_hashes list."""
    return {
        "colors": {"approved": ["#1450F5"], "remap": {}},
        "fonts": {"approved": ["Inter"], "remap": {}},
        "typography_rules": {},
        "logo": {"old_logo_hashes": [old_hash], "new_logo_path": str(new_logo_path)},
        "layout": {},
        "audit": {"fail_on": []},
    }


def test_fix_replaces_old_logo_placed_directly_on_a_slide_layout(tmp_path):
    """Regression test for a real report: a logo placed on a slide LAYOUT
    (inherited by every slide that uses it, a common way to put a logo on
    every slide at once) rather than any individual slide is invisible to
    ordinary slide-level scanning (inventory.slides[].shapes never visits
    layout shapes) -- so it was never detected or fixed."""
    old_logo = make_pattern_png(tmp_path / "old.png", seed=1)
    new_logo = make_pattern_png(tmp_path / "new.png", seed=9)
    old_hash = logo_mod.compute_phash(old_logo.read_bytes())

    prs = new_deck()
    layout = prs.slide_layouts[1]
    add_picture_to_container(layout, str(old_logo), name="Logo")
    add_slide(prs, layout_idx=1)

    report = fix_deck(prs, _logo_config(old_hash, new_logo), source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "old_logo" and c.scope == "layout" for c in report.changes)
    pic = next(s for s in layout.shapes if s.name == "Logo")
    assert pic.image.blob == new_logo.read_bytes()


def test_fix_replaces_old_logo_placed_directly_on_a_slide_master(tmp_path):
    old_logo = make_pattern_png(tmp_path / "old.png", seed=1)
    new_logo = make_pattern_png(tmp_path / "new.png", seed=9)
    old_hash = logo_mod.compute_phash(old_logo.read_bytes())

    prs = new_deck()
    master = prs.slide_masters[0]
    add_picture_to_container(master, str(old_logo), name="Logo")
    add_slide(prs)

    report = fix_deck(prs, _logo_config(old_hash, new_logo), source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "old_logo" and c.scope == "master" for c in report.changes)
    pic = next(s for s in master.shapes if s.name == "Logo")
    assert pic.image.blob == new_logo.read_bytes()


def test_fix_replaces_a_vector_logo_mark_via_old_logo_region_in(tmp_path):
    """Regression test for a real report: an old logo that isn't a raster
    image at all -- a wordmark built from freeform vector shapes directly
    on a slide master -- can never be found by old_logo_hashes (nothing
    to perceptual-hash). old_logo_region_in identifies it by position
    instead."""
    from pptx.util import Inches

    new_logo = make_pattern_png(tmp_path / "new.png", seed=9)

    prs = new_deck()
    master = prs.slide_masters[0]
    spTree = master.shapes._spTree
    id_ = master.shapes._next_shape_id
    spTree.add_textbox(id_, "OldVectorMark", Inches(11), Inches(0.2), Inches(1), Inches(0.5))
    add_slide(prs)

    config = {
        "colors": {"approved": ["#1450F5"], "remap": {}},
        "fonts": {"approved": ["Inter"], "remap": {}},
        "typography_rules": {},
        "logo": {"old_logo_region_in": [10.5, 0.0, 2.8, 1.2], "new_logo_path": str(new_logo)},
        "layout": {},
        "audit": {"fail_on": []},
    }

    report = fix_deck(prs, config, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "old_logo_region" and c.scope == "master" for c in report.changes)
    assert not any(s.name == "OldVectorMark" for s in master.shapes)
    pictures = [s for s in master.shapes if s.shape_type is not None and s.shape_type.name == "PICTURE"]
    assert len(pictures) == 1
    assert pictures[0].image.blob == new_logo.read_bytes()


def test_fix_sizes_the_replacement_logo_from_the_template_not_the_search_region(tmp_path):
    """Regression test for a real report: the replacement logo rendered
    'super huge' because it was sized to fill old_logo_region_in (drawn
    deliberately generous, to reliably catch the old mark regardless of
    its own size) instead of the org template's own actual logo size.
    The configured region here is intentionally much bigger than the
    template's real logo -- the output picture must still come out at
    the template's own size, not the region's."""
    from pptx.util import Inches

    from deckguard import logo as logo_mod
    from deckguard.slide_import import default_template_path

    if not default_template_path().exists():
        pytest.skip("bundled template asset not present")

    new_logo = make_pattern_png(tmp_path / "new.png", seed=11)

    prs = new_deck()
    master = prs.slide_masters[0]
    spTree = master.shapes._spTree
    id_ = logo_mod._next_shape_id_in_tree(spTree)
    spTree.add_textbox(id_, "OldVectorMark", Inches(11), Inches(0.2), Inches(1), Inches(0.5))
    add_slide(prs)

    config = {
        "colors": {"approved": ["#1450F5"], "remap": {}},
        "fonts": {"approved": ["Inter"], "remap": {}},
        "typography_rules": {},
        # deliberately much larger than the template's real ~0.85 x 0.33in logo
        "logo": {"old_logo_region_in": [9.0, 0.0, 4.0, 3.0], "new_logo_path": str(new_logo)},
        "layout": {},
        "audit": {"fail_on": []},
    }

    fix_deck(prs, config, source_path="in.pptx", output_path=None, dry_run=True)

    pic = next(s for s in master.shapes if s.shape_type is not None and s.shape_type.name == "PICTURE")
    expected = logo_mod.reference_logo_geometry(default_template_path())
    assert expected is not None
    _e_left, _e_top, e_width, e_height = expected
    # aspect-fit within the template's own box -- at least one axis should
    # land close to it (not blown up to fill the much bigger search region)
    assert pic.width <= e_width + 1000  # a few EMU of rounding slack
    assert pic.height <= e_height + 1000


def test_fix_old_logo_region_in_unset_is_a_no_op():
    """Unset must never touch a master -- deleting shapes by position
    alone is only safe once a human has confirmed the region, so no
    region configured means no shapes are ever removed."""
    from pptx.util import Inches

    prs = new_deck()
    master = prs.slide_masters[0]
    spTree = master.shapes._spTree
    id_ = master.shapes._next_shape_id
    spTree.add_textbox(id_, "SomeShape", Inches(11.9), Inches(0.0), Inches(1), Inches(0.5))
    add_slide(prs)

    config = {
        "colors": {"approved": ["#1450F5"], "remap": {}},
        "fonts": {"approved": ["Inter"], "remap": {}},
        "typography_rules": {},
        "logo": {"old_logo_region_in": None, "new_logo_path": str(default_config_path().parent / "assets" / "kone_logo.png")},
        "layout": {},
        "audit": {"fail_on": []},
    }

    report = fix_deck(prs, config, source_path="in.pptx", output_path=None, dry_run=True)

    assert not any(c.rule == "old_logo_region" for c in report.changes)
    assert any(s.name == "SomeShape" for s in master.shapes)  # left untouched


def test_default_config_replaces_the_known_old_logo_region(tmp_path):
    """The shipped brand_rules.yaml now ships old_logo_region_in with a
    region confirmed against a real legacy KONE deck (see logo.py's
    module comment) -- exercise it end to end against the same shape
    layout that deck's master actually has, using the real default
    config rather than a hand-built one."""
    from pptx.util import Inches

    prs = new_deck()
    master = prs.slide_masters[0]
    spTree = master.shapes._spTree
    id_ = master.shapes._next_shape_id
    spTree.add_textbox(id_, "OldWordmark", Inches(12.0), Inches(0.1), Inches(1), Inches(0.5))
    add_slide(prs)

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "old_logo_region" and c.scope == "master" for c in report.changes)
    assert not any(s.name == "OldWordmark" for s in master.shapes)


def _duplicate_a_shape_id(container_element):
    """Simulate a real, confirmed defect: an old deck with the same
    shape id used twice within one part (e.g. an embedded object
    duplicated at some point in its edit history, both copies keeping
    the original id). Deep-copies the first shape element found and
    re-appends it with the SAME id -- python-pptx has no API for this
    (it always assigns a fresh id), so it's built directly via lxml,
    exactly the way a real malformed file would already have it."""
    import copy

    from pptx.oxml.ns import qn

    first_sp = container_element.find(f".//{qn('p:sp')}")
    clone = copy.deepcopy(first_sp)
    container_element.find(f".//{qn('p:spTree')}").append(clone)
    return clone.find(f".//{qn('p:cNvPr')}").get("id")


def test_fix_deck_renumbers_a_duplicate_shape_id_on_a_slide():
    from collections import Counter

    from pptx.oxml.ns import qn

    prs = new_deck()
    slide = add_slide(prs)
    dup_id = _duplicate_a_shape_id(slide._element)

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "duplicate_shape_id" and c.scope == "slide" for c in report.changes)
    ids = [el.get("id") for el in slide._element.findall(f".//{qn('p:cNvPr')}")]
    assert len(ids) == len(set(ids))  # all unique now
    assert ids.count(dup_id) == 1  # the original kept its id; only the clone moved


def test_fix_deck_renumbers_a_duplicate_shape_id_on_a_master():
    from pptx.oxml.ns import qn

    prs = new_deck()
    master = prs.slide_masters[0]
    _duplicate_a_shape_id(master._element)
    add_slide(prs)

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "duplicate_shape_id" and c.scope == "master" for c in report.changes)
    ids = [el.get("id") for el in master._element.findall(f".//{qn('p:cNvPr')}")]
    assert len(ids) == len(set(ids))


def test_fix_deck_leaves_unique_shape_ids_untouched():
    prs = new_deck()
    add_slide(prs)

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert not any(c.rule == "duplicate_shape_id" for c in report.changes)


def test_fix_replaces_old_logo_baked_into_a_slide_background_fill(tmp_path):
    """A logo can also be baked into a page-level background-FILL image
    (<p:cSld><p:bg>) rather than being a picture shape at all -- outside
    the shape tree entirely, so even a full recursive shape scan can't
    see it. This is the "hard-coded, not an element" case."""
    old_logo = make_pattern_png(tmp_path / "old.png", seed=1)
    new_logo = make_pattern_png(tmp_path / "new.png", seed=9)
    old_hash = logo_mod.compute_phash(old_logo.read_bytes())

    prs = new_deck()
    slide = add_slide(prs)
    blip = set_background_image(slide, str(old_logo))

    report = fix_deck(prs, _logo_config(old_hash, new_logo), source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "old_logo" and c.field == "background image" and c.scope == "slide" for c in report.changes)
    rid = blip.get(f"{{{logo_mod.R_NS}}}embed")
    assert slide.part.related_part(rid).blob == new_logo.read_bytes()


def test_fix_replaces_old_logo_baked_into_a_layout_background_fill(tmp_path):
    old_logo = make_pattern_png(tmp_path / "old.png", seed=1)
    new_logo = make_pattern_png(tmp_path / "new.png", seed=9)
    old_hash = logo_mod.compute_phash(old_logo.read_bytes())

    prs = new_deck()
    layout = prs.slide_layouts[1]
    blip = set_background_image(layout, str(old_logo))
    add_slide(prs, layout_idx=1)

    report = fix_deck(prs, _logo_config(old_hash, new_logo), source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "old_logo" and c.field == "background image" and c.scope == "layout" for c in report.changes)
    rid = blip.get(f"{{{logo_mod.R_NS}}}embed")
    assert layout.part.related_part(rid).blob == new_logo.read_bytes()


def test_fix_old_logo_replacement_is_a_no_op_without_configured_hashes():
    """Documents the OTHER real gap: even a logo sitting in plain sight as
    a normal slide picture is never touched if logo.old_logo_hashes is
    empty (nothing to match against) -- config population, not detection
    scope, is the remaining blocker for that case."""
    prs = new_deck()
    slide = add_slide(prs)

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert not any(c.rule == "old_logo" for c in report.changes)


def _add_layout_placeholder_to_slide(slide, ph_type_name: str):
    """Clone a DATE/FOOTER/SLIDE_NUMBER placeholder from the slide's own
    layout onto the slide, the way PowerPoint's Insert Header & Footer
    does -- python-pptx's add_slide() doesn't clone these by default
    (only content placeholders), so a test needs to add one explicitly
    to reproduce a real deck's footer chrome."""
    import copy

    layout = slide.slide_layout
    src = next(
        ph for ph in layout.placeholders
        if ph.placeholder_format.type is not None and ph.placeholder_format.type.name == ph_type_name
    )
    clone = copy.deepcopy(src._element)
    slide.shapes._spTree.append(clone)
    return next(
        ph for ph in slide.placeholders
        if ph.placeholder_format.type is not None and ph.placeholder_format.type.name == ph_type_name
    )


def test_fix_forces_brand_font_and_color_onto_unstyled_footer_text():
    """Regression test for a real report: footer/date/slide-number text
    on a slide kept from an old deck (brand mode's own 'skipped slide'
    path) stayed on whatever font/color the OLD master defined --
    confirmed on a real deck: Arial, no color override -- because the
    per-violation audit can't flag a run with no explicit color AND no
    resolvable background to compute contrast against (footer text sits
    directly on the page canvas, not a fill)."""
    prs = new_deck()
    slide = add_slide(prs)
    footer_ph = _add_layout_placeholder_to_slide(slide, "FOOTER")
    footer_ph.text_frame.text = "Some Corp Internal"
    run = footer_ph.text_frame.paragraphs[0].runs[0]
    assert run.font.name is None  # sanity: no explicit override to start with
    assert run.font.color.type is None

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "footer_chrome_default" for c in report.changes)
    run = footer_ph.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Inter"
    assert str(run.font.color.rgb) == "141414"


def test_fix_never_overrides_a_footer_run_with_an_explicit_color_already():
    """Font isn't a fair test of "leave a deliberate choice alone" --
    every non-approved font gets corrected to Inter regardless (that's
    the whole point of the tool), so an already-Calibri run would be
    fixed by the ordinary font-compliance pass either way. Color is
    the real test: white is an equally brand-approved choice to black
    for footer text with no resolvable background (see
    typography_rules.text_colors' "Inter": [black, white] list), so an
    explicit white must survive, not get silently forced to the "black
    first" default."""
    prs = new_deck()
    slide = add_slide(prs)
    footer_ph = _add_layout_placeholder_to_slide(slide, "FOOTER")
    footer_ph.text_frame.text = "Some Corp Internal"
    run = footer_ph.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter SemiBold"  # already approved -- the general font pass has nothing to fix
    from pptx.dml.color import RGBColor

    run.font.color.rgb = RGBColor.from_string("FFFFFF")

    fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    run = footer_ph.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Inter SemiBold"
    assert str(run.font.color.rgb) == "FFFFFF"


def test_fix_forces_brand_font_and_color_onto_date_and_slide_number_fields():
    """Date/slide-number placeholders are PowerPoint auto-fields
    (<a:fld>), not <a:r> runs -- python-pptx's own paragraph.runs never
    returns them, so they need their own handling entirely separate
    from ordinary text runs."""
    prs = new_deck()
    slide = add_slide(prs)
    date_ph = _add_layout_placeholder_to_slide(slide, "DATE")
    from pptx.oxml.ns import qn

    fld = etree.SubElement(date_ph.text_frame.paragraphs[0]._p, qn("a:fld"))
    fld.set("id", "{00000000-0000-0000-0000-000000000000}")
    fld.set("type", "datetime1")
    rPr = etree.SubElement(fld, qn("a:rPr"))
    rPr.set("lang", "en-US")
    t = etree.SubElement(fld, qn("a:t"))
    t.text = "1 January 2026"

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "footer_chrome_default" and c.shape_name == date_ph.name for c in report.changes)
    latin = fld.find(qn("a:rPr")).find(qn("a:latin"))
    solid_fill = fld.find(qn("a:rPr")).find(qn("a:solidFill"))
    assert latin.get("typeface") == "Inter"
    assert solid_fill.find(qn("a:srgbClr")).get("val") == "141414"


def test_fix_footer_chrome_leaves_ordinary_body_placeholders_alone():
    prs = new_deck()
    slide = add_slide(prs)
    body_run(slide).text = "Ordinary body text"

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert not any(c.rule == "footer_chrome_default" for c in report.changes)


def test_fix_repositions_date_and_slide_number_to_brand_position():
    """Regression test for a real report: an old deck's own layout had
    slide-number at the LEFT edge and its confidentiality/footer text
    at the RIGHT -- backwards from the org template's own convention
    (date left, slide number right, confirmed by direct inspection)."""
    from deckguard.fixer import _reference_placeholder_geometry
    from deckguard.slide_import import default_template_path

    if not default_template_path().exists():
        pytest.skip("bundled template asset not present")

    prs = new_deck()
    slide = add_slide(prs)
    layout = slide.slide_layout
    date_ph = next(ph for ph in layout.placeholders if ph.placeholder_format.type.name == "DATE")
    slidenum_ph = next(ph for ph in layout.placeholders if ph.placeholder_format.type.name == "SLIDE_NUMBER")
    # deliberately swapped from brand -- slide number on the left, date on the right
    date_ph.left, slidenum_ph.left = Inches(9), Inches(0.2)

    fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    expected_date = _reference_placeholder_geometry(default_template_path(), "DATE")
    expected_slidenum = _reference_placeholder_geometry(default_template_path(), "SLIDE_NUMBER")
    assert (date_ph.left, date_ph.top, date_ph.width, date_ph.height) == expected_date
    assert (slidenum_ph.left, slidenum_ph.top, slidenum_ph.width, slidenum_ph.height) == expected_slidenum
    assert date_ph.left < slidenum_ph.left  # date ends up left of slide number, as brand requires


def test_fix_removes_confidentiality_footer_text():
    prs = new_deck()
    slide = add_slide(prs)
    footer_ph = _add_layout_placeholder_to_slide(slide, "FOOTER")
    footer_ph.text_frame.text = "Confidential  |  © KONE Corporation"

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert any(c.rule == "confidentiality_footer_removed" for c in report.changes)
    assert footer_ph.text_frame.text == ""


def test_fix_leaves_an_ordinary_footer_line_alone():
    prs = new_deck()
    slide = add_slide(prs)
    footer_ph = _add_layout_placeholder_to_slide(slide, "FOOTER")
    footer_ph.text_frame.text = "Q3 2026 Investor Update"

    report = fix_deck(prs, CONFIG, source_path="in.pptx", output_path=None, dry_run=True)

    assert not any(c.rule == "confidentiality_footer_removed" for c in report.changes)
    assert footer_ph.text_frame.text == "Q3 2026 Investor Update"
