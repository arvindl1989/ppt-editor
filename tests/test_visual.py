"""Tests for the visual audit (visual.py).

Two layers, deliberately: the CHECKS run on synthetic measurement
records and always execute, so the rules are pinned regardless of what
is installed; the MEASUREMENT tests need a real headless Chromium and
skip cleanly without one. That split is the point -- the browser is a
measuring instrument, not the thing under test.
"""

import pytest

from deckguard.visual import (
    MIN_LEGIBLE_PT,
    PT_TO_PX,
    RENDER_WIDTH_PX,
    VisualReport,
    audit_previews,
    check_measurements,
    measure_html,
    playwright_available,
    to_json,
)

needs_browser = pytest.mark.skipif(
    not playwright_available(), reason="Playwright + Chromium not available here"
)


def _frame(shapes, width=RENDER_WIDTH_PX, height=720, background="rgb(255, 255, 255)"):
    return {"width": width, "height": height, "background": background, "shapes": shapes}


def _shape(name="Shape 1", left=0, top=0, width=400, height=200, runs=(),
           overflowX=0, overflowY=0, background="rgba(0, 0, 0, 0)"):
    return {
        "name": name, "left": left, "top": top, "width": width, "height": height,
        "overflowX": overflowX, "overflowY": overflowY, "background": background,
        "runs": list(runs),
    }


def _run(text="Hello", pt=24.0, color="rgb(20, 20, 20)"):
    return {"text": text, "fontPx": pt * PT_TO_PX, "color": color}


# --------------------------------------------------------------------------
# the checks
# --------------------------------------------------------------------------


def test_a_clean_slide_produces_no_findings():
    findings = check_measurements([_frame([_shape(runs=[_run()])])])
    assert findings == []


def test_text_taller_than_its_box_is_reported_as_overflow():
    findings = check_measurements([_frame([_shape(overflowY=41, runs=[_run()])])])
    assert [f.rule for f in findings] == ["text_overflow"]
    assert findings[0].severity == "major"
    assert "41px taller" in findings[0].message


def test_a_shape_hanging_off_the_slide_is_reported():
    """Verified against a real deck: OLD_TEST slide 5 has a text box
    whose right edge is at 13.444in on a 13.333in slide."""
    off_right = _shape(name="TextBox 17", left=RENDER_WIDTH_PX - 50, width=200, runs=[_run()])
    findings = check_measurements([_frame([off_right])])
    assert [f.rule for f in findings] == ["outside_frame"]
    assert findings[0].shape_name == "TextBox 17"


def test_a_shape_flush_to_the_edge_is_not_reported():
    flush = _shape(left=RENDER_WIDTH_PX - 200, width=200, runs=[_run()])
    assert check_measurements([_frame([flush])]) == []


def test_text_below_the_legible_floor_is_reported():
    findings = check_measurements([_frame([_shape(runs=[_run(pt=5.0)])])])
    assert [f.rule for f in findings] == ["tiny_text"]
    assert f"{MIN_LEGIBLE_PT:.0f}pt" in findings[0].message

    ok = check_measurements([_frame([_shape(runs=[_run(pt=MIN_LEGIBLE_PT + 0.5)])])])
    assert ok == []


def test_low_contrast_text_is_reported_against_its_own_background():
    """The grey-on-grey pills in the real OLD_TEST deck: #BABDBF text on
    an #EDEFF0 fill is 1.6:1 -- legal brand colors, unreadable slide.
    Structural rules pass it; only measuring the pair catches it."""
    pill = _shape(
        name="Rectangle 17", background="rgb(237, 239, 240)",
        runs=[_run("Code", color="rgb(186, 189, 191)")],
    )
    findings = check_measurements([_frame([pill])])
    assert [f.rule for f in findings] == ["low_contrast"]
    assert "1.6:1" in findings[0].message


def test_contrast_falls_back_to_the_frame_background_when_the_shape_is_transparent():
    """A transparent shape shows the slide behind it -- white text on a
    white slide has to be caught, not skipped."""
    ghost = _shape(background="rgba(0, 0, 0, 0)", runs=[_run(color="rgb(255, 255, 255)")])
    findings = check_measurements([_frame([ghost], background="rgb(255, 255, 255)")])
    assert [f.rule for f in findings] == ["low_contrast"]


def test_white_on_kone_blue_passes():
    blue = _shape(background="rgb(20, 80, 245)", runs=[_run(color="rgb(255, 255, 255)")])
    assert check_measurements([_frame([blue])]) == []


def test_a_frame_that_drew_nothing_is_reported():
    findings = check_measurements([_frame([])])
    assert [f.rule for f in findings] == ["empty_frame"]


def test_findings_carry_the_slide_they_came_from():
    frames = [_frame([_shape(runs=[_run()])]), _frame([_shape(name="B", overflowY=20, runs=[_run()])])]
    findings = check_measurements(frames)
    assert len(findings) == 1
    assert findings[0].frame_index == 1  # rendered 1-based by the CLI/JSON


def test_report_json_is_serialisable_and_summarised():
    import json

    report = VisualReport(findings=check_measurements([_frame([_shape(overflowY=20, runs=[_run()])])]),
                          frames_measured=1, ran=True)
    payload = json.loads(to_json(report))
    assert payload["ran"] is True
    assert payload["summary"]["major"] == 1
    assert payload["findings"][0]["slide"] == 1


# --------------------------------------------------------------------------
# the measurement (needs a browser)
# --------------------------------------------------------------------------


@needs_browser
def test_measurement_reports_type_at_true_slide_scale():
    """The regression that motivated this module: a 60pt run must render
    at 60pt-equivalent px in a slide-width frame. It was landing ~2.4x
    oversized because `cqh` resolved against the viewport."""
    from deckguard.preview import _box

    frag = (
        '<div data-dg-frame="1" style="position:relative;aspect-ratio:1280/720;'
        'container-type:inline-size;background:#FFFFFF;overflow:hidden;width:100%;">'
        + _box(0, 0, 60, 30, '<span style="font-size:6.25cqw;color:#141414;">Ascenseur</span>',
               label="Title 1")
        + "</div>"
    )
    frames = measure_html([frag])

    assert len(frames) == 1
    assert abs(frames[0]["width"] - RENDER_WIDTH_PX) < 2
    run = frames[0]["shapes"][0]["runs"][0]
    assert abs(run["fontPx"] / PT_TO_PX - 60.0) < 1.0


@needs_browser
def test_measurement_catches_real_overflow_in_a_narrow_box():
    from deckguard.preview import _box

    frag = (
        '<div data-dg-frame="1" style="position:relative;aspect-ratio:1280/720;'
        'container-type:inline-size;background:#FFFFFF;overflow:hidden;width:100%;">'
        + _box(0, 0, 8, 4, '<span style="font-size:5cqw;">overflowing text here</span>',
               label="Tiny box")
        + "</div>"
    )
    findings = check_measurements(measure_html([frag]))
    assert any(f.rule == "text_overflow" and f.shape_name == "Tiny box" for f in findings)


@needs_browser
def test_a_filled_shape_does_not_render_wider_than_the_shape():
    """Regression: filled boxes carry padding, and under the default
    content-box that padding was ADDED to the declared width, so a
    full-width filled shape drew past the slide edge. Caught by this
    module reporting off-slide shapes on a deck that has none."""
    from pptx import Presentation

    from deckguard.inventory import build_inventory
    from deckguard.preview import slide_preview_html
    from tests.helpers import add_rectangle, add_slide, new_deck

    import pathlib
    import tempfile

    with tempfile.TemporaryDirectory() as td:
        prs = new_deck()
        s = add_slide(prs)
        # a band spanning the full slide width, exactly to the edges
        add_rectangle(s, name="Band", fill_hex="1450F5", left_in=0, top_in=3,
                      width_in=13.333, height_in=1)
        path = pathlib.Path(td) / "d.pptx"
        prs.save(str(path))
        inv = build_inventory(Presentation(str(path)))

    findings = check_measurements(measure_html([slide_preview_html(inv.slides[0])]))
    assert not [f for f in findings if f.rule == "outside_frame"], (
        "a shape flush with the slide edge must not render past it"
    )


@needs_browser
def test_audit_previews_end_to_end():
    from deckguard.preview import slide_preview_html
    from pptx import Presentation

    from deckguard.inventory import build_inventory
    from tests.helpers import add_slide, new_deck, set_run, title_run

    import pathlib
    import tempfile

    with tempfile.TemporaryDirectory() as td:
        prs = new_deck()
        s = add_slide(prs)
        set_run(title_run(s), text="Readable title", font="Inter", color_hex="141414")
        path = pathlib.Path(td) / "d.pptx"
        prs.save(str(path))
        inv = build_inventory(Presentation(str(path)))

    report = audit_previews([slide_preview_html(rec) for rec in inv.slides])
    assert report.ran and report.frames_measured == 1
    assert not [f for f in report.findings if f.severity == "major"]


def test_audit_degrades_cleanly_with_no_browser(monkeypatch):
    """No Playwright must mean "not measured", never a crash -- this
    runs inside a transform."""
    import deckguard.visual as visual_mod

    monkeypatch.setattr(visual_mod, "playwright_available", lambda: False)
    report = visual_mod.audit_previews(["<div data-dg-frame='1'></div>"])
    assert report.ran is False and report.findings == []
    assert visual_mod.measure_html(["<div></div>"]) == []


def test_contrast_uses_the_panel_behind_a_transparent_shape():
    """A text placeholder has no fill of its own; if a painted panel
    sits behind it, THAT is what its text must be legible against.
    Measured on a real deck: the AMP layout puts a KONE Blue panel
    behind its right-hand column, so its white body text is correct --
    reading only the shape's own (transparent) fill reported every one
    of those slides as white-on-white."""
    on_blue = _shape(background="rgba(0, 0, 0, 0)", runs=[_run(color="rgb(255, 255, 255)")])
    on_blue["effectiveBg"] = "rgb(20, 80, 245)"
    assert check_measurements([_frame([on_blue])]) == []


@needs_browser
def test_a_layout_colour_panel_is_drawn_behind_the_slide():
    """Regression: previews drew only the slide's own shapes, so a slide
    leaning on a layout panel previewed as text floating on white."""
    from deckguard.preview import slide_preview_html

    class _Colour:
        hex = "1450F5"

    class _Fill:
        type = "solid"
        colors = [_Colour()]

    class _BgShape:
        name = "Rectangle 4"
        left_in, top_in, width_in, height_in = 5.34, 1.23, 7.55, 4.93
        fill = _Fill()

    class _Record:
        index = 1
        shapes: list = []
        layout_background_shapes = [_BgShape()]

    frames = measure_html([slide_preview_html(_Record())])
    painted = [s for s in frames[0]["shapes"] if s["name"] == "Rectangle 4"]
    assert painted, "the layout's colour panel must be drawn"
    assert painted[0]["background"] == "rgb(20, 80, 245)"


@needs_browser
def test_a_box_too_small_for_its_label_is_drawn_as_a_plain_swatch():
    """The chevron marks in a real reference deck are 0.15in wide. The
    preview's own "GROUP" label doesn't fit and reported 50 phantom
    overflows -- preview chrome must never become a finding."""
    from deckguard.preview import _placeholder_box

    frag = (
        '<div data-dg-frame="1" style="position:relative;aspect-ratio:1280/720;'
        'container-type:inline-size;background:#FFFFFF;overflow:hidden;width:100%;">'
        + _placeholder_box(10, 10, 1.1, 1.8, "GROUP", name="Group 31")
        + _placeholder_box(30, 10, 30, 20, "GROUP", name="Group 99")
        + "</div>"
    )
    findings = check_measurements(measure_html([frag]))
    assert not [f for f in findings if f.rule == "text_overflow"]

    frames = measure_html([frag])
    by_name = {s["name"]: s for s in frames[0]["shapes"]}
    assert by_name["Group 31"]["runs"] == []  # too small: swatch only
    assert by_name["Group 99"]["runs"], "a roomy placeholder keeps its label"
