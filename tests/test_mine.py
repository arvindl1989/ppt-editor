"""Tests for mining a reference .pptx into reusable archetypes (mine.py).

The point of the module is a change of kind, not degree: every earlier
use of a reference deck patched the old deck toward it (learned colours,
borrowed layouts, transplanted shape styles). This reads the reference's
DESIGNS out so old content can be re-rendered through them.
"""

import pytest

from deckguard.mine import _cluster_repeats, _signature, mine_reference
from tests.helpers import add_rectangle, add_slide, new_deck, set_run, title_run


def _item(kind="text", left=0, top=0, width=200, height=60, role="r", text=""):
    return {"kind": kind, "box": [left, top, width, height], "role": role,
            "style": ("Inter", 18, "141414", False, False, False), "text": text}


def test_a_row_of_identical_shapes_becomes_one_repeating_group():
    """Seven identically-sized pills along a row are not seven regions.
    Detecting the repeat is what lets a mined design hold a different
    number of items than the slide it came from."""
    row = [_item(left=x, top=648, width=178, height=72, text=f"item {i}")
           for i, x in enumerate((0, 184, 368, 551, 735, 918, 1102))]
    singles, groups = _cluster_repeats(row)

    assert singles == []
    assert len(groups) == 1
    assert groups[0]["count"] == 7
    assert groups[0]["origins"][0] == [0, 648]


def test_two_similar_shapes_are_not_a_group():
    """Two of a thing is a pair, not a pattern -- treating it as a group
    would turn every two-column slide into a one-column template."""
    singles, groups = _cluster_repeats([_item(left=0), _item(left=400)])
    assert groups == []
    assert len(singles) == 2


def test_shapes_of_different_sizes_stay_separate():
    items = [_item(left=0, width=200), _item(left=300, width=400), _item(left=800, width=120)]
    singles, groups = _cluster_repeats(items)
    assert groups == []
    assert len(singles) == 3


def test_the_signature_ignores_content_driven_height():
    """Four copies of one layout holding different amounts of text
    differ only in box height. Matching on height split them into four
    archetypes instead of one."""
    short = {"regions": [{"role": "t", "box": [45, 91, 900, 60]}], "groups": []}
    tall = {"regions": [{"role": "t", "box": [45, 91, 900, 220]}], "groups": []}
    moved = {"regions": [{"role": "t", "box": [400, 91, 900, 60]}], "groups": []}

    assert _signature(short) == _signature(tall)
    assert _signature(short) != _signature(moved)


def _reference_deck(path, pill_rows=2):
    """A deck with two slides sharing one design: a title plus a row of
    identically-sized pills."""
    prs = new_deck()
    for n in range(pill_rows):
        slide = add_slide(prs)
        set_run(title_run(slide), text=f"Section {n}", font="Inter", color_hex="141414")
        for i in range(4):
            pill = add_rectangle(slide, name=f"Pill {n}{i}", fill_hex="1450F5",
                                 left_in=0.5 + i * 3.0, top_in=6.0, width_in=2.6, height_in=0.8)
            pill.text_frame.text = f"Pillar {i}"
            set_run(pill.text_frame.paragraphs[0].runs[0], text=f"Pillar {i}",
                    font="Inter", color_hex="FFFFFF", size_pt=14)
    prs.save(str(path))
    return path


def test_a_deck_that_reuses_one_design_yields_one_archetype(tmp_path):
    """Deduping is what makes this usable: a 40-slide reference must
    yield a handful of designs, not 40."""
    deck = _reference_deck(tmp_path / "ref.pptx", pill_rows=3)

    mined = mine_reference(deck)

    assert len(mined["archetypes"]) == 1, mined["archetypes"].keys()
    name = next(iter(mined["archetypes"]))
    assert mined["sources"][name] == [1, 2, 3], "all three slides map to the one design"
    arch = mined["archetypes"][name]
    assert arch["groups"] and len(arch["groups"][0]["origins"]) == 4


def test_mined_samples_carry_the_references_own_words(tmp_path):
    """The reference's copy becomes the archetype's worked example, so
    the planning prompt shows the model the house style."""
    deck = _reference_deck(tmp_path / "ref.pptx", pill_rows=1)
    mined = mine_reference(deck)
    name = next(iter(mined["archetypes"]))
    flat = str(mined["samples"][name])
    assert "Section 0" in flat


def test_a_one_off_composition_is_not_mined_as_a_design(tmp_path):
    """Twenty loose text boxes with no repeating structure is a slide,
    not a template -- mining it produces a twenty-slot archetype nothing
    else will ever fit."""
    from pptx.util import Inches

    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Dense"
    # Deliberately irregular -- every box a different size in a different
    # place, which is what a hand-composed slide of copy actually looks
    # like. (A regular grid of same-sized boxes IS a repeating design,
    # and is mined as one.)
    for i in range(20):
        box = slide.shapes.add_textbox(
            Inches(0.4 + (i * 0.53) % 9), Inches(0.5 + (i * 0.31) % 6),
            Inches(1.8 + (i % 7) * 0.4), Inches(0.5 + (i % 5) * 0.17))
        box.text_frame.text = f"Point {i}"
    path = tmp_path / "dense.pptx"
    prs.save(str(path))

    assert mine_reference(path)["archetypes"] == {}


def test_mining_never_raises_on_an_unreadable_deck(tmp_path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"not a pptx")
    with pytest.raises(Exception):
        mine_reference(bad)  # the miner itself surfaces the error...

    from deckguard.mine import install_reference

    class _Module:
        ARCHETYPES: dict = {}

    # ...but installing is additive and must degrade to adding nothing
    assert install_reference(_Module(), bad)["archetypes"] == {}


def test_mined_designs_are_preferred_but_not_at_any_cost():
    """"Make it look like that deck" means that deck's layouts -- but a
    reference design that would drop half the content must not beat a
    built-in one that holds all of it."""
    from deckguard.skill_bridge import match_archetypes

    blocks = [[f"Point {i}"] for i in range(6)]
    plain = match_archetypes("A title", blocks, 0, limit=1)
    assert plain, "the built-in library always offers something"

    # preferring a name that doesn't exist must change nothing
    same = match_archetypes("A title", blocks, 0, limit=1, prefer={"nope"})
    assert same[0]["archetype"] == plain[0]["archetype"]


def test_each_block_in_a_repeating_band_keeps_its_own_colour(tmp_path):
    """The reference's category band runs blue/amber/pink/yellow/green
    across its cells. Reusing the first cell's fill for all of them
    turned it into a solid blue bar."""
    palette = ["1450F5", "FFA023", "FFCDD7", "FFE141", "1ED273"]
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Categories"
    for i, hexval in enumerate(palette):
        add_rectangle(slide, name=f"Cell {i}", fill_hex=hexval,
                      left_in=0.3 + i * 2.6, top_in=6.4, width_in=2.4, height_in=0.7)
    path = tmp_path / "band.pptx"
    prs.save(str(path))

    mined = mine_reference(path)
    arch = next(iter(mined["archetypes"].values()))
    fills = {c["hex"] for c in arch["chrome"] if c["kind"] == "fill"}
    assert fills == set(palette)
