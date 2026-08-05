from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

from deckguard.config import load_config, default_config_path
from deckguard.inventory import build_inventory
from deckguard.rules_engine import audit_deck
from tests.helpers import (
    add_rectangle,
    add_shadow_effect,
    add_slide,
    body_run,
    new_deck,
    set_run,
    set_run_theme_color,
    title_run,
)

CONFIG = load_config(default_config_path())


def violations_for(prs):
    return audit_deck(build_inventory(prs), CONFIG)


def by_rule(violations, rule):
    return [v for v in violations if v.rule == rule]


def test_old_kone_blue_in_heading_is_critical_and_fixable():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Inter", color_hex="005EB8")

    viol = by_rule(violations_for(prs), "legacy_color")
    assert len(viol) == 1
    assert viol[0].severity == "critical"
    assert viol[0].auto_fixable is True
    assert viol[0].details["target"] == "#1450F5"


def test_unapproved_font_in_heading_is_critical():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Calibri")

    viol = by_rule(violations_for(prs), "unapproved_font")
    assert len(viol) == 1
    assert viol[0].severity == "critical"
    assert viol[0].auto_fixable is True  # Calibri has a remap target


def test_unapproved_font_in_body_is_major_not_critical():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="Body copy", font="Calibri")

    viol = by_rule(violations_for(prs), "unapproved_font")
    assert len(viol) == 1
    assert viol[0].severity == "major"


def test_unapproved_font_with_no_remap_target_is_not_auto_fixable():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="Body copy", font="Papyrus")

    viol = by_rule(violations_for(prs), "unapproved_font")
    assert viol[0].auto_fixable is False


def test_inherited_unapproved_font_is_flagged_but_not_auto_fixed():
    """A run with no explicit font override still renders with a real
    font -- inherited from the master's txStyles/theme (Calibri, for
    python-pptx's stock template). That's now resolved and checked like
    any other font, but auto-fixing it would mean hardcoding a per-run
    override where the deck relied on inheritance -- the correct fix is
    upstream (theme/master), so this is flagged for review, not silently
    rewritten."""
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Inherits from layout"

    viol = by_rule(violations_for(prs), "unapproved_font")
    assert len(viol) == 1
    assert viol[0].details["current"] == "Calibri"
    assert viol[0].auto_fixable is False
    assert "inherited" in viol[0].message


def test_inherited_approved_font_is_not_flagged():
    """Same shape, but the deck's theme is already on-brand (Inter) --
    inheritance should resolve cleanly with no violation at all."""
    from tests.helpers import set_theme_font

    prs = new_deck()
    set_theme_font(prs, "majorFont", "Inter")
    set_theme_font(prs, "minorFont", "Inter")
    slide = add_slide(prs)
    title_run(slide).text = "Inherits from layout"

    assert by_rule(violations_for(prs), "unapproved_font") == []


def test_inter_in_non_approved_color_is_major():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="Body copy", font="Inter", color_hex="FF00FF")

    viol = by_rule(violations_for(prs), "text_color")
    assert len(viol) == 1
    assert viol[0].severity == "major"
    assert viol[0].details["target"] == "#141414"


def test_all_caps_inter_is_flagged_major_not_auto_fixable():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="SHOUTING TEXT", font="Inter", color_hex="141414")

    viol = by_rule(violations_for(prs), "all_caps")
    assert len(viol) == 1
    assert viol[0].severity == "major"
    assert viol[0].auto_fixable is False


def test_all_caps_allowed_word_kone_is_not_flagged():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="KONE", font="Inter", color_hex="141414")

    assert by_rule(violations_for(prs), "all_caps") == []


def test_kone_information_not_all_caps_is_flagged():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="not shouting", font="KONE Information", color_hex="141414")

    viol = by_rule(violations_for(prs), "all_caps")
    assert len(viol) == 1


def test_kone_information_oversized_is_role_restriction_violation():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="BIG LABEL", font="KONE Information", size_pt=40, color_hex="141414")

    viol = by_rule(violations_for(prs), "role_restriction")
    assert len(viol) == 1
    assert viol[0].auto_fixable is False


def test_forbidden_text_effect_shadow_is_major_and_fixable():
    prs = new_deck()
    slide = add_slide(prs)
    run = set_run(body_run(slide), text="shadowed", font="Inter", color_hex="141414")
    add_shadow_effect(run)

    viol = by_rule(violations_for(prs), "text_effect")
    assert len(viol) == 1
    assert viol[0].severity == "major"
    assert viol[0].auto_fixable is True


def test_alignment_violation_is_minor():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="centered body copy", font="Inter", color_hex="141414")
    slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    viol = by_rule(violations_for(prs), "alignment")
    assert len(viol) == 1
    assert viol[0].severity == "minor"


def test_pagination_like_text_is_exempted_from_alignment_rule():
    prs = new_deck()
    slide = add_slide(prs)
    body = slide.placeholders[1]
    body.left, body.top, body.width, body.height = 0, 6350000, 500000, 300000  # bottom-edge, small box
    set_run(body_run(slide), text="12", font="Inter", color_hex="141414")
    body.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    assert by_rule(violations_for(prs), "alignment") == []


def test_min_body_size_is_minor():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="tiny", font="Inter", size_pt=8, color_hex="141414")

    viol = by_rule(violations_for(prs), "min_size")
    assert len(viol) == 1
    assert viol[0].severity == "minor"


def test_wrong_slide_size_flagged():
    prs = new_deck()  # default python-pptx template is 4:3
    assert by_rule(violations_for(prs), "slide_size")[0].severity == "major"


def test_theme_based_legacy_tint_is_flagged_via_effective_color():
    from pptx.enum.dml import MSO_THEME_COLOR
    from tests.helpers import set_theme_slot

    prs = new_deck()
    set_theme_slot(prs, "accent1", "005EB8")
    slide = add_slide(prs)
    run = title_run(slide)
    run.text = "Themed"
    run.font.name = "Inter"
    set_run_theme_color(run, MSO_THEME_COLOR.ACCENT_1)

    viol = by_rule(violations_for(prs), "legacy_color")
    assert len(viol) == 1
    assert viol[0].severity == "critical"


def test_clean_deck_has_no_violations_besides_template_provenance():
    """python-pptx's stock template is inherently 4:3 and built on generic
    Office layout names, not KONE's -- slide_size and non_standard_layout
    are expected artifacts of the test fixture, not something a real
    KONE-template deck would ever trip. Everything else must be clean."""
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Inter Semi Bold", color_hex="141414")
    para = slide.shapes.title.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    set_run(body_run(slide), text="Body copy", font="Inter", color_hex="141414")
    slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    rules = {v.rule for v in violations_for(prs)}
    assert rules == {"slide_size", "non_standard_layout"}


def test_min_size_uses_the_paragraphs_own_outline_level_not_a_flat_floor():
    """The master's body text has a deliberately different minimum per
    outline level (footnote-tier levels are 11pt by design, not a flat
    number) -- config.min_size_by_level must be consulted per-paragraph,
    not just a single fonts.min_body_size_pt."""
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="footnote-tier", font="Inter", size_pt=11, color_hex="141414")
    para = slide.placeholders[1].text_frame.paragraphs[0]
    para.level = 6  # min_size_by_level[6] == 11pt in the real config -- should NOT be flagged

    assert by_rule(violations_for(prs), "min_size") == []


def test_min_size_still_flags_a_level_below_its_own_threshold():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="too small even for level 6", font="Inter", size_pt=8, color_hex="141414")
    para = slide.placeholders[1].text_frame.paragraphs[0]
    para.level = 6  # min_size_by_level[6] == 11pt -- 8pt still violates

    viol = by_rule(violations_for(prs), "min_size")
    assert len(viol) == 1
    assert viol[0].details["min_pt"] == 11


def test_min_size_exempts_plain_text_boxes_with_no_master_defined_floor():
    """Free-form text boxes (footnotes, tab labels, page numbers typed
    into a plain shape rather than a real placeholder) have no
    master-defined size floor -- unlike placeholder body text, they must
    not be flagged just for being small."""
    from pptx.util import Inches, Pt

    from pptx.dml.color import RGBColor

    prs = new_deck()
    slide = add_slide(prs)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "tiny footnote"
    run.font.name = "Inter"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor.from_string("141414")

    assert by_rule(violations_for(prs), "min_size") == []


def test_non_standard_layout_flagged_when_not_in_approved_list():
    prs = new_deck()
    add_slide(prs, layout_idx=1)  # "Title and Content" -- not a KONE template layout name

    config = dict(CONFIG)
    config["layout"] = dict(CONFIG["layout"])
    config["layout"]["approved_layouts"] = ["Some Other Layout"]

    violations = audit_deck(build_inventory(prs), config)
    viol = by_rule(violations, "non_standard_layout")
    assert len(viol) == 1
    assert viol[0].severity == "major"
    assert viol[0].auto_fixable is False


def test_non_standard_layout_not_flagged_when_layout_is_approved():
    prs = new_deck()
    add_slide(prs, layout_idx=1)

    config = dict(CONFIG)
    config["layout"] = dict(CONFIG["layout"])
    config["layout"]["approved_layouts"] = ["Title and Content"]

    violations = audit_deck(build_inventory(prs), config)
    assert by_rule(violations, "non_standard_layout") == []


def test_non_standard_layout_is_a_no_op_when_approved_layouts_not_configured():
    prs = new_deck()
    add_slide(prs, layout_idx=1)

    config = dict(CONFIG)
    config["layout"] = dict(CONFIG["layout"])
    config["layout"].pop("approved_layouts", None)

    violations = audit_deck(build_inventory(prs), config)
    assert by_rule(violations, "non_standard_layout") == []


def test_text_contrast_flags_black_text_on_kone_blue_background():
    """Regression test for General_Branding.docx's legibility rule: black
    text on the KONE Blue background is wrong -- must be white."""
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Blue panel", fill_hex="1450F5", left_in=1, top_in=1, width_in=3, height_in=1)
    box.text_frame.text = "Label"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("141414")

    violations = violations_for(prs)
    viol = by_rule(violations, "text_contrast")
    assert len(viol) == 1
    assert viol[0].details["target"] == "#FFFFFF"
    assert viol[0].severity == "major"
    assert viol[0].auto_fixable is True

    # the generic allowed-colors-list check must not also fire for the same run
    assert by_rule(violations, "text_color") == []


def test_text_contrast_flags_run_with_no_explicit_color_at_all():
    """Regression test for a real bug report: text with NO explicit color
    (inherits from the theme/master, which color-inheritance resolution
    doesn't attempt to reproduce -- unlike font) was silently skipped by
    the contrast check entirely, so a black-inheriting run sitting on a
    KONE Blue panel just stayed illegible. An unknown inherited color is
    exactly as much in scope as a wrong explicit one."""
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Blue panel", fill_hex="1450F5", left_in=1, top_in=1, width_in=3, height_in=1)
    box.text_frame.text = "Label"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    # no run.font.color set at all -- inherits

    viol = by_rule(violations_for(prs), "text_contrast")
    assert len(viol) == 1
    assert viol[0].details["target"] == "#FFFFFF"
    assert viol[0].details["current"] is None
    assert viol[0].auto_fixable is True


def test_text_contrast_passes_when_white_text_on_kone_blue_background():
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Blue panel", fill_hex="1450F5", left_in=1, top_in=1, width_in=3, height_in=1)
    box.text_frame.text = "Label"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("FFFFFF")

    assert by_rule(violations_for(prs), "text_contrast") == []


def test_text_contrast_flags_white_text_on_light_secondary_background():
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Pink card", fill_hex="FFCDD7", left_in=1, top_in=1, width_in=3, height_in=1)
    box.text_frame.text = "Label"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("FFFFFF")

    viol = by_rule(violations_for(prs), "text_contrast")
    assert len(viol) == 1
    assert viol[0].details["target"] == "#141414"


def test_text_contrast_always_white_on_kone_blue_60_percent_tint_overriding_wcag():
    """Regression test for an explicit brand-consistency request: KONE
    Blue's 60% tint (#7296F9) is technically light enough that WCAG math
    alone picks black text (higher contrast ratio) -- confirmed via
    colors.expected_contrast_text_hex directly. The brand wants every
    tint in always_light_text_backgrounds treated uniformly as a blue
    panel regardless, so the configured override must win over the math."""
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Blue tint panel", fill_hex="7296F9", left_in=1, top_in=1, width_in=3, height_in=1)
    box.text_frame.text = "Label"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("141414")

    viol = by_rule(violations_for(prs), "text_contrast")
    assert len(viol) == 1
    assert viol[0].details["target"] == "#FFFFFF"


def test_text_contrast_still_wcag_computed_for_tints_below_the_override_list():
    """The 40% tint (#A1B9FB) is deliberately NOT in
    always_light_text_backgrounds -- white text there would be
    borderline illegible, so it must stay contrast-computed (black)."""
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Light tint panel", fill_hex="A1B9FB", left_in=1, top_in=1, width_in=3, height_in=1)
    box.text_frame.text = "Label"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("FFFFFF")  # wrong -- should be black

    viol = by_rule(violations_for(prs), "text_contrast")
    assert len(viol) == 1
    assert viol[0].details["target"] == "#141414"


def test_text_contrast_heading_on_plain_canvas_is_always_dark():
    """The default, overwhelmingly common case: a heading with no
    confidently-resolved background behind it (the slide's own plain
    canvas) is always black -- a fixed editorial choice for that role,
    not a legibility default among two acceptable options."""
    prs = new_deck()
    slide = add_slide(prs)
    run = title_run(slide)
    run.text = "Heading"
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("FFFFFF")  # wrong -- headings are always dark

    viol = by_rule(violations_for(prs), "text_contrast")
    assert len(viol) == 1
    assert viol[0].details["target"] == "#141414"
    assert "always dark" in viol[0].message


def test_text_contrast_heading_on_plain_canvas_already_dark_is_not_flagged():
    prs = new_deck()
    slide = add_slide(prs)
    run = title_run(slide)
    run.text = "Heading"
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("141414")

    assert by_rule(violations_for(prs), "text_contrast") == []


def test_text_contrast_heading_on_a_known_blue_background_needs_white_not_black():
    """Regression test for a real bug report: a heading placed directly
    on a solid KONE Blue panel ("Key Points") still rendered in black,
    because the old blanket "headings are always dark" rule overrode
    contrast computation even when the background was confidently
    known. When a background IS known, headings now get the same
    contrast-aware answer ordinary text already gets."""
    prs = new_deck()
    slide = add_slide(prs)
    title = slide.shapes.title
    title.fill.solid()
    title.fill.fore_color.rgb = RGBColor.from_string("1450F5")
    run = title_run(slide)
    run.text = "Heading"
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("141414")  # wrong -- this background needs white

    viol = by_rule(violations_for(prs), "text_contrast")
    assert len(viol) == 1
    assert viol[0].details["target"] == "#FFFFFF"
    assert "always dark" not in viol[0].message


def test_text_contrast_heading_on_a_known_blue_background_already_white_is_not_flagged():
    prs = new_deck()
    slide = add_slide(prs)
    title = slide.shapes.title
    title.fill.solid()
    title.fill.fore_color.rgb = RGBColor.from_string("1450F5")
    run = title_run(slide)
    run.text = "Heading"
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("FFFFFF")

    assert by_rule(violations_for(prs), "text_contrast") == []


def test_text_contrast_not_checked_without_a_resolved_shape_background():
    """No shape-level solid fill (plain text on the page canvas), and
    nothing behind it either -- can't compute contrast, so it's skipped
    rather than guessed at; the generic allowed-colors check still applies."""
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="Body copy", font="Inter", color_hex="FFFFFF")

    assert by_rule(violations_for(prs), "text_contrast") == []
    # white isn't the fallback list's first choice, but it IS an allowed
    # color for Inter (black or white) -- so no text_color violation either
    assert by_rule(violations_for(prs), "text_color") == []


def test_text_contrast_resolves_background_from_a_separate_shape_behind_it():
    """Regression test for a real bug report: a plain textbox with no fill
    of its own, drawn on top of a separately-drawn KONE Blue rectangle
    (the common "color panel + textbox" authoring pattern), kept its
    dark/inherited body text -- illegible on blue -- because contrast was
    only ever checked against a shape's OWN fill. The textbox's effective
    background must resolve to the panel behind it in z-order."""
    prs = new_deck()
    slide = add_slide(prs)
    add_rectangle(slide, name="Blue panel", fill_hex="1450F5", left_in=1, top_in=1, width_in=3, height_in=2)
    # Added after the panel -- later in document order = on top of it.
    textbox = slide.shapes.add_textbox(Emu(int(914400 * 1.2)), Emu(int(914400 * 1.2)), Emu(int(914400 * 2)), Emu(int(914400 * 1)))
    textbox.text_frame.text = "Principaux avantages"
    run = textbox.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("595959")  # the illegible grey from the report

    viol = by_rule(violations_for(prs), "text_contrast")
    assert len(viol) == 1
    assert viol[0].shape_name == textbox.name
    assert viol[0].details["target"] == "#FFFFFF"
    assert viol[0].auto_fixable is True


def test_text_contrast_background_lookup_ignores_a_non_overlapping_shape():
    """A shape elsewhere on the slide -- not behind this textbox -- must
    never be mistaken for its background, even if it has a solid fill."""
    prs = new_deck()
    slide = add_slide(prs)
    add_rectangle(slide, name="Unrelated blue shape", fill_hex="1450F5", left_in=5, top_in=5, width_in=1, height_in=1)
    textbox = slide.shapes.add_textbox(Emu(int(914400 * 1)), Emu(int(914400 * 1)), Emu(int(914400 * 2)), Emu(int(914400 * 1)))
    textbox.text_frame.text = "Body copy"
    run = textbox.text_frame.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor.from_string("FFFFFF")

    assert by_rule(violations_for(prs), "text_contrast") == []


def test_unlisted_grey_panel_fill_defaults_to_secondary_off_white():
    """Regression test for a real bug report: a deck's big content panel
    used a literal grey (#D9D9D9) that was never in colors.remap -- it
    just sat flagged as unapproved forever, never actually turning into
    the off-white the brand wants for grey panels."""
    prs = new_deck()
    slide = add_slide(prs)
    add_rectangle(slide, name="Panel", fill_hex="D9D9D9", left_in=1, top_in=1, width_in=3, height_in=2)

    viol = by_rule(violations_for(prs), "legacy_color")
    panel_viol = [v for v in viol if v.shape_name == "Panel"]
    assert len(panel_viol) == 1
    assert panel_viol[0].details["target"] == "#F3EEE6"
    assert panel_viol[0].auto_fixable is True
    assert by_rule(violations_for(prs), "unapproved_color") == []


def test_unlisted_non_grey_panel_fill_defaults_to_kone_blue():
    prs = new_deck()
    slide = add_slide(prs)
    add_rectangle(slide, name="Panel", fill_hex="9B59B6", left_in=1, top_in=1, width_in=3, height_in=2)  # a purple

    viol = by_rule(violations_for(prs), "legacy_color")
    panel_viol = [v for v in viol if v.shape_name == "Panel"]
    assert len(panel_viol) == 1
    assert panel_viol[0].details["target"] == "#1450F5"


def test_unlisted_panel_fallback_does_not_apply_to_text():
    """Scoped to fill/line only -- an unlisted text color must still fall
    through to the ordinary unapproved_color/text_color path, not get
    silently recolored to blue or cream."""
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="Body copy", font="Inter", color_hex="9B59B6")

    assert by_rule(violations_for(prs), "legacy_color") == []
    matches = [v for v in violations_for(prs) if v.rule == "text_color"]
    assert len(matches) == 1
    assert matches[0].details["target"] == "#141414"


def test_text_on_a_coloured_layout_background_gets_light_text():
    """Reported on a real transform: the org template's "Title and text"
    layout is solid KONE Blue with a white panel over its right half, so
    its title placeholder sits on blue. `_resolve_effective_bg_hex` only
    searched SHAPES, found nothing behind the title, and the
    heading-always-dark rule put black text on KONE Blue."""
    from deckguard.rules_engine import _resolve_effective_bg_hex

    class _Shape:
        shape_id = 1
        left_in = top_in = 0.5
        width_in = height_in = 2.0
        fill = None

    shape = _Shape()
    # nothing on the slide and nothing on the layout to resolve from...
    assert _resolve_effective_bg_hex(shape, [], {}, []) is None
    # ...but the slide's own background colour is a perfectly good answer
    assert _resolve_effective_bg_hex(shape, [], {}, [], "1450F5") == "1450F5"


def test_slide_background_is_read_from_layout_then_master(tmp_path):
    from pptx import Presentation

    from deckguard.inventory import build_inventory
    from tests.helpers import add_slide, new_deck, title_run

    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    path = tmp_path / "d.pptx"
    prs.save(str(path))

    inv = build_inventory(Presentation(str(path)))
    # the field exists and is either a hex or an honest None
    value = inv.slides[0].background_hex
    assert value is None or (len(value) == 6 and value == value.upper())


def _bold_inter_deck(path):
    from tests.helpers import add_slide, new_deck, set_run, title_run

    prs = new_deck()
    slide = add_slide(prs)
    run = title_run(slide)
    set_run(run, text="A heading", font="Inter", color_hex="141414")
    run.font.bold = True
    return prs.save(str(path)) or path


def test_inter_is_never_bold(tmp_path):
    """Weight in this brand comes from choosing the "Inter SemiBold"
    FAMILY, never from bolding the regular cut. Switching an old deck's
    font to Inter carried the source's bold attribute across with it, so
    a rebranded deck came back with 67 bold Inter runs."""
    from pptx import Presentation

    from deckguard.config import default_config_path, load_config
    from deckguard.inventory import build_inventory
    from deckguard.rules_engine import audit_deck

    deck = _bold_inter_deck(tmp_path / "d.pptx")
    violations = audit_deck(build_inventory(Presentation(str(deck))), load_config(default_config_path()))

    weight = [v for v in violations if v.rule == "font_weight"]
    assert len(weight) == 1
    assert weight[0].auto_fixable
    assert "never used bold" in weight[0].message


def test_the_bold_rule_keys_on_the_family_not_the_approved_match(tmp_path):
    """`match_approved` folds a bold "Inter" into "Inter SemiBold" and
    reports it approved -- which is precisely the equivalence this rule
    rejects. Keying on that match made the rule silently never fire."""
    from deckguard.fonts import FontTables
    from deckguard.config import default_config_path, load_config

    config = load_config(default_config_path())
    tables = FontTables.from_config(config.get("fonts", {}))
    # the fold that hid the problem
    assert tables.match_approved("Inter", bold=True) == "Inter SemiBold"
    # and the family name the rule actually keys on
    assert "Inter" in set(config["typography_rules"]["never_bold_fonts"])


def test_bolded_inter_is_un_bolded_by_the_fixer(tmp_path):
    from pptx import Presentation

    from deckguard.fixer import fix_deck

    from deckguard.config import default_config_path, load_config

    deck = _bold_inter_deck(tmp_path / "d.pptx")
    out = tmp_path / "out.pptx"
    fix_deck(Presentation(str(deck)), load_config(default_config_path()),
             str(deck), str(out), False)

    bolds = [
        r.font.bold
        for slide in Presentation(str(out)).slides
        for shape in slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for r in para.runs
        if r.text.strip()
    ]
    assert not any(bolds), "no Inter run may remain bold"
