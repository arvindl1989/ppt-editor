"""What each archetype needs, and whether the renderer can take it.

The repetition problem was never the prompt. A third of the library
advertised nothing but `title` and `body`, so a planner reaching for the
timeline found it could put two paragraphs there and went back to the
handful of layouts with real slots. These tests hold the contract and
the renderer to each other in both directions, because either one
drifting reopens that hole silently.
"""

import pytest

from deckguard import brandmode as bm
from deckguard import contracts as C


def _built():
    from deckguard.registry import _load_archetypes

    return set(_load_archetypes().ARCHETYPES)


# --------------------------------------------------------------------------
# reading the handoff's own contract column
# --------------------------------------------------------------------------


def test_the_external_table_carries_a_contract_for_all_twenty_five():
    """It always did. `brandmode.jobs()` matched the column and threw it
    away, which is why the planner never saw what a slide needs."""
    table = C.external()
    assert len(table) == 25
    assert all(c.slots for c in table.values())


def test_the_contract_mini_language_parses_lists_with_their_fields():
    slots = C.parse("title:title · items[3]:{heading:heading, text:body}")
    assert [s.key for s in slots] == ["title", "items"]
    title, items = slots
    assert not title.is_list and title.role == "title"
    assert items.is_list and items.maximum == 3
    assert [f.key for f in items.fields] == ["heading", "text"]
    assert [f.role for f in items.fields] == ["heading", "body"]


def test_a_comma_inside_the_braces_is_not_a_slot_separator():
    slots = C.parse("a:title · b[2]:{x:body, y:body} · c:body")
    assert [s.key for s in slots] == ["a", "b", "c"]


def test_a_minimum_never_exceeds_what_the_layout_holds():
    for audience in ("internal", "external"):
        for contract in C.table(audience).values():
            for slot in contract.slots:
                assert slot.minimum <= slot.maximum, (contract.archetype, slot.key)


def test_the_handoff_stays_the_source_for_the_external_set():
    """Transcribing it by hand would let the two drift. Only the slot
    NAMES the built layout disagrees with are corrected in code."""
    assert C.external()["three_content"].slots == C.parse(
        "title:title · items[3]:{heading:heading, text:body}")


# --------------------------------------------------------------------------
# the contract and the renderer, held to each other
# --------------------------------------------------------------------------


def test_every_archetype_the_picker_offers_has_a_contract():
    missing = []
    for audience in bm.set_names():
        for slide in bm.slides_in(audience):
            name = slide["archetype"]
            if name in _built() and C.for_archetype(name, audience) is None:
                missing.append(name)
    assert not missing, missing


def test_no_contract_promises_a_slot_the_renderer_cannot_take():
    """The punch list, generated rather than written down: it shrinks as
    the registry is fixed and can never claim a gap has been closed when
    it has not. It started at thirteen archetypes."""
    found = C.gaps()
    assert not found, "\n".join(
        f"{name} [{aud}]: {'; '.join(miss)}"
        for name, per in sorted(found.items()) for aud, miss in per.items())


def test_a_broken_contract_is_actually_reported():
    """A check that passes because it looks at nothing is worse than no
    check, and this one is generated -- so prove it can still fail."""
    original = dict(C._INTERNAL_SPECS)
    C._INTERNAL_SPECS["three_content"] = "title:title · nonesuch[9]:{a, b}"
    C.internal.cache_clear()
    try:
        found = C.gaps("internal")
        assert "three_content" in found
        assert any("no such slot" in m for m in found["three_content"]["internal"])
    finally:
        C._INTERNAL_SPECS.clear()
        C._INTERNAL_SPECS.update(original)
        C.internal.cache_clear()


def test_a_field_name_the_layout_does_not_read_is_reported():
    """The half of the check that was missing at first: `quarters[4]`
    existed, but its fields were named `period, heading` against a group
    built as `label, items` -- so the planner would have been told to
    emit two names nothing reads, and the labels came out blank."""
    original = dict(C._INTERNAL_SPECS)
    C._INTERNAL_SPECS["three_content"] = "title:title · items[3]:{wrong:body, alsowrong:body}"
    C.internal.cache_clear()
    try:
        found = C.gaps("internal")
        assert any("wrong" in m for m in found["three_content"]["internal"])
    finally:
        C._INTERNAL_SPECS.clear()
        C._INTERNAL_SPECS.update(original)
        C.internal.cache_clear()


# --------------------------------------------------------------------------
# the three layouts that could not carry structured content
# --------------------------------------------------------------------------


@pytest.mark.parametrize("name,keys", [
    ("agenda_c_split", {"lead", "items"}),
    ("timeline_quarter_axis", {"lead", "bullets", "events"}),
    ("picture_intro", {"eyebrow", "points"}),
])
def test_the_starved_layouts_now_take_structured_content(name, keys):
    """All three offered `title` and a paragraph. An agenda with no
    agenda items is not a slide anyone would choose on purpose."""
    have = set(C._registry_slots(name))
    assert keys <= have, f"{name} is missing {keys - have}"


def test_the_split_agenda_draws_five_numbered_rows(tmp_path):
    from pptx import Presentation

    from deckguard import assemble

    out = tmp_path / "a.pptx"
    checks = assemble.build({"title": "T", "date": "1 March 2026", "slides": [{
        "archetype": "agenda_c_split", "title": "Today", "lead": "The shape of it.",
        "items": [{"number": f"0{n}", "label": f"Item {n}"} for n in range(1, 6)],
    }]}, str(out))
    assert not checks["findings"], checks["findings"]

    text = " ".join(
        sh.text_frame.text for sh in list(Presentation(str(out)).slides)[1].shapes
        if getattr(sh, "has_text_frame", False))
    for n in range(1, 6):
        assert f"Item {n}" in text
        assert f"0{n}" in text


def test_the_quarter_timeline_draws_its_events(tmp_path):
    from pptx import Presentation

    from deckguard import assemble

    out = tmp_path / "t.pptx"
    assemble.build({"title": "T", "date": "1 March 2026", "slides": [{
        "archetype": "timeline_quarter_axis", "title": "Milestones",
        "lead": "One intake, four quarters.", "bullets": ["Migration complete"],
        "events": [{"period": f"Q{n} 2026", "text": f"Step {n}."} for n in range(1, 5)],
    }]}, str(out))
    text = " ".join(
        sh.text_frame.text for sh in list(Presentation(str(out)).slides)[1].shapes
        if getattr(sh, "has_text_frame", False))
    for n in range(1, 5):
        assert f"Q{n} 2026" in text and f"Step {n}." in text


def test_a_full_height_colour_field_is_not_a_floor_violation(tmp_path):
    """The split agenda's mint column runs edge to edge by design. Read
    as content it reaches y=720 and preflight cried wolf on every deck
    that used the slide."""
    from deckguard import assemble

    out = tmp_path / "f.pptx"
    checks = assemble.build({"title": "T", "date": "1 March 2026", "slides": [
        {"archetype": "agenda_c_split", "title": "Today",
         "items": [{"number": "01", "label": "One"}]}]}, str(out))
    assert not [m for _n, m in checks["findings"] if "floor" in m], checks["findings"]


# --------------------------------------------------------------------------
# what the planner is told
# --------------------------------------------------------------------------


def test_the_menu_says_what_each_slide_needs():
    guide = C.guide("internal")
    assert "needs:" in guide
    assert "items (5 × {number, label})" in guide      # the re-specced agenda
    assert "events (3-4 × {period, text})" in guide    # a range, not a single number


def test_the_planner_is_told_what_to_have_before_choosing():
    from deckguard import assemble, planner

    seen = {}

    def _capture(brief, **kw):
        seen["notes"] = kw.get("notes")
        return {"slides": [{"archetype": "title_content", "title": "x"}]}, None

    original = planner.call_claude_for_kone_spec
    planner.call_claude_for_kone_spec = _capture
    try:
        assemble.plan(brief="An email.", audience="internal")
    finally:
        planner.call_claude_for_kone_spec = original

    notes = seen["notes"]
    assert "says what it NEEDS" in notes
    assert "needs:" in notes


def test_the_system_prompt_rules_an_archetype_out_rather_than_only_in():
    from deckguard.planner import _kone_archetype_guide

    guide = _kone_archetype_guide()
    assert "Do not choose this unless the source gives you" in guide


# --------------------------------------------------------------------------
# sample content, so a preview shows the slide full
# --------------------------------------------------------------------------


def test_sample_content_fills_a_list_to_what_the_layout_holds():
    """Guessed from key names it produced three rows for a five-row
    agenda, and the preview showed two empty sand blocks."""
    from deckguard.preview import sample_content

    assert len(sample_content("agenda_c_split")["items"]) == 5
    assert len(sample_content("timeline_quarter_axis")["events"]) == 4


def test_sample_content_never_names_an_icon():
    """Unnamed, the engine runs its own rotation -- which is what an
    unspecified slide actually looks like."""
    from deckguard.preview import sample_content

    for name in ("picture_intro", "icon_columns_5", "resource_links"):
        for item in sample_content(name).get("items") or sample_content(name).get("tiles") or []:
            if isinstance(item, dict):
                assert "icon" not in item, name


# --------------------------------------------------------------------------
# the floor, which has no exceptions
# --------------------------------------------------------------------------


def test_no_set_archetype_draws_below_the_floor():
    """BRAND_MODE §7 is unconditional. Thirteen of the built archetypes
    broke it, one by 83px, and nobody had noticed because nothing built
    the whole library at once and read the preflight back."""
    from deckguard.registry import _load_archetypes

    over = []
    # The canonical set only. A design mined from someone's own deck is
    # theirs -- it is reproduced as drawn, and preflight is what says
    # where it sits off-brand. Holding it to §7 here would fail this
    # test whenever another module had uploaded a deck first.
    canonical = bm.canonical_archetypes()
    for name in sorted(_load_archetypes().ARCHETYPES):
        spec = _load_archetypes().ARCHETYPES[name]
        if not isinstance(spec, dict) or name not in canonical:
            continue
        from deckguard.layouts import _is_full_bleed, _PICTURE_REGION_ROLES

        for region in spec.get("regions") or []:
            if region.get("role") in _PICTURE_REGION_ROLES or region.get("role") == "image":
                continue
            box = region["box"]
            if not _is_full_bleed(box) and box[1] + box[3] > bm.FLOOR:
                over.append(f"{name}.{region.get('content')} → {box[1] + box[3]:.0f}")
        for group in spec.get("groups") or []:
            lowest = max((o[1] for o in group.get("origins") or []), default=0)
            for region in group.get("regions") or []:
                if region.get("role") in _PICTURE_REGION_ROLES:
                    continue
                bottom = lowest + region["box"][1] + region["box"][3]
                if bottom > bm.FLOOR:
                    over.append(f"{name}.{group.get('content')} → {bottom:.0f}")
    assert not over, over


def test_the_whole_library_builds_and_preflights_clean(tmp_path):
    """One deck holding every archetype the picker offers. This is the
    test that found the floor violations, the duplicated footers and the
    table read as a bulleted dash -- none of which any single-archetype
    test could see."""
    from deckguard import assemble
    from deckguard.preview import sample_content
    from deckguard.registry import _load_archetypes

    built = set(_load_archetypes().ARCHETYPES)
    names = []
    for audience in bm.set_names():
        for slide in bm.slides_in(audience):
            if slide["archetype"] in built and slide["archetype"] not in names:
                names.append(slide["archetype"])

    checks = assemble.build(
        {"title": "All", "date": "1 March 2026",
         "slides": [{"archetype": n, **sample_content(n)} for n in names]},
        str(tmp_path / "all.pptx"))
    assert not checks["findings"], [
        (names[n - 2] if 1 < n <= len(names) + 1 else n, m)
        for n, m in checks["findings"]]


def test_the_archetypes_own_footer_lines_are_gone():
    """`stamp_chrome` draws the date and page number on every body
    slide. An archetype-level footer was a second one underneath it, at
    y=680, in a slot the planner was invited to fill."""
    for name, key in (("kone_numbers", "footer"), ("credits", "footer"),
                      ("milestone_slide", "classification")):
        assert key not in C._registry_slots(name), f"{name}.{key}"
        assert C.for_archetype(name, "internal").slot(key) is None


def test_a_dash_in_a_table_cell_is_a_value_not_a_bullet(tmp_path):
    from deckguard import assemble

    checks = assemble.build({"title": "T", "date": "1 March 2026", "slides": [{
        "archetype": "comparison_table", "title": "Tiers",
        "table": {"headers": ["", "A", "B"],
                  "rows": [["Monitoring", "—", "Yes"]]},
    }]}, str(tmp_path / "c.pptx"))
    assert not [m for _n, m in checks["findings"] if "dash" in m], checks["findings"]


def test_the_table_slot_advertises_its_real_shape():
    """Annotated `(text)`, the one archetype built to hold a table was
    advertised as taking a paragraph -- and got one."""
    from deckguard.registry import _derived_content_keys

    keys = _derived_content_keys("comparison_table")
    assert any("headers" in k and "rows" in k for k in keys), keys
