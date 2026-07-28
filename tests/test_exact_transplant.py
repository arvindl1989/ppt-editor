from pptx import Presentation

from deckguard.exact_transplant import transplant_exact_treatment
from tests.helpers import add_rectangle, set_run


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # "Blank" -- no placeholders to confuse shape lookups


def _by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(name)


def test_transplant_copies_fill_from_name_matched_reference_shape():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Box A", fill_hex="FF0000")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Box A", fill_hex="0000FF")

    result = transplant_exact_treatment(old_prs, new_prs)

    box = _by_name(old_prs.slides[0], "Box A")
    assert str(box.fill.fore_color.rgb) == "0000FF"
    assert len(result.changes) == 1
    assert result.changes[0].field == "fill"
    assert result.changes[0].old == "FF0000"
    assert result.changes[0].new == "0000FF"


def test_transplant_canonicalizes_reference_color_through_remap():
    """The reference deck can still carry a stale literal hex
    brand_rules.yaml already knows the approved replacement for --
    copying it verbatim would regress a color the generic brand pass
    already fixed, so it's canonicalized through colors.remap first."""
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Box A", fill_hex="F3EEE6")  # already the approved target

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Box A", fill_hex="EDEFF0")  # the legacy source hex

    config = {"colors": {"remap": {"#EDEFF0": "#F3EEE6"}}}
    result = transplant_exact_treatment(old_prs, new_prs, rules_config=config)

    box = _by_name(old_prs.slides[0], "Box A")
    assert str(box.fill.fore_color.rgb) == "F3EEE6"
    assert result.changes == []  # already matched the canonical value -- no-op


def test_transplant_leaves_shape_alone_with_no_name_match_in_reference():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Box A", fill_hex="FF0000")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Box B", fill_hex="0000FF")  # different name -- no correspondence

    result = transplant_exact_treatment(old_prs, new_prs)

    box = _by_name(old_prs.slides[0], "Box A")
    assert str(box.fill.fore_color.rgb) == "FF0000"
    assert result.changes == []


def test_transplant_flags_slide_with_low_shape_name_match_ratio():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Matched", fill_hex="FF0000", left_in=1)
    add_rectangle(old_slide, name="Unmatched 1", fill_hex="FF0000", left_in=2)
    add_rectangle(old_slide, name="Unmatched 2", fill_hex="FF0000", left_in=3)

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Matched", fill_hex="0000FF", left_in=1)
    add_rectangle(new_slide, name="Different A", fill_hex="00FF00", left_in=2)
    add_rectangle(new_slide, name="Different B", fill_hex="00FF00", left_in=3)

    result = transplant_exact_treatment(old_prs, new_prs)

    assert result.flagged_slides == [0]


def test_transplant_does_not_flag_slide_with_high_shape_name_match_ratio():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Box A", fill_hex="FF0000")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Box A", fill_hex="0000FF")

    result = transplant_exact_treatment(old_prs, new_prs)

    assert result.flagged_slides == []


def test_transplant_copies_font_name_bold_and_color_from_matched_run():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    old_box = add_rectangle(old_slide, name="Box A")
    old_run = old_box.text_frame.paragraphs[0].add_run()
    set_run(old_run, text="Hello", font="Arial", bold=False, color_hex="141414")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A")
    new_run = new_box.text_frame.paragraphs[0].add_run()
    set_run(new_run, text="Bonjour", font="Inter", bold=True, color_hex="1450F5")

    config = {"fonts": {"approved": ["Inter"]}}
    result = transplant_exact_treatment(old_prs, new_prs, rules_config=config)

    run = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0].runs[0]
    assert run.text == "Hello"  # content stays verbatim
    assert run.font.name == "Inter"
    assert run.font.bold is True
    assert str(run.font.color.rgb) == "1450F5"
    fields = {c.field for c in result.changes}
    assert {"font_name", "font_bold", "font_color"} <= fields


def test_transplant_never_copies_a_font_the_reference_deck_uses_that_is_not_approved():
    """Hard rule, independent of the reference deck's own quality: this
    deck must only ever end up with KONE-approved fonts. If the reference
    itself carries something off-brand (never fixed, or hand-edited after
    the fact) with no known remap for it, the font_name copy is skipped
    entirely rather than propagating it -- unlike colors, there's no
    canonicalization fallback here, just a hard no."""
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    old_box = add_rectangle(old_slide, name="Box A")
    set_run(old_box.text_frame.paragraphs[0].add_run(), text="Hello", font="Inter")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A")
    set_run(new_box.text_frame.paragraphs[0].add_run(), text="Bonjour", font="Comic Sans MS")

    config = {"fonts": {"approved": ["Inter"], "remap": {}}}
    result = transplant_exact_treatment(old_prs, new_prs, rules_config=config)

    run = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Inter"  # left alone, never regressed to the reference's off-brand font
    assert not [c for c in result.changes if c.field == "font_name"]


def test_transplant_canonicalizes_reference_font_through_remap():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    old_box = add_rectangle(old_slide, name="Box A")
    set_run(old_box.text_frame.paragraphs[0].add_run(), text="Hello", font="Inter")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A")
    set_run(new_box.text_frame.paragraphs[0].add_run(), text="Bonjour", font="Arial")  # legacy, but remapped

    config = {"fonts": {"approved": ["Inter"], "remap": {"Arial": "Inter"}}}
    result = transplant_exact_treatment(old_prs, new_prs, rules_config=config)

    run = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Inter"
    assert not [c for c in result.changes if c.field == "font_name"]  # already matched the canonical value


def test_transplant_copies_paragraph_alignment_from_matched_reference_shape():
    from pptx.enum.text import PP_ALIGN

    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    old_box = add_rectangle(old_slide, name="Box A")
    old_box.text_frame.paragraphs[0].add_run().text = "Code"
    old_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A")
    new_box.text_frame.paragraphs[0].add_run().text = "Code"
    new_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    result = transplant_exact_treatment(old_prs, new_prs)

    para = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0]
    assert para.alignment == PP_ALIGN.CENTER
    assert any(c.field == "alignment" for c in result.changes)


def test_transplant_skips_text_color_that_would_match_shapes_own_fill():
    """A duplicate/decoy shape in the reference deck can pair a text
    color with a fill that canonicalize to the SAME approved token --
    applying it verbatim would make the text invisible against its own
    background, so this is refused even though the fill copy itself
    (evaluated first) is applied normally."""
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    box = add_rectangle(old_slide, name="Box A", fill_hex="AAAAAA")
    set_run(box.text_frame.paragraphs[0].add_run(), text="Code", color_hex="141414")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A", fill_hex="F3EEE6")
    set_run(new_box.text_frame.paragraphs[0].add_run(), text="Code", color_hex="F3EEE6")

    result = transplant_exact_treatment(old_prs, new_prs)

    run = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == "141414"  # left alone, not made invisible
    fill_changes = [c for c in result.changes if c.field == "fill"]
    assert fill_changes and fill_changes[0].new == "F3EEE6"  # fill copy still applied
    assert not [c for c in result.changes if c.field == "font_color"]
