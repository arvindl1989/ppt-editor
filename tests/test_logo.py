import pytest
from pptx.util import Emu, Inches

from deckguard import logo as logo_mod
from tests.helpers import add_picture, make_pattern_png, make_solid_png, new_deck, set_background_image


def _kone_master():
    """The kone-design master, vendored or installed -- None if neither."""
    from pathlib import Path

    candidates = [
        Path(__file__).resolve().parents[1] / "src" / "deckguard" / "assets" / "kone-design"
        / "uploads" / "master_ppt-1784774200983.pptx",
        Path.home() / ".claude" / "skills" / "kone-design" / "uploads"
        / "master_ppt-1784774200983.pptx",
    ]
    return next((c for c in candidates if c.is_file()), None)


def test_compute_phash_stable_for_identical_image(tmp_path):
    p = make_solid_png(tmp_path / "a.png", rgb=(0, 94, 184))
    h1 = logo_mod.compute_phash(p.read_bytes())
    h2 = logo_mod.compute_phash(p.read_bytes())
    assert h1 == h2
    assert len(h1) == 16  # 8x8 phash -> 64 bits -> 16 hex chars


def test_hamming_distance_zero_for_same_hash():
    h = "abcdef1234567890"
    assert logo_mod.hamming_distance(h, h) == 0


def test_hamming_distance_large_for_different_patterns(tmp_path):
    a = make_pattern_png(tmp_path / "a.png", seed=1)
    b = make_pattern_png(tmp_path / "b.png", seed=2)
    ha = logo_mod.compute_phash(a.read_bytes())
    hb = logo_mod.compute_phash(b.read_bytes())
    assert logo_mod.hamming_distance(ha, hb) > 5


def test_find_old_logo_matches_by_similar_image(tmp_path):
    old_logo = make_pattern_png(tmp_path / "old_logo.png", seed=1)
    unrelated = make_pattern_png(tmp_path / "unrelated.png", seed=2)
    old_hash = logo_mod.compute_phash(old_logo.read_bytes())

    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    logo_shape = add_picture(slide, str(old_logo))
    add_picture(slide, str(unrelated), left_in=4)

    matches = logo_mod.find_old_logo_matches(slide.shapes, [old_hash], threshold=5)
    assert len(matches) == 1
    # python-pptx hands back a fresh proxy object per traversal, so compare
    # the underlying XML element rather than object identity.
    assert matches[0].shape._element is logo_shape._element


def test_replace_logo_image_preserves_position_and_size(tmp_path):
    old_logo = make_pattern_png(tmp_path / "old_logo.png", seed=1)
    new_logo = make_pattern_png(tmp_path / "new_logo.png", seed=3)

    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = add_picture(slide, str(old_logo), left_in=2, top_in=3, width_in=1.5, height_in=0.75)
    left, top, width, height = shape.left, shape.top, shape.width, shape.height

    logo_mod.replace_logo_image(shape, str(new_logo))

    assert shape.left == left
    assert shape.top == top
    assert shape.width == width
    assert shape.height == height
    assert shape.image.blob == new_logo.read_bytes()


def test_find_shapes_in_region_matches_only_fully_contained_shapes():
    prs = new_deck()
    master = prs.slide_masters[0]
    spTree = master.shapes._spTree
    id_ = master.shapes._next_shape_id
    inside = spTree.add_textbox(id_, "Inside", Inches(11), Inches(0.2), Inches(1), Inches(0.5))
    id_ += 1
    straddling = spTree.add_textbox(id_, "Straddling", Inches(9), Inches(0.2), Inches(3), Inches(0.5))

    region = (Emu(Inches(10.5)), Emu(Inches(0)), Emu(Inches(2.8)), Emu(Inches(1.2)))
    matches = logo_mod.find_shapes_in_region(master.shapes, region)

    matched_elements = {s._element for s in matches}
    assert inside in matched_elements
    assert straddling not in matched_elements  # left edge (9in) is outside the region -- not fully contained


def test_find_shapes_in_region_empty_when_nothing_is_inside():
    prs = new_deck()
    master = prs.slide_masters[0]
    region = (Emu(Inches(10.5)), Emu(Inches(0)), Emu(Inches(2.8)), Emu(Inches(1.2)))
    # the master's own default title/body/date/footer placeholders all live elsewhere
    assert logo_mod.find_shapes_in_region(master.shapes, region) == []


def test_replace_shapes_in_region_with_logo_removes_matches_and_inserts_sized_picture(tmp_path):
    new_logo = make_pattern_png(tmp_path / "new_logo.png", seed=4)

    prs = new_deck()
    master = prs.slide_masters[0]
    spTree = master.shapes._spTree
    id_ = logo_mod._next_shape_id_in_tree(spTree)
    old_mark = spTree.add_textbox(id_, "OldMark", Inches(11), Inches(0.2), Inches(1), Inches(0.5))
    before_count = len(master.shapes)

    region = (Emu(Inches(10.5)), Emu(Inches(0)), Emu(Inches(2.8)), Emu(Inches(1.2)))
    matches = logo_mod.find_shapes_in_region(master.shapes, region)
    assert len(matches) == 1

    logo_mod.replace_shapes_in_region_with_logo(master.shapes, matches, str(new_logo), region)

    assert len(master.shapes) == before_count  # one removed, one added -- net unchanged
    assert old_mark.getparent() is None  # actually removed from the tree, not just unreferenced

    pictures = [s for s in master.shapes if s.shape_type is not None and s.shape_type.name == "PICTURE"]
    assert len(pictures) == 1
    pic = pictures[0]
    # fits inside the region on both axes, aspect ratio preserved (not stretched)
    r_left, r_top, r_width, r_height = region
    assert r_left <= pic.left and pic.left + pic.width <= r_left + r_width
    assert r_top <= pic.top and pic.top + pic.height <= r_top + r_height


def test_replace_shapes_in_region_with_logo_on_a_master_never_assigns_a_layout_id(tmp_path):
    """Regression test for a real report: PowerPoint outright refused to
    open a .pptx this produced. Root cause -- confirmed by inspecting the
    output XML directly -- was calling python-pptx's own `_next_shape_id`
    on a MASTER's shape tree: a slide master always has a sibling
    `p:sldLayoutIdLst` (listing the layouts that belong to it) whose
    `p:sldLayoutId` elements use a completely different id namespace
    starting at 2**31 by OOXML convention, and `_next_shape_id`'s
    document-wide `//@id` XPath scan picks those up too, handing back an
    id like 2147483687 -- past the signed-32-bit range PowerPoint's own
    parser tolerates. Every stock python-pptx master has this sibling
    list (it's how a master knows which layouts belong to it), so this
    isn't specific to any one deck's quirks."""
    new_logo = make_pattern_png(tmp_path / "new_logo.png", seed=5)

    prs = new_deck()
    master = prs.slide_masters[0]
    assert len(master.shapes._spTree.xpath("//@id")) > 0  # sanity: sldLayoutIdLst ids are visible in this scan
    spTree = master.shapes._spTree
    id_ = logo_mod._next_shape_id_in_tree(spTree)
    spTree.add_textbox(id_, "OldMark", Inches(11), Inches(0.2), Inches(1), Inches(0.5))

    region = (Emu(Inches(10.5)), Emu(Inches(0)), Emu(Inches(2.8)), Emu(Inches(1.2)))
    matches = logo_mod.find_shapes_in_region(master.shapes, region)

    logo_mod.replace_shapes_in_region_with_logo(master.shapes, matches, str(new_logo), region)

    pic = next(s for s in master.shapes if s.shape_type is not None and s.shape_type.name == "PICTURE")
    assert pic.shape_id < 2**31


def test_reference_logo_geometry_matches_the_bundled_template():
    """The bundled KONE template's own logo shapes agree on one
    size/position across the overwhelming majority of its layouts --
    reference_logo_geometry should return exactly that, in EMU."""
    from deckguard.legacy.slide_import import default_template_path

    if not default_template_path().exists():
        import pytest

        pytest.skip("bundled template asset not present")

    geometry = logo_mod.reference_logo_geometry(default_template_path())
    assert geometry is not None
    left, top, width, height = geometry
    # known-good values confirmed by direct inspection of the bundled template
    assert abs(Emu(left).inches - 12.013) < 0.05
    assert abs(Emu(top).inches - 0.472) < 0.05
    assert abs(Emu(width).inches - 0.846) < 0.05
    assert abs(Emu(height).inches - 0.328) < 0.05


def test_reference_logo_geometry_returns_none_for_a_template_with_no_logo_shape():
    prs = new_deck()  # a plain default python-pptx template, no shape named "logo"
    import tempfile
    from pathlib import Path

    with tempfile.TemporaryDirectory() as d:
        path = Path(d) / "blank.pptx"
        prs.save(str(path))
        assert logo_mod.reference_logo_geometry(path) is None


def test_next_shape_id_in_tree_ignores_sibling_sldLayoutIdLst_ids():
    """Unit-level version of the same regression: a master's own
    _next_shape_id_in_tree must stay scoped to p:cNvPr ids and ignore
    the 2**31+ range p:sldLayoutId elements use, unlike python-pptx's
    own _next_shape_id (a document-wide XPath scan)."""
    prs = new_deck()
    master = prs.slide_masters[0]
    spTree = master.shapes._spTree

    buggy = master.shapes._next_shape_id  # python-pptx's own property
    fixed = logo_mod._next_shape_id_in_tree(spTree)

    assert buggy >= 2**31
    assert fixed < 100  # a stock master's own placeholder ids are all small


def test_find_background_blip_returns_none_when_no_bg_element(tmp_path):
    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert logo_mod.find_background_blip(slide._element) is None


def test_find_background_blip_finds_a_picture_fill_background(tmp_path):
    """Regression coverage for a real gap: a logo baked into a slide's
    page-level background-FILL image (<p:cSld><p:bg>) rather than being a
    picture shape at all -- outside the shape tree entirely, so no
    shape-based scan (however thorough) can see it without this."""
    logo_img = make_pattern_png(tmp_path / "logo.png", seed=1)
    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    blip = set_background_image(slide, str(logo_img))

    assert logo_mod.find_background_blip(slide._element) is blip


def test_find_old_logo_background_match_by_similar_image(tmp_path):
    old_logo = make_pattern_png(tmp_path / "old_logo.png", seed=1)
    old_hash = logo_mod.compute_phash(old_logo.read_bytes())

    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_image(slide, str(old_logo))

    match = logo_mod.find_old_logo_background_match(slide.part, slide._element, [old_hash], threshold=5)
    assert match is not None
    assert match.matched_hash == old_hash


def test_find_old_logo_background_match_none_for_unrelated_image(tmp_path):
    old_logo = make_pattern_png(tmp_path / "old_logo.png", seed=1)
    old_hash = logo_mod.compute_phash(old_logo.read_bytes())
    unrelated = make_pattern_png(tmp_path / "unrelated.png", seed=2)

    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_image(slide, str(unrelated))

    match = logo_mod.find_old_logo_background_match(slide.part, slide._element, [old_hash], threshold=5)
    assert match is None


def test_replace_background_image_swaps_the_blob(tmp_path):
    old_logo = make_pattern_png(tmp_path / "old_logo.png", seed=1)
    new_logo = make_pattern_png(tmp_path / "new_logo.png", seed=3)

    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    blip = set_background_image(slide, str(old_logo))

    logo_mod.replace_background_image(slide.part, blip, str(new_logo))

    rid = blip.get(f"{{{logo_mod.R_NS}}}embed")
    assert slide.part.related_part(rid).blob == new_logo.read_bytes()


def test_the_masters_empty_logo_frames_are_repaired(tmp_path):
    """The KONE master ships 45 `Logo` and 7 `Tagline` picture shapes --
    across 47 of its 63 layouts -- whose `<a:blip>` carries no
    relationship at all. PowerPoint draws a picture frame with no
    picture as a dotted rectangle, so every deck built on that master
    shows dotted boxes where its logo should be."""
    import shutil

    from pptx import Presentation

    from pathlib import Path

    from deckguard.logo import repair_empty_logo_frames

    # The defect is in the kone-design MASTER the archetype engine builds
    # from, not in the org template.
    candidates = [
        Path(__file__).resolve().parents[1] / "src" / "deckguard" / "assets" / "kone-design"
        / "uploads" / "master_ppt-1784774200983.pptx",
        Path.home() / ".claude" / "skills" / "kone-design" / "uploads"
        / "master_ppt-1784774200983.pptx",
    ]
    master = next((c for c in candidates if c.is_file()), None)
    if master is None:
        pytest.skip("kone-design master template not available")

    deck = tmp_path / "d.pptx"
    shutil.copy(str(master), str(deck))

    def empty_frames(path):
        count = 0
        prs = Presentation(str(path))
        for master in prs.slide_masters:
            for container in [master, *master.slide_layouts]:
                for shape in container.shapes:
                    if shape._element.tag.split("}")[-1] != "pic":
                        continue
                    try:
                        shape.image.blob
                    except Exception:  # noqa: BLE001
                        count += 1
        return count

    # The vendored master is now shipped already repaired, so the
    # invariant to hold is that it has nothing left to fix -- and that
    # a fresh pass over it is a no-op rather than a re-write.
    assert empty_frames(deck) == 0, (
        "the vendored master must ship with every logo frame filled; "
        "re-syncing it from the skill means re-running the repair"
    )
    assert repair_empty_logo_frames(deck) == 0


def test_an_empty_logo_frame_is_filled_with_the_right_variant(tmp_path):
    """The repair itself, on a deck built to have the defect -- the
    real master no longer does, and the repair still has to work for
    every user deck built on an unrepaired copy.

    A dark page gets the white mark; a light page gets the blue one.
    """
    from pptx import Presentation

    from deckguard.logo import _brand_asset, repair_empty_logo_frames

    if _brand_asset("logo", light=False) is None:
        pytest.skip("vendored KONE marks not available")

    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    picture = slide.shapes.add_picture(
        _brand_asset("logo", light=False), Emu(0), Emu(0), Emu(914400), Emu(914400)
    )
    picture.name = "Logo"
    # strip the relationship, reproducing the master's defect exactly
    picture._element.blipFill.blip.rEmbed = None

    deck = tmp_path / "broken.pptx"
    prs.save(str(deck))

    assert repair_empty_logo_frames(deck) == 1
    reopened = Presentation(str(deck))
    repaired = next(
        s for s in reopened.slides[0].shapes
        if s._element.tag.split("}")[-1] == "pic"
    )
    assert repaired.image.blob, "the frame must carry a real image afterwards"


def test_a_repaired_mark_carries_its_vector_original(tmp_path):
    """PowerPoint keeps SVG as an extension on a raster blip. The mark
    goes in as PNG so python-pptx accepts it, and the SVG rides along so
    it stays crisp at any zoom -- which is exactly what the KONE master
    does for 46 of its own images."""
    import re
    import zipfile
    from pathlib import Path

    from pptx import Presentation
    from pptx.util import Emu

    from deckguard.logo import _SVG_EXT_URI, _brand_asset, attach_svg

    png = _brand_asset("logo", light=False)
    if png is None or not Path(png).with_suffix(".svg").is_file():
        pytest.skip("vendored KONE marks not available as SVG + PNG")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    picture = slide.shapes.add_picture(png, Emu(0), Emu(0), Emu(914400), Emu(914400))
    assert attach_svg(picture, Path(png).with_suffix(".svg")) is True

    deck = tmp_path / "vector.pptx"
    prs.save(str(deck))

    with zipfile.ZipFile(str(deck)) as bundle:
        assert "image/svg+xml" in bundle.read("[Content_Types].xml").decode()
        xml = bundle.read("ppt/slides/slide1.xml").decode()
        assert _SVG_EXT_URI in xml
        rels = bundle.read("ppt/slides/_rels/slide1.xml.rels").decode()
        svg_rel = re.search(r'Id="([^"]+)"[^>]*Target="[^"]*\.svg"', rels)
        assert svg_rel, "the SVG must be a real relationship, not a dangling reference"
        assert f'r:embed="{svg_rel.group(1)}"' in xml

    # and the package is still readable -- a malformed extension would
    # make PowerPoint offer to repair the file
    assert len(Presentation(str(deck)).slides[0].shapes) == 1


def test_attaching_a_missing_svg_leaves_the_raster_picture_alone(tmp_path):
    from pptx import Presentation
    from pptx.util import Emu

    from deckguard.logo import _brand_asset, attach_svg

    png = _brand_asset("logo", light=False)
    if png is None:
        pytest.skip("vendored KONE marks not available")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    picture = slide.shapes.add_picture(png, Emu(0), Emu(0), Emu(914400), Emu(914400))
    assert attach_svg(picture, tmp_path / "nope.svg") is False
    assert picture.image.blob


def test_a_layout_whose_logo_is_a_placeholder_gets_one_stamped():
    """The master supplies its logo three ways and only one survives
    `add_slide`. 52 layouts carry a `<p:pic>` that paints through by
    itself; 15 carry `Logo Placeholder 9`, a BODY placeholder python-pptx
    never clones, so those slides -- the covers and every
    text-and-picture layout -- come out with no mark at all."""
    import posixpath

    from pptx import Presentation

    from deckguard.logo import stamp_logo_chrome

    master = _kone_master()
    if master is None:
        pytest.skip("kone-design master template not available")

    prs = Presentation(str(master))
    by_key = {
        posixpath.basename(layout.part.partname).replace(".xml", ""): layout
        for layout in prs.slide_layouts
    }

    # Cover A's logo is the placeholder kind; Title and content A's is a picture
    placeholder_layout, picture_layout = by_key["slideLayout1"], by_key["slideLayout16"]

    on_cover = prs.slides.add_slide(placeholder_layout)
    assert stamp_logo_chrome(on_cover) >= 1, "a cover must not come out logo-less"
    assert [s for s in on_cover.shapes if s.name == "Logo"]

    on_content = prs.slides.add_slide(picture_layout)
    assert stamp_logo_chrome(on_content) == 0, (
        "the layout already paints this one -- stamping again is the double-logo defect"
    )


def test_stamping_is_idempotent():
    import posixpath

    from pptx import Presentation

    from deckguard.logo import stamp_logo_chrome

    master = _kone_master()
    if master is None:
        pytest.skip("kone-design master template not available")
    prs = Presentation(str(master))
    layout = next(
        l for l in prs.slide_layouts
        if posixpath.basename(l.part.partname) == "slideLayout1.xml"
    )
    slide = prs.slides.add_slide(layout)
    first = stamp_logo_chrome(slide)
    assert first >= 1
    assert stamp_logo_chrome(slide) == 0


def test_a_full_bleed_photo_is_not_mistaken_for_a_logo():
    """Regression, caught by the gallery's own cut-cover test: a
    1280x422 banner completely contains the 81x31 logo slot, so scoring
    the overlap against the SMALLER box read the banner as a duplicate
    mark and deleted it. Against the larger box the same pair is 0.008.
    """
    from deckguard.logo import _boxes_overlap

    banner = (0, 0, 1280 * 9525, 422 * 9525)
    logo_slot = (45 * 9525, 45 * 9525, 81 * 9525, 31 * 9525)
    assert _boxes_overlap(banner, logo_slot) is False

    a_mark_in_the_slot = (46 * 9525, 46 * 9525, 80 * 9525, 30 * 9525)
    assert _boxes_overlap(a_mark_in_the_slot, logo_slot) is True
