from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR

from deckguard.legacy.exact_transplant import transplant_exact_treatment
from tests.helpers import add_rectangle, set_run, set_run_theme_color, set_theme_slot


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # "Blank" -- no placeholders to confuse shape lookups


def _by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(name)


def test_transplant_copies_fill_from_name_matched_reference_shape():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Box A", fill_hex="FF0000")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Box A", fill_hex="0000FF")

    result = transplant_exact_treatment(old_prs, new_prs)

    box = _by_name(old_prs.slides[0], "Box A")
    assert str(box.fill.fore_color.rgb) == "0000FF"
    assert len(result.changes) == 1
    assert result.changes[0].field == "fill"
    assert result.changes[0].old == "FF0000"
    assert result.changes[0].new == "0000FF"


def test_transplant_canonicalizes_reference_color_through_remap():
    """The reference deck can still carry a stale literal hex
    brand_rules.yaml already knows the approved replacement for --
    copying it verbatim would regress a color the generic brand pass
    already fixed, so it's canonicalized through colors.remap first."""
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Box A", fill_hex="F3EEE6")  # already the approved target

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Box A", fill_hex="EDEFF0")  # the legacy source hex

    config = {"colors": {"remap": {"#EDEFF0": "#F3EEE6"}}}
    result = transplant_exact_treatment(old_prs, new_prs, rules_config=config)

    box = _by_name(old_prs.slides[0], "Box A")
    assert str(box.fill.fore_color.rgb) == "F3EEE6"
    assert result.changes == []  # already matched the canonical value -- no-op


def test_transplant_leaves_shape_alone_with_no_name_match_in_reference():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Box A", fill_hex="FF0000", left_in=1)

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    # Different name AND different position -- no correspondence by
    # either signal (contrast this with the position-fallback tests
    # below, where a shared position DOES count as a match).
    add_rectangle(new_slide, name="Box B", fill_hex="0000FF", left_in=5)

    result = transplant_exact_treatment(old_prs, new_prs)

    box = _by_name(old_prs.slides[0], "Box A")
    assert str(box.fill.fore_color.rgb) == "FF0000"
    assert result.changes == []


def test_transplant_matches_by_position_when_name_does_not_survive():
    """Real-world case: PowerPoint's own auto-generated shape names are
    sequential by creation order, not stable across independently-edited
    copies of the same repeating design element (a footer "category
    chip" row copy-pasted across a whole catalog deck, say) -- the exact
    same visual chip can end up "Rectangle 27" on one slide and
    "Rectangle 10" on another. Falling back to position keeps the
    transplant working for exactly this pattern instead of silently
    finding no match at all."""
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Rectangle 27", fill_hex="FF0000", left_in=2, top_in=3)

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Rectangle 10", fill_hex="0000FF", left_in=2, top_in=3)

    result = transplant_exact_treatment(old_prs, new_prs)

    box = _by_name(old_prs.slides[0], "Rectangle 27")
    assert str(box.fill.fore_color.rgb) == "0000FF"
    assert len(result.changes) == 1


def test_transplant_flags_slide_with_low_shape_name_match_ratio():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Matched", fill_hex="FF0000", left_in=1)
    add_rectangle(old_slide, name="Unmatched 1", fill_hex="FF0000", left_in=2)
    add_rectangle(old_slide, name="Unmatched 2", fill_hex="FF0000", left_in=3)

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Matched", fill_hex="0000FF", left_in=1)
    # Different name AND different position from anything on the old
    # slide -- position fallback must not manufacture a match either.
    add_rectangle(new_slide, name="Different A", fill_hex="00FF00", left_in=6)
    add_rectangle(new_slide, name="Different B", fill_hex="00FF00", left_in=7)

    result = transplant_exact_treatment(old_prs, new_prs)

    assert result.flagged_slides == [0]


def test_transplant_does_not_flag_slide_with_high_shape_name_match_ratio():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    add_rectangle(old_slide, name="Box A", fill_hex="FF0000")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    add_rectangle(new_slide, name="Box A", fill_hex="0000FF")

    result = transplant_exact_treatment(old_prs, new_prs)

    assert result.flagged_slides == []


def test_transplant_copies_font_name_bold_and_color_from_matched_run():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    old_box = add_rectangle(old_slide, name="Box A")
    old_run = old_box.text_frame.paragraphs[0].add_run()
    set_run(old_run, text="Hello", font="Arial", bold=False, color_hex="141414")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A")
    new_run = new_box.text_frame.paragraphs[0].add_run()
    set_run(new_run, text="Bonjour", font="Inter", bold=True, color_hex="1450F5")

    config = {"fonts": {"approved": ["Inter"]}}
    result = transplant_exact_treatment(old_prs, new_prs, rules_config=config)

    run = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0].runs[0]
    assert run.text == "Hello"  # content stays verbatim
    assert run.font.name == "Inter"
    assert run.font.bold is True
    assert str(run.font.color.rgb) == "1450F5"
    fields = {c.field for c in result.changes}
    assert {"font_name", "font_bold", "font_color"} <= fields


def test_transplant_never_copies_a_font_the_reference_deck_uses_that_is_not_approved():
    """Hard rule, independent of the reference deck's own quality: this
    deck must only ever end up with KONE-approved fonts. If the reference
    itself carries something off-brand (never fixed, or hand-edited after
    the fact) with no known remap for it, the font_name copy is skipped
    entirely rather than propagating it -- unlike colors, there's no
    canonicalization fallback here, just a hard no."""
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    old_box = add_rectangle(old_slide, name="Box A")
    set_run(old_box.text_frame.paragraphs[0].add_run(), text="Hello", font="Inter")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A")
    set_run(new_box.text_frame.paragraphs[0].add_run(), text="Bonjour", font="Comic Sans MS")

    config = {"fonts": {"approved": ["Inter"], "remap": {}}}
    result = transplant_exact_treatment(old_prs, new_prs, rules_config=config)

    run = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Inter"  # left alone, never regressed to the reference's off-brand font
    assert not [c for c in result.changes if c.field == "font_name"]


def test_transplant_canonicalizes_reference_font_through_remap():
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    old_box = add_rectangle(old_slide, name="Box A")
    set_run(old_box.text_frame.paragraphs[0].add_run(), text="Hello", font="Inter")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A")
    set_run(new_box.text_frame.paragraphs[0].add_run(), text="Bonjour", font="Arial")  # legacy, but remapped

    config = {"fonts": {"approved": ["Inter"], "remap": {"Arial": "Inter"}}}
    result = transplant_exact_treatment(old_prs, new_prs, rules_config=config)

    run = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Inter"
    assert not [c for c in result.changes if c.field == "font_name"]  # already matched the canonical value


def test_transplant_copies_paragraph_alignment_from_matched_reference_shape():
    from pptx.enum.text import PP_ALIGN

    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    old_box = add_rectangle(old_slide, name="Box A")
    old_box.text_frame.paragraphs[0].add_run().text = "Code"
    old_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A")
    new_box.text_frame.paragraphs[0].add_run().text = "Code"
    new_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    result = transplant_exact_treatment(old_prs, new_prs)

    para = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0]
    assert para.alignment == PP_ALIGN.CENTER
    assert any(c.field == "alignment" for c in result.changes)


def test_transplant_skips_text_color_that_would_match_shapes_own_fill():
    """A duplicate/decoy shape in the reference deck can pair a text
    color with a fill that canonicalize to the SAME approved token --
    applying it verbatim would make the text invisible against its own
    background, so this is refused even though the fill copy itself
    (evaluated first) is applied normally."""
    old_prs = Presentation()
    old_slide = _blank_slide(old_prs)
    box = add_rectangle(old_slide, name="Box A", fill_hex="AAAAAA")
    set_run(box.text_frame.paragraphs[0].add_run(), text="Code", color_hex="141414")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Box A", fill_hex="F3EEE6")
    set_run(new_box.text_frame.paragraphs[0].add_run(), text="Code", color_hex="F3EEE6")

    result = transplant_exact_treatment(old_prs, new_prs)

    run = _by_name(old_prs.slides[0], "Box A").text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == "141414"  # left alone, not made invisible
    fill_changes = [c for c in result.changes if c.field == "fill"]
    assert fill_changes and fill_changes[0].new == "F3EEE6"  # fill copy still applied
    assert not [c for c in result.changes if c.field == "font_color"]


def test_transplant_does_not_match_a_real_placeholder_to_a_same_named_unrelated_shape():
    """Real-world bug: a deck's actual TITLE placeholder and an unrelated
    free-standing textbox on the reference deck's (differently-laid-out)
    slide can coincidentally share PowerPoint's auto-generated name
    ("Title 1") -- and even land at the same position, since the
    reference's decoy textbox often visually resembles a title. Matching
    by name (or position fallback) alone would transplant the decoy's
    color onto the real title; the placeholder-role check refuses this."""
    old_prs = Presentation()
    old_slide = old_prs.slides.add_slide(old_prs.slide_layouts[5])  # "Title Only"
    title = old_slide.shapes.title
    title.name = "Title 1"
    set_run(title.text_frame.paragraphs[0].add_run(), text="Elevator Plan", color_hex="141414")

    new_prs = Presentation()
    new_slide = _blank_slide(new_prs)  # "Blank" -- no placeholders at all
    decoy = add_rectangle(
        new_slide, name="Title 1",
        left_in=title.left / 914400, top_in=title.top / 914400,
        width_in=title.width / 914400, height_in=title.height / 914400,
    )
    set_run(decoy.text_frame.paragraphs[0].add_run(), text="Decoy", color_hex="1450F5")

    result = transplant_exact_treatment(old_prs, new_prs)

    run = old_prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == "141414"  # untouched -- not the decoy's blue
    assert not [c for c in result.changes if c.field == "font_color"]


def test_transplant_copies_scheme_color_uniformly_from_reference():
    """Real-world case: a generic contrast-fix pass upstream (fix_deck)
    forced an explicit literal color onto some, but not all, of a row of
    identically-styled "category chip" boxes -- each colored differently,
    each meant to keep the SAME theme-relative text color (e.g.
    background1) regardless of its own fill, per the reference deck's own
    design. Copying the reference's scheme reference itself (not just its
    currently-resolved RGB) restores that uniform, correct treatment."""
    old_prs = Presentation()
    set_theme_slot(old_prs, "lt1", "FFFFFF")
    old_slide = _blank_slide(old_prs)
    box = add_rectangle(old_slide, name="Chip", fill_hex="FFA023")
    set_run(box.text_frame.paragraphs[0].add_run(), text="Safety", color_hex="141414")  # wrongly forced explicit black

    new_prs = Presentation()
    set_theme_slot(new_prs, "lt1", "FFFFFF")
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Chip", fill_hex="FFA023")
    new_run = new_box.text_frame.paragraphs[0].add_run()
    new_run.text = "Safety"
    set_run_theme_color(new_run, MSO_THEME_COLOR.BACKGROUND_1)

    result = transplant_exact_treatment(old_prs, new_prs)

    run = _by_name(old_prs.slides[0], "Chip").text_frame.paragraphs[0].runs[0]
    assert run.font.color.theme_color == MSO_THEME_COLOR.BACKGROUND_1
    assert any(c.field == "font_color" for c in result.changes)


def test_transplant_scheme_color_no_op_when_already_matching():
    old_prs = Presentation()
    set_theme_slot(old_prs, "lt1", "FFFFFF")
    old_slide = _blank_slide(old_prs)
    box = add_rectangle(old_slide, name="Chip", fill_hex="FFA023")
    old_run = box.text_frame.paragraphs[0].add_run()
    old_run.text = "Safety"
    set_run_theme_color(old_run, MSO_THEME_COLOR.BACKGROUND_1)

    new_prs = Presentation()
    set_theme_slot(new_prs, "lt1", "FFFFFF")
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Chip", fill_hex="FFA023")
    new_run = new_box.text_frame.paragraphs[0].add_run()
    new_run.text = "Safety"
    set_run_theme_color(new_run, MSO_THEME_COLOR.BACKGROUND_1)

    result = transplant_exact_treatment(old_prs, new_prs)

    assert not [c for c in result.changes if c.field == "font_color"]


def test_transplant_skips_scheme_color_that_would_match_shapes_own_fill():
    old_prs = Presentation()
    set_theme_slot(old_prs, "lt1", "F3EEE6")
    old_slide = _blank_slide(old_prs)
    box = add_rectangle(old_slide, name="Chip", fill_hex="F3EEE6")
    set_run(box.text_frame.paragraphs[0].add_run(), text="Code", color_hex="141414")

    new_prs = Presentation()
    set_theme_slot(new_prs, "lt1", "F3EEE6")
    new_slide = _blank_slide(new_prs)
    new_box = add_rectangle(new_slide, name="Chip", fill_hex="F3EEE6")
    new_run = new_box.text_frame.paragraphs[0].add_run()
    new_run.text = "Code"
    set_run_theme_color(new_run, MSO_THEME_COLOR.BACKGROUND_1)  # resolves to F3EEE6 -- same as its own fill

    result = transplant_exact_treatment(old_prs, new_prs)

    run = _by_name(old_prs.slides[0], "Chip").text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == "141414"  # left alone, not made invisible
    assert not [c for c in result.changes if c.field == "font_color"]
