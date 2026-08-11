from copy import deepcopy

from pptx import Presentation

from deckguard.legacy import learn
from deckguard.config import default_config_path, load_config
from tests.helpers import add_rectangle, add_slide, body_run, new_deck, set_run, title_run

BASE_CONFIG = {
    "colors": {"approved": ["#1450F5", "#FFFFFF"], "remap": {}, "tolerance": 0},
    "fonts": {"approved": ["Inter", "Inter Semi Bold"], "remap": {}},
}


def _deck_with_run(font, color_hex, text="Body copy"):
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text=text, font=font, color_hex=color_hex)
    return prs


def test_high_confidence_color_and_font_match():
    old_prs = _deck_with_run("Arial", "AABBCC")
    new_prs = _deck_with_run("Inter", "1450F5")

    result = learn.learn(old_prs, new_prs, BASE_CONFIG)

    assert len(result.color_proposals) == 1
    cp = result.color_proposals[0]
    assert cp.old_hex == "AABBCC" and cp.new_hex == "1450F5"
    assert cp.confidence == "high"

    assert len(result.font_proposals) == 1
    fp = result.font_proposals[0]
    assert fp.old_font == "Arial" and fp.new_font == "Inter"
    assert fp.confidence == "high"


def test_low_confidence_when_counts_diverge():
    prs_old = new_deck()
    slide = add_slide(prs_old)
    set_run(title_run(slide), text="Title", font="Inter", color_hex="AABBCC")

    prs_new = new_deck()
    slide2 = add_slide(prs_new)
    for i in range(10):
        set_run(body_run(slide2) if i == 0 else title_run(slide2), text=f"t{i}", font="Inter", color_hex="1450F5")

    result = learn.learn(prs_old, prs_new, BASE_CONFIG)
    assert len(result.color_proposals) == 1
    assert result.color_proposals[0].confidence == "low"


def test_already_approved_or_remapped_colors_are_not_proposed():
    old_prs = _deck_with_run("Inter", "1450F5")  # already approved
    new_prs = _deck_with_run("Inter", "FFFFFF")
    result = learn.learn(old_prs, new_prs, BASE_CONFIG)
    assert result.color_proposals == []

    config_with_remap = {
        "colors": {"approved": ["#1450F5", "#FFFFFF"], "remap": {"#AABBCC": "#1450F5"}, "tolerance": 0},
        "fonts": {"approved": ["Inter"], "remap": {}},
    }
    old_prs2 = _deck_with_run("Inter", "AABBCC")
    new_prs2 = _deck_with_run("Inter", "FFFFFF")
    result2 = learn.learn(old_prs2, new_prs2, config_with_remap)
    assert result2.color_proposals == []  # already has a remap entry, not re-proposed


def test_unmatched_color_when_no_candidate_role_in_new_deck():
    old_prs = new_deck()
    slide = add_slide(old_prs)
    set_run(title_run(slide), text="Title", font="Inter", color_hex="1450F5")
    title_run(slide).text = "Title"
    # give the title shape a line color with no counterpart role in the new deck
    from pptx.dml.color import RGBColor

    old_prs.slides[0].shapes.title.line.color.rgb = RGBColor.from_string("AABBCC")

    new_prs = new_deck()
    add_slide(new_prs)  # no line color set anywhere

    result = learn.learn(old_prs, new_prs, BASE_CONFIG)
    assert any(role == "line" and hexval == "AABBCC" for role, hexval, _count in result.unmatched_old_colors)


def test_apply_learned_merges_high_confidence_only_by_default():
    old_prs = _deck_with_run("Arial", "AABBCC")
    new_prs = _deck_with_run("Inter", "1450F5")
    result = learn.learn(old_prs, new_prs, BASE_CONFIG)

    updated = learn.apply_learned(BASE_CONFIG, result)
    assert updated["colors"]["remap"]["#AABBCC"] == "#1450F5"
    assert "#1450F5" in updated["colors"]["approved"]  # already there, no dup
    assert updated["fonts"]["remap"]["Arial"] == "Inter"
    # original config untouched
    assert BASE_CONFIG["colors"]["remap"] == {}


def test_apply_learned_min_confidence_low_includes_low_confidence():
    prs_old = new_deck()
    slide = add_slide(prs_old)
    set_run(title_run(slide), text="Title", font="Inter", color_hex="AABBCC")
    prs_new = new_deck()
    slide2 = add_slide(prs_new)
    for i in range(10):
        set_run(body_run(slide2) if i == 0 else title_run(slide2), text=f"t{i}", font="Inter", color_hex="1450F5")

    result = learn.learn(prs_old, prs_new, BASE_CONFIG)
    updated_default = learn.apply_learned(BASE_CONFIG, result)
    assert updated_default["colors"]["remap"] == {}  # low confidence, not applied by default

    updated_low = learn.apply_learned(BASE_CONFIG, result, min_confidence="low")
    assert updated_low["colors"]["remap"]["#AABBCC"] == "#1450F5"


def test_write_learned_to_yaml_preserves_comments_and_dedupes(tmp_path):
    rules_path = tmp_path / "brand_rules.yaml"
    rules_path.write_text(
        "colors:\n"
        "  approved:\n"
        "    - \"#1450F5\"   # KONE Blue\n"
        "    - \"#FFFFFF\"   # White\n"
        "  remap: {}\n"
        "  tolerance: 0\n"
        "fonts:\n"
        "  approved:\n"
        "    - \"Inter\"   # brand font\n"
        "  remap: {}\n",
        encoding="utf-8",
    )

    old_prs = _deck_with_run("Arial", "AABBCC")
    new_prs = _deck_with_run("Inter", "1450F5")
    result = learn.learn(old_prs, new_prs, BASE_CONFIG)

    applied = learn.write_learned_to_yaml(rules_path, result)
    assert applied == 2  # one color, one font

    text = rules_path.read_text(encoding="utf-8")
    assert "# KONE Blue" in text
    assert "# brand font" in text
    assert "#AABBCC" in text and "#1450F5" in text.upper() or "1450F5" in text
    assert "Arial" in text

    # re-applying is idempotent -- nothing new to add
    applied_again = learn.write_learned_to_yaml(rules_path, result)
    assert applied_again == 0


def test_default_config_learn_smoke():
    # sanity: learn() runs cleanly against the shipped default config with
    # two trivially-identical decks (nothing off-brand to propose).
    config = load_config(default_config_path())
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Inter Semi Bold", color_hex="141414")
    result = learn.learn(prs, prs, config)
    assert result.color_proposals == []
    assert result.font_proposals == []


def test_exact_shape_match_beats_ambiguous_count_correlation():
    """Regression test for a real bug: two old colors (A, B) at similar
    overall counts to two new colors (A', B') can be matched *backwards*
    by whole-deck count correlation alone when the counts are ambiguous.
    Exact same-slide/same-shape-name evidence must resolve this correctly
    even when the naive count correlation would prefer the wrong pairing.
    """
    old_prs = new_deck()
    slide = add_slide(old_prs)
    # 3 shapes with color A, 2 with color B -- overall counts 3 vs 2
    add_rectangle(slide, name="Box 1", fill_hex="AAAAAA", left_in=0, top_in=0)
    add_rectangle(slide, name="Box 2", fill_hex="AAAAAA", left_in=1, top_in=0)
    add_rectangle(slide, name="Box 3", fill_hex="AAAAAA", left_in=2, top_in=0)
    add_rectangle(slide, name="Box 4", fill_hex="BBBBBB", left_in=0, top_in=1)
    add_rectangle(slide, name="Box 5", fill_hex="BBBBBB", left_in=1, top_in=1)

    new_prs = new_deck()
    slide2 = add_slide(new_prs)
    # Same shape names, but deliberately close overall counts for the
    # *other* pairing too (2 of A''s new color, 3 of B's), so a naive
    # global count-correlation matcher could easily pick A->new_count=2
    # and B->new_count=3 -- the WRONG pairing. Shape identity makes the
    # correct pairing unambiguous regardless.
    add_rectangle(slide2, name="Box 1", fill_hex="1450F5", left_in=0, top_in=0)
    add_rectangle(slide2, name="Box 2", fill_hex="1450F5", left_in=1, top_in=0)
    add_rectangle(slide2, name="Box 3", fill_hex="1450F5", left_in=2, top_in=0)
    add_rectangle(slide2, name="Box 4", fill_hex="7296F9", left_in=0, top_in=1)
    add_rectangle(slide2, name="Box 5", fill_hex="7296F9", left_in=1, top_in=1)

    result = learn.learn(old_prs, new_prs, BASE_CONFIG)
    by_old_hex = {p.old_hex: p for p in result.color_proposals}

    assert by_old_hex["AAAAAA"].new_hex == "1450F5"
    assert by_old_hex["AAAAAA"].confidence == "high"
    assert by_old_hex["BBBBBB"].new_hex == "7296F9"
    assert by_old_hex["BBBBBB"].confidence == "high"


def test_write_learned_to_yaml_preserves_sexagesimal_like_quoted_strings(tmp_path):
    """Regression test: ruamel interprets an unquoted `16:9`-style value as
    a base-60 YAML 1.1 integer (16*60+9=969) unless quotes are preserved.
    brand_rules.yaml's layout.slide_size is exactly this shape."""
    rules_path = tmp_path / "brand_rules.yaml"
    rules_path.write_text(
        "colors:\n"
        "  approved: [\"#1450F5\"]\n"
        "  remap: {}\n"
        "  tolerance: 0\n"
        "fonts:\n"
        "  approved: [\"Inter\"]\n"
        "  remap: {}\n"
        "layout:\n"
        "  slide_size: \"16:9\"\n",
        encoding="utf-8",
    )

    old_prs = _deck_with_run("Arial", "AABBCC")
    new_prs = _deck_with_run("Inter", "1450F5")
    result = learn.learn(old_prs, new_prs, BASE_CONFIG)
    learn.write_learned_to_yaml(rules_path, result)

    text = rules_path.read_text(encoding="utf-8")
    assert 'slide_size: "16:9"' in text
    assert "969" not in text


def _add_rectangle_to_layout(layout, name, fill_hex, left_in, top_in, width_in, height_in):
    """LayoutShapes has no add_shape() -- build on a scratch slide, then
    move the shape's XML element onto the layout's shape tree."""
    scratch = Presentation()
    scratch_slide = scratch.slides.add_slide(scratch.slide_layouts[6])
    shape = add_rectangle(
        scratch_slide, name=name, fill_hex=fill_hex,
        left_in=left_in, top_in=top_in, width_in=width_in, height_in=height_in,
    )
    elem = deepcopy(shape._element)
    layout.shapes._spTree.append(elem)
    return layout.shapes[-1]


def test_learn_proposes_layout_panel_fix_for_large_renamed_shape():
    """Regression test for a real bug: a shape's fill can live entirely on
    a slide LAYOUT (not any slide), and its name can drift between deck
    revisions (e.g. 'Rectangle 19' -> 'Rectangle 4') -- so it must be
    matched by position+size, not name, and proposed as a distinct
    layout_panel_remap entry rather than folded into ordinary colors.remap."""
    old_prs = new_deck()
    _add_rectangle_to_layout(old_prs.slide_layouts[1], "Rectangle 19", "EDEFF0", left_in=5, top_in=1, width_in=4, height_in=5)
    add_slide(old_prs, layout_idx=1)

    new_prs = new_deck()
    _add_rectangle_to_layout(new_prs.slide_layouts[1], "Rectangle 4", "1450F5", left_in=5, top_in=1, width_in=4, height_in=5)
    add_slide(new_prs, layout_idx=1)

    result = learn.learn(old_prs, new_prs, BASE_CONFIG)

    assert len(result.layout_panel_proposals) == 1
    p = result.layout_panel_proposals[0]
    assert p.old_hex == "EDEFF0" and p.new_hex == "1450F5"
    assert p.old_shape_name == "Rectangle 19" and p.new_shape_name == "Rectangle 4"
    assert p.confidence == "high"
    # Not a colors.remap proposal -- EDEFF0 is small-shape-safe elsewhere.
    assert not any(cp.old_hex == "EDEFF0" for cp in result.color_proposals)


def test_learn_ignores_small_layout_shapes_below_area_threshold():
    old_prs = new_deck()
    _add_rectangle_to_layout(old_prs.slide_layouts[1], "Tab", "EDEFF0", left_in=1, top_in=1, width_in=1, height_in=0.5)
    add_slide(old_prs, layout_idx=1)

    new_prs = new_deck()
    _add_rectangle_to_layout(new_prs.slide_layouts[1], "Tab", "1450F5", left_in=1, top_in=1, width_in=1, height_in=0.5)
    add_slide(new_prs, layout_idx=1)

    result = learn.learn(old_prs, new_prs, BASE_CONFIG)

    assert result.layout_panel_proposals == []


def test_apply_learned_writes_layout_panel_remap_separately_from_colors_remap():
    old_prs = new_deck()
    _add_rectangle_to_layout(old_prs.slide_layouts[1], "Rectangle 19", "EDEFF0", left_in=5, top_in=1, width_in=4, height_in=5)
    add_slide(old_prs, layout_idx=1)

    new_prs = new_deck()
    _add_rectangle_to_layout(new_prs.slide_layouts[1], "Rectangle 4", "1450F5", left_in=5, top_in=1, width_in=4, height_in=5)
    add_slide(new_prs, layout_idx=1)

    result = learn.learn(old_prs, new_prs, BASE_CONFIG)
    merged = learn.apply_learned(BASE_CONFIG, result, min_confidence="high")

    assert merged["colors"]["layout_panel_remap"] == {"#EDEFF0": "#1450F5"}
    assert "EDEFF0" not in merged["colors"]["remap"]
