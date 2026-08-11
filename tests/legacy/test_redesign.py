"""Tests for deckguard's AI-assisted redesign (redesign.py): extracting
eligible content from an arbitrary deck, asking Claude to map it onto
compose.py's outline schema, then building it through the same
deterministic pipeline `deckguard create` uses.

No real network/API calls -- every test injects a fake Anthropic-shaped
client, consistent with the rest of this suite's "works with no API key"
philosophy for anything that doesn't require a live account.
"""

import json

import anthropic
import httpx
import pytest
from pptx import Presentation
from pptx.util import Inches

from deckguard.legacy.redesign import (
    RedesignError,
    Usage,
    call_claude_for_brand_review,
    call_claude_for_outline,
    extract_eligible_slides,
    partition_skipped,
    redesign_deck,
)
from deckguard.legacy.skill_bridge import _skill_dir
from deckguard.legacy.slide_import import default_template_path
from tests.helpers import add_picture, add_rectangle, add_slide, body_run, make_pattern_png, new_deck, title_run

TEMPLATE_PATH = default_template_path()
pytestmark = pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="bundled template asset not present")


def _kone_skill_available() -> bool:
    """The kone-deck-generator skill is an out-of-repo, per-machine
    install (see skill_bridge.py's own docstring) -- tests exercising it
    for real (fake AI response, real deterministic render) skip cleanly
    where it isn't present, same as this file already does for the
    bundled org template."""
    return (_skill_dir() / "kone_deck_creator.py").is_file()


# --------------------------------------------------------------------------
# fake Anthropic client
# --------------------------------------------------------------------------


class _FakeTextBlock:
    type = "text"

    def __init__(self, text):
        self.text = text


class _FakeUsage:
    def __init__(self, input_tokens, output_tokens):
        self.input_tokens = input_tokens
        self.output_tokens = output_tokens


class _FakeResponse:
    def __init__(self, content_text, stop_reason="end_turn", input_tokens=1000, output_tokens=500):
        self.content = [_FakeTextBlock(content_text)] if content_text is not None else []
        self.stop_reason = stop_reason
        self.usage = _FakeUsage(input_tokens, output_tokens)


class _FakeStreamCM:
    def __init__(self, response):
        self._response = response

    def __enter__(self):
        return self

    def __exit__(self, *exc):
        return False

    def get_final_message(self):
        return self._response


class _FakeMessages:
    def __init__(self, response):
        self._response = response
        self.calls = []

    def stream(self, **kwargs):
        self.calls.append(kwargs)
        return _FakeStreamCM(self._response)


class _FakeClient:
    def __init__(self, response):
        self.messages = _FakeMessages(response)


class _SequencedFakeMessages:
    """Returns a different canned response per call, in order -- for
    tests exercising redesign_deck's two model calls (the outline plan,
    then the archetype-override pass), which a single-response
    _FakeClient can't distinguish between. Repeats the last response if
    called more times than provided."""
    def __init__(self, responses: list):
        self._responses = responses
        self.calls = []

    def stream(self, **kwargs):
        self.calls.append(kwargs)
        response = self._responses[min(len(self.calls) - 1, len(self._responses) - 1)]
        return _FakeStreamCM(response)


class _SequencedFakeClient:
    def __init__(self, responses: list):
        self.messages = _SequencedFakeMessages(responses)


def _fake_api_status_error(status_code: int, message: str) -> anthropic.APIStatusError:
    body = {"type": "error", "error": {"type": "overloaded_error", "message": message}}
    resp = httpx.Response(
        status_code, request=httpx.Request("POST", "https://api.anthropic.com/v1/messages"), json=body
    )
    return anthropic.APIStatusError(message, response=resp, body=body)


class _RaisingMessages:
    def __init__(self, exc):
        self._exc = exc

    def stream(self, **kwargs):
        raise self._exc


class _RaisingClient:
    def __init__(self, exc):
        self.messages = _RaisingMessages(exc)


def _outline_json(*items):
    return json.dumps({"slides": list(items)})


def _slide_item(index, kind="content", **overrides):
    item = {
        "source_slide_index": index,
        "kind": kind,
        "title": "A title",
        "subtitle": None,
        "bullets": ["Point one", "Point two"],
        "columns": [],
        "quote_text": None,
        "quote_author": None,
        "quote_label": None,
        "stats": [],
        "milestones": [],
        "variant": None,
    }
    item.update(overrides)
    return item


def _kone_spec_json(title, *slides):
    return json.dumps({"title": title, "slides": list(slides)})


def _kone_slide(archetype, **content):
    """A kone-deck-generator spec slide: {"archetype": <name>, ...that
    archetype's own content fields...} -- content shape varies per
    archetype (see skill_bridge.py's own docstring), so this takes
    whatever content keys the caller passes, unlike the old fixed
    layout schema's always-present, often-null fields."""
    return {"archetype": archetype, **content}


# --------------------------------------------------------------------------
# extract_eligible_slides
# --------------------------------------------------------------------------


def test_extract_eligible_slides_separates_ordinary_from_disqualified():
    prs = new_deck()
    ok_slide = add_slide(prs)
    title_run(ok_slide).text = "Ordinary slide"
    body_run(ok_slide).text = "Some content"

    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))

    eligible, skipped = extract_eligible_slides(prs)

    assert [i for i, _p in eligible] == [1]
    assert [s.slide_index for s in skipped] == [2]
    assert "table" in skipped[0].reason


def test_extract_eligible_slides_empty_deck_slide_is_skipped():
    prs = new_deck()
    add_slide(prs)  # title/body placeholders present but empty
    _eligible, skipped = extract_eligible_slides(prs)
    assert len(skipped) == 1


def test_extract_eligible_slides_accepts_dense_text_retemplate_would_skip():
    """Regression test: a hand-built slide with more text boxes than any
    layout has placeholders used to be skipped by redesign too (it
    reused retemplate.classify_slide's MAX_TEXT_BLOCKS=3 cap verbatim),
    even though redesign is allowed to condense wording and retemplate
    is not. redesign should accept this and let the model condense it,
    not skip it."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    for i in range(5):  # one more than retemplate.MAX_TEXT_BLOCKS (3)
        box = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(3), Inches(0.5))
        box.text_frame.text = f"Block {i}"

    eligible, skipped = extract_eligible_slides(prs)

    assert len(eligible) == 1
    assert skipped == []
    assert len(eligible[0][1].text_blocks) == 5


def test_extract_eligible_slides_accepts_many_small_boxes_regardless_of_count():
    """Block COUNT alone, however high, is never a reason to skip --
    only total text volume is (see test below). Twenty tiny caption
    boxes have far less content than a handful of paragraph-sized ones,
    and redesign condenses either way."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    for i in range(20):
        box = slide.shapes.add_textbox(Inches(1), Inches(1.0 + 0.3 * i), Inches(3), Inches(0.25))
        box.text_frame.text = f"Block {i}"

    eligible, skipped = extract_eligible_slides(prs)

    assert skipped == []
    assert len(eligible[0][1].text_blocks) == 20


def test_extract_eligible_slides_still_skips_truly_excessive_text_volume():
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    # A handful of boxes, but each stuffed with enough text to blow past
    # REDESIGN_MAX_TEXT_CHARS -- volume, not count, is what's capped.
    huge_text = "Lorem ipsum dolor sit amet. " * 800  # ~23,000 characters
    for i in range(3):
        box = slide.shapes.add_textbox(Inches(1), Inches(1.0 + 2.0 * i), Inches(3), Inches(1.5))
        box.text_frame.text = huge_text

    _eligible, skipped = extract_eligible_slides(prs)

    assert len(skipped) == 1
    assert "reasonably split" in skipped[0].reason


def test_extract_eligible_slides_still_hard_skips_tables_regardless_of_text_cap():
    """A brief -- or redesign's higher text cap -- never overrides the
    shape-type safety rules shared with retemplate."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))

    _eligible, skipped = extract_eligible_slides(prs)

    assert len(skipped) == 1
    assert "table" in skipped[0].reason


# --------------------------------------------------------------------------
# call_claude_for_outline
# --------------------------------------------------------------------------


def test_call_claude_for_outline_parses_valid_response():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Quarterly review"
    body_run(slide).text = "Revenue up"
    eligible, _skipped = extract_eligible_slides(prs)

    response_json = _outline_json(_slide_item(1, title="Quarterly review"))
    client = _FakeClient(_FakeResponse(response_json, input_tokens=1200, output_tokens=600))

    raw_slides, usage = call_claude_for_outline(eligible, client=client)

    assert raw_slides[0]["title"] == "Quarterly review"
    assert usage.input_tokens == 1200
    assert usage.output_tokens == 600
    assert usage.model == "claude-opus-5"
    # (1200/1e6)*5 + (600/1e6)*25 = 0.006 + 0.015 = 0.021
    assert usage.estimated_cost_usd == pytest.approx(0.021)


def test_call_claude_for_outline_passes_model_and_effort_through():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _FakeClient(_FakeResponse(_outline_json(_slide_item(1))))
    call_claude_for_outline(eligible, model="claude-sonnet-5", effort="medium", client=client)

    call_kwargs = client.messages.calls[0]
    assert call_kwargs["model"] == "claude-sonnet-5"
    assert call_kwargs["output_config"]["effort"] == "medium"
    assert call_kwargs["output_config"]["format"]["type"] == "json_schema"


def test_call_claude_for_outline_raises_on_refusal():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _FakeClient(_FakeResponse(None, stop_reason="refusal"))
    with pytest.raises(RedesignError, match="declined"):
        call_claude_for_outline(eligible, client=client)


def test_call_claude_for_outline_raises_clean_message_on_overloaded_error():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _RaisingClient(_fake_api_status_error(529, "Overloaded"))
    with pytest.raises(RedesignError, match="temporarily rate-limited or overloaded"):
        call_claude_for_outline(eligible, client=client)


def test_call_claude_for_outline_raises_clean_message_on_bad_request_error():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _RaisingClient(_fake_api_status_error(400, "adaptive thinking is not supported on this model"))
    with pytest.raises(RedesignError, match="Claude API error"):
        call_claude_for_outline(eligible, client=client)


def test_call_claude_for_outline_raises_on_invalid_json():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _FakeClient(_FakeResponse("not json"))
    with pytest.raises(RedesignError, match="not valid JSON"):
        call_claude_for_outline(eligible, client=client)


def test_call_claude_for_outline_raises_on_missing_slides_key():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _FakeClient(_FakeResponse(json.dumps({"oops": []})))
    with pytest.raises(RedesignError, match="'slides'"):
        call_claude_for_outline(eligible, client=client)


# --------------------------------------------------------------------------
# redesign_deck (end to end)
# --------------------------------------------------------------------------


def test_redesign_deck_builds_a_valid_composed_deck(tmp_path):
    prs = new_deck()
    cover = add_slide(prs, layout_idx=0)
    title_run(cover).text = "Annual Review"

    content = add_slide(prs)
    title_run(content).text = "Highlights"
    body_run(content).text = "Grew nicely"

    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))

    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    response_json = _outline_json(
        _slide_item(1, kind="cover", title="Annual Review", subtitle="A great year", bullets=[]),
        _slide_item(2, kind="content", title="Highlights", bullets=["Grew nicely"]),
    )
    client = _FakeClient(_FakeResponse(response_json, input_tokens=2000, output_tokens=800))

    out_path = tmp_path / "redesigned.pptx"
    compose_result, redesign_result = redesign_deck(str(src_path), str(out_path), client=client)

    assert compose_result.slide_count == 2
    assert compose_result.layouts_used[0] == "Cover B"
    assert len(redesign_result.skipped) == 1
    assert redesign_result.skipped[0].slide_index == 3
    assert redesign_result.usage.input_tokens == 2000

    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 2


@pytest.mark.skipif(not _kone_skill_available(), reason="kone-deck-generator skill not installed")
def test_redesign_deck_applies_an_archetype_override_for_a_strong_fit_slide(tmp_path):
    """The archetype coexistence pass: a second model call gets a shot
    at any content/stat/timeline slide, and when it names a real
    archetype that slide renders through kone_engine instead of an
    org-template layout -- verified here both via the reported
    layouts_used and by checking the archetype's own content actually
    landed in the output file."""
    prs = new_deck()
    cover = add_slide(prs, layout_idx=0)
    title_run(cover).text = "Annual Review"
    content = add_slide(prs)
    title_run(content).text = "Resolution rate"
    body_run(content).text = "91.2% of requests resolved"
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    outline_response = _outline_json(
        _slide_item(1, kind="cover", title="Annual Review", subtitle="A great year", bullets=[]),
        _slide_item(2, kind="content", title="Resolution rate", bullets=["91.2% of requests resolved"]),
    )
    override_response = json.dumps({"overrides": [
        {"outline_index": 1, "archetype": "hero_stat", "eyebrow": "Resolution rate", "value": "91.2%",
         "caption": "of all requests cleared within the focus period", "support": "91.2% of requests resolved"},
    ]})
    client = _SequencedFakeClient([_FakeResponse(outline_response), _FakeResponse(override_response)])

    out_path = tmp_path / "redesigned.pptx"
    compose_result, _redesign_result = redesign_deck(str(src_path), str(out_path), client=client)

    assert compose_result.slide_count == 2
    assert compose_result.layouts_used == ["Cover B", "hero_stat"]
    assert len(client.messages.calls) == 2

    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 2
    body_texts = [
        shape.text_frame.text for shape in out_prs.slides[1].shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    ]
    assert any("91.2%" in t for t in body_texts)


@pytest.mark.skipif(not _kone_skill_available(), reason="kone-deck-generator skill not installed")
def test_redesign_deck_ignores_an_archetype_override_call_that_errors(tmp_path):
    """Fail-closed: if the second (archetype-override) call blows up,
    the redesign must still succeed on the outline plan alone -- this
    pass is a pure quality enhancement, never a hard requirement."""
    prs = new_deck()
    content = add_slide(prs)
    title_run(content).text = "Highlights"
    body_run(content).text = "Grew nicely"
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    outline_response = _outline_json(
        _slide_item(1, kind="content", title="Highlights", bullets=["Grew nicely"]),
    )

    class _FailOnSecondCallMessages:
        def __init__(self, first_response):
            self._first_response = first_response
            self.calls = []

        def stream(self, **kwargs):
            self.calls.append(kwargs)
            if len(self.calls) == 1:
                return _FakeStreamCM(self._first_response)
            raise RuntimeError("simulated transient failure")

    class _FailOnSecondCallClient:
        def __init__(self, first_response):
            self.messages = _FailOnSecondCallMessages(first_response)

    client = _FailOnSecondCallClient(_FakeResponse(outline_response))

    out_path = tmp_path / "redesigned.pptx"
    compose_result, _redesign_result = redesign_deck(str(src_path), str(out_path), client=client)

    assert compose_result.slide_count == 1
    assert compose_result.layouts_used == ["Title and content A"]
    assert len(client.messages.calls) == 2  # the override call was attempted, and its failure was swallowed


def test_redesign_deck_allows_one_source_slide_to_split_across_multiple_output_slides(tmp_path):
    """Direct response to explicit direction: AI mode must never reword
    or condense real source content -- when it doesn't fit one layout,
    it should split across multiple output slides instead (all sharing
    the same source_slide_index), never dropping or paraphrasing
    anything. This is a prompt-level instruction the model follows, not
    something the schema enforces, so this test drives the pipeline
    with a fake response that already reflects that -- the real
    assertion is that redesign_deck's plumbing (outline building,
    image backfill) tolerates and correctly handles multiple entries
    sharing one source index."""
    img_path = make_pattern_png(tmp_path / "img.png", seed=6)

    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Dense slide"
    for i in range(5):
        box = slide.shapes.add_textbox(Inches(1), Inches(0.3 * i), Inches(3), Inches(0.25))
        box.text_frame.text = f"Point {i}"
    add_picture(slide, str(img_path))
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    response_json = _outline_json(
        _slide_item(1, kind="content", title="Dense slide (1/2)", bullets=["Point 0", "Point 1", "Point 2"]),
        _slide_item(1, kind="content", title="Dense slide (2/2)", bullets=["Point 3", "Point 4"]),
    )
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "redesigned.pptx"
    compose_result, _redesign_result = redesign_deck(str(src_path), str(out_path), client=client)

    assert compose_result.slide_count == 2
    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 2
    all_text = "\n".join(
        shp.text_frame.text for slide in out_prs.slides for shp in slide.shapes
        if shp.has_text_frame and shp.text_frame.text.strip()
    )
    for i in range(5):
        assert f"Point {i}" in all_text

    # the source slide's one image is attached to only the FIRST split
    # entry, never duplicated onto both
    def _image_blobs(slide):
        blobs = []
        for shp in slide.shapes:
            try:
                blobs.append(shp.image.blob)
            except (AttributeError, ValueError):
                continue
        return blobs

    slide_image_counts = [len(_image_blobs(s)) for s in out_prs.slides]
    assert sorted(slide_image_counts) == [0, 1]


def test_redesign_deck_carries_source_images_into_the_redesigned_slide(tmp_path):
    """Regression test: the model's outline schema has no field for images
    at all (it's never shown the pixels), so images used to be silently
    dropped for every redesigned slide -- a real bug reported against a
    genuinely image-heavy source deck. Images should now survive, backfilled
    deterministically by source_slide_index after the model call."""
    prs = new_deck()
    content = add_slide(prs)
    title_run(content).text = "Highlights"
    body_run(content).text = "Grew nicely"
    img_path = make_pattern_png(tmp_path / "img.png", seed=3)
    add_picture(content, str(img_path))

    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    response_json = _outline_json(
        _slide_item(1, kind="content", title="Highlights", bullets=["Grew nicely"]),
    )
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "redesigned.pptx"
    compose_result, _redesign_result = redesign_deck(str(src_path), str(out_path), client=client)

    assert compose_result.slide_count == 1
    out_prs = Presentation(str(out_path))
    blobs = []
    for shp in out_prs.slides[0].shapes:
        try:
            blobs.append(shp.image.blob)
        except (AttributeError, ValueError):
            continue
    assert img_path.read_bytes() in blobs


def test_redesign_deck_raises_when_nothing_is_eligible_and_no_brief(tmp_path):
    prs = new_deck()
    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    client = _FakeClient(_FakeResponse(_outline_json()))
    with pytest.raises(RedesignError, match="nothing to work with"):
        redesign_deck(str(src_path), str(tmp_path / "out.pptx"), client=client)


def test_redesign_deck_raises_with_no_deck_and_no_brief(tmp_path):
    client = _FakeClient(_FakeResponse(_outline_json()))
    with pytest.raises(RedesignError, match="nothing to work with"):
        redesign_deck(None, str(tmp_path / "out.pptx"), client=client)


# --------------------------------------------------------------------------
# partition_skipped -- blank slides vs genuinely-unsafe content
# --------------------------------------------------------------------------


def test_partition_skipped_splits_blank_from_real_skips():
    prs = new_deck()
    add_slide(prs)  # empty -> blank

    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))  # real skip

    _eligible, skipped = extract_eligible_slides(prs)
    blank_indices, real_skipped = partition_skipped(skipped)

    assert blank_indices == [1]
    assert [s.slide_index for s in real_skipped] == [2]
    assert "table" in real_skipped[0].reason


# --------------------------------------------------------------------------
# redesign_deck -- brief-driven modes (fill blanks / build from scratch)
# --------------------------------------------------------------------------


@pytest.mark.skipif(not _kone_skill_available(), reason="kone-deck-generator skill not installed")
def test_redesign_deck_builds_from_scratch_with_no_source_deck(tmp_path):
    """The pure "Claude, design me a deck" mode -- no deck_path at all --
    routes through the kone-deck-generator skill (skill_bridge.py), not
    compose.py's outline path: see redesign_deck's own branch for why."""
    response_json = _kone_spec_json(
        "Predictive Maintenance",
        _kone_slide(
            "agenda_contents", title="Agenda",
            items=[{"number": "01", "item": "Why predictive maintenance"}, {"number": "02", "item": "What changes"}],
        ),
        _kone_slide(
            "three_stats", title="Predictive maintenance pays for itself fast.",
            stats=[
                {"label": "Downtime", "value": "-30%", "desc": "fewer unplanned callouts."},
                {"label": "Cost", "value": "-18%", "desc": "lower maintenance spend."},
                {"label": "Payback", "value": "9 mo", "desc": "typical time to break even."},
            ],
        ),
    )
    client = _FakeClient(_FakeResponse(response_json, input_tokens=1800, output_tokens=900))

    out_path = tmp_path / "from_scratch.pptx"
    compose_result, redesign_result = redesign_deck(
        None, str(out_path), brief="A short deck on predictive maintenance for facilities managers.", client=client
    )

    assert compose_result.slide_count == 4  # retained Cover F + 2 body + retained Outro
    assert redesign_result.skipped == []
    assert len(Presentation(str(out_path)).slides) == 4

    sent_content = client.messages.calls[0]["messages"][0]["content"]
    assert "predictive maintenance" in sent_content.lower()


def test_redesign_deck_fills_blank_slides_alongside_real_content(tmp_path):
    prs = new_deck()
    content_slide = add_slide(prs)
    title_run(content_slide).text = "Existing highlight"
    body_run(content_slide).text = "Real content here"

    add_slide(prs)  # blank -- should be authored from the brief, not skipped

    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))  # real skip, brief or not

    src_path = tmp_path / "half_empty.pptx"
    prs.save(str(src_path))

    response_json = _outline_json(
        _slide_item(1, kind="content", title="Existing highlight", bullets=["Real content here"]),
        _slide_item(None, kind="stat", title="By the numbers", stats=[{"number": "3", "label": "sites live"}], bullets=[]),
    )
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "filled.pptx"
    compose_result, redesign_result = redesign_deck(
        str(src_path), str(out_path), brief="Add a stats slide about rollout progress.", client=client
    )

    assert compose_result.slide_count == 2
    # the table slide is still reported -- a brief never overrides "unsafe to touch"
    assert [s.slide_index for s in redesign_result.skipped] == [3]

    call_kwargs = client.messages.calls[0]
    sent_content = call_kwargs["messages"][0]["content"]
    assert "1 blank slide" in sent_content
    assert "rollout progress" in sent_content.lower()


def test_call_claude_for_outline_target_slides_guidance_is_included():
    client = _FakeClient(_FakeResponse(_outline_json(_slide_item(None))))
    call_claude_for_outline([], blank_count=0, brief="A topic", target_slides=7, client=client)
    sent_content = client.messages.calls[0]["messages"][0]["content"]
    assert "approximately 7 total slides" in sent_content


def test_usage_estimated_cost_unknown_model_is_zero():
    usage = Usage(input_tokens=1000, output_tokens=1000, model="some-future-model")
    assert usage.estimated_cost_usd == 0.0


# --------------------------------------------------------------------------
# redesign_deck(mode="brand") -- deterministic, no LLM call at all
# --------------------------------------------------------------------------


def test_redesign_deck_brand_mode_makes_no_client_call(tmp_path):
    """The whole point of brand mode: it never touches the LLM path at
    all. A client that would raise if ever called proves it."""
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Highlights"
    body_run(slide).text = "Grew nicely"
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    class _ExplodingClient:
        @property
        def messages(self):
            raise AssertionError("brand mode must never touch the Anthropic client")

    out_path = tmp_path / "rebranded.pptx"
    compose_result, redesign_result = redesign_deck(
        str(src_path), str(out_path), mode="brand", client=_ExplodingClient()
    )

    assert compose_result.slide_count == 1
    assert redesign_result.usage.model == "none"
    assert redesign_result.usage.estimated_cost_usd == 0.0

    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 1


def test_redesign_deck_brand_mode_requires_a_deck(tmp_path):
    with pytest.raises(RedesignError, match="needs a source deck"):
        redesign_deck(None, str(tmp_path / "out.pptx"), mode="brand")


def test_redesign_deck_brand_mode_rejects_a_brief(tmp_path):
    prs = new_deck()
    title_run(add_slide(prs)).text = "Deck"
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    with pytest.raises(RedesignError, match="never authors content"):
        redesign_deck(str(src_path), str(tmp_path / "out.pptx"), mode="brand", brief="a brief")


def test_redesign_deck_rejects_unknown_mode(tmp_path):
    with pytest.raises(RedesignError, match="mode must be"):
        redesign_deck(str(tmp_path / "x.pptx"), str(tmp_path / "out.pptx"), mode="bogus")


def test_redesign_deck_rewrite_mode_rejects_reference_path(tmp_path):
    with pytest.raises(RedesignError, match="reference_path only applies to mode='brand'"):
        redesign_deck(str(tmp_path / "x.pptx"), str(tmp_path / "out.pptx"), mode="rewrite", reference_path="ref.pptx")


def _add_grouped_chip(slide, name: str, fill_hex: str):
    """A rectangle inside a group -- GROUP is one of apply_rebrand's
    disqualifying shape types, so a slide built this way is left
    untouched (its original shape names intact) rather than rebuilt onto
    a template layout, matching the real-world "category chip" case
    exact_transplant.py's own docstring describes."""
    from pptx.util import Emu

    group = slide.shapes.add_group_shape()
    group.name = f"{name} group"
    rect = add_rectangle(group, name=name, fill_hex=fill_hex, left_in=1, top_in=1)
    group.left, group.top, group.width, group.height = Emu(0), Emu(0), rect.width, rect.height
    return rect


def test_redesign_deck_brand_mode_reports_exact_reference_match(tmp_path):
    """`reference_path` runs the exact-transplant pass (see
    exact_transplant.py) as brand mode's last step -- a per-run override
    that copies the reference's own per-shape treatment, never touching
    brand_rules.yaml. Findings land in `reference_match_notes`."""
    from tests.helpers import add_rectangle

    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Highlights"
    body_run(slide).text = "Grew nicely"
    _add_grouped_chip(slide, "Chip", "FF0000")
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    ref_prs = new_deck()
    ref_slide = add_slide(ref_prs)
    title_run(ref_slide).text = "Highlights"
    body_run(ref_slide).text = "Grew nicely"
    _add_grouped_chip(ref_slide, "Chip", "0000FF")
    ref_path = tmp_path / "reference.pptx"
    ref_prs.save(str(ref_path))

    out_path = tmp_path / "rebranded.pptx"
    compose_result, redesign_result = redesign_deck(
        str(src_path), str(out_path), mode="brand", reference_path=str(ref_path)
    )

    assert any("Reference match:" in n for n in redesign_result.reference_match_notes)
    out_prs = Presentation(str(out_path))
    group = next(s for s in out_prs.slides[0].shapes if s.shape_type == 6)
    chip = next(s for s in group.shapes if s.name == "Chip")
    assert str(chip.fill.fore_color.rgb) == "0000FF"


# --------------------------------------------------------------------------
# redesign_deck(mode="brand", review=True) -- the small, optional AI pass
# --------------------------------------------------------------------------


def _review_json(*items):
    return json.dumps({"slides": list(items)})


def _decorative_overload_deck(tmp_path, title="Appendix"):
    """A slide with too many free-form decorative shapes to verbatim-carry
    (retemplate's own MAX_DECORATIVE_SHAPES cap) -- disqualified regardless
    of how little real content it has, exactly the shape a genuine
    divider/transition slide takes in a real deck: a short title and
    little else besides decoration."""
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = title
    for _ in range(6):
        add_rectangle(slide)
    path = tmp_path / "source.pptx"
    prs.save(str(path))
    return path


def test_call_claude_for_brand_review_parses_response():
    from deckguard.legacy.retemplate import SlideProposal

    proposal = SlideProposal(
        slide_index=3, eligible=False, reason="too many free-form shapes to safely reflow",
        title_preview="Appendix", body_preview=None, image_count=0,
    )
    response_json = _review_json({"slide_index": 3, "is_divider": True, "divider_title": "Appendix", "note": None})
    client = _FakeClient(_FakeResponse(response_json, input_tokens=200, output_tokens=50))

    raw_slides, usage = call_claude_for_brand_review([proposal], client=client)

    assert raw_slides == [{"slide_index": 3, "is_divider": True, "divider_title": "Appendix", "note": None}]
    assert usage.input_tokens == 200


def test_call_claude_for_brand_review_no_op_with_nothing_to_review():
    raw_slides, usage = call_claude_for_brand_review([], client=_FakeClient(_FakeResponse("{}")))
    assert raw_slides == []
    assert usage.input_tokens == 0


def test_redesign_deck_brand_review_rebuilds_a_divider_slide_with_ai_suggested_title(tmp_path):
    src_path = _decorative_overload_deck(tmp_path)
    response_json = _review_json({"slide_index": 1, "is_divider": True, "divider_title": "Appendix", "note": None})
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "reviewed.pptx"
    compose_result, redesign_result = redesign_deck(
        str(src_path), str(out_path), mode="brand", review=True, client=client
    )

    assert compose_result.slide_count == 1
    assert "Section divider A" in compose_result.layouts_used
    assert redesign_result.skipped == []  # rebuilt, no longer reported as skipped

    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 1
    assert out_prs.slides[0].shapes.title.text_frame.text == "Appendix"


@pytest.mark.skipif(not _kone_skill_available(), reason="kone-deck-generator skill not installed")
def test_redesign_deck_brand_review_applies_an_archetype_override_to_an_accepted_slide(tmp_path):
    """mode='brand', review=True: an already-accepted, ordinary-content
    slide can get swapped for a KONE archetype -- the same coexistence
    idea as redesign_deck's AI-rewrite path, reused here behind the
    same opt-in --review flag brand mode already has, so brand mode's
    own "fully deterministic by default" identity stays intact."""
    prs = new_deck()
    cover = add_slide(prs)
    title_run(cover).text = "Annual Review"
    content = add_slide(prs)
    title_run(content).text = "Resolution rate"
    body_run(content).text = "91.2% of requests resolved"
    end = add_slide(prs)
    title_run(end).text = "Thank you"
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    # No slide here is ineligible, so call_claude_for_brand_review is
    # never invoked -- this is the ONLY model call in this run.
    override_response = json.dumps({"overrides": [
        {"outline_index": 2, "archetype": "hero_stat", "eyebrow": "Resolution rate", "value": "91.2%",
         "caption": "of all requests cleared within the focus period", "support": "91.2% of requests resolved"},
    ]})
    client = _FakeClient(_FakeResponse(override_response))

    out_path = tmp_path / "reviewed.pptx"
    compose_result, _redesign_result = redesign_deck(
        str(src_path), str(out_path), mode="brand", review=True, client=client
    )

    assert "hero_stat" in compose_result.layouts_used
    assert len(client.messages.calls) == 1

    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 3
    body_texts = [
        shape.text_frame.text for shape in out_prs.slides[1].shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    ]
    assert any("91.2%" in t for t in body_texts)
    # The cover/end positions are never offered up for an archetype
    # override (see compose.py's own "don't second-guess" rule for
    # cover/end/agenda/etc.) -- confirm they kept their swapped layouts.
    assert compose_result.layouts_used[0] == "Cover B"
    assert compose_result.layouts_used[2] == "Outro"


def test_redesign_deck_brand_review_ignores_an_archetype_override_call_that_errors(tmp_path):
    """Fail-closed here too: a broken archetype-override call must never
    break a brand-mode --review run that otherwise has nothing else to
    do (no ineligible slides, so this really is the only call made)."""
    prs = new_deck()
    cover = add_slide(prs)
    title_run(cover).text = "Annual Review"
    content = add_slide(prs)
    title_run(content).text = "Highlights"
    body_run(content).text = "Grew nicely"
    end = add_slide(prs)
    title_run(end).text = "Thank you"
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    class _RaisingMessagesForThisTest:
        def __init__(self):
            self.calls = []

        def stream(self, **kwargs):
            self.calls.append(kwargs)
            raise RuntimeError("simulated transient failure")

    class _RaisingClientForThisTest:
        def __init__(self):
            self.messages = _RaisingMessagesForThisTest()

    client = _RaisingClientForThisTest()

    out_path = tmp_path / "reviewed.pptx"
    compose_result, _redesign_result = redesign_deck(
        str(src_path), str(out_path), mode="brand", review=True, client=client
    )

    assert compose_result.slide_count == 3
    assert len(client.messages.calls) == 1  # the call was attempted, and its failure was swallowed
    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 3


def test_redesign_deck_brand_review_leaves_a_non_divider_slide_skipped(tmp_path):
    """A slide with real substantial content, still ineligible for
    verbatim carryover -- the model correctly says it's not a divider,
    so it should stay reported as skipped, untouched."""
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Quarterly results"
    for i in range(20):
        box = slide.shapes.add_textbox(Inches(1), Inches(0.3 * i), Inches(3), Inches(0.25))
        box.text_frame.text = f"Line {i}: real detailed content"
    from deckguard.legacy.retemplate import MAX_TEXT_BLOCKS

    assert 20 > MAX_TEXT_BLOCKS  # sanity: this really is ineligible for verbatim carryover
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    response_json = _review_json(
        {"slide_index": 1, "is_divider": False, "divider_title": None, "note": "looks like a real content slide"}
    )
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "reviewed.pptx"
    compose_result, redesign_result = redesign_deck(
        str(src_path), str(out_path), mode="brand", review=True, client=client
    )

    assert compose_result.slide_count == 0
    assert [s.slide_index for s in redesign_result.skipped] == [1]
    assert redesign_result.review_notes == ["slide 1: looks like a real content slide"]


def test_redesign_deck_brand_review_skips_genuinely_blank_slides():
    """A truly empty slide has nothing for the model to judge (and
    nothing to derive a title from without inventing one), so it's
    excluded from the review call entirely rather than wasting a slide
    on a guaranteed-uninteresting judgment."""
    prs = new_deck()
    add_slide(prs)  # genuinely blank -- EMPTY_SLIDE_REASON

    class _ExplodingClient:
        @property
        def messages(self):
            raise AssertionError("a genuinely blank slide should never reach the review call")

    import tempfile
    from pathlib import Path as _Path

    with tempfile.TemporaryDirectory() as d:
        src_path = _Path(d) / "source.pptx"
        prs.save(str(src_path))
        out_path = _Path(d) / "out.pptx"
        redesign_deck(str(src_path), str(out_path), mode="brand", review=True, client=_ExplodingClient())


def test_redesign_deck_review_rejects_rewrite_mode(tmp_path):
    with pytest.raises(RedesignError, match="only applies to mode='brand'"):
        redesign_deck(str(tmp_path / "x.pptx"), str(tmp_path / "out.pptx"), mode="rewrite", review=True)
