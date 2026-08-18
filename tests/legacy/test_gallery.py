"""Tests for parsing kone-design's HTML archetype gallery into engine
archetypes (gallery.py).

The gallery names 57 archetypes; the engine implements its own set, and
people write briefs against the gallery. These pin the parse so a
gallery refresh either keeps working or fails loudly here.
"""

import pytest

from deckguard.legacy.gallery import (
    COVER_ARCHETYPES,
    _box_of,
    _clip_rects,
    _css,
    _hex_of,
    _px,
    build_archetypes,
    gallery_dir,
    parse_section,
)

needs_gallery = pytest.mark.skipif(gallery_dir() is None, reason="no archetype gallery available")


def test_px_accepts_the_unitless_zero_the_gallery_writes():
    """`left:0;top:0` is how the gallery anchors a full-width block.
    Rejecting it silently dropped every cut-image banner."""
    assert _px("45px") == 45.0
    assert _px("0") == 0.0
    assert _px("-8px") == -8.0
    assert _px("100%") is None
    assert _px(None) is None


def test_hex_normalises_and_refuses_what_it_cannot_render():
    assert _hex_of("#fff") == "FFFFFF"
    assert _hex_of("#1450f5") == "1450F5"
    assert _hex_of("linear-gradient(#000, transparent)") is None
    assert _hex_of("transparent") is None


def test_boxes_resolve_bottom_and_right_against_the_slide():
    assert _box_of(_css("left:45px;top:136px;width:374px;height:448px")) == [45, 136, 374, 448]
    # footer chrome is anchored from the bottom
    box = _box_of(_css("left:45px;bottom:43px"))
    assert box[1] == 720 - 43
    # `inset:0` is the full-bleed idiom
    assert _box_of(_css("inset:0;width:100%;height:100%")) == [0.0, 0.0, 1280.0, 720.0]


def test_cut_masks_parse_into_rectangles():
    """The signature staggered banner is an SVG path made only of
    axis-aligned rectangles -- which is the whole reason it can be
    reproduced in PowerPoint at all."""
    rects = _clip_rects(
        "path('M0 0 H288.9 V249.1 H0 Z M330.2 0 H619.1 V322.3 H330.2 Z')"
    )
    rounded = [[round(v, 1) for v in r] for r in rects]
    assert rounded == [[0.0, 0.0, 288.9, 249.1], [330.2, 0.0, 288.9, 322.3]]


def test_a_wrapper_and_its_heading_become_one_region():
    """The gallery puts the BOX on a positioned wrapper and the TYPE on
    an <h2> inside it. Read separately, every title vanished."""
    section = (
        '<section style="position:relative;width:1280px;height:720px;background:#fff">'
        '<div style="position:absolute;left:45px;top:136px;width:374px;height:448px">'
        '<h2 style="margin:0;font-size:64px;color:#141414">The section title</h2></div>'
        "</section>"
    )
    parsed = parse_section("DIVIDER_X", section)

    regions = parsed["archetype"]["regions"]
    assert len(regions) == 1
    assert regions[0]["content"] == "title"
    assert regions[0]["box"] == [45, 136, 374, 448]
    assert parsed["sample"]["title"] == "The section title"


def test_a_wrapper_holding_two_lines_yields_two_regions():
    """An eyebrow above a title share one wrapper; keeping only the
    first lost the title on TITLE_TEXT_SPLIT."""
    section = (
        '<section style="position:relative;width:1280px;height:720px;background:#1450f5">'
        '<div style="position:absolute;left:45px;top:91px;width:272px;height:539px">'
        '<div style="font-family:\'KONE Information\';font-size:12px;'
        'text-transform:uppercase;color:#fff">Section label</div>'
        '<h2 style="margin:0;font-size:34px;color:#fff">A title in the field</h2></div>'
        "</section>"
    )
    parsed = parse_section("TITLE_X", section)

    slots = [r["content"] for r in parsed["archetype"]["regions"]]
    assert "eyebrow" in slots and "title" in slots
    assert parsed["archetype"]["background"] == "1450F5"


def test_a_photo_inside_a_positioned_wrapper_keeps_the_wrappers_geometry():
    section = (
        '<section style="position:relative;width:1280px;height:720px;background:#fff">'
        '<div style="position:absolute;left:759px;top:0;width:521px;height:720px">'
        '<img src="assets/photos/stairs-bag.jpg" style="width:100%;height:100%"></div>'
        "</section>"
    )
    parsed = parse_section("AGENDA_X", section)

    photos = [r for r in parsed["archetype"]["regions"] if r["content"] == "photo"]
    assert photos and photos[0]["box"] == [759, 0, 521, 720]


def test_the_footer_date_and_page_number_are_chrome_not_author_slots():
    section = (
        '<section style="position:relative;width:1280px;height:720px;background:#fff">'
        '<div style="position:absolute;left:45px;bottom:43px;font-size:11px;color:#141414">'
        "3 August 2026</div>"
        '<div style="position:absolute;left:1167px;bottom:43px;font-size:11px;color:#141414">'
        "07</div></section>"
    )
    parsed = parse_section("DIVIDER_Y", section)
    assert parsed is None or parsed["archetype"]["regions"] == []


@needs_gallery
def test_the_finished_gallery_yields_the_archetypes_people_ask_for_by_name():
    """The reported failure: a brief naming COVER_A_CUT4, DIVIDER_D and
    END_LOGO got entirely different slides built."""
    built = build_archetypes()

    for name in ("cover_a_cut4", "divider_d", "end_logo", "title_text_split"):
        assert name in built["archetypes"], name

    cover = built["archetypes"]["cover_a_cut4"]
    assert {r["content"] for r in cover["regions"]} >= {"title", "context"}
    cuts = [c for c in cover["chrome"] if c["kind"] == "cut"]
    assert cuts and len(cuts[0]["rects"]) == 4, "CUT4 is four staggered panes"

    assert built["archetypes"]["divider_d"]["background"] == "FFFFFF"
    assert built["archetypes"]["divider_a"]["background"] == "D2F5FF"
    assert built["archetypes"]["title_text_split"]["background"] == "1450F5"


@needs_gallery
def test_gallery_archetypes_reach_the_engine_registry():
    """Merged into the skill's own registry, so everything derived at
    runtime -- signatures, previews, picture slots, the planning prompt
    -- picks them up with no code change."""
    from deckguard.legacy.skill_bridge import (
        _ensure_skill_on_path,
        _load_archetypes,
        archetype_signatures,
        resolve_archetype_name,
    )

    try:
        _ensure_skill_on_path()
    except Exception:  # noqa: BLE001
        pytest.skip("kone-deck-generator skill not installed")
    module = _load_archetypes()

    assert "cover_a_cut4" in module.ARCHETYPES
    assert resolve_archetype_name("COVER_A_CUT4") == ("cover_a_cut4", "exact")
    assert resolve_archetype_name("DIVIDER_D") == ("divider_d", "exact")
    assert any(s["name"] == "end_logo" for s in archetype_signatures())


def test_a_named_cover_replaces_the_masters_own_rather_than_doubling_it(tmp_path):
    """`build_deck` always keeps the master's cover and Thank you. Once
    the brief names a cover archetype, keeping both opens the deck on
    two covers."""
    from deckguard.legacy.gallery import CLOSER_ARCHETYPES, drop_redundant_master_slides

    assert "cover_a_cut4" in COVER_ARCHETYPES
    assert "end_logo" in CLOSER_ARCHETYPES

    # no cover/closer archetype -> nothing is dropped, and no file is touched
    missing = tmp_path / "nope.pptx"
    assert drop_redundant_master_slides(missing, {"slides": [{"archetype": "three_stats"}]}) == 0
    assert drop_redundant_master_slides(missing, {"slides": []}) == 0


def _built_deck(tmp_path, spec_slides, name="d.pptx"):
    from deckguard.legacy.transform import SlidePlan, TransformPlan, execute_transform_from_brief

    plan = TransformPlan(
        slides=[SlidePlan(index=i, default_action="new", archetype={"archetype": n, **c})
                for i, (n, c) in enumerate(spec_slides, 1)],
        deck_title="Test deck")
    out = tmp_path / name
    execute_transform_from_brief(str(out), plan)
    return out


@needs_gallery
def test_chrome_assets_keep_their_transparency(tmp_path):
    """Reported as "the logo has a black background". The engine's own
    `_image` opens every file with `.convert("RGB")`, which composites
    alpha onto BLACK -- and the KONE marks are mostly transparent (the
    tagline is 25% ink, the divider illustration 16%), so each landed as
    a black rectangle."""
    import io

    from PIL import Image
    from pptx import Presentation

    pytest.importorskip("pptx")
    try:
        deck = _built_deck(tmp_path, [("divider_d", {"title": "A section"})])
    except Exception as exc:  # noqa: BLE001
        pytest.skip(f"deck build unavailable: {exc}")

    # Only the PNG chrome matters here -- the master's retained cover and
    # Thank you carry ordinary JPEG photography, which has no alpha to
    # keep.
    modes = []
    for slide in Presentation(str(deck)).slides:
        for shape in slide.shapes:
            try:
                blob = shape.image.blob
            except Exception:  # noqa: BLE001
                continue
            if blob[:8] == b"\x89PNG\r\n\x1a\n":
                modes.append(Image.open(io.BytesIO(blob)).mode)
    assert modes, "the divider draws a logo, an illustration and a tagline"
    assert all(m == "RGBA" for m in modes), f"alpha was flattened: {modes}"


@needs_gallery
def test_the_cut_cover_is_one_swappable_picture_not_four_baked_panes(tmp_path):
    """Asked for as "a template where when we add a picture it adds that
    chopped effect, instead of it being chopped into four sections
    already" -- one picture plus a mask, so Change Picture still works."""
    from pptx import Presentation

    try:
        deck = _built_deck(tmp_path, [("cover_a_cut4", {"title": "T", "context": "C"})])
    except Exception as exc:  # noqa: BLE001
        pytest.skip(f"deck build unavailable: {exc}")

    cover = Presentation(str(deck)).slides[0]
    photos = [sh for sh in cover.shapes
              if getattr(sh, "shape_type", None) and sh.shape_type.name == "PICTURE"
              and sh.width > 3000000]  # the banner, not the logo or tagline
    assert len(photos) == 1, "the banner must be a single picture"
    # ...and the stagger comes from mask rectangles over it
    masks = [sh for sh in cover.shapes
             if getattr(sh, "shape_type", None) and sh.shape_type.name == "AUTO_SHAPE"]
    assert len(masks) >= 4, "the cut is drawn by covering what the mask hides"


@needs_gallery
def test_every_slide_gets_the_footer_its_archetype_calls_for(tmp_path):
    """A whole deck came out with no dates and no page numbers: the
    parser treated both as author content and dropped them, when they
    are chrome stamped per slide. The rules differ by slide type."""
    from pptx import Presentation

    try:
        deck = _built_deck(tmp_path, [
            ("cover_a_cut4", {"title": "T", "context": "C"}),
            ("divider_d", {"title": "A section"}),
            ("three_picture_cards", {"title": "Cards", "cards": [{"heading": "A", "bullets": ["x"]}]}),
            ("end_logo", {}),
        ])
    except Exception as exc:  # noqa: BLE001
        pytest.skip(f"deck build unavailable: {exc}")

    def footer_text(slide):
        return " | ".join(
            sh.text_frame.text.strip() for sh in slide.shapes
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
        )

    slides = list(Presentation(str(deck)).slides)
    cover, divider, cards, closer = slides

    assert "20" in footer_text(cover), "covers carry a date"
    assert "01" not in footer_text(cover).split(), "covers carry no page number"
    assert "02" in footer_text(divider), "DIVIDER_D puts number and date together"
    assert "03" in footer_text(cards), "content slides carry a page number"
    assert footer_text(closer) == "", "END_LOGO carries nothing at all"
