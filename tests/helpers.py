"""Fixture-deck builders shared across the test suite.

Every test builds its .pptx in-memory with python-pptx rather than
shipping binary fixture files, so each known violation is traceable to
the exact code that created it.
"""

from __future__ import annotations

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu, Pt

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def a(tag: str) -> str:
    return f"{{{A_NS}}}{tag}"


def new_deck() -> Presentation:
    return Presentation()


def add_slide(prs: Presentation, layout_idx: int = 1):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_run(
    run,
    text: str | None = None,
    font: str | None = None,
    bold: bool | None = None,
    size_pt: float | None = None,
    color_hex: str | None = None,
):
    if text is not None:
        run.text = text
    if font is not None:
        run.font.name = font
    if bold is not None:
        run.font.bold = bold
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if color_hex is not None:
        run.font.color.rgb = RGBColor.from_string(color_hex)
    return run


def title_run(slide):
    tf = slide.shapes.title.text_frame
    if not tf.paragraphs[0].runs:
        tf.text = ""
    return tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()


def body_run(slide):
    body = slide.placeholders[1]
    tf = body.text_frame
    if not tf.paragraphs[0].runs:
        tf.text = ""
    return tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()


def _theme_part(prs):
    master = prs.slide_masters[0]
    return master.part.part_related_by(RT.THEME)


def set_theme_slot(prs, slot: str, hex_val: str):
    """Directly rewrite a <a:clrScheme> slot (e.g. 'accent1') to hex_val."""
    theme_part = _theme_part(prs)
    root = etree.fromstring(theme_part.blob)
    clr_scheme = root.find(f".//{a('clrScheme')}")
    slot_el = clr_scheme.find(a(slot))
    srgb = slot_el.find(a("srgbClr"))
    sys_clr = slot_el.find(a("sysClr"))
    if srgb is not None:
        srgb.set("val", hex_val)
    elif sys_clr is not None:
        slot_el.remove(sys_clr)
        new_el = etree.SubElement(slot_el, a("srgbClr"))
        new_el.set("val", hex_val)
    theme_part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def set_theme_font(prs, slot: str, typeface: str):
    """slot: 'majorFont' or 'minorFont'."""
    theme_part = _theme_part(prs)
    root = etree.fromstring(theme_part.blob)
    font_scheme = root.find(f".//{a('fontScheme')}")
    latin = font_scheme.find(a(slot)).find(a("latin"))
    latin.set("typeface", typeface)
    theme_part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def set_run_theme_color(run, theme_color: MSO_THEME_COLOR, brightness: float = 0.0):
    run.font.color.theme_color = theme_color
    if brightness:
        run.font.color.brightness = brightness


def add_shadow_effect(run):
    """Attach a minimal <a:outerShdw> to a run's rPr, as PowerPoint's
    Format Text Effects > Shadow would."""
    rPr = run.font._rPr  # get_or_add is fine here — we're building fixtures
    effect_lst = etree.SubElement(rPr, a("effectLst"))
    shdw = etree.SubElement(effect_lst, a("outerShdw"))
    shdw.set("blurRad", "40000")
    shdw.set("dist", "20000")
    shdw.set("dir", "5400000")
    clr = etree.SubElement(shdw, a("srgbClr"))
    clr.set("val", "000000")
    return rPr


def add_text_outline(run, hex_val: str = "FF0000"):
    rPr = run.font._rPr
    ln = etree.SubElement(rPr, a("ln"))
    fill = etree.SubElement(ln, a("solidFill"))
    clr = etree.SubElement(fill, a("srgbClr"))
    clr.set("val", hex_val)
    return rPr


def add_shape_3d(shape):
    spPr = shape._element.spPr
    etree.SubElement(spPr, a("sp3d"))


def add_picture(slide, image_path, left_in=1, top_in=1, width_in=2, height_in=1):
    return slide.shapes.add_picture(
        image_path, Emu(int(914400 * left_in)), Emu(int(914400 * top_in)),
        width=Emu(int(914400 * width_in)), height=Emu(int(914400 * height_in)),
    )


def make_solid_png(path, rgb=(0, 94, 184), size=(64, 64)):
    from PIL import Image

    Image.new("RGB", size, color=rgb).save(path)
    return path


def make_pattern_png(path, seed=0, size=(64, 64)):
    """A deterministic, non-uniform image — perceptual hashing degenerates
    to a constant value for solid-color images (no DCT detail to hash), so
    fixtures that need distinguishable phashes use this instead."""
    import random

    from PIL import Image, ImageDraw

    rng = random.Random(seed)
    img = Image.new("RGB", size, color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    for _ in range(12):
        x0, y0 = rng.randint(0, size[0] - 1), rng.randint(0, size[1] - 1)
        x1, y1 = rng.randint(0, size[0] - 1), rng.randint(0, size[1] - 1)
        color = (rng.randint(0, 255), rng.randint(0, 255), rng.randint(0, 255))
        xa, xb = sorted((x0, x1))
        ya, yb = sorted((y0, y1))
        draw.rectangle([xa, ya, xb, yb], fill=color)
    img.save(path)
    return path
