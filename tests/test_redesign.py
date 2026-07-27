"""Tests for deckguard's AI-assisted redesign (redesign.py): extracting
eligible content from an arbitrary deck, asking Claude to map it onto
compose.py's outline schema, then building it through the same
deterministic pipeline `deckguard create` uses.

No real network/API calls -- every test injects a fake Anthropic-shaped
client, consistent with the rest of this suite's "works with no API key"
philosophy for anything that doesn't require a live account.
"""

import json

import pytest
from pptx import Presentation
from pptx.util import Inches

from deckguard.redesign import (
    RedesignError,
    Usage,
    call_claude_for_outline,
    extract_eligible_slides,
    partition_skipped,
    redesign_deck,
)
from deckguard.slide_import import default_template_path
from tests.helpers import add_picture, add_slide, body_run, make_pattern_png, new_deck, title_run

TEMPLATE_PATH = default_template_path()
pytestmark = pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="bundled template asset not present")


# --------------------------------------------------------------------------
# fake Anthropic client
# --------------------------------------------------------------------------


class _FakeTextBlock:
    type = "text"

    def __init__(self, text):
        self.text = text


class _FakeUsage:
    def __init__(self, input_tokens, output_tokens):
        self.input_tokens = input_tokens
        self.output_tokens = output_tokens


class _FakeResponse:
    def __init__(self, content_text, stop_reason="end_turn", input_tokens=1000, output_tokens=500):
        self.content = [_FakeTextBlock(content_text)] if content_text is not None else []
        self.stop_reason = stop_reason
        self.usage = _FakeUsage(input_tokens, output_tokens)


class _FakeStreamCM:
    def __init__(self, response):
        self._response = response

    def __enter__(self):
        return self

    def __exit__(self, *exc):
        return False

    def get_final_message(self):
        return self._response


class _FakeMessages:
    def __init__(self, response):
        self._response = response
        self.calls = []

    def stream(self, **kwargs):
        self.calls.append(kwargs)
        return _FakeStreamCM(self._response)


class _FakeClient:
    def __init__(self, response):
        self.messages = _FakeMessages(response)


def _outline_json(*items):
    return json.dumps({"slides": list(items)})


def _slide_item(index, kind="content", **overrides):
    item = {
        "source_slide_index": index,
        "kind": kind,
        "title": "A title",
        "subtitle": None,
        "bullets": ["Point one", "Point two"],
        "columns": [],
        "quote_text": None,
        "quote_author": None,
        "quote_label": None,
        "stats": [],
        "milestones": [],
        "variant": None,
    }
    item.update(overrides)
    return item


# --------------------------------------------------------------------------
# extract_eligible_slides
# --------------------------------------------------------------------------


def test_extract_eligible_slides_separates_ordinary_from_disqualified():
    prs = new_deck()
    ok_slide = add_slide(prs)
    title_run(ok_slide).text = "Ordinary slide"
    body_run(ok_slide).text = "Some content"

    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))

    eligible, skipped = extract_eligible_slides(prs)

    assert [i for i, _p in eligible] == [1]
    assert [s.slide_index for s in skipped] == [2]
    assert "table" in skipped[0].reason


def test_extract_eligible_slides_empty_deck_slide_is_skipped():
    prs = new_deck()
    add_slide(prs)  # title/body placeholders present but empty
    _eligible, skipped = extract_eligible_slides(prs)
    assert len(skipped) == 1


def test_extract_eligible_slides_accepts_dense_text_retemplate_would_skip():
    """Regression test: a hand-built slide with more text boxes than any
    layout has placeholders used to be skipped by redesign too (it
    reused retemplate.classify_slide's MAX_TEXT_BLOCKS=3 cap verbatim),
    even though redesign is allowed to condense wording and retemplate
    is not. redesign should accept this and let the model condense it,
    not skip it."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    for i in range(5):  # one more than retemplate.MAX_TEXT_BLOCKS (3)
        box = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(3), Inches(0.5))
        box.text_frame.text = f"Block {i}"

    eligible, skipped = extract_eligible_slides(prs)

    assert len(eligible) == 1
    assert skipped == []
    assert len(eligible[0][1].text_blocks) == 5


def test_extract_eligible_slides_accepts_many_small_boxes_regardless_of_count():
    """Block COUNT alone, however high, is never a reason to skip --
    only total text volume is (see test below). Twenty tiny caption
    boxes have far less content than a handful of paragraph-sized ones,
    and redesign condenses either way."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    for i in range(20):
        box = slide.shapes.add_textbox(Inches(1), Inches(1.0 + 0.3 * i), Inches(3), Inches(0.25))
        box.text_frame.text = f"Block {i}"

    eligible, skipped = extract_eligible_slides(prs)

    assert skipped == []
    assert len(eligible[0][1].text_blocks) == 20


def test_extract_eligible_slides_still_skips_truly_excessive_text_volume():
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    # A handful of boxes, but each stuffed with enough text to blow past
    # REDESIGN_MAX_TEXT_CHARS -- volume, not count, is what's capped.
    huge_text = "Lorem ipsum dolor sit amet. " * 800  # ~23,000 characters
    for i in range(3):
        box = slide.shapes.add_textbox(Inches(1), Inches(1.0 + 2.0 * i), Inches(3), Inches(1.5))
        box.text_frame.text = huge_text

    _eligible, skipped = extract_eligible_slides(prs)

    assert len(skipped) == 1
    assert "reasonably condensed" in skipped[0].reason


def test_extract_eligible_slides_still_hard_skips_tables_regardless_of_text_cap():
    """A brief -- or redesign's higher text cap -- never overrides the
    shape-type safety rules shared with retemplate."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))

    _eligible, skipped = extract_eligible_slides(prs)

    assert len(skipped) == 1
    assert "table" in skipped[0].reason


# --------------------------------------------------------------------------
# call_claude_for_outline
# --------------------------------------------------------------------------


def test_call_claude_for_outline_parses_valid_response():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Quarterly review"
    body_run(slide).text = "Revenue up"
    eligible, _skipped = extract_eligible_slides(prs)

    response_json = _outline_json(_slide_item(1, title="Quarterly review"))
    client = _FakeClient(_FakeResponse(response_json, input_tokens=1200, output_tokens=600))

    raw_slides, usage = call_claude_for_outline(eligible, client=client)

    assert raw_slides[0]["title"] == "Quarterly review"
    assert usage.input_tokens == 1200
    assert usage.output_tokens == 600
    assert usage.model == "claude-opus-5"
    # (1200/1e6)*5 + (600/1e6)*25 = 0.006 + 0.015 = 0.021
    assert usage.estimated_cost_usd == pytest.approx(0.021)


def test_call_claude_for_outline_passes_model_and_effort_through():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _FakeClient(_FakeResponse(_outline_json(_slide_item(1))))
    call_claude_for_outline(eligible, model="claude-sonnet-5", effort="medium", client=client)

    call_kwargs = client.messages.calls[0]
    assert call_kwargs["model"] == "claude-sonnet-5"
    assert call_kwargs["output_config"]["effort"] == "medium"
    assert call_kwargs["output_config"]["format"]["type"] == "json_schema"


def test_call_claude_for_outline_raises_on_refusal():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _FakeClient(_FakeResponse(None, stop_reason="refusal"))
    with pytest.raises(RedesignError, match="declined"):
        call_claude_for_outline(eligible, client=client)


def test_call_claude_for_outline_raises_on_invalid_json():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _FakeClient(_FakeResponse("not json"))
    with pytest.raises(RedesignError, match="not valid JSON"):
        call_claude_for_outline(eligible, client=client)


def test_call_claude_for_outline_raises_on_missing_slides_key():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "T"
    body_run(slide).text = "B"
    eligible, _ = extract_eligible_slides(prs)

    client = _FakeClient(_FakeResponse(json.dumps({"oops": []})))
    with pytest.raises(RedesignError, match="'slides'"):
        call_claude_for_outline(eligible, client=client)


# --------------------------------------------------------------------------
# redesign_deck (end to end)
# --------------------------------------------------------------------------


def test_redesign_deck_builds_a_valid_composed_deck(tmp_path):
    prs = new_deck()
    cover = add_slide(prs, layout_idx=0)
    title_run(cover).text = "Annual Review"

    content = add_slide(prs)
    title_run(content).text = "Highlights"
    body_run(content).text = "Grew nicely"

    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))

    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    response_json = _outline_json(
        _slide_item(1, kind="cover", title="Annual Review", subtitle="A great year", bullets=[]),
        _slide_item(2, kind="content", title="Highlights", bullets=["Grew nicely"]),
    )
    client = _FakeClient(_FakeResponse(response_json, input_tokens=2000, output_tokens=800))

    out_path = tmp_path / "redesigned.pptx"
    compose_result, redesign_result = redesign_deck(str(src_path), str(out_path), client=client)

    assert compose_result.slide_count == 2
    assert compose_result.layouts_used[0] == "Cover B"
    assert len(redesign_result.skipped) == 1
    assert redesign_result.skipped[0].slide_index == 3
    assert redesign_result.usage.input_tokens == 2000

    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 2


def test_redesign_deck_carries_source_images_into_the_redesigned_slide(tmp_path):
    """Regression test: the model's outline schema has no field for images
    at all (it's never shown the pixels), so images used to be silently
    dropped for every redesigned slide -- a real bug reported against a
    genuinely image-heavy source deck. Images should now survive, backfilled
    deterministically by source_slide_index after the model call."""
    prs = new_deck()
    content = add_slide(prs)
    title_run(content).text = "Highlights"
    body_run(content).text = "Grew nicely"
    img_path = make_pattern_png(tmp_path / "img.png", seed=3)
    add_picture(content, str(img_path))

    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    response_json = _outline_json(
        _slide_item(1, kind="content", title="Highlights", bullets=["Grew nicely"]),
    )
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "redesigned.pptx"
    compose_result, _redesign_result = redesign_deck(str(src_path), str(out_path), client=client)

    assert compose_result.slide_count == 1
    out_prs = Presentation(str(out_path))
    blobs = []
    for shp in out_prs.slides[0].shapes:
        try:
            blobs.append(shp.image.blob)
        except (AttributeError, ValueError):
            continue
    assert img_path.read_bytes() in blobs


def test_redesign_deck_raises_when_nothing_is_eligible_and_no_brief(tmp_path):
    prs = new_deck()
    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))
    src_path = tmp_path / "source.pptx"
    prs.save(str(src_path))

    client = _FakeClient(_FakeResponse(_outline_json()))
    with pytest.raises(RedesignError, match="nothing to work with"):
        redesign_deck(str(src_path), str(tmp_path / "out.pptx"), client=client)


def test_redesign_deck_raises_with_no_deck_and_no_brief(tmp_path):
    client = _FakeClient(_FakeResponse(_outline_json()))
    with pytest.raises(RedesignError, match="nothing to work with"):
        redesign_deck(None, str(tmp_path / "out.pptx"), client=client)


# --------------------------------------------------------------------------
# partition_skipped -- blank slides vs genuinely-unsafe content
# --------------------------------------------------------------------------


def test_partition_skipped_splits_blank_from_real_skips():
    prs = new_deck()
    add_slide(prs)  # empty -> blank

    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))  # real skip

    _eligible, skipped = extract_eligible_slides(prs)
    blank_indices, real_skipped = partition_skipped(skipped)

    assert blank_indices == [1]
    assert [s.slide_index for s in real_skipped] == [2]
    assert "table" in real_skipped[0].reason


# --------------------------------------------------------------------------
# redesign_deck -- brief-driven modes (fill blanks / build from scratch)
# --------------------------------------------------------------------------


def test_redesign_deck_builds_from_scratch_with_no_source_deck(tmp_path):
    """The pure "Claude, design me a deck" mode -- no deck_path at all."""
    response_json = _outline_json(
        _slide_item(None, kind="cover", title="Predictive Maintenance", subtitle="A new era", bullets=[]),
        _slide_item(None, kind="content", title="Why it matters", bullets=["Less downtime", "Lower cost"]),
        _slide_item(None, kind="end", title="Thank you", bullets=[]),
    )
    client = _FakeClient(_FakeResponse(response_json, input_tokens=1800, output_tokens=900))

    out_path = tmp_path / "from_scratch.pptx"
    compose_result, redesign_result = redesign_deck(
        None, str(out_path), brief="A short deck on predictive maintenance for facilities managers.", client=client
    )

    assert compose_result.slide_count == 3
    assert redesign_result.skipped == []

    # the model saw "no source deck" framing, not per-slide extracted content
    sent_content = client.messages.calls[0]["messages"][0]["content"]
    assert "no source deck" in sent_content.lower()
    assert "predictive maintenance" in sent_content.lower()


def test_redesign_deck_fills_blank_slides_alongside_real_content(tmp_path):
    prs = new_deck()
    content_slide = add_slide(prs)
    title_run(content_slide).text = "Existing highlight"
    body_run(content_slide).text = "Real content here"

    add_slide(prs)  # blank -- should be authored from the brief, not skipped

    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))  # real skip, brief or not

    src_path = tmp_path / "half_empty.pptx"
    prs.save(str(src_path))

    response_json = _outline_json(
        _slide_item(1, kind="content", title="Existing highlight", bullets=["Real content here"]),
        _slide_item(None, kind="stat", title="By the numbers", stats=[{"number": "3", "label": "sites live"}], bullets=[]),
    )
    client = _FakeClient(_FakeResponse(response_json))

    out_path = tmp_path / "filled.pptx"
    compose_result, redesign_result = redesign_deck(
        str(src_path), str(out_path), brief="Add a stats slide about rollout progress.", client=client
    )

    assert compose_result.slide_count == 2
    # the table slide is still reported -- a brief never overrides "unsafe to touch"
    assert [s.slide_index for s in redesign_result.skipped] == [3]

    call_kwargs = client.messages.calls[0]
    sent_content = call_kwargs["messages"][0]["content"]
    assert "1 blank slide" in sent_content
    assert "rollout progress" in sent_content.lower()


def test_call_claude_for_outline_target_slides_guidance_is_included():
    client = _FakeClient(_FakeResponse(_outline_json(_slide_item(None))))
    call_claude_for_outline([], blank_count=0, brief="A topic", target_slides=7, client=client)
    sent_content = client.messages.calls[0]["messages"][0]["content"]
    assert "approximately 7 total slides" in sent_content


def test_usage_estimated_cost_unknown_model_is_zero():
    usage = Usage(input_tokens=1000, output_tokens=1000, model="some-future-model")
    assert usage.estimated_cost_usd == 0.0
