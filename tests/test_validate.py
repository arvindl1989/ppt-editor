"""Gate 1: the assertions, and the exceptions that keep them credible.

The risk with a check suite is not that it misses things -- it is that
it reports things nobody agrees are faults, and then everyone stops
reading the panel. Most of these tests are about the false positive
rather than the catch.
"""

import os

import pytest

from deckguard import validate as V


@pytest.fixture(autouse=True)
def _clean_severity():
    before = os.environ.get("DECKGUARD_SEVERITY")
    os.environ.pop("DECKGUARD_SEVERITY", None)
    yield
    if before is None:
        os.environ.pop("DECKGUARD_SEVERITY", None)
    else:
        os.environ["DECKGUARD_SEVERITY"] = before


# --------------------------------------------------------------------------
# the mechanism
# --------------------------------------------------------------------------


def test_every_check_ships_at_report():
    """Open question 5 -- what preflight should refuse -- is not settled,
    and the architecture must not settle it. A check that shipped at
    `refuse` would decide the policy by being written."""
    assert set(V.DEFAULT_SEVERITY.values()) == {V.REPORT}


def test_severity_is_configuration_not_code():
    os.environ["DECKGUARD_SEVERITY"] = "copy_fits=refuse,cardinality=block"
    assert V.severity("copy_fits") == V.REFUSE
    assert V.severity("cardinality") == V.BLOCK
    assert V.severity("no_source_echo") == V.REPORT      # untouched
    assert V.severity("not_a_check") == V.REPORT         # unknown names are safe


def test_a_nonsense_severity_is_ignored_rather_than_obeyed():
    os.environ["DECKGUARD_SEVERITY"] = "copy_fits=maybe"
    assert V.severity("copy_fits") == V.REPORT


def test_worst_reads_the_strongest_severity_present():
    os.environ["DECKGUARD_SEVERITY"] = "cardinality=refuse"
    findings = [V.Finding("copy_fits", "a"), V.Finding("cardinality", "b")]
    assert V.worst(findings) == V.REFUSE
    assert V.worst([]) == V.REPORT


# --------------------------------------------------------------------------
# gate 1a
# --------------------------------------------------------------------------


def _plan(slide: dict, **extra) -> dict:
    return {"title": "T", "audience": "internal", "slides": [slide], **extra}


def test_copy_that_does_not_fit_is_reported_rather_than_trimmed():
    """The fault this exists for: a cover title arrived as the brief's
    opening sentence, hard-cut mid-word, set at 76px. The cut is gone;
    this is the check that stops the copy being written that way."""
    room = V.budget("title_content", "title")
    assert room > 0, "no budget for a slot that plainly has one"

    long = _plan({"archetype": "title_content", "title": "x " * room})
    checks = [f for f in V.before_build(long) if f.check == "copy_fits"]
    assert checks and "title" in checks[0].note

    short = _plan({"archetype": "title_content", "title": "Channels in scope"})
    assert not [f for f in V.before_build(short) if f.check == "copy_fits"]


def test_the_budget_follows_a_migrated_archetype():
    """The size comes from the region's own type block where it has one
    and from the brand where it does not, so an archetype moved off
    baked type does not silently lose its budget."""
    assert V.budget("divider_numbering", "title") > 0     # fully migrated
    assert V.budget("title_content", "bullets") > 0       # still baked
    assert V.budget("title_content", "nonexistent") == 0


def test_copy_lifted_from_the_brief_is_reported():
    brief = ("I would like to share our plan of ONE Week MOD deployment with "
             "you regarding the September launch across six European frontlines.")
    lifted = _plan({"archetype": "title_content"},
                   title="I would like to share our plan of ONE Week MOD deployment")
    found = [f for f in V.before_build(lifted, brief) if f.check == "no_source_echo"]
    assert found, "the cover title is the brief's own first line"

    written = _plan({"archetype": "title_content"}, title="ONE Week MOD deployment")
    assert not [f for f in V.before_build(written, brief)
                if f.check == "no_source_echo"]


def test_a_short_brief_cannot_trigger_the_echo_check():
    """Below the threshold any slide title would match by coincidence."""
    assert not V.before_build(_plan({"archetype": "title_content"}, title="Scope"),
                              brief="Scope")


def test_a_bulleted_list_is_not_a_cardinality_fault():
    """`bullets` is one slot holding many lines and the contract
    notation cannot say so -- it records cardinality 1 in all fourteen
    archetypes that have one. Reading that literally reported every
    correctly-filled bullet slide as a fault."""
    plan = _plan({"archetype": "title_content", "title": "Scope",
                  "bullets": ["one", "two", "three"]})
    assert not [f for f in V.before_build(plan) if f.check == "cardinality"]


def test_a_list_longer_than_its_region_is_a_cardinality_fault():
    from deckguard import contracts as C

    contract = C.for_archetype("numbered_icon_row_6", "internal")
    slot = next(s for s in contract.slots if s.is_list)
    plan = _plan({"archetype": "numbered_icon_row_6", "title": "Six",
                  slot.key: [{"text": "x"}] * (slot.maximum + 3)})
    found = [f for f in V.before_build(plan) if f.check == "cardinality"]
    assert found and str(slot.maximum) in found[0].note


def test_a_slide_picked_from_the_gallery_is_not_incomplete():
    """A bare `{archetype: name}` is a pick -- "give me this layout" --
    and the renderer fills it with the archetype's own sample.
    Reporting those turned a hand-picked deck into 92 findings."""
    picked = {"title": "P", "audience": "internal",
              "slides": [{"archetype": "numbered_icon_row_6"},
                         {"archetype": "kone_numbers"}]}
    assert not [f for f in V.before_build(picked)
                if f.check == "slot_completeness"]


def test_a_half_filled_slide_still_is():
    half = _plan({"archetype": "numbered_icon_row_6", "title": "Six workstreams"})
    found = [f for f in V.before_build(half) if f.check == "slot_completeness"]
    assert found, "a slide with a title and no items is genuinely incomplete"


def test_variety_counts_content_layouts_and_forgives_dividers():
    slides = [{"archetype": "divider_numbering", "number": f"0{n}", "title": "s"}
              for n in range(1, 5)]
    slides += [{"archetype": "title_content", "title": "t", "bullets": ["a"]}
               for _ in range(6)]
    found = [f.note for f in V.before_build({"title": "T", "audience": "internal",
                                             "slides": slides})
             if f.check == "archetype_variety"]
    assert any("6 times" in n for n in found), found
    assert not any("divider" in n for n in found)


def test_two_consecutive_content_slides_of_one_layout_are_reported():
    slides = [{"archetype": "title_content", "title": "a", "bullets": ["x"]},
              {"archetype": "title_content", "title": "b", "bullets": ["y"]}]
    found = [f for f in V.before_build({"title": "T", "audience": "internal",
                                        "slides": slides})
             if f.check == "archetype_variety"]
    assert any("consecutive" in f.note for f in found)


# --------------------------------------------------------------------------
# gate 1b
# --------------------------------------------------------------------------


def test_real_output_passes_everything_except_the_notes(tmp_path):
    """The gate has to be quiet on a deck that is actually fine, or the
    panel stops being read. `deck-13.pptx` is real output: the only
    thing it is missing is speaker notes, which it genuinely is."""
    from pathlib import Path

    fixture = Path(__file__).parent / "fixtures" / "deck-13.pptx"
    checks = {f.check for f in V.after_build(str(fixture))}
    assert checks == {"notes_present"}, checks


def test_a_stale_page_number_is_caught(tmp_path):
    """A page number reading 11 on slide 12 shipped on a real deck, and
    nothing in the pipeline looked at it."""
    from pptx import Presentation
    from pptx.util import Emu, Pt

    out = tmp_path / "chrome.pptx"
    prs = Presentation()
    for _ in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Emu(400000), Emu(6300000), Emu(600000), Emu(200000))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "11"
        run.font.size = Pt(9)
    prs.save(str(out))

    found = [f for f in V.after_build(str(out)) if f.check == "chrome"]
    assert any("reads 11 on slide 1" in f.note for f in found), found


def test_two_different_dates_in_one_deck_are_caught(tmp_path):
    from pptx import Presentation
    from pptx.util import Emu, Pt

    out = tmp_path / "dates.pptx"
    prs = Presentation()
    for text in ("21 AUGUST 2026", "23 JULY 2026"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Emu(400000), Emu(6300000), Emu(3000000), Emu(200000))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(9)
    prs.save(str(out))

    found = [f for f in V.after_build(str(out)) if f.check == "chrome"]
    assert any("more than one date" in f.note for f in found), found


def test_casing_is_checked_in_both_directions(tmp_path):
    """A KONE Information label in sentence case is how the divider
    fault was first reported -- "the fonts look whack" -- and Inter in
    caps is the same mistake mirrored."""
    from pptx import Presentation
    from pptx.util import Emu, Pt

    out = tmp_path / "case.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def run_at(top, text, face):
        box = slide.shapes.add_textbox(Emu(400000), Emu(top), Emu(4000000), Emu(300000))
        piece = box.text_frame.paragraphs[0].add_run()
        piece.text, piece.font.name, piece.font.size = text, face, Pt(12)

    run_at(400000, "Boundaries", "KONE Information")     # should be caps
    run_at(1000000, "WHAT SITS OUTSIDE SCOPE", "Inter")  # should not be
    run_at(1600000, "SCOPE", "KONE Information")         # correct
    run_at(2200000, "What sits outside scope", "Inter")  # correct
    run_at(2800000, "Anything", "Comic Sans MS")         # not an approved face
    prs.save(str(out))

    notes = [f.note for f in V.after_build(str(out)) if f.check == "font_role"]
    assert any("not in caps" in n and "Boundaries" in n for n in notes), notes
    assert any("Inter set in caps" in n for n in notes), notes
    assert any("Comic Sans MS" in n for n in notes), notes
    # Three faults among five runs: the two correct ones are silent.
    assert len(notes) == 3, notes
