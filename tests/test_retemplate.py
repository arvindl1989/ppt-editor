"""Tests for deckguard's content-preserving re-layout (retemplate):
classifying an old deck's slides, matching each to an org-template
layout, and rebuilding accepted slides on it while carrying over
title/body text and images.
"""

import io
import zipfile

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PROG_ID
from pptx.util import Inches

from deckguard.retemplate import (
    CONTENT_LAYOUT_CANDIDATES,
    LayoutProfile,
    SlideProfile,
    _freeze_placeholder_geometry,
    _reparent_slide_layout,
    apply_rebrand,
    apply_retemplate,
    classify_slide,
    match_layout,
    propose_retemplate,
    rebuild_slides_as_dividers,
)
from deckguard.slide_import import default_template_path
from tests.helpers import add_picture, add_rectangle, add_slide, body_run, make_pattern_png, new_deck, title_run

TEMPLATE_PATH = default_template_path()
pytestmark = pytest.mark.skipif(not TEMPLATE_PATH.exists(), reason="bundled template asset not present")


def _no_duplicate_zip_entries(path) -> bool:
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
    return len(names) == len(set(names))


def _no_malformed_xml(path) -> list:
    bad = []
    with zipfile.ZipFile(path) as z:
        for name in z.namelist():
            if name.endswith((".xml", ".rels")):
                try:
                    etree.fromstring(z.read(name))
                except Exception as exc:  # noqa: BLE001
                    bad.append((name, str(exc)))
    return bad


# --------------------------------------------------------------------------
# classify_slide
# --------------------------------------------------------------------------


def test_classify_slide_extracts_title_and_body_text():
    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "My Title"
    body_run(slide).text = "First line"

    profile = classify_slide(slide)

    assert profile.eligible
    assert profile.title == "My Title"
    assert profile.text_blocks == [[(0, "First line")]]
    assert profile.images == []


def test_classify_slide_extracts_images(tmp_path):
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    img_path = make_pattern_png(tmp_path / "img.png", seed=1)
    add_picture(slide, str(img_path))

    profile = classify_slide(slide)

    assert profile.eligible
    assert len(profile.images) == 1
    assert profile.images[0] == img_path.read_bytes()


def test_classify_slide_ignores_small_footer_text_near_slide_edge():
    """A "Confidential | (c) KONE Corporation"-style footer must not be
    treated as real body content -- it isn't something a human wants
    migrated, and would otherwise crowd out real content against the
    max-text-blocks cap and pollute the preview text."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(4), Inches(0.3))
    footer.text_frame.text = "Confidential | © KONE Corporation"

    profile = classify_slide(slide, slide_height_in=7.5)

    assert profile.text_blocks == []
    assert not profile.eligible
    assert profile.reason == "no title, text, or images to migrate"


def test_classify_slide_keeps_a_similarly_small_textbox_away_from_the_edge():
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(0.4))
    box.text_frame.text = "Real short content"

    profile = classify_slide(slide, slide_height_in=7.5)

    assert profile.eligible
    assert profile.text_blocks == [[(0, "Real short content")]]


@pytest.mark.parametrize(
    "build,expected_reason",
    [
        (lambda slide: slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2)), "contains a table"),
        (lambda slide: slide.shapes.add_group_shape(), "contains a group"),
        (
            lambda slide: slide.shapes.add_ole_object(
                io.BytesIO(b"fake xlsx bytes"), PROG_ID.XLSX, Inches(1), Inches(1)
            ),
            "contains an embedded OLE object",
        ),
    ],
)
def test_classify_slide_disqualifies_unsafe_shape_types(build, expected_reason):
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    build(slide)

    profile = classify_slide(slide)

    assert not profile.eligible
    assert profile.reason == expected_reason


def test_classify_slide_disqualifies_too_many_decorative_shapes():
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    for i in range(6):
        add_rectangle(slide, left_in=i, top_in=0, width_in=0.5, height_in=0.5)

    profile = classify_slide(slide)

    assert not profile.eligible
    assert profile.reason == "too many free-form shapes to safely reflow"


def test_classify_slide_disqualifies_more_text_blocks_than_any_layout_holds():
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    for i in range(5):
        box = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(3), Inches(0.5))
        box.text_frame.text = f"Block {i}"

    profile = classify_slide(slide)

    assert not profile.eligible
    assert profile.reason == "more body text blocks than any template layout can hold"


def test_classify_slide_empty_slide_is_ineligible():
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)

    profile = classify_slide(slide)

    assert not profile.eligible
    assert profile.reason == "no title, text, or images to migrate"


# --------------------------------------------------------------------------
# match_layout
# --------------------------------------------------------------------------


def test_match_layout_picks_the_tightest_fitting_candidate():
    profiles = [
        LayoutProfile(name="Loose", has_title=True, n_body=3, n_picture=2),
        LayoutProfile(name="Tight", has_title=True, n_body=1, n_picture=0),
        LayoutProfile(name="TooSmall", has_title=True, n_body=0, n_picture=0),
    ]
    profile = SlideProfile(title="T", text_blocks=[[(0, "x")]], images=[], eligible=True)

    match = match_layout(profile, profiles)

    assert match is not None
    assert match.layout_name == "Tight"


def test_match_layout_requires_a_title_placeholder_when_slide_has_a_title():
    profiles = [LayoutProfile(name="NoTitle", has_title=False, n_body=5, n_picture=5)]
    profile = SlideProfile(title="T", text_blocks=[], images=[], eligible=True)

    assert match_layout(profile, profiles) is None


def test_match_layout_returns_none_when_nothing_fits():
    profiles = [LayoutProfile(name="Small", has_title=True, n_body=1, n_picture=0)]
    profile = SlideProfile(title="T", text_blocks=[[(0, "a")], [(0, "b")], [(0, "c")]], images=[], eligible=True)

    assert match_layout(profile, profiles) is None


def test_match_layout_usage_counts_only_break_ties_never_override_a_better_fit():
    """usage_counts is an anti-repeat tie-break for apply_rebrand's layout
    variety, not a second scoring signal that can outrank actual fit: a
    heavily-reused Tight layout must still beat a fresh, but wasteful,
    Loose one."""
    profiles = [
        LayoutProfile(name="Loose", has_title=True, n_body=3, n_picture=2),
        LayoutProfile(name="Tight", has_title=True, n_body=1, n_picture=0),
    ]
    profile = SlideProfile(title="T", text_blocks=[[(0, "x")]], images=[], eligible=True)

    match = match_layout(profile, profiles, usage_counts={"Tight": 10, "Loose": 0})

    assert match.layout_name == "Tight"


def test_match_layout_usage_counts_break_ties_among_equally_good_fits():
    profiles = [
        LayoutProfile(name="A", has_title=True, n_body=1, n_picture=0),
        LayoutProfile(name="B", has_title=True, n_body=1, n_picture=0),
    ]
    profile = SlideProfile(title="T", text_blocks=[[(0, "x")]], images=[], eligible=True)

    match_no_usage = match_layout(profile, profiles)
    assert match_no_usage.layout_name == "A"  # ties resolve by candidate-list order, unchanged from before

    match_with_usage = match_layout(profile, profiles, usage_counts={"A": 3, "B": 0})
    assert match_with_usage.layout_name == "B"  # least-used wins the tie


def test_candidate_layouts_all_exist_in_the_bundled_template():
    tmpl_prs = Presentation(str(TEMPLATE_PATH))
    names = {layout.name for master in tmpl_prs.slide_masters for layout in master.slide_layouts}
    missing = [n for n in CONTENT_LAYOUT_CANDIDATES if n not in names]
    assert missing == []


# --------------------------------------------------------------------------
# propose_retemplate / apply_retemplate (end to end)
# --------------------------------------------------------------------------


def _simple_deck(tmp_path):
    prs = new_deck()
    slide1 = add_slide(prs)
    title_run(slide1).text = "First Slide"
    body_run(slide1).text = "Some content"

    slide2 = add_slide(prs)
    title_run(slide2).text = "Second Slide"
    body_run(slide2).text = "More content"

    slide3 = add_slide(prs, layout_idx=6)  # ineligible: has a table
    slide3.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))

    path = tmp_path / "src.pptx"
    prs.save(str(path))
    return path


def test_propose_retemplate_classifies_every_slide(tmp_path):
    path = _simple_deck(tmp_path)
    prs = Presentation(str(path))

    proposals = propose_retemplate(prs, template_path=TEMPLATE_PATH)

    assert [p.slide_index for p in proposals] == [1, 2, 3]
    assert proposals[0].eligible and proposals[0].layout_name
    assert proposals[1].eligible and proposals[1].layout_name
    assert not proposals[2].eligible
    assert proposals[2].reason == "contains a table"


def test_apply_retemplate_defaults_to_every_eligible_slide(tmp_path):
    path = _simple_deck(tmp_path)
    out_path = tmp_path / "out.pptx"

    result = apply_retemplate(str(path), str(out_path), template_path=TEMPLATE_PATH)

    assert result.transformed == [1, 2]
    assert result.skipped == [3]
    assert _no_duplicate_zip_entries(out_path)
    assert _no_malformed_xml(out_path) == []

    prs = Presentation(str(out_path))
    assert len(prs.slides) == 3
    assert prs.slides[0].shapes.title.text_frame.text == "First Slide"
    assert prs.slides[1].shapes.title.text_frame.text == "Second Slide"
    assert any(s.shape_type is not None and s.shape_type.name == "TABLE" for s in prs.slides[2].shapes)  # slide 3 untouched


def test_apply_retemplate_honors_accepted_indexes_subset(tmp_path):
    path = _simple_deck(tmp_path)
    out_path = tmp_path / "out.pptx"

    result = apply_retemplate(str(path), str(out_path), accepted_indexes={1}, template_path=TEMPLATE_PATH)

    assert result.transformed == [1]
    assert result.skipped == [2, 3]

    prs = Presentation(str(out_path))
    assert prs.slides[0].shapes.title.text_frame.text == "First Slide"
    assert prs.slides[1].shapes.title.text_frame.text == "Second Slide"  # untouched original, not rebuilt


def test_rebuild_non_standard_layout_slides_only_rebuilds_flagged_slides(tmp_path):
    """Fix's opt-in "also rebuild non-standard layouts" checkbox: finds
    non_standard_layout via a real audit, rebuilds only those slides
    (verbatim), leaves everything else -- including OTHER eligible
    slides that just happen to already be fine -- completely alone."""
    from deckguard.retemplate import rebuild_non_standard_layout_slides

    path = _simple_deck(tmp_path)  # all 3 slides use the test fixture's own non-org-template layout
    out_path = tmp_path / "out.pptx"

    result = rebuild_non_standard_layout_slides(str(path), str(out_path), template_path=TEMPLATE_PATH)

    # Every slide here is on a non-standard layout (the test fixture's
    # own, unrelated to the org template) -- the table slide is still
    # excluded since it's ineligible for verbatim carryover regardless.
    assert result.transformed == [1, 2]
    assert result.skipped == [3]

    prs = Presentation(str(out_path))
    approved_names = {layout.name for master in Presentation(str(TEMPLATE_PATH)).slide_masters for layout in master.slide_layouts}
    assert prs.slides[0].slide_layout.name in approved_names
    assert prs.slides[1].slide_layout.name in approved_names
    assert prs.slides[0].shapes.title.text_frame.text == "First Slide"


def test_rebuild_non_standard_layout_slides_is_a_no_op_when_everything_is_already_approved(tmp_path):
    from deckguard.retemplate import apply_rebrand, rebuild_non_standard_layout_slides

    path = _simple_deck(tmp_path)
    rebranded_path = tmp_path / "rebranded.pptx"
    apply_rebrand(str(path), str(rebranded_path), template_path=TEMPLATE_PATH)  # lands every eligible slide on an approved layout

    out_path = tmp_path / "out.pptx"
    result = rebuild_non_standard_layout_slides(str(rebranded_path), str(out_path), template_path=TEMPLATE_PATH)

    assert result.transformed == []


def test_apply_retemplate_is_a_no_op_when_nothing_is_eligible(tmp_path):
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))
    path = tmp_path / "src.pptx"
    prs.save(str(path))
    out_path = tmp_path / "out.pptx"

    result = apply_retemplate(str(path), str(out_path), template_path=TEMPLATE_PATH)

    assert result.transformed == []
    assert result.skipped == [1]
    reopened = Presentation(str(out_path))
    assert any(s.shape_type is not None and s.shape_type.name == "TABLE" for s in reopened.slides[0].shapes)


def test_apply_retemplate_replacing_three_or_more_slides_does_not_collide_partnames(tmp_path):
    """Regression test for a real bug: python-pptx's own add_slide()
    computes each new slide's partname as just "current slide count + 1"
    (not a scan for a free name -- see PresentationPart._next_slide_partname).
    Interleaving a delete between two add_slide() calls shrinks that count
    and can make a LATER add_slide() reuse a partname a PRIOR new slide
    from the same run is still using, silently merging two slides onto one
    XML part (a real zipfile "Duplicate name" corruption on save, not just
    clutter). Needs >= 3 replaced slides to actually exercise the shrink/
    reuse cycle a single replacement can't trigger."""
    import warnings

    prs = new_deck()
    titles = ["Alpha", "Bravo", "Charlie", "Delta"]
    for t in titles:
        slide = add_slide(prs)
        title_run(slide).text = t
        body_run(slide).text = f"{t} body"
    path = tmp_path / "src.pptx"
    prs.save(str(path))
    out_path = tmp_path / "out.pptx"

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        result = apply_retemplate(str(path), str(out_path), template_path=TEMPLATE_PATH)
        dupes = [str(w.message) for w in caught if "Duplicate name" in str(w.message)]

    assert dupes == []
    assert result.transformed == [1, 2, 3, 4]
    assert _no_duplicate_zip_entries(out_path)
    assert _no_malformed_xml(out_path) == []

    prs2 = Presentation(str(out_path))
    assert len(prs2.slides) == 4
    assert [s.shapes.title.text_frame.text for s in prs2.slides] == titles


def test_apply_retemplate_carries_images_into_the_new_layouts_picture_placeholder(tmp_path):
    prs = new_deck()
    slide = add_slide(prs, layout_idx=6)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    title_box.text_frame.text = "irrelevant"  # not a real TITLE placeholder -- just body-like text
    img_path = make_pattern_png(tmp_path / "img.png", seed=2)
    add_picture(slide, str(img_path), left_in=1, top_in=2)
    path = tmp_path / "src.pptx"
    prs.save(str(path))
    out_path = tmp_path / "out.pptx"

    result = apply_retemplate(str(path), str(out_path), template_path=TEMPLATE_PATH)

    assert result.transformed == [1]
    prs2 = Presentation(str(out_path))
    # A filled picture placeholder reports shape_type PLACEHOLDER, not
    # PICTURE (python-pptx's PlaceholderPicture) -- check by trying to
    # read an image blob from every shape rather than filtering by type.
    blobs = []
    for s in prs2.slides[0].shapes:
        try:
            blobs.append(s.image.blob)
        except (AttributeError, ValueError):
            continue
    assert img_path.read_bytes() in blobs


# --------------------------------------------------------------------------
# apply_rebrand -- verbatim carryover + layout variety + cover/end swap
# --------------------------------------------------------------------------


def _cover_content_end_deck(tmp_path):
    prs = new_deck()
    cover = add_slide(prs)
    title_run(cover).text = "Annual Review"

    content1 = add_slide(prs)
    title_run(content1).text = "Highlights"
    body_run(content1).text = "Grew nicely"

    content2 = add_slide(prs)
    title_run(content2).text = "Next Steps"
    body_run(content2).text = "Keep going"

    end = add_slide(prs)
    title_run(end).text = "Thank You"

    path = tmp_path / "src.pptx"
    prs.save(str(path))
    return path


def test_apply_rebrand_swaps_cover_and_end_slides_onto_current_brand_layout(tmp_path):
    path = _cover_content_end_deck(tmp_path)
    out_path = tmp_path / "out.pptx"

    result = apply_rebrand(str(path), str(out_path), template_path=TEMPLATE_PATH)

    assert result.transformed == [1, 2, 3, 4]
    by_index = {p.slide_index: p for p in result.proposals}
    assert by_index[1].layout_name == "Cover B"
    assert by_index[4].layout_name == "Outro"

    prs2 = Presentation(str(out_path))
    texts = [s.shapes.title.text_frame.text for s in prs2.slides]
    assert texts == ["Annual Review", "Highlights", "Next Steps", "Thank You"]


def test_apply_rebrand_never_changes_wording(tmp_path):
    """The whole point of brand mode: title/body text survives character
    for character -- no condensing, no rewording, regardless of what
    match_layout or the cover/end override picks."""
    path = _cover_content_end_deck(tmp_path)
    out_path = tmp_path / "out.pptx"

    apply_rebrand(str(path), str(out_path), template_path=TEMPLATE_PATH)

    prs2 = Presentation(str(out_path))
    body_texts = {
        shp.text_frame.text for slide in prs2.slides for shp in slide.shapes
        if shp.has_text_frame and shp.text_frame.text.strip()
    }
    assert "Grew nicely" in body_texts
    assert "Keep going" in body_texts


def test_apply_rebrand_cover_with_no_source_image_gets_a_real_editable_picture(tmp_path):
    """Regression test for a real bug report: a cover/end slide with no
    image of its own on the source deck left its picture placeholder
    completely empty on the rebranded slide too -- <p:spPr/>, no
    <p:blipFill> -- which PowerPoint renders by inheriting the LAYOUT's
    own baked-in default photo, but only ever offers "Save as Picture"
    for (never "Change Picture", since there's no picture object owned
    by the slide itself). The swapped-in Cover B/Outro picture
    placeholder must come out with a real, slide-owned picture --
    materialized from the layout's own default image -- so it's
    independently editable like any other picture."""
    path = _cover_content_end_deck(tmp_path)
    out_path = tmp_path / "out.pptx"

    apply_rebrand(str(path), str(out_path), template_path=TEMPLATE_PATH)

    prs2 = Presentation(str(out_path))
    for slide in (prs2.slides[0], prs2.slides[-1]):
        pic = next(s for s in slide.shapes if "Picture" in s.name)
        xml = pic._element.xml
        assert "blipFill" in xml and "r:embed" in xml, f"{slide.slide_layout.name}'s picture placeholder was left empty"


def test_apply_rebrand_cover_borrows_reference_decks_own_photo_over_the_generic_default(tmp_path):
    """The "Learn from a reference" flow's per-run image borrow: when the
    source deck's cover has no picture of its own, a reference deck's own
    cover photo (when given) is the deck-specific answer for THIS run --
    it should win over the org template's generic stock photo fallback."""
    path = _cover_content_end_deck(tmp_path)
    out_path = tmp_path / "out.pptx"

    ref_prs = new_deck()
    ref_cover = add_slide(ref_prs)
    title_run(ref_cover).text = "Reference Cover"
    img_path = tmp_path / "ref_cover.png"
    make_pattern_png(img_path, seed=3)
    add_picture(ref_cover, str(img_path))
    ref_path = tmp_path / "reference.pptx"
    ref_prs.save(str(ref_path))

    apply_rebrand(str(path), str(out_path), template_path=TEMPLATE_PATH, reference_path=str(ref_path))

    prs2 = Presentation(str(out_path))
    pic = next(s for s in prs2.slides[0].shapes if "Picture" in s.name)
    assert pic.image.blob == img_path.read_bytes()


def test_apply_rebrand_reparents_a_middle_slide_onto_the_reference_decks_own_layout(tmp_path):
    """The "Learn from a reference" flow's layout carryover: when an old
    slide and its reference counterpart at the SAME index already sit on
    a layout of the exact same name -- here "Title and Content", which
    (like the real deck pair this was built against) isn't in the org
    template's own candidate list at all -- that name is ground truth for
    what the slide should look like. Content stays 100% untouched (proven
    here by title text AND a non-placeholder decorative shape surviving
    verbatim, since this bypasses classify_slide/match_layout's content
    caps entirely); only the layout/master is refreshed from the
    reference's own copy."""
    old = new_deck()
    title_run(add_slide(old, layout_idx=0)).text = "Cover"
    middle = add_slide(old, layout_idx=1)  # "Title and Content"
    title_run(middle).text = "Middle slide"
    add_rectangle(middle, name="Chip", fill_hex="1450F5")  # already-approved KONE blue -- fix_deck leaves it alone
    title_run(add_slide(old, layout_idx=0)).text = "The End"
    old_path = tmp_path / "old.pptx"
    old.save(str(old_path))

    ref = new_deck()
    add_slide(ref, layout_idx=0)
    add_slide(ref, layout_idx=1)  # same layout name, no content needed -- only its NAME matters
    add_slide(ref, layout_idx=0)
    ref_path = tmp_path / "reference.pptx"
    ref.save(str(ref_path))

    out_path = tmp_path / "out.pptx"
    result = apply_rebrand(str(old_path), str(out_path), template_path=TEMPLATE_PATH, reference_path=str(ref_path))

    assert result.reference_layout_indices == [2]
    assert 2 in result.transformed
    assert 2 not in result.skipped

    prs2 = Presentation(str(out_path))
    reparented = prs2.slides[1]
    assert reparented.slide_layout.name == "Title and Content"
    assert reparented.shapes.title.text_frame.text == "Middle slide"
    chip = next(s for s in reparented.shapes if s.name == "Chip")
    assert str(chip.fill.fore_color.rgb) == "1450F5"


def test_apply_rebrand_does_not_reparent_when_layout_names_differ(tmp_path):
    """Negative case: an old slide and its reference counterpart at the
    same index sitting on DIFFERENTLY-named layouts have no confirmed
    correspondence -- falls back to ordinary classify_slide/match_layout
    against the org template, same as with no reference_path at all."""
    old = new_deck()
    title_run(add_slide(old, layout_idx=0)).text = "Cover"
    middle = add_slide(old, layout_idx=1)  # "Title and Content"
    title_run(middle).text = "Middle slide"
    title_run(add_slide(old, layout_idx=0)).text = "The End"
    old_path = tmp_path / "old.pptx"
    old.save(str(old_path))

    ref = new_deck()
    add_slide(ref, layout_idx=0)
    add_slide(ref, layout_idx=2)  # "Section Header" -- different name at the same index
    add_slide(ref, layout_idx=0)
    ref_path = tmp_path / "reference.pptx"
    ref.save(str(ref_path))

    out_path = tmp_path / "out.pptx"
    result = apply_rebrand(str(old_path), str(out_path), template_path=TEMPLATE_PATH, reference_path=str(ref_path))

    assert result.reference_layout_indices == []


def test_apply_rebrand_reference_layout_carryover_survives_full_deckguard_fix_afterward(tmp_path):
    """The reparented slide's own package (rels/media/master) must be
    completely valid -- e.g. no dangling relationships or duplicate zip
    entries -- since fix_deck touches every shape/run in the deck
    immediately afterward as part of apply_rebrand itself."""
    old = new_deck()
    title_run(add_slide(old, layout_idx=0)).text = "Cover"
    middle = add_slide(old, layout_idx=1)
    title_run(middle).text = "Middle slide"
    title_run(add_slide(old, layout_idx=0)).text = "The End"
    old_path = tmp_path / "old.pptx"
    old.save(str(old_path))

    ref = new_deck()
    add_slide(ref, layout_idx=0)
    add_slide(ref, layout_idx=1)
    add_slide(ref, layout_idx=0)
    ref_path = tmp_path / "reference.pptx"
    ref.save(str(ref_path))

    out_path = tmp_path / "out.pptx"
    apply_rebrand(str(old_path), str(out_path), template_path=TEMPLATE_PATH, reference_path=str(ref_path))

    assert _no_duplicate_zip_entries(out_path)
    assert _no_malformed_xml(out_path) == []
    Presentation(str(out_path))  # reopens cleanly


def test_reparent_slide_layout_freezes_inherited_geometry(tmp_path):
    """Regression test for a real bug: a placeholder with no explicit
    position of its own (purely inherited from its layout) silently
    followed the NEW layout's own (differently positioned) placeholder
    after a reparent -- while a placeholder that already had an explicit
    position on the slide stayed put, since reparenting never touches
    shape content. On a real deck this produced an internally
    inconsistent slide: a title jumped to the reference layout's spot
    while the body text below it stayed at the old slide's position,
    and the two overlapped into unreadable garbled text. Freezing the
    inherited geometry BEFORE the layout swap keeps the slide
    self-consistent regardless of where the new layout's own version of
    that placeholder sits."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=1)  # "Title and Content"
    title = slide.shapes.title
    old_pos = (title.left, title.top, title.width, title.height)

    new_layout = prs.slide_layouts[2]  # "Section Header" -- a different title position
    new_title_pos = None
    for ph in new_layout.placeholders:
        if ph.placeholder_format.idx == 0:
            new_title_pos = (ph.left, ph.top, ph.width, ph.height)
    assert new_title_pos != old_pos  # the two layouts must actually differ for this test to mean anything

    _reparent_slide_layout(slide, new_layout)

    assert slide.slide_layout.name == "Section Header"
    reopened_title = slide.shapes.title
    assert (reopened_title.left, reopened_title.top, reopened_title.width, reopened_title.height) == old_pos


def test_freeze_placeholder_geometry_leaves_an_explicit_position_untouched(tmp_path):
    """A placeholder that already has its own explicit position (was
    manually adjusted at some point) must not be touched -- freezing is
    only for placeholders that were purely inheriting."""
    prs = new_deck()
    slide = add_slide(prs, layout_idx=1)
    title = slide.shapes.title
    title.left, title.top, title.width, title.height = 111111, 222222, 333333, 444444

    _freeze_placeholder_geometry(slide)

    assert (title.left, title.top, title.width, title.height) == (111111, 222222, 333333, 444444)


def test_apply_rebrand_picks_varied_layouts_instead_of_repeating_one(tmp_path):
    """Several ordinary content slides, identically shaped (title + one
    body block, no images) -- without anti-repeat scoring they'd all pick
    the same single tightest-fitting layout; with it, match_layout should
    spread them across the tied candidates instead."""
    prs = new_deck()
    for i in range(4):
        slide = add_slide(prs)
        title_run(slide).text = f"Slide {i}"
        body_run(slide).text = f"Body {i}"
    path = tmp_path / "src.pptx"
    prs.save(str(path))
    out_path = tmp_path / "out.pptx"

    result = apply_rebrand(str(path), str(out_path), template_path=TEMPLATE_PATH)

    layouts_used = [p.layout_name for p in result.proposals if p.eligible]
    assert len(set(layouts_used)) > 1


def test_apply_rebrand_still_hard_skips_a_table_and_leaves_it_untouched(tmp_path):
    prs = new_deck()
    cover = add_slide(prs)
    title_run(cover).text = "Deck"
    table_slide = add_slide(prs)
    table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))
    path = tmp_path / "src.pptx"
    prs.save(str(path))
    out_path = tmp_path / "out.pptx"

    result = apply_rebrand(str(path), str(out_path), template_path=TEMPLATE_PATH)

    assert result.skipped == [2]
    proposal = next(p for p in result.proposals if p.slide_index == 2)
    assert proposal.reason == "contains a table"


def test_apply_rebrand_runs_fix_deck_for_color_and_font_compliance(tmp_path):
    """Brand mode's whole promise is 'follow the color/font guidelines' --
    `apply_retemplate` alone never runs `fix_deck`, so a rebuilt slide's
    freshly-written text (no explicit run color, by design -- see
    compose.py's own module docstring) is left as an inherited, therefore
    audit-unresolvable color. `apply_rebrand` finishes that off, the same
    way `deckguard fix` would."""
    from deckguard.config import default_config_path, load_config
    from deckguard.inventory import build_inventory
    from deckguard.rules_engine import audit_deck

    path = _cover_content_end_deck(tmp_path)
    retemplate_only = tmp_path / "retemplate_only.pptx"
    apply_retemplate(str(path), str(retemplate_only), template_path=TEMPLATE_PATH)
    rebrand_out = tmp_path / "rebrand.pptx"
    apply_rebrand(str(path), str(rebrand_out), template_path=TEMPLATE_PATH)

    config = load_config(default_config_path())
    retemplate_contrast = [
        v for v in audit_deck(build_inventory(Presentation(str(retemplate_only))), config) if v.rule == "text_contrast"
    ]
    rebrand_contrast = [
        v for v in audit_deck(build_inventory(Presentation(str(rebrand_out))), config) if v.rule == "text_contrast"
    ]

    assert retemplate_contrast  # confirms the premise: retemplate alone leaves this unresolved
    assert not rebrand_contrast  # apply_rebrand's fix_deck pass resolves it


# --------------------------------------------------------------------------
# rebuild_slides_as_dividers -- brand mode's --review divider rebuild
# --------------------------------------------------------------------------


def test_rebuild_slides_as_dividers_replaces_content_with_just_a_title(tmp_path):
    prs = new_deck()
    slide1 = add_slide(prs)
    title_run(slide1).text = "Highlights"
    body_run(slide1).text = "Real content that must not survive"
    slide2 = add_slide(prs)
    title_run(slide2).text = "Appendix"
    path = tmp_path / "src.pptx"
    prs.save(str(path))
    out_path = tmp_path / "out.pptx"

    rebuilt = rebuild_slides_as_dividers(str(path), str(out_path), TEMPLATE_PATH, {2: "Appendix"})

    assert rebuilt == [2]
    prs2 = Presentation(str(out_path))
    assert len(prs2.slides) == 2
    # slide 1 completely untouched
    assert prs2.slides[0].shapes.title.text_frame.text == "Highlights"
    body_texts = [
        s.text_frame.text for s in prs2.slides[0].shapes
        if s.has_text_frame and s.text_frame.text.strip()
    ]
    assert "Real content that must not survive" in body_texts
    # slide 2 rebuilt as a divider with only the given title
    assert prs2.slides[1].slide_layout.name in ("Section divider A", "Section divider B")
    assert prs2.slides[1].shapes.title.text_frame.text == "Appendix"


def test_rebuild_slides_as_dividers_alternates_variants_across_multiple_dividers(tmp_path):
    prs = new_deck()
    for title in ("Appendix", "Q&A", "Thank You"):
        slide = add_slide(prs)
        title_run(slide).text = title
    path = tmp_path / "src.pptx"
    prs.save(str(path))
    out_path = tmp_path / "out.pptx"

    rebuilt = rebuild_slides_as_dividers(str(path), str(out_path), TEMPLATE_PATH, {1: "Appendix", 2: "Q&A", 3: "Thank You"})

    assert rebuilt == [1, 2, 3]
    prs2 = Presentation(str(out_path))
    layouts_used = [s.slide_layout.name for s in prs2.slides]
    assert layouts_used == ["Section divider A", "Section divider B", "Section divider A"]
