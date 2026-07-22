from deckguard import learn
from deckguard.config import default_config_path, load_config
from tests.helpers import add_slide, body_run, new_deck, set_run, title_run

BASE_CONFIG = {
    "colors": {"approved": ["#1450F5", "#FFFFFF"], "remap": {}, "tolerance": 0},
    "fonts": {"approved": ["Inter", "Inter Semi Bold"], "remap": {}},
}


def _deck_with_run(font, color_hex, text="Body copy"):
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text=text, font=font, color_hex=color_hex)
    return prs


def test_high_confidence_color_and_font_match():
    old_prs = _deck_with_run("Arial", "AABBCC")
    new_prs = _deck_with_run("Inter", "1450F5")

    result = learn.learn(old_prs, new_prs, BASE_CONFIG)

    assert len(result.color_proposals) == 1
    cp = result.color_proposals[0]
    assert cp.old_hex == "AABBCC" and cp.new_hex == "1450F5"
    assert cp.confidence == "high"

    assert len(result.font_proposals) == 1
    fp = result.font_proposals[0]
    assert fp.old_font == "Arial" and fp.new_font == "Inter"
    assert fp.confidence == "high"


def test_low_confidence_when_counts_diverge():
    prs_old = new_deck()
    slide = add_slide(prs_old)
    set_run(title_run(slide), text="Title", font="Inter", color_hex="AABBCC")

    prs_new = new_deck()
    slide2 = add_slide(prs_new)
    for i in range(10):
        set_run(body_run(slide2) if i == 0 else title_run(slide2), text=f"t{i}", font="Inter", color_hex="1450F5")

    result = learn.learn(prs_old, prs_new, BASE_CONFIG)
    assert len(result.color_proposals) == 1
    assert result.color_proposals[0].confidence == "low"


def test_already_approved_or_remapped_colors_are_not_proposed():
    old_prs = _deck_with_run("Inter", "1450F5")  # already approved
    new_prs = _deck_with_run("Inter", "FFFFFF")
    result = learn.learn(old_prs, new_prs, BASE_CONFIG)
    assert result.color_proposals == []

    config_with_remap = {
        "colors": {"approved": ["#1450F5", "#FFFFFF"], "remap": {"#AABBCC": "#1450F5"}, "tolerance": 0},
        "fonts": {"approved": ["Inter"], "remap": {}},
    }
    old_prs2 = _deck_with_run("Inter", "AABBCC")
    new_prs2 = _deck_with_run("Inter", "FFFFFF")
    result2 = learn.learn(old_prs2, new_prs2, config_with_remap)
    assert result2.color_proposals == []  # already has a remap entry, not re-proposed


def test_unmatched_color_when_no_candidate_role_in_new_deck():
    old_prs = new_deck()
    slide = add_slide(old_prs)
    set_run(title_run(slide), text="Title", font="Inter", color_hex="1450F5")
    title_run(slide).text = "Title"
    # give the title shape a line color with no counterpart role in the new deck
    from pptx.dml.color import RGBColor

    old_prs.slides[0].shapes.title.line.color.rgb = RGBColor.from_string("AABBCC")

    new_prs = new_deck()
    add_slide(new_prs)  # no line color set anywhere

    result = learn.learn(old_prs, new_prs, BASE_CONFIG)
    assert any(role == "line" and hexval == "AABBCC" for role, hexval, _count in result.unmatched_old_colors)


def test_apply_learned_merges_high_confidence_only_by_default():
    old_prs = _deck_with_run("Arial", "AABBCC")
    new_prs = _deck_with_run("Inter", "1450F5")
    result = learn.learn(old_prs, new_prs, BASE_CONFIG)

    updated = learn.apply_learned(BASE_CONFIG, result)
    assert updated["colors"]["remap"]["#AABBCC"] == "#1450F5"
    assert "#1450F5" in updated["colors"]["approved"]  # already there, no dup
    assert updated["fonts"]["remap"]["Arial"] == "Inter"
    # original config untouched
    assert BASE_CONFIG["colors"]["remap"] == {}


def test_apply_learned_min_confidence_low_includes_low_confidence():
    prs_old = new_deck()
    slide = add_slide(prs_old)
    set_run(title_run(slide), text="Title", font="Inter", color_hex="AABBCC")
    prs_new = new_deck()
    slide2 = add_slide(prs_new)
    for i in range(10):
        set_run(body_run(slide2) if i == 0 else title_run(slide2), text=f"t{i}", font="Inter", color_hex="1450F5")

    result = learn.learn(prs_old, prs_new, BASE_CONFIG)
    updated_default = learn.apply_learned(BASE_CONFIG, result)
    assert updated_default["colors"]["remap"] == {}  # low confidence, not applied by default

    updated_low = learn.apply_learned(BASE_CONFIG, result, min_confidence="low")
    assert updated_low["colors"]["remap"]["#AABBCC"] == "#1450F5"


def test_write_learned_to_yaml_preserves_comments_and_dedupes(tmp_path):
    rules_path = tmp_path / "brand_rules.yaml"
    rules_path.write_text(
        "colors:\n"
        "  approved:\n"
        "    - \"#1450F5\"   # KONE Blue\n"
        "    - \"#FFFFFF\"   # White\n"
        "  remap: {}\n"
        "  tolerance: 0\n"
        "fonts:\n"
        "  approved:\n"
        "    - \"Inter\"   # brand font\n"
        "  remap: {}\n",
        encoding="utf-8",
    )

    old_prs = _deck_with_run("Arial", "AABBCC")
    new_prs = _deck_with_run("Inter", "1450F5")
    result = learn.learn(old_prs, new_prs, BASE_CONFIG)

    applied = learn.write_learned_to_yaml(rules_path, result)
    assert applied == 2  # one color, one font

    text = rules_path.read_text(encoding="utf-8")
    assert "# KONE Blue" in text
    assert "# brand font" in text
    assert "#AABBCC" in text and "#1450F5" in text.upper() or "1450F5" in text
    assert "Arial" in text

    # re-applying is idempotent -- nothing new to add
    applied_again = learn.write_learned_to_yaml(rules_path, result)
    assert applied_again == 0


def test_default_config_learn_smoke():
    # sanity: learn() runs cleanly against the shipped default config with
    # two trivially-identical decks (nothing off-brand to propose).
    config = load_config(default_config_path())
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Title", font="Inter Semi Bold", color_hex="141414")
    result = learn.learn(prs, prs, config)
    assert result.color_proposals == []
    assert result.font_proposals == []
