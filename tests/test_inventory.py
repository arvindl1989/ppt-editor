from pptx.dml.color import RGBColor

from deckguard.inventory import build_inventory
from tests.helpers import add_picture, add_rectangle, add_slide, body_run, make_pattern_png, new_deck, set_run, title_run


def test_build_inventory_basic_shape_fields():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(title_run(slide), text="Hello", font="Inter", bold=True, size_pt=32, color_hex="141414")

    inv = build_inventory(prs)
    assert len(inv.slides) == 1
    shapes = inv.slides[0].shapes
    title_shape = next(s for s in shapes if s.placeholder_type in ("TITLE", "CENTER_TITLE"))
    run = title_shape.paragraphs[0].runs[0]
    assert run.text == "Hello"
    assert run.font_raw == "Inter"
    assert run.bold is True
    assert run.size_pt == 32
    assert run.color.hex == "141414"


def test_build_inventory_picks_up_image_phash(tmp_path):
    prs = new_deck()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = make_pattern_png(tmp_path / "img.png", seed=5)
    add_picture(slide, str(img))

    inv = build_inventory(prs)
    shape = inv.slides[0].shapes[0]
    assert shape.image is not None
    assert shape.image.phash is not None
    assert len(shape.image.phash) == 16


def test_build_inventory_all_upper_detection():
    prs = new_deck()
    slide = add_slide(prs)
    set_run(body_run(slide), text="SHOUT", font="Inter")
    set_run(title_run(slide), text="Mixed Case", font="Inter")

    inv = build_inventory(prs)
    runs = {r.text: r for s in inv.slides[0].shapes for p in s.paragraphs for r in p.runs}
    assert runs["SHOUT"].all_upper is True
    assert runs["Mixed Case"].all_upper is False


def _line_element(shape):
    from pptx.oxml.ns import qn

    spPr = shape._element.find(qn("p:spPr"))
    return spPr.find(qn("a:ln")) if spPr is not None else None


def test_build_inventory_never_adds_a_line_element_as_a_side_effect():
    """Regression test for a real deck-corrupting bug found in production:
    python-pptx's LineFormat.color forces the line's fill type to SOLID as
    a documented side effect of merely being read, turning "no <a:ln> at
    all" (the normal state of a plain text box) into an empty
    `<a:ln><a:solidFill/></a:ln>` with no color specified -- which
    PowerPoint renders as a visible default-colored border around every
    text box in the deck. Plain inspection/audit must never touch shapes
    this way."""
    from pptx.util import Inches

    prs = new_deck()
    slide = add_slide(prs)
    title_run(slide).text = "Title"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    textbox.text_frame.text = "Plain text box, no explicit line"
    assert _line_element(textbox) is None  # sanity: matches real-world default shape state

    build_inventory(prs)  # read-only -- must not mutate anything

    assert _line_element(textbox) is None, "inspection added a line element that wasn't there before"


def test_build_inventory_reports_existing_line_color_correctly():
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Colored box", fill_hex="1450F5")
    box.line.color.rgb = RGBColor.from_string("FF5F28")

    inv = build_inventory(prs)
    shape = next(s for s in inv.slides[0].shapes if s.name == "Colored box")
    assert shape.line is not None
    assert shape.line.color.hex == "FF5F28"


def _run_element(run):
    return run._r


def test_build_inventory_never_adds_a_solid_fill_to_run_text_color_as_a_side_effect():
    """Regression test for a real, severe deck-corrupting bug found in
    production: python-pptx's Font.color forces the run's fill type to
    SOLID as a documented side effect of merely being read (it calls
    self.fill.solid() whenever the type isn't already SOLID), turning "no
    color override, inherits from theme/master" into an explicit-but-empty
    <a:solidFill/> with no color specified -- which PowerPoint renders as
    black regardless of what the run actually inherited. This corrupted
    every run without an explicit color in every deck ever run through
    inspect/audit/fix, silently overriding inherited/theme text colors to
    black. Plain inspection must never touch a run this way."""
    from lxml import etree

    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Panel", fill_hex="1450F5")
    box.text_frame.text = "Plain text, no explicit color"
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.bold = True  # a real rPr exists, but with no solidFill child
    before = etree.tostring(_run_element(run)).decode()
    assert "solidFill" not in before  # sanity: matches real-world unstyled-color state

    build_inventory(prs)  # read-only -- must not mutate anything

    after = etree.tostring(_run_element(run)).decode()
    assert "solidFill" not in after, "inspection added a solidFill that wasn't there before"


def test_build_inventory_reports_run_with_no_explicit_color_as_none():
    prs = new_deck()
    slide = add_slide(prs)
    box = add_rectangle(slide, name="Panel", fill_hex="1450F5")
    box.text_frame.text = "No explicit color"

    inv = build_inventory(prs)
    shape = next(s for s in inv.slides[0].shapes if s.name == "Panel")
    run = shape.paragraphs[0].runs[0]
    assert run.color.kind == "none"


def test_an_image_in_a_picture_placeholder_is_inventoried(tmp_path):
    """A filled picture placeholder reports its type as PLACEHOLDER, not
    PICTURE. Keying on the type made those photos invisible to the
    inventory -- and so to previews, the visual audit, and the logo
    checks that read it. On a real deck the images were present in the
    .pptx and absent from every view of it."""
    from pptx import Presentation

    from deckguard.inventory import build_inventory
    from deckguard.legacy.slide_import import default_template_path
    from tests.helpers import make_solid_png

    photo = tmp_path / "p.png"
    make_solid_png(photo)

    prs = Presentation(str(default_template_path()))
    layout = next(
        (l for m in prs.slide_masters for l in m.slide_layouts
         if any(ph.placeholder_format.type and ph.placeholder_format.type.name == "PICTURE"
                for ph in l.placeholders)),
        None,
    )
    assert layout is not None, "the bundled template has picture layouts"

    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.type and ph.placeholder_format.type.name == "PICTURE":
            ph.insert_picture(str(photo))
            break
    path = tmp_path / "d.pptx"
    prs.save(str(path))

    inv = build_inventory(Presentation(str(path)))
    with_images = [s for s in inv.slides[0].shapes if s.image is not None]
    assert len(with_images) == 1
    assert with_images[0].image.size_bytes > 0
