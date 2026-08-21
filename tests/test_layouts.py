"""Tests for layouts.py -- the master's geometry read as data.

The parsing tests run on inline fixtures so the rules are pinned
regardless of which spec revision is installed. The tests that read the
shipped `LAYOUTS.md` / `ARCHETYPES.md` are there to catch a spec update
that silently stops parsing -- the failure mode this module exists to
prevent is not a crash, it is quietly generating nothing.
"""

import pytest

from deckguard import layouts as L

LAYOUTS_FIXTURE = """\
# geometry

## Covers

### Cover A
`slideLayout1`

- **Picture** — 0, 0, 1280 × 421 · white · image
- **Logo** — 45, 45, 81 × 31
- **Title** — 45, 429, 578 × 155
- **Footer** — 215, 658, 408 × 19

### Three content A
`slideLayout24`

- **Text/body** — 45, 227, 374 × 403
- **Text/body** — 453, 227, 374 × 403
- **Text/body** — 861, 227, 374 × 403
- **Logo** — 1153, 45, 81 × 31 · image

### Quote A
`slideLayout40`

- **Background** — 453, 136, 782 × 493
- **Title** — 45, 136, 272 × 104
- **Text/body** — 510, 212, 657 × 349 · white
"""

ARCHETYPES_FIXTURE = """\
# glossary

## Grade A — most used

| Archetype | What it's for | Master | Status |
| --- | --- | --- | --- |
| `COVER_A_CUT4` | Title slide. Alias: `cover_cut` | `slideLayout1` | built · 01 |
| `THREE_CONTENT` | Three columns. | `slideLayout24` | built · 09 |
| `TIMELINE` | Roadmap. | no master | built · 10 |

## Grade D — good to have

| Archetype | What it's for | Master | Status |
| --- | --- | --- | --- |
| `THREE_CONTENT_B` | Twin of `THREE_CONTENT`. | `slideLayout25` | twin |
"""


@pytest.fixture
def spec(monkeypatch, tmp_path):
    (tmp_path / "LAYOUTS.md").write_text(LAYOUTS_FIXTURE)
    (tmp_path / "ARCHETYPES.md").write_text(ARCHETYPES_FIXTURE)
    monkeypatch.setattr(L, "spec_dir", lambda: tmp_path)
    return tmp_path


# --------------------------------------------------------------------------
# parsing
# --------------------------------------------------------------------------


def test_geometry_parses_into_boxes():
    layouts = L.parse_layouts(LAYOUTS_FIXTURE)
    assert set(layouts) == {"slideLayout1", "slideLayout24", "slideLayout40"}
    cover = layouts["slideLayout1"]
    assert cover.name == "Cover A"
    picture = next(b for b in cover.boxes if b.role == "Picture")
    assert (picture.x, picture.y, picture.w, picture.h) == (0, 0, 1280, 421)
    assert "image" in picture.mods


def test_chrome_is_recognised_and_excluded_from_content():
    """An archetype that places its own logo produces two of them once
    the master's own frames are repaired. Chrome belongs to the layout."""
    cover = L.parse_layouts(LAYOUTS_FIXTURE)["slideLayout1"]
    assert [b.role for b in cover.content_boxes()] == ["Picture", "Title"]


def test_archetype_rows_carry_their_grade_and_binding():
    archetypes = L.parse_archetypes(ARCHETYPES_FIXTURE)
    cover = archetypes["COVER_A_CUT4"]
    assert cover.grade == "A"
    assert cover.master == "slideLayout1"
    assert cover.aliases == ("cover_cut",)
    assert cover.is_built and cover.engine_key == "cover_a_cut4"
    assert archetypes["TIMELINE"].master is None
    assert archetypes["THREE_CONTENT_B"].twin_of == "THREE_CONTENT"


def test_an_ungraded_archetypes_file_yields_nothing_rather_than_guessing():
    assert L.parse_archetypes("| `FOO` | bar | `slideLayout1` | built · 01 |") == {}


# --------------------------------------------------------------------------
# role binding
# --------------------------------------------------------------------------


def test_a_simple_layout_becomes_regions(spec):
    built = L.build_archetypes()
    cover = built["cover_a_cut4"]
    assert [r["role"] for r in cover["regions"]] == ["picture", "title"]
    assert cover["regions"][0]["content"] == "image"
    assert "groups" not in cover


def test_equal_boxes_at_one_y_become_one_repeating_group():
    """Three 374x403 boxes at y=227 are not three regions -- they are
    one group of three, which is the only form the engine can expand
    over a content list of a different length.

    Tested on `_bind_roles` rather than `build_archetypes` because
    THREE_CONTENT carries a reference refinement that replaces this
    binding; the binding still has to be right underneath it."""
    layouts = L.parse_layouts(LAYOUTS_FIXTURE)
    three = L._bind_roles(layouts["slideLayout24"])
    assert three["regions"] == []
    (group,) = three["groups"]
    assert group["origins"] == [[45, 227], [453, 227], [861, 227]]
    assert [r["box"] for r in group["regions"]] == [[0, 0, 374, 403]]


def test_a_background_box_becomes_a_panel_not_a_region(spec):
    """The engine's own `panel` role is hardcoded to KONE Blue and these
    come in five colours, so panels are carried separately and painted
    by `render` before the engine draws anything."""
    layouts = L.parse_layouts(LAYOUTS_FIXTURE)
    quote = L._bind_roles(layouts["slideLayout40"])
    assert quote["panels"] == [{"box": [453, 136, 782, 493], "fill": "1450F5"}]
    assert not any(r["role"] == "panel" for r in quote["regions"])


def test_white_marks_ink_on_a_text_box_and_fill_on_a_panel():
    """The same modifier means different things either side of the
    role: `white` text on Quote A's blue panel, a white FILL on Title
    and Text's field."""
    layouts = L.parse_layouts(LAYOUTS_FIXTURE)
    quote = L._bind_roles(layouts["slideLayout40"])
    on_panel = next(r for r in quote["regions"] if r["box"][0] == 510)
    assert on_panel["role"] == "on_panel_body"


def test_a_twin_renders_as_its_parent(spec):
    """`ARCHETYPES.md` says to prefer the parent -- but a brief that
    names the twin still has to render, and a twin is by definition
    geometrically identical."""
    built = L.build_archetypes()
    assert built["three_content_b"] == built["three_content"]


def test_timeline_has_geometry_despite_having_no_master(spec):
    built = L.build_archetypes()
    assert "timeline" in built
    assert any(r["role"] == "axis" for r in built["timeline"]["regions"])


# --------------------------------------------------------------------------
# installing
# --------------------------------------------------------------------------


class _Registry:
    def __init__(self, existing=None):
        self.ARCHETYPES = dict(existing or {})


def test_install_never_overwrites_a_hand_built_archetype(spec):
    """Hand-built archetypes were tuned against a real rendering; these
    are derived from a coarser description and must lose."""
    hand = {"regions": [{"role": "title", "box": [1, 2, 3, 4], "content": "title"}]}
    module = _Registry({"three_content": hand})
    added = L.install(module)
    assert "three_content" not in added
    assert module.ARCHETYPES["three_content"] is hand


def test_install_defers_to_an_alias_too(spec):
    """`COVER_A_CUT4` aliases `cover_cut`; implementing either means the
    archetype is covered and must not be regenerated under the other
    name."""
    module = _Registry({"cover_cut": {"regions": []}})
    assert "cover_a_cut4" not in L.install(module)


def test_install_is_idempotent(spec):
    module = _Registry()
    first = L.install(module)
    assert first and L.install(module) == []


def test_coverage_counts_aliases_as_covered(spec):
    module = _Registry({"cover_cut": {"regions": []}})
    assert L.coverage(module)["A"] == (1, 3)


# --------------------------------------------------------------------------
# against the shipped spec
# --------------------------------------------------------------------------


def test_the_shipped_spec_parses_completely():
    """The failure mode this guards is silence: a spec revision that
    stops matching the regex generates nothing and reports nothing.
    Every geometry line in the file must turn into a box."""
    text = (L.spec_dir() / "LAYOUTS.md").read_text()
    # `- **Grade A** — most used...` is prose in the preamble; a
    # geometry line is the one carrying a `w × h`.
    geometry_lines = [ln for ln in text.splitlines() if ln.startswith("- **") and " × " in ln]
    parsed = sum(len(v.boxes) for v in L.parse_layouts(text).values())
    assert parsed == len(geometry_lines) > 300


def test_every_canonical_archetype_is_renderable():
    """The whole point of the module. If a spec update adds archetypes
    faster than this can bind them, this is where it shows up."""
    import sys

    from deckguard.legacy import gallery
    from deckguard.registry import _ensure_skill_on_path

    _ensure_skill_on_path()
    archetypes = __import__("archetypes")
    try:
        gallery.install(archetypes)
    except Exception:  # noqa: BLE001 -- the gallery is optional
        pass
    L.install(archetypes)

    _, meta = L.load_spec()
    missing = sorted(
        a.name for a in meta.values()
        if a.engine_key not in archetypes.ARCHETYPES
        and a.engine_key not in L._UNREACHABLE
        and not any(x in archetypes.ARCHETYPES for x in a.aliases)
    )
    assert missing == [], f"no geometry for: {missing}"


def test_the_excluded_archetypes_really_are_unreachable():
    """`_UNREACHABLE` is the one place the library gets smaller, so it
    has to keep earning it. An archetype the meter lists, or a set
    names, can be chosen -- excluding it would silently delete a layout
    someone can ask for, which is the opposite of the intended trade."""
    from deckguard import brandmode as bm
    from deckguard import meter

    declared = set(meter.spec().get("archetypes") or {})
    in_sets = {s["archetype"] for audience in ("external", "internal")
               for s in bm.slides_in(audience)}
    for name in L._UNREACHABLE:
        assert name not in declared, f"{name} is in meter.json -- a plan can pick it"
        assert name not in in_sets, f"{name} is in slide-sets.json -- pickable by hand"


def test_the_pictogram_set_is_rasterised():
    """A .pptx cannot embed an SVG directly, so the icons have to exist
    as raster beside their vector source."""
    marks = L.pictograms()
    assert len(marks) >= 3
    assert all(m.endswith(".png") for m in marks)


def test_the_two_most_used_layouts_carry_an_icon_grid():
    """Measured, not assumed: across two on-brand KONE decks
    `Text and picture A` is used 18 times and `Text and picture G` 8 --
    more than every other layout combined -- and both carry a grid of
    icon-plus-short-text cells that the placeholder map has no form for.
    Slides using them average 7-10 text blocks against 3 bound regions.
    """
    for key in ("text_picture_a", "text_picture_g"):
        refinement = L._REFINEMENTS[key]
        group = refinement["groups"][0]
        assert any(r.get("role") == "icon" for r in group["regions"]), key
        assert len(group["origins"]) >= 4, key


def test_an_archetype_serves_both_its_plain_and_grid_forms():
    """One spec, two shapes. The engine skips a region whose content key
    is absent, so supplying `body` gives the paragraph version and
    supplying `items` gives the grid -- which is how the real decks use
    this layout, sometimes one way and sometimes the other."""
    regions = L._REFINEMENTS["text_picture_a"]["regions"]
    assert any(r.get("content") == "body" for r in regions)
    assert L._REFINEMENTS["text_picture_a"]["groups"][0]["content"] == "items"


def test_a_photo_banner_declares_protection_for_its_reversed_type():
    """White type on a pale photo is unreadable, and the brand specifies
    a gradient for exactly this. Declared only when there IS a photo --
    a scrim over white is just a grey band."""
    scrims = L._REFINEMENTS["text_picture_g"]["scrims"]
    assert scrims and scrims[0]["content"] == "image"
    assert scrims[0]["box"] == [0, 0, 1280, 440]


def test_the_scrim_is_a_gradient_that_spares_the_middle_of_the_picture():
    """A flat tint greys out the subject, which is why the spec calls
    for a gradient: dark at the edges where the type sits, clear through
    the middle where the photograph does its work."""
    import sys

    from pptx import Presentation
    from pptx.util import Emu

    sys.path.insert(0, "/root/.claude/skills/kone-deck-generator")
    try:
        import kone_engine as engine
    except Exception:  # pragma: no cover
        pytest.skip("archetype engine not available")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = L._draw_scrim(slide, engine, [0, 0, 1280, 440])

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    grad = shape._element.spPr.find(f"{{{A}}}gradFill")
    assert grad is not None, "a flat fill would grey out the photograph"
    stops = [(int(gs.get("pos")), int(a.get("val")))
             for gs in grad.iter(f"{{{A}}}gs")
             for a in gs.iter(f"{{{A}}}alpha")]
    assert stops == sorted(stops), "gradient stops must run in order"
    by_pos = dict(stops)
    assert by_pos[0] > 0 and by_pos[100000] > 0, "the edges carry the type"
    middle = [alpha for pos, alpha in stops if 30000 <= pos <= 70000]
    assert middle and max(middle) == 0, "the subject of the photograph is spared"


def test_icons_are_counted_against_content_not_geometry():
    """An archetype serving both a plain and a grid form declares the
    grid's origins either way. Counting those drew a full set of icons
    onto the plain form -- four of them, two straddling the body
    paragraph, on a slide that asked for none."""
    spec = {
        "regions": [{"role": "body", "box": [0, 0, 10, 10], "content": "body"}],
        "groups": [{
            "content": "items",
            "origins": [[0, 0], [1, 0], [2, 0], [3, 0]],
            "regions": [{"role": "icon", "box": [0, 0, 60, 60]}],
        }],
    }
    assert L._icon_slots(spec, {"body": "just a paragraph"}) == 0
    assert L._icon_slots(spec, {"items": [{}, {}]}) == 2
    assert L._icon_slots(spec, {"items": [{}] * 9}) == 4   # never past the origins
    assert L._icon_slots(spec) == 4                        # geometry, when asked


def test_the_six_up_row_stays_inside_the_right_margin():
    """The grid puts the last column at x=1065 and the page margin is
    45, so a cell wider than 170 runs off the edge -- it was 187."""
    group = L._REFINEMENTS["text_picture_g"]["groups"][0]
    last_x = max(x for x, _ in group["origins"])
    widest = max(r["box"][2] for r in group["regions"])
    assert last_x + widest <= 1280 - 45


def test_a_theme_referenced_background_resolves_to_a_colour():
    """Four KONE layouts set their background by reference --
    `<p:bgRef><a:schemeClr val="bg2"/>` -- which names a theme slot, not
    a colour. Reading only a literal srgbClr reported them as having no
    background, so archetypes drew onto them blind."""
    import posixpath

    from pptx import Presentation

    from deckguard.logo import _page_background_hex

    master = (L._VENDORED_SKILL_DIR / "uploads" / "master_ppt-1784774200983.pptx")
    if not master.is_file():
        pytest.skip("master template not available")
    prs = Presentation(str(master))
    by_key = {posixpath.basename(l.part.partname).replace(".xml", ""): l
              for l in prs.slide_layouts}

    divider = by_key["slideLayout10"]
    assert _page_background_hex(divider.element) is None      # no literal colour
    assert _page_background_hex(divider.element, divider) == "F3EEEA"   # resolved

    # a layout stating a literal colour still reads directly
    assert _page_background_hex(by_key["slideLayout39"].element) == "1450F5"


def test_the_divider_puts_the_number_left_and_the_title_right():
    """Corrected against a real deck that uses this layout five times.
    The gallery port and the HTML reference both had it the other way;
    the master's own boxes agree with the real slides."""
    regions = {r["content"]: r["box"] for r in L._REFINEMENTS["divider_numbering"]["regions"]}
    assert regions["number"][0] == 45
    assert regions["title"][0] == 453
    assert regions["number"][0] < regions["title"][0]


def test_a_named_colour_floods_the_divider_and_carries_its_own_ink():
    """Every divider in sand made a deck monotonous, and the brand puts
    the secondary palette on exactly this slide. The ink follows the
    brand's own rule -- white out of blue, black out of a secondary."""
    assert L._field_for({"field": True}, {"colour": "blue"}) == ("1450F5", "FFFFFF")
    assert L._field_for({"field": True}, {"colour": "light-blue"}) == ("D2F5FF", "141414")
    assert L._field_for({"field": True}, {"color": "pink"}) == ("FFCDD7", "141414")
    # unnamed falls back to KONE Blue rather than to nothing
    assert L._field_for({"field": True}, {}) == L.BRAND_FIELDS[L.DEFAULT_FIELD]
    # an archetype that does not declare a field never takes one
    assert L._field_for({}, {"colour": "pink"}) is None
    for fill, ink in L.BRAND_FIELDS.values():
        assert ink in ("FFFFFF", "141414") and len(fill) == 6


def test_the_colour_field_replaces_the_engines_background_rather_than_hiding_under_it():
    """`render_archetype` opens by flooding the slide with the
    archetype's default fill. A field inserted at the bottom of the
    shape tree is covered by that sand, and only the ink survives --
    which is how a blue divider rendered as white type on sand."""
    import sys

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    sys.path.insert(0, "/root/.claude/skills/kone-deck-generator")
    try:
        import kone_engine as engine
    except Exception:  # pragma: no cover
        pytest.skip("archetype engine not available")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    engine._rect(slide, [0, 0, 1280, 720], engine._hex("F3EEEA"))   # the engine's own
    L._paint_field(slide, engine, "1450F5")

    fills = [s for s in slide.shapes if s.name == "Colour field"]
    assert len(fills) == 1
    assert str(fills[0].fill.fore_color.rgb) == "1450F5"
    # and nothing full-slide is left painted over it
    assert len(L._full_slide_fills(slide, engine)) == 1

    # a shape that merely happens to be large is not mistaken for one
    slide.shapes.add_shape(MSO_SHAPE.OVAL, 0, 0, Emu(12192000), Emu(6858000))
    assert len(L._full_slide_fills(slide, engine)) == 1


def test_the_scrim_darkens_where_the_white_type_actually_sits():
    """A fixed dark-at-the-edges ramp guessed wrong on TEXT_PICTURE_G,
    whose headline starts 60% of the way down its banner -- exactly
    where a ramp clearing at the midpoint has recovered almost nothing.
    The bands come from the archetype's own reversed-out boxes."""
    spec = {"regions": [
        {"box": [45, 39, 917, 36], "content": "eyebrow",
         "dg": {"kind": "text", "color": "FFFFFF"}},
        {"box": [45, 262, 697, 125], "content": "title",
         "dg": {"kind": "text", "color": "FFFFFF"}},
        {"box": [45, 470, 1190, 40], "content": "body",     # below the banner
         "dg": {"kind": "text", "color": "141414"}},
    ]}
    content = {"eyebrow": "WHAT TENANTS EXPECT", "title": "Six expectations", "body": "x"}
    bands = L._reversed_bands(spec, content, [0, 0, 1280, 440])
    assert bands == [(39, 75), (262, 387)]

    stops = dict(L._scrim_ramp([0, 0, 1280, 440], bands, 78))
    # the headline band is fully protected, the middle of the picture is not
    assert stops[65000] == 78 and stops[80000] == 78
    # and clears completely between the eyebrow and the headline
    assert max(stops[pos] for pos in (30000, 35000, 40000, 45000)) == 0

    # an empty block is not protected: nothing is written there
    assert L._reversed_bands(spec, {"title": "Six expectations"}, [0, 0, 1280, 440]) \
        == [(262, 387)]


def test_the_cut_cover_banner_crops_its_photograph_instead_of_stretching_it():
    """`add_picture` given both a width and a height scales each axis
    independently. Every photo in the set is between 1.3 and 1.9 wide and
    the banner is 2.25, so the cover came out squashed by half again."""
    import sys

    from pptx import Presentation
    from pptx.util import Emu

    sys.path.insert(0, "/root/.claude/skills/kone-deck-generator")
    try:
        import kone_engine as engine
    except Exception:  # pragma: no cover
        pytest.skip("archetype engine not available")

    from deckguard import photos

    tall = next((p for p in photos.load_photos().values()
                 if photos.crop_severity(p, 950 / 422) > 0.1), None)
    if tall is None:  # pragma: no cover
        pytest.skip("no photo narrower than the banner")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    picture = photos.place_cover(slide, engine, (330, 0, 950, 422), tall.path)

    assert (picture.width, picture.height) == (engine.X(950), engine.X(422))
    # narrower than the slot, so a centre BAND is kept and nothing is
    # taken off the sides
    assert picture.crop_top > 0 and picture.crop_top == pytest.approx(picture.crop_bottom)
    assert picture.crop_left == 0 and picture.crop_right == 0

    kept = (1 - picture.crop_top - picture.crop_bottom)
    from PIL import Image
    with Image.open(str(tall.path)) as image:
        width, height = image.size
    assert (width / (height * kept)) == pytest.approx(950 / 422, rel=1e-3)


def test_the_slot_shape_breaks_ties_between_equally_matching_photos():
    """The cut-cover banner is 2.25:1 and the set holds photographs at
    exactly that ratio. Handing it a 3:2 portrait threw away a third of
    the frame while one that fitted sat unused."""
    from deckguard import photos

    wide = photos.choose("elevator women people", slot_aspect=950 / 422)
    assert wide is not None
    assert photos.crop_severity(wide, 950 / 422) < 0.15

    # the subject still outranks the shape: a query with no aspect given
    # is unchanged by the new argument
    assert photos.choose("skyline city") == photos.choose("skyline city", slot_aspect=None)


def test_a_picture_carrying_white_type_gets_a_scrim_nobody_declared():
    """`LAYOUTS.md` marks the type white and leaves the protection to
    the designer, so the derived archetypes shipped without any -- the
    full-bleed cover put a white headline onto a sunlit treeline."""
    regions = [
        {"role": "picture", "box": [0, 0, 1280, 720], "content": "photo"},
        {"role": "title_light", "box": [45, 136, 578, 448], "content": "title"},
    ]
    assert L._implied_scrims(regions) == [
        {"box": [0, 0, 1280, 720], "content": "photo"}]

    # gallery ports carry the ink in the role name instead
    assert L._is_light_role("gal_i64_FFFFFF") and L._is_light_role("title_light")
    assert not L._is_light_role("gal_i19_141414")

    # black type over a picture needs no protection, and neither does a
    # white block that misses the picture entirely
    assert L._implied_scrims([
        regions[0], {"role": "body", "box": [45, 136, 578, 40], "content": "body"}]) == []
    assert L._implied_scrims([
        {"role": "picture", "box": [0, 0, 640, 300], "content": "photo"},
        {"role": "title_light", "box": [700, 400, 400, 60], "content": "title"}]) == []


def test_the_scrim_protects_the_type_not_the_frame_it_was_given():
    """The full-bleed cover's title frame is 448px tall for three lines
    of type. Protecting the frame darkened the picture end to end and
    the cover came back a uniform grey."""
    region = {"role": "gal_i64_FFFFFF", "box": [45, 136, 578, 448], "content": "title"}
    height = L._typeset_height(region, "The lift is the last thing you decarbonise")
    assert 150 < height < 300, height
    assert height < region["box"][3]

    # a long enough string still fills the frame and never overruns it
    assert L._typeset_height(region, "word " * 400) == region["box"][3]


def test_the_scrim_sits_above_the_photograph_and_below_the_type():
    """A cover carries its logo and tagline as pictures too, and they
    are added after the type. Seating the scrim above the LAST picture
    put it over the headline and dimmed the words it exists to make
    readable."""
    import sys

    from pptx import Presentation
    from pptx.util import Emu

    sys.path.insert(0, "/root/.claude/skills/kone-deck-generator")
    try:
        import kone_engine as engine
    except Exception:  # pragma: no cover
        pytest.skip("archetype engine not available")

    from deckguard import photos

    any_photo = next(iter(photos.load_photos().values()), None)
    if any_photo is None:  # pragma: no cover
        pytest.skip("no photographs installed")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(any_photo.path), 0, 0,
                             Emu(1280 * 9525), Emu(720 * 9525))          # the photograph
    slide.shapes.add_textbox(0, Emu(136 * 9525), Emu(578 * 9525), Emu(200 * 9525))
    slide.shapes.add_picture(str(any_photo.path), Emu(1102 * 9525), Emu(633 * 9525),
                             Emu(133 * 9525), Emu(45 * 9525))            # the tagline

    L._draw_scrim(slide, engine, [0, 0, 1280, 720], protect=[(136, 300)])
    order = [s.shape_type for s in slide.shapes]
    names = [s.name for s in slide.shapes]
    assert names.index("Photo protection") == 1, order
    assert names.index("Photo protection") < names.index("TextBox 2")


def test_the_engines_caption_grey_is_corrected_to_black_at_install():
    """The brand allows three inks for type -- black, white and, for
    KONE Information only, blue. The engine sets captions, muted body
    and quote attributions in a #727272 that is in no palette, and it
    shows on any deck carrying a quote or a captioned statistic."""
    import sys

    sys.path.insert(0, "/root/.claude/skills/kone-deck-generator")
    try:
        import archetypes as module
        import kone_engine as engine
    except Exception:  # pragma: no cover
        pytest.skip("archetype engine not available")

    # restore the engine's shipped greys first: another test in the
    # session may already have installed and corrected them
    greyed = {}
    for role in ("caption", "body_muted", "attribution"):
        style = engine.ROLE_STYLE[role]
        greyed[role] = style
        engine.ROLE_STYLE[role] = tuple(
            engine.GREY if i == 2 else v for i, v in enumerate(style))

    corrected = L._correct_grey_ink(module)
    assert {"caption", "body_muted", "attribution"} <= set(corrected)
    for role in corrected:
        assert engine.ROLE_STYLE[role][2] == engine.BLACK
    # and nothing anywhere is left painting type in the caption grey
    assert not [r for r, s in engine.ROLE_STYLE.items()
                if len(s) > 2 and s[2] == engine.GREY]
    # idempotent: a second pass finds nothing left to correct
    assert L._correct_grey_ink(module) == []
    # everything else about the role is untouched -- font, size, caps
    for role, before in greyed.items():
        after = engine.ROLE_STYLE[role]
        assert [v for i, v in enumerate(after) if i != 2] \
            == [v for i, v in enumerate(before) if i != 2]


def _skill_modules():
    import sys

    sys.path.insert(0, "/root/.claude/skills/kone-deck-generator")
    try:
        import archetypes
        import kone_engine  # noqa: F401
    except Exception:  # pragma: no cover
        pytest.skip("archetype engine not available")
    from deckguard.legacy import gallery

    # NOTE: this installs the LEGACY gallery, which `registry` does not
    # use -- runtime merges its own `_install_gallery`. `ARCHETYPES` is a
    # module global shared by every test in the session, so from here on
    # other tests measure a registry the app never builds. Anything
    # asserting on the shipped library has to be robust to that.
    L.install(archetypes)
    gallery.install(archetypes)
    return archetypes


def test_a_deck_built_from_a_brief_puts_each_archetype_on_its_own_layout(tmp_path):
    """The skill's own `build_deck` puts every archetype on BLANK and
    calls `archetypes.render` directly, which bypasses this module
    entirely. The first deck a user built through the web tool came back
    with the engine's rasterised PNGs and blue placeholder chips instead
    of KONE pictograms, and no scrim on the cover."""
    archetypes = _skill_modules()
    out = tmp_path / "brief.pptx"

    L.build_deck({"title": "Service review", "slides": [
        {"archetype": "icon_columns_5", "title": "What we need", "items": [
            {"icon": "clock", "text": "Confirm the forecast"},
            {"icon": "people", "text": "Name an owner"},
            {"icon": "wrench", "text": "Apply the pricing"},
        ]},
    ]}, str(out), archetypes)

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides._sldIdLst) == 3      # master cover, body, master outro

    body = prs.slides[1]
    kinds = [str(s.shape_type).split(" ")[0] for s in body.shapes]
    # native editable pictograms, not rasters and not the engine's chips
    assert kinds.count("FREEFORM") >= 3
    assert "PICTURE" not in kinds
    assert not [s for s in body.shapes if "Rounded Rectangle" in s.name]


def test_the_deck_opens_on_the_four_pane_cut_not_the_masters_own_cover(tmp_path):
    """The master retains a full-bleed photograph with a white title on
    it, and on a sunlit frame the title all but vanished -- decks were
    opening on a single legible letter. The reference deck opens on the
    cut, so that is the default now, and its title sits on white where
    nothing can swallow it.

    This replaces a test that asserted the retained cover got a scrim.
    That cover is no longer in the file, so the scrim it needed is moot.
    """
    from pptx import Presentation

    from deckguard import assemble

    # Through `assemble.build`, which is the path the tool takes: the
    # legacy skill renderer draws its own cut and names the masks
    # differently, and that is not what ships.
    out = tmp_path / "cover.pptx"
    assemble.build({"title": "Service business review", "date": "1 March 2026",
                    "slides": []}, str(out))

    slides = list(Presentation(str(out)).slides)
    cover = slides[0]
    names = [s.name for s in cover.shapes]
    # One picture, masked into panes. The masks are named by whichever
    # renderer drew them -- the gallery's chrome path calls them
    # `Rectangle N` -- so the check is that the cut exists, not what it
    # is called: one banner and at least four rectangles over it.
    banners = [sh for sh in cover.shapes
               if sh.shape_type and sh.shape_type.name == "PICTURE"
               and sh.width > 3000000]
    assert len(banners) == 1, names
    covers = [sh for sh in cover.shapes
              if sh.shape_type and sh.shape_type.name == "AUTO_SHAPE"]
    assert len(covers) >= 4, names
    # and the KONE marks land whichever path drew the cut
    assert "Logo" in names and "Tagline" in names, names
    # and the master's own cover is gone rather than sitting behind it
    assert len(slides) == 2, "cut cover plus the retained Thank you"


def test_a_chart_slot_nobody_filled_is_dropped_not_invented(tmp_path):
    """The engine keeps sample artwork and stamps it onto every slide of
    certain archetypes whether or not anyone asked. `segment_breakdown`
    gets a donut reading 53% against satisfaction bands -- invented data
    in the company's own chart styling, and it went out in a business
    review."""
    archetypes = _skill_modules()
    assert "chart" in archetypes.FIGURES["segment_breakdown"], "guard is about this map"

    region = {"content": "chart", "role": "figure"}
    assert L._is_unsupplied_figure(region, {}, "segment_breakdown", archetypes)
    # supplied by the author, it is theirs and must be drawn
    assert not L._is_unsupplied_figure(
        region, {"chart": "/tmp/mine.png"}, "segment_breakdown", archetypes)
    # a slot that is not a figure is never touched by this rule
    assert not L._is_unsupplied_figure(
        {"content": "title"}, {}, "segment_breakdown", archetypes)

    out = tmp_path / "nochart.pptx"
    L.build_deck({"title": "Review", "slides": [
        {"archetype": "segment_breakdown", "title": "Where we stand",
         "highlight_value": "88%", "highlight_caption": "first-time fix"},
    ]}, str(out), archetypes)

    from pptx import Presentation

    body = Presentation(str(out)).slides[1]
    assert not [s for s in body.shapes if str(s.shape_type).startswith("PICTURE")]
    assert any("88%" in s.text_frame.text
               for s in body.shapes if getattr(s, "has_text_frame", False))


def test_build_deck_takes_the_registry_from_the_loader_not_a_bare_import(tmp_path):
    """`gallery.install` OVERWRITES what it finds; `layouts.install`
    defers to it except for a short override list. So the gallery has to
    be installed first and this module second, and only
    `skill_bridge._load_archetypes` guarantees that order.

    Getting it backwards is silent and expensive: the refined agenda
    reverted to a port with no bullets, and the divider to one with no
    number and no colour field. Both looked like rendering bugs."""
    _skill_modules()
    from deckguard.registry import _load_archetypes

    registry = _load_archetypes().ARCHETYPES

    agenda = registry["agenda_a_table"]
    assert [r.get("content") for r in agenda["regions"]] == [
        "photo", "eyebrow", "title", "items"]
    assert any(r.get("role") == "dg_bullets" for r in agenda["regions"])

    divider = registry["divider_numbering"]
    assert divider.get("field") is True
    assert {r.get("content") for r in divider["regions"]} == {"number", "eyebrow", "title"}
    # the number sits left of the title, as five real KONE slides do
    boxes = {r["content"]: r["box"] for r in divider["regions"]}
    assert boxes["number"][0] < boxes["title"][0]

    # and a deck built without naming a module picks up exactly that
    out = tmp_path / "ordered.pptx"
    L.build_deck({"title": "T", "slides": [
        {"archetype": "divider_numbering", "number": "2", "eyebrow": "Section 02",
         "title": "What's working", "colour": "light-blue"},
    ]}, str(out))

    from pptx import Presentation

    body = Presentation(str(out)).slides[1]
    field = next(s for s in body.shapes if s.name == "Colour field")
    assert str(field.fill.fore_color.rgb) == "D2F5FF"
    assert any("2" == s.text_frame.text.strip()
               for s in body.shapes if getattr(s, "has_text_frame", False))


def test_the_planner_is_told_the_keys_the_renderer_actually_reads():
    """`catalog.json` describes 22 of 80 archetypes and `SAMPLES` 41, so
    most were named and nothing more. The keys come off the live
    registry instead, and cannot drift from it."""
    _skill_modules()
    from deckguard.registry import _derived_content_keys, _load_archetypes

    registry = _load_archetypes().ARCHETYPES
    described = [n for n in registry if _derived_content_keys(n)]
    assert len(described) >= len(registry) - 2, "nearly every archetype must describe itself"

    agenda = dict(k.split(" (")[0:1] + [k] for k in
                  L.content_keys(registry["agenda_a_table"]))
    assert set(agenda) == {"photo", "eyebrow", "title", "items"}
    assert "list of strings" in agenda["items"]
    assert "do not supply" in agenda["photo"]

    # a repeating group states its capacity and its per-item fields,
    # including the icon, which carries no content key of its own
    (items,) = [k for k in L.content_keys(registry["text_picture_g"])
                if k.startswith("items ")]
    assert "up to 6" in items and "icon" in items and "text" in items


def test_a_worked_example_that_no_longer_matches_is_not_shown():
    """The model follows a concrete example over an abstract slot list.
    `agenda_a_table`'s advertised `text1..text4` long after the renderer
    was rebuilt to read `items` -- so a planner emitted four keys nothing
    reads and the agenda came back as a title on an empty half."""
    _skill_modules()
    from deckguard.legacy.skill_bridge import _kone_archetype_guide, _load_archetypes, _sample_agrees

    archetypes = _load_archetypes()
    stale = {"title": "x", "text1": "a", "text2": "b"}
    assert not _sample_agrees("agenda_a_table", stale)
    assert _sample_agrees("agenda_a_table", {"title": "x", "items": ["a"]})
    # an archetype we cannot describe keeps its example rather than losing both
    assert _sample_agrees("nothing_we_know_about", {"anything": 1})

    guide = _kone_archetype_guide()
    section = guide[guide.index("### agenda_a_table"):]
    section = section[:section.index("\n\n")]
    assert "Content keys (authoritative" in section
    assert "text1" not in section
    assert archetypes.SAMPLES["agenda_a_table"].get("text1"), "guard is about this sample"


def test_content_the_archetype_cannot_hold_is_reported_rather_than_dropped(tmp_path):
    """Silent loss is the failure this catches. The deck came back with
    less in it than the brief had, and nothing said so -- not the build,
    not the review page, not the audit."""
    archetypes = _skill_modules()
    agenda = archetypes.ARCHETYPES["agenda_a_table"]

    assert L.unread_keys(agenda, {"title": "x", "text1": "a", "text2": "b"}) == ["text1", "text2"]
    assert L.unread_keys(agenda, {"title": "x", "items": ["a"]}) == []
    # an empty field was never content, and `colour` is read by the
    # renderer rather than by the archetype's regions
    assert L.unread_keys(agenda, {"title": "x", "text1": "", "colour": "blue"}) == []

    report: dict = {}
    L.build_deck({"title": "T", "slides": [
        {"archetype": "agenda_a_table", "title": "Agenda",
         "text1": "one", "text2": "two"},
    ]}, str(tmp_path / "d.pptx"), archetypes, report=report)
    # position 2, not 1: every deck now opens on a generated four-pane
    # cut cover, so the first slide the caller asked for is the second
    # slide of the file.
    assert report["dropped"] == {2: ("agenda_a_table", ["text1", "text2"])}


def test_the_photo_divider_gets_a_scrim_like_every_other_full_bleed_layout():
    """`image_section_divider` calls its full-bleed picture `image_band`
    rather than `picture`, and matching on one role name meant the ONE
    archetype whose entire design is white type over a photograph was
    the one that never got protected. A Q2 review came back with its
    section titles lost in a sunlit stairwell."""
    archetypes = _skill_modules()
    spec = archetypes.ARCHETYPES["image_section_divider"]

    roles = [r.get("role") for r in spec["regions"]]
    assert "image_band" in roles, "guard is about this role name"
    assert L._implied_scrims(spec["regions"]) == [{"box": [0, 0, 1280, 720], "content": "image"}]

    # both names count, and a role that is not a picture still does not
    assert L._PICTURE_REGION_ROLES == frozenset({"picture", "image_band"})
    assert L._implied_scrims([
        {"role": "body", "box": [0, 0, 1280, 720], "content": "image"},
        {"role": "title_light", "box": [60, 545, 1050, 130], "content": "title"}]) == []


def test_a_repeating_group_tells_the_planner_it_can_name_the_icons():
    """The icon region carries no content key of its own, so the guide
    never mentioned it and every deck came back with the default
    cloud/people/clock cycle regardless of what the text said."""
    _skill_modules()
    from deckguard.registry import _derived_content_keys

    (items,) = [k for k in _derived_content_keys("icon_columns_5") if k.startswith("items ")]
    assert "icon" in items and "text" in items and "up to 5" in items


def test_an_item_row_is_pulled_up_when_nothing_fills_the_band_above_it():
    """Four archetypes place their row in the bottom third with only a
    title above -- geometry from real KONE slides where a paragraph of
    body copy filled the gap. These have no such paragraph, so a Q2
    review came back with 248px of blank sand between "Plan of action"
    and the six things it listed."""
    spec = {
        "regions": [{"role": "title", "box": [45, 91, 1189, 104], "content": "title"}],
        "groups": [{"content": "items", "origins": [[45, 443], [249, 443]],
                    "regions": [{"role": "icon", "box": [0, 0, 40, 40]},
                                {"role": "body", "box": [0, 104, 170, 100], "content": "label"}]}],
    }
    assert L.lift_low_rows(spec) == 1
    assert [o[1] for o in spec["groups"][0]["origins"]] == [264, 264]
    assert [o[0] for o in spec["groups"][0]["origins"]] == [45, 249], "columns must not move"

    # idempotent: the gap is now ordinary breathing room
    assert L.lift_low_rows(spec) == 0


def test_a_row_is_left_alone_when_the_band_is_spoken_for():
    """The guard that matters: two of the four candidates have a SECOND
    row in the band, and lifting either would have collided them."""
    occupied = {
        "regions": [{"role": "title", "box": [45, 91, 900, 60], "content": "title"}],
        "groups": [
            {"content": "pictures", "origins": [[45, 181]],
             "regions": [{"role": "picture", "box": [0, 0, 280, 216]}]},
            {"content": "items", "origins": [[45, 420]],
             "regions": [{"role": "body", "box": [0, 0, 280, 100], "content": "text"}]},
        ],
    }
    before = [[list(o) for o in g["origins"]] for g in occupied["groups"]]
    assert L.lift_low_rows(occupied) == 0
    assert [[list(o) for o in g["origins"]] for g in occupied["groups"]] == before

    # a modest gap is breathing room, not a hole
    modest = {
        "regions": [{"role": "title", "box": [45, 91, 900, 104], "content": "title"}],
        "groups": [{"content": "items", "origins": [[45, 300]],
                    "regions": [{"role": "body", "box": [0, 0, 170, 100], "content": "t"}]}],
    }
    assert L.lift_low_rows(modest) == 0

    # and a full-bleed photograph is a background, not a block to clear
    banner = {
        "regions": [{"role": "picture", "box": [0, 0, 1280, 720], "content": "image"},
                    {"role": "title", "box": [45, 91, 900, 104], "content": "title"}],
        "groups": [{"content": "items", "origins": [[45, 476]],
                    "regions": [{"role": "body", "box": [0, 0, 170, 100], "content": "t"}]}],
    }
    assert L.lift_low_rows(banner) == 1
    assert banner["groups"][0]["origins"][0][1] == 264


def test_installing_closes_the_dead_bands_in_the_shipped_registry():
    archetypes = _skill_modules()
    row = archetypes.ARCHETYPES["numbered_icon_row_6"]
    top = min(o[1] for o in row["groups"][0]["origins"])
    title = next(r for r in row["regions"] if r.get("content") == "title")
    assert top - (title["box"][1] + title["box"][3]) == pytest.approx(69, abs=1)

    # the two whose band is occupied keep their geometry
    for name in ("four_point_value", "quarterly_plan_4col"):
        spec = archetypes.ARCHETYPES[name]
        assert len(spec.get("groups") or []) >= 2 or L.lift_low_rows(dict(spec)) == 0


# --------------------------------------------------------------------------
# milestone slide (kone-milestone-slide)
# --------------------------------------------------------------------------


def test_the_milestone_slide_matches_its_published_geometry():
    """One slide from one announcement email -- the recognition slide
    the master does not have. Transcribed from kone-milestone-slide's
    own spec rather than eyeballed, so it stays comparable with the
    reference it came from."""
    spec = L._EXTRAS["milestone_slide"]
    boxes = {r["content"]: r["box"] for r in spec["regions"]}

    assert boxes["eyebrow"][:2] == [45, 47]
    assert boxes["title"][:3] == [45, 82, 790]
    assert boxes["lede"][:3] == [45, 186, 700]
    assert boxes["classification"][0] == 45
    # the sand band is the only secondary colour on the slide
    assert spec["panels"] == [{"box": [0, 276, 1280, 196], "fill": "F3EEEA"}]
    assert {p["fill"] for p in spec["panels"]} == {"F3EEEA"}

    # five stat cells, 206 wide on a 246 pitch across 1190 with a 40 gap
    stats = next(g for g in spec["groups"] if g["content"] == "stats")
    assert [o[0] for o in stats["origins"]] == [45, 291, 537, 783, 1029]
    assert all(o[1] == 307 for o in stats["origins"])
    value = next(r for r in stats["regions"] if r["content"] == "value")
    assert value["dg"]["px"] == 62 and value["dg"]["font"] == "KONE Information"

    # Inter is never blue; KONE Information is the only blue text
    for region in spec["regions"]:
        style = region["dg"]
        if style.get("color") == "1450F5":
            assert style.get("font") == "KONE Information", region["content"]
    # and never grey, anywhere
    # `_text` leaves the colour unset where black is the default
    inks = {r["dg"].get("color") or "141414" for r in spec["regions"]}
    assert inks <= {"141414", "FFFFFF", "1450F5"}


def test_the_zero_stat_is_black_because_it_is_a_different_kind_of_claim():
    """The brand's own instruction: the number that is deliberately zero
    -- no disruption, no downtime -- is the strongest claim on the slide
    and must not sit in the same blue as the counts."""
    assert L._reads_as_zero("0") and L._reads_as_zero(" 0 ") and L._reads_as_zero("zero")
    assert not L._reads_as_zero("6") and not L._reads_as_zero("100+")
    assert not L._reads_as_zero("3+3")

    stats = next(g for g in L._EXTRAS["milestone_slide"]["groups"]
                 if g["content"] == "stats")
    value = next(r for r in stats["regions"] if r["content"] == "value")
    assert value["dg"]["zero_is_black"] is True


def test_the_milestone_slide_renders_its_own_worked_example(tmp_path):
    archetypes = _skill_modules()
    out = tmp_path / "milestone.pptx"
    L.build_deck({"title": "Milestone", "slides": [{
        "archetype": "milestone_slide",
        "eyebrow": "Marketing Hub · Request Management",
        "title": "From Monday.com to ServiceNow in six weeks",
        "lede": "Delivered with full data continuity.",
        "done": [{"text": "MVP pilot-tested and now live"},
                 {"text": "100% transition completed"}],
        "stats": [{"value": "6", "label": "Weeks end to end"},
                  {"value": "100+", "label": "Users migrated"},
                  {"value": "12", "label": "Frontlines"},
                  {"value": "3+3", "label": "Regions + global teams"},
                  {"value": "0", "label": "Business disruption"}],
        "scope_label": "The frontlines", "scope": "KSEA · KMTA · KANZ",
        "next_label": "What's next", "next": ["Hypercare ongoing"],
        "credits_label": "Thank you", "credits": "Arvind and the Hub specialists",
        "classification": "KONE Internal",
    }]}, str(out), archetypes)

    from pptx import Presentation

    slide = Presentation(str(out)).slides[1]
    text = " ".join(s.text_frame.text for s in slide.shapes
                    if getattr(s, "has_text_frame", False))
    # stat labels are set in caps by the brand, so match that form
    for fragment in ("ServiceNow", "100+", "3+3", "BUSINESS DISRUPTION",
                     "KSEA", "Hypercare", "Hub specialists"):
        assert fragment in text, fragment

    # the sand band, the three hairlines and two tick discs are all shapes
    assert len([s for s in slide.shapes if s.name == "Hairline"]) == 3
    assert len([s for s in slide.shapes if s.name == "Tick"]) == 2
    ticks = [s for s in slide.shapes if s.name == "Tick"]
    assert all(str(t.fill.fore_color.rgb) == "1450F5" for t in ticks)
    assert all("✓" in t.text_frame.text for t in ticks)

    # the zero renders black while the counts stay blue
    inks = {}
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() in ("6", "100+", "12", "3+3", "0") and run.font.color.type:
                    inks[run.text.strip()] = str(run.font.color.rgb)
    assert inks.get("0") == "141414", inks
    assert inks.get("6") == "1450F5" and inks.get("100+") == "1450F5", inks


def test_three_completion_ticks_clear_the_sand_band():
    """Three is the published ceiling, and the tick column has 186->276
    to do it in before the band starts. A roomier pitch fits two and
    puts the third INSIDE the band -- which is what shipped, and what a
    real deck came back with."""
    spec = L._EXTRAS["milestone_slide"]
    band_top = spec["panels"][0]["box"][1]
    ticks = next(g for g in spec["groups"] if g["content"] == "done")
    row = ticks["regions"][0]["box"][3]

    assert len(ticks["origins"]) == 3
    assert max(o[1] for o in ticks["origins"]) + row <= band_top
    # evenly pitched, so the column reads as one list
    ys = sorted(o[1] for o in ticks["origins"])
    assert len({b - a for a, b in zip(ys, ys[1:])}) == 1


def test_the_guide_teaches_the_rules_a_planner_has_to_act_on():
    """They live in the skill's SKILL.md, which the planner never sees.
    A deck came back with `6 weeks` as a stat value against an
    `End-to-end migration` label -- the unit said twice -- and a scope
    line restating counts already in the band."""
    from deckguard.legacy import skill_bridge

    entry = [b for b in skill_bridge._kone_archetype_guide().split("\n\n")
             if b.startswith("### milestone_slide")][0]
    rules = " ".join(l for l in entry.splitlines() if l.startswith("Rule:"))
    assert "BARE number" in rules
    assert "55 characters" in rules
    assert "one line" in rules and "done" in rules
    assert "saying nothing twice" in rules


# --------------------------------------------------------------------------
# the recognition deck (kone-recognition-deck)
# --------------------------------------------------------------------------


MILESTONE = {
    "eyebrow": "Marketing Hub", "footer": "Marketing Hub · April 2026",
    "title": "From Monday.com to ServiceNow in six weeks",
    "context": "The Hub moved its entire Request Management framework across.",
    "stats": [{"value": "6", "label": "Weeks"}, {"value": "0", "label": "Disruption"}],
    "scope_label": "The frontlines", "scope": "KSEA · KMTA · KANZ",
    "done": [{"text": "Pilot-tested and live"}, {"text": "100% transitioned"}],
    "groups": [{"heading": "Frontlines", "text": "12 frontlines."}],
    "quote": "A benchmark in easy to work with.",
    "next": ["Hypercare ongoing", "Power BI integration targeted for Q2"],
    "credit_names": ["Arvind", "Suresh Kumar", "Rupesh", "Golda"],
    "credits_note": "And the Hub specialists.",
}


def test_the_arc_does_not_emit_a_cover_or_a_closing_slide():
    """`build_deck` retains the master's Cover F and Thank you. The
    first build of this arc emitted its own too, so the deck opened on
    two covers and closed on three thank-yous."""
    slides = L.recognition_deck(MILESTONE)["slides"]
    used = [s["archetype"] for s in slides]
    assert "cover_a_cut4" not in used and "outro" not in used
    assert used[0] == "agenda_contents"
    assert used[-1] == "credits"


def test_a_section_divider_precedes_every_section_that_has_material():
    slides = L.recognition_deck(MILESTONE)["slides"]
    used = [s["archetype"] for s in slides]
    assert used.count("divider_numbering") == 4
    numbers = [s["number"] for s in slides if s["archetype"] == "divider_numbering"]
    assert numbers == ["01", "02", "03", "04"]


def test_sections_without_material_are_dropped_rather_than_padded():
    """The skill is explicit: nine slides that carry weight beat twelve
    with three filler ones."""
    thin = {"title": "A thing shipped", "stats": [{"value": "1", "label": "Thing"}]}
    used = [s["archetype"] for s in L.recognition_deck(thin)["slides"]]
    assert "three_content" not in used   # no groups given
    assert "quote_b" not in used         # no quote given
    assert "credits" not in used         # no names given
    assert used.count("divider_numbering") == 1
    assert "kone_numbers" in used


def test_a_period_in_the_sentence_becomes_the_timeline_period():
    items = L._as_timeline(["Hypercare ongoing",
                            "Power BI integration targeted for Q2",
                            "Demo at the Hub call in April 2026"])
    assert [i["period"] for i in items] == ["Now", "Q2", "April 2026"]


def test_the_numbers_slide_and_the_milestone_band_are_the_same_geometry():
    """A number that moved between the one-slide and the deck-length
    form would read as a different number."""
    band = L._EXTRAS["kone_numbers"]
    slide = L._EXTRAS["milestone_slide"]
    assert band["panels"] == slide["panels"]
    for spec in (band, slide):
        stats = next(g for g in spec["groups"] if g["content"] == "stats")
        assert [o[0] for o in stats["origins"]] == [45, 291, 537, 783, 1029]


def test_credit_rows_rule_only_where_there_are_names():
    """Eight names fill two rows of four; a third row of rules drawn
    statically would hang under them ruling nothing. The rule belongs to
    the cell, so it appears only where content does."""
    names = next(g for g in L._EXTRAS["credits"]["groups"] if g["content"] == "names")
    assert len(names["origins"]) == 12
    assert names["regions"][0]["dg"]["kind"] == "ruled"
    # three rows of four, and the last row clears the closing line
    assert sorted({o[1] for o in names["origins"]}) == [240, 321, 402]
    note = next(r for r in L._EXTRAS["credits"]["regions"] if r["content"] == "note")
    assert 402 + 81 <= note["box"][1]


def test_the_recognition_deck_renders(tmp_path):
    archetypes = _skill_modules()
    out = tmp_path / "recognition.pptx"
    L.build_deck(L.recognition_deck(MILESTONE), str(out), archetypes)

    from pptx import Presentation

    prs = Presentation(str(out))
    text = " ".join(
        sh.text_frame.text for slide in prs.slides for sh in slide.shapes
        if getattr(sh, "has_text_frame", False)
    )
    for fragment in ("What we'll cover", "Suresh Kumar", "Hub specialists",
                     "benchmark", "KSEA"):
        assert fragment in text, fragment


def test_a_chosen_shape_is_planner_guidance_and_an_unknown_one_changes_nothing():
    notes, target = L.shape_notes("recognition")
    assert target == 12
    assert "kone_numbers" in notes and "agenda_contents" in notes
    assert "do not emit" in notes.lower()

    notes, target = L.shape_notes("milestone")
    assert target == 1 and "milestone_slide" in notes

    assert L.shape_notes("auto") == (None, None)
    assert L.shape_notes("") == (None, None)
    assert L.shape_notes("something-else") == (None, None)


def test_the_arc_notes_are_built_from_the_arc_itself():
    """A hand-written copy would be a second place to update."""
    notes = L._arc_notes()
    for name, _why in L.RECOGNITION_ARC:
        assert name in notes


# --------------------------------------------------------------------------
# chrome and the cut cover
# --------------------------------------------------------------------------


def test_every_content_slide_carries_a_date_and_a_page_number(tmp_path):
    """BRAND_MODE section 3 puts chrome in the layout and tells the
    archetype to draw none. Every archetype duly drew none -- and so did
    the layout, so generated body slides came out with nothing at all
    below y=629 while the master's retained cover and Thank you kept
    theirs, which made it read as a quirk rather than as chrome missing."""
    from pptx import Presentation

    out = tmp_path / "chrome.pptx"
    L.build_deck({"title": "T", "date": "12 March 2026", "slides": [
        {"archetype": "cover_a_cut4", "title": "Cover"},
        {"archetype": "title_content", "title": "Findings", "bullets": ["One"]},
        {"archetype": "divider_numbering", "number": "01", "title": "Section"},
        {"archetype": "three_content", "title": "Three",
         "items": [{"heading": "A", "text": "a"}]},
    ]}, str(out), _skill_modules())

    prs = Presentation(str(out))
    px = prs.slide_width / 1280
    footers = []
    for slide in prs.slides:
        low = [(round(sh.left / px), sh.text_frame.text.strip())
               for sh in slide.shapes
               if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
               and sh.top / px > 640]
        footers.append(low)

    # deck order: master cover, cover, title_content, divider, three_content, master outro
    # Indices shifted by one: the master's retained cover is gone and the
    # spec's own `cover_a_cut4` is slide 0, so the first content slide is
    # slide 1 and carries page 02.
    assert any("12 MARCH 2026" in t for _x, t in footers[1]), footers[1]
    assert any(t == "02" for _x, t in footers[1]), "page number missing"
    assert any(t == "04" for _x, t in footers[3]), footers[3]
    assert footers[0] == [], "a cover takes no footer"
    assert footers[2] == [], "a divider takes no footer"


def test_the_page_number_sits_right_and_the_date_left():
    from deckguard import brandmode as bm

    assert bm.FOOTER_DATE_X == 45
    assert bm.FOOTER_PAGE_X == 1167
    assert bm.FOOTER_Y == 658 > bm.FLOOR


def test_the_cut_cover_masks_one_photograph_rather_than_slicing_it():
    """One picture behind background-coloured masks, so dropping in a
    new image reproduces the chopped effect and Change Picture works."""
    for name, cut in L._CUT_COVERS.items():
        masks = L._cut_masks(cut["band"], cut["panes"])
        assert len(masks) >= len(cut["panes"]), name
        bx, by, bw, bh = cut["band"]
        for x, y, w, h in masks:
            assert w > 0 and h > 0, (name, x, y, w, h)
            assert bx <= x and x + w <= bx + bw + 1, (name, "mask outside the band")
            assert by <= y and y + h <= by + bh + 1, (name, "mask outside the band")
        # and no mask covers a pane's visible area
        for px_, pw, ph in cut["panes"]:
            for x, y, w, h in masks:
                overlap_x = max(0, min(px_ + pw, x + w) - max(px_, x))
                overlap_y = max(0, min(ph, y + h) - max(by, y))
                assert not (overlap_x and overlap_y), (name, "mask covers a pane")


def test_a_mask_is_painted_in_the_slides_own_field_colour():
    """A mask that misses the field by a shade shows as a seam."""
    from deckguard import brandmode as bm

    assert L._bg_hex(None) == bm.WHITE
    assert L._bg_hex("sand") == bm.SAND
    assert L._bg_hex("#F3EEE6") == "F3EEE6"
    assert L._bg_hex("light_blue") == bm.LIGHT_BLUE
    assert L._bg_hex("nonsense") == bm.WHITE


def test_every_registered_archetype_has_a_style_for_every_role_it_uses():
    """`KeyError: 'gal_i43_141414'` mid-draw.

    `gallery.install` registered these into the engine's ROLE_STYLE as
    it parsed. Snapshotting its output to JSON captured the archetypes,
    the backgrounds and the samples -- and not this -- so the first
    slide using one crashed the build. Checking every archetype rather
    than the eleven that regressed, because the next omission will be a
    different eleven.
    """
    archetypes = _skill_modules()
    known = set(archetypes.E.ROLE_STYLE)
    # roles the layout layer draws itself, never through the engine
    # roles handled before the ROLE_STYLE lookup is ever reached
    drawn_here = {"dg_text", "dg_bullets", "dg_tick", "picture", "image_band",
                  "panel", "panel_sand", "icon", "figure", "table", "axis",
                  "bullets", None}

    missing = {}
    for name, spec in archetypes.ARCHETYPES.items():
        roles = {r.get("role") for r in spec.get("regions") or []}
        for group in spec.get("groups") or []:
            roles |= {r.get("role") for r in group.get("regions") or []}
        unknown = sorted(r for r in roles if r not in known and r not in drawn_here)
        if unknown:
            missing[name] = unknown
    assert not missing, missing


def test_an_encoded_role_name_is_decoded_to_the_type_it_names():
    """`gal_i43_141414` is "Inter 43, #141414" -- an output used as a
    name, which is why BRAND_MODE retires it. It is decoded literally
    rather than mapped to its brand role, because the box was laid out
    for 43px and `hero_value` is 280px: retiring the role means moving
    the type AND the geometry, which belongs with the re-spec."""
    from deckguard.registry import _decode_role

    assert _decode_role("gal_i43_141414")[:2] == ("Inter", 43)
    assert _decode_role("ref_i53_141414")[:2] == ("Inter", 53)
    kone = _decode_role("gal_k12_FFFFFF_c")
    assert kone[0] == "KONE Information" and kone[3] is True, "KONE Information is caps"
    assert kone[5] is True, "the _c suffix means centred"
    assert _decode_role("not_a_role") is None
    assert _decode_role("gal_bad") is None


def test_bullets_are_real_markers_not_a_dash(tmp_path):
    """BRAND_MODE section 6 calls a dash standing in for a bullet a
    brand violation and names the archetypes doing it. The engine's own
    `_dash_bullets` emits "—  " as the marker, so every archetype with a
    plain `bullets` region was shipping one -- including inside groups,
    which is where all four named archetypes keep theirs."""
    from pptx import Presentation

    archetypes = _skill_modules()
    names = ("statement_links", "lifecycle_4stage", "two_picture_compare",
             "three_pictures_text", "org_functions")
    out = tmp_path / "bullets.pptx"
    L.build_deck({"title": "T", "slides": [
        dict(archetypes.SAMPLES.get(n) or {}, archetype=n) for n in names
    ]}, str(out), archetypes)

    prs = Presentation(str(out))
    dashes, markers = 0, 0
    for slide in prs.slides:
        markers += slide.shapes._spTree.xml.count("buChar")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() in ("—", "-", "–"):
                        dashes += 1
    assert dashes == 0, f"{dashes} dash markers still drawn"
    assert markers > 0, "no real list markers were drawn either"


def test_a_supplied_grid_is_not_drawn_over_a_supplied_paragraph(tmp_path):
    """Found on a real deck: `text_picture_a` serves two shapes from one
    spec -- a paragraph when given `body`, a grid of icon-plus-text
    cells when given `items` -- and the engine skips a region whose key
    is missing, so one spec covers both.

    Nothing said what happens when a planner supplies BOTH, and it does,
    because each key is advertised separately. What happened was a 402px
    paragraph drawn straight through the icon grid: a lead sentence
    crossing a cloud pictogram and two cell captions. Preflight caught
    it as an overlap, which is the right answer far too late.

    The group wins and the paragraph is reported, not silently lost.
    """
    from deckguard import assemble

    slide = {
        "archetype": "text_picture_a",
        "eyebrow": "Ways of working", "title": "How the work reaches you",
        "body": "The routing does not change for this deployment. Tickets "
                "follow the same path they always do.",
        "items": [{"text": f"Step {n}"} for n in range(1, 5)],
    }
    out = tmp_path / "both.pptx"
    report: dict = {}
    L.build_deck({"title": "T", "date": "21 August 2026", "audience": "internal",
                  "slides": [slide]}, str(out), report=report)

    assert report.get("crowded_out"), "the drop has to be reported, not silent"
    assert report["crowded_out"][2] == ("text_picture_a", ["body"])
    assert not assemble.preflight(str(out))["findings"]

    # The title box runs 14px into the icon row -- box against box, ink
    # nowhere near -- and an any-intersection rule threw the headline
    # away. Only a region the grid genuinely sits on counts.
    from deckguard.registry import _load_archetypes

    arch = _load_archetypes().ARCHETYPES["text_picture_a"]
    assert L.alternate_form_conflicts(arch, {"title": "x", "items": [{"text": "y"}]}) == []


def test_the_grid_form_alone_is_untouched():
    """The check must not fire when only one form is supplied -- that is
    the normal case for every archetype that has two."""
    from deckguard.registry import _load_archetypes

    arch = _load_archetypes().ARCHETYPES["text_picture_a"]
    assert L.alternate_form_conflicts(arch, {"items": [{"text": "y"}]}) == []
    assert L.alternate_form_conflicts(arch, {"body": "just a paragraph"}) == []
