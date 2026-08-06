"""Tests for illustrations.py -- KONE line art as editable shape groups.

Conversion rules run on inline fixtures; the tests that read the
shipped SVGs guard the failure mode that matters, which is silence -- a
re-export that stops parsing yields empty geometry and draws nothing at
all rather than raising.
"""

import pytest

from deckguard import illustrations as I

needs_art = pytest.mark.skipif(
    not I.illustration_names(), reason="illustration set not installed"
)

SVG = """<?xml version="1.0"?>
<svg viewBox="0 0 1250 1250">
<style type="text/css">
    .st0{fill:none;stroke:#141414;stroke-width:2;}
    .st1{fill:#1450F5;}
    .st2{fill:#FFFFFF;stroke:#141414;stroke-width:6;}
</style>
<rect class="st1" x="100" y="200" width="300" height="150"/>
<circle class="st0" cx="600" cy="600" r="50"/>
<path class="st2" d="M10 10 h100 v100 z"/>
<line class="st0" x1="0" y1="0" x2="50" y2="50"/>
</svg>"""


# --------------------------------------------------------------------------
# style resolution
# --------------------------------------------------------------------------


def test_classes_resolve_to_fill_stroke_and_width():
    rules = I._css(SVG)
    assert rules["st1"]["fill"] == "#1450F5"
    assert I._resolve({"class": "st2"}, rules) == ("FFFFFF", "141414", 6.0)


def test_fill_none_means_no_fill_not_black():
    """`fill:none` with a stroke is most of this artwork -- it is line
    art. Treating none as the SVG default black would fill every
    outline solid."""
    fill, stroke, _ = I._resolve({"class": "st0"}, I._css(SVG))
    assert fill is None and stroke == "141414"


def test_an_element_with_no_class_gets_svgs_own_default():
    assert I._resolve({}, {})[0] == "000000"


def test_a_presentation_attribute_beats_the_class():
    fill, _, _ = I._resolve({"class": "st1", "fill": "#FFFFFF"}, I._css(SVG))
    assert fill == "FFFFFF"


# --------------------------------------------------------------------------
# primitives
# --------------------------------------------------------------------------


def test_a_rect_becomes_a_closed_path():
    ops = I._primitive_ops("rect", {"x": "10", "y": "20", "width": "30", "height": "40"})
    assert [op for op, _ in ops] == ["moveTo", "lnTo", "lnTo", "lnTo", "close"]
    assert ops[2][1] == [(40.0, 60.0)]


def test_a_circle_becomes_four_beziers():
    """DrawingML has no circle. Four cubics with the kappa constant is
    the standard approximation and stays round under scaling."""
    ops = I._primitive_ops("circle", {"cx": "100", "cy": "100", "r": "50"})
    assert [op for op, _ in ops] == ["moveTo"] + ["cubicBezTo"] * 4 + ["close"]
    assert ops[0][1] == [(150.0, 100.0)]


def test_a_polyline_stays_open_and_a_polygon_closes():
    points = {"points": "0,0 10,10 20,0"}
    assert I._primitive_ops("polyline", points)[-1][0] == "lnTo"
    assert I._primitive_ops("polygon", points)[-1][0] == "close"


def test_a_line_is_two_points():
    ops = I._primitive_ops("line", {"x1": "1", "y1": "2", "x2": "3", "y2": "4"})
    assert ops == [("moveTo", [(1.0, 2.0)]), ("lnTo", [(3.0, 4.0)])]


def test_a_matrix_transform_moves_the_geometry():
    ops = [("moveTo", [(10.0, 20.0)])]
    moved = I._apply(ops, (1, 0, 0, 1, 100, 200))
    assert moved == [("moveTo", [(110.0, 220.0)])]


# --------------------------------------------------------------------------
# assembly
# --------------------------------------------------------------------------


def test_pieces_are_read_in_paint_order_with_their_own_styling(tmp_path, monkeypatch):
    path = tmp_path / "KONE_Illustrations_RGB_Test_art.svg"
    path.write_text(SVG)
    monkeypatch.setattr(I, "illustrations_dir", lambda: tmp_path)
    I._files.cache_clear()
    I.load_illustration.cache_clear()
    try:
        assert I.illustration_names() == ["test-art"]
        pieces = I.load_illustration("test-art")
        assert len(pieces) == 4
        assert pieces[0].fill == "1450F5" and pieces[0].stroke is None
        assert pieces[1].fill is None and pieces[1].stroke_width == 2.0
        assert pieces[2].stroke_width == 6.0
    finally:
        I._files.cache_clear()
        I.load_illustration.cache_clear()


def test_an_unknown_illustration_returns_none():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert I.add_illustration(slide, "not-a-real-drawing", (0, 0, 100, 100)) is None


# --------------------------------------------------------------------------
# the shipped set
# --------------------------------------------------------------------------


@needs_art
def test_the_shipped_set_loads_completely():
    names = I.illustration_names()
    assert len(names) >= 50
    for name in names:
        pieces = I.load_illustration(name)
        assert pieces, f"{name} produced no geometry"
        assert all(piece.ops for piece in pieces), f"{name} has an empty piece"


@needs_art
def test_no_illustration_needs_an_elliptical_arc():
    """Same load-bearing fact as the icons: DrawingML has no arc, so one
    would be silently dropped rather than fail."""
    import re

    offenders = []
    for name, path in I._files().items():
        text = path.read_text(errors="replace")
        for data in re.findall(r'\sd="([^"]+)"', text):
            if re.search(r"[Aa]", data):
                offenders.append(name)
                break
    assert offenders == []


@needs_art
def test_an_illustration_lands_as_a_group_of_shapes_with_no_pictures():
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = I.add_illustration(slide, "cloud", (100, 100, 300, 300))

    assert group is not None and group.shape_type.name == "GROUP"
    assert len(group.shapes) >= 1
    assert not [s for s in slide.shapes if s.shape_type == 13], "no pictures"
    child = list(group.shapes)[0]
    assert child._element.spPr.find(f"{{{I.A_NS}}}custGeom") is not None


@needs_art
def test_a_groups_child_space_matches_its_own_extents():
    """A group has two coordinate systems and PowerPoint scales one onto
    the other. Sizing the group without sizing the child space made that
    ratio arbitrary -- a correctly written 30x13 cloud rendered eight
    times too large, and the file measured perfect throughout."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = I.add_illustration(slide, "cloud", (100, 100, 300, 300))

    xfrm = group._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr"
    ).find(f"{{{I.A_NS}}}xfrm")
    off, ch_off = xfrm.find(f"{{{I.A_NS}}}off"), xfrm.find(f"{{{I.A_NS}}}chOff")
    ext, ch_ext = xfrm.find(f"{{{I.A_NS}}}ext"), xfrm.find(f"{{{I.A_NS}}}chExt")
    assert (off.get("x"), off.get("y")) == (ch_off.get("x"), ch_off.get("y"))
    assert (ext.get("cx"), ext.get("cy")) == (ch_ext.get("cx"), ch_ext.get("cy"))


@needs_art
def test_trimming_fills_the_box_and_the_untrimmed_form_does_not():
    """The art uses a fraction of its 1250 canvas -- the cloud is 146x64
    units in the middle -- so honouring the viewBox puts a postage stamp
    in the middle of whatever box was asked for."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    trimmed = I.add_illustration(slide, "cloud", (0, 0, 400, 400), trim=True)
    framed = I.add_illustration(slide, "cloud", (0, 0, 400, 400), trim=False)
    assert max(trimmed.width, trimmed.height) > max(framed.width, framed.height) * 2


@needs_art
def test_stroke_width_scales_with_the_drawing():
    """A 6-unit outline drawn at a quarter size has to become 1.5, or
    the line art thickens into a blob."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def widest(size):
        group = I.add_illustration(slide, "transport-bus", (0, 0, size, size))
        return max((s.line.width or 0) for s in group.shapes)

    assert widest(600) > widest(150)
