"""Deck-specific, per-run visual transplant for the "Learn from a
reference" flow: copies the EXACT fill/line/font treatment from a
reference deck's shapes onto the corresponding shapes of the deck being
transformed, wherever shape identity survives between the two decks (same
slide index + same shape name, including group ancestry).

This is deliberately separate from `learn.py`'s color/font proposals,
which infer a GENERIC hex->hex / font->font rule and merge it into
brand_rules.yaml for reuse across any deck: that's the right tool when
the same source color always means the same thing, but it's provably
wrong when the reference deck assigns two DIFFERENT target colors to two
shapes that both started from a similarly-classified source color -- on
a real deck pair, a generic remap got 4 of 7 "category chip" colors right
and 3 wrong, because two different source blues/greens/browns each had to
become two different target colors depending on which chip they were on,
which no deck-wide hex->hex rule can express. Where an old shape's
IDENTITY survives into the reference deck (its name, at the same slide
index), the reference's own answer for THAT shape is unambiguous ground
truth -- this reads it directly and applies it only to that one shape,
for this one run. It never writes to brand_rules.yaml and never mutates
the config used elsewhere in the pipeline.

Shapes with no name match fall back to a POSITION match against whatever
in the reference deck's same slide is still unclaimed (see
`_match_by_position`) -- real decks routinely lose shape-name stability
for a design element repeated across many slides (PowerPoint's own
auto-generated names are sequential by creation order, not stable across
independently-edited copies), and without this fallback those shapes
silently fall through to whatever the deck-wide color/font remap already
did, which is exactly the "same source color, different target per
shape" case this whole module exists to fix correctly.

Shapes with no match at all (most often a hand-built diagram redrawn
from scratch, with different connector shapes/counts) are left
completely alone -- this never guesses at a correspondence, only acts on
an exact one. A slide where fewer than half its shapes have a match at
all is reported back via `TransplantResult.flagged_slides` so a human
knows that specific slide's look needs to be reproduced by hand rather
than by this pass.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckguard.colors import ThemeColorScheme, effective_rgb, get_theme_scheme
from deckguard.fonts import normalize_key

# Below this fraction of an old slide's (leaf) shapes finding a name-match
# in the reference deck's same-index slide, the slide is treated as having
# no real structural correspondence to the reference -- most of it is
# unrelated, so applying the handful of coincidental name matches (if any)
# would read as arbitrary rather than a faithful transplant.
LOW_MATCH_THRESHOLD = 0.5

# Summed left/top/width/height difference (EMU) within which two leaf
# shapes are treated as "the same shape, name lost" for the position
# fallback below -- deliberately tight (a small fraction of an inch):
# this is for a design element repeated at a genuinely FIXED spot across
# many slides (a footer chip row, say), not a loose "roughly similar
# shape" match, so it should only ever catch a real repeat, never two
# merely-nearby unrelated shapes on a busy slide.
POSITION_MATCH_TOLERANCE_EMU = 91440


@dataclass
class TransplantChange:
    slide_index: int  # 0-based
    shape_path: str
    field: str  # "fill" | "line" | "font_name" | "font_color" | "font_bold" | "alignment"
    old: Optional[str]
    new: Optional[str]


@dataclass
class TransplantResult:
    changes: list = field(default_factory=list)  # list[TransplantChange]
    flagged_slides: list = field(default_factory=list)  # list[int], 0-based slide indices


def _shape_path_map(shapes, prefix: str = "") -> dict:
    out = {}
    for shape in shapes:
        path = f"{prefix}{shape.name}"
        out[path] = shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            out.update(_shape_path_map(shape.shapes, prefix=path + "/"))
    return out


def _leaf_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _leaf_shapes(shape.shapes)
        else:
            yield shape


def _shape_pos(shape) -> Optional[tuple]:
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
    except Exception:
        return None
    if l is None or t is None or w is None or h is None:
        return None
    return (l, t, w, h)


def _match_by_position(old_items: list, ref_shapes: list) -> list:
    """Greedy nearest-position match for leaf shapes whose NAME didn't
    survive between decks -- a real, common case for a design element
    repeated at a fixed spot across many slides (a footer "category chip"
    row copy-pasted across a whole catalog deck, say): PowerPoint's own
    auto-generated shape names are sequential by creation order within
    each slide, not stable across independently-touched copies, so the
    exact same visual chip can end up "Rectangle 27" on one slide and
    "Rectangle 10" on another. Name matching alone silently misses these
    -- not an error, just no match found -- so the shape falls back to
    whatever the deck-wide color/font remap already did, which is
    provably wrong for exactly this "same source color, different target
    per chip" pattern (see this module's own docstring). Position is the
    reliable identity signal here instead: same design, same fixed spot,
    every time.

    `old_items`: [(path, shape), ...] for old shapes with no name match.
    `ref_shapes`: reference leaf shapes with no name match, at the same
    slide index. Returns [(path, old_shape, new_shape), ...].
    """
    candidates = []
    for oi, (_path, old_shape) in enumerate(old_items):
        old_pos = _shape_pos(old_shape)
        if old_pos is None:
            continue
        for ni, new_shape in enumerate(ref_shapes):
            new_pos = _shape_pos(new_shape)
            if new_pos is None:
                continue
            dist = sum(abs(a - b) for a, b in zip(old_pos, new_pos))
            if dist <= POSITION_MATCH_TOLERANCE_EMU:
                candidates.append((dist, oi, ni))
    candidates.sort(key=lambda c: c[0])

    used_old: set = set()
    used_new: set = set()
    pairs = []
    for _dist, oi, ni in candidates:
        if oi in used_old or ni in used_new:
            continue
        used_old.add(oi)
        used_new.add(ni)
        path, old_shape = old_items[oi]
        pairs.append((path, old_shape, ref_shapes[ni]))
    return pairs


def _canonical_hex(hexval: str, color_remap: dict) -> str:
    """The reference deck is a real-world file, not a guaranteed-pristine
    source of truth -- it can still carry a stale/legacy literal hex
    brand_rules.yaml already knows the approved replacement for (e.g. an
    old neutral grey the deck's own author never got around to refreshing
    to the current sand token). Copying such a value VERBATIM would
    regress a color the generic brand-compliance pass already got right,
    so any hex already covered by `colors.remap` is canonicalized to its
    approved target before it's used as a transplant source or compared
    against the old shape's current value."""
    return color_remap.get(hexval.upper(), hexval.upper())


def _canonical_font(font_name: str, font_remap: dict, font_approved_keys: set) -> Optional[str]:
    """Unlike colors (where a stale-but-known legacy hex still canonicalizes
    to something usable), a font the brand rules have never heard of gets a
    hard NO here, not a best-effort pass-through: this deck must never end
    up carrying a non-KONE font just because the reference deck happened to
    use one (e.g. it was never brand-fixed itself, or was hand-edited after
    the fact). Returns the canonical approved name to use, or None to skip
    the font_name copy entirely and leave the old run's own (already
    brand-fixed) font alone."""
    key = normalize_key(font_name)
    if key in font_remap:
        return font_remap[key]
    if key in font_approved_keys:
        return font_name
    return None


def _copy_fill(old_shape, new_shape, slide_index: int, path: str, color_remap: dict, changes: list) -> None:
    try:
        new_fill = new_shape.fill
        if new_fill.type != MSO_FILL_TYPE.SOLID or new_fill.fore_color.type != MSO_COLOR_TYPE.RGB:
            return
        new_hex = _canonical_hex(str(new_fill.fore_color.rgb), color_remap)
    except Exception:
        return
    old_hex = None
    try:
        old_fill = old_shape.fill
        if old_fill.type == MSO_FILL_TYPE.SOLID and old_fill.fore_color.type == MSO_COLOR_TYPE.RGB:
            old_hex = str(old_fill.fore_color.rgb)
    except Exception:
        pass
    if old_hex == new_hex:
        return
    try:
        old_shape.fill.solid()
        old_shape.fill.fore_color.rgb = RGBColor.from_string(new_hex)
    except Exception:
        return
    changes.append(TransplantChange(slide_index, path, "fill", old_hex, new_hex))


def _copy_line(old_shape, new_shape, slide_index: int, path: str, color_remap: dict, changes: list) -> None:
    try:
        new_line = new_shape.line
        if new_line.color.type != MSO_COLOR_TYPE.RGB:
            return
        new_hex = _canonical_hex(str(new_line.color.rgb), color_remap)
    except Exception:
        return
    old_hex = None
    try:
        old_line = old_shape.line
        if old_line.color.type == MSO_COLOR_TYPE.RGB:
            old_hex = str(old_line.color.rgb)
    except Exception:
        pass
    if old_hex == new_hex:
        return
    try:
        old_shape.line.color.rgb = RGBColor.from_string(new_hex)
    except Exception:
        return
    changes.append(TransplantChange(slide_index, path, "line", old_hex, new_hex))


def _hex_str(rgb: Optional[tuple]) -> Optional[str]:
    return "%02X%02X%02X" % rgb if rgb is not None else None


def _copy_font_color(
    old_run, new_run, own_fill_hex: Optional[str], color_remap: dict,
    old_theme_scheme: Optional[ThemeColorScheme], ref_theme_scheme: Optional[ThemeColorScheme],
    slide_index: int, path: str, changes: list,
) -> None:
    try:
        new_color_type = new_run.font.color.type
    except Exception:
        return
    if new_color_type == MSO_COLOR_TYPE.RGB:
        new_hex = _canonical_hex(str(new_run.font.color.rgb), color_remap)
        old_hex = str(old_run.font.color.rgb) if old_run.font.color.type == MSO_COLOR_TYPE.RGB else None
        # Never transplant a text color that would land exactly on the
        # shape's own (just-transplanted) fill -- e.g. a hidden duplicate/
        # decoy shape in the reference deck with a stale text/fill pairing
        # that happens to canonicalize to the same approved token for both
        # roles. Applying it would make the text invisible; better to
        # leave the old color as-is than regress contrast for a value
        # this pass isn't confident about.
        if old_hex != new_hex and new_hex != own_fill_hex:
            old_run.font.color.rgb = RGBColor.from_string(new_hex)
            changes.append(TransplantChange(slide_index, path, "font_color", old_hex, new_hex))
        return

    if new_color_type == MSO_COLOR_TYPE.SCHEME:
        # A THEME-relative reference (e.g. "this text is always background1,
        # whatever that currently renders as") -- the reference deck's own
        # design intent, not a literal color, and specifically how real
        # decks keep text uniformly on-brand across several differently
        # colored "category chip" boxes (accent1/accent2/.../background1
        # each pick their own contrast-appropriate slot). A generic
        # contrast-fix pass upstream (fix_deck) has no notion of this --
        # it only ever sees/sets literal RGB -- and can end up "fixing"
        # some chips in a uniformly-scheme-colored row to an explicit
        # literal color while leaving others alone, which is exactly the
        # kind of inconsistency this whole module exists to eliminate.
        # Copying the theme reference itself (not just its current
        # resolved RGB) keeps it correct even if the theme changes again.
        new_theme_color = new_run.font.color.theme_color
        new_brightness = new_run.font.color.brightness or 0.0
        new_effective = effective_rgb(new_run.font.color, ref_theme_scheme)
        new_hex = _hex_str(new_effective)
        if new_hex is not None and new_hex == own_fill_hex:
            return
        old_effective = effective_rgb(old_run.font.color, old_theme_scheme)
        old_hex = _hex_str(old_effective)
        if old_hex == new_hex and old_run.font.color.type == MSO_COLOR_TYPE.SCHEME:
            return
        try:
            old_run.font.color.theme_color = new_theme_color
            if new_brightness:
                old_run.font.color.brightness = new_brightness
        except Exception:
            return
        changes.append(TransplantChange(slide_index, path, "font_color", old_hex, new_hex))


def _copy_text_style(
    old_shape, new_shape, slide_index: int, path: str, color_remap: dict,
    font_remap: dict, font_approved_keys: set,
    old_theme_scheme: Optional[ThemeColorScheme], ref_theme_scheme: Optional[ThemeColorScheme],
    changes: list,
) -> None:
    if not (getattr(old_shape, "has_text_frame", False) and getattr(new_shape, "has_text_frame", False)):
        return
    own_fill_hex = None
    try:
        fill = old_shape.fill
        if fill.type == MSO_FILL_TYPE.SOLID and fill.fore_color.type == MSO_COLOR_TYPE.RGB:
            own_fill_hex = str(fill.fore_color.rgb)
    except Exception:
        pass
    for old_para, new_para in zip(old_shape.text_frame.paragraphs, new_shape.text_frame.paragraphs):
        if new_para.alignment is not None and new_para.alignment != old_para.alignment:
            old_val, new_val = old_para.alignment, new_para.alignment
            old_para.alignment = new_val
            changes.append(TransplantChange(slide_index, path, "alignment", str(old_val), str(new_val)))
        for old_run, new_run in zip(old_para.runs, new_para.runs):
            if not old_run.text.strip():
                continue
            canonical_font = _canonical_font(new_run.font.name, font_remap, font_approved_keys) if new_run.font.name else None
            if canonical_font and canonical_font != old_run.font.name:
                old_val, new_val = old_run.font.name, canonical_font
                old_run.font.name = new_val
                changes.append(TransplantChange(slide_index, path, "font_name", old_val, new_val))
            if new_run.font.bold is not None and new_run.font.bold != old_run.font.bold:
                old_val, new_val = old_run.font.bold, new_run.font.bold
                old_run.font.bold = new_val
                changes.append(TransplantChange(slide_index, path, "font_bold", str(old_val), str(new_val)))
            try:
                _copy_font_color(
                    old_run, new_run, own_fill_hex, color_remap,
                    old_theme_scheme, ref_theme_scheme, slide_index, path, changes,
                )
            except Exception:
                pass


def transplant_exact_treatment(prs, reference_prs, rules_config: Optional[dict] = None) -> TransplantResult:
    """Mutates `prs` in place, slide by slide, up to however many slides
    the two decks have in common (by index). `reference_prs` is read-only.

    `rules_config`: the SAME brand_rules.yaml-shaped dict the rest of the
    pipeline uses -- read-only here too, only for its `colors.remap`/
    `fonts.remap`/`fonts.approved` tables (see `_canonical_hex` and
    `_canonical_font`). Never written back to.
    """
    colors_cfg = (rules_config or {}).get("colors", {}) or {}
    color_remap = {
        k.lstrip("#").upper(): v.lstrip("#").upper()
        for k, v in (colors_cfg.get("remap", {}) or {}).items()
    }
    fonts_cfg = (rules_config or {}).get("fonts", {}) or {}
    font_remap = {
        normalize_key(k): v
        for k, v in (fonts_cfg.get("remap", {}) or {}).items()
    }
    font_approved_keys = {normalize_key(n) for n in (fonts_cfg.get("approved", []) or [])}
    result = TransplantResult()

    theme_scheme_cache: dict = {}

    def theme_scheme_for(slide):
        master = slide.slide_layout.slide_master
        key = id(master.part)
        if key not in theme_scheme_cache:
            try:
                theme_scheme_cache[key] = get_theme_scheme(master)
            except Exception:
                theme_scheme_cache[key] = None
        return theme_scheme_cache[key]

    n = min(len(prs.slides), len(reference_prs.slides))
    for i in range(n):
        old_slide = prs.slides[i]
        old_shapes = list(old_slide.shapes)
        old_map = _shape_path_map(old_shapes)
        ref_map = _shape_path_map(list(reference_prs.slides[i].shapes))
        old_leaf_count = sum(1 for _ in _leaf_shapes(old_shapes))
        if old_leaf_count == 0:
            continue

        old_theme_scheme = theme_scheme_for(old_slide)
        ref_theme_scheme = theme_scheme_for(reference_prs.slides[i])

        def apply_copies(old_shape, new_shape, path):
            _copy_fill(old_shape, new_shape, i, path, color_remap, result.changes)
            _copy_line(old_shape, new_shape, i, path, color_remap, result.changes)
            _copy_text_style(
                old_shape, new_shape, i, path, color_remap, font_remap, font_approved_keys,
                old_theme_scheme, ref_theme_scheme, result.changes,
            )

        matched_leaf_count = 0
        name_matched_paths: set = set()
        used_ref_ids: set = set()
        for path, old_shape in old_map.items():
            if old_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            new_shape = ref_map.get(path)
            if new_shape is None or new_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue
            matched_leaf_count += 1
            name_matched_paths.add(path)
            used_ref_ids.add(id(new_shape))
            apply_copies(old_shape, new_shape, path)

        # Position fallback for leaf shapes whose name didn't survive --
        # see _match_by_position's own docstring for why this matters
        # (repeating design elements copy-pasted across many slides).
        unmatched_old = [
            (path, s) for path, s in old_map.items()
            if s.shape_type != MSO_SHAPE_TYPE.GROUP and path not in name_matched_paths
        ]
        unmatched_ref = [
            s for path, s in ref_map.items()
            if s.shape_type != MSO_SHAPE_TYPE.GROUP and id(s) not in used_ref_ids
        ]
        for path, old_shape, new_shape in _match_by_position(unmatched_old, unmatched_ref):
            matched_leaf_count += 1
            apply_copies(old_shape, new_shape, path)

        if matched_leaf_count / old_leaf_count < LOW_MATCH_THRESHOLD:
            result.flagged_slides.append(i)

    return result
