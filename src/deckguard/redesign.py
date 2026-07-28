"""AI-assisted deck redesign: turn ANY starting point -- a completely
off-brand deck with real content, a mostly-empty deck, or no deck at
all, just a topic -- into a fresh outline of the kind `compose.py`
already knows how to build. The realization of the "Phase 2" judgment
layer this project's own README has flagged as not-yet-implemented
since Phase 1.

Deliberately narrow in scope, and deliberately NOT where any new brand
logic lives:

- Content extraction reuses `retemplate.py`'s own shape-walk
  (`_extract_slide_content`) verbatim, so a table, chart, embedded
  object, media, group shape, or overload of free-form decorative
  shapes is left alone and reported here exactly as it is for
  `retemplate` -- the model never sees a slide this project's own rules
  already consider unsafe to reinterpret, no matter what mode redesign
  is running in.
- What's DIFFERENT from `retemplate` on purpose: retemplate caps a
  slide at `MAX_TEXT_BLOCKS` (3) separate text boxes, because it
  carries content over verbatim onto exactly ONE new slide -- that's
  the literal max placeholder count any single layout offers, a real
  ceiling. redesign also carries content over verbatim (never
  rewording, never condensing -- see REDESIGN_RULES), but isn't
  confined to one output slide for it: a source slide with more text
  than any layout holds gets SPLIT across multiple output slides
  instead, so nothing is ever dropped or paraphrased just because it
  was originally typed into a lot of text boxes. What redesign caps is
  total text *volume* (`REDESIGN_MAX_TEXT_CHARS`), and generously -- a
  sanity ceiling against a pathological/corrupted file (an unbounded
  number of split slides is still a real cost/quality concern), not a
  second-guess of an ordinary dense slide.
- The LLM call decides two things, and only two, for a slide WITH
  source content: which `compose.py` slide *kind* best fits it, and --
  when it doesn't fit on one slide -- how to split its existing text
  across as many output slides as it takes. It is never asked to, and
  the prompt explicitly forbids it from, changing a single word of
  that text. Separately, for a blank slide or a bare topic brief with
  no slide behind it at all, it decides what content to AUTHOR from
  the brief to make the deck whole -- the opposite rule
  (never-invent vs. please-invent-from-the-brief), and the prompt keeps
  the two explicitly separate so the model never blurs them.
- The model's output is validated against a JSON schema shaped exactly
  like `compose.py`'s own outline dict format (see `outline_from_list`),
  so a human-written YAML outline, a redesigned deck, and a
  from-scratch AI-authored one are indistinguishable from that point
  on -- they all run through the identical `build_deck` (same layout
  selection, same final `fix_deck` pass, same brand guarantees).
- Nothing about color, font, or layout-approval judgment is delegated to
  the model, in any of the three modes. That's exactly the split this
  project has used since Phase 1: deterministic first, AI second, and
  AI only for the one judgment call a shape-count heuristic can't make.
- A source slide's own images are carried into its redesigned replacement
  by `_attach_source_images`, keyed off `source_slide_index` -- this is
  deliberately NOT something the model decides. The outline schema the
  model fills in has no field for images at all: the model is never shown
  the actual pixels (only an `image_count`), so it has no basis to choose
  among them. Which images survive is a deterministic backfill after the
  model call, capped per slide at `REDESIGN_IMAGES_PER_SLIDE`.

Three ways to call `redesign_deck` in its default `mode="rewrite"`,
matched to the three starting points:

    redesign_deck("old_deck.pptx", out_path)                  # redesign
    redesign_deck("half_empty.pptx", out_path, brief="...")   # fill gaps
    redesign_deck(None, out_path, brief="...")                # from scratch

A fourth, orthogonal option -- `mode="brand"` -- skips the LLM
entirely: it's `retemplate.apply_rebrand`'s fully deterministic
verbatim-carryover-plus-layout-variety-plus-cover/end-swap path (see
its own docstring), for a deck whose wording is fine as written and
just needs to land on brand. No API key needed.

    redesign_deck("old_deck.pptx", out_path, mode="brand")

Requires the `anthropic` package (a base dependency -- see
requirements.txt / pyproject.toml) and an `ANTHROPIC_API_KEY` at
runtime. Everything else in `deckguard` works without either; this is
the one command in the whole tool that makes an outbound API call.
"""

from __future__ import annotations

import json
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import anthropic
from pptx import Presentation

from deckguard.compose import ComposeError, ComposeResult, Outline, build_deck, outline_from_list
from deckguard.config import default_config_path, load_config
from deckguard.fixer import fix_deck
from deckguard.retemplate import EMPTY_SLIDE_REASON, SlideProfile, _extract_slide_content

DEFAULT_MODEL = "claude-opus-5"

# Deliberately NOT a cap on text-block *count* (a first cut at this used
# one, at 10 -- still wrong, and for the same reason a cap of 3 was
# wrong: block count doesn't measure how much there is to place. A
# slide split across 20 tiny caption boxes has less actual content than
# one with 3 paragraph-sized text boxes, and redesign carries every one
# of them over verbatim regardless of how the original was carved up
# -- splitting across multiple output slides rather than condensing,
# see REDESIGN_RULES. What actually bounds the work (and the cost,
# since an unbounded split is still real output) is total text VOLUME,
# so that's what's capped here -- generously, since even a genuinely
# dense slide is a few thousand characters at most; this exists to
# catch a pathological or corrupted file, not to second-guess an
# ordinary hand-built slide.
REDESIGN_MAX_TEXT_CHARS = 20_000
REDESIGN_MAX_IMAGES = 12

# How many of a source slide's own images get carried into its redesigned
# replacement. Not the same cap as REDESIGN_MAX_IMAGES above (that's an
# eligibility ceiling -- how many images before a slide is refused
# entirely); this is a per-slide output cap, set to the most any
# picture-carrying candidate layout in compose.py's CONTENT_LAYOUT_CANDIDATES
# actually has PICTURE placeholders for ("Two pictures and text *" tops out
# at 2). Keeping every image the model never even sees the pixels of and
# has no way to curate isn't the goal here -- just not silently discarding
# all of them, which is what happened before this existed.
REDESIGN_IMAGES_PER_SLIDE = 2

# As of this writing (see the claude-api skill's cached pricing table) --
# used only to give the caller a rough, clearly-labeled cost estimate
# alongside the real usage numbers the API returns. Verify against
# platform.claude.com/docs/en/pricing before trusting this for billing.
PRICE_PER_MTOK_USD = {
    "claude-opus-5": {"input": 5.00, "output": 25.00},
    "claude-sonnet-5": {"input": 3.00, "output": 15.00},
    "claude-haiku-4-5": {"input": 1.00, "output": 5.00},
}


def _classify_slide_for_redesign(slide, slide_height_in: Optional[float] = None) -> SlideProfile:
    """Same shape-safety rules as retemplate.classify_slide (a table,
    chart, embedded object, media, group, or decorative-shape overload
    is still a hard skip -- never touched, brief or no brief). No cap on
    text-block *count* at all -- redesign carries every block over
    verbatim, splitting across multiple output slides rather than
    condensing when one won't hold it all (see REDESIGN_RULES), so how
    many boxes the original text happened to be split across is
    irrelevant; only total text *volume* is capped (generously, as a
    sanity ceiling against an unbounded split on a pathological file),
    and image count gets a higher ceiling than retemplate's, since
    redesign is allowed to select a subset."""
    title, text_blocks, images, reason = _extract_slide_content(slide, slide_height_in)
    if reason:
        return SlideProfile(None, [], [], False, reason)
    if title is None and not text_blocks and not images:
        return SlideProfile(None, [], [], False, EMPTY_SLIDE_REASON)
    total_chars = sum(len(text) for block in text_blocks for _level, text in block)
    if total_chars > REDESIGN_MAX_TEXT_CHARS:
        return SlideProfile(None, [], [], False, "far more text than could be reasonably split across a few slides")
    if len(images) > REDESIGN_MAX_IMAGES:
        return SlideProfile(None, [], [], False, "far more images than any layout could hold")
    return SlideProfile(title=title, text_blocks=text_blocks, images=images, eligible=True)

KIND_GUIDE = """\
Available slide kinds and when to use each:

- cover: the deck's own title slide -- exactly one, first.
- agenda: a short list of upcoming topics/sections (an outline of the
  deck itself). Only worth including if the deck has enough distinct
  sections to preview.
- section: a chapter/divider slide -- just a short heading, marking a
  transition between topics. No body content.
- content: the default for an ordinary slide with a title and 1-3 blocks
  of bullet points. Use `columns` (one list of bullets per column) when
  the content visually reads as 2-3 side-by-side blocks; otherwise use a
  single `bullets` list.
- quote: a slide whose main content is a single attributed quotation.
  Needs quote_text and, if known, quote_author.
- statement: one short, unmissable, single-sentence message with no
  supporting bullets -- a big declarative claim, not a list.
- stat: a row of 2-3 numeric callouts (e.g. "40M / people moved daily").
  Needs `stats`, each with a short `number` and a `label`.
- timeline: a sequence of 2-3 dated/labeled milestones. Needs
  `milestones`, each with a `label` (date/phase) and `text` (what
  happens).
- end: the deck's closing/thank-you slide -- exactly one, last.
- blank: use only if a slide's content genuinely doesn't fit any of the
  above and you cannot responsibly summarize it (rare -- prefer content).

General rules:
- `source_slide_index` must be the 1-based index of the ORIGINAL slide
  (given below) an entry was built from, or `null` if you authored this
  slide from scratch (see below) with no original slide behind it.
- A coherent deck reads as a narrative, not a pile of unrelated slides:
  cover, optional agenda, then content grouped under section dividers
  where it naturally clusters, closing with end. Don't force an agenda
  or section dividers onto a short/simple deck that doesn't need them.
"""

REDESIGN_RULES = """\
Rules for slides that have ORIGINAL source content (see the extracted
text below) -- this is RE-LAYOUT, not a rewrite. Your only job for
these slides is deciding which `kind` and layout best presents content
that already exists -- never editing what it says:
- Never invent facts, numbers, names, or claims that are not present in
  that slide's own extracted source text.
- Never paraphrase, tighten, summarize, or condense the wording either
  -- carry every bullet/line over VERBATIM, character for character
  (light structural cleanup only: e.g. dropping a "Subject:" prefix
  from something that's clearly a title, is fine; rewriting a sentence
  is not). If the source text is genuinely messy, it stays messy --
  fixing prose is not this tool's job.
- Every layout has a real capacity limit (at most 3 body columns, a
  handful of bullets each). A source slide may have far more raw text
  than that -- it was likely built by hand with many separate text
  boxes. When that happens, DO NOT drop or condense anything: return
  MULTIPLE outline entries for that one source slide instead, each
  with the SAME `source_slide_index`, splitting its bullets/columns
  across them so every line ends up somewhere. Use as few slides as
  keep each one readable (usually 2-3), and keep them adjacent in your
  output so the split reads as one continuous section. Use `columns`
  to organize genuinely multi-part content into up to 3 side-by-side
  groups on any one of those slides if that reads better than a single
  list.
- Every source slide must appear in your output at least once, in
  order, each entry carrying its correct `source_slide_index`.
"""

AUTHORING_RULES = """\
Rules for content you are AUTHORING FROM SCRATCH (blank slides in the
source deck, and/or the brief below) -- this is the opposite of the
redesign rules above: there is no original text to preserve, so write
real, substantive, specific content grounded in the brief. Don't pad
with vague filler ("drives value", "best-in-class") -- if the brief
doesn't give you enough to make a claim specific, write a more general
but still concrete statement rather than invent a fake number or name.
Give these entries `source_slide_index: null`.
"""


def _target_slides_guidance(target_slides: Optional[int]) -> str:
    if target_slides is not None:
        return f"Aim for approximately {target_slides} total slides in the finished deck."
    return (
        "No specific slide count was requested -- use your judgment for a deck length "
        "that suits the brief's scope (a narrow topic might be 5-6 slides; a broad one 10-12)."
    )


OUTLINE_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "source_slide_index": {"type": ["integer", "null"]},
        "kind": {
            "type": "string",
            "enum": [
                "cover", "agenda", "section", "content", "quote",
                "statement", "stat", "timeline", "end", "blank",
            ],
        },
        "title": {"type": ["string", "null"]},
        "subtitle": {"type": ["string", "null"]},
        "bullets": {"type": "array", "items": {"type": "string"}},
        "columns": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}},
        "quote_text": {"type": ["string", "null"]},
        "quote_author": {"type": ["string", "null"]},
        "quote_label": {"type": ["string", "null"]},
        "stats": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {"number": {"type": "string"}, "label": {"type": "string"}},
                "required": ["number", "label"],
                "additionalProperties": False,
            },
        },
        "milestones": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {"label": {"type": "string"}, "text": {"type": "string"}},
                "required": ["label", "text"],
                "additionalProperties": False,
            },
        },
        "variant": {"type": ["string", "null"]},
    },
    "required": [
        "source_slide_index", "kind", "title", "subtitle", "bullets", "columns",
        "quote_text", "quote_author", "quote_label", "stats", "milestones", "variant",
    ],
    "additionalProperties": False,
}

OUTLINE_SCHEMA = {
    "type": "object",
    "properties": {"slides": {"type": "array", "items": OUTLINE_ITEM_SCHEMA}},
    "required": ["slides"],
    "additionalProperties": False,
}


class RedesignError(ComposeError):
    """Raised when the AI redesign step itself fails (missing API key,
    missing dependency, nothing to work from, or a malformed/unusable
    model response)."""


def _stream_final_message(client, **stream_kwargs):
    """Call `client.messages.stream(...)` and return the final message,
    translating Anthropic's own error types into a clean `RedesignError`
    instead of letting a raw API error dict reach the user -- a
    `rate_limit_error`/`overloaded_error`/5xx is transient on Anthropic's
    end and worth retrying as-is; anything else is a real request problem.
    """
    try:
        with client.messages.stream(**stream_kwargs) as stream:
            return stream.get_final_message()
    except anthropic.APIStatusError as exc:
        if exc.status_code == 429 or exc.status_code >= 500:
            raise RedesignError(
                f"Claude's API is temporarily rate-limited or overloaded (HTTP {exc.status_code}) -- "
                "this isn't a problem with your deck, just Anthropic's servers being busy right now. "
                "Wait a moment and try again."
            ) from exc
        raise RedesignError(f"Claude API error (HTTP {exc.status_code}): {exc.message}") from exc
    except anthropic.APIConnectionError as exc:
        raise RedesignError(f"Could not reach the Anthropic API: {exc}") from exc


@dataclass
class SkippedSlide:
    slide_index: int  # 1-based, original deck numbering
    reason: str


@dataclass
class Usage:
    input_tokens: int
    output_tokens: int
    model: str

    @property
    def estimated_cost_usd(self) -> float:
        prices = PRICE_PER_MTOK_USD.get(self.model)
        if not prices:
            return 0.0
        return (self.input_tokens / 1_000_000) * prices["input"] + (self.output_tokens / 1_000_000) * prices["output"]


@dataclass
class RedesignResult:
    outline: Outline
    skipped: list  # list[SkippedSlide] -- unsafe-to-touch slides, never sent to the model
    usage: Usage
    review_notes: list = field(default_factory=list)  # brand mode --review's free-text findings, "slide N: ..."
    reference_match_notes: list = field(default_factory=list)  # "Learn" mode's exact-transplant pass findings


def _slide_height_in(prs) -> Optional[float]:
    return prs.slide_height / 914400 if prs.slide_height else None


def _describe_slide(index: int, profile: SlideProfile) -> dict:
    return {
        "source_slide_index": index,
        "title": profile.title,
        "text_blocks": [[text for _level, text in block] for block in profile.text_blocks],
        "image_count": len(profile.images),
    }


def extract_eligible_slides(prs) -> tuple[list, list]:
    """Classify every slide with `_classify_slide_for_redesign` --
    retemplate's own shape-safety rules, but with redesign's own
    (higher) cap on text/image count, since redesign condenses rather
    than carrying content over verbatim. Returns (eligible:
    list[(index, SlideProfile)], skipped: list[SkippedSlide]) --
    `skipped` here includes blank slides; callers that need blanks
    split out separately should use `partition_skipped`."""
    slide_height_in = _slide_height_in(prs)
    eligible = []
    skipped = []
    for i, slide in enumerate(prs.slides, start=1):
        profile = _classify_slide_for_redesign(slide, slide_height_in)
        if profile.eligible:
            eligible.append((i, profile))
        else:
            skipped.append(SkippedSlide(slide_index=i, reason=profile.reason or "not eligible for redesign"))
    return eligible, skipped


def partition_skipped(skipped: list) -> tuple[list, list]:
    """Split `extract_eligible_slides`'s skip list into (blank slide
    indices, everything else). A blank slide has nothing to protect --
    it's safe to author content for it from a brief. Every other skip
    reason (table/chart/media/overfull) means real content this project
    has already decided is unsafe to reinterpret, brief or no brief."""
    blank_indices = [s.slide_index for s in skipped if s.reason == EMPTY_SLIDE_REASON]
    real_skipped = [s for s in skipped if s.reason != EMPTY_SLIDE_REASON]
    return blank_indices, real_skipped


def _build_messages(
    eligible: list, blank_count: int, brief: Optional[str], notes: Optional[str], target_slides: Optional[int]
) -> list:
    sections = [KIND_GUIDE]

    if eligible:
        sections.append(REDESIGN_RULES)
    if blank_count or brief or not eligible:
        sections.append(AUTHORING_RULES)
        sections.append(_target_slides_guidance(target_slides))

    if notes:
        sections.append(f"Additional guidance from the operator running this tool:\n{notes}")
    if brief:
        sections.append(f"Brief describing the deck to build (use this to author any new content needed):\n{brief}")

    if eligible:
        slides_json = json.dumps([_describe_slide(i, p) for i, p in eligible], indent=2)
        sections.append(
            f"Here is the extracted content of every slide with original source content, in order "
            f"(redesign each of these per the rules above):\n\n{slides_json}"
        )
    if blank_count:
        sections.append(
            f"The source deck also has {blank_count} blank slide(s) with no content at all -- "
            "author new slides for these (or fold them into a better overall structure) using the brief."
        )
    if not eligible and not blank_count:
        sections.append("There is no source deck -- author the entire outline from the brief above.")

    return [{"role": "user", "content": "\n\n".join(sections)}]


def call_claude_for_outline(
    eligible: list,
    blank_count: int = 0,
    brief: Optional[str] = None,
    target_slides: Optional[int] = None,
    model: str = DEFAULT_MODEL,
    effort: str = "high",
    notes: Optional[str] = None,
    api_key: Optional[str] = None,
    client=None,
) -> tuple[list, Usage]:
    """Send whatever's available -- extracted content, a count of blank
    slides to fill, and/or a topic brief -- to Claude and get back a
    list of outline-item dicts (compose.py's own schema) plus the call's
    token usage. `client` is an injection point for tests -- any object
    exposing `.messages.stream(...)` in the Anthropic SDK shape.
    """
    if client is None:
        key = api_key or os.environ.get("ANTHROPIC_API_KEY")
        if not key:
            raise RedesignError("ANTHROPIC_API_KEY is not set -- redesign needs an Anthropic API key")
        client = anthropic.Anthropic(api_key=key)

    messages = _build_messages(eligible, blank_count, brief, notes, target_slides)
    response = _stream_final_message(
        client,
        model=model,
        max_tokens=32000,
        thinking={"type": "adaptive"},
        output_config={"effort": effort, "format": {"type": "json_schema", "schema": OUTLINE_SCHEMA}},
        messages=messages,
    )

    if response.stop_reason == "refusal":
        raise RedesignError("Claude declined the redesign request (safety refusal) -- try again or adjust the brief/source deck")

    text = next((b.text for b in response.content if b.type == "text"), None)
    if not text:
        raise RedesignError(f"no text content in the model response (stop_reason={response.stop_reason!r})")
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError as exc:
        raise RedesignError(f"model response was not valid JSON: {exc}") from exc

    raw_slides = parsed.get("slides")
    if not isinstance(raw_slides, list):
        raise RedesignError("model response had no 'slides' array")
    if not raw_slides:
        raise RedesignError("model returned zero slides")

    usage = Usage(input_tokens=response.usage.input_tokens, output_tokens=response.usage.output_tokens, model=model)
    return raw_slides, usage


# --------------------------------------------------------------------------
# Brand-mode's optional --review pass: a deliberately small, narrowly-scoped
# AI call, NOT the full redesign pipeline above. It looks only at slides
# brand mode already left untouched (real content it can't safely condense,
# not blank/empty ones), and only ever answers one question per slide: does
# this read as a short divider/transition/section-break page (an Appendix,
# a Q&A break, etc.) that the org template has a purpose-built layout for,
# and if so what should its title say. It never rewrites content, never
# authors new slides, and never touches a slide with real body content --
# see REVIEW_RULES below for the exact instruction.
# --------------------------------------------------------------------------

REVIEW_MODEL = "claude-haiku-4-5"

REVIEW_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "slide_index": {"type": "integer"},
        "is_divider": {
            "type": "boolean",
            "description": "true only if this slide is a short section-divider/transition page (a single "
            "short heading and little/no other real content) -- never true for a slide with real body content.",
        },
        "divider_title": {
            "type": ["string", "null"],
            "description": "Required if is_divider is true: a short title (a few words) taken from this "
            "slide's own text -- never invented. Null if is_divider is false.",
        },
        "note": {
            "type": ["string", "null"],
            "description": "A short (one sentence) note flagging unreplaced placeholder/template copy (e.g. "
            "\"Lorem ipsum\", \"[bracketed placeholder]\", \"PRODUCT NAME\") or a confidentiality/proprietary "
            "notice worded differently than plain \"Confidential\". Null if neither is visible in the text. "
            "Text-only judgment (no visual rendering is shown) -- never guess at colors, fonts, or layout.",
        },
    },
    "required": ["slide_index", "is_divider", "divider_title", "note"],
    "additionalProperties": False,
}

REVIEW_SCHEMA = {
    "type": "object",
    "properties": {"slides": {"type": "array", "items": REVIEW_ITEM_SCHEMA}},
    "required": ["slides"],
    "additionalProperties": False,
}

REVIEW_RULES = """\
Each slide below was left untouched by a deterministic brand-compliance pass
because it has real content that can't be safely rewritten (see its "reason"
field for why -- e.g. more text than any layout could hold, or a genuinely
short heading with nothing else on the slide). You have two jobs, both
narrow, both text-only (you're given a preview, not a rendering, so never
guess at colors, fonts, or layout you can't see):

1. Spot a slide that is actually a short section-divider or transition
   page -- like an "Appendix", "Q&A", "Next Steps", or "Thank You" break
   between sections -- which the org template has a purpose-built divider
   layout for, currently unused because this slide instead kept its old,
   off-brand structure.
2. Separately, in `note`, flag anything else the text preview itself
   plainly reveals as off-brand or leftover, specifically:
   - Placeholder/template copy that was never replaced with real content
     (e.g. "[bracketed placeholder]", "Lorem ipsum...", "PRODUCT NAME",
     "TBD", "XXX", "Insert text here").
   - A confidentiality/proprietary/internal-use notice worded differently
     than plain "Confidential" (which a separate deterministic pass
     already removes on its own) -- e.g. "Internal use only",
     "Proprietary and confidential", "Do not distribute".
   Don't speculate beyond what the text shows -- if nothing like this is
   visible, note is null.

Rules:
- Mark is_divider true ONLY for a slide whose entire real content is a
  single short heading/label with nothing else of substance -- if it has
  any real body paragraph, bullet list, or table/chart, is_divider is
  always false, no matter how short the reason field's preview looks.
- divider_title must be taken verbatim (or lightly trimmed, e.g. dropping
  a "Subject:" prefix) from that slide's own extracted text -- never
  invent or guess a title for a slide that doesn't already have one.
- Return exactly one entry per slide listed, in the same order.
"""


def _describe_skipped_slide_for_review(proposal) -> dict:
    return {
        "slide_index": proposal.slide_index,
        "reason": proposal.reason,
        "title_preview": proposal.title_preview,
        "body_preview": proposal.body_preview,
    }


def _build_review_messages(skipped_proposals: list) -> list:
    slides_json = json.dumps([_describe_skipped_slide_for_review(p) for p in skipped_proposals], indent=2)
    content = f"{REVIEW_RULES}\n\nSlides to review:\n\n{slides_json}"
    return [{"role": "user", "content": content}]


def call_claude_for_brand_review(
    skipped_proposals: list,
    model: str = REVIEW_MODEL,
    api_key: Optional[str] = None,
    client=None,
) -> tuple[list, "Usage"]:
    """Send brand mode's skipped-slide previews to Claude for the narrow
    divider-detection judgment call described above. Returns (list of
    per-slide review dicts, Usage). `client` is the same test-injection
    point `call_claude_for_outline` uses.
    """
    if not skipped_proposals:
        return [], Usage(input_tokens=0, output_tokens=0, model=model)

    if client is None:
        key = api_key or os.environ.get("ANTHROPIC_API_KEY")
        if not key:
            raise RedesignError("ANTHROPIC_API_KEY is not set -- --review needs an Anthropic API key")
        client = anthropic.Anthropic(api_key=key)

    messages = _build_review_messages(skipped_proposals)
    response = _stream_final_message(
        client,
        model=model,
        max_tokens=4000,
        # REVIEW_MODEL (Haiku 4.5) doesn't support adaptive thinking or the
        # `effort` knob -- both 400 on this model -- so this call only sets
        # the structured-output schema, unlike the outline call above.
        output_config={"format": {"type": "json_schema", "schema": REVIEW_SCHEMA}},
        messages=messages,
    )

    if response.stop_reason == "refusal":
        raise RedesignError("Claude declined the review request (safety refusal)")

    text = next((b.text for b in response.content if b.type == "text"), None)
    if not text:
        raise RedesignError(f"no text content in the review response (stop_reason={response.stop_reason!r})")
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError as exc:
        raise RedesignError(f"review response was not valid JSON: {exc}") from exc

    raw_slides = parsed.get("slides")
    if not isinstance(raw_slides, list):
        raise RedesignError("review response had no 'slides' array")

    usage = Usage(input_tokens=response.usage.input_tokens, output_tokens=response.usage.output_tokens, model=model)
    return raw_slides, usage


def _strip_ai_only_fields(raw_slides: list) -> list:
    """source_slide_index is for our own traceability; compose.py's
    SlideSpec doesn't take it, and unknown keys are otherwise harmless
    (ignored via .get()), but stripping keeps the outline dict honestly
    scoped to what compose.py actually reads."""
    return [{k: v for k, v in s.items() if k != "source_slide_index"} for s in raw_slides]


def _attach_source_images(raw_slides: list, eligible: list) -> list:
    """Carry each source slide's own images into the outline entry built
    from it. Which images belong to which output slide is a fact already
    known from source_slide_index (compose.py's outline schema has no
    field for the model to report images itself -- it's never shown the
    pixels, so it isn't in a position to choose among them anyway); this
    is a deterministic backfill, not a judgment call, done after the
    model call rather than asking the model to round-trip image data it
    never needed to see. Capped per slide at REDESIGN_IMAGES_PER_SLIDE so
    a slide that had many images still lands on a layout match_layout can
    actually find (see that constant's own comment).

    A source slide's content can now legally split across MULTIPLE
    output entries sharing the same source_slide_index (see
    REDESIGN_RULES) -- its images are attached only to the FIRST such
    entry, never duplicated onto every split piece, since the same
    picture appearing on 2-3 consecutive slides would look like a
    copy-paste mistake, not a deliberate choice.
    """
    images_by_source_index = {i: profile.images[:REDESIGN_IMAGES_PER_SLIDE] for i, profile in eligible}
    attached_indices: set = set()
    for slide in raw_slides:
        src = slide.get("source_slide_index")
        if src in attached_indices:
            continue
        images = images_by_source_index.get(src)
        if images:
            slide["images"] = images
            attached_indices.add(src)
    return raw_slides


def _rebrand_deck(
    deck_path, out_path, template_path=None, rules_config: Optional[dict] = None,
    review: bool = False, review_model: str = REVIEW_MODEL,
    api_key: Optional[str] = None, client=None, reference_path: Optional[str] = None,
):
    """mode='brand' path: no LLM call at all by default, so wrap
    `apply_rebrand`'s result into the same `(ComposeResult,
    RedesignResult)` shape the rewrite path returns, for a single
    consistent return type regardless of mode. `Usage` is all zeros
    with `model="none"` -- there's no API call to report, and
    `Usage.estimated_cost_usd` is already defined to read as $0 for an
    unrecognized model.

    `review=True` adds the one deliberately small, optional AI step:
    see `call_claude_for_brand_review`'s own module-level comment for
    what it does and doesn't do. It only ever looks at slides
    `apply_rebrand` already left untouched (never a rebuilt slide), and
    its only possible structural action is rebuilding a divider-like
    slide onto the org template's own Section Divider layout with a
    short, verbatim-derived title -- nothing it finds is silently
    discarded either: anything else worth a human's attention comes
    back in `RedesignResult.review_notes`.

    `reference_path` (the "Learn from a reference" flow only): after
    everything above, run `exact_transplant.transplant_exact_treatment`
    against this reference deck as one final pass -- see that module's
    own docstring for why this is a DIFFERENT, more precise mechanism
    than the generic `colors.remap`/`fonts.remap` tables `apply_rebrand`'s
    own `fix_deck` pass already applied: it copies the reference's exact
    per-shape answer wherever shape identity survives (same slide index +
    name), rather than one hex/font substitution applied deck-wide. Never
    writes to brand_rules.yaml. Findings land in
    `RedesignResult.reference_match_notes`.
    """
    from deckguard.retemplate import EMPTY_SLIDE_REASON, apply_rebrand, rebuild_slides_as_dividers
    from deckguard.slide_import import default_template_path

    rebrand_result = apply_rebrand(
        str(deck_path), out_path, template_path=template_path, rules_config=rules_config,
        reference_path=reference_path,
    )

    layouts_used = [p.layout_name for p in rebrand_result.proposals if p.eligible and p.layout_name]
    ineligible = [p for p in rebrand_result.proposals if not p.eligible]
    skipped_indices = {p.slide_index for p in ineligible}
    manual_review = rebrand_result.manual_review
    usage = Usage(input_tokens=0, output_tokens=0, model="none")
    review_notes: list = []
    rebuilt_divider_count = 0

    if review:
        reviewable = [p for p in ineligible if p.reason != EMPTY_SLIDE_REASON]
        if reviewable:
            raw_review, usage = call_claude_for_brand_review(
                reviewable, model=review_model, api_key=api_key, client=client,
            )
            title_by_index: dict = {}
            for item in raw_review:
                idx = item.get("slide_index")
                if not isinstance(idx, int):
                    continue
                if item.get("is_divider") and item.get("divider_title"):
                    title_by_index[idx] = str(item["divider_title"]).strip()
                note = item.get("note")
                if note:
                    review_notes.append(f"slide {idx}: {note}")

            if title_by_index:
                effective_template = Path(template_path) if template_path else default_template_path()
                rebuilt = rebuild_slides_as_dividers(out_path, out_path, effective_template, title_by_index)
                rebuilt_divider_count = len(rebuilt)
                variants = ["Section divider A", "Section divider B"]
                layouts_used += [variants[i % len(variants)] for i in range(rebuilt_divider_count)]
                skipped_indices -= set(rebuilt)

                # The new divider slides' freshly-written title text has no
                # explicit run color yet (same "inherited, therefore
                # unresolved" situation apply_rebrand's own fix_deck pass
                # already resolved once for its first rebuild round) --
                # re-run it so these get the same brand-color guarantee.
                config = rules_config if rules_config is not None else load_config(default_config_path())
                reopened = Presentation(str(out_path))
                fix_report = fix_deck(reopened, config, source_path=str(out_path), output_path=str(out_path), dry_run=False)
                manual_review = fix_report.manual_review

    reference_match_notes: list = []
    if reference_path is not None:
        from deckguard.exact_transplant import transplant_exact_treatment

        config = rules_config if rules_config is not None else load_config(default_config_path())
        out_prs = Presentation(str(out_path))
        ref_prs = Presentation(str(reference_path))
        transplant_result = transplant_exact_treatment(out_prs, ref_prs, rules_config=config)
        if transplant_result.changes:
            out_prs.save(str(out_path))
            touched_slides = sorted({c.slide_index + 1 for c in transplant_result.changes})
            reference_match_notes.append(
                f"Reference match: {len(transplant_result.changes)} shape style(s) copied exactly from the "
                f"reference deck across {len(touched_slides)} slide(s) ({', '.join(map(str, touched_slides))})."
            )
        for slide_index in transplant_result.flagged_slides:
            # A REBUILT slide (e.g. the cover/end swap) is expected to have
            # entirely new shape names from the org template and will
            # always show as "no match" here -- that's not a diagram this
            # pass failed to reconcile, so only flag slides `apply_rebrand`
            # itself left untouched (still carrying their OLD shape names,
            # where a low match ratio is actually diagnostic).
            if (slide_index + 1) not in skipped_indices:
                continue
            reference_match_notes.append(
                f"slide {slide_index + 1}: reference deck uses a different diagram/shape layout here -- "
                "not auto-matched, review manually."
            )

    compose_result = ComposeResult(
        slide_count=len(rebrand_result.transformed) + rebuilt_divider_count,
        layouts_used=layouts_used, manual_review=manual_review,
    )
    skipped = [
        SkippedSlide(slide_index=p.slide_index, reason=p.reason or "not eligible")
        for p in ineligible if p.slide_index in skipped_indices
    ]
    redesign_result = RedesignResult(
        outline=Outline(slides=[]), skipped=skipped, usage=usage, review_notes=review_notes,
        reference_match_notes=reference_match_notes,
    )
    return compose_result, redesign_result


def redesign_deck(
    deck_path=None,
    out_path=None,
    brief: Optional[str] = None,
    target_slides: Optional[int] = None,
    model: str = DEFAULT_MODEL,
    effort: str = "high",
    notes: Optional[str] = None,
    template_path=None,
    rules_config: Optional[dict] = None,
    api_key: Optional[str] = None,
    client=None,
    mode: str = "rewrite",
    review: bool = False,
    review_model: str = REVIEW_MODEL,
    reference_path=None,
) -> tuple[ComposeResult, RedesignResult]:
    """One entry point for all three starting points:

    - `deck_path` set, `brief` unset: redesign an existing deck's
      content onto the org template (blank slides in it are skipped,
      same as any other unsafe-to-touch content, since there's no brief
      to author them from).
    - `deck_path` set, `brief` set: redesign the deck's real content
      AND author its blank slides from the brief, as one coherent deck.
    - `deck_path` unset, `brief` set: no source deck at all -- author
      the entire outline from the brief, exactly like asking a designer
      to build a deck on a topic from nothing.

    Whichever mode, everything downstream is identical: the resulting
    outline runs through the same `build_deck`/`fix_deck` pipeline
    `deckguard create` uses, so layout selection and brand compliance
    are guaranteed the same way regardless of which mode produced it.
    Slides with real content this project's rules already consider
    unsafe to reinterpret (a table, chart, embedded object, or more
    content than any layout can hold) are never sent to the model in
    any mode; they're reported back in `RedesignResult.skipped`.

    `mode`: "rewrite" (default) is everything described above -- the AI
    edits/condenses wording and picks a kind. "brand" is a completely
    different, fully deterministic path with no LLM call and no API key
    needed at all: `retemplate.apply_rebrand` carries every eligible
    slide's text/images over VERBATIM, picks layouts with an anti-repeat
    tie-break for visual variety, and swaps a confidently-detected cover/
    closing slide onto the current brand layout -- see that function's
    own docstring. `brief`/`target_slides`/`model`/`effort`/`notes`/
    `api_key`/`client` are all rewrite-mode-only and are rejected if
    `mode="brand"`; `deck_path` is required in brand mode since there's
    no content to rework without one.

    `review` (mode='brand' only): adds one small, optional AI call --
    see `call_claude_for_brand_review`'s own comment for exactly what
    it does (and deliberately doesn't do). Needs an API key, unlike the
    rest of brand mode; `api_key`/`client` are reused for it the same
    way rewrite mode uses them.

    `reference_path` (mode='brand' only, the "Learn from a reference"
    flow): an already-on-brand deck to copy EXACT per-shape styling from,
    on top of everything else brand mode already does -- see
    `_rebrand_deck`'s own comment. No API key needed; never persisted to
    brand_rules.yaml, purely a per-call source for this one run.
    """
    if out_path is None:
        raise RedesignError("out_path is required")

    if mode not in ("rewrite", "brand"):
        raise RedesignError(f"mode must be 'rewrite' or 'brand', got {mode!r}")

    if mode == "brand":
        if deck_path is None:
            raise RedesignError("mode='brand' needs a source deck (it never authors content, only re-lays-out what's there)")
        if brief is not None:
            raise RedesignError("mode='brand' never authors content, so --brief doesn't apply -- use mode='rewrite' instead")
        return _rebrand_deck(
            deck_path, out_path, template_path=template_path, rules_config=rules_config,
            review=review, review_model=review_model, api_key=api_key, client=client,
            reference_path=reference_path,
        )
    if review:
        raise RedesignError("--review only applies to mode='brand' -- mode='rewrite' already sends every eligible slide to the model")
    if reference_path is not None:
        raise RedesignError("reference_path only applies to mode='brand' -- mode='rewrite' doesn't have a matching concept")

    eligible: list = []
    blank_indices: list = []
    skipped: list = []
    if deck_path is not None:
        prs = Presentation(str(deck_path))
        eligible, all_skipped = extract_eligible_slides(prs)
        blank_indices, skipped = partition_skipped(all_skipped)

    if not eligible and not blank_indices and not brief:
        raise RedesignError(
            "nothing to work with: the deck has no eligible content to redesign and no --brief was given "
            "to generate from (pass --brief to fill blank slides, or to build a new deck with no source at all)"
        )

    effective_target = target_slides
    if effective_target is None and (eligible or blank_indices):
        effective_target = len(eligible) + len(blank_indices)

    raw_slides, usage = call_claude_for_outline(
        eligible,
        blank_count=len(blank_indices),
        brief=brief,
        target_slides=effective_target,
        model=model,
        effort=effort,
        notes=notes,
        api_key=api_key,
        client=client,
    )
    raw_slides = _attach_source_images(raw_slides, eligible)
    outline = outline_from_list(_strip_ai_only_fields(raw_slides))

    compose_result = build_deck(outline, out_path, template_path=template_path, rules_config=rules_config)
    redesign_result = RedesignResult(outline=outline, skipped=skipped, usage=usage)
    return compose_result, redesign_result
