"""Content-preserving re-layout: migrate an outdated deck's slides onto
the org template's own sanctioned layouts, via the slide master.

Unlike `fix` (recolors/refonts shapes in place) and `migrate` (replaces
only the first/last slide), this replaces a slide's entire structure --
useful for a genuinely old deck built on ad hoc, non-template layouts,
where the goal isn't to patch colors/fonts but to rebuild each slide on
an approved layout and carry the content over.

Deliberately scoped to TEXT and IMAGES only (titles, body/bullet text,
pictures) -- the content types that map onto a new layout's placeholders
reliably. A slide with a table, chart, embedded object, media, or a
grouped shape is never guessed at: it's left completely untouched and
flagged, rather than silently dropping or mangling content. Matching a
slide to a layout is deterministic (shape-count/type profile match, no
AI), same as the rest of this engine.

Two-phase by design (`propose_retemplate` / `apply_retemplate`) so a
caller -- the web UI in particular -- can show the proposed
slide->layout mapping and let a human accept or override it before
anything is generated, rather than committing to a one-shot guess.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckguard import effects as effects_mod
from deckguard.slide_import import (
    _delete_slide,
    _move_slide,
    default_template_path,
    import_layouts,
)

# Shapes types that make a slide's content unsafe to reflow automatically
# -- tables/charts/embedded objects/media don't map onto a placeholder,
# and a group can hide arbitrary complexity behind a single shape.
DISQUALIFYING_SHAPE_TYPES = {"TABLE", "CHART", "EMBEDDED_OLE_OBJECT", "MEDIA", "GROUP"}

# Readable overrides for the generic "contains a <type>" skip-reason
# phrasing below -- OLE is an acronym (stays uppercase), and the
# indefinite article needs to be "an" before it.
_SHAPE_TYPE_PHRASES = {"EMBEDDED_OLE_OBJECT": "an embedded OLE object"}


def _disqualifying_shape_phrase(type_name: str) -> str:
    if type_name in _SHAPE_TYPE_PHRASES:
        return _SHAPE_TYPE_PHRASES[type_name]
    label = type_name.lower().replace("_", " ")
    article = "an" if label[:1] in "aeiou" else "a"
    return f"{article} {label}"

TITLE_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}
BODY_PLACEHOLDER_TYPES = {"BODY", "OBJECT", "SUBTITLE"}

MAX_DECORATIVE_SHAPES = 4  # e.g. divider lines/accent rectangles with no text
MAX_TEXT_BLOCKS = 3
MAX_IMAGES = 4
MAX_PREVIEW_CHARS = 120

# classify_slide's ineligibility reason for a genuinely blank slide --
# distinct from every OTHER ineligibility reason (table/chart/media/
# overfull), which mean "this has real content we can't safely
# reinterpret." A blank slide has nothing to protect; redesign.py uses
# this exact string to tell "empty, safe to author from a brief" apart
# from "has content, must never be touched or guessed at."
EMPTY_SLIDE_REASON = "no title, text, or images to migrate"

# The org template's "ordinary content" layouts -- deliberately excludes
# Cover/Section/Agenda/Statement/Quote/Outro/End/REPORT-*/Blank, which are
# special-purpose or structural, not sensible targets for an arbitrary old
# body slide. Listed simplest-first: ties in `match_layout` favor whichever
# candidate appears earliest, so ordering here is a real tie-break, not
# just documentation.
CONTENT_LAYOUT_CANDIDATES = [
    "Title (24pt) and footers",
    "Title and text",
    "Title and content A",
    "Title and content B",
    "Two content A",
    "Two content B",
    "Two content C",
    "Two content D",
    "Three content A",
    "Three content B",
    "Three content C",
    "Three content D",
    "Text and picture A",
    "Text and picture B",
    "Text and picture C",
    "Text and picture F",
    "Text and picture G",
    "Text and picture H",
    "Two pictures and text A",
    "Two pictures and text B",
    "Two pictures and text C",
    "Three pictures and text",
    "Four pictures and text",
    "Fullslide picture",
]


@dataclass
class SlideProfile:
    """What a slide is made of, extracted independent of any target layout."""

    title: Optional[str]
    text_blocks: list  # list[list[tuple[int, str]]] -- one list per non-title text shape, each a (level, text) per paragraph
    images: list  # list[bytes] -- raw image blobs, in shape order
    eligible: bool
    reason: Optional[str] = None


def _is_footer_like(shape, slide_height_in: Optional[float]) -> bool:
    """Small text box tucked in the top/bottom margin -- a footer,
    confidentiality line, date, or page number, not real slide content.

    Mirrors the pagination heuristic rules_engine.py already uses for
    alignment exceptions, broadened beyond numeric-only text: a
    "Confidential | (c) KONE Corporation" footer isn't numeric, but is
    exactly as much chrome as a lone page number is, and was showing up
    as a bogus "body text block" (crowding out real content against
    MAX_TEXT_BLOCKS, and as a misleading preview) before this existed.
    """
    try:
        width_in = shape.width.inches if shape.width is not None else None
        height_in = shape.height.inches if shape.height is not None else None
        top_in = shape.top.inches if shape.top is not None else None
    except AttributeError:
        return False
    small_box = (width_in if width_in is not None else 99) < 4.5 and (height_in if height_in is not None else 99) < 0.6
    if not small_box or top_in is None or not slide_height_in:
        return False
    return top_in < 0.6 or top_in > slide_height_in - 0.75


def _extract_slide_content(slide, slide_height_in: Optional[float] = None):
    """Walk a slide's shapes once and pull out title/text-blocks/images.

    Returns (title, text_blocks, images, disqualify_reason). A non-None
    `disqualify_reason` means the slide has content that's unsafe to
    reinterpret in ANY context -- a table/chart/embedded-object/media/
    group shape, or too many free-form decorative shapes to trust a
    reflow -- and callers should treat the slide as ineligible outright,
    ignoring the other three (empty/undefined) return values.

    Deliberately does NOT apply a cap on text-block or image *count* --
    that's caller-side policy, not a property of the shapes themselves.
    `classify_slide` below applies retemplate's own cap (the org
    template's actual per-layout placeholder capacity, since retemplate
    carries content over VERBATIM); a caller willing to condense or
    rewrite content, like redesign.py, can reasonably choose a higher
    one against this exact same extraction.
    """
    title = None
    text_blocks: list = []
    images: list = []
    decorative = 0

    for shape in slide.shapes:
        try:
            type_name = shape.shape_type.name if shape.shape_type is not None else "UNKNOWN"
        except AttributeError:
            type_name = "UNKNOWN"

        if type_name in DISQUALIFYING_SHAPE_TYPES:
            return None, [], [], f"contains {_disqualifying_shape_phrase(type_name)}"

        ph_type_name = None
        if getattr(shape, "is_placeholder", False):
            try:
                ph_type = shape.placeholder_format.type
                ph_type_name = ph_type.name if ph_type else None
            except (AttributeError, ValueError):
                ph_type_name = None

        has_text = bool(getattr(shape, "has_text_frame", False)) and shape.text_frame.text.strip()

        if ph_type_name in TITLE_PLACEHOLDER_TYPES and has_text and title is None:
            title = shape.text_frame.text.strip()
            continue

        if type_name == MSO_SHAPE_TYPE.PICTURE.name:
            try:
                images.append(shape.image.blob)
            except Exception:  # noqa: BLE001 -- unreadable/corrupt embedded image
                decorative += 1
            continue

        if has_text:
            if _is_footer_like(shape, slide_height_in):
                continue  # boilerplate chrome, not real content -- doesn't count toward eligibility either
            paras = [(p.level, "".join(r.text for r in p.runs)) for p in shape.text_frame.paragraphs]
            paras = [(lvl, t) for lvl, t in paras if t.strip()]
            if paras:
                text_blocks.append(paras)
            continue

        decorative += 1
        if decorative > MAX_DECORATIVE_SHAPES:
            return None, [], [], "too many free-form shapes to safely reflow"

    return title, text_blocks, images, None


def classify_slide(slide, slide_height_in: Optional[float] = None) -> SlideProfile:
    """Extract a slide's title/body-text/image content and decide
    whether it's safe to auto-map onto a new layout at all, under
    retemplate's own (verbatim-carryover) capacity rules."""
    title, text_blocks, images, reason = _extract_slide_content(slide, slide_height_in)
    if reason:
        return SlideProfile(None, [], [], False, reason)
    if title is None and not text_blocks and not images:
        return SlideProfile(None, [], [], False, EMPTY_SLIDE_REASON)
    # Unlike the two branches above (genuinely nothing usable was ever
    # extracted), title/text_blocks/images ARE preserved here even though
    # the slide is ineligible -- a caller reporting *why* a slide was
    # skipped (e.g. a title/body preview for a human, or brand mode's
    # --review judgment call) needs to see what's actually on it, not a
    # blanked-out profile. Nothing downstream mistakes this for
    # "safe to transplant": every caller gates on `.eligible` first.
    if len(text_blocks) > MAX_TEXT_BLOCKS:
        return SlideProfile(title, text_blocks, images, False, "more body text blocks than any template layout can hold")
    if len(images) > MAX_IMAGES:
        return SlideProfile(title, text_blocks, images, False, "more images than any template layout can hold")

    return SlideProfile(title=title, text_blocks=text_blocks, images=images, eligible=True)


@dataclass
class LayoutProfile:
    name: str
    has_title: bool
    n_body: int
    n_picture: int


def _layout_profile(layout) -> LayoutProfile:
    has_title = False
    n_body = 0
    n_picture = 0
    for ph in layout.placeholders:
        ph_type = ph.placeholder_format.type
        name = ph_type.name if ph_type else None
        if name in TITLE_PLACEHOLDER_TYPES:
            has_title = True
        elif name in BODY_PLACEHOLDER_TYPES:
            n_body += 1
        elif name == "PICTURE":
            n_picture += 1
    return LayoutProfile(name=layout.name, has_title=has_title, n_body=n_body, n_picture=n_picture)


@dataclass
class LayoutMatch:
    layout_name: str
    score: int


def match_layout(
    profile: SlideProfile, layout_profiles: list[LayoutProfile], usage_counts: Optional[dict] = None
) -> Optional[LayoutMatch]:
    """Pick the tightest-fitting candidate layout that can hold every
    piece of this slide's content -- never one that would force dropping
    a text block or image. None if nothing qualifies.

    `usage_counts` (optional, `{layout_name: times already used}`) only
    ever breaks a TIE between equally-good fits -- it can never make a
    worse-fitting layout win over a better one. Passed by `apply_rebrand`
    so a deck's slides don't all land on the same "obvious" layout; every
    other caller omits it and gets the exact fit-only behavior this
    function has always had (every count reads as 0, so ties resolve by
    candidate-list order same as before).
    """
    title_needed = profile.title is not None
    n_body_needed = len(profile.text_blocks)
    n_pic_needed = len(profile.images)

    best: Optional[tuple[tuple[int, int], str]] = None
    for lp in layout_profiles:
        if title_needed and not lp.has_title:
            continue
        if lp.n_body < n_body_needed or lp.n_picture < n_pic_needed:
            continue
        waste = (lp.n_body - n_body_needed) + (lp.n_picture - n_pic_needed)
        if lp.has_title and not title_needed:
            waste += 1
        used = usage_counts.get(lp.name, 0) if usage_counts else 0
        score = (-waste, -used)
        if best is None or score > best[0]:
            best = (score, lp.name)

    if best is None:
        return None
    return LayoutMatch(layout_name=best[1], score=best[0][0])


@dataclass
class SlideProposal:
    slide_index: int  # 1-based
    eligible: bool
    reason: Optional[str] = None
    layout_name: Optional[str] = None
    title_preview: Optional[str] = None
    body_preview: Optional[str] = None
    image_count: int = 0


def _preview(text: Optional[str]) -> Optional[str]:
    if not text:
        return None
    text = " ".join(text.split())
    return text if len(text) <= MAX_PREVIEW_CHARS else text[: MAX_PREVIEW_CHARS - 1] + "…"


def _candidate_layout_profiles(template_path) -> list[LayoutProfile]:
    tmpl_prs = Presentation(str(template_path))
    layout_by_name = {layout.name: layout for master in tmpl_prs.slide_masters for layout in master.slide_layouts}
    return [
        _layout_profile(layout_by_name[name])
        for name in CONTENT_LAYOUT_CANDIDATES
        if name in layout_by_name
    ]


def propose_retemplate(prs, template_path=None) -> list[SlideProposal]:
    """Classify every slide and propose a target layout for each, without
    changing anything. Safe to call repeatedly / for preview purposes."""
    template_path = template_path or default_template_path()
    layout_profiles = _candidate_layout_profiles(template_path)
    slide_height_in = prs.slide_height / 914400 if prs.slide_height else None

    proposals = []
    for i, slide in enumerate(prs.slides, start=1):
        profile = classify_slide(slide, slide_height_in)
        if not profile.eligible:
            body_preview = _preview(profile.text_blocks[0][0][1]) if profile.text_blocks and profile.text_blocks[0] else None
            proposals.append(
                SlideProposal(
                    slide_index=i, eligible=False, reason=profile.reason,
                    title_preview=_preview(profile.title), body_preview=body_preview, image_count=len(profile.images),
                )
            )
            continue
        match = match_layout(profile, layout_profiles)
        if match is None:
            proposals.append(
                SlideProposal(
                    slide_index=i, eligible=False, reason="no template layout fits this slide's content",
                    image_count=len(profile.images),
                )
            )
            continue
        body_preview = _preview(profile.text_blocks[0][0][1]) if profile.text_blocks and profile.text_blocks[0] else None
        proposals.append(
            SlideProposal(
                slide_index=i, eligible=True, layout_name=match.layout_name,
                title_preview=_preview(profile.title), body_preview=body_preview,
                image_count=len(profile.images),
            )
        )
    return proposals


def _set_text_block(text_frame, block: list) -> None:
    text_frame.clear()
    first = True
    for level, text in block:
        p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        p.text = text
        p.level = max(0, min(level, 8))


_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _layout_placeholder_image_blob(layout, idx: int) -> Optional[bytes]:
    """The image bytes behind a layout's own picture placeholder at
    `idx`, if it has one baked in via `<a:blipFill><a:blip>` -- e.g. the
    org template's default cover-photo stock image, which any slide-
    level placeholder of the same idx shows through inheritance as long
    as its own `<p:spPr>` stays empty. Used to materialize a REAL,
    independently editable picture on the slide when the source slide
    had none of its own to carry over (see `_transplant_content`):
    PowerPoint only offers "Change Picture" for an actual picture object
    belonging to the slide itself, never for one it's merely inheriting
    from its layout -- an unfilled placeholder left that way only ever
    offers "Save as Picture" on the rendered (layout-owned) image, same
    limitation the raw org template itself has for an untouched Cover/
    Outro/End slide.
    """
    for ph in layout.placeholders:
        if ph.placeholder_format.idx != idx:
            continue
        blip = ph._element.find(".//" + effects_mod.a_qn("blip"))
        if blip is None:
            return None
        rid = blip.get(f"{{{_R_NS}}}embed")
        if not rid:
            return None
        try:
            return layout.part.related_part(rid).blob
        except KeyError:
            return None
    return None


def _first_picture_blob(slide) -> Optional[bytes]:
    """The image bytes of the first real picture (plain or placeholder)
    found directly on `slide`, if any -- used to borrow a reference
    deck's own cover/end photo (see `apply_rebrand`'s `reference_path`)
    ahead of falling back to the org template's generic stock photo."""
    for shape in slide.shapes:
        blip = shape._element.find(".//" + effects_mod.a_qn("blip"))
        if blip is None:
            continue
        rid = blip.get(f"{{{_R_NS}}}embed")
        if not rid:
            continue
        try:
            return shape.part.related_part(rid).blob
        except KeyError:
            continue
    return None


def _transplant_content(
    new_slide, profile: SlideProfile, fallback_to_layout_default: bool = False,
    reference_image_blob: Optional[bytes] = None,
) -> None:
    body_iter = iter(profile.text_blocks)
    image_iter = iter(profile.images)
    layout = new_slide.slide_layout if fallback_to_layout_default else None
    for ph in new_slide.placeholders:
        ph_type = ph.placeholder_format.type
        name = ph_type.name if ph_type else None
        if name in TITLE_PLACEHOLDER_TYPES:
            if profile.title:
                ph.text_frame.text = profile.title
        elif name in BODY_PLACEHOLDER_TYPES:
            block = next(body_iter, None)
            if block:
                _set_text_block(ph.text_frame, block)
        elif name == "PICTURE":
            blob = next(image_iter, None)
            if blob is not None:
                ph.insert_picture(BytesIO(blob))
            elif reference_image_blob is not None:
                # No source image of the deck's own to carry over, but a
                # reference deck was given for this run (see the "Learn
                # from a reference" flow's per-run image borrow) and it
                # has its own photo in this same cover/end role -- use
                # THAT ahead of the org template's generic stock photo,
                # since it's the deck-specific answer, not a placeholder.
                ph.insert_picture(BytesIO(reference_image_blob))
            elif layout is not None:
                # No source image on this slide to carry over -- give
                # the picture placeholder the layout's own default photo
                # as a real, editable picture instead of leaving it
                # empty (see `_layout_placeholder_image_blob`). Scoped to
                # the cover/end swap (see call site) -- an ordinary
                # content slide with no image of its own should just
                # have no picture, not get padded out with a stock photo.
                default_blob = _layout_placeholder_image_blob(layout, ph.placeholder_format.idx)
                if default_blob is not None:
                    ph.insert_picture(BytesIO(default_blob))


@dataclass
class RetemplateResult:
    proposals: list  # list[SlideProposal], every slide (eligible or not)
    transformed: list  # slide_index (1-based, ORIGINAL numbering) actually rebuilt
    skipped: list  # slide_index left untouched (ineligible, or not accepted)

    def to_dict(self) -> dict:
        return {
            "proposals": [vars(p) for p in self.proposals],
            "transformed": self.transformed,
            "skipped": self.skipped,
        }


def apply_retemplate(
    deck_path, out_path, accepted_indexes: Optional[set] = None, template_path=None
) -> RetemplateResult:
    """Rebuild every ACCEPTED, eligible slide on its matched template
    layout, carrying its title/body-text/images over; every other slide
    (ineligible, or simply not accepted) is left byte-for-byte untouched.

    `accepted_indexes`: 1-based original slide indices to actually
    transform. None (the default) means "every eligible slide" -- the
    fully-automatic path. A caller doing guided review passes exactly the
    indices a human approved; anything else is filtered out even if it
    was eligible.
    """
    template_path = template_path or default_template_path()
    prs = Presentation(deck_path)
    proposals = propose_retemplate(prs, template_path)
    proposal_by_index = {p.slide_index: p for p in proposals}
    eligible_indexes = {p.slide_index for p in proposals if p.eligible}

    accepted = eligible_indexes if accepted_indexes is None else (set(accepted_indexes) & eligible_indexes)
    all_indexes = set(range(1, len(prs.slides) + 1))
    result = RetemplateResult(proposals=proposals, transformed=sorted(accepted), skipped=sorted(all_indexes - accepted))

    if not accepted:
        import shutil

        shutil.copy(deck_path, out_path)
        return result

    layout_by_index = {i: proposal_by_index[i].layout_name for i in accepted}
    _rebuild_accepted_slides(deck_path, out_path, template_path, accepted, layout_by_index)
    return result


def _resolve_imported_layouts(prs, partname_by_name: dict) -> dict:
    """{layout_name: SlideLayout} for every entry in `partname_by_name`
    (as returned by `import_layouts`) resolved against `prs`'s CURRENT
    part graph -- robust regardless of which/how-many slide masters the
    import(s) landed on (`import_layouts` may be called more than once,
    from different source decks, in the same rebuild -- see
    `reference_layout_by_index` below)."""
    all_layouts = {
        str(layout.part.partname).lstrip("/"): layout
        for master in prs.slide_masters for layout in master.slide_layouts
    }
    return {name: all_layouts[partname] for name, partname in partname_by_name.items()}


def _freeze_placeholder_geometry(slide) -> None:
    """Snapshot each of `slide`'s placeholders' CURRENT effective
    position (its own explicit xfrm if it has one, else whatever it
    currently inherits from its layout) as an explicit xfrm on the
    placeholder itself.

    Called before `_reparent_slide_layout` swaps which layout the slide
    inherits from, because a REAL slide routinely has a MIX: one
    placeholder was never touched (no xfrm of its own, purely inherited)
    while another was manually adjusted at some point (has its own
    xfrm) -- both perfectly fine under the OLD layout, where the two
    happened to line up without overlapping. If the reference deck's
    version of a same-NAMED layout was itself redesigned (not just
    re-branded) and moved that inherited placeholder, only the
    inheriting one would silently follow -- producing an internally
    INCONSISTENT slide where the moved placeholder now overlaps the one
    that stayed put. This is exactly what surfaced on a real slide: an
    inherited title jumped to the reference layout's own (differently
    positioned) title placeholder while an explicitly-positioned body
    placeholder stayed at the old slide's position, and the two
    overlapped into an unreadable mess. Freezing every placeholder's
    geometry first keeps the whole slide internally consistent with
    itself, unaffected by whatever the new layout's own placeholders
    are positioned at -- this only ever touches geometry, never the
    chrome (background art, logo, footer format) the reparent is for.
    """
    for ph in slide.placeholders:
        spPr = effects_mod.get_spPr(ph)
        if spPr is not None and spPr.find(effects_mod.a_qn("xfrm")) is not None:
            continue  # already has its own explicit position
        left, top, width, height = ph.left, ph.top, ph.width, ph.height
        if left is None or top is None or width is None or height is None:
            continue  # nothing resolvable to freeze
        ph.left, ph.top, ph.width, ph.height = left, top, width, height


def _reparent_slide_layout(slide, new_layout) -> None:
    """Re-point `slide`'s single required slideLayout relationship at
    `new_layout`, leaving every one of the slide's own shapes' CONTENT
    untouched (their geometry is frozen first -- see
    `_freeze_placeholder_geometry`).

    Used for the "Learn from a reference" flow's layout carryover (see
    `apply_rebrand`'s `reference_path`): when an old slide and its
    reference counterpart already sit on a layout of the SAME NAME (most
    often a deck-specific custom layout with no equivalent in the org
    template at all -- e.g. a big internal catalog deck's own one-off
    layout), the right fix isn't a content reflow onto something else,
    it's just pointing the slide at a fresh, on-brand copy of the exact
    layout it already used -- carrying over whatever chrome (logo,
    footer/date format) lives on that layout's own master, without
    touching the slide's body content or position at all.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    _freeze_placeholder_geometry(slide)

    part = slide.part
    old_rid = next(rid for rid, rel in part.rels.items() if rel.reltype == RT.SLIDE_LAYOUT)
    part.drop_rel(old_rid)
    part.relate_to(new_layout.part, RT.SLIDE_LAYOUT)


def _rebuild_accepted_slides(
    deck_path, out_path, template_path, accepted: set, layout_by_index: dict,
    profile_overrides: Optional[dict] = None, reference_image_by_index: Optional[dict] = None,
    reference_path=None, reference_layout_by_index: Optional[dict] = None,
) -> None:
    """Shared slide-surgery for `apply_retemplate` and `apply_rebrand`:
    rebuild every slide index in `accepted` on the layout named in
    `layout_by_index[index]`, carrying its content over verbatim, in the
    original slide's own position -- every other slide in the deck is
    left byte-for-byte untouched. Saves to `out_path` itself.

    `profile_overrides` (optional, `{index: SlideProfile}`): for an
    index present here, this SlideProfile is used verbatim instead of
    the one `classify_slide` would derive from that slide's own
    content -- the one case that needs this today is rebuilding a
    divider-like slide onto a Section Divider layout with a short
    AI-suggested title, where the ORIGINAL content is deliberately not
    what should end up on the new slide.

    `reference_image_by_index` (optional, `{index: bytes}`): see
    `apply_rebrand`'s `reference_path` -- passed straight through to
    `_transplant_content` for the one index it applies to.

    `reference_layout_by_index` (optional, `{index: layout_name}`, disjoint
    from `accepted`): indices to re-parent onto a layout of that name
    IMPORTED FROM `reference_path` instead of `template_path` -- see
    `_reparent_slide_layout`. These slides' own content is never touched.
    """
    profile_overrides = profile_overrides or {}
    reference_image_by_index = reference_image_by_index or {}
    reference_layout_by_index = reference_layout_by_index or {}
    needed_layouts = sorted(set(layout_by_index[i] for i in accepted))
    tmp_path = Path(out_path).with_suffix(".layouts.pptx")
    partname_by_name = import_layouts(deck_path, str(template_path), str(tmp_path), needed_layouts) if needed_layouts else {}

    reference_partname_by_name: dict = {}
    if reference_layout_by_index:
        needed_reference_layouts = sorted(set(reference_layout_by_index.values()))
        source_path = str(tmp_path) if needed_layouts else deck_path
        reference_partname_by_name = import_layouts(
            source_path, str(reference_path), str(tmp_path), needed_reference_layouts
        )

    try:
        prs2 = Presentation(str(tmp_path))
        layout_by_name = _resolve_imported_layouts(prs2, partname_by_name)
        reference_layout_objs = _resolve_imported_layouts(prs2, reference_partname_by_name)
        slide_height_in = prs2.slide_height / 914400 if prs2.slide_height else None

        # Captured once, before any add/move/delete -- a slide's OPC
        # partname never changes, so this is a stable way to keep finding
        # "the original slide at position N" as the deck's slide list
        # shifts underneath us with every replacement.
        original_partnames = [str(s.part.partname).lstrip("/") for s in prs2.slides]

        # Reference-layout carryover: a straight relationship swap, no
        # content surgery at all, so it's safe to do before the
        # add/move/delete dance below (touches a disjoint set of slides).
        for orig_idx, layout_name in reference_layout_by_index.items():
            old_partname = original_partnames[orig_idx - 1]
            old_slide = next(s for s in prs2.slides if str(s.part.partname).lstrip("/") == old_partname)
            _reparent_slide_layout(old_slide, reference_layout_objs[layout_name])

        # Phase A: create every new slide and transplant its content FIRST,
        # with no deletions interleaved. python-pptx's own new-slide
        # partname isn't a scan for a free name -- it's just "current
        # slide count + 1" (PresentationPart._next_slide_partname) -- so a
        # delete between two add_slide() calls shrinks that count and can
        # make the next add_slide() hand out a partname still in use by an
        # earlier new slide from this same run, silently colliding two
        # live slides onto one XML part (confirmed by direct repro: a
        # second replacement reused the first's "slide4.xml"). Doing every
        # add first, before any delete, keeps the count strictly growing
        # for the whole phase, so that collision can never trigger.
        new_partname_by_old: dict[str, str] = {}
        for orig_idx in sorted(accepted):
            old_partname = original_partnames[orig_idx - 1]
            old_slide = next(s for s in prs2.slides if str(s.part.partname).lstrip("/") == old_partname)
            if orig_idx in profile_overrides:
                profile = profile_overrides[orig_idx]
            else:
                profile = classify_slide(old_slide, slide_height_in)  # re-derive from the still-intact original (bytes are unchanged from `prs`)
            layout = layout_by_name[layout_by_index[orig_idx]]

            new_slide = prs2.slides.add_slide(layout)
            is_cover_or_end_layout = layout_by_index[orig_idx] in (REBRAND_COVER_LAYOUT, REBRAND_END_LAYOUT)
            _transplant_content(
                new_slide, profile, fallback_to_layout_default=is_cover_or_end_layout,
                reference_image_blob=reference_image_by_index.get(orig_idx),
            )
            new_partname_by_old[old_partname] = str(new_slide.part.partname).lstrip("/")

        # Phase B: now that every new slide safely exists, move each into
        # its old slide's position and remove the old slide -- one pair at
        # a time, re-resolving both positions by partname each step since
        # they shift as earlier pairs are processed.
        for orig_idx in sorted(accepted):
            old_partname = original_partnames[orig_idx - 1]
            new_partname = new_partname_by_old[old_partname]
            old_pos = next(i for i, s in enumerate(prs2.slides) if str(s.part.partname).lstrip("/") == old_partname)
            new_pos = next(i for i, s in enumerate(prs2.slides) if str(s.part.partname).lstrip("/") == new_partname)
            _move_slide(prs2, new_pos, old_pos)
            _delete_slide(prs2, old_pos + 1)

        prs2.save(out_path)
    finally:
        tmp_path.unlink(missing_ok=True)


def rebuild_slides_as_dividers(deck_path, out_path, template_path, title_by_index: dict) -> list:
    """Rebuild specific slides (already chosen by the caller -- this
    does no eligibility checking of its own) onto the org template's
    Section Divider layout, bearing only a short title, discarding
    whatever the original slide's own content was. Alternates between
    "Section divider A" and "Section divider B" for a little visual
    variety across multiple dividers in one deck, same anti-repeat
    spirit as `apply_rebrand`'s own layout matching.

    `title_by_index`: `{1-based slide index: title text}`. Returns the
    sorted list of indices actually rebuilt (every key, since this
    function doesn't second-guess the caller's choice of which slides
    are divider-like -- that judgment already happened upstream).
    """
    variants = ["Section divider A", "Section divider B"]
    layout_by_index = {idx: variants[i % len(variants)] for i, idx in enumerate(sorted(title_by_index))}
    profile_overrides = {
        idx: SlideProfile(title=title, text_blocks=[], images=[], eligible=True)
        for idx, title in title_by_index.items()
    }
    accepted = set(title_by_index)
    _rebuild_accepted_slides(deck_path, out_path, template_path, accepted, layout_by_index, profile_overrides)
    return sorted(accepted)


# The org template's current title/closing layouts -- what a confidently
# detected cover or closing slide gets force-swapped onto in
# apply_rebrand, in place of whatever ordinary-content layout would
# otherwise fit it. Same names `compose.py` uses for a from-scratch
# cover/end slide (its default cover variant, and its only end layout),
# kept as plain literals here rather than imported to avoid a circular
# import (compose.py already imports FROM this module).
REBRAND_COVER_LAYOUT = "Cover B"
REBRAND_END_LAYOUT = "Outro"


def _looks_like_cover_or_end_slide(profile: SlideProfile) -> bool:
    """A confident-only heuristic for a genuine title/closing slide, as
    opposed to an ordinary content slide that merely happens to be
    sparse: a real cover/closing slide has a title and, at most, one
    short supporting line -- nothing more. Anything busier is real
    content and must never be guessed at, which is why this only ever
    overrides which layout an ALREADY-eligible slide lands on (see its
    one call site in `apply_rebrand`) -- it never changes eligibility
    itself."""
    return profile.title is not None and len(profile.text_blocks) <= 1 and len(profile.images) <= 1


@dataclass
class RebrandResult:
    proposals: list  # list[SlideProposal], every slide (eligible or not)
    transformed: list  # slide_index (1-based, ORIGINAL numbering) actually rebuilt
    skipped: list  # slide_index left untouched (ineligible, or no layout fit)
    manual_review: list  # fix_deck's own post-pass findings it couldn't auto-resolve (e.g. all_caps)
    reference_layout_indices: list = field(default_factory=list)  # subset of `transformed` handled by reference layout carryover, not org-template rebuild


def apply_rebrand(
    deck_path, out_path, template_path=None, rules_config: Optional[dict] = None,
    reference_path: Optional[str] = None,
) -> RebrandResult:
    """A fully deterministic alternative to `redesign`'s AI content-rewrite
    mode, for a deck whose wording is fine as written and just needs to
    land on brand: every eligible slide's title/body text/images are
    carried over VERBATIM -- same eligibility rules `apply_retemplate`
    already uses (nothing here condenses, so there's no reason to accept
    more than a layout can literally hold, unlike `redesign`'s permissive
    caps). Matched to a layout with an anti-repeat tie-break
    (`match_layout`'s `usage_counts`) so the deck doesn't read as one
    layout stamped out slide after slide, and the deck's own cover/
    closing slide -- when confidently identifiable as one, see
    `_looks_like_cover_or_end_slide` -- is swapped onto the org
    template's current `Cover B`/`Outro` layout instead of whatever
    ordinary-content layout would otherwise fit it (its own image, if it
    has one, is carried into that layout's own picture placeholder same
    as any other -- a normal, independently editable placeholder in the
    output file, not baked into anything). No LLM call, no API key, no
    wording changes anywhere. Color/font brand compliance is finished by
    running `fix_deck` over the whole result before returning -- the same
    engine/config `deckguard fix` uses -- so this is equivalent to
    `deckguard retemplate` immediately followed by `deckguard fix`, plus
    the cover/end swap and layout variety neither of those two do alone.

    `reference_path` (optional -- the "Learn from a reference" flow
    only): when the swapped-in cover/end slide has no picture of its own
    to carry over, its picture placeholder is normally padded out with
    the org template's own generic stock photo (see
    `_layout_placeholder_image_blob`). If a reference deck is given here,
    its own first/last slide's picture (if it has one) is tried FIRST --
    the deck-specific answer for THIS run beats the org template's
    generic default. Never persisted anywhere; purely a per-call source
    for this one picture.

    `reference_path` also drives a SECOND, independent mechanism for
    every slide EXCEPT the cover/end positions (which keep their own
    dedicated handling above): wherever an old slide and its reference
    counterpart at the SAME index already sit on a layout of the exact
    same name, that's ground truth for what this slide should look like
    -- often a deck-specific custom layout the org template has no
    equivalent for at all (a large internal catalog deck's own one-off
    layout, say), which `classify_slide`/`match_layout` below can never
    correctly handle since they only ever draw from the org template.
    These slides are re-parented onto a fresh copy of that exact layout
    IMPORTED FROM THE REFERENCE DECK (see `_reparent_slide_layout`) --
    content stays 100% untouched, only the layout/master (and therefore
    its chrome: logo, footer/date format) is refreshed. They're excluded
    from `classify_slide`/`match_layout` entirely -- there's no content
    reflow to evaluate eligibility for.
    """
    template_path = template_path or default_template_path()
    prs = Presentation(deck_path)
    layout_profiles = _candidate_layout_profiles(template_path)
    slide_height_in = prs.slide_height / 914400 if prs.slide_height else None

    tmpl_prs = Presentation(str(template_path))
    tmpl_layout_names = {layout.name for master in tmpl_prs.slide_masters for layout in master.slide_layouts}

    n_slides = len(prs.slides)

    reference_layout_by_index: dict = {}
    ref_prs = None
    if reference_path is not None:
        ref_prs = Presentation(str(reference_path))
        for i in range(2, n_slides):  # excludes slide 1 (cover) and n_slides (end)
            if i > len(ref_prs.slides):
                break
            old_layout_name = prs.slides[i - 1].slide_layout.name
            if old_layout_name == ref_prs.slides[i - 1].slide_layout.name:
                reference_layout_by_index[i] = old_layout_name

    proposals: list = []
    usage_counts: dict = {}
    cover_index: Optional[int] = None
    end_index: Optional[int] = None
    for i, slide in enumerate(prs.slides, start=1):
        if i in reference_layout_by_index:
            title, text_blocks, images, _reason = _extract_slide_content(slide, slide_height_in)
            body_preview = _preview(text_blocks[0][0][1]) if text_blocks and text_blocks[0] else None
            proposals.append(
                SlideProposal(
                    slide_index=i, eligible=True, layout_name=reference_layout_by_index[i],
                    title_preview=_preview(title), body_preview=body_preview, image_count=len(images),
                )
            )
            continue

        profile = classify_slide(slide, slide_height_in)
        if not profile.eligible:
            body_preview = _preview(profile.text_blocks[0][0][1]) if profile.text_blocks and profile.text_blocks[0] else None
            proposals.append(
                SlideProposal(
                    slide_index=i, eligible=False, reason=profile.reason,
                    title_preview=_preview(profile.title), body_preview=body_preview, image_count=len(profile.images),
                )
            )
            continue

        layout_name = None
        is_cover = i == 1 and REBRAND_COVER_LAYOUT in tmpl_layout_names
        is_end = i == n_slides and n_slides > 1 and REBRAND_END_LAYOUT in tmpl_layout_names
        if (is_cover or is_end) and _looks_like_cover_or_end_slide(profile):
            layout_name = REBRAND_COVER_LAYOUT if is_cover else REBRAND_END_LAYOUT
            if is_cover:
                cover_index = i
            else:
                end_index = i
        else:
            match = match_layout(profile, layout_profiles, usage_counts)
            if match is not None:
                layout_name = match.layout_name

        if layout_name is None:
            proposals.append(
                SlideProposal(
                    slide_index=i, eligible=False, reason="no template layout fits this slide's content",
                    image_count=len(profile.images),
                )
            )
            continue

        usage_counts[layout_name] = usage_counts.get(layout_name, 0) + 1
        body_preview = _preview(profile.text_blocks[0][0][1]) if profile.text_blocks and profile.text_blocks[0] else None
        proposals.append(
            SlideProposal(
                slide_index=i, eligible=True, layout_name=layout_name,
                title_preview=_preview(profile.title), body_preview=body_preview,
                image_count=len(profile.images),
            )
        )

    proposal_by_index = {p.slide_index: p for p in proposals}
    accepted = {
        p.slide_index for p in proposals if p.eligible and p.slide_index not in reference_layout_by_index
    }
    all_indexes = set(range(1, n_slides + 1))
    transformed = sorted(accepted | set(reference_layout_by_index))
    skipped = sorted(all_indexes - accepted - set(reference_layout_by_index))

    if not accepted and not reference_layout_by_index:
        import shutil

        shutil.copy(deck_path, out_path)
        return RebrandResult(
            proposals=proposals, transformed=transformed, skipped=skipped, manual_review=[],
            reference_layout_indices=[],
        )

    layout_by_index = {i: proposal_by_index[i].layout_name for i in accepted}

    reference_image_by_index: dict = {}
    if reference_path is not None:
        if ref_prs is None:
            ref_prs = Presentation(str(reference_path))
        if cover_index in accepted and len(ref_prs.slides) > 0:
            blob = _first_picture_blob(ref_prs.slides[0])
            if blob is not None:
                reference_image_by_index[cover_index] = blob
        if end_index in accepted and len(ref_prs.slides) > 0:
            blob = _first_picture_blob(ref_prs.slides[-1])
            if blob is not None:
                reference_image_by_index[end_index] = blob

    _rebuild_accepted_slides(
        deck_path, out_path, template_path, accepted, layout_by_index,
        reference_image_by_index=reference_image_by_index,
        reference_path=reference_path, reference_layout_by_index=reference_layout_by_index,
    )

    from deckguard.config import default_config_path, load_config
    from deckguard.fixer import fix_deck

    config = rules_config if rules_config is not None else load_config(default_config_path())
    out_prs = Presentation(str(out_path))
    fix_report = fix_deck(out_prs, config, source_path=str(out_path), output_path=str(out_path), dry_run=False)

    return RebrandResult(
        proposals=proposals, transformed=transformed, skipped=skipped, manual_review=fix_report.manual_review,
        reference_layout_indices=sorted(reference_layout_by_index),
    )
