"""Derive brand_rules.yaml updates by diffing an old deck against an
already-on-brand reference deck.

This formalizes, as a repeatable tool feature, the workflow used to grow
brand_rules.yaml from real decks throughout this project. Two passes,
in priority order:

1. Exact evidence from shapes that survive (same slide index + shape
   name) between the two decks — unambiguous ground truth where it's
   available. This matters: on a real deck pair, two old colors used at
   similar overall counts (e.g. 16 and 17 instances) were both plausible
   matches for two new colors at similar counts (15 and 17) -- whole-deck
   correlation alone picked the WRONG pairing (swapped), because nothing
   about aggregate counts alone could tell them apart. Exact shape-level
   correspondence resolves it unambiguously.
2. Whole-deck usage-count correlation for anything pass 1 has no exact
   evidence for (e.g. a color only used in text on a shape that didn't
   survive structurally between deck versions). Correlates which old
   value disappeared while a new one appeared at a similar count.
   Ambiguous/low-confidence matches are flagged rather than guessed —
   same "never destructive, prefer flagging over guessing" principle as
   the rest of the engine.

Not covered here: layout/structural differences. This only reasons about
color and font *usage*, not shape position/size, and pass 1 only helps
for shapes whose name happens to survive between deck revisions (see the
shape-name drift observed between real old/new deck pairs — many
survive, but not all).
"""

from __future__ import annotations

import copy
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable

from pptx.enum.shapes import MSO_SHAPE_TYPE
from ruamel.yaml import YAML

from deckguard import colors as colors_mod
from deckguard import effects as effects_mod
from deckguard.fonts import normalize_key
from deckguard.inventory import build_inventory, iter_shapes_recursive

HIGH_CONFIDENCE_SCORE = 0.7
CONFIDENCE_LEVELS = {"low": 0, "high": 1}
DEFAULT_LAYOUT_PANEL_MIN_AREA_SQ_IN = 8.0
# Two layout shapes are treated as "the same shape, recolored" only if their
# position+size differ by less than this (summed across left/top/width/
# height, in EMU) -- 1 inch per dimension. Prevents matching two unrelated
# large panels on a busy layout just because they're both above the area
# threshold.
_LAYOUT_PANEL_MATCH_TOLERANCE_EMU = 914400


@dataclass
class ColorProposal:
    role: str  # fill | line | text
    old_hex: str
    new_hex: str
    old_count: int
    new_count: int
    confidence: str  # high | low


@dataclass
class FontProposal:
    old_font: str
    old_bold: bool
    new_font: str
    new_bold: bool
    old_count: int
    new_count: int
    confidence: str


@dataclass
class LayoutPanelProposal:
    """A large background-panel fill that differs between the old and new
    decks' same-named slide layout, matched by position+size since the
    shape itself may have been renamed between deck revisions (e.g.
    'Rectangle 19' -> 'Rectangle 4')."""

    layout_name: str
    old_shape_name: str
    new_shape_name: str
    old_hex: str
    new_hex: str
    area_sq_in: float
    confidence: str  # high | low


@dataclass
class LearnResult:
    color_proposals: list[ColorProposal] = field(default_factory=list)
    font_proposals: list[FontProposal] = field(default_factory=list)
    layout_panel_proposals: list[LayoutPanelProposal] = field(default_factory=list)
    unmatched_old_colors: list[tuple[str, str, int]] = field(default_factory=list)  # role, hex, count
    unmatched_old_fonts: list[tuple[str, bool, int]] = field(default_factory=list)  # name, bold, count


def _color_usage(prs) -> Counter:
    inv = build_inventory(prs)
    usage: Counter = Counter()
    for slide in inv.slides:
        for shape in iter_shapes_recursive(slide.shapes):
            if shape.fill and shape.fill.type in ("solid", "gradient"):
                for c in shape.fill.colors:
                    if c.hex:
                        usage[("fill", c.hex)] += 1
            if shape.line and shape.line.color and shape.line.color.hex:
                usage[("line", shape.line.color.hex)] += 1
            for para in shape.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if run.color and run.color.hex:
                        usage[("text", run.color.hex)] += 1
    return usage


def _shape_fill_line_colors(prs) -> dict[tuple[int, str], dict[str, str]]:
    """{(slide_index, shape_name): {"fill": hex, "line": hex}} for every shape."""
    inv = build_inventory(prs)
    out: dict[tuple[int, str], dict[str, str]] = {}
    for slide in inv.slides:
        for shape in iter_shapes_recursive(slide.shapes):
            colors: dict[str, str] = {}
            if shape.fill and shape.fill.type == "solid":
                for c in shape.fill.colors:
                    if c.hex:
                        colors["fill"] = c.hex
            if shape.line and shape.line.color and shape.line.color.hex:
                colors["line"] = shape.line.color.hex
            if colors:
                out[(slide.index, shape.name)] = colors
    return out


def _shape_matched_color_correspondence(old_prs, new_prs) -> Counter:
    """Exact (role, old_hex) -> new_hex evidence from shapes present in both
    decks at the same (slide_index, shape_name) key.

    This is unambiguous ground truth where it's available, and is
    strictly more trustworthy than whole-deck usage-count correlation:
    when two old colors are both candidates for two new colors at similar
    overall counts, count correlation alone can pick the wrong pairing
    (this happened in practice — see learn.py's module docstring history).
    Falls back to count correlation (in `learn()`) only for colors that
    never appear on a name-matched shape.
    """
    old_map = _shape_fill_line_colors(old_prs)
    new_map = _shape_fill_line_colors(new_prs)
    correspondence: Counter = Counter()
    for key, old_colors in old_map.items():
        new_colors = new_map.get(key)
        if not new_colors:
            continue
        for role, old_hex in old_colors.items():
            new_hex = new_colors.get(role)
            if new_hex and new_hex != old_hex:
                correspondence[(role, old_hex, new_hex)] += 1
    return correspondence


def _iter_shapes_flat(shapes):
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_flat(shape.shapes)
        else:
            yield shape


def _layout_panel_shapes(prs, min_area_sq_in: float) -> dict[str, list[dict]]:
    """{layout_name: [{"shape_name", "left", "top", "width", "height", "hex"}, ...]}
    for shapes on the deck's slide layouts (first layout seen per name) with
    a literal solid fill and area >= min_area_sq_in.

    Raw-XML reads only (never touching python-pptx's higher-level fill/color
    properties) -- consistent with the rest of the codebase's non-mutating
    read pattern, since some ColorFormat properties have documented
    mutate-on-read side effects.
    """
    min_area_emu2 = min_area_sq_in * 914400 * 914400
    out: dict[str, list[dict]] = {}
    seen_names: set[str] = set()
    for master in colors_mod.iter_slide_masters(prs):
        for layout in master.slide_layouts:
            if layout.name in seen_names:
                continue
            seen_names.add(layout.name)
            shapes_info = []
            for shape in _iter_shapes_flat(layout.shapes):
                width, height = getattr(shape, "width", None), getattr(shape, "height", None)
                left, top = getattr(shape, "left", None), getattr(shape, "top", None)
                if not width or not height or left is None or top is None:
                    continue
                if width * height < min_area_emu2:
                    continue
                spPr = effects_mod.get_spPr(shape)
                if spPr is None:
                    continue
                solid_fill = spPr.find(effects_mod.a_qn("solidFill"))
                if solid_fill is None:
                    continue
                srgb = solid_fill.find(effects_mod.a_qn("srgbClr"))
                if srgb is None:
                    continue
                shapes_info.append(
                    {
                        "shape_name": shape.name,
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                        "hex": colors_mod.normalize_hex(srgb.get("val")),
                    }
                )
            if shapes_info:
                out[layout.name] = shapes_info
    return out


def _layout_panel_correspondence(
    old_prs, new_prs, min_area_sq_in: float = DEFAULT_LAYOUT_PANEL_MIN_AREA_SQ_IN
) -> list[LayoutPanelProposal]:
    """Match large background panels between old/new decks' same-named
    layouts by closest position+size (not by shape name, which can drift
    between deck revisions), and propose a fix for any whose fill differs.
    """
    old_layouts = _layout_panel_shapes(old_prs, min_area_sq_in)
    new_layouts = _layout_panel_shapes(new_prs, min_area_sq_in)

    proposals: list[LayoutPanelProposal] = []
    for layout_name, old_shapes in old_layouts.items():
        new_shapes = new_layouts.get(layout_name)
        if not new_shapes:
            continue
        used_new: set[int] = set()
        for old_shape in old_shapes:
            best_ni, best_dist = None, None
            for ni, new_shape in enumerate(new_shapes):
                if ni in used_new:
                    continue
                dist = (
                    abs(old_shape["left"] - new_shape["left"])
                    + abs(old_shape["top"] - new_shape["top"])
                    + abs(old_shape["width"] - new_shape["width"])
                    + abs(old_shape["height"] - new_shape["height"])
                )
                if best_dist is None or dist < best_dist:
                    best_dist, best_ni = dist, ni
            if best_ni is None or best_dist > _LAYOUT_PANEL_MATCH_TOLERANCE_EMU * 4:
                continue
            used_new.add(best_ni)
            new_shape = new_shapes[best_ni]
            if old_shape["hex"] == new_shape["hex"]:
                continue
            proposals.append(
                LayoutPanelProposal(
                    layout_name=layout_name,
                    old_shape_name=old_shape["shape_name"],
                    new_shape_name=new_shape["shape_name"],
                    old_hex=old_shape["hex"],
                    new_hex=new_shape["hex"],
                    area_sq_in=round((old_shape["width"] * old_shape["height"]) / (914400 * 914400), 2),
                    confidence="high" if best_dist <= _LAYOUT_PANEL_MATCH_TOLERANCE_EMU else "low",
                )
            )
    return proposals


def _font_usage(prs) -> Counter:
    inv = build_inventory(prs)
    usage: Counter = Counter()
    for slide in inv.slides:
        for shape in iter_shapes_recursive(slide.shapes):
            for para in shape.paragraphs:
                for run in para.runs:
                    if not run.text.strip() or run.font_raw is None:
                        continue
                    usage[(run.font_raw, run.bold)] += 1
    return usage


def _score(old_count: int, new_count: int) -> float:
    """1.0 = identical counts; decays toward 0 as they diverge."""
    lo, hi = sorted((old_count, new_count))
    return lo / hi if hi else 0.0


def _greedy_match(
    old_items: list[tuple], new_items: list[tuple], key_fn: Callable
) -> tuple[list[tuple[int, int, float]], list[tuple]]:
    """old_items/new_items: [(identity, count), ...]. Greedily pairs the
    highest-scoring (old, new) combinations first, restricted to pairs
    where key_fn(old_identity) == key_fn(new_identity) (role for colors,
    bold-flag for fonts). Returns (matches, unmatched_old_items)."""
    candidates = []
    for oi, (o_id, o_count) in enumerate(old_items):
        for ni, (n_id, n_count) in enumerate(new_items):
            if key_fn(o_id) != key_fn(n_id):
                continue
            candidates.append((_score(o_count, n_count), oi, ni))
    candidates.sort(key=lambda c: -c[0])

    used_old: set[int] = set()
    used_new: set[int] = set()
    matches = []
    for score, oi, ni in candidates:
        if oi in used_old or ni in used_new:
            continue
        used_old.add(oi)
        used_new.add(ni)
        matches.append((oi, ni, score))
    unmatched = [item for idx, item in enumerate(old_items) if idx not in used_old]
    return matches, unmatched


def learn(old_prs, new_prs, config: dict) -> LearnResult:
    colors_cfg = config.get("colors", {}) or {}
    approved = {colors_mod.normalize_hex(h) for h in colors_cfg.get("approved", []) or []}
    remap = {colors_mod.normalize_hex(k) for k in (colors_cfg.get("remap", {}) or {}).keys()}

    old_colors = _color_usage(old_prs)
    new_colors = _color_usage(new_prs)

    # Pass 1: exact evidence from shapes that survive (same slide index +
    # shape name) between the two decks. Unambiguous where available, so
    # it's resolved here directly rather than handed to the count-based
    # matcher below, which can be fooled when several old colors are all
    # plausible matches for several new colors at similar overall counts.
    exact_evidence = _shape_matched_color_correspondence(old_prs, new_prs)
    exact_target_votes: dict[tuple[str, str], Counter] = {}
    for (role, old_hex, new_hex), count in exact_evidence.items():
        exact_target_votes.setdefault((role, old_hex), Counter())[new_hex] = count

    color_proposals = []
    exactly_resolved: set[tuple[str, str]] = set()
    for (role, old_hex), votes in exact_target_votes.items():
        if old_hex in approved or old_hex in remap:
            continue
        new_hex, count = votes.most_common(1)[0]
        color_proposals.append(
            ColorProposal(
                role=role, old_hex=old_hex, new_hex=new_hex,
                old_count=count, new_count=count, confidence="high",
            )
        )
        exactly_resolved.add((role, old_hex))

    # Pass 2: whole-deck usage-count correlation, for colors with no exact
    # shape-level evidence (e.g. only used in text on shapes that didn't
    # survive structurally between deck versions).
    old_items = [
        ((role, hexval), count)
        for (role, hexval), count in old_colors.items()
        if hexval not in approved and hexval not in remap and (role, hexval) not in exactly_resolved
    ]
    new_items = [((role, hexval), count) for (role, hexval), count in new_colors.items()]

    matches, unmatched = _greedy_match(old_items, new_items, key_fn=lambda ident: ident[0])

    for oi, ni, score in matches:
        (role, old_hex), old_count = old_items[oi]
        (_, new_hex), new_count = new_items[ni]
        if old_hex == new_hex:
            continue
        color_proposals.append(
            ColorProposal(
                role=role,
                old_hex=old_hex,
                new_hex=new_hex,
                old_count=old_count,
                new_count=new_count,
                confidence="high" if score >= HIGH_CONFIDENCE_SCORE else "low",
            )
        )
    unmatched_old_colors = [(role, hexval, count) for (role, hexval), count in unmatched]

    fonts_cfg = config.get("fonts", {}) or {}
    font_approved_keys = {normalize_key(n) for n in fonts_cfg.get("approved", []) or []}
    font_remap_keys = {normalize_key(k) for k in (fonts_cfg.get("remap", {}) or {}).keys()}

    old_fonts = _font_usage(old_prs)
    new_fonts = _font_usage(new_prs)
    old_font_items = [
        ((name, bold), count)
        for (name, bold), count in old_fonts.items()
        if normalize_key(name) not in font_approved_keys and normalize_key(name) not in font_remap_keys
    ]
    new_font_items = [((name, bold), count) for (name, bold), count in new_fonts.items()]

    font_matches, unmatched_fonts = _greedy_match(old_font_items, new_font_items, key_fn=lambda ident: ident[1])
    font_proposals = []
    for oi, ni, score in font_matches:
        (old_name, old_bold), old_count = old_font_items[oi]
        (new_name, new_bold), new_count = new_font_items[ni]
        if normalize_key(old_name) == normalize_key(new_name):
            continue
        font_proposals.append(
            FontProposal(
                old_font=old_name,
                old_bold=old_bold,
                new_font=new_name,
                new_bold=new_bold,
                old_count=old_count,
                new_count=new_count,
                confidence="high" if score >= HIGH_CONFIDENCE_SCORE else "low",
            )
        )
    unmatched_old_fonts = [(name, bold, count) for (name, bold), count in unmatched_fonts]

    # Pass 3: large background panels defined on a slide LAYOUT rather than
    # any slide -- invisible to passes 1/2, which only ever look at shapes
    # actually placed on a slide (see fixer.py's layout_panel_remap, which
    # applies whatever this proposes).
    layout_panel_min_area = colors_cfg.get("layout_panel_min_area_sq_in", DEFAULT_LAYOUT_PANEL_MIN_AREA_SQ_IN)
    existing_layout_panel_remap = {
        colors_mod.normalize_hex(k) for k in (colors_cfg.get("layout_panel_remap", {}) or {}).keys()
    }
    layout_panel_proposals = [
        p
        for p in _layout_panel_correspondence(old_prs, new_prs, layout_panel_min_area)
        if p.old_hex not in existing_layout_panel_remap
    ]

    return LearnResult(
        color_proposals=sorted(color_proposals, key=lambda p: -p.old_count),
        font_proposals=sorted(font_proposals, key=lambda p: -p.old_count),
        layout_panel_proposals=sorted(layout_panel_proposals, key=lambda p: -p.area_sq_in),
        unmatched_old_colors=sorted(unmatched_old_colors, key=lambda t: -t[2]),
        unmatched_old_fonts=sorted(unmatched_old_fonts, key=lambda t: -t[2]),
    )


def apply_learned(config: dict, result: LearnResult, min_confidence: str = "high") -> dict:
    """Return a NEW config dict with proposals at/above min_confidence merged
    in. Never mutates the passed-in config. Later (higher-count) proposals
    never overwrite an earlier one already set in this same call -- for
    colors specifically, colors.remap is keyed by hex only (not role-aware),
    so if the same hex was proposed with two different targets under
    different roles, the higher-count one wins and the other is dropped."""
    config = copy.deepcopy(config)
    threshold = CONFIDENCE_LEVELS[min_confidence]

    colors_cfg = config.setdefault("colors", {})
    remap = colors_cfg.setdefault("remap", {})
    approved = colors_cfg.setdefault("approved", [])
    approved_set = {colors_mod.normalize_hex(h) for h in approved}

    for p in result.color_proposals:
        if CONFIDENCE_LEVELS[p.confidence] < threshold:
            continue
        remap.setdefault(f"#{p.old_hex}", f"#{p.new_hex}")
        if p.new_hex not in approved_set:
            approved.append(f"#{p.new_hex}")
            approved_set.add(p.new_hex)

    layout_panel_remap = colors_cfg.setdefault("layout_panel_remap", {})
    colors_cfg.setdefault("layout_panel_min_area_sq_in", DEFAULT_LAYOUT_PANEL_MIN_AREA_SQ_IN)

    for p in result.layout_panel_proposals:
        if CONFIDENCE_LEVELS[p.confidence] < threshold:
            continue
        layout_panel_remap.setdefault(f"#{p.old_hex}", f"#{p.new_hex}")
        if p.new_hex not in approved_set:
            approved.append(f"#{p.new_hex}")
            approved_set.add(p.new_hex)

    fonts_cfg = config.setdefault("fonts", {})
    font_remap = fonts_cfg.setdefault("remap", {})
    font_approved = fonts_cfg.setdefault("approved", [])
    font_approved_keys = {normalize_key(n) for n in font_approved}

    for p in result.font_proposals:
        if CONFIDENCE_LEVELS[p.confidence] < threshold:
            continue
        font_remap.setdefault(p.old_font, p.new_font)
        if normalize_key(p.new_font) not in font_approved_keys:
            font_approved.append(p.new_font)
            font_approved_keys.add(normalize_key(p.new_font))

    return config


def write_learned_to_yaml(path: str | Path, result: LearnResult, min_confidence: str = "high") -> int:
    """Apply proposals directly onto the YAML file at `path`, in place,
    preserving its existing comments/structure (via ruamel's round-trip
    mode) rather than round-tripping through a plain dict, which would
    flatten brand_rules.yaml's hand-curated comments and grouping.

    Returns the number of proposals actually applied (skips anything
    already present).
    """
    path = Path(path)
    yaml = YAML()
    yaml.preserve_quotes = True
    yaml.indent(mapping=2, sequence=4, offset=2)  # matches brand_rules.yaml's existing style
    with path.open("r", encoding="utf-8") as f:
        data = yaml.load(f)

    threshold = CONFIDENCE_LEVELS[min_confidence]
    applied = 0

    colors_cfg = data.setdefault("colors", {})
    remap = colors_cfg.setdefault("remap", {})
    approved = colors_cfg.setdefault("approved", [])
    approved_set = {colors_mod.normalize_hex(str(h)) for h in approved}
    remap_keys = {colors_mod.normalize_hex(str(k)) for k in remap.keys()}

    for p in result.color_proposals:
        if CONFIDENCE_LEVELS[p.confidence] < threshold:
            continue
        if p.old_hex in remap_keys:
            continue
        remap[f"#{p.old_hex}"] = f"#{p.new_hex}"
        remap_keys.add(p.old_hex)
        applied += 1
        if p.new_hex not in approved_set:
            approved.append(f"#{p.new_hex}")
            approved_set.add(p.new_hex)

    layout_panel_remap = colors_cfg.setdefault("layout_panel_remap", {})
    colors_cfg.setdefault("layout_panel_min_area_sq_in", DEFAULT_LAYOUT_PANEL_MIN_AREA_SQ_IN)
    layout_panel_remap_keys = {colors_mod.normalize_hex(str(k)) for k in layout_panel_remap.keys()}

    for p in result.layout_panel_proposals:
        if CONFIDENCE_LEVELS[p.confidence] < threshold:
            continue
        if p.old_hex in layout_panel_remap_keys:
            continue
        layout_panel_remap[f"#{p.old_hex}"] = f"#{p.new_hex}"
        layout_panel_remap_keys.add(p.old_hex)
        applied += 1
        if p.new_hex not in approved_set:
            approved.append(f"#{p.new_hex}")
            approved_set.add(p.new_hex)

    fonts_cfg = data.setdefault("fonts", {})
    font_remap = fonts_cfg.setdefault("remap", {})
    font_approved = fonts_cfg.setdefault("approved", [])
    font_approved_keys = {normalize_key(str(n)) for n in font_approved}
    font_remap_keys = {normalize_key(str(k)) for k in font_remap.keys()}

    for p in result.font_proposals:
        if CONFIDENCE_LEVELS[p.confidence] < threshold:
            continue
        if normalize_key(p.old_font) in font_remap_keys:
            continue
        font_remap[p.old_font] = p.new_font
        font_remap_keys.add(normalize_key(p.old_font))
        applied += 1
        if normalize_key(p.new_font) not in font_approved_keys:
            font_approved.append(p.new_font)
            font_approved_keys.add(normalize_key(p.new_font))

    if applied:
        with path.open("w", encoding="utf-8") as f:
            yaml.dump(data, f)

    return applied
