"""Logo detection (perceptual hash) and replacement.

Old logos are identified by perceptual hash rather than exact byte
match, since the same logo is routinely re-exported/re-compressed across
a legacy deck estate. Replacement swaps the embedded image part and
rewrites the `<a:blip r:embed>` reference only — the shape's `spPr`/
`xfrm` (position, size, rotation, crop) is left untouched.
"""

from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from functools import lru_cache
from io import BytesIO
from pathlib import Path
from typing import Optional

import imagehash
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE

DEFAULT_MATCH_THRESHOLD = 10  # max Hamming distance to count as a logo match

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _a(tag: str) -> str:
    return f"{{{A_NS}}}{tag}"


def _p(tag: str) -> str:
    return f"{{{P_NS}}}{tag}"


def _r_embed_attr() -> str:
    return f"{{{R_NS}}}embed"


def compute_phash(image_bytes: bytes) -> str:
    with Image.open(BytesIO(image_bytes)) as img:
        return str(imagehash.phash(img))


def hamming_distance(hash_a: str, hash_b: str) -> int:
    return imagehash.hex_to_hash(hash_a) - imagehash.hex_to_hash(hash_b)


def iter_picture_shapes(shapes):
    """Recurse into groups, yielding every picture shape."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_picture_shapes(shape.shapes)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def _walk_shapes(shapes):
    """Every shape, groups flattened -- and unlike `iter_picture_shapes`
    this filters nothing. A `<p:pic>` whose blip has lost its
    relationship cannot always be recognised by `shape_type`, so the
    caller matches on the XML tag instead and needs to see everything."""
    for shape in shapes:
        try:
            is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:  # noqa: BLE001 -- damaged shapes still need walking
            is_group = False
        if is_group:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


@dataclass
class LogoMatch:
    shape: object
    phash: str
    matched_hash: str
    distance: int


def find_old_logo_matches(
    shapes, old_hashes: list[str], threshold: int = DEFAULT_MATCH_THRESHOLD
) -> list[LogoMatch]:
    matches = []
    if not old_hashes:
        return matches
    for shape in iter_picture_shapes(shapes):
        try:
            phash = compute_phash(shape.image.blob)
        except Exception:  # noqa: BLE001 — unreadable/corrupt image, skip rather than crash
            continue
        best = min(old_hashes, key=lambda h: hamming_distance(phash, h))
        distance = hamming_distance(phash, best)
        if distance <= threshold:
            matches.append(LogoMatch(shape=shape, phash=phash, matched_hash=best, distance=distance))
    return matches


def replace_logo_image(shape, new_image_path: str) -> None:
    """Swap the embedded image, preserving the shape's position/size/crop."""
    image_part, rId = shape.part.get_or_add_image_part(new_image_path)
    blip = shape._element.blipFill.blip
    blip.rEmbed = rId


# PowerPoint 2016+ stores a vector image as a raster blip carrying an
# extension that points at the SVG. The URI is Microsoft's fixed
# identifier for that extension -- the master template uses this exact
# mechanism for 46 of its own image parts.
_SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
_SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"


def attach_svg(shape, svg_path) -> bool:
    """Give a picture shape a vector version of itself.

    A .pptx cannot hold an SVG on its own -- `add_picture` refuses one,
    because python-pptx identifies images through Pillow and Pillow does
    not read SVG. What PowerPoint actually does is keep BOTH: the
    `<a:blip>` points at a raster fallback, and an extension on that
    blip points at the SVG. Readers that understand the extension draw
    the vector and stay crisp at any zoom; readers that don't fall back
    to the PNG. Nothing is lost either way.

    So a KONE pictogram or logo goes in as the PNG we rasterised, plus
    this. Returns False if the part could not be added, leaving the
    raster picture perfectly usable.
    """
    from pathlib import Path

    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    svg_path = Path(svg_path)
    if not svg_path.is_file():
        return False
    part = shape.part
    package = part.package

    used = {p.partname for p in package.iter_parts()}
    for n in range(1, 10000):
        partname = PackURI(f"/ppt/media/dg-vector{n}.svg")
        if partname not in used:
            break
    else:  # pragma: no cover -- 10k vector parts in one deck
        return False

    try:
        svg_part = Part(partname, "image/svg+xml", package, svg_path.read_bytes())
        rId = part.relate_to(svg_part, RT.IMAGE)
        blip = shape._element.blipFill.blip
        ext_lst = blip.find(_a("extLst"))
        if ext_lst is None:
            ext_lst = etree.SubElement(blip, _a("extLst"))
        ext = etree.SubElement(ext_lst, _a("ext"))
        ext.set("uri", _SVG_EXT_URI)
        svg_blip = etree.SubElement(ext, f"{{{_SVG_NS}}}svgBlip")
        svg_blip.set(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed", rId
        )
    except Exception:  # noqa: BLE001 -- the raster picture is already valid
        return False
    return True


def find_background_blip(container_element):
    """Return the `<a:blip>` of a slide/layout/master's own page-level
    background-fill image (`<p:cSld><p:bg><p:bgPr><a:blipFill><a:blip>`),
    or None.

    A logo is sometimes not a picture *shape* at all -- it's baked into
    the page background fill (a banner/header image, or the whole page
    background) via `<p:bg>`, which lives outside the shape tree
    entirely. `iter_picture_shapes` (shape-tree only) can never see this,
    which is exactly why such a logo appears invisible to deckguard even
    though it's visibly present in the deck.
    """
    bg = container_element.find(f"{_p('cSld')}/{_p('bg')}")
    if bg is None:
        return None
    return bg.find(f".//{_a('blip')}")


def find_old_logo_background_match(
    part, container_element, old_hashes: list[str], threshold: int = DEFAULT_MATCH_THRESHOLD
) -> "LogoMatch | None":
    """Same idea as `find_old_logo_matches`, scoped to a page-level
    background-fill image rather than a shape."""
    if not old_hashes:
        return None
    blip = find_background_blip(container_element)
    if blip is None:
        return None
    rid = blip.get(_r_embed_attr())
    if not rid:
        return None
    try:
        image_part = part.related_part(rid)
        phash = compute_phash(image_part.blob)
    except Exception:  # noqa: BLE001 -- unreadable/corrupt/missing image, skip rather than crash
        return None
    best = min(old_hashes, key=lambda h: hamming_distance(phash, h))
    distance = hamming_distance(phash, best)
    if distance <= threshold:
        return LogoMatch(shape=blip, phash=phash, matched_hash=best, distance=distance)
    return None


def replace_background_image(part, blip_element, new_image_path: str) -> None:
    """Swap a page-level background-fill image found via `find_background_blip`."""
    image_part, rId = part.get_or_add_image_part(new_image_path)
    blip_element.set(_r_embed_attr(), rId)


def find_shapes_in_region(shapes, region_emu: tuple) -> list:
    """Top-level shapes (deliberately NOT recursing into groups -- a
    logo lockup built as a group is one match, not each of its individual
    paths) whose own bounding box is fully contained within
    `region_emu = (left, top, width, height)`.

    Exists for the case `find_old_logo_matches` can't handle at all: an
    old brand mark that isn't a raster picture, so there's no image to
    perceptual-hash -- e.g. a wordmark drawn as vector freeform shapes
    directly on a slide master. A logo's one reliable property across
    however differently it was constructed is that it sits in a fixed,
    small corner region of every slide -- so unlike hash matching (which
    identifies WHAT the old logo looks like), this identifies WHERE it
    lives, config-driven and opt-in (`logo.old_logo_region_in`) for the
    same reason `old_logo_hashes` is opt-in: guessing at removing
    shapes from a master, without a human confirming the region first,
    is exactly the kind of silent guess this project avoids everywhere
    else.
    """
    r_left, r_top, r_width, r_height = region_emu
    r_right, r_bottom = r_left + r_width, r_top + r_height
    matches = []
    for shape in shapes:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        if None in (left, top, width, height):
            continue
        if left >= r_left and top >= r_top and (left + width) <= r_right and (top + height) <= r_bottom:
            matches.append(shape)
    return matches


@lru_cache(maxsize=8)
def reference_logo_geometry(template_path) -> Optional[tuple]:
    """The current org template's own standard corner-logo size and
    position (left, top, width, height, in EMU) -- the ground truth for
    how big and where a REPLACEMENT logo should render, as opposed to
    however generously `old_logo_region_in` was drawn to reliably catch
    the old mark regardless of its own (possibly very different) size.
    Confirmed necessary the hard way: sizing a replacement to fill a
    deliberately generous search region rendered a logo several times
    too large next to every other layout's own actual logo.

    Computed as the majority (left, top, width, height) across every
    layout's own "Logo"/"Logo Placeholder"-named shape -- nearly all of
    them agree by design (cover/section/agenda/outro/end layouts are the
    outliers, using a different corner or a large centered mark, so
    letting the vote naturally favor the far more numerous "ordinary
    content layout" position is correct, not an arbitrary tiebreak).
    Returns None if the template can't be read or has no such shape at
    all -- callers should fall back to their own prior behavior.
    Cached (`lru_cache`) since this reparses a 60-layout template from
    scratch and `fix_deck`/`apply_rebrand` may call it many times in
    one process -- the bundled template's own geometry can't change
    mid-run.
    """
    from pptx import Presentation

    try:
        prs = Presentation(str(template_path))
    except Exception:  # noqa: BLE001 -- missing/corrupt template, not fatal to the caller
        return None

    sizes: Counter = Counter()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for shp in layout.shapes:
                name = (shp.name or "").lower()
                if "logo" in name and None not in (shp.left, shp.top, shp.width, shp.height):
                    sizes[(shp.left, shp.top, shp.width, shp.height)] += 1
    if not sizes:
        return None
    return sizes.most_common(1)[0][0]


def _next_shape_id_in_tree(spTree) -> int:
    """Same contract as python-pptx's own `_next_shape_id` (smallest
    positive integer not already used), but scoped to THIS shape tree's
    own `p:cNvPr/@id` values only -- confirmed necessary the hard way: a
    slide MASTER has a `p:sldLayoutIdLst` sibling (outside the shape
    tree, listing the layouts that belong to it) whose `p:sldLayoutId`
    elements use a completely different id namespace that starts at
    2^31 by OOXML convention. python-pptx's own `_next_shape_id`
    property does a document-WIDE `//@id` XPath scan (an absolute path,
    so it searches from the document root no matter which element it's
    called on) and picks those up too, handing back a shape id above
    2^31 -- which overflows a signed 32-bit int and produced a .pptx
    PowerPoint outright refused to open, confirmed by inspecting the
    output directly. Safe to call on a slide's shape tree (nothing else
    on a slide uses a bare, unprefixed `id` attribute); only actually
    matters for a master, but used everywhere here for one code path
    that's correct in both cases rather than two that silently diverge.
    """
    used_ids = {
        int(el.get("id")) for el in spTree.iter(_p("cNvPr"))
        if el.get("id") is not None and el.get("id").isdigit()
    }
    n = 1
    while n in used_ids:
        n += 1
    return n


def replace_shapes_in_region_with_logo(
    shape_container, matches: list, new_image_path: str, region_emu: tuple, target_emu: Optional[tuple] = None
) -> None:
    """Delete every shape in `matches` from `shape_container` (a
    master/layout/slide's own `.shapes`) and insert the new logo image
    in their place, sized to fit within `target_emu` (or `region_emu`
    when `target_emu` isn't given) preserving its native aspect ratio
    (never stretched), centered within that box.

    `region_emu` (the search region old shapes were matched against) is
    usually NOT the right size to render the new logo at -- it's drawn
    generously to reliably catch an old mark regardless of its own
    size, so sizing the replacement to fill it renders an
    oversized logo. Pass `target_emu` (see `reference_logo_geometry`)
    for the actual correct size/position to use; the `region_emu`
    fallback exists only for a caller with no better answer.
    """
    r_left, r_top, r_width, r_height = target_emu or region_emu
    spTree = shape_container._spTree
    for shape in matches:
        spTree.remove(shape._element)

    image_part, rId = shape_container.part.get_or_add_image_part(new_image_path)
    native_width, native_height = image_part.image.size
    scale = min(r_width / native_width, r_height / native_height)
    cx, cy = int(native_width * scale), int(native_height * scale)
    left = r_left + (r_width - cx) // 2
    top = r_top + (r_height - cy) // 2

    id_ = _next_shape_id_in_tree(spTree)
    name = "Picture %d" % id_
    spTree.add_pic(id_, name, image_part.desc, rId, left, top, cx, cy)


# --------------------------------------------------------------------------
# Repairing the master template's own empty logo frames
# --------------------------------------------------------------------------

_LOGO_ASSET_BY_NAME = {
    "logo": ("kone-logo.png", "kone-logo-white.png"),
    "tagline": ("kone-tagline.png", "kone-tagline-white.png"),
}


def _brand_asset(name: str, light: bool) -> Optional[str]:
    """Path to a vendored KONE mark: (standard, white) by shape name."""
    from pathlib import Path

    key = (name or "").strip().lower()
    pair = _LOGO_ASSET_BY_NAME.get(key)
    if pair is None:
        return None
    filename = pair[1] if light else pair[0]
    for base in (
        Path(__file__).parent / "assets" / "kone-design" / "logo",
        Path.home() / ".claude" / "skills" / "kone-design" / "assets" / "logo",
    ):
        candidate = base / filename
        if candidate.is_file():
            return str(candidate)
    return None


def _is_dark(hex_value: Optional[str]) -> bool:
    if not hex_value or len(hex_value) != 6:
        return False
    try:
        rgb = tuple(int(hex_value[i: i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return False
    return (0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]) < 140


def repair_empty_logo_frames(deck_path) -> int:
    """Give the master template's empty logo frames their image back.

    The KONE master ships 47 `Logo` and 7 `Tagline` picture shapes --
    across 49 of its 63 layouts -- whose `<a:blip>` carries no
    relationship at all. PowerPoint draws a picture frame with no
    picture as a dotted rectangle, so every deck built on that master
    shows dotted boxes where its logo should be. Reported from a slide.

    This is a defect in the template, not in anything deckguard writes,
    and it survived a full re-export of the master (`All Slides.pptx`
    has the same 54). The marks are rasterised and vendored here, so
    filling them in is a two-line repair rather than something the user
    has to fix by hand in 63 layouts.

    Returns how many frames were filled.
    """
    from pptx import Presentation

    prs = Presentation(str(deck_path))
    filled = 0
    containers = []
    for master in prs.slide_masters:
        containers += [master, *master.slide_layouts]
    # Slides too: a template that ships one instantiated slide per
    # layout (the current master does, 67 of them) carries the same
    # empty frames on the slides themselves, and those are what a user
    # actually opens.
    containers += list(prs.slides)

    for container in containers:
        background = _page_background_hex(container.element)
        for shape in _walk_shapes(container.shapes):
            if shape._element.tag != _p("pic"):
                continue
            try:
                shape.image.blob
                continue  # already has its picture
            except Exception:  # noqa: BLE001 -- the empty frames we're here for
                pass
            path = _brand_asset(shape.name, light=_is_dark(background))
            if path is None:
                continue
            try:
                replace_logo_image(shape, path)
                # the marks ship as SVG beside the PNG; prefer vector
                attach_svg(shape, Path(path).with_suffix(".svg"))
                filled += 1
            except Exception:  # noqa: BLE001 -- cosmetic repair, never fatal
                continue
    if filled:
        prs.save(str(deck_path))
    return filled


def _logo_kind(name: str) -> Optional[str]:
    """`Logo Placeholder 9` -> "logo", `Tagline` -> "tagline"."""
    key = (name or "").strip().lower()
    for kind in ("logo", "tagline"):
        if key.startswith(kind):
            return kind
    return None


def stamp_logo_chrome(slide) -> int:
    """Put the mark on a slide whose layout cannot supply it.

    The master carries its logo three different ways, and they do not
    behave the same when a slide is built programmatically:

    - 52 are `<p:pic>` shapes on the layout. A layout's non-placeholder
      shapes paint behind every slide using it, so these just work
      (once their images are restored -- see `repair_empty_logo_frames`).
    - 15 are `Logo Placeholder 9` / `Tagline Placeholder 9`, which are
      BODY placeholders. python-pptx never clones a layout placeholder
      the slide does not fill, so these render nothing at all -- the
      same latent-placeholder trap that loses the date and page number.
      They sit on Cover A, Cover C, Cover F, Fullslide picture and every
      Text-and-picture layout: the covers, in other words.
    - 2 are buried in a group and empty.

    So the first kind needs nothing and the other two need the mark
    stamped onto the slide itself, at the box the layout declares, in
    the variant the page background calls for. Returns how many marks
    were added; skips anything the slide or its layout already draws,
    because two logos on one slide is its own reported defect.
    """
    layout = slide.slide_layout

    already = {_logo_kind(s.name) for s in slide.shapes if _logo_kind(s.name)}
    for shape in layout.shapes:
        kind = _logo_kind(shape.name)
        if kind is None or shape.is_placeholder:
            continue
        try:  # a filled picture on the layout paints through by itself
            shape.image.blob
            already.add(kind)
        except Exception:  # noqa: BLE001 -- an empty frame supplies nothing
            pass

    # Name is not enough. An archetype that draws its own mark names it
    # whatever python-pptx auto-assigns -- `Picture 9` -- so matching on
    # the name alone stamps a second logo directly on top of the first.
    # What identifies a mark is where it is, so anything already sitting
    # in the box counts as that box being taken.
    occupied = [
        (s.left, s.top, s.width, s.height)
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and None not in
        (s.left, s.top, s.width, s.height)
    ]

    background = _page_background_hex(layout.element) or _page_background_hex(slide.element)
    added = 0
    for shape in layout.shapes:
        kind = _logo_kind(shape.name)
        if kind is None or kind in already:
            continue
        path = _brand_asset(kind, light=_is_dark(background))
        if path is None or not all(
            getattr(shape, side, None) for side in ("left", "top", "width", "height")
        ):
            continue
        if any(_boxes_overlap(box, (shape.left, shape.top, shape.width, shape.height))
               for box in occupied):
            already.add(kind)
            continue
        try:
            picture = _fit_picture(slide, path, shape.left, shape.top, shape.width, shape.height)
            attach_svg(picture, Path(path).with_suffix(".svg"))
            picture.name = shape.name.split(" Placeholder")[0]
        except Exception:  # noqa: BLE001 -- chrome, never fatal
            continue
        already.add(kind)
        added += 1
    return added


def drop_duplicate_logo_marks(slide) -> int:
    """Remove a slide-drawn mark that the layout already provides.

    `ARCHETYPES.md` settled this: the layout owns the chrome, an
    archetype declares which variant it needs and never places the mark
    itself. Several ported archetypes still draw their own, which was
    harmless while the master's frames were empty and became a visible
    double the moment they were repaired -- two logos stacked in the
    same 81x31 box on every divider and the outro.

    The layout's copy is the one that stays, because it is the one the
    spec says owns the slot. Returns how many were removed.
    """
    provided = []
    for shape in slide.slide_layout.shapes:
        if _logo_kind(shape.name) is None or shape.is_placeholder:
            continue
        try:
            shape.image.blob
        except Exception:  # noqa: BLE001 -- an empty frame provides nothing
            continue
        provided.append((shape.left, shape.top, shape.width, shape.height))

    removed = 0
    for shape in list(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        box = (shape.left, shape.top, shape.width, shape.height)
        if None in box:
            continue
        if any(_boxes_overlap(box, slot) for slot in provided):
            shape._element.getparent().remove(shape._element)
            removed += 1
    return removed


def restore_logo_chrome(deck_path) -> tuple[int, int]:
    """One mark per logo slot per slide, across a finished deck.

    Two failure modes meet here and they pull in opposite directions:
    a layout whose logo is a placeholder supplies nothing and needs one
    stamped, while an archetype that draws its own on a layout that
    already has one produces two. Both are resolved against the same
    invariant, in that order -- deduplicate first so a removal cannot
    leave a slot empty, then fill whatever is still bare.

    Returns (marks added, duplicates removed).
    """
    from pptx import Presentation

    prs = Presentation(str(deck_path))
    removed = sum(drop_duplicate_logo_marks(slide) for slide in prs.slides)
    added = sum(stamp_logo_chrome(slide) for slide in prs.slides)
    if added or removed:
        prs.save(str(deck_path))
    return added, removed


def _boxes_overlap(a, b, threshold: float = 0.5) -> bool:
    """Do two EMU boxes occupy substantially the same place?

    The ratio is taken against the LARGER box, which is the whole point.
    A 1280x422 cover photo completely contains the 81x31 logo slot, so
    measuring against the smaller box scores that pair 1.0 and reads the
    cover photo as a duplicate logo -- which deleted the banner off
    every cut cover until a test caught it. Against the larger box the
    same pair scores 0.008, while a mark sitting in its own slot scores
    about 1.
    """
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    overlap = (max(0, min(ax + aw, bx + bw) - max(ax, bx))
               * max(0, min(ay + ah, by + bh) - max(ay, by)))
    larger = max(aw * ah, bw * bh)
    return bool(larger) and overlap / larger >= threshold


def _fit_picture(slide, path, left, top, width, height):
    """Place an image inside a box without distorting it -- the marks
    are wider than they are tall and the declared boxes are not always
    the same ratio."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(width / iw, height / ih) if iw and ih else 1.0
    w, h = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(
        path, int(left + (width - w) / 2), int(top + (height - h) / 2), w, h
    )


def _page_background_hex(element) -> Optional[str]:
    """The solid `<p:bg>` colour of a slide/layout/master, if it has one."""
    import re

    match = re.search(r"<p:bg>.*?</p:bg>", element.xml, re.S)
    if not match:
        return None
    colour = re.search(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', match.group(0))
    return colour.group(1).upper() if colour else None
