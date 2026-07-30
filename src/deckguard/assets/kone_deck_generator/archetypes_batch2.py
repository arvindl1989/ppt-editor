"""
archetypes_batch2.py — eight more archetypes as DATA (from the uploaded KONE slides),
on the same engine. Builds a combined gallery with batch 1.
"""
import os
from pptx import Presentation
from pptx.oxml.ns import qn
import kone_engine as E
import archetypes_batch1 as B1

def photo(n): return os.path.join(E.PHOTO_DIR, n)

ARCHETYPES = {
 # slide 3 — text + stacked stats + full-height picture right
 "text_stats_picture_right":{
   "regions":[
     {"role":"picture","box":[725,0,555,720],"content":"image"},
     {"role":"title","box":[45,91,577,104],"content":"title"},
     {"role":"body_muted","box":[45,227,354,402],"content":"body"}],
   "groups":[{"content":"stats","origins":[[453,227],[453,320],[453,413]],
     "regions":[
       {"role":"stat_value_md","box":[0,0,251,44],"content":"value"},
       {"role":"caption","box":[0,46,251,44],"content":"desc"}]}]},

 # slide 10 — three picture cards
 "three_picture_cards":{
   "regions":[{"role":"title","box":[45,91,915,104],"content":"title"}],
   "groups":[{"content":"cards","origins":[[45,224],[453,224],[861,224]],
     "regions":[
       {"role":"picture","box":[0,0,374,213],"content":"image"},
       {"role":"heading","box":[0,247,374,40],"content":"heading"},
       {"role":"bullets","box":[0,297,374,170],"content":"bullets"}]}]},

 # slide 14 — how it works, 3 steps + header image
 "how_it_works_3step":{
   "regions":[
     {"role":"image_band","box":[0,0,1280,382],"content":"image"},
     {"role":"title","box":[45,435,272,194],"content":"title"}],
   "groups":[{"content":"steps","origins":[[351,435],[657,435],[963,435]],
     "regions":[
       {"role":"number","box":[0,0,272,60],"content":"number"},
       {"role":"body","box":[0,73,272,120],"content":"text"}]}]},

 # slide 15 — six numbered icon items in a row
 "numbered_icon_row_6":{
   "regions":[{"role":"title","box":[45,91,1189,104],"content":"title"}],
   "groups":[{"content":"items",
     "origins":[[45,443],[249,443],[452,443],[658,443],[862,443],[1065,443]],
     "regions":[
       {"role":"icon","box":[0,0,40,40]},
       {"role":"number","box":[0,48,170,50],"content":"number"},
       {"role":"body","box":[0,104,170,100],"content":"label"}]}]},

 # slide 13 — four-point value grid with picture row
 "four_point_value":{
   "regions":[
     {"role":"eyebrow","box":[45,43,585,32],"content":"eyebrow"},
     {"role":"title","box":[45,91,883,61],"content":"title"}],
   "groups":[
     {"content":"pictures","origins":[[45,181],[350,181],[657,181],[964,181]],
      "regions":[{"role":"picture","box":[0,0,272,216],"content":"image"}]},
     {"content":"points","origins":[[46,420],[352,420],[657,420],[963,420]],
      "regions":[
        {"role":"heading","box":[0,0,271,44],"content":"heading"},
        {"role":"body","box":[0,50,271,188],"content":"text"}]}]},

 # slide 6 — statement + three link columns + icon
 "statement_links":{
   "regions":[
     {"role":"statement","box":[46,92,985,155],"content":"statement"},
     {"role":"icon","box":[1140,96,90,90]}],
   "groups":[{"content":"columns","origins":[[46,473],[453,473],[861,473]],
     "regions":[
       {"role":"heading","box":[0,0,374,30],"content":"heading"},
       {"role":"bullets","box":[0,40,374,120],"content":"links"}]}]},

 # slide 2 — resource link tiles (2x2)
 "resource_links":{
   "regions":[
     {"role":"title","box":[28,48,600,45],"content":"title"},
     {"role":"body_muted","box":[49,654,500,23],"content":"contact"}],
   "groups":[{"content":"tiles","origins":[[46,230],[346,230],[46,348],[346,348]],
     "regions":[
       {"role":"icon","box":[2,0,28,28]},
       {"role":"body","box":[0,34,290,46],"content":"label"}]}]},

 # slide 12 — offer / price + CTA panel + header image
 "offer_cta":{
   "regions":[
     {"role":"image_band","box":[0,0,1280,460],"content":"image"},
     {"role":"title","box":[45,362,476,150],"content":"title"},
     {"role":"price","box":[45,496,476,60],"content":"price"},
     {"role":"caption","box":[45,558,476,72],"content":"fineprint"},
     {"role":"panel","box":[725,522,510,153]},
     {"role":"on_panel_heading","box":[759,543,350,30],"content":"cta_heading"},
     {"role":"on_panel_body","box":[759,578,350,80],"content":"cta_text"},
     {"role":"icon","box":[1133,573,73,73]}]},
}

SAMPLES = {
 "text_stats_picture_right":{"title":"Service that keeps buildings moving",
   "body":"Our maintenance model pairs skilled technicians with connected data, so issues are caught before they interrupt people flow.",
   "image":photo("technician-van-branded.jpg"),
   "stats":[{"value":"13 000","desc":"connected units under contract"},
            {"value":"50%","desc":"faster fault diagnosis"},
            {"value":"+34","desc":"markets served"}]},
 "three_picture_cards":{"title":"Three ways modernization pays off",
   "cards":[{"image":photo("elevator-women.jpg"),"heading":"Lower carbon footprint",
             "bullets":["Cut energy use","No full replacement needed"]},
            {"image":photo("stairs-bag.jpg"),"heading":"Keep tenants happy",
             "bullets":["Smoother rides","Less downtime"]},
            {"image":photo("elevator-bike.jpg"),"heading":"No new equipment",
             "bullets":["Reuse the shaft","Staged upgrades"]}]},
 "how_it_works_3step":{"image":photo("product-signalization.jpg"),
   "title":"KONE Destination — how it works",
   "steps":[{"number":"01","text":"Select your destination on the panel or swipe your access card."},
            {"number":"02","text":"Check which elevator is yours from the digital guides."},
            {"number":"03","text":"Enjoy a non-crowded ride with fewer intermediate stops."}]},
 "numbered_icon_row_6":{"title":"Faultless maintenance & repair",
   "items":[{"number":"01","label":"Tailored maintenance plan"},
            {"number":"02","label":"Professional service"},
            {"number":"03","label":"Skilled technicians"},
            {"number":"04","label":"Maintenance data online"},
            {"number":"05","label":"Preventive maintenance"},
            {"number":"06","label":"Regular equipment checks"}]},
 "four_point_value":{"eyebrow":"KONE Journey introduction","title":"KONE Journey supports building value",
   "pictures":[{"image":photo("elevator-women.jpg")},{"image":photo("stairs-phone.jpg")},
               {"image":photo("handrail-hands.jpg")},{"image":photo("elevator-bike.jpg")}],
   "points":[{"heading":"1. Premium impression","text":"Elevate building positioning from the first ride."},
             {"heading":"2. Optimized performance","text":"Maximize handling capacity for smooth people flow."},
             {"heading":"3. Built-in safety","text":"Integrate elevators with KONE access and security."},
             {"heading":"4. Future-proof value","text":"Protect long-term value with modernization."}]},
 "statement_links":{"statement":"The Customer Loyalty Survey helps us strengthen client relationships and improve service.",
   "columns":[{"heading":"Learn more","links":["Program overview","Why it matters"]},
              {"heading":"Sharepoint","links":["Customer Insight site","2025 results"]},
              {"heading":"Qlik dashboards","links":["Loyalty survey","Transaction analytics"]}]},
 "resource_links":{"title":"Find out more","contact":"customerinsight@kone.com",
   "tiles":[{"label":"Customer Loyalty Survey results"},{"label":"CLS results in Qlik"},
            {"label":"Transaction survey materials"},{"label":"Transaction analytics in Qlik"}]},
 "offer_cta":{"image":photo("stairs-phone.jpg"),"title":"Upgrade to KONE 24/7 Connected Services",
   "price":"1234€ / unit","fineprint":"Indicative price. Requirements may apply, e.g. an active KONE service contract.",
   "cta_heading":"Book a call","cta_text":"Schedule a call to find out how our solution fits your building."},
}

def build_full_gallery(out_path):
    A={**B1.ARCHETYPES, **ARCHETYPES}; S={**B1.SAMPLES, **SAMPLES}
    prs=Presentation(os.path.join(E.KONE_DESIGN,"uploads","master_ppt-1784774200983.pptx"))
    lst=prs.slides._sldIdLst
    for el in list(lst): prs.part.drop_rel(el.get(qn('r:id'))); lst.remove(el)
    blank=next(l for l in prs.slide_layouts if l.name.strip().lower()=="blank")
    for name,arch in A.items():
        s=prs.slides.add_slide(blank)
        cap=E._tf(s,[45,676,900,24]); E._run(cap.paragraphs[0],f"archetype: {name}",E.KINFO,11,E.GREY,caps=True)
        E.render_archetype(s, arch, S[name])
    prs.save(out_path); return out_path, list(A.keys())

if __name__=="__main__":
    path,names=build_full_gallery("/home/claude/KONE_Archetype_Gallery.pptx")
    print("saved", path, "with", len(names), "archetypes")
