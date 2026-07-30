"""
archetypes_batch1.py — first 6 archetypes as DATA, defined from the geometry extracted
from the uploaded KONE slides. Plus sample content and a gallery builder so each can be
eyeballed. Coordinates are px on the 1280 grid (straight from the source slides).
"""
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
import kone_engine as E

PH = E.PHOTO_DIR
def photo(name): return os.path.join(PH, name)

ARCHETYPES = {
 # --- slide 6-style stats (already proven; now as data) ---
 "three_stats": {
   "regions":[{"role":"title","box":[46,92,985,150],"content":"title"}],
   "groups":[{"content":"stats","origins":[[46,364],[453,364],[861,364]],
     "regions":[
       {"role":"stat_label","box":[0,0,374,24],"content":"label"},
       {"role":"stat_value","box":[0,30,374,90],"content":"value"},
       {"role":"caption","box":[0,128,374,130],"content":"desc"}]}]},

 # --- slide 4: 5 icon columns / pillars ---
 "icon_columns_5": {
   "regions":[
     {"role":"eyebrow","box":[45,45,917,30],"content":"eyebrow"},
     {"role":"title","box":[45,91,917,104],"content":"title"},
     {"role":"body_muted","box":[45,227,917,90],"content":"intro"}],
   "groups":[{"content":"items","origins":[[45,376],[291,376],[537,376],[783,376],[1031,376]],
     "regions":[
       {"role":"icon","box":[6,4,45,45]},
       {"role":"body","box":[0,64,203,150],"content":"text"}]}]},

 # --- slide 8: two picture columns compare ---
 "two_picture_compare": {
   "regions":[{"role":"title","box":[45,91,883,61],"content":"title"}],
   "groups":[{"content":"items","origins":[[45,181],[657,181]],
     "regions":[
       {"role":"picture","box":[0,0,272,448],"content":"image"},
       {"role":"heading","box":[307,0,271,44],"content":"heading"},
       {"role":"bullets","box":[307,54,271,394],"content":"bullets"}]}]},

 # --- slide 9: numbered summary + full-height picture right ---
 "numbered_summary_picture": {
   "regions":[
     {"role":"picture","box":[419,0,861,720],"content":"image"},
     {"role":"title","box":[45,91,374,104],"content":"title"}],
   "groups":[{"content":"points","origins":[[45,162],[45,277],[45,391]],
     "regions":[
       {"role":"number","box":[0,4,90,60],"content":"number"},
       {"role":"heading","box":[103,10,300,90],"content":"text"}]}]},

 # --- slide 7: 4-stage lifecycle with header image + axis ---
 "lifecycle_4stage": {
   "regions":[
     {"role":"image_band","box":[0,0,1280,382],"content":"image"},
     {"role":"eyebrow","box":[45,404,917,30],"content":"eyebrow"},
     {"role":"title","box":[45,432,883,90],"content":"title"},
     {"role":"axis","box":[47,548,1190,0]}],
   "groups":[{"content":"stages","origins":[[47,556],[351,556],[656,556],[965,556]],
     "regions":[
       {"role":"icon","box":[0,0,40,40]},
       {"role":"heading","box":[0,52,272,36],"content":"heading"},
       {"role":"bullets","box":[0,92,272,64],"content":"bullets"}]}]},

 # --- slide 11: quote with context (no panel — faithful to Quote A) ---
 "quote_context": {
   "regions":[
     {"role":"label","box":[45,136,272,104],"content":"label"},
     {"role":"body_muted","box":[45,272,272,357],"content":"context"},
     {"role":"quote","box":[510,212,657,349],"content":"quote"},
     {"role":"attribution","box":[510,561,657,30],"content":"attribution"}]},
}

SAMPLES = {
 "three_stats":{"title":"The Hub cleared most of what came in, even as demand climbed.",
   "stats":[{"label":"Requests in","value":"739","desc":"Across WCM, DEA and Graphic Design."},
            {"label":"Resolution rate","value":"91.2%","desc":"674 resolved in the focus period."},
            {"label":"YoY growth","value":"~2×","desc":"Inflow roughly doubled year on year."}]},
 "icon_columns_5":{"eyebrow":"Strategic shifts","title":"Five shifts shaping how we work",
   "intro":"Each shift moves the Hub from reactive delivery toward proactive, analytics-led service.",
   "items":[{"text":"Standardised intake across every service line."},
            {"text":"Automation for repeatable request types."},
            {"text":"Live dashboards for queue visibility."},
            {"text":"Predictive capacity planning."},
            {"text":"Self-serve reporting for stakeholders."}]},
 "two_picture_compare":{"title":"To summarize",
   "items":[{"image":photo("elevator-women.jpg"),"heading":"How to differentiate",
             "bullets":["Highest scores in ratings","Best sustainable innovation","Clearest people-flow story"]},
            {"image":photo("technician-van-branded.jpg"),"heading":"How to deliver",
             "bullets":["The best partners in the industry","Customer-language communication","Delivery focused on outcomes"]}]},
 "numbered_summary_picture":{"title":"Let’s summarize","image":photo("escalator-station.jpg"),
   "points":[{"number":"01","text":"We know our carbon footprint"},
             {"number":"02","text":"We have ambitious climate goals"},
             {"number":"03","text":"We have actions in place to reach them"}]},
 "lifecycle_4stage":{"image":photo("stairs-phone.jpg"),"eyebrow":"Inside out",
   "title":"Reducing footprint across the whole lifecycle",
   "stages":[{"heading":"Materials","bullets":["Less and lighter","Lower-carbon steel","Recycled content"]},
             {"heading":"Construction","bullets":["Optimized packaging","Optimized logistics"]},
             {"heading":"Use phase","bullets":["Energy-saving solutions","Regenerative drive"]},
             {"heading":"Maintenance","bullets":["Predictive maintenance","Fewer site visits"]}]},
 "quote_context":{"label":"Customer voice",
   "context":"We are setting new standards for smarter, more sustainable buildings with a ground-breaking approach.",
   "quote":"“Our goal was to increase energy efficiency while maintaining tenant comfort — a great project for us.”",
   "attribution":"Rene Klesment, CTO, Mainor Ülemiste"},
}

def build_gallery(out_path):
    prs=Presentation(os.path.join(E.KONE_DESIGN,"uploads","master_ppt-1784774200983.pptx"))
    lst=prs.slides._sldIdLst
    from pptx.oxml.ns import qn
    for el in list(lst):                       # start from a clean deck (drop master examples)
        prs.part.drop_rel(el.get(qn('r:id'))); lst.remove(el)
    blank=next(l for l in prs.slide_layouts if l.name.strip().lower()=="blank")
    for name,arch in ARCHETYPES.items():
        s=prs.slides.add_slide(blank)
        # tiny caption so we know which archetype we're looking at
        cap=E._tf(s,[45,676,900,24]); E._run(cap.paragraphs[0],f"archetype: {name}",E.KINFO,11,E.GREY,caps=True)
        E.render_archetype(s, arch, SAMPLES[name])
    prs.save(out_path); return out_path

if __name__=="__main__":
    print("saved", build_gallery("/home/claude/KONE_Archetype_Gallery.pptx"))
