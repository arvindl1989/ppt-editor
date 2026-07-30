"""
kone_deck_creator.py — assemble a finished deck from a spec, using the archetype library.
  master Cover F (kept) -> generated body archetypes -> master "Thank you" Outro (kept)

spec = {"title": "...", "slides": [{"archetype": "<name>", ...content...}, ...]}
Each slide names an archetype from archetypes.ARCHETYPES; the rest of the dict is its content.
"""
import os
from pptx import Presentation
from pptx.oxml.ns import qn
import kone_engine as E
import archetypes as A

MASTER = os.path.join(E.KONE_DESIGN, "uploads", "master_ppt-1784774200983.pptx")
INTRO_IDX, OUTRO_IDX = 1, 10          # Cover F, Outro "Thank you"

def _set_title(slide, text):
    tgt=next((sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.idx==0), None)
    if not tgt or not tgt.has_text_frame: return
    p=tgt.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text=text
        for r in p.runs[1:]: r.text=""
    else:
        p.add_run().text=text

def build_deck(spec, out_path):
    prs=Presentation(MASTER)
    lst=prs.slides._sldIdLst; originals=list(lst)
    intro, outro = originals[INTRO_IDX], originals[OUTRO_IDX]
    if spec.get("title"): _set_title(list(prs.slides)[INTRO_IDX], spec["title"])
    blank=next(l for l in prs.slide_layouts if l.name.strip().lower()=="blank")
    for s in spec["slides"]:
        name=s["archetype"]; content={k:v for k,v in s.items() if k!="archetype"}
        slide=prs.slides.add_slide(blank)
        A.render(slide, name, content)
    body=[el for el in list(lst) if el not in originals]
    keep={id(intro),id(outro),*(id(b) for b in body)}
    for el in originals:
        if id(el) not in keep: prs.part.drop_rel(el.get(qn('r:id'))); lst.remove(el)
    for el in list(lst): lst.remove(el)
    lst.append(intro); [lst.append(b) for b in body]; lst.append(outro)
    prs.save(out_path); return out_path

if __name__=="__main__":
    import json, sys
    spec=json.load(open(sys.argv[1]))
    print("built:", build_deck(spec, sys.argv[2]))
