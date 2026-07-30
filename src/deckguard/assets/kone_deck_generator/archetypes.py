"""
archetypes.py — the full library + slot-based real art.
Icons fill each archetype's icon slots in order (sized to the slot, never overlapping text);
charts/diagrams fill a dedicated 'figure' slot (contain-fit, never cropped). Backgrounds from
the source deck. Single render entry point for the deck creator and gallery.
"""
import os, json
import archetypes_batch1 as B1, archetypes_batch2 as B2, archetypes_batch3 as B3
import kone_engine as E

ARCHETYPES = {**B1.ARCHETYPES, **B2.ARCHETYPES, **B3.ARCHETYPES}
SAMPLES    = {**B1.SAMPLES,    **B2.SAMPLES,    **B3.SAMPLES}

_icondir = os.path.join(E.SKILL_DIR, "assets", "icons")
def _load(n):
    p=os.path.join(_icondir, n);  return json.load(open(p)) if os.path.exists(p) else {}
BG = _load("bg_map.json")

ICONS = {
 "icon_columns_5":  ["ic0.png","ic1.png","ic2.png","ic3.png","ic4.png"],
 "lifecycle_4stage":["lc0.png","lc1.png","lc2.png","lc3.png"],
 "resource_links":  ["rl0.png","rl1.png","rl2.png","rl3.png"],
 "segment_breakdown":["seg_ill0.png","seg_ill1.png","seg_ill2.png"],
 "statement_links": ["sl_icon.png"],
 "offer_cta":       ["oc_icon.png"],
}
FIGURES = {
 "segment_breakdown":{"chart":"seg_chart.png"},
 "chart_commentary": {"chart":"cc_chart.png"},
 "org_functions":    {"diagram":"org_diagram.png"},
}

def render(slide, name, content):
    c=dict(content)
    for key,fn in FIGURES.get(name,{}).items():
        c[key]=os.path.join(_icondir, fn)
    E.render_archetype(slide, ARCHETYPES[name], c, icons=ICONS.get(name), bg=BG.get(name))

def build_gallery(out_path):
    from pptx import Presentation
    from pptx.oxml.ns import qn
    prs=Presentation(os.path.join(E.KONE_DESIGN,"uploads","master_ppt-1784774200983.pptx"))
    lst=prs.slides._sldIdLst
    for el in list(lst): prs.part.drop_rel(el.get(qn('r:id'))); lst.remove(el)
    blank=next(l for l in prs.slide_layouts if l.name.strip().lower()=="blank")
    for name in ARCHETYPES:
        s=prs.slides.add_slide(blank)
        cap=E._tf(s,[45,690,900,20]); E._run(cap.paragraphs[0],f"archetype: {name}",E.KINFO,10,E.GREY,caps=True)
        render(s, name, SAMPLES[name])
    prs.save(out_path); return out_path, list(ARCHETYPES.keys())

if __name__=="__main__":
    p,names=build_gallery("/home/claude/KONE_Archetype_Gallery.pptx")
    print("built", len(names), "archetypes")
