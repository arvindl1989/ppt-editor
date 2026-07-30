"""
kone_engine.py — the declarative spine.
An archetype is DATA (a list of regions + repeating groups), not code. One renderer
interprets any archetype: each region has a ROLE (title, stat_value, caption, picture,
axis, ...) whose KONE styling comes from ROLE_STYLE, and a box in px on the 1280 grid.
Repeating groups (3 stats, 4 stages, 5 columns) expand over a content list.
Adding an archetype = adding a dict entry, not a function.
"""
import os
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image, ImageFont

SKILL_DIR   = os.path.dirname(os.path.abspath(__file__))
KONE_DESIGN = os.environ.get("KONE_DESIGN_DIR") or os.path.join(SKILL_DIR, os.pardir, "kone-design")
PHOTO_DIR   = os.path.join(KONE_DESIGN, "assets", "photos")
FONT_DIR    = os.path.join(SKILL_DIR, "fonts")
ICON_DIR    = os.path.join(SKILL_DIR, "assets", "icons")
def _hex(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

def X(px):  return Emu(int(round(px*9525)))
def PT(px): return Pt(px*0.75)
BLUE=RGBColor(0x14,0x50,0xF5); BLACK=RGBColor(0x14,0x14,0x14); WHITE=RGBColor(0xFF,0xFF,0xFF)
SAND=RGBColor(0xF3,0xEE,0xEA); GREY=RGBColor(0x72,0x72,0x72); HAIR=RGBColor(0xD0,0xD0,0xD0)
INTER="Inter"; INTER_SB="Inter SemiBold"; KINFO="KONE Information"
FONT_FILE={INTER:os.path.join(FONT_DIR,"Inter-Regular.ttf"),
           INTER_SB:os.path.join(FONT_DIR,"Inter-SemiBold.ttf"),
           KINFO:os.path.join(FONT_DIR,"KONE_Information.ttf")}

# role -> (font, px, color, caps, bold, fit_single_line)
ROLE_STYLE={
 "eyebrow":  (KINFO,12,BLUE,True,False,True),
 "subtitle": (KINFO,14,GREY,True,False,True),
 "title":    (INTER,40,BLACK,False,False,False),
 "title_lg": (INTER,54,BLACK,False,False,False),
 "statement":(INTER,40,BLACK,False,False,False),
 "body":     (INTER,16,BLACK,False,False,False),
 "body_muted":(INTER,16,GREY,False,False,False),
 "heading":  (INTER_SB,22,BLACK,False,True,False),
 "stat_value":(INTER,64,BLACK,False,False,True),
 "stat_value_md":(INTER,32,BLACK,False,False,True),
 "price":    (INTER,40,BLACK,False,False,True),
 "on_panel_heading":(INTER_SB,20,WHITE,False,True,False),
 "on_panel_body":(INTER,14,WHITE,False,False,False),
 "stat_label":(KINFO,14,BLUE,True,False,True),
 "caption":  (INTER,16,GREY,False,False,False),
 "number":   (INTER,40,BLUE,False,False,True),
 "quote":    (INTER,28,BLACK,False,False,False),
 "attribution":(KINFO,14,GREY,True,False,True),
 "label":    (INTER,22,BLACK,False,True,False),
 "hero_value":(INTER,120,BLACK,False,False,True),
 "title_light":(INTER,54,WHITE,False,False,False),
 "heading_light":(INTER_SB,20,WHITE,False,True,False),
 "eyebrow_light":(KINFO,14,WHITE,True,False,True),
}

def _fit_px(text,font,base_px,box_w,min_px=9):
    path=FONT_FILE.get(font)
    if not path or not os.path.exists(path) or not text.strip(): return base_px
    try:
        for s in range(int(base_px),min_px-1,-1):
            if ImageFont.truetype(path,s).getlength(text)<=box_w*0.97: return s
        return min_px
    except Exception: return base_px

def _tf(slide,box,anchor=MSO_ANCHOR.TOP):
    x,y,w,h=box; tb=slide.shapes.add_textbox(X(x),X(y),X(w),X(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor; tf.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for m in("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    return tf
def _run(p,text,font,px,color,caps=False,bold=False,fit_w=None):
    disp=str(text).upper() if caps else str(text)
    if fit_w: px=_fit_px(disp,font,px,fit_w)
    r=p.add_run(); r.text=disp; f=r.font
    f.name=font; f.size=PT(px); f.color.rgb=color; f.bold=bold; return r

def _rect(slide,box,color):
    x,y,w,h=box; s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,X(x),X(y),X(w),X(h))
    s.fill.solid(); s.fill.fore_color.rgb=color; s.line.fill.background(); s.shadow.inherit=False; return s
def _hline(slide,box):
    x,y,w,h=box; ln=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,X(x),X(y),X(x+w),X(y))
    ln.line.color.rgb=HAIR; ln.line.width=Pt(1); return ln
def _icon(slide,box):                       # placeholder icon chip (swap for real KONE icons)
    x,y,w,h=box; s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,X(x),X(y),X(w),X(h))
    s.fill.solid(); s.fill.fore_color.rgb=BLUE; s.line.fill.background(); s.shadow.inherit=False; return s
def _image(slide,box,path,mode="cover"):
    x,y,w,h=box
    if path and os.path.exists(path):
        im=Image.open(path).convert("RGB"); iw,ih=im.size
        if mode=="contain":                       # fit whole image inside box (charts, icons)
            sc=min(w/iw,h/ih); nw,nh=max(1,int(iw*sc)),max(1,int(ih*sc))
            tmp=f"/tmp/_img_{abs(hash((path,w,h,'c')))%99999}.png"
            im.resize((nw,nh),Image.LANCZOS).save(tmp)
            slide.shapes.add_picture(tmp,X(x+(w-nw)/2),X(y+(h-nh)/2),X(nw),X(nh))
        else:                                     # cover-crop to fill (photos)
            sc=max(w/iw,h/ih); im=im.resize((int(iw*sc),int(ih*sc)),Image.LANCZOS)
            l=(im.width-w)//2; t=(im.height-h)//2; tmp=f"/tmp/_img_{abs(hash((path,w,h,'v')))%99999}.jpg"
            im.crop((l,t,l+w,t+h)).save(tmp,quality=90); slide.shapes.add_picture(tmp,X(x),X(y),X(w),X(h))
    else:
        _rect(slide,box,SAND)

def _dash_bullets(slide,box,items):
    tf=_tf(slide,box)
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(6)
        _run(p,"—  ",INTER,16,BLUE); _run(p,it,INTER,16,BLACK)

def _table(slide,box,data):
    """simple KONE table: header row (blue text) + rows, hairline separators."""
    x,y,w,h=box; headers=data.get("headers",[]); rows=data.get("rows",[])
    ncol=max(1,len(headers)); cw=w/ncol; nrow=len(rows)+1; rh=h/nrow
    _hline(slide,[x,y,w,0])
    for c,htxt in enumerate(headers):
        tf=_tf(slide,[x+c*cw+6,y+8,cw-12,rh-8]); _run(tf.paragraphs[0],htxt,KINFO,14,BLUE,caps=True,fit_w=cw-12)
    for r,row in enumerate(rows,1):
        ry=y+r*rh; _hline(slide,[x,ry,w,0])
        for c,cell in enumerate(row):
            tf=_tf(slide,[x+c*cw+6,ry+8,cw-12,rh-8]); _run(tf.paragraphs[0],cell,INTER,16,BLACK)
    _hline(slide,[x,y+h,w,0])

def _region(slide,role,box,content):
    if role in ("picture","image_band"): _image(slide,box,content); return
    if role=="figure":      _image(slide,box,content,mode="contain"); return
    if role=="panel":       _rect(slide,box,BLUE); return
    if role=="panel_sand":  _rect(slide,box,SAND); return
    if role=="table":       _table(slide,box,content or {}); return
    if role=="axis":        _hline(slide,box); return
    if role=="icon":        _icon(slide,box); return
    if role=="bullets":     _dash_bullets(slide,box,content if isinstance(content,list) else [content]); return
    font,px,color,caps,bold,fit=ROLE_STYLE[role]
    anchor=MSO_ANCHOR.MIDDLE if role=="quote" else MSO_ANCHOR.TOP
    tf=_tf(slide,box,anchor); _run(tf.paragraphs[0],content,font,px,color,caps=caps,bold=bold,
                                   fit_w=(box[2] if fit else None))

_BG={"sand":SAND,"blue":BLUE,"white":WHITE}
def render_archetype(slide, arch, content, icons=None, bg=None):
    """arch: {'background'?, 'regions':[...], 'groups':[...]}; icons: [img,...] fill icon slots in order."""
    color = bg or arch.get("background")
    if color:
        rgb = _BG.get(color) or (_hex(color) if len(str(color))==6 else None)
        if rgb is not None and color!="FFFFFF": _rect(slide,[0,0,1280,720],rgb)
    iq=list(icons or [])
    def draw(role, box, val):
        if role=="icon":
            if iq:
                p=os.path.join(ICON_DIR, iq.pop(0))
                _image(slide,box,p,mode="contain") if os.path.exists(p) else _icon(slide,box)
            else:
                _icon(slide,box)
            return
        _region(slide, role, box, val)
    for reg in arch.get("regions",[]):
        key=reg.get("content")
        val=content.get(key) if key else reg.get("value","")
        if key and key not in content and "value" not in reg: continue
        draw(reg["role"], reg["box"], val)
    for grp in arch.get("groups",[]):
        items=content.get(grp["content"],[])
        for (ox,oy),item in zip(grp["origins"], items):
            for reg in grp["regions"]:
                rx,ry,rw,rh=reg["box"]
                cval=item.get(reg["content"]) if reg.get("content") else reg.get("value","")
                draw(reg["role"], [ox+rx,oy+ry,rw,rh], cval)
