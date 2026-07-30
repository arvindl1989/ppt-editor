"""
archetypes_batch3.py — the final nine: reusable cores for the three complex slides
(chart/diagram left as a picture slot to drop in), the sand quarterly plan, and the
five NEW archetypes from the master (agenda, comparison table, hero stat, image
divider, 2x2 matrix). Builds the complete gallery with batches 1 and 2.
"""
import os
from pptx import Presentation
from pptx.oxml.ns import qn
import kone_engine as E, archetypes_batch1 as B1, archetypes_batch2 as B2
def photo(n): return os.path.join(E.PHOTO_DIR, n)

ARCHETYPES = {
 # slide 1 — customer/segment breakdown (reusable core; waffle chart -> picture slot)
 "segment_breakdown":{
   "regions":[
     {"role":"title","box":[28,48,600,40],"content":"title"},
     {"role":"stat_value","box":[645,435,297,80],"content":"highlight_value"},
     {"role":"caption","box":[645,520,297,60],"content":"highlight_caption"},
     {"role":"figure","box":[946,150,290,240],"content":"chart"}],
   "groups":[{"content":"categories","origins":[[47,152],[250,152],[452,152]],
     "regions":[
       {"role":"icon","box":[0,0,100,100]},
       {"role":"heading","box":[0,112,190,30],"content":"heading"},
       {"role":"bullets","box":[0,150,190,320],"content":"items"}]}]},

 # slide 5 — chart with commentary columns (reusable core; chart -> picture slot)
 "chart_commentary":{
   "regions":[
     {"role":"eyebrow","box":[45,39,917,30],"content":"eyebrow"},
     {"role":"title","box":[45,91,917,61],"content":"title"},
     {"role":"figure","box":[40,300,380,330],"content":"chart"}],
   "groups":[{"content":"columns","origins":[[453,227],[861,227]],
     "regions":[
       {"role":"heading","box":[0,0,374,40],"content":"heading"},
       {"role":"bullets","box":[0,50,374,353],"content":"bullets"}]}]},

 # slide 16 — org / functions (reusable core; org diagram -> picture slot)
 "org_functions":{
   "regions":[
     {"role":"title","box":[45,91,1189,104],"content":"title"},
     {"role":"bullets","box":[45,240,360,400],"content":"functions"},
     {"role":"figure","box":[430,175,805,455],"content":"diagram"}]},

 # slide 17 (sand) — quarterly plan, 4 columns + 4 quarter blocks
 "quarterly_plan_4col":{
   "background":"sand",
   "regions":[
     {"role":"eyebrow","box":[45,39,917,32],"content":"eyebrow"},
     {"role":"title","box":[45,91,917,104],"content":"title"},
     {"role":"body_muted","box":[45,205,917,60],"content":"intro"}],
   "groups":[
     {"content":"columns","origins":[[45,300],[351,300],[657,300],[963,300]],
      "regions":[{"role":"body","box":[0,0,272,140],"content":"text"}]},
     {"content":"quarters","origins":[[45,470],[351,470],[657,470],[963,470]],
      "regions":[
        {"role":"stat_label","box":[0,0,272,26],"content":"label"},
        {"role":"bullets","box":[0,32,272,150],"content":"items"}]}]},

 # NEW-1 — agenda / contents
 "agenda_contents":{
   "regions":[{"role":"title","box":[45,91,917,104],"content":"title"}],
   "groups":[{"content":"items","origins":[[45,240],[45,310],[45,380],[45,450],[45,520]],
     "regions":[
       {"role":"number","box":[0,0,70,50],"content":"number"},
       {"role":"heading","box":[90,4,900,50],"content":"item"}]}]},

 # NEW-2 — comparison table
 "comparison_table":{
   "regions":[
     {"role":"title","box":[45,91,1189,104],"content":"title"},
     {"role":"table","box":[45,240,1190,380],"content":"table"}]},

 # NEW-3 — hero single stat
 "hero_stat":{
   "regions":[
     {"role":"eyebrow","box":[45,210,900,30],"content":"eyebrow"},
     {"role":"hero_value","box":[45,248,1190,210],"content":"value"},
     {"role":"heading","box":[45,470,900,44],"content":"caption"},
     {"role":"body_muted","box":[45,520,900,60],"content":"support"}]},

 # NEW-4 — image section divider (full-bleed image + overlaid white title)
 "image_section_divider":{
   "regions":[
     {"role":"image_band","box":[0,0,1280,720],"content":"image"},
     {"role":"eyebrow_light","box":[60,510,900,30],"content":"eyebrow"},
     {"role":"title_light","box":[60,545,1050,130],"content":"title"}]},

 # NEW-5 — 2x2 matrix / quadrants
 "matrix_2x2":{
   "regions":[
     {"role":"title","box":[45,91,900,60],"content":"title"},
     {"role":"eyebrow","box":[340,640,560,24],"content":"xlabel"},
     {"role":"eyebrow","box":[45,155,300,24],"content":"ylabel"}],
   "groups":[{"content":"quadrants",
     "origins":[[340,190],[792,190],[340,415],[792,415]],
     "regions":[
       {"role":"panel_sand","box":[0,0,440,215]},
       {"role":"heading","box":[20,18,400,36],"content":"heading"},
       {"role":"bullets","box":[20,60,400,140],"content":"items"}]}]},
}

SAMPLES = {
 "segment_breakdown":{"title":"Customer profile",
   "highlight_value":"53%","highlight_caption":"rate us easy to work with (2024)",
   "chart":None,
   "categories":[
     {"heading":"By segment","items":["44% Residential","37% Office & retail","15% Infrastructure"]},
     {"heading":"By role","items":["40% Owners","20% Management","10% Builders"]},
     {"heading":"By offering","items":["29% New building","59% Service","12% Modernization"]}]},
 "chart_commentary":{"eyebrow":"Survey results","title":"What the loyalty scores tell us",
   "chart":None,"highlight":"53%",
   "columns":[
     {"heading":"Strengths","bullets":["Responsiveness","Technician professionalism","Safety record"]},
     {"heading":"Areas to improve","bullets":["Digital self-service","Proactive updates","Speed of quoting"]}]},
 "org_functions":{"title":"How the Marketing Hub is organised","diagram":None,
   "functions":["Commercial & Operations","Technology & Innovation","Supply Chain",
                "Strategy & Transformation","People & Communications","Finance","Legal"]},
 "quarterly_plan_4col":{"eyebrow":"Core","title":"Our plan across 2025",
   "intro":"Four workstreams, sequenced across the year, each with clear quarterly milestones.",
   "columns":[{"text":"Standardise intake and reporting on ServiceNow."},
              {"text":"Automate repeat request types across services."},
              {"text":"Roll out dashboards for live queue visibility."},
              {"text":"Introduce self-serve reporting for stakeholders."}],
   "quarters":[{"label":"Q1 2025","items":["Migration complete","Baseline set"]},
               {"label":"Q2 2025","items":["First automations live"]},
               {"label":"Q3 2025","items":["Dashboards adopted"]},
               {"label":"Q4 2025","items":["Self-serve pilot"]}]},
 "agenda_contents":{"title":"Agenda",
   "items":[{"number":"01","item":"Where the Hub stands today"},
            {"number":"02","item":"Quarter in review"},
            {"number":"03","item":"What changed and why"},
            {"number":"04","item":"The 2025–2027 roadmap"},
            {"number":"05","item":"What we're asking for"}]},
 "comparison_table":{"title":"Service tiers at a glance",
   "table":{"headers":["","Standard","Connected","Premium"],
            "rows":[["Response time","Next day","Same day","Priority"],
                    ["Remote monitoring","—","Yes","Yes"],
                    ["Predictive maintenance","—","—","Yes"],
                    ["Reporting","Quarterly","Monthly","Live dashboard"]]}},
 "hero_stat":{"eyebrow":"Resolution rate","value":"91.2%",
   "caption":"of all requests cleared within the focus period",
   "support":"674 of 739 tickets resolved across WCM, DEA and Graphic Design."},
 "image_section_divider":{"image":photo("escalator-station.jpg"),
   "eyebrow":"Part two","title":"Where the Hub goes next"},
 "matrix_2x2":{"title":"Where to focus next","xlabel":"Effort →","ylabel":"Impact →",
   "quadrants":[
     {"heading":"Quick wins","items":["Automate intake tags","Templated replies"]},
     {"heading":"Major projects","items":["Predictive capacity model"]},
     {"heading":"Fill-ins","items":["Tidy archived boards"]},
     {"heading":"Reconsider","items":["Bespoke one-off decks"]}]},
}

def build_full_gallery(out_path):
    A={**B1.ARCHETYPES,**B2.ARCHETYPES,**ARCHETYPES}
    S={**B1.SAMPLES,**B2.SAMPLES,**SAMPLES}
    prs=Presentation(os.path.join(E.KONE_DESIGN,"uploads","master_ppt-1784774200983.pptx"))
    lst=prs.slides._sldIdLst
    for el in list(lst): prs.part.drop_rel(el.get(qn('r:id'))); lst.remove(el)
    blank=next(l for l in prs.slide_layouts if l.name.strip().lower()=="blank")
    import json
    icondir=os.path.join(E.SKILL_DIR,"assets","icons")
    decor_map=json.load(open(os.path.join(icondir,"decor_map.json"))) if os.path.exists(os.path.join(icondir,"decor_map.json")) else {}
    bg_map=json.load(open(os.path.join(icondir,"bg_map.json"))) if os.path.exists(os.path.join(icondir,"bg_map.json")) else {}
    for name,arch in A.items():
        s=prs.slides.add_slide(blank)
        cap=E._tf(s,[45,676,900,24]); E._run(cap.paragraphs[0],f"archetype: {name}",E.KINFO,11,E.GREY,caps=True)
        E.render_archetype(s, arch, S[name], decor=decor_map.get(name), bg=bg_map.get(name))
    prs.save(out_path); return out_path, list(A.keys())

if __name__=="__main__":
    path,names=build_full_gallery("/home/claude/KONE_Archetype_Gallery.pptx")
    print("saved", path, "with", len(names), "archetypes:"); [print("  -",n) for n in names]
