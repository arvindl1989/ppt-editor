"""Turn kone-design's HTML archetype gallery into archetypes the
kone-deck-generator engine can render.

Two archetype vocabularies exist. `kone-design/templates/kone-deck/`
documents 57 slide archetypes as HTML; `kone-deck-generator` implements
23 as declarative region/group data. Only 17 names appear in both, so a
brief asking for COVER_A_CUT4 or DIVIDER_D -- names people read off the
gallery -- got something else built. Reported by a user whose prompt
named six archetypes and got two.

Rather than hand-transcribing geometry (57 slides of it, stale the
moment the gallery moves), this PARSES the gallery. Every block in it is
absolutely positioned in px on the same 1280x720 grid the engine uses,
so `left:45px;top:136px;width:374px;height:448px` maps directly onto a
region's `box`. Re-run `deckguard sync-gallery` after a gallery update
and the archetypes follow.

What comes out is ordinary engine data -- `{"background", "regions",
"groups"}` -- plus one deckguard-only key, `chrome`, for the primitives
the engine has no role for: an arbitrary-colour fill, a logo/tagline
raster, and the signature staggered cut-image banner. `install()` merges
everything into the loaded skill module and wraps its `render` so those
draw too, which keeps ONE renderer and one registry: archetype
signatures, previews, picture-slot detection and the planning prompt all
read the registry at runtime, so they pick these up for free.

Deliberately not handled, and skipped rather than approximated:

- Gradients (only COVER_F_FULLBLEED's photo-protection wash). Faking it
  with a solid would darken a photo that should stay clean.
- `<ul>` bullet markup, which the engine already renders its own way.

Only gallery 1 (covers, dividers, agendas, closers) plus
TITLE_TEXT_SPLIT are installed by default: the gallery's own header
calls gallery 1 "rebuilt 1:1 against the KONE master template" while
galleries 2-4 are "pending the same rework", and most of their contents
already have engine equivalents under the same name.
"""

from __future__ import annotations

import html as html_mod
import os
import re
from pathlib import Path
from typing import Optional

SLIDE_W = 1280.0
SLIDE_H = 720.0

# Gallery 1 is the finished one; TITLE_TEXT_SPLIT is pulled from gallery
# 2 because people ask for it by name. Everything else in galleries 2-4
# either has an engine equivalent already or is pending rework.
DEFAULT_GALLERIES = ("archetypes-1-covers-dividers.html",)
EXTRA_ARCHETYPES = {"archetypes-2-content.html": ("TITLE_TEXT_SPLIT",)}

_SECTION_RE = re.compile(r"<!--\s*(\d+)\s+([A-Z][A-Z0-9_]+)\s*-->(.*?)(?=<!--\s*\d+\s+[A-Z]|\Z)", re.S)
_TAG_RE = re.compile(r"<(section|div|img|h1|h2|p|ul)\b([^>]*)>", re.I)
_STYLE_RE = re.compile(r'style="([^"]*)"', re.I)
_SRC_RE = re.compile(r'src="([^"]*)"', re.I)
_TEXT_RE = re.compile(r">([^<>]+)<")


def gallery_dir() -> Optional[Path]:
    """Where the gallery lives: an explicit override, the installed
    kone-design skill, then deckguard's vendored copy."""
    candidates = []
    env = os.environ.get("KONE_DESIGN_DIR")
    if env:
        candidates.append(Path(env) / "templates" / "kone-deck")
        candidates.append(Path(env))
    candidates.append(Path.home() / ".claude" / "skills" / "kone-design" / "templates" / "kone-deck")
    candidates.append(Path(__file__).parent / "assets" / "kone-design" / "templates" / "kone-deck")
    # Keyed on the gallery FILE, not on ARCHETYPES.md: an older install
    # can carry the spec under the previous gallery filenames, and
    # picking it would parse nothing while looking like it worked.
    for directory in candidates:
        if (directory / DEFAULT_GALLERIES[0]).is_file():
            return directory
    return None


def _css(style: str) -> dict:
    out = {}
    for part in style.split(";"):
        if ":" not in part:
            continue
        key, _, value = part.partition(":")
        out[key.strip().lower()] = value.strip()
    return out


def _px(value: Optional[str]) -> Optional[float]:
    """`45px`, and also a bare `0` -- the gallery writes `left:0;top:0`
    unitless, and rejecting those silently dropped every block anchored
    at the slide origin, the cut-image banners among them."""
    if value is None:
        return None
    m = re.match(r"^(-?[\d.]+)(px)?$", value.strip())
    return float(m.group(1)) if m else None


def _hex_of(value: Optional[str]) -> Optional[str]:
    """`#fff` / `#1450f5` -> uppercase 6-digit hex. None for gradients,
    `transparent`, and anything else this can't render faithfully."""
    if not value:
        return None
    m = re.match(r"^#([0-9a-fA-F]{3}|[0-9a-fA-F]{6})\b", value.strip())
    if not m:
        return None
    digits = m.group(1)
    if len(digits) == 3:
        digits = "".join(c * 2 for c in digits)
    return digits.upper()


def _box_of(css: dict) -> Optional[list]:
    """A block's box in px. `bottom`/`right` are resolved against the
    slide, so footer chrome anchored from the bottom lands correctly."""
    if _px(css.get("inset")) == 0:
        return [0.0, 0.0, SLIDE_W, SLIDE_H]  # `inset:0` is the full-bleed idiom
    left, top = _px(css.get("left")), _px(css.get("top"))
    width, height = _px(css.get("width")), _px(css.get("height"))
    bottom, right = _px(css.get("bottom")), _px(css.get("right"))
    if left is None and right is not None and width is not None:
        left = SLIDE_W - right - width
    if top is None and bottom is not None:
        top = SLIDE_H - bottom - (height or 0)
    if left is None or top is None:
        return None
    return [left, top, width if width is not None else SLIDE_W - left,
            height if height is not None else 0]


def _clip_rects(clip: str) -> list:
    """The cut-image masks are SVG paths made only of axis-aligned
    rectangles (`M0 0 H288.9 V249.1 H0 Z` ...), which is what makes the
    signature staggered banner reproducible in PowerPoint at all: each
    rectangle becomes its own cropped picture."""
    rects = []
    for chunk in re.findall(r"M\s*([\d.]+)\s+([\d.]+)\s*H\s*([\d.]+)\s*V\s*([\d.]+)", clip):
        x0, y0, x1, y1 = (float(v) for v in chunk)
        rects.append([x0, y0, x1 - x0, y1 - y0])
    return rects


def _elements(section_html: str) -> list:
    """Flat list of (tag, css, src, text) for every positioned block."""
    out = []
    for m in _TAG_RE.finditer(section_html):
        tag, attrs = m.group(1).lower(), m.group(2)
        style_m = _STYLE_RE.search(attrs)
        css = _css(style_m.group(1)) if style_m else {}
        src_m = _SRC_RE.search(attrs)
        tail = section_html[m.end(): m.end() + 600]
        text_m = _TEXT_RE.search(">" + tail)
        text = html_mod.unescape(text_m.group(1)).strip() if text_m else ""
        out.append({
            "tag": tag, "css": css, "attrs": attrs,
            "src": src_m.group(1) if src_m else None,
            "text": text,
        })

    # The gallery puts the BOX on a positioned wrapper and the TYPE on an
    # <h2>/<p> inside it. Read separately, the wrapper looks like an
    # empty rectangle and the heading like text with nowhere to go, so
    # every title vanished. Fold the child's typography and copy into
    # its wrapper.
    merged, skip = [], set()
    for i, el in enumerate(out):
        if i in skip:
            continue
        # No font-size on a positioned block means the type lives on a
        # child, even when the naive text scan already picked the
        # child's words up through the wrapper.
        if _box_of(el["css"]) and not el["css"].get("font-size"):
            stack = []
            for j in range(i + 1, min(i + 5, len(out))):
                child = out[j]
                if _box_of(child["css"]):
                    break
                if child["tag"] == "img" and child["src"]:
                    # a photo/illustration sized 100% inside a positioned
                    # wrapper -- the wrapper holds the geometry
                    el = {**el, "src": child["src"]}
                    skip.add(j)
                    break
                if child["css"].get("font-size"):
                    stack.append(child)
                    skip.add(j)
            if stack:
                # One wrapper can hold several styled lines (an eyebrow
                # above a title). Keep them all; parse_section stacks
                # them down the wrapper in proportion to their type size.
                first = stack[0]
                el = {**el, "css": {**first["css"], **{
                    k: v for k, v in el["css"].items()
                    if k in ("left", "top", "width", "height", "bottom", "right")
                }}, "text": first["text"] or el["text"],
                    "stack": stack if len(stack) > 1 else None}
        merged.append(el)
    return merged


# Placeholder copy in the gallery -> the content key an author fills.
_SLOT_BY_TEXT = {
    "presentation title": "title",
    "the section title": "title",
    "section title": "title",
    "thank you": "title",
    "a title held in the coloured field": "title",
    "presenter name · occasion": "context",
    "section label": "eyebrow",
    "agenda": "title",
}


def _role_and_slot(el: dict, index: int) -> Optional[tuple]:
    """Derive (role_key, content_key, style) for a text block from its
    own CSS -- the gallery states font, size, colour and case inline, so
    nothing here is guessed at."""
    css = el["css"]
    size = _px(css.get("font-size"))
    if size is None:
        return None
    color = _hex_of(css.get("color")) or "141414"
    family = css.get("font-family", "")
    caps = "uppercase" in css.get("text-transform", "")
    font = "KONE Information" if "KONE Information" in family else "Inter"
    if _px(css.get("font-weight")) is None and css.get("font-weight") in ("600", "700", "bold"):
        font = "Inter SemiBold"
    role = f"gal_{'k' if font.startswith('KONE') else 'i'}{int(size)}_{color}{'_c' if caps else ''}"
    style = (font, size, color, caps, font == "Inter SemiBold", False)

    text = (el["text"] or "").lower()
    slot = _SLOT_BY_TEXT.get(text)
    if slot is None:
        if re.match(r"^\d+ \w+ \d{4}$", el["text"] or ""):
            return None  # the footer date is chrome, not an author slot
        if re.match(r"^\d{1,2}$", (el["text"] or "").strip()):
            return None  # page number, likewise
        slot = "title" if size >= 40 else ("eyebrow" if caps and font.startswith("KONE") else f"text{index}")
    return role, slot, style


_LOGO_ASSETS = {
    "kone-pictogram-arrow.svg": "kone-pictogram-arrow.png",
    "kone-numeral-3.svg": "kone-numeral-3.png",
    "kone-logo.svg": "kone-logo.png",
    "kone-logo-white.svg": "kone-logo-white.png",
    "kone-tagline.svg": "kone-tagline.png",
    "kone-tagline-white.svg": "kone-tagline-white.png",
    "kone-illustration-technician.png": "kone-illustration-technician.png",
}


def parse_section(name: str, section_html: str) -> Optional[dict]:
    """One gallery `<section>` -> an engine archetype dict (+ `chrome`)."""
    start = section_html.find("<section")
    if start < 0:
        return None
    body = section_html[start:]
    elements = _elements(body)
    if not elements:
        return None

    section_css = elements[0]["css"]
    background = _hex_of(section_css.get("background")) or "FFFFFF"

    regions, chrome, role_styles, sample = [], [], {}, {}
    text_index = 0
    used_slots: set = set()

    def _claim(slot: str) -> str:
        """One `title` per slide. A second wrapper whose own biggest line
        happens to be a lead sentence must not overwrite the real one --
        content keys are what an author fills, so they have to be
        distinct."""
        if slot not in used_slots:
            used_slots.add(slot)
            return slot
        n = 2
        while f"{slot}{n}" in used_slots:
            n += 1
        used_slots.add(f"{slot}{n}")
        return f"{slot}{n}"
    for el in elements[1:]:
        css = el["css"]
        box = _box_of(css)

        clip = css.get("clip-path") or ""
        if "path(" in clip and box:
            rects = _clip_rects(clip)
            if rects:
                chrome.append({"kind": "cut", "box": box, "rects": rects,
                               "content": "photo", "background": background})
                sample.setdefault("photo", "")
                continue

        if el["src"] and box:
            asset = _LOGO_ASSETS.get(Path(el["src"]).name)
            if asset:
                chrome.append({"kind": "asset", "box": box, "asset": asset})
            elif "photos/" in el["src"]:
                regions.append({"role": "picture", "content": "photo", "box": box})
                sample.setdefault("photo", "")
            continue

        if box and not el["text"]:
            fill = _hex_of(css.get("background"))
            if fill and fill != background and box[2] and box[3]:
                chrome.append({"kind": "fill", "box": box, "hex": fill})
            continue

        if not el["text"] or not box:
            continue
        derived = _role_and_slot(el, text_index)
        if derived is None:
            continue
        stack = el.get("stack")
        if stack:
            top = box[1]
            # Within one wrapper the biggest line IS the title, whatever
            # its absolute size: a panel title is 34px where a cover
            # title is 64px, and a flat size threshold called the first
            # of those a nameless text block.
            sizes = [(_px(c["css"].get("font-size")) or 0) for c in stack]
            biggest = max(sizes)
            # Share out the wrapper's REAL height in proportion to type
            # size. Giving each line a flat 1.5x its own size ignored the
            # column it sits in -- a 19px lead paragraph in a 539px
            # column got a 28px box, and shrink-to-fit crushed it to
            # 8.5pt to make it fit a box it never needed to fit.
            total_size = sum(sizes) or 1
            # ...but only when one line actually stands out, and only
            # once: a table's rows all share a size, and calling every
            # one of them "title" is worse than naming none.
            title_at = sizes.index(biggest) if sizes.count(biggest) == 1 else -1
            for child_i, child in enumerate(stack):
                child_el = {"css": child["css"], "text": child["text"]}
                child_derived = _role_and_slot(child_el, text_index)
                if child_derived is None:
                    continue
                c_role, c_slot, c_style = child_derived
                if child_i == title_at and c_slot.startswith("text"):
                    c_slot = "title"
                c_slot = _claim(c_slot)
                text_index += 1
                role_styles[c_role] = c_style
                share = (box[3] or 0) * ((_px(child["css"].get("font-size")) or 0) / total_size)
                height = max(c_style[1] * 1.5, share, 24)
                regions.append({"role": c_role, "content": c_slot,
                                "box": [box[0], top, box[2], height]})
                sample[c_slot] = child["text"]
                top += height
            continue

        role, slot, style = derived
        slot = _claim(slot)
        text_index += 1
        role_styles[role] = style
        if box[3] in (None, 0):
            box = [box[0], box[1], box[2], max(style[1] * 1.6, 24)]
        regions.append({"role": role, "content": slot, "box": box})
        sample[slot] = el["text"]

    if not regions and not chrome:
        return None
    return {
        "archetype": {"background": background, "regions": regions, "groups": [], "chrome": chrome},
        "role_styles": role_styles,
        "sample": sample,
        "name": name,
    }


def build_archetypes(directory: Optional[Path] = None) -> dict:
    """Parse the gallery into `{archetypes, samples, role_styles}`."""
    directory = directory or gallery_dir()
    result = {"archetypes": {}, "samples": {}, "role_styles": {}}
    if directory is None:
        return result

    wanted = {f: None for f in DEFAULT_GALLERIES}
    wanted.update(EXTRA_ARCHETYPES)
    for filename, only in wanted.items():
        path = directory / filename
        try:
            text = path.read_text(encoding="utf-8")
        except OSError:
            continue
        for _number, name, chunk in _SECTION_RE.findall(text):
            if only is not None and name not in only:
                continue
            parsed = parse_section(name, chunk)
            if parsed is None:
                continue
            key = name.lower()
            result["archetypes"][key] = parsed["archetype"]
            result["samples"][key] = parsed["sample"]
            result["role_styles"].update(parsed["role_styles"])
    return result


# --------------------------------------------------------------------------
# Installing into the engine
# --------------------------------------------------------------------------


def _asset_path(name: str) -> Optional[str]:
    for base in (
        Path(__file__).parent / "assets" / "kone-design" / "logo",
        Path.home() / ".claude" / "skills" / "kone-design" / "assets" / "logo",
    ):
        candidate = base / name
        if candidate.is_file():
            return str(candidate)
    return None


def _photo_for(content: dict) -> Optional[str]:
    from deckguard.skill_bridge import _photo_library

    supplied = content.get("photo")
    if supplied and Path(str(supplied)).is_file():
        return str(supplied)
    photos = _photo_library()
    if not photos:
        return None
    seed = str(content.get("title") or content.get("eyebrow") or "")
    return photos[sum(ord(c) for c in seed) % len(photos)]


# Chrome splits by depth, not by type. A cut banner and a colour fill
# sit BEHIND the archetype's text; a logo or tagline sits in FRONT of
# it. Drawing them in one pass put the mark down before the engine
# painted the agenda's full-height photo over the top of it, so that
# slide came out with no logo at all.
BACKGROUND_CHROME = frozenset({"fill", "cut"})
FOREGROUND_CHROME = frozenset({"asset"})


def draw_chrome(slide, arch: dict, content: dict, kinds=None) -> None:
    """Draw the primitives the engine has no role for: colour fills,
    logo/tagline rasters, and the staggered cut-image banner.

    `kinds` restricts the pass to one depth -- see `BACKGROUND_CHROME`
    and `FOREGROUND_CHROME`. Passing nothing draws everything, which is
    only correct for an archetype that paints no picture of its own.
    """
    import importlib

    engine = importlib.import_module("kone_engine")
    for item in arch.get("chrome") or []:
        kind = item.get("kind")
        box = item.get("box")
        if kinds is not None and kind not in kinds:
            continue
        if kind == "fill":
            engine._rect(slide, box, engine._hex(item["hex"]))
        elif kind == "asset":
            path = _asset_path(item["asset"])
            if path:
                _add_transparent_picture(slide, engine, box, path)
        elif kind == "cut":
            _draw_cut(slide, engine, box, item.get("rects") or [],
                      _photo_for(content), item.get("background") or "FFFFFF")


def _add_transparent_picture(slide, engine, box, path) -> None:
    """Place a PNG without going through the engine's `_image`.

    That helper opens every image with `.convert("RGB")`, which drops the
    alpha channel and composites transparency onto BLACK. The KONE marks
    are mostly transparent -- the tagline is 25% ink, the divider
    illustration 16% -- so each one landed as a black rectangle. Reported
    as "the logo has a black background". PowerPoint renders PNG alpha
    natively; the only thing needed is to stop destroying it.

    Fits inside `box` preserving aspect, like `mode="contain"` did.
    """
    from PIL import Image

    x, y, w, h = box
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:  # noqa: BLE001 -- decorative
        return
    scale = min(w / iw, h / ih) if iw and ih else 1.0
    nw, nh = max(1.0, iw * scale), max(1.0, ih * scale)
    slide.shapes.add_picture(str(path), engine.X(x + (w - nw) / 2), engine.X(y + (h - nh) / 2),
                             engine.X(nw), engine.X(nh))


def _draw_cut(slide, engine, box, rects, photo, background="FFFFFF") -> None:
    """The staggered banner is ONE photo with a mask over it, and that is
    how it is built here: a single full-width picture, then
    background-coloured rectangles covering everything the mask cuts
    away.

    Baking four pre-cropped panes instead produced the same picture but
    left the author four separate images to manage -- swapping the cover
    photo meant replacing all four. Reported as wanting "a template
    where when we add a picture it adds that chopped effect, instead of
    it being chopped into four sections already". One picture, one
    Change Picture, the effect intact.
    """
    if not photo or not rects:
        return
    bx, by, bw, bh = box
    slide.shapes.add_picture(str(photo), engine.X(bx), engine.X(by), engine.X(bw), engine.X(bh))

    # Everything the mask does NOT reveal, covered in the slide colour:
    # the gutters between panes, and the area below each pane's own
    # height (which is what staggers them).
    mask = engine._hex(background if len(str(background)) == 6 else "FFFFFF")
    edges = sorted(rects, key=lambda r: r[0])
    cursor = 0.0
    for rx, ry, rw, rh in edges:
        if rx > cursor:
            engine._rect(slide, [bx + cursor, by, rx - cursor, bh], mask)
        if ry > 0:
            engine._rect(slide, [bx + rx, by, rw, ry], mask)
        bottom = ry + rh
        if bottom < bh:
            engine._rect(slide, [bx + rx, by + bottom, rw, bh - bottom], mask)
        cursor = rx + rw
    if cursor < bw:
        engine._rect(slide, [bx + cursor, by, bw - cursor, bh], mask)


_installed = False


def install(archetypes_module) -> int:
    """Merge the gallery archetypes into the loaded skill module and wrap
    its `render` so `chrome` draws too. Idempotent; returns how many
    archetypes were added."""
    global _installed
    if _installed:
        return 0
    import importlib

    built = build_archetypes()
    if not built["archetypes"]:
        return 0

    engine = importlib.import_module("kone_engine")
    engine.ROLE_STYLE.update({
        role: (font, px, engine._hex(color), caps, bold, fit)
        for role, (font, px, color, caps, bold, fit) in built["role_styles"].items()
    })

    archetypes_module.ARCHETYPES.update(built["archetypes"])
    from deckguard.skill_bridge import invalidate_archetype_caches

    invalidate_archetype_caches()
    if hasattr(archetypes_module, "SAMPLES"):
        archetypes_module.SAMPLES.update(built["samples"])
    if hasattr(archetypes_module, "BG"):
        for key, arch in built["archetypes"].items():
            archetypes_module.BG[key] = arch.get("background")

    original_render = archetypes_module.render

    def render_with_chrome(slide, name, content):
        """Paint order is the whole point here, in three layers:
        background, then the regions, then the foreground chrome.

        Both halves were learned from a broken slide. Drawing chrome
        before delegating put TITLE_TEXT_SPLIT's white field down first
        and let the engine's own full-slide background paint over it --
        solid blue slide. And drawing the logo in that same early pass
        put it under the agenda's full-height photo -- no logo at all.
        So the background is drawn here and suppressed for the delegated
        call, and the marks go on last."""
        arch = archetypes_module.ARCHETYPES.get(name) or {}
        if not arch.get("chrome"):
            return original_render(slide, name, content)

        background = archetypes_module.BG.get(name) or arch.get("background")
        if background and str(background).upper() != "FFFFFF":
            rgb = engine._BG.get(background) or (
                engine._hex(background) if len(str(background)) == 6 else None
            )
            if rgb is not None:
                engine._rect(slide, [0, 0, 1280, 720], rgb)
        draw_chrome(slide, arch, content, kinds=BACKGROUND_CHROME)

        saved = archetypes_module.BG.get(name)
        archetypes_module.BG[name] = None
        stripped = {k: v for k, v in arch.items() if k != "background"}
        archetypes_module.ARCHETYPES[name] = stripped
        try:
            return original_render(slide, name, content)
        finally:
            archetypes_module.ARCHETYPES[name] = arch
            archetypes_module.BG[name] = saved
            draw_chrome(slide, arch, content, kinds=FOREGROUND_CHROME)

    archetypes_module.render = render_with_chrome
    _installed = True
    return len(built["archetypes"])


__all__ = ["build_archetypes", "draw_chrome", "gallery_dir", "install", "parse_section"]


# --------------------------------------------------------------------------
# Cover / closer archetypes vs the master's retained slides
# --------------------------------------------------------------------------

COVER_ARCHETYPES = frozenset({
    "cover_a_cut4", "cover_b_cut3", "cover_c_cut4_wide", "cover_d_cut3_wide",
    "cover_e_side", "cover_f_fullbleed",
})
CLOSER_ARCHETYPES = frozenset({"end_logo", "outro"})


def drop_redundant_master_slides(out_path, spec: dict) -> int:
    """`kone_deck_creator.build_deck` always keeps the master's own cover
    and "Thank you" around the body it renders. That is right until the
    brief names a cover or closer ARCHETYPE -- then the deck opens on the
    master's Cover F followed immediately by the COVER_A_CUT4 the author
    actually asked for, and closes on END_LOGO followed by the master's
    Thank you. Two covers, two endings.

    Drops the retained slide on whichever end the author supplied their
    own. Touches nothing when they didn't. Returns how many were removed.
    """
    slides = spec.get("slides") or []
    if not slides:
        return 0
    first = str(slides[0].get("archetype") or "").lower()
    last = str(slides[-1].get("archetype") or "").lower()
    drop_first = first in COVER_ARCHETYPES
    drop_last = last in CLOSER_ARCHETYPES
    if not (drop_first or drop_last):
        return 0

    from pptx import Presentation as _Presentation

    prs = _Presentation(str(out_path))
    id_list = prs.slides._sldIdLst
    elements = list(id_list)
    removed = 0
    for element in ([elements[0]] if drop_first else []) + ([elements[-1]] if drop_last else []):
        try:
            prs.part.drop_rel(element.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
            id_list.remove(element)
            removed += 1
        except Exception:  # noqa: BLE001 -- never fail a build over this
            continue
    if removed:
        prs.save(str(out_path))
    return removed


# --------------------------------------------------------------------------
# Footer chrome: date, page number, tagline
# --------------------------------------------------------------------------

# The gallery's chrome table, which differs by slide type -- covers carry
# no page number, DIVIDER_D puts the number and date together at bottom
# left, END_LOGO carries nothing at all. Extraction dropped all of it,
# because a date and a page number look like author content and are not:
# they are stamped per slide, and a deck came out with no dates or page
# numbers anywhere.
FOOTER_NONE = frozenset({"end_logo"})
FOOTER_NO_PAGE = frozenset({  # date only; covers and the outro
    "cover_a_cut4", "cover_b_cut3", "cover_c_cut4_wide", "cover_d_cut3_wide",
    "cover_e_side", "cover_f_fullbleed",
})
FOOTER_MERGED_LEFT = frozenset({"divider_d"})  # page number then date, both bottom-left
LIGHT_FOOTER = frozenset({"cover_f_fullbleed"})  # white, because it sits on a photo

FOOTER_Y = 677.0  # 43px up from a 720px slide, per the spec
FOOTER_PX = 11.0


def stamp_footers(out_path, spec: dict, date_text: Optional[str] = None) -> int:
    """Stamp the date and page number onto gallery-archetype slides.

    Done as a post-pass rather than inside the renderer because the page
    number is a property of the slide's position in the finished deck,
    which the per-slide renderer has no way to know.
    """
    import datetime
    import importlib

    slides = spec.get("slides") or []
    if not slides:
        return 0
    engine = importlib.import_module("kone_engine")
    from pptx import Presentation as _Presentation

    if date_text is None:
        today = datetime.date.today()
        date_text = f"{today.day} {today.strftime('%B %Y')}"

    prs = _Presentation(str(out_path))
    stamped = 0
    for slide, planned in zip(prs.slides, slides):
        name = str(planned.get("archetype") or "").lower()
        if name in FOOTER_NONE:
            continue
        page = str(prs.slides.index(slide) + 1).zfill(2)
        colour = engine.WHITE if name in LIGHT_FOOTER else engine.BLACK

        def _stamp(box, text):
            frame = engine._tf(slide, box)
            engine._run(frame.paragraphs[0], text, engine.KINFO, FOOTER_PX, colour, caps=True)

        if name in FOOTER_MERGED_LEFT:
            _stamp([45, FOOTER_Y, 300, 20], f"{page}   {date_text}")
        else:
            _stamp([45, FOOTER_Y, 300, 20], date_text)
            if name not in FOOTER_NO_PAGE:
                _stamp([1167, FOOTER_Y, 68, 20], page)
        stamped += 1
    if stamped:
        prs.save(str(out_path))
    return stamped
