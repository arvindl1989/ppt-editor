"""Gate 1: assertions, before and after the deck is drawn.

Claude Design's `VALIDATION.md`, built. The principle worth restating
because it decides the whole shape: **assertions first, rendering
second.** Thirteen of the fourteen faults found in a real deck are
catchable with no model tokens and a deterministic failure. Building a
render gate first is paying vision prices for a `len()`.

Two phases:

    before_build(plan, brief)   on the plan, nothing drawn yet
    after_build(path)           read the .pptx back

`preflight` keeps the checks it already had -- colour, dashes, the
floor, overlaps, the logo count. This module adds what it had no way to
see, and both are merged into one list for the page.

## Severity

Open question 5 -- what preflight should REFUSE -- is not settled, and
it is deliberately not settled here. Every check ships at `report`, and
promotion is a config change rather than a code change:

    report   recorded, deck is returned. The default for every check.
    refuse   the plan is rejected; the caller may retry with the
             finding fed back.
    block    the deck is not returned at all.

Override with `DECKGUARD_SEVERITY="copy_fits=refuse,cardinality=refuse"`.
Nothing in this module decides policy; it only provides the mechanism.
"""

from __future__ import annotations

import os
import re
import zipfile
from dataclasses import dataclass
from typing import Optional

from deckguard import brandmode as bm

REPORT = "report"
REFUSE = "refuse"
BLOCK = "block"

# Every check, at its shipped severity. The names are the contract with
# `DECKGUARD_SEVERITY` and with the tests, so they do not change casually.
DEFAULT_SEVERITY: dict[str, str] = {
    # gate 1a -- on the plan
    "copy_fits": REPORT,
    "no_source_echo": REPORT,
    "cardinality": REPORT,
    "slot_completeness": REPORT,
    "archetype_variety": REPORT,
    # gate 1b -- on the built file
    "font_role": REPORT,
    "chrome": REPORT,
    "mask_integrity": REPORT,
    "notes_present": REPORT,
}

# How much of a slot's estimated capacity may be used before it is
# reported. Not 1.0: the estimate is a character count against a box,
# and a headline that fills its box exactly is already too long.
FIT_TOLERANCE = 1.0

# The shortest run of the brief that counts as lifted rather than
# coincidental. Design's number.
ECHO_CHARS = 40

_GLYPH_EM = 0.52        # average glyph width as a fraction of the size
_LINE_HEIGHT = 1.3

# Roles that are one slot holding many values. The contract mini-language
# cannot express this, so it records them at cardinality 1.
MULTI_VALUE_ROLES = frozenset({"bullets", "table"})


@dataclass(frozen=True)
class Finding:
    check: str
    note: str
    where: str = ""          # archetype, or "slide 4"
    slide: Optional[int] = None

    @property
    def severity(self) -> str:
        return severity(self.check)

    def as_line(self) -> str:
        head = f"{self.where}: " if self.where else ""
        return f"{head}{self.note}"


def severity(check: str) -> str:
    """The configured severity for a check, or its shipped default.

    Read from the environment on every call rather than cached, so a
    deployment can change policy without a restart and a test can set
    one check without leaking into the next.
    """
    raw = os.environ.get("DECKGUARD_SEVERITY") or ""
    for pair in raw.split(","):
        name, _, value = pair.partition("=")
        if name.strip() == check and value.strip() in (REPORT, REFUSE, BLOCK):
            return value.strip()
    return DEFAULT_SEVERITY.get(check, REPORT)


def worst(findings) -> str:
    """The strongest severity present, so a caller can act once."""
    order = [REPORT, REFUSE, BLOCK]
    return max((f.severity for f in findings), key=order.index, default=REPORT)


# --------------------------------------------------------------------------
# gate 1a -- the plan, before anything is drawn
# --------------------------------------------------------------------------


def before_build(plan: dict, brief: str = "") -> list:
    """Every pre-build check, in one pass over the plan."""
    findings: list = []
    slides = plan.get("slides") or []
    audience = str(plan.get("audience") or "internal")

    findings += _copy_fits(slides, plan)
    findings += _no_source_echo(slides, plan, brief)
    findings += _cardinality(slides, audience)
    findings += _slot_completeness(slides, audience)
    findings += _archetype_variety(slides)
    return findings


def _texts(slide: dict):
    """Every string a slide carries, with the slot it came from.

    Yields nested list items too -- `items: [{heading, text}, ...]` is
    where most of a deck's copy actually lives, and checking only the
    top level would miss all of it.
    """
    for key, value in (slide or {}).items():
        if key in ("archetype", "colour", "color", "notes"):
            continue
        if isinstance(value, str):
            yield key, value
        elif isinstance(value, (list, tuple)):
            for item in value:
                if isinstance(item, str):
                    yield key, item
                elif isinstance(item, dict):
                    for field, text in item.items():
                        if isinstance(text, str):
                            yield f"{key}.{field}", text


def budget(archetype: str, key: str) -> int:
    """Roughly how many characters the slot's box holds at its own size.

    Estimated from the registry rather than measured off a render,
    because this runs before anything is drawn -- which is the point of
    gate 1a. The size comes from the region's own type block where it
    has one and from the brand where it does not, so it follows a
    migrated archetype rather than going stale.
    """
    from deckguard.registry import _load_archetypes

    spec = _load_archetypes().ARCHETYPES.get(archetype)
    if not isinstance(spec, dict):
        return 0
    stem, _, field = key.partition(".")

    region = None
    if field:
        for group in spec.get("groups") or []:
            if group.get("content") != stem:
                continue
            region = next((r for r in group.get("regions") or []
                           if r.get("content") == field), None)
    else:
        region = next((r for r in spec.get("regions") or []
                       if r.get("content") == stem), None)
    if not region:
        return 0

    _, _, width, height = region["box"]
    px = (region.get("dg") or {}).get("px")
    if not px:
        settled = bm.resolve(region.get("role") or "", width=width)
        px = settled["px"] if settled else 0
    if not (width and height and px):
        return 0
    per_line = max(int(width / (px * _GLYPH_EM)), 1)
    lines = max(int(height / (px * _LINE_HEIGHT)), 1)
    return per_line * lines


def _copy_fits(slides, plan) -> list:
    """Over-length copy is information, never something to trim.

    A cover title arrived on a real deck as the brief's opening
    sentence hard-cut mid-word at 76px. The cut is gone, but the check
    that would have stopped it being written that way is this one: if
    the string does not fit the slot, the copy was not written for the
    slot.
    """
    out = []
    for index, slide in enumerate(slides, start=1):
        name = slide.get("archetype") or ""
        for key, text in _texts(slide):
            room = budget(name, key)
            if room and len(text) > room * FIT_TOLERANCE:
                out.append(Finding(
                    "copy_fits",
                    f"{key} is {len(text)} characters into room for about "
                    f"{room}. Write for the slot rather than trimming to it.",
                    where=name, slide=index))
    return out


def _normalise(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").lower()).strip()


def _no_source_echo(slides, plan, brief: str) -> list:
    """No slot may be a long run lifted straight out of the brief.

    The highest-value check in the document, and the one deck_13 still
    needs: its cover reads "I would like to share our plan of ONE Week
    MOD deployment with you" -- the email's first line, no longer cut
    mid-word but still not a title anybody wrote.
    """
    source = _normalise(brief)
    if len(source) < ECHO_CHARS:
        return []

    out = []
    checked = [("title", plan.get("title") or "", None)]
    checked += [(key, text, index)
                for index, slide in enumerate(slides, start=1)
                for key, text in _texts(slide)]
    for key, text, index in checked:
        flat = _normalise(text)
        if len(flat) < ECHO_CHARS:
            continue
        if flat[:ECHO_CHARS] in source or flat in source:
            out.append(Finding(
                "no_source_echo",
                f"{key} is lifted from the brief verbatim: {text.strip()[:60]!r}. "
                "Copy for a slot has to be written for it.",
                where=(slides[index - 1].get("archetype") if index else "deck"),
                slide=index))
    return out


def _cardinality(slides, audience: str) -> list:
    """A list of N items needs a region that holds N.

    The fault this names is a planner joining five owners into one
    paragraph in a 270pt column with two thirds of the slide empty
    beside it. If nothing in the library holds the shape of the
    material, that is worth saying rather than flattening.
    """
    from deckguard import contracts as C

    out = []
    for index, slide in enumerate(slides, start=1):
        name = slide.get("archetype") or ""
        contract = C.for_archetype(name, audience)
        if not contract:
            continue
        for slot in contract.slots:
            value = slide.get(slot.key)
            if not isinstance(value, (list, tuple)) or not value:
                continue
            if slot.role in MULTI_VALUE_ROLES:
                # A bulleted list is one slot holding many lines, and the
                # contract notation has no way to say that -- it records
                # `bullets` at cardinality 1 in all fourteen archetypes
                # that have one. Treating that literally reported every
                # correctly-filled bullet slide as a fault, which is how
                # a check gets switched off.
                continue
            if not slot.is_list:
                out.append(Finding(
                    "cardinality",
                    f"{slot.key} was given {len(value)} items but holds one value.",
                    where=name, slide=index))
            elif not (slot.minimum <= len(value) <= slot.maximum):
                out.append(Finding(
                    "cardinality",
                    f"{slot.key} was given {len(value)} items; it holds "
                    f"{slot.minimum}-{slot.maximum}.",
                    where=name, slide=index))
    return out


def _slot_completeness(slides, audience: str) -> list:
    """Every required slot has something in it.

    A slide carrying NO content at all is skipped. That is a pick from
    the slide gallery -- "give me this layout" -- and the renderer fills
    it with the archetype's own sample. Reporting those turned a
    hand-picked deck into 92 findings, which teaches people that the
    panel is noise. A HALF-filled slide is the real fault and is still
    reported.
    """
    from deckguard import contracts as C

    out = []
    for index, slide in enumerate(slides, start=1):
        name = slide.get("archetype") or ""
        contract = C.for_archetype(name, audience)
        if not contract:
            continue
        if not any(v not in (None, "", [], {}) for _k, v in _texts(slide)):
            continue
        for slot in contract.needs:
            if slot.is_picture:
                continue          # filled by the renderer, not by the planner
            value = slide.get(slot.key)
            if value in (None, "", [], {}):
                out.append(Finding(
                    "slot_completeness",
                    f"{slot.key} is required and empty.",
                    where=name, slide=index))
    return out


def _archetype_variety(slides) -> list:
    """The layout-reuse problem, as an assertion.

    A divider repeating is right -- that is what a divider is for -- so
    only content archetypes are counted, the same exception
    `assemble.variety` already makes.
    """
    out = []
    names = [s.get("archetype") or "" for s in slides]
    content = [n for n in names if bm.slide_kind(n) not in ("divider", "cover", "outro")]

    if len(slides) >= 10:
        seen: dict[str, int] = {}
        for name in content:
            seen[name] = seen.get(name, 0) + 1
        for name, count in sorted(seen.items()):
            if count > 2:
                out.append(Finding(
                    "archetype_variety",
                    f"used {count} times in a deck of {len(slides)}.", where=name))

    for index in range(1, len(names)):
        here, before = names[index], names[index - 1]
        if here and here == before and bm.slide_kind(here) not in (
                "divider", "cover", "outro"):
            out.append(Finding(
                "archetype_variety",
                "two consecutive slides use this layout.",
                where=here, slide=index + 1))
    return out


# --------------------------------------------------------------------------
# gate 1b -- the built file
# --------------------------------------------------------------------------


APPROVED_FACES = frozenset({"Inter", "Inter SemiBold", bm.KONE_INFO})


def after_build(deck_path: str) -> list:
    """The post-build checks `preflight` had no way to see.

    Deliberately NOT a replacement for it. `preflight` reads colour,
    dashes, the floor, overlaps and the logo count off the same file
    and is already trusted; this adds the faults a review found that it
    was silent about.
    """
    from pptx import Presentation

    prs = Presentation(deck_path)
    slides = list(prs.slides)
    findings: list = []
    findings += _font_role(slides)
    findings += _chrome(slides)
    findings += _mask_integrity(slides)
    findings += _notes_present(deck_path, len(slides))
    return findings


def _font_role(slides) -> list:
    """The right face, and the casing that goes with it.

    KONE Information is always caps; Inter never is. Both directions
    matter -- a label in sentence case reads as a stray sentence, which
    is exactly how the divider fault was first reported.
    """
    out = []
    for number, slide in enumerate(slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = (run.text or "").strip()
                    face = (run.font.name or "").strip()
                    if not text or not face:
                        continue
                    if face not in APPROVED_FACES:
                        out.append(Finding(
                            "font_role", f"{face!r} is not an approved face.",
                            where=f"slide {number}", slide=number))
                        continue
                    letters = [c for c in text if c.isalpha()]
                    if not letters:
                        continue
                    if face == bm.KONE_INFO and not text.upper() == text:
                        out.append(Finding(
                            "font_role",
                            f"KONE Information not in caps: {text[:40]!r}",
                            where=f"slide {number}", slide=number))
                    # A PHRASE in caps, not a word. Single tokens in
                    # capitals are acronyms and the deck is full of
                    # legitimate ones -- DACH, KONE, AME, EUR, UTM, DAM,
                    # MOD. Requiring two words is what separates
                    # "WHAT SITS OUTSIDE SCOPE" from "DACH".
                    if face.startswith("Inter") and text.upper() == text \
                            and len(text.split()) > 1 and len(letters) > 3:
                        out.append(Finding(
                            "font_role", f"Inter set in caps: {text[:40]!r}",
                            where=f"slide {number}", slide=number))
    return out


def _chrome(slides) -> list:
    """The footer is present exactly where it belongs and says the
    truth.

    All three faults here were live on a real deck at once: a page
    number reading 11 on slide 12, a date reading 23 July 2026 in a
    deck dated 21 August, and a footer on an outro entitled to none.
    """
    out = []
    floor = bm.FLOOR * 9525
    dates: set = set()
    for number, slide in enumerate(slides, start=1):
        band = [sh for sh in slide.shapes
                if getattr(sh, "has_text_frame", False)
                and sh.top is not None and sh.top >= floor
                and sh.text_frame.text.strip()]
        for shape in band:
            text = shape.text_frame.text.strip()
            if re.fullmatch(r"\d{1,3}", text):
                if int(text) != number:
                    out.append(Finding(
                        "chrome",
                        f"page number reads {text} on slide {number}.",
                        where=f"slide {number}", slide=number))
            elif re.search(r"\d", text):
                dates.add(text.upper())
    if len(dates) > 1:
        out.append(Finding(
            "chrome",
            "more than one date in the deck: " + ", ".join(sorted(dates)),
            where="deck"))
    return out


def _mask_integrity(slides) -> list:
    """A cut cover is a mask, not a photo band.

    One picture frame with background-coloured rectangles cut into it.
    A single unbroken band means the mask never drew, which is what a
    reviewed deck opened on.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if not slides:
        return []
    cover = slides[0]
    pictures = [sh for sh in cover.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                and sh.width and sh.width > 3_000_000]
    if not pictures:
        return []            # a cover with no banner is a different design
    rectangles = [sh for sh in cover.shapes
                  if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    if len(rectangles) < 3:
        return [Finding(
            "mask_integrity",
            f"the cover has {len(rectangles)} mask shapes over its banner; "
            "a cut cover needs at least three.",
            where="slide 1", slide=1)]
    return []


def _notes_present(deck_path: str, count: int) -> list:
    """Speaker notes exist at all.

    A deck handing scope to another team carries its non-negotiables in
    the notes. The package currently has no `notesSlide` part anywhere,
    so this reports once rather than once per slide.
    """
    with zipfile.ZipFile(deck_path) as bundle:
        notes = [n for n in bundle.namelist() if "notesSlide" in n]
    if not notes:
        return [Finding("notes_present",
                        f"no speaker notes on any of the {count} slides.",
                        where="deck")]
    return []
