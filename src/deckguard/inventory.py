"""Deck inventory: walks every slide/shape/run once into typed records.

This is the single source of truth `inspect`, `audit`, and `fix` all read
from. Each record keeps a reference to the live python-pptx object
(`obj`) so the fixer can mutate in place without a second XML walk.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

from deckguard import colors as colors_mod
from deckguard import effects as effects_mod
from deckguard import fonts as fonts_mod
from deckguard import logo as logo_mod

HEADING_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}

ALIGN_NAMES = {
    PP_ALIGN.LEFT: "left",
    PP_ALIGN.CENTER: "center",
    PP_ALIGN.RIGHT: "right",
    PP_ALIGN.JUSTIFY: "justify",
    PP_ALIGN.JUSTIFY_LOW: "justify",
    PP_ALIGN.DISTRIBUTE: "distribute",
    PP_ALIGN.THAI_DISTRIBUTE: "distribute",
}
# Reverse lookup for the fixer (e.g. "left" -> PP_ALIGN.LEFT). Several enum
# members collapse to the same name (justify/distribute variants); dict
# construction order keeps the first (canonical) one.
ALIGN_BY_NAME = {name: enum for enum, name in reversed(ALIGN_NAMES.items())}


def _emu_to_in(v) -> Optional[float]:
    if v is None:
        return None
    return round(Emu(v).inches, 4)


@dataclass
class ColorInfo:
    kind: str  # "rgb" | "scheme" | "none"
    hex: Optional[str] = None  # effective resolved hex (no '#'), normalized upper
    theme_color: Optional[str] = None  # MSO_THEME_COLOR name, if scheme-typed
    brightness: float = 0.0


def _color_info(color_format, theme_scheme) -> Optional[ColorInfo]:
    try:
        ctype = color_format.type
    except AttributeError:
        return None
    if ctype is None:
        return ColorInfo(kind="none")
    from pptx.enum.dml import MSO_COLOR_TYPE

    if ctype == MSO_COLOR_TYPE.RGB:
        rgb = color_format.rgb
        hexval = colors_mod.normalize_hex(str(rgb)) if rgb is not None else None
        return ColorInfo(kind="rgb", hex=hexval)
    if ctype == MSO_COLOR_TYPE.SCHEME:
        eff = colors_mod.effective_rgb(color_format, theme_scheme)
        hexval = "".join(f"{c:02X}" for c in eff) if eff else None
        return ColorInfo(
            kind="scheme",
            hex=hexval,
            theme_color=color_format.theme_color.name if color_format.theme_color else None,
            brightness=color_format.brightness or 0.0,
        )
    return None


@dataclass
class FillInfo:
    type: str  # "solid" | "gradient" | "pattern" | "picture" | "background" | "none" | "unknown"
    colors: list[ColorInfo] = field(default_factory=list)


def _fill_info(shape, theme_scheme) -> Optional[FillInfo]:
    try:
        fill = shape.fill
    except (AttributeError, TypeError):
        return None
    try:
        ftype = fill.type
    except Exception:  # noqa: BLE001 — unsupported fill XML shape
        return None
    from pptx.enum.dml import MSO_FILL_TYPE

    if ftype is None:
        return FillInfo(type="inherited")
    if ftype == MSO_FILL_TYPE.SOLID:
        info = _color_info(fill.fore_color, theme_scheme)
        return FillInfo(type="solid", colors=[info] if info else [])
    if ftype == MSO_FILL_TYPE.GRADIENT:
        cinfos = []
        try:
            for stop in fill.gradient_stops:
                info = _color_info(stop.color, theme_scheme)
                if info:
                    cinfos.append(info)
        except Exception:  # noqa: BLE001
            pass
        return FillInfo(type="gradient", colors=cinfos)
    if ftype == MSO_FILL_TYPE.PATTERNED:
        return FillInfo(type="pattern")
    if ftype == MSO_FILL_TYPE.BACKGROUND:
        return FillInfo(type="background")
    if ftype == MSO_FILL_TYPE.PICTURE:
        return FillInfo(type="picture")
    return FillInfo(type="unknown")


@dataclass
class LineInfo:
    color: Optional[ColorInfo] = None


def _has_solid_line_color(shape) -> bool:
    """Read-only check for an actual line color, via raw XML.

    `shape.line.color` looks like a harmless read but isn't: python-pptx's
    LineFormat.color forces the line's fill type to SOLID as a documented
    side effect of merely accessing it, turning "no <a:ln> at all" into
    `<a:ln><a:solidFill/></a:ln>` with no color specified -- an empty fill
    that PowerPoint renders as a visible default-colored line. Every shape
    without an explicit line would silently gain one the moment this
    module inspected it. Only touch the high-level accessor once this
    check has confirmed a real solid color already exists, so plain
    inspection/audit can never corrupt a shape that never had a line.
    """
    spPr = effects_mod.get_spPr(shape)
    if spPr is None:
        return False
    ln = spPr.find(effects_mod.a_qn("ln"))
    if ln is None:
        return False
    solid_fill = ln.find(effects_mod.a_qn("solidFill"))
    if solid_fill is None:
        return False
    return solid_fill.find(effects_mod.a_qn("srgbClr")) is not None or solid_fill.find(effects_mod.a_qn("schemeClr")) is not None


def _line_info(shape, theme_scheme) -> Optional[LineInfo]:
    if not _has_solid_line_color(shape):
        return LineInfo(color=ColorInfo(kind="none"))
    try:
        info = _color_info(shape.line.color, theme_scheme)
    except (AttributeError, TypeError):
        return None
    except Exception:  # noqa: BLE001
        info = None
    return LineInfo(color=info)


@dataclass
class FontContext:
    """What a run with no explicit font override should resolve against:
    its shape's role (governs which of the master's titleStyle/bodyStyle/
    otherStyle applies) plus the master + its theme's font scheme."""

    category: str  # "title" | "body" | "other"
    master: object
    theme_fonts: dict


@dataclass
class RunRecord:
    text: str
    font_raw: Optional[str]
    font_effective: Optional[str]
    bold: bool
    italic: bool
    size_pt: Optional[float]
    color: Optional[ColorInfo]
    effects: set
    all_upper: bool
    obj: object = field(repr=False, compare=False, default=None)


def _has_explicit_run_color(run) -> bool:
    """Read-only check for an actual run text color, via raw XML.

    `run.font.color` looks like a harmless read but isn't: python-pptx's
    Font.color forces the run's fill type to SOLID as a documented side
    effect of merely being accessed (it calls `self.fill.solid()`
    whenever the current type isn't already SOLID), turning "no color
    override, inherits from the theme/master" into an explicit-but-empty
    `<a:solidFill/>` with no color specified. PowerPoint renders that as
    black regardless of what the run actually inherited -- silently
    breaking every run's color the moment this module inspected it.
    Mirrors `_has_solid_line_color` for the exact same reason: only touch
    the high-level accessor once this check confirms a real color
    already exists.
    """
    rPr = effects_mod.get_rPr(run)
    if rPr is None:
        return False
    solid_fill = rPr.find(effects_mod.a_qn("solidFill"))
    if solid_fill is None:
        return False
    return solid_fill.find(effects_mod.a_qn("srgbClr")) is not None or solid_fill.find(effects_mod.a_qn("schemeClr")) is not None


def _run_record(run, theme_scheme, font_ctx: Optional[FontContext], level: int) -> RunRecord:
    font = run.font
    size_pt = font.size.pt if font.size is not None else None
    color_info = _color_info(font.color, theme_scheme) if _has_explicit_run_color(run) else ColorInfo(kind="none")
    text = run.text or ""
    all_upper = any(c.isalpha() for c in text) and text == text.upper()
    font_raw = font.name
    font_effective = font_raw
    if font_effective is None and font_ctx is not None and font_ctx.master is not None:
        font_effective = fonts_mod.resolve_effective_font(
            font_raw, font_ctx.category, level, font_ctx.master, font_ctx.theme_fonts
        )
    return RunRecord(
        text=text,
        font_raw=font_raw,
        font_effective=font_effective,
        bold=bool(font.bold),
        italic=bool(font.italic),
        size_pt=size_pt,
        color=color_info,
        effects=effects_mod.detect_run_effects(run),
        all_upper=all_upper,
        obj=run,
    )


@dataclass
class ParagraphRecord:
    alignment: Optional[str]
    alignment_explicit: bool
    level: int
    runs: list
    obj: object = field(repr=False, compare=False, default=None)


def _paragraph_record(paragraph, theme_scheme, font_ctx: Optional[FontContext] = None) -> ParagraphRecord:
    raw_align = paragraph.alignment
    level = paragraph.level
    return ParagraphRecord(
        alignment=ALIGN_NAMES.get(raw_align) if raw_align is not None else None,
        alignment_explicit=raw_align is not None,
        level=level,
        runs=[_run_record(r, theme_scheme, font_ctx, level) for r in paragraph.runs],
        obj=paragraph,
    )


@dataclass
class ImageInfo:
    phash: Optional[str]
    filename: Optional[str]
    ext: Optional[str]
    size_bytes: Optional[int]


def _image_info(shape) -> Optional[ImageInfo]:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    try:
        image = shape.image
    except Exception:  # noqa: BLE001 — unreadable/corrupt embedded image
        return ImageInfo(phash=None, filename=None, ext=None, size_bytes=None)
    try:
        phash = logo_mod.compute_phash(image.blob)
    except Exception:  # noqa: BLE001
        phash = None
    return ImageInfo(phash=phash, filename=image.filename, ext=image.ext, size_bytes=len(image.blob))


def _detect_wordart(shape) -> bool:
    body_pr = shape._element.find(".//" + effects_mod.a_qn("bodyPr"))
    if body_pr is None:
        return False
    warp = body_pr.find(effects_mod.a_qn("prstTxWarp"))
    if warp is None:
        return False
    prst = warp.get("prst")
    return prst not in (None, "textNoShape")


@dataclass
class ShapeRecord:
    shape_id: int
    name: str
    shape_type: str
    left_in: Optional[float]
    top_in: Optional[float]
    width_in: Optional[float]
    height_in: Optional[float]
    rotation: float
    fill: Optional[FillInfo]
    line: Optional[LineInfo]
    effects: set
    paragraphs: list
    image: Optional[ImageInfo]
    is_placeholder: bool
    placeholder_type: Optional[str]
    is_wordart: bool
    children: list = field(default_factory=list)
    obj: object = field(repr=False, compare=False, default=None)


def _shape_record(shape, theme_scheme, master=None, theme_fonts: Optional[dict] = None) -> ShapeRecord:
    placeholder_type = None
    is_placeholder = False
    try:
        is_placeholder = shape.is_placeholder
        if is_placeholder:
            placeholder_type = shape.placeholder_format.type.name if shape.placeholder_format.type else None
    except (AttributeError, ValueError):
        pass

    if placeholder_type in HEADING_PLACEHOLDER_TYPES:
        category = "title"
    elif is_placeholder:
        category = "body"
    else:
        category = "other"
    font_ctx = FontContext(category=category, master=master, theme_fonts=theme_fonts or {}) if master is not None else None

    children = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        children = [_shape_record(child, theme_scheme, master, theme_fonts) for child in shape.shapes]

    paragraphs = []
    if getattr(shape, "has_text_frame", False):
        paragraphs = [_paragraph_record(p, theme_scheme, font_ctx) for p in shape.text_frame.paragraphs]

    try:
        rotation = shape.rotation
    except (AttributeError, NotImplementedError):
        rotation = 0.0

    try:
        shape_type_name = shape.shape_type.name if shape.shape_type is not None else "UNKNOWN"
    except AttributeError:
        shape_type_name = "UNKNOWN"

    return ShapeRecord(
        shape_id=shape.shape_id,
        name=shape.name,
        shape_type=shape_type_name,
        left_in=_emu_to_in(getattr(shape, "left", None)),
        top_in=_emu_to_in(getattr(shape, "top", None)),
        width_in=_emu_to_in(getattr(shape, "width", None)),
        height_in=_emu_to_in(getattr(shape, "height", None)),
        rotation=rotation,
        fill=_fill_info(shape, theme_scheme),
        line=_line_info(shape, theme_scheme),
        effects=effects_mod.detect_shape_effects(shape),
        paragraphs=paragraphs,
        image=_image_info(shape),
        is_placeholder=is_placeholder,
        placeholder_type=placeholder_type,
        is_wordart=_detect_wordart(shape),
        children=children,
        obj=shape,
    )


@dataclass
class SlideRecord:
    index: int
    layout_name: str
    master_name: str
    shapes: list
    # Non-placeholder decorative shapes drawn directly on this slide's
    # LAYOUT (e.g. a full-height color panel behind a "text" placeholder
    # column) -- these render behind every slide built on that layout but
    # are never copied into the slide's own shape tree, so a plain
    # placeholder textbox sitting on one has no fill of its own and no
    # same-slide shape to resolve it from either. See
    # `rules_engine._resolve_effective_bg_hex`.
    layout_background_shapes: list = field(default_factory=list)
    obj: object = field(repr=False, compare=False, default=None)


@dataclass
class DeckInventory:
    slide_width_emu: int
    slide_height_emu: int
    aspect_ratio: str
    slides: list
    prs: object = field(repr=False, compare=False, default=None)


def _aspect_ratio(width_emu: int, height_emu: int) -> str:
    if not width_emu or not height_emu:
        return "unknown"
    ratio = width_emu / height_emu
    # 16:9 ~= 1.778, 4:3 ~= 1.333 — use a tolerance to absorb rounding.
    if abs(ratio - 16 / 9) < 0.02:
        return "16:9"
    if abs(ratio - 4 / 3) < 0.02:
        return "4:3"
    return f"{ratio:.3f}:1"


def build_inventory(prs: Presentation) -> DeckInventory:
    theme_by_master_part = {}
    theme_fonts_by_master_part = {}
    for master in colors_mod.iter_slide_masters(prs):
        theme_by_master_part[id(master.part)] = colors_mod.get_theme_scheme(master)
        theme_fonts_by_master_part[id(master.part)] = fonts_mod.read_theme_font_scheme(master)

    layout_bg_by_layout_part: dict = {}

    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        layout = slide.slide_layout
        master = layout.slide_master
        theme_scheme = theme_by_master_part.get(id(master.part))
        theme_fonts = theme_fonts_by_master_part.get(id(master.part), {})
        shape_records = [_shape_record(s, theme_scheme, master, theme_fonts) for s in slide.shapes]

        layout_key = id(layout.part)
        if layout_key not in layout_bg_by_layout_part:
            layout_shape_records = [_shape_record(s, theme_scheme, master, theme_fonts) for s in layout.shapes]
            layout_bg_by_layout_part[layout_key] = [rec for rec in layout_shape_records if not rec.is_placeholder]

        slides.append(
            SlideRecord(
                index=i,
                layout_name=layout.name or "",
                master_name=master.name or "",
                shapes=shape_records,
                layout_background_shapes=layout_bg_by_layout_part[layout_key],
                obj=slide,
            )
        )

    return DeckInventory(
        slide_width_emu=prs.slide_width,
        slide_height_emu=prs.slide_height,
        aspect_ratio=_aspect_ratio(prs.slide_width, prs.slide_height),
        slides=slides,
    )


def iter_shapes_recursive(shape_records):
    """Flatten a list of ShapeRecord (including group children) depth-first."""
    for rec in shape_records:
        yield rec
        if rec.children:
            yield from iter_shapes_recursive(rec.children)
