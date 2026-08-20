"""Turn whatever the user gave us into a deck.

One function decides what a deck is. The interface collects inputs and
shows results; this is where a brief, a set of picks, and a mined deck
become an ordered list of slides with content in them.

The rules, in order:

  1. Slides the user picked are honoured exactly, in the order picked.
  2. A brief with no picks is planned -- the model chooses archetypes
     from the audience's set and writes the content.
  3. A brief WITH picks fills those picks rather than choosing its own.
  4. With neither, a mined deck's own templates are the deck.

Nothing here asks a question. A missing slot stays empty and shows as
empty; a slide that cannot be built is dropped and reported. The caller
edits afterwards.
"""

from __future__ import annotations

import re
from typing import Optional

from deckguard import brandmode as bm
from deckguard.registry import _derived_content_keys, _load_archetypes


def plan(
    *,
    brief: str = "",
    audience: str = "internal",
    picks: Optional[list] = None,
    mined: Optional[dict] = None,
    title: str = "",
    sections: Optional[list] = None,
) -> dict:
    """The deck, as a spec, before anything is drawn."""
    picks = list(picks or [])
    mined = mined or {}
    built = set(_load_archetypes().ARCHETYPES) | set(mined.get("archetypes") or {})

    chosen = _chosen_archetypes(picks, audience, mined, built)
    samples = mined.get("samples") or {}
    if brief.strip():
        slides = _from_brief(brief, audience, chosen, sections or [])
    elif chosen:
        slides = [_seed(name, samples) for name in chosen]
    else:
        # A deck and nothing else: rebuild it from its own designs,
        # carrying the words that were on those slides. An archetype
        # with no content renders as an empty page, which is not a
        # template anyone can judge.
        slides = [_seed(name, samples) for name in _mined_names(mined)][:25]

    return {
        "title": title or _title_from(brief) or "Untitled deck",
        "audience": audience,
        "date": bm_date(),
        "slides": slides,
    }


def _seed(name: str, samples: dict) -> dict:
    """A slide, pre-filled from its mined sample where there is one."""
    sample = samples.get(name)
    return {"archetype": name, **(dict(sample) if isinstance(sample, dict) else {})}


def _chosen_archetypes(picks, audience, mined, built) -> list:
    """Picks are `set:n` for a KONE slide or `mined:name` for one of
    yours. Anything not actually built is dropped rather than silently
    producing a blank slide."""
    names = []
    for pick in picks:
        source, _, rest = pick.partition(":")
        if source == "mined":
            if rest in (mined.get("archetypes") or {}):
                names.append(rest)
            continue
        try:
            number = int(rest)
        except ValueError:
            continue
        for slide in bm.slides_in(source if source in bm.set_names() else audience):
            if slide["n"] == number and slide["archetype"] in built:
                names.append(slide["archetype"])
                break
    return names


def _mined_names(mined: dict) -> list:
    order = mined.get("sources") or {}
    return sorted((mined.get("archetypes") or {}),
                  key=lambda n: min(order.get(n) or [999]))


def _menu(audience: str) -> str:
    """The audience's set, each entry carrying what it needs.

    `brandmode.menu` gave the job and nothing else, so an archetype
    whose material the brief does not contain looked exactly as
    choosable as one whose material it does. The contract is the part
    that lets a layout be ruled OUT.
    """
    try:
        from deckguard import contracts

        return contracts.guide(audience)
    except Exception:  # noqa: BLE001 -- a menu without contracts still works
        return bm.menu(audience)


def _from_brief(brief: str, audience: str, chosen: list,
                sections: Optional[list] = None) -> list:
    """Plan the content. With picks already made, the planner fills
    those; without, it chooses from the audience's set too."""
    from deckguard.planner import call_claude_for_kone_spec

    notes = None
    if chosen:
        notes = (
            "Use EXACTLY these archetypes, in this order, one slide each:\n"
            + "\n".join(f"  {i:02d}. {name}" for i, name in enumerate(chosen, 1))
            + "\nDo not add, drop or reorder them."
        )
    else:
        # Two things made briefs come out using the same first handful
        # every time. "Keep them in the set's order" read as "march down
        # this list from the top". And the list was bare names -- nothing
        # said when `matrix_2x2` beats `segment_breakdown`, so there was
        # no merit to choose on and position was the only signal left.
        notes = (
            f"This is a KONE {audience} deck.\n\n"
            "Below is the menu you may choose from. It is a MENU, not a "
            "running order: pick the archetypes that suit what the source "
            "actually says, in whatever order tells the story best.\n\n"
            "How to choose:\n"
            "  - Match the slide to the CONTENT. A number that carries the "
            "whole point wants hero_stat; three or more figures want "
            "kone_numbers or statement_b; a comparison wants two_content or "
            "comparison_table; a sequence wants timeline or "
            "how_it_works_3step; someone's words want a quote.\n"
            "  - Use an archetype ONCE unless the content genuinely repeats. "
            "A deck that reuses one layout for five slides reads as a "
            "template nobody chose.\n"
            "  - Reach across the whole menu. The interesting archetypes are "
            "in the middle and at the end, not only at the top.\n"
            "  - Open with a cover and put a divider before each section. "
            "Do NOT emit an outro: the master's own Thank you is retained.\n"
            "  - Drop a slide rather than pad it. Eight slides that each say "
            "something beat fifteen that do not.\n"
            "  - Every entry below says what it NEEDS. If the source does "
            "not give you that material, the answer is a different "
            "archetype, not the same one half filled.\n\n"
            f"The {audience} menu:\n" + _menu(audience)
        )
        wanted = bm.sections_brief(sections or [], audience)
        if wanted:
            # The brief says what happened; this says what the deck has
            # to cover. Together they narrow the choice from fifty
            # archetypes to a shortlist per section, which is the whole
            # point of asking.
            notes += (
                "\n\nThe deck MUST cover these, roughly in this order. For "
                "each one, choose the archetype from its shortlist that best "
                "fits what the source actually says -- and if the source has "
                "nothing for a section, drop it rather than invent:\n"
                + wanted
            )
    spec, _usage = call_claude_for_kone_spec(brief, notes=notes)
    slides = spec.get("slides") or []
    return slides or [{"archetype": name} for name in chosen]


def variety(plan_dict: dict) -> dict:
    """How many distinct archetypes the deck uses, and the worst repeat.

    Reported rather than enforced. A divider repeating once per section
    is correct; the same content layout five times is a planner that
    fell back on position instead of choosing. Only the reader can tell
    those apart, so show the number and let them.
    """
    slides = plan_dict.get("slides") or []
    counts: dict = {}
    for slide in slides:
        name = slide.get("archetype", "")
        counts[name] = counts.get(name, 0) + 1
    repeated = [(n, c) for n, c in counts.items() if c > 1 and "divider" not in n]
    repeated.sort(key=lambda pair: -pair[1])
    return {
        "distinct": len(counts),
        "total": len(slides),
        "repeats": repeated,
    }


def _title_from(brief: str) -> str:
    line = (brief or "").strip().splitlines()[:1]
    if not line:
        return ""
    text = re.sub(r"^(subject|re|fwd)\s*:\s*", "", line[0], flags=re.I).strip()
    return text[:70]


def bm_date() -> str:
    from datetime import date

    return date.today().strftime("%d %B %Y").lstrip("0")


# --------------------------------------------------------------------------
# building
# --------------------------------------------------------------------------


def build(plan_dict: dict, out_path: str, mined: Optional[dict] = None) -> dict:
    """Draw the deck and run preflight over it. Returns the findings."""
    from deckguard import layouts
    from deckguard.registry import fill_empty_photo_slots, register_mined

    # Mined designs have to be in the registry before the renderer looks
    # one up, or it finds nothing and draws an empty slide.
    if mined:
        register_mined(mined)

    spec = {
        "title": plan_dict.get("title") or "Untitled deck",
        "date": plan_dict.get("date"),
        # Carried through because the renderer needs it: `slide-sets.json`
        # declares a field per (archetype, audience), and without the
        # audience every slide fell back to its layout's own background.
        # That is how a customer deck kept a sand band on `kone_numbers`
        # where the external policy is white and blue only.
        "audience": plan_dict.get("audience") or "",
        "slides": [dict(s) for s in plan_dict.get("slides") or []],
    }
    fill_empty_photo_slots(spec)
    layouts.build_deck(spec, out_path)
    return preflight(out_path)


def _is_dash_marker(text: str, rest: str = "") -> bool:
    """Is this run a dash pretending to be a bullet?

    Checking `text.strip().startswith("— ")` missed the real case: the
    marker is its own run holding exactly `"—  "`, and stripping it
    leaves a bare dash with no trailing space. So the check passed on
    every deck while the violation it was written for went out the
    door.

    `rest` is whatever follows in the same paragraph, and a bare dash
    needs it: a marker is a dash with something after it. A dash ALONE
    in its paragraph is a value -- the "not in this tier" cell of a
    comparison table -- and reading it as a §6 violation made the one
    archetype built to hold a table always report one.
    """
    stripped = (text or "").strip()
    if stripped in ("-", "—", "–", "*", "•-"):
        return bool(rest.strip())
    return stripped.startswith(("- ", "— ", "– "))


def _below_the_floor(shape, px) -> bool:
    """Does this text region hang below the floor?

    The floor is about content, not about the slide. A full-bleed
    photograph reaches y=720 by design, and so does a cover's own
    banner -- flagging those made preflight cry wolf on every deck,
    which is the fastest way to teach someone to ignore it.
    """
    try:
        top, height, width = shape.top, shape.height, shape.width
    except Exception:  # noqa: BLE001
        return False
    if None in (top, height, width):
        return False
    if width <= 0 or height <= 0:
        # The master's latent DATE / FOOTER / SLIDE_NUMBER placeholders
        # come back as 0x0 boxes parked at y=720. They draw nothing, and
        # flagging them made every deck report two findings it could not
        # act on.
        return False
    if width / px > 1000 and height / px > 300:
        return False                      # full-bleed art, not a text region
    if top / px <= 1 and (top + height) / px >= 719:
        # A colour field running edge to edge top to bottom -- the mint
        # column on the split agenda, the blue column on a quote. It is
        # the slide's ground, not something standing on the floor, and
        # it is meant to reach the bottom edge.
        return False
    if (getattr(shape, "name", "") or "").startswith(("Chrome", "Hairline")):
        return False                      # chrome belongs below the floor
    return (top + height) / px > bm.FLOOR + 1


def _ink_size(shape, width: float) -> tuple:
    """Roughly what the text in a shape occupies, in px.

    Same estimate the previews use: a sans glyph advances ~0.52em and a
    line occupies 1.25x its size. Approximate on purpose -- it only has
    to tell a one-line eyebrow in a tall box from a paragraph that fills
    one.
    """
    used, widest = 0.0, 0.0
    for para in shape.text_frame.paragraphs:
        text = "".join(r.text or "" for r in para.runs).strip()
        sizes = [r.font.size.pt for r in para.runs if r.font.size]
        px_size = (max(sizes) if sizes else 12.0) * (1280.0 / 960.0)
        if not text:
            used += px_size * 1.25
            continue
        run = len(text) * px_size * 0.52
        lines = max(1, int(run // width) + (1 if run % width else 0)) if width else 1
        used += lines * px_size * 1.25
        widest = max(widest, min(run, width))
    return (widest or width, used or 16.0)


def _overlaps(slide, px) -> list:
    """Pairs of shapes that sit on top of each other.

    Preflight measured one thing about position -- whether a region
    crossed the floor -- so an icon drawn over its own heading passed
    every check while being the most visible fault on the slide.
    `lifecycle_4stage` had eight overlapping pairs and `text_picture_a`
    nine, and both shipped.

    Only text against text or text against a picture counts. Panels,
    scrims, colour fields and the logo are meant to sit under things.
    """
    ignore = ("Chrome", "Colour field", "Logo", "Hairline", "Scrim", "Panel")
    boxes = []
    for shape in slide.shapes:
        name = getattr(shape, "name", "") or ""
        if name.startswith(ignore):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        if not shape.text_frame.text.strip():
            continue
        try:
            left, top = shape.left / px, shape.top / px
            width, height = shape.width / px, shape.height / px
        except (TypeError, ZeroDivisionError):
            continue
        if width <= 0 or height <= 0 or width > 1000 and height > 300:
            continue
        # Ink, not box. Boxes are drawn generously -- a divider's eyebrow
        # sits in a 92px box holding one 12px line -- so comparing boxes
        # reported the eyebrow as colliding with the title below it on
        # every numbered divider in the library. What matters is whether
        # the TEXT lands on other text.
        ink_w, ink_h = _ink_size(shape, width)
        boxes.append((left, top, min(width, ink_w), min(height, ink_h),
                      " ".join(shape.text_frame.text.split())[:28]))

    found = []
    for index, one in enumerate(boxes):
        for other in boxes[index + 1:]:
            across = min(one[0] + one[2], other[0] + other[2]) - max(one[0], other[0])
            down = min(one[1] + one[3], other[1] + other[3]) - max(one[1], other[1])
            # 4px of slack: boxes are drawn to touch, and a shared edge
            # is a layout, not a collision.
            if across > 4 and down > 4:
                found.append((one[4], other[4], across, down))
    return found


def preflight(deck_path: str) -> dict:
    """The checks from BRAND_MODE section 10 that can be read back off a
    built file. A deck that fails one is still returned -- the point is
    to say so, not to withhold it."""
    from pptx import Presentation

    findings: list = []
    prs = Presentation(deck_path)
    px = prs.slide_width / 1280
    allowed = {bm.BLACK, bm.WHITE, bm.BLUE}

    slides = list(prs.slides)
    for number, slide in enumerate(slides, start=1):
        # The first and last slides are the master's own retained cover
        # and Thank you. The tool does not place their content and
        # cannot move it, so holding them to the floor reports something
        # nobody can act on -- the same reason the 0x0 placeholders are
        # skipped.
        retained = number in (1, len(slides))
        logos = 0
        for shape in slide.shapes:
            name = (getattr(shape, "name", "") or "").lower()
            if "logo" in name:
                logos += 1
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                runs = list(para.runs)
                for index, run in enumerate(runs):
                    text = run.text or ""
                    after = "".join(r.text or "" for r in runs[index + 1:])
                    if not text.strip():
                        continue
                    colour = None
                    try:
                        if run.font.color and run.font.color.type is not None:
                            colour = str(run.font.color.rgb)
                    except Exception:  # noqa: BLE001
                        colour = None
                    if colour and colour.upper() not in allowed:
                        findings.append((number, f"type in #{colour}, which is not "
                                                 "black, white or KONE Blue"))
                    if _is_dash_marker(text, after):
                        findings.append((number, "a dash standing in for a bullet"))
            if not retained and _below_the_floor(shape, px):
                bottom = (shape.top + shape.height) / px
                findings.append((number, f"a region reaching y={bottom:.0f}, past the floor"))
        if logos > 1:
            findings.append((number, f"{logos} logos; there must be exactly one"))
        if not retained:
            for one, other, across, down in _overlaps(slide, px)[:3]:
                findings.append((
                    number,
                    f"{one!r} and {other!r} overlap by {across:.0f}x{down:.0f}px"))

    seen, unique = set(), []
    for item in findings:
        if item not in seen:
            seen.add(item)
            unique.append(item)
    return {"slides": len(prs.slides), "findings": unique}


# --------------------------------------------------------------------------
# edits from the result page
# --------------------------------------------------------------------------


def apply_edits(plan_dict: dict, form) -> dict:
    """Text, order, drop and duplicate, read back off the result page.

    Order is a number that only has to sort, which is what makes
    inserting a slide between two others easy.
    """
    slides = []
    for index, slide in enumerate(plan_dict.get("slides") or []):
        sid = str(index)
        if form.get(f"drop:{sid}"):
            continue
        updated = dict(slide)
        for key in list(updated):
            if key == "archetype":
                continue
            field = form.get(f"v:{sid}:{key}")
            if field is None:
                continue
            updated[key] = _retype(updated[key], str(field))
        try:
            order = float(form.get(f"o:{sid}") or index)
        except ValueError:
            order = float(index)
        slides.append((order, updated))
        if form.get(f"dup:{sid}"):
            slides.append((order + 0.5, dict(updated)))

    return {**plan_dict,
            "slides": [s for _o, s in sorted(slides, key=lambda pair: pair[0])]}


def _retype(original, edited: str):
    """Keep a slot's shape. A list edited as lines stays a list; a list
    of dicts keeps its keys, split on pipes in the order they appear."""
    text = edited.strip()
    if isinstance(original, list):
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        if original and isinstance(original[0], dict):
            fields = list(original[0])
            out = []
            for line in lines:
                parts = [p.strip() for p in line.split("|")]
                out.append({f: (parts[i] if i < len(parts) else "")
                            for i, f in enumerate(fields)})
            return out
        return lines
    return text


def slots_for(archetype: str, content: dict) -> list:
    """The editable slots for a slide: what the renderer reads, in the
    order it reads them, with what is currently in each."""
    slots = []
    for raw in _derived_content_keys(archetype):
        key = raw.split(" (")[0]
        if "filled automatically" in raw:
            continue
        slots.append((key, raw[len(key):].strip(" ()"), content.get(key)))
    return slots
