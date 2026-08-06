"""KONE pictograms as native PowerPoint shapes.

The brand ships 609 icons as an icon font plus an SVG sprite. Neither
form goes into a .pptx as-is:

- An icon FONT renders the glyph as text, so it needs `Kone-icons`
  installed on every machine that opens the deck. Anyone without it
  sees tofu, and the "icon" is a character -- it cannot be recoloured
  independently of its text run and every brand rule about type applies
  to it by accident.
- An SVG cannot be handed to `add_picture`; python-pptx identifies
  images through Pillow, which does not read SVG. It can be smuggled in
  as an extension on a raster blip (see `logo.attach_svg`), and that
  renders as vector -- but PowerPoint still treats it as a picture. You
  cannot pick the shape apart, and recolouring means editing the asset.
- Rasterising loses the point entirely.

So these are converted to DrawingML `<a:custGeom>` instead: a real
PowerPoint freeform shape, drawn from the icon's own outline. It scales
without blurring, takes a fill like any other shape, needs no font, no
image part and no external asset, and a designer can select it and edit
the points.

That conversion is only tractable because of what the sprite happens to
be. All 609 symbols share one `0 0 1024 1024` viewBox, every one is a
single `<path>` with no groups, transforms or primitives, and -- the
part that actually matters -- **not one of them uses an elliptical
arc**. Arcs have no DrawingML equivalent and would each need
decomposing into beziers. Without them the command set is
move/line/horizontal/vertical/cubic/smooth-cubic/close, and every one
of those maps directly onto a DrawingML path element.
"""

from __future__ import annotations

import re
from functools import lru_cache
from pathlib import Path
from typing import Iterator, Optional

from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

_SPRITE_NAME = "kone-icons-sprite.svg"
_VENDORED = Path(__file__).with_name("assets") / "kone-design" / "icons"
_INTERACTIVE = Path("~/.claude/skills/kone-design/assets/icons").expanduser()

# The sprite's own coordinate space. Every symbol declares it, and the
# loader rejects any that doesn't rather than silently mis-scaling.
VIEWBOX = 1024

_SYMBOL_RE = re.compile(
    r'<symbol\s+id="(?P<id>[^"]+)"\s+viewBox="(?P<box>[^"]+)"[^>]*>\s*<path\s+d="(?P<d>[^"]+)"',
    re.S,
)
_NUMBER_RE = re.compile(r"[-+]?(?:\d*\.\d+|\d+)(?:[eE][-+]?\d+)?")
_TOKEN_RE = re.compile(r"([MmLlHhVvCcSsZz])|([-+]?(?:\d*\.\d+|\d+)(?:[eE][-+]?\d+)?)")


def sprite_path() -> Optional[Path]:
    for base in (_VENDORED, _INTERACTIVE):
        candidate = base / _SPRITE_NAME
        if candidate.is_file():
            return candidate
    return None


@lru_cache(maxsize=1)
def load_icons() -> dict[str, str]:
    """{name: path data}, keyed without the sprite's `i-` prefix.

    The sprite ships UTF-16 encoded; the vendored copy is re-encoded to
    UTF-8, but an installed skill may still have the original, so both
    are handled.
    """
    path = sprite_path()
    if path is None:
        return {}
    raw = path.read_bytes()
    text = raw.decode("utf-16") if raw[:2] in (b"\xff\xfe", b"\xfe\xff") else raw.decode("utf-8")

    icons: dict[str, str] = {}
    for match in _SYMBOL_RE.finditer(text):
        box = [float(v) for v in match.group("box").replace(",", " ").split()]
        if box != [0.0, 0.0, float(VIEWBOX), float(VIEWBOX)]:
            continue  # a symbol in another coordinate space would scale wrong
        name = match.group("id")
        icons[name[2:] if name.startswith("i-") else name] = match.group("d")
    return icons


def icon_names() -> list[str]:
    return sorted(load_icons())


def find_icons(term: str, limit: int = 12) -> list[str]:
    """Icon names containing `term` -- how a caller goes from "elevator"
    to a real name without knowing the catalogue."""
    needle = term.strip().lower().replace(" ", "-")
    if not needle:
        return []
    names = icon_names()
    exact = [n for n in names if n == needle]
    prefix = [n for n in names if n.startswith(needle) and n not in exact]
    rest = [n for n in names if needle in n and n not in exact and n not in prefix]
    return (exact + prefix + rest)[:limit]


# --------------------------------------------------------------------------
# SVG path -> DrawingML path
# --------------------------------------------------------------------------


def _tokens(data: str) -> Iterator[tuple[Optional[str], Optional[float]]]:
    for command, number in _TOKEN_RE.findall(data):
        yield (command, None) if command else (None, float(number))


def path_to_drawingml(data: str, scale: float = 1.0) -> list[tuple]:
    """SVG path data -> a flat list of DrawingML path operations.

    Each operation is `("moveTo"|"lnTo"|"cubicBezTo"|"close", points)`,
    points being absolute (x, y) pairs in the SOURCE's own coordinate
    space. Everything is resolved here -- relative commands, the
    horizontal/vertical shorthands, and the reflected control point of
    a smooth curve -- because DrawingML has none of those forms.

    `scale` is a plain multiplier and defaults to none. It was once
    "the viewBox you want, relative to the icon sprite's 1024", which
    silently multiplied the illustrations by 1.22 when they passed
    their own 1250 -- paths oversized, circles and rects untouched,
    every drawing scattered. Only a side-by-side render caught it.
    """
    numbers: list[float] = []
    command: Optional[str] = None
    ops: list[tuple] = []
    x = y = start_x = start_y = 0.0
    # last cubic control point, for reflecting into a following S/s
    ctrl_x = ctrl_y = None

    def flush() -> None:
        nonlocal numbers, x, y, start_x, start_y, ctrl_x, ctrl_y
        if command is None:
            numbers = []
            return
        relative = command.islower()
        upper = command.upper()
        values = numbers

        def take(n: int) -> Iterator[list[float]]:
            for i in range(0, len(values) - n + 1, n):
                yield values[i:i + n]

        if upper == "M":
            first = True
            for px, py in take(2):
                x, y = (x + px, y + py) if relative else (px, py)
                if first:
                    start_x, start_y = x, y
                    ops.append(("moveTo", [(x, y)]))
                    first = False
                else:  # extra pairs after a moveto are implicit linetos
                    ops.append(("lnTo", [(x, y)]))
            ctrl_x = ctrl_y = None
        elif upper == "L":
            for px, py in take(2):
                x, y = (x + px, y + py) if relative else (px, py)
                ops.append(("lnTo", [(x, y)]))
            ctrl_x = ctrl_y = None
        elif upper == "H":
            for (px,) in take(1):
                x = x + px if relative else px
                ops.append(("lnTo", [(x, y)]))
            ctrl_x = ctrl_y = None
        elif upper == "V":
            for (py,) in take(1):
                y = y + py if relative else py
                ops.append(("lnTo", [(x, y)]))
            ctrl_x = ctrl_y = None
        elif upper == "C":
            for c1x, c1y, c2x, c2y, px, py in take(6):
                if relative:
                    c1x, c1y, c2x, c2y, px, py = (
                        x + c1x, y + c1y, x + c2x, y + c2y, x + px, y + py)
                ops.append(("cubicBezTo", [(c1x, c1y), (c2x, c2y), (px, py)]))
                ctrl_x, ctrl_y = c2x, c2y
                x, y = px, py
        elif upper == "S":
            for c2x, c2y, px, py in take(4):
                if relative:
                    c2x, c2y, px, py = x + c2x, y + c2y, x + px, y + py
                # first control mirrors the previous one through the
                # current point; with no previous curve it IS the point
                if ctrl_x is None:
                    c1x, c1y = x, y
                else:
                    c1x, c1y = 2 * x - ctrl_x, 2 * y - ctrl_y
                ops.append(("cubicBezTo", [(c1x, c1y), (c2x, c2y), (px, py)]))
                ctrl_x, ctrl_y = c2x, c2y
                x, y = px, py
        elif upper == "Z":
            ops.append(("close", []))
            x, y = start_x, start_y
            ctrl_x = ctrl_y = None
        numbers = []

    for token_command, number in _tokens(data):
        if token_command is not None:
            flush()
            command = token_command
            if command in "Zz":
                flush()
                command = None
        else:
            numbers.append(number)
    flush()

    if scale != 1.0:
        ops = [(op, [(px * scale, py * scale) for px, py in pts]) for op, pts in ops]
    return ops


def custgeom_xml(ops: list[tuple], scale: int = VIEWBOX) -> etree._Element:
    """DrawingML `<a:custGeom>` for a converted path.

    Every subpath goes into ONE `<a:path>`. That is what makes holes
    work: an icon's counter is a subpath wound the other way, and
    splitting them into separate `<a:path>` elements fills the hole in.
    """
    def el(tag: str, parent=None, **attrs):
        node = etree.SubElement(parent, f"{{{A_NS}}}{tag}") if parent is not None \
            else etree.Element(f"{{{A_NS}}}{tag}")
        for key, value in attrs.items():
            node.set(key, str(value))
        return node

    geom = el("custGeom")
    for empty in ("avLst", "gdLst", "ahLst", "cxnLst"):
        el(empty, geom)
    el("rect", geom, l="0", t="0", r="r", b="b")
    path_list = el("pathLst", geom)
    path = el("path", path_list, w=scale, h=scale)

    for op, points in ops:
        node = el(op, path)
        for px, py in points:
            point = el("pt", node)
            point.set("x", str(int(round(px))))
            point.set("y", str(int(round(py))))
    return geom


# --------------------------------------------------------------------------
# placing one on a slide
# --------------------------------------------------------------------------

KONE_BLUE = "1450F5"


def add_icon(slide, name: str, box, colour: str = KONE_BLUE):
    """Draw a KONE pictogram as an editable shape. Returns it, or None
    if the icon does not exist.

    `box` is (x, y, w, h) in px on the 1280x720 grid, matching every
    other geometry in this project. The icon keeps its aspect ratio --
    they are square, so it is centred in a non-square box.
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    data = load_icons().get(name)
    if data is None:
        return None

    x, y, w, h = box
    side = min(w, h)
    left, top = x + (w - side) / 2, y + (h - side) / 2

    def emu(px):
        return Emu(int(round(px * 9525)))

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(left), emu(top), emu(side), emu(side))
    shape.name = f"Icon {name}"
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(colour)

    # `add_shape` attaches a `<p:style>` referencing the theme's accent
    # fill, line and -- the visible one -- `effectRef idx="2"`, a drop
    # shadow. An empty `<a:effectLst/>` in spPr does not reliably beat
    # it, so the style block goes: everything this shape needs is set
    # explicitly above, and a pictogram with a drop shadow is off-brand.
    element = shape._element
    style = element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style"
    )
    if style is not None:
        element.remove(style)

    # swap the preset rectangle for the icon's own outline
    spPr = shape._element.spPr
    prst = spPr.find(f"{{{A_NS}}}prstGeom")
    geom = custgeom_xml(path_to_drawingml(data))
    if prst is not None:
        spPr.replace(prst, geom)
    else:  # pragma: no cover -- every autoshape has one
        spPr.append(geom)
    return shape
