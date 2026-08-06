"""KONE line-art illustrations as editable PowerPoint shape groups.

Same argument as `icons`, one step harder. The brand ships 60
illustrations as Illustrator-exported SVG, and they must arrive in a
deck as shapes a designer can select, recolour and pull apart -- not as
a flattened picture.

An icon was one path in one colour, so it became one shape. An
illustration is a drawing: a median of 13 paths plus circles, rects,
lines, polylines and polygons, each carrying its own fill and stroke
from a CSS class. So each piece becomes its own shape and the set is
grouped, which is what makes it editable in the way that matters --
ungroup it and the technician's jacket is a shape you can recolour.

What makes it tractable, measured across all 60 before committing to
the approach:

- one shared `0 0 1250 1250` viewBox
- no gradients anywhere, so every fill is a flat colour
- **no elliptical arcs**, exactly as with the icons -- the command set
  is move/line/h/v/cubic/smooth/close, all of which map to DrawingML
- the only transform used is `matrix(...)`, on 16 elements
- no rounded rectangles

The palette is the brand's: black `#141414` and KONE Blue `#1450F5`
fills, black strokes at 2/3/6 units, with white and a yellow accent.

Colour lives in a `<style>` block rather than on the elements, so the
classes are resolved here; a `class` that does not resolve gets SVG's
own defaults (black fill, no stroke) rather than vanishing.
"""

from __future__ import annotations

import math
import re
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import Optional

from deckguard.icons import A_NS, custgeom_xml, path_to_drawingml

VIEWBOX = 1250

_VENDORED = Path(__file__).with_name("assets") / "kone-design" / "illustrations"
_INTERACTIVE = Path("~/.claude/skills/kone-design/assets/illustrations").expanduser()

# A circle drawn as four cubic beziers needs this much control-point
# reach to stay circular; the classic magic number.
_KAPPA = 0.5522847498307936

_STYLE_RE = re.compile(r"<style[^>]*>(.*?)</style>", re.S)
_RULE_RE = re.compile(r"\.([\w-]+)\s*\{([^}]*)\}")
_ELEMENT_RE = re.compile(
    r"<(path|circle|ellipse|rect|line|polyline|polygon|g|/g)\b([^>]*?)(/?)>", re.S
)
_ATTR_RE = re.compile(r'([\w-]+)\s*=\s*"([^"]*)"')
_NUM_RE = re.compile(r"[-+]?(?:\d*\.\d+|\d+)(?:[eE][-+]?\d+)?")


@dataclass(frozen=True)
class Piece:
    """One drawn element of an illustration, in viewBox coordinates."""

    ops: tuple
    fill: Optional[str]  # hex, or None for `fill:none`
    stroke: Optional[str]
    stroke_width: float


def illustrations_dir() -> Optional[Path]:
    for base in (_VENDORED, _INTERACTIVE):
        if base.is_dir() and any(base.glob("*.svg")):
            return base
    return None


@lru_cache(maxsize=1)
def _files() -> dict[str, Path]:
    base = illustrations_dir()
    if base is None:
        return {}
    out = {}
    for path in sorted(base.glob("*.svg")):
        name = path.stem
        for prefix in ("KONE_Illustrations_RGB_", "KONE_Illustrations_"):
            if name.startswith(prefix):
                name = name[len(prefix):]
        out[name.replace("_", "-").replace(" ", "-").lower()] = path
    return out


def illustration_names() -> list[str]:
    return sorted(_files())


def find_illustrations(term: str, limit: int = 12) -> list[str]:
    needle = term.strip().lower().replace(" ", "-").replace("_", "-")
    if not needle:
        return []
    names = illustration_names()
    exact = [n for n in names if n == needle]
    rest = [n for n in names if needle in n and n not in exact]
    return (exact + rest)[:limit]


# --------------------------------------------------------------------------
# parsing
# --------------------------------------------------------------------------


def _css(text: str) -> dict[str, dict]:
    """`.st0{fill:none;stroke:#141414;stroke-width:2}` -> declarations."""
    rules: dict[str, dict] = {}
    for block in _STYLE_RE.findall(text):
        for name, body in _RULE_RE.findall(block):
            declarations = {}
            for declaration in body.split(";"):
                if ":" in declaration:
                    key, _, value = declaration.partition(":")
                    declarations[key.strip().lower()] = value.strip()
            rules.setdefault(name, {}).update(declarations)
    return rules


def _numbers(value: str) -> list[float]:
    return [float(n) for n in _NUM_RE.findall(value or "")]


def _matrix(value: str) -> Optional[tuple]:
    """Only `matrix(a b c d e f)` appears in these files."""
    if not value or "matrix" not in value:
        return None
    parts = _numbers(value)
    return tuple(parts[:6]) if len(parts) >= 6 else None


def _apply(ops, matrix):
    if matrix is None:
        return ops
    a, b, c, d, e, f = matrix
    return [
        (op, [(a * x + c * y + e, b * x + d * y + f) for x, y in points])
        for op, points in ops
    ]


def _ellipse_ops(cx, cy, rx, ry):
    ox, oy = rx * _KAPPA, ry * _KAPPA
    return [
        ("moveTo", [(cx + rx, cy)]),
        ("cubicBezTo", [(cx + rx, cy + oy), (cx + ox, cy + ry), (cx, cy + ry)]),
        ("cubicBezTo", [(cx - ox, cy + ry), (cx - rx, cy + oy), (cx - rx, cy)]),
        ("cubicBezTo", [(cx - rx, cy - oy), (cx - ox, cy - ry), (cx, cy - ry)]),
        ("cubicBezTo", [(cx + ox, cy - ry), (cx + rx, cy - oy), (cx + rx, cy)]),
        ("close", []),
    ]


def _primitive_ops(tag: str, attrs: dict):
    """A shape primitive as path operations. DrawingML has no circle or
    polyline, so everything becomes a path -- which is also what makes
    them editable point-by-point once in PowerPoint."""
    number = lambda key, default=0.0: float(attrs.get(key, default) or 0)  # noqa: E731

    if tag == "circle":
        radius = number("r")
        return _ellipse_ops(number("cx"), number("cy"), radius, radius)
    if tag == "ellipse":
        return _ellipse_ops(number("cx"), number("cy"), number("rx"), number("ry"))
    if tag == "rect":
        x, y, w, h = number("x"), number("y"), number("width"), number("height")
        return [("moveTo", [(x, y)]), ("lnTo", [(x + w, y)]), ("lnTo", [(x + w, y + h)]),
                ("lnTo", [(x, y + h)]), ("close", [])]
    if tag == "line":
        return [("moveTo", [(number("x1"), number("y1"))]),
                ("lnTo", [(number("x2"), number("y2"))])]
    if tag in ("polyline", "polygon"):
        values = _numbers(attrs.get("points", ""))
        points = list(zip(values[0::2], values[1::2]))
        if not points:
            return []
        ops = [("moveTo", [points[0]])] + [("lnTo", [p]) for p in points[1:]]
        if tag == "polygon":
            ops.append(("close", []))
        return ops
    return []


def _resolve(attrs: dict, rules: dict) -> tuple:
    """Fill, stroke and stroke width for an element, from its class,
    then any presentation attributes, then SVG's own defaults."""
    style: dict = {}
    for name in (attrs.get("class") or "").split():
        style.update(rules.get(name, {}))
    for key in ("fill", "stroke", "stroke-width"):
        if key in attrs:
            style[key] = attrs[key]
    inline = attrs.get("style", "")
    for declaration in inline.split(";"):
        if ":" in declaration:
            key, _, value = declaration.partition(":")
            style[key.strip().lower()] = value.strip()

    def colour(value):
        value = (value or "").strip().lower()
        if not value or value == "none":
            return None
        if value.startswith("#"):
            digits = value[1:]
            if len(digits) == 3:
                digits = "".join(c * 2 for c in digits)
            return digits.upper() if len(digits) == 6 else None
        return None

    # SVG's default fill is black; only an explicit `none` clears it
    fill = colour(style["fill"]) if "fill" in style else "000000"
    stroke = colour(style.get("stroke"))
    try:
        width = float(_NUM_RE.findall(style.get("stroke-width", "1"))[0])
    except (IndexError, ValueError):
        width = 1.0
    return fill, stroke, width


@lru_cache(maxsize=64)
def load_illustration(name: str) -> tuple:
    """An illustration as an ordered tuple of `Piece`, in paint order."""
    path = _files().get(name)
    if path is None:
        return ()
    text = path.read_text(errors="replace")
    rules = _css(text)

    pieces: list[Piece] = []
    stack: list = []  # accumulated <g> transforms
    for match in _ELEMENT_RE.finditer(text):
        tag, raw, self_closing = match.group(1), match.group(2), match.group(3)
        if tag == "/g":
            if stack:
                stack.pop()
            continue
        attrs = dict(_ATTR_RE.findall(raw))
        if tag == "g":
            transform = _matrix(attrs.get("transform", ""))
            if not self_closing:
                stack.append(transform)
            continue

        ops = (path_to_drawingml(attrs["d"])
               if tag == "path" and attrs.get("d") else _primitive_ops(tag, attrs))
        if not ops:
            continue
        for transform in [_matrix(attrs.get("transform", ""))] + list(reversed(stack)):
            ops = _apply(ops, transform)

        fill, stroke, width = _resolve(attrs, rules)
        if fill is None and stroke is None:
            continue  # invisible either way
        pieces.append(Piece(tuple(ops), fill, stroke, width))
    return tuple(pieces)


# --------------------------------------------------------------------------
# placing one on a slide
# --------------------------------------------------------------------------


def _bounds(ops, pad: float):
    xs = [x for _, points in ops for x, _ in points]
    ys = [y for _, points in ops for _, y in points]
    if not xs:
        return None
    return (min(xs) - pad, min(ys) - pad, max(xs) + pad, max(ys) + pad)


def artwork_bounds(name: str):
    """Where the drawing actually sits inside the 1250 canvas."""
    pieces = load_illustration(name)
    return _bounds([op for piece in pieces for op in piece.ops], pad=0) if pieces else None


def common_bounds(names) -> Optional[tuple]:
    """One box covering the artwork of several illustrations.

    Pass it to `add_illustration` and a row of drawings shares a single
    scale. Neither default does this: trimming each to its own box makes
    a runner tower over a walking figure, because one drawing happens to
    fill more of the canvas than the other; keeping the full viewBox
    holds them consistent but renders everything tiny, since the art
    uses a fraction of the 1250 square. The union of the set is both
    consistent and close-cropped.
    """
    boxes = [b for b in (artwork_bounds(name) for name in names) if b]
    if not boxes:
        return None
    return (min(b[0] for b in boxes), min(b[1] for b in boxes),
            max(b[2] for b in boxes), max(b[3] for b in boxes))


def add_illustration(slide, name: str, box, group: bool = True, trim: bool = True,
                     bounds: Optional[tuple] = None):
    """Draw a KONE illustration as editable shapes. Returns the group
    (or the list of shapes when `group` is False), or None if unknown.

    `box` is (x, y, w, h) in px on the 1280x720 grid. The drawing keeps
    its aspect ratio and is centred, and stroke widths scale with it --
    a 6-unit outline drawn at a quarter size has to become 1.5, or the
    line art thickens into a blob.

    `trim` scales the DRAWING to the box rather than its viewBox. These
    files are Illustrator exports on a shared 1250 canvas and the art
    uses a fraction of it -- the cloud is 146x64 units in the middle of
    it -- so honouring the viewBox puts a postage stamp in the middle of
    whatever box you asked for. Pass `trim=False` to keep the original
    framing.

    For a ROW of illustrations, neither default is right: trimming each
    to its own box makes a runner tower over a walking figure, and the
    full viewBox holds them consistent but tiny. Pass `bounds` from
    `common_bounds(names)` and the whole set shares one scale.
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu, Pt

    pieces = load_illustration(name)
    if not pieces:
        return None

    x, y, w, h = box
    art = bounds or (
        _bounds([op for piece in pieces for op in piece.ops], pad=0) if trim else None
    )
    if art:
        art_x, art_y, art_right, art_bottom = art
        art_w, art_h = max(art_right - art_x, 1e-6), max(art_bottom - art_y, 1e-6)
        scale = min(w / art_w, h / art_h)
        origin_x = x + (w - art_w * scale) / 2 - art_x * scale
        origin_y = y + (h - art_h * scale) / 2 - art_y * scale
    else:
        scale = min(w / VIEWBOX, h / VIEWBOX)
        side = VIEWBOX * scale
        origin_x, origin_y = x + (w - side) / 2, y + (h - side) / 2

    def emu(px):
        return Emu(int(round(px * 9525)))

    container = slide.shapes.add_group_shape() if group else None
    target = container.shapes if container is not None else slide.shapes

    drawn = []
    for index, piece in enumerate(pieces):
        bounds = _bounds(piece.ops, pad=piece.stroke_width / 2 if piece.stroke else 0)
        if bounds is None:
            continue
        left, top, right, bottom = bounds
        width, height = max(right - left, 1e-6), max(bottom - top, 1e-6)

        shape = target.add_shape(
            MSO_SHAPE.RECTANGLE,
            emu(origin_x + left * scale), emu(origin_y + top * scale),
            emu(width * scale), emu(height * scale),
        )
        shape.name = f"{name}-{index}"

        # geometry expressed in this piece's own box, so the shape's
        # selection handles match what it actually draws
        local = [(op, [(px - left, py - top) for px, py in points]) for op, points in piece.ops]
        spPr = shape._element.spPr
        prst = spPr.find(f"{{{A_NS}}}prstGeom")
        geom = custgeom_xml(local, scale=int(math.ceil(max(width, height))))
        # custgeom_xml squares the path box; restore the true extents
        path_el = geom.find(f".//{{{A_NS}}}path")
        path_el.set("w", str(int(round(width))))
        path_el.set("h", str(int(round(height))))
        if prst is not None:
            spPr.replace(prst, geom)
        else:  # pragma: no cover
            spPr.append(geom)

        if piece.fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(piece.fill)
        else:
            shape.fill.background()
        if piece.stroke:
            shape.line.color.rgb = RGBColor.from_string(piece.stroke)
            shape.line.width = Pt(max(piece.stroke_width * scale * 0.75, 0.25))
        else:
            shape.line.fill.background()

        shape.shadow.inherit = False
        style = shape._element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}style"
        )
        if style is not None:
            shape._element.remove(style)
        drawn.append(shape)

    if container is not None:
        container.name = f"Illustration {name}"
        _fit_group(container, drawn)
        return container
    return drawn


def _fit_group(container, children) -> None:
    """Give a group extents its children actually fit.

    A group carries TWO coordinate systems: `off`/`ext` -- where it sits
    on the slide -- and `chOff`/`chExt`, the space its children's
    coordinates are expressed in. PowerPoint scales one onto the other.
    Setting the group's size without setting the child space makes that
    ratio arbitrary: a cloud written correctly as a 30x13 shape rendered
    eight times too large, because `ext` said 250x250 while `chExt`
    still described the children's own bounding box.

    Setting both to the same rectangle makes the ratio exactly 1, so
    children land where their coordinates say. The file was numerically
    perfect either way -- only a render showed it.
    """
    if not children:
        return
    left = min(child.left for child in children)
    top = min(child.top for child in children)
    right = max(child.left + child.width for child in children)
    bottom = max(child.top + child.height for child in children)

    xfrm = container._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr"
    ).find("{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
    for tag, (x, y) in (
        ("off", (left, top)), ("chOff", (left, top)),
        ("ext", (right - left, bottom - top)), ("chExt", (right - left, bottom - top)),
    ):
        node = xfrm.find(f"{{{A_NS}}}{tag}")
        if node is None:
            continue
        if tag in ("off", "chOff"):
            node.set("x", str(int(x)))
            node.set("y", str(int(y)))
        else:
            node.set("cx", str(int(x)))
            node.set("cy", str(int(y)))
