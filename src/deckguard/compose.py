"""Outline-driven deck composition: generate a new, on-brand .pptx --
or append fresh, on-brand slides to an existing one -- from a
structured content outline, using the org master template's own
sanctioned layouts.

This is the forward complement to `retemplate` (which rebuilds an old
deck's slides onto approved layouts after the fact): content built here
never touches an unapproved layout, color, or font in the first place,
so a composed deck needs no `fix` pass -- `audit` on its own output
should come back clean by construction, the same guarantee `retemplate`
gives its rebuilt slides. The layout-fit algorithm is reused as-is
(`SlideProfile`/`LayoutProfile`/`match_layout` from `retemplate.py`): a
content/stat/timeline slide's shape is described the same way whether
it's being read out of an old deck or built from scratch, so the same
matcher that picks a rebuild target for a legacy slide also picks the
tightest-fitting layout for a brand-new one.

Slide "kinds":
    cover / agenda / section / quote / statement / end / blank -- one
        fixed, configurable-variant layout each (a natural single best
        shape -- no fitting problem to solve).
    content / stat / timeline -- the number of body blocks needed
        varies with the content supplied, so these go through
        `match_layout` against a per-kind candidate list, exactly like
        `retemplate`'s ordinary-content slides.

Logo/tagline placeholders are never written to. A placeholder shape
left with an empty `<p:spPr>`/`<p:txBody>` fully inherits the layout's
own fill -- confirmed against the shipped template: Cover A's Logo
Placeholder carries its `<a:blipFill>` at the LAYOUT level, and a
slide-level placeholder of the same idx with no properties of its own
renders that same image. So leaving a chrome placeholder alone is what
makes the logo show up, not a gap to work around. Placeholders are
matched to this rule by NAME (containing "logo" or "tagline",
case-insensitive) rather than by idx, since idx numbering isn't stable
across layouts but these names are.

DATE/FOOTER/SLIDE_NUMBER placeholders are deliberately left off
generated slides entirely -- same as what python-pptx's own
`add_slide()` does, since those are the presentation-level
Insert-Header/Footer placeholders, not per-slide content. Restoring
per-slide date/footer chrome is a known gap, not attempted here.

Every generated slide's text is written with NO explicit run color --
it inherits the theme/placeholder default, same as typing straight
into a template in PowerPoint. That inheritance is exactly what
`rules_engine`'s `text_contrast` rule can't see through (it requires an
explicit, resolvable color to consider a run compliant -- confirmed by
running `audit` against the shipped master template's OWN sample
slides: 91 pre-existing "major" findings, entirely of this shape, on
content nobody but this project's own template authors wrote). Rather
than duplicate the org's dark/white/tint contrast rules a second time
here, `build_deck` runs its own output through the exact `fix_deck`
engine `deckguard fix` uses, once, before saving -- the same function
that already turns "inherited, therefore unresolved" into an explicit,
correct color for a legacy deck does the identical job for a brand-new
one. `all_caps` findings (e.g. a quarter label like "Q3 2026" reads as
all-caps to that rule, same as it would in any hand-typed deck) are
never auto-rewritten by `fix_deck` either, by design -- they surface in
`ComposeResult.manual_review` same as `deckguard fix`'s own changelog.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Optional

import yaml
from pptx import Presentation
from pptx.util import Pt

from deckguard import effects as effects_mod
from deckguard.config import default_config_path, load_config
from deckguard.fixer import fix_deck
from deckguard.retemplate import SlideProfile, _layout_profile, _set_text_block, match_layout
from deckguard.slide_import import _delete_slide, default_template_path, import_layouts

CHROME_NAME_RE = re.compile(r"logo|tagline", re.IGNORECASE)

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

TITLE_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}
BODY_PLACEHOLDER_TYPES = {"BODY", "OBJECT", "SUBTITLE"}

KNOWN_KINDS = {
    "cover", "agenda", "section", "content", "quote", "statement",
    "stat", "timeline", "end", "blank",
}

# Kinds with one natural layout choice -- no content-driven fitting needed.
FIXED_LAYOUT_BY_KIND = {
    "agenda": "Agenda A",
    "statement": "Statement or quote",
    "end": "Outro",
    "blank": "Blank",
}

COVER_VARIANTS = {chr(c): f"Cover {chr(c)}" for c in range(ord("A"), ord("F") + 1)}
QUOTE_VARIANTS = {chr(c): f"Quote {chr(c)}" for c in range(ord("A"), ord("E") + 1)}
SECTION_VARIANTS = {
    "plain": "Section divider (just title)",
    "numbered": "Section divider (numbering)",
    "A": "Section divider A",
    "B": "Section divider B",
    "C": "Section divider C",
    "D": "Section divider D",
}

# Ordered simplest-first, same convention retemplate.py uses for its own
# candidate list -- ties in match_layout favor whichever appears earliest.
CONTENT_LAYOUT_CANDIDATES = [
    "Title and content A", "Title and content B",
    "Two content A", "Two content B", "Two content C", "Two content D",
    "Three content A", "Three content B", "Three content C", "Three content D",
    # Picture-carrying candidates -- only ever picked over the text-only
    # ones above when the slide actually needs images: match_layout scores
    # by wasted body/picture slots, and a text-only slide with n_images=0
    # always scores strictly better on a text-only layout (0 unused
    # picture slots) than on any of these, so this list is additive, not a
    # change in behavior for slides that don't carry images.
    "Text and picture B", "Text and picture C", "Text and picture G", "Text and picture H",
    "Text and picture A", "Text and picture F",
    "Two pictures and text B", "Two pictures and text C", "Two pictures and text A",
]

STAT_LAYOUT_CANDIDATES = [
    "Statement B",
    "Two content A", "Three content A",
    "REPORT Three content", "REPORT Five content", "REPORT Seven content",
]

TIMELINE_LAYOUT_CANDIDATES = [
    "Two content A", "Three content A",
    "REPORT Three content", "REPORT Five content", "REPORT Seven content",
]


@dataclass
class SlideSpec:
    kind: str
    title: Optional[str] = None
    subtitle: Optional[str] = None
    bullets: list = field(default_factory=list)  # one text block: list[str | {"level": int, "text": str}]
    columns: list = field(default_factory=list)  # multiple text blocks, one per body placeholder, same shape as bullets
    quote_text: Optional[str] = None
    quote_author: Optional[str] = None
    quote_label: Optional[str] = None
    stats: list = field(default_factory=list)  # [{"number": "40M", "label": "..."}]
    milestones: list = field(default_factory=list)  # [{"label": "2024", "text": "..."}]
    images: list = field(default_factory=list)  # file paths or raw image bytes, filled into PICTURE placeholders in reading order
    variant: Optional[str] = None  # picks among Cover A-F / Quote A-E / Section variants
    layout: Optional[str] = None  # escape hatch: force an exact layout name, bypassing kind-based selection


@dataclass
class Outline:
    slides: list  # list[SlideSpec]


@dataclass
class ComposeResult:
    slide_count: int
    layouts_used: list  # list[str], one per slide, in order
    manual_review: list  # violations fix_deck's own final pass couldn't auto-resolve (e.g. all_caps)


class ComposeError(ValueError):
    """Raised for an invalid outline or a slide that no layout can hold."""


def _normalize_block(raw: list) -> list:
    """list[str | {"level": int, "text": str}] -> list[(level, text)]."""
    block = []
    for item in raw:
        if isinstance(item, str):
            block.append((0, item))
        elif isinstance(item, dict):
            block.append((int(item.get("level", 0)), str(item.get("text", ""))))
        else:
            raise ComposeError(f"invalid bullet item: {item!r} (expected a string or a {{level, text}} mapping)")
    return block


def _slide_spec_from_dict(raw, position: int) -> SlideSpec:
    if not isinstance(raw, dict) or "kind" not in raw:
        raise ComposeError(f"slide {position}: each entry needs a 'kind'")
    kind = raw["kind"]
    if kind not in KNOWN_KINDS:
        raise ComposeError(f"slide {position}: unknown kind {kind!r} (expected one of {sorted(KNOWN_KINDS)})")
    return SlideSpec(
        kind=kind,
        title=raw.get("title"),
        subtitle=raw.get("subtitle"),
        bullets=raw.get("bullets", []) or [],
        columns=raw.get("columns", []) or [],
        quote_text=raw.get("quote_text"),
        quote_author=raw.get("quote_author"),
        quote_label=raw.get("quote_label"),
        stats=raw.get("stats", []) or [],
        milestones=raw.get("milestones", []) or [],
        images=raw.get("images", []) or [],
        variant=raw.get("variant"),
        layout=raw.get("layout"),
    )


def outline_from_list(raw_slides: list) -> Outline:
    """Build an Outline from already-parsed slide dicts -- the shared tail
    end of `load_outline` (after YAML parsing) and of `redesign.py`'s AI
    path (after parsing the model's structured-output JSON). Both producers
    speak the identical schema, so a human-written outline and an
    AI-generated one are indistinguishable to `build_deck` from here on.
    """
    if not isinstance(raw_slides, list) or not raw_slides:
        raise ComposeError("outline.slides must be a non-empty list")
    return Outline(slides=[_slide_spec_from_dict(s, i) for i, s in enumerate(raw_slides, start=1)])


def load_outline(path) -> Outline:
    """Parse a YAML content outline: a top-level 'slides' list, each
    entry a mapping with at least a 'kind'. See README for the schema."""
    try:
        text = Path(path).read_text(encoding="utf-8")
    except FileNotFoundError as exc:
        raise ComposeError(f"outline not found: {path}") from exc
    try:
        data = yaml.safe_load(text)
    except yaml.YAMLError as exc:
        raise ComposeError(f"invalid YAML in {path}: {exc}") from exc

    if not isinstance(data, dict) or "slides" not in data:
        raise ComposeError("outline must be a YAML mapping with a top-level 'slides' list")
    return outline_from_list(data["slides"])


def _chrome_idxs(layout) -> set:
    """Placeholder idxs whose NAME marks them as logo/tagline chrome --
    these are never written to, so their empty <p:spPr>/<p:txBody>
    keeps inheriting the layout's own baked-in image. See module
    docstring for why this is correct rather than an omission."""
    return {
        ph.placeholder_format.idx
        for ph in layout.placeholders
        if CHROME_NAME_RE.search(ph.name or "")
    }


def _content_bodies(slide, chrome_idxs) -> list:
    """Non-chrome BODY/OBJECT/SUBTITLE placeholders, reading-order (top then left)."""
    bodies = [
        ph for ph in slide.placeholders
        if ph.placeholder_format.idx not in chrome_idxs
        and ph.placeholder_format.type is not None
        and ph.placeholder_format.type.name in BODY_PLACEHOLDER_TYPES
    ]
    bodies.sort(key=lambda ph: (ph.top or 0, ph.left or 0))
    return bodies


def _content_pictures(slide, chrome_idxs) -> list:
    pics = [
        ph for ph in slide.placeholders
        if ph.placeholder_format.idx not in chrome_idxs
        and ph.placeholder_format.type is not None
        and ph.placeholder_format.type.name == "PICTURE"
    ]
    pics.sort(key=lambda ph: (ph.top or 0, ph.left or 0))
    return pics


def _fill_title_subtitle(slide, chrome_idxs, title, subtitle) -> None:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in chrome_idxs:
            continue
        ph_type = ph.placeholder_format.type
        name = ph_type.name if ph_type else None
        if name in TITLE_PLACEHOLDER_TYPES and title:
            ph.text_frame.text = title
        elif name == "SUBTITLE" and subtitle:
            ph.text_frame.text = subtitle


def _layout_placeholder_image_blob(layout, idx: int) -> Optional[bytes]:
    """The image bytes behind a layout's own picture placeholder at
    `idx`, if it has one baked in via `<a:blipFill><a:blip>` -- e.g. the
    org template's default cover-photo stock image, which any slide-level
    placeholder of the same idx shows through inheritance as long as its
    own `<p:spPr>` stays empty. Used to materialize a REAL, independently
    editable picture on the slide when no source image is available to
    carry over (see `_fill_images`): PowerPoint only offers "Change
    Picture" for an actual picture object belonging to the slide itself,
    never for one it's merely inheriting from its layout -- an unfilled
    placeholder left that way only ever offers "Save as Picture" on the
    rendered (layout-owned) image, same limitation the raw org template
    itself has for an untouched Cover/Outro/End slide.
    """
    for ph in layout.placeholders:
        if ph.placeholder_format.idx != idx:
            continue
        blip = ph._element.find(".//" + effects_mod.a_qn("blip"))
        if blip is None:
            return None
        rid = blip.get(f"{{{_R_NS}}}embed")
        if not rid:
            return None
        try:
            return layout.part.related_part(rid).blob
        except KeyError:
            return None
    return None


def _fill_images(slide, chrome_idxs, images, fallback_to_layout_default: bool = False) -> None:
    pics = _content_pictures(slide, chrome_idxs)
    layout = slide.slide_layout if fallback_to_layout_default else None
    for i, ph in enumerate(pics):
        if i < len(images):
            image = images[i]
            ph.insert_picture(BytesIO(image) if isinstance(image, (bytes, bytearray)) else str(image))
        elif layout is not None:
            # No source image for this placeholder -- give it the
            # layout's own default photo as a real, editable picture
            # instead of leaving it empty (see
            # `_layout_placeholder_image_blob`'s own docstring for why
            # that would silently break "Change Picture" in PowerPoint).
            # Scoped to cover/end (see call site) -- a content slide
            # with fewer supplied images than its layout's picture slots
            # should just show fewer pictures, not get padded out with
            # repeated stock photos in every leftover slot.
            default_blob = _layout_placeholder_image_blob(layout, ph.placeholder_format.idx)
            if default_blob is not None:
                ph.insert_picture(BytesIO(default_blob))


def _fill_columns(slide, chrome_idxs, blocks: list) -> None:
    for ph, block in zip(_content_bodies(slide, chrome_idxs), blocks):
        if block:
            _set_text_block(ph.text_frame, block)


def _fill_quote(slide, chrome_idxs, quote_text, author, label) -> None:
    bodies = _content_bodies(slide, chrome_idxs)
    if not bodies:
        return
    # The big quote is whichever body placeholder has the most area; the
    # attribution line is whichever is shortest; anything left over (an
    # eyebrow/label placeholder in the Quote A-E layouts) takes `label`.
    quote_ph = max(bodies, key=lambda ph: (ph.width or 0) * (ph.height or 0))
    remaining = [ph for ph in bodies if ph is not quote_ph]
    attribution_ph = min(remaining, key=lambda ph: (ph.height or 0)) if remaining else None
    label_ph = next((ph for ph in remaining if ph is not attribution_ph), None)

    if quote_text:
        quote_ph.text_frame.text = quote_text
    if attribution_ph is not None and author:
        attribution_ph.text_frame.text = author
    if label_ph is not None and label:
        label_ph.text_frame.text = label


def _fill_number_blocks(slide, chrome_idxs, items: list, number_key: str, label_key: str, number_size_pt: int) -> None:
    """Shared by stat/timeline: one block per item, a bold/enlarged
    first line (the number or milestone label) plus a plain second
    line -- never a hand-set color, so it stays theme-inherited."""
    blocks = []
    for item in items:
        headline = str(item.get(number_key, "")).strip()
        detail = str(item.get(label_key, "")).strip()
        blocks.append([(0, headline), (1, detail)] if detail else [(0, headline)])
    bodies = _content_bodies(slide, chrome_idxs)
    for ph, block in zip(bodies, blocks):
        if not block:
            continue
        _set_text_block(ph.text_frame, block)
        headline_para = ph.text_frame.paragraphs[0]
        for run in headline_para.runs:
            run.font.size = Pt(number_size_pt)
            run.font.bold = True


def _fill_stats(slide, chrome_idxs, stats: list) -> None:
    _fill_number_blocks(slide, chrome_idxs, stats, "number", "label", number_size_pt=40)


def _fill_timeline(slide, chrome_idxs, milestones: list) -> None:
    _fill_number_blocks(slide, chrome_idxs, milestones, "label", "text", number_size_pt=20)


def _fix_new_slides(prs, config: dict, slide_indices: set) -> list:
    """Run just the per-run/shape auto-fix pass -- no theme/master/logo
    remapping -- over ONLY the given 1-based slide indices.

    Used for the append path instead of the full `fix_deck`: the newly
    imported layouts already carry the org template's own
    already-compliant theme verbatim (import_layouts copies it as-is),
    so there's no legacy theme/master color or font to remap in the
    first place -- and running that global pass over the WHOLE file
    would also touch the existing deck's own pre-existing theme/master,
    breaking the "existing content is left completely untouched"
    guarantee. Scoping to just the new slides' own violations keeps
    that guarantee while still resolving their inherited-not-explicit
    run colors the same way `fix_deck` would.
    """
    from deckguard.fixer import _apply_violation
    from deckguard.inventory import build_inventory
    from deckguard.rules_engine import audit_deck, sort_violations

    for _ in range(5):
        violations = sort_violations(audit_deck(build_inventory(prs), config))
        scoped = [v for v in violations if v.slide_index in slide_indices and v.auto_fixable]
        applied = [_apply_violation(v, config) for v in scoped]
        if not any(applied):
            break
    return [
        v for v in sort_violations(audit_deck(build_inventory(prs), config))
        if v.slide_index in slide_indices
    ]


def _populate(new_slide, layout, spec: SlideSpec) -> None:
    chrome_idxs = _chrome_idxs(layout)
    _fill_title_subtitle(new_slide, chrome_idxs, spec.title, spec.subtitle)
    _fill_images(new_slide, chrome_idxs, spec.images, fallback_to_layout_default=spec.kind in ("cover", "end"))

    if spec.kind == "quote":
        _fill_quote(new_slide, chrome_idxs, spec.quote_text, spec.quote_author, spec.quote_label)
    elif spec.kind == "stat":
        _fill_stats(new_slide, chrome_idxs, spec.stats)
    elif spec.kind == "timeline":
        _fill_timeline(new_slide, chrome_idxs, spec.milestones)
    elif spec.kind == "content":
        blocks = [_normalize_block(col) for col in spec.columns] if spec.columns else (
            [_normalize_block(spec.bullets)] if spec.bullets else []
        )
        _fill_columns(new_slide, chrome_idxs, blocks)
    elif spec.bullets:
        # agenda / section / statement / end / cover: a single supporting
        # text block, if the outline gave one.
        _fill_columns(new_slide, chrome_idxs, [_normalize_block(spec.bullets)])


def _content_candidates(spec: SlideSpec) -> tuple:
    if spec.kind == "content":
        n_blocks = max(1, len(spec.columns) or (1 if spec.bullets else 0))
        return CONTENT_LAYOUT_CANDIDATES, n_blocks, len(spec.images)
    if spec.kind == "stat":
        return STAT_LAYOUT_CANDIDATES, max(1, len(spec.stats)), 0
    if spec.kind == "timeline":
        return TIMELINE_LAYOUT_CANDIDATES, max(1, len(spec.milestones)), 0
    raise ComposeError(f"unknown slide kind {spec.kind!r}")  # pragma: no cover -- guarded by caller


def _select_layout(spec: SlideSpec, layouts_by_name: dict, layout_profile_cache: dict) -> str:
    if spec.layout:
        if spec.layout not in layouts_by_name:
            raise ComposeError(f"layout {spec.layout!r} does not exist in the template")
        return spec.layout

    if spec.kind == "cover":
        variant = (spec.variant or "B").upper()
        if variant not in COVER_VARIANTS:
            raise ComposeError(f"cover variant {spec.variant!r} is not one of A-F")
        return COVER_VARIANTS[variant]
    if spec.kind == "quote":
        variant = (spec.variant or "A").upper()
        if variant not in QUOTE_VARIANTS:
            raise ComposeError(f"quote variant {spec.variant!r} is not one of A-E")
        return QUOTE_VARIANTS[variant]
    if spec.kind == "section":
        variant = spec.variant or "plain"
        if variant not in SECTION_VARIANTS:
            raise ComposeError(f"section variant {spec.variant!r} is not one of {sorted(SECTION_VARIANTS)}")
        return SECTION_VARIANTS[variant]
    if spec.kind in FIXED_LAYOUT_BY_KIND:
        return FIXED_LAYOUT_BY_KIND[spec.kind]

    candidates, n_blocks, n_images = _content_candidates(spec)
    cache_key = tuple(candidates)
    layout_profiles = layout_profile_cache.get(cache_key)
    if layout_profiles is None:
        layout_profiles = [_layout_profile(layouts_by_name[name]) for name in candidates if name in layouts_by_name]
        layout_profile_cache[cache_key] = layout_profiles

    profile = SlideProfile(title=spec.title, text_blocks=[[(0, "x")]] * n_blocks, images=[b""] * n_images, eligible=True)
    match = match_layout(profile, layout_profiles)
    if match is None:
        raise ComposeError(
            f"no layout fits a {spec.kind!r} slide needing {n_blocks} text block(s) and {n_images} image(s)"
        )
    return match.layout_name


def build_deck(
    outline: Outline, out_path, template_path=None, existing_deck_path=None, rules_config: Optional[dict] = None
) -> ComposeResult:
    """Build every slide in `outline` onto the org template's own
    layouts, writing the result to `out_path`.

    With `existing_deck_path` unset (the default), starts fresh: a copy
    of the template itself with its own sample slides stripped out (the
    template already carries every layout/master/theme it needs, so no
    cross-file import is required). With `existing_deck_path` set, the
    outline's slides are APPENDED to a copy of that deck instead --
    needed layouts are imported into it first (via the same
    `import_layouts` retemplate.py uses), so the existing content is
    left completely untouched.

    The result is run through `fix_deck` (same engine/config `deckguard
    fix` uses, defaulting to the packaged KONE rules) before being
    saved, so every run's inherited-not-explicit text color is resolved
    to the correct brand color in one pass -- see the module docstring.
    """
    template_path = Path(template_path) if template_path else default_template_path()
    template_prs = Presentation(str(template_path))
    layouts_by_name = {layout.name: layout for master in template_prs.slide_masters for layout in master.slide_layouts}

    layout_profile_cache: dict = {}
    chosen_layouts = [_select_layout(spec, layouts_by_name, layout_profile_cache) for spec in outline.slides]
    config = rules_config if rules_config is not None else load_config(default_config_path())

    if existing_deck_path:
        needed_layouts = sorted(set(chosen_layouts))
        tmp_path = Path(out_path).with_suffix(".layouts.pptx")
        import_layouts(str(existing_deck_path), str(template_path), str(tmp_path), needed_layouts)
        try:
            prs = Presentation(str(tmp_path))
            base_slide_count = len(prs.slides)
            imported_master = prs.slide_masters[-1]
            target_layouts_by_name = {layout.name: layout for layout in imported_master.slide_layouts}
            for spec, layout_name in zip(outline.slides, chosen_layouts):
                new_slide = prs.slides.add_slide(target_layouts_by_name[layout_name])
                _populate(new_slide, target_layouts_by_name[layout_name], spec)
            new_slide_indices = set(range(base_slide_count + 1, len(prs.slides) + 1))
            manual_review = _fix_new_slides(prs, config, new_slide_indices)
            prs.save(out_path)
        finally:
            tmp_path.unlink(missing_ok=True)
    else:
        prs = Presentation(str(template_path))
        for i in range(len(prs.slides) - 1, -1, -1):
            _delete_slide(prs, i)
        own_layouts_by_name = {layout.name: layout for master in prs.slide_masters for layout in master.slide_layouts}
        for spec, layout_name in zip(outline.slides, chosen_layouts):
            new_slide = prs.slides.add_slide(own_layouts_by_name[layout_name])
            _populate(new_slide, own_layouts_by_name[layout_name], spec)
        fix_report = fix_deck(prs, config, source_path=str(out_path), output_path=str(out_path), dry_run=False)
        manual_review = fix_report.manual_review

    return ComposeResult(slide_count=len(outline.slides), layouts_used=chosen_layouts, manual_review=manual_review)
