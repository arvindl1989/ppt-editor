"""The master template's own geometry, read as data.

Every archetype in this tool used to be ported by hand -- read the
reference, eyeball the boxes, type a region dict, repeat. That got us 22
of the 61 archetypes the brand actually documents, and the 39 missing
ones were not hard, just numerous.

They do not have to be typed at all. `LAYOUTS.md` already lists the px
geometry of all 62 master layouts in a rigidly regular form::

    ### Title and content A
    `slideLayout16`

    - **Title** — 45, 91, 917 × 104
    - **Text/body** — 45, 227, 917 × 402
    - **Footer** — 215, 658, 747 × 19
    - **Logo** — 1153, 45, 81 × 31 · image

and `ARCHETYPES.md` binds each canonical `UPPER_SNAKE` archetype name to
one of those layouts. Between them that is 51 of the 61 archetypes,
specified, in files we already ship. This module reads both and emits
engine archetype dicts, so adding the rest of the vocabulary is a
parsing problem rather than a typing one.

Two measurements justify trusting the files this much. Every one of the
321 geometry lines in `LAYOUTS.md` parses -- there is no ragged tail to
hand-correct. And 316 of those 321 boxes are present at exactly that
pixel in the master .pptx itself; the five that differ are two logo
rectangles off by a pixel and three shapes on the user-guide layout.
`LAYOUTS.md` is not a description of the master, it *is* the master.

What this module deliberately does NOT do:

- It never emits chrome. Logo, tagline, date and page number are the
  layout's job, not the archetype's -- an archetype that draws its own
  logo produces two of them the moment the master's frames are
  repaired. `CHROME_ROLES` is dropped on the floor during parsing.
- It never overwrites a hand-built archetype. `install()` skips any
  name (or alias) that the engine or the gallery port already
  implements, because those were tuned against a real rendering and
  this module is working from a coarser description. Generated geometry
  is the fallback, not the authority.

The role vocabulary in `LAYOUTS.md` is coarse on purpose -- `Text/body`
covers 109 of the 321 boxes and stands for eyebrows, bullet lists,
captions and stat columns alike. `_bind_roles` recovers the structure
that the flat list loses, mainly by noticing that equal-sized boxes at
one y with evenly spaced x are not three regions but one repeating
group of three.
"""

from __future__ import annotations

import copy
import os
import re
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import Iterable, Optional

_INTERACTIVE_SKILL_DIR = "~/.claude/skills/kone-design"
_VENDORED_SKILL_DIR = Path(__file__).with_name("assets") / "kone-design"

# Roles the LAYOUT owns. An archetype must never place these itself.
CHROME_ROLES = frozenset({"Logo", "Tagline", "Footer", "Date", "Page number"})

# `- **Title** — 45, 91, 917 × 104 · black`
_BOX_RE = re.compile(
    r"^- \*\*(?P<role>.+?)\*\* — (?P<x>-?\d+), (?P<y>-?\d+), (?P<w>\d+) × (?P<h>\d+)(?P<mods>.*)$",
    re.M,
)
_LAYOUT_KEY_RE = re.compile(r"`(slideLayout\d+)`")

# `| `TITLE_CONTENT` | The workhorse... | `slideLayout16` | built · 04 |`
_ARCH_ROW_RE = re.compile(
    r"^\| `(?P<name>[A-Z][A-Z0-9_]+)` \| (?P<purpose>.*?) \| (?P<master>.*?) \| (?P<status>[^|]+) \|$",
    re.M,
)
_ALIAS_RE = re.compile(r"Alias: `([a-z0-9_]+)`")
_GRADE_RE = re.compile(r"^## Grade ([A-D])", re.M)
_TWIN_RE = re.compile(r"Twin of `([A-Z][A-Z0-9_]+)`")


@dataclass(frozen=True)
class Box:
    """One positioned shape on the 1280x720 grid."""

    role: str
    x: int
    y: int
    w: int
    h: int
    mods: tuple[str, ...] = ()

    @property
    def is_chrome(self) -> bool:
        return self.role in CHROME_ROLES

    @property
    def is_decoration(self) -> bool:
        """`Rectangle 9`, `Fast overskrift` -- shapes the master names
        after itself rather than after a role. Nothing can be bound to
        them, so they are not archetype content."""
        return bool(re.match(r"^(Rectangle|Picture Placeholder|Fast) ", self.role))

    @property
    def on_dark(self) -> bool:
        return "white" in self.mods


@dataclass(frozen=True)
class Layout:
    """A master layout: `slideLayout16`, its title, and its boxes."""

    key: str
    name: str
    boxes: tuple[Box, ...]

    def content_boxes(self) -> list[Box]:
        return [b for b in self.boxes if not b.is_chrome and not b.is_decoration]


@dataclass(frozen=True)
class Archetype:
    """One row of `ARCHETYPES.md`."""

    name: str
    grade: str
    purpose: str
    master: Optional[str]
    aliases: tuple[str, ...]
    status: str
    twin_of: Optional[str] = None

    @property
    def is_built(self) -> bool:
        return self.status.startswith("built")

    @property
    def is_twin(self) -> bool:
        return self.status.startswith("twin")

    @property
    def engine_key(self) -> str:
        return self.name.lower()


def _is_graded(spec: Path) -> bool:
    """Does this directory hold the CURRENT spec?

    An installed kone-design skill may predate the rework and still
    carry an `ARCHETYPES.md` that lists the old gallery vocabulary with
    no grades -- parsing that yields zero archetypes and silently
    disables this whole module. Both files must be present and
    `ARCHETYPES.md` must be the graded one, or we fall back.
    """
    layouts, archetypes = spec / "LAYOUTS.md", spec / "ARCHETYPES.md"
    if not (layouts.is_file() and archetypes.is_file()):
        return False
    return bool(_GRADE_RE.search(archetypes.read_text()))


@lru_cache(maxsize=1)
def spec_dir() -> Path:
    """Where `LAYOUTS.md` / `ARCHETYPES.md` live -- same three-step
    resolution as `skill_bridge._skill_dir`, for the same reason: a
    deployed deckguard has no `~/.claude/skills/` at all. The vendored
    copy is the last candidate AND the fallback, so a stale installed
    skill degrades to the bundled spec rather than to nothing."""
    candidates = []
    env = os.environ.get("KONE_DESIGN_DIR")
    if env:
        candidates.append(Path(env).expanduser())
    candidates.append(Path(_INTERACTIVE_SKILL_DIR).expanduser())
    candidates.append(_VENDORED_SKILL_DIR)

    for base in candidates:
        spec = base / "templates" / "kone-deck"
        if _is_graded(spec):
            return spec
    return _VENDORED_SKILL_DIR / "templates" / "kone-deck"


# --------------------------------------------------------------------------
# parsing
# --------------------------------------------------------------------------


def _mods(raw: str) -> tuple[str, ...]:
    return tuple(p.strip() for p in raw.split("·") if p.strip())


def parse_layouts(text: str) -> dict[str, Layout]:
    """`LAYOUTS.md` -> {slideLayoutN: Layout}. Sections without a
    `slideLayoutN` marker (the grid preamble, the grades table) are
    prose, not geometry, and are skipped."""
    out: dict[str, Layout] = {}
    for section in re.split(r"^### ", text, flags=re.M)[1:]:
        key_match = _LAYOUT_KEY_RE.search(section)
        if not key_match:
            continue
        boxes = tuple(
            Box(m.group("role").strip(), int(m.group("x")), int(m.group("y")),
                int(m.group("w")), int(m.group("h")), _mods(m.group("mods")))
            for m in _BOX_RE.finditer(section)
        )
        name = section.splitlines()[0].strip()
        out[key_match.group(1)] = Layout(key_match.group(1), name, boxes)
    return out


def parse_archetypes(text: str) -> dict[str, Archetype]:
    """`ARCHETYPES.md` -> {NAME: Archetype}, carrying the grade down
    from whichever `## Grade X` heading the row sits under."""
    out: dict[str, Archetype] = {}
    parts = _GRADE_RE.split(text)
    for i in range(1, len(parts), 2):
        grade, body = parts[i], parts[i + 1]
        for m in _ARCH_ROW_RE.finditer(body):
            purpose = m.group("purpose")
            master = _LAYOUT_KEY_RE.search(m.group("master"))
            twin = _TWIN_RE.search(purpose)
            out[m.group("name")] = Archetype(
                name=m.group("name"),
                grade=grade,
                purpose=purpose,
                master=master.group(1) if master else None,
                aliases=tuple(_ALIAS_RE.findall(purpose)),
                status=m.group("status").strip(),
                twin_of=twin.group(1) if twin else None,
            )
    return out


def load_spec() -> tuple[dict[str, Layout], dict[str, Archetype]]:
    base = spec_dir()
    return (
        parse_layouts((base / "LAYOUTS.md").read_text()),
        parse_archetypes((base / "ARCHETYPES.md").read_text()),
    )


# --------------------------------------------------------------------------
# role binding
# --------------------------------------------------------------------------

# LAYOUTS.md role -> (engine role, content key stem). `Title` and
# `Text/body` are resolved positionally in `_bind_roles`, since the same
# label covers a headline and a caption.
_PICTURE_ROLES = frozenset({"Picture"})
_TITLE_ROLES = frozenset({"Title"})
_BODY_ROLES = frozenset({"Text/body"})
_PANEL_ROLES = frozenset({"Background"})

# Below this height a `Title` box is a line, not a headline -- the
# master uses the same placeholder type for the 36px subtitle on
# slideLayout18 as for the 104px headline above it.
_SUBTITLE_MAX_H = 60

# A modifier means different things either side of the role: on a
# `Background` box it is the FILL (`pink`, `#D2F5FF`); on a text box it
# is the INK (`white` text on the blue panel of Quote A). Both readings
# are unambiguous given the role, so they share the field.
_PANEL_FILL = {
    "white": "FFFFFF",
    "sand": "F3EEEA",
    "pink": "FFCDD7",
    "yellow": "FFE141",
    "mint": "AAE1C8",
    "blue": "1450F5",
}
# Quote A's panel carries no modifier at all -- an unqualified panel is
# KONE Blue, which is why its body text is separately marked `white`.
_DEFAULT_PANEL_FILL = "1450F5"


def _key(stem: str, n: int) -> str:
    return stem if n == 0 else f"{stem}{n + 1}"


def _columns(boxes: list[Box]) -> tuple[list[list[Box]], list[Box]]:
    """Split boxes into repeating column groups and singletons.

    Three 374x403 boxes at y=227 sitting at x=45/453/861 are not three
    regions, they are one group rendered three times -- and only the
    group form lets the engine expand it over a content list of
    whatever length the caller supplies. Boxes are a column set when
    they share role, y, w and h and there are at least two of them;
    sets that share the same x origins (the photo row and the caption
    row of `Three pictures and text`) are then merged into one group so
    each column carries both its picture and its text.
    """
    rows: dict[tuple, list[Box]] = {}
    for b in boxes:
        rows.setdefault((b.role, b.y, b.w, b.h), []).append(b)

    column_rows = [sorted(v, key=lambda b: b.x) for v in rows.values() if len(v) > 1]
    singles = [v[0] for v in rows.values() if len(v) == 1]

    merged: dict[tuple, list[list[Box]]] = {}
    for row in column_rows:
        merged.setdefault(tuple(b.x for b in row), []).append(row)

    groups: list[list[Box]] = []
    for row_set in merged.values():
        row_set.sort(key=lambda r: r[0].y)
        groups.append([b for row in row_set for b in row])
    return groups, singles


def _engine_role(box: Box, kind: str, ordinal: int) -> str:
    if kind == "title":
        if box.h <= _SUBTITLE_MAX_H or ordinal > 0:
            return "heading"
        return "title_light" if box.on_dark else "title"
    if kind == "body":
        return "on_panel_body" if box.on_dark else "body"
    return kind


def _bind_roles(layout: Layout) -> dict:
    """A `Layout` -> an engine archetype dict.

    Chrome is dropped (the layout owns it). Remaining boxes are read in
    reading order, so the first `Title` is the headline and a later,
    shorter one is the line beneath it. Repeating columns become a
    `groups` entry; everything else becomes a flat region.
    """
    content = layout.content_boxes()
    groups_boxes, singles = _columns(content)

    regions: list[dict] = []
    counters: dict[str, int] = {}

    def add(box: Box, kind: str, stem: str) -> None:
        n = counters.get(stem, 0)
        counters[stem] = n + 1
        regions.append({
            "role": _engine_role(box, kind, n),
            "box": [box.x, box.y, box.w, box.h],
            "content": _key(stem, n),
        })

    # Colour panels are carried outside `regions` and painted by
    # `render` before the engine draws anything -- the engine's own
    # `panel` role is hardcoded to KONE Blue, and these come in five
    # colours. Keeping them out of `regions` also keeps the archetype
    # dict something the unmodified engine can still render.
    panels = [
        {
            "box": [b.x, b.y, b.w, b.h],
            "fill": next((_PANEL_FILL.get(m, m.lstrip("#")) for m in b.mods),
                         _DEFAULT_PANEL_FILL),
        }
        for b in sorted(singles, key=lambda b: (b.y, b.x))
        if b.role in _PANEL_ROLES
    ]

    for box in sorted(singles, key=lambda b: (b.y, b.x)):
        if box.role in _TITLE_ROLES:
            add(box, "title", "title")
        elif box.role in _BODY_ROLES:
            add(box, "body", "body")
        elif box.role in _PICTURE_ROLES:
            add(box, "picture", "image")

    groups: list[dict] = []
    for column_boxes in groups_boxes:
        origins = sorted({b.x for b in column_boxes})
        top = min(b.y for b in column_boxes)
        first_column = sorted((b for b in column_boxes if b.x == origins[0]),
                              key=lambda b: b.y)
        sub: list[dict] = []
        seen: dict[str, int] = {}
        for box in first_column:
            if box.role in _PICTURE_ROLES:
                kind, stem = "picture", "image"
            elif box.role in _TITLE_ROLES:
                kind, stem = "title", "heading"
            else:
                kind, stem = "body", "text"
            n = seen.get(stem, 0)
            seen[stem] = n + 1
            sub.append({
                "role": _engine_role(box, kind, n),
                "box": [0, box.y - top, box.w, box.h],
                "content": _key(stem, n),
            })
        groups.append({
            "content": "items",
            "origins": [[x, top] for x in origins],
            "regions": sub,
        })

    spec: dict = {"regions": regions}
    if groups:
        spec["groups"] = groups
    if panels:
        spec["panels"] = panels
    scrims = _implied_scrims(regions)
    if scrims:
        spec["scrims"] = scrims
    return spec


# Roles the engine draws in white. A layout that puts one of these over a
# photograph has reversed type out of a photograph, whatever else it does.
_LIGHT_ROLES = frozenset({"title_light", "on_panel_body", "eyebrow_light"})

# Every role the engine renders as a photograph. `image_band` is not a
# synonym nobody uses -- it is what `image_section_divider` calls its
# full-bleed picture, and matching only on "picture" meant that
# archetype, whose whole design is white type over a photograph, was the
# one full-bleed layout that never got a scrim.
_PICTURE_REGION_ROLES = frozenset({"picture", "image_band"})


def _is_light_role(role) -> bool:
    """Whether a region's ink is white, whichever path produced it.

    Bound layouts name the role (`title_light`); gallery ports carry the
    colour in the name itself (`gal_i64_FFFFFF`). Both reverse out of the
    picture underneath and both need protecting.
    """
    role = str(role or "")
    return role in _LIGHT_ROLES or role.upper().endswith("_FFFFFF")


def _implied_scrims(regions: list[dict]) -> list[dict]:
    """A picture carrying white type needs a scrim, and no layout says so.

    `LAYOUTS.md` marks the type white and leaves the protection to the
    designer, so the derived archetypes shipped without any: the
    full-bleed cover put a white headline straight onto a sunlit
    treeline. Any picture region a light-ink box sits on gets one.
    """
    pictures = [r for r in regions if r.get("role") in _PICTURE_REGION_ROLES]
    if not pictures:
        return []
    light = [r for r in regions if _is_light_role(r.get("role"))]
    scrims = []
    for picture in pictures:
        px, py, pw, ph = picture["box"]
        over = [
            r for r in light
            if r["box"][1] < py + ph and r["box"][1] + r["box"][3] > py
            and r["box"][0] < px + pw and r["box"][0] + r["box"][2] > px
        ]
        if over:
            scrims.append({"box": [px, py, pw, ph], "content": picture.get("content")})
    return scrims


# Sensible generic pictograms for an archetype whose caller named none.
# Checked against the sprite at use time, so a renamed icon degrades to
# the ones that do exist rather than drawing nothing.
_DEFAULT_ICONS = ("cloud", "people", "clock", "wrench", "calendar", "elevator")


def pictograms() -> list[str]:
    """The rasterised marks, kept only as a fallback.

    Superseded by `deckguard.icons`, which draws the real KONE
    pictograms as native editable shapes. This survives for the case
    where the icon sprite is not installed.
    """
    base = spec_dir().parent.parent / "icons"
    return [str(p) for p in sorted(base.glob("*.png"))]


def _icon_names_for(spec: dict, content: dict, slots: int) -> list[str]:
    """Which pictogram goes in each icon slot.

    A caller names one per item (`{"icon": "elevator", ...}`) and gets
    it; anything unnamed falls back to the generic rotation. With 609
    icons in the sprite an author can finally be specific, but nothing
    breaks if they are not.
    """
    from deckguard import icons as icon_mod

    available = icon_mod.load_icons()
    named: list[str] = []
    for group in spec.get("groups", []):
        if not any(r.get("role") == "icon" for r in group["regions"]):
            continue
        for item in (content.get(group["content"]) or [])[:len(group["origins"])]:
            named.append((item or {}).get("icon") if isinstance(item, dict) else None)

    fallback = [n for n in _DEFAULT_ICONS if n in available] or sorted(available)[:1]
    chosen: list[str] = []
    for index in range(slots):
        want = named[index] if index < len(named) else None
        chosen.append(want if want in available else fallback[index % len(fallback)])
    return chosen


def _icon_slots(spec: dict, content: Optional[dict] = None) -> int:
    """How many icons this slide actually needs.

    Counted against the content, not the geometry. An archetype that
    serves both a plain and a grid form declares the grid's origins
    either way, so counting those drew a full set of icons onto the
    plain form -- four of them, two straddling the body paragraph.
    """
    n = sum(1 for r in spec.get("regions", []) if r.get("role") == "icon")
    for group in spec.get("groups", []):
        per = sum(1 for r in group["regions"] if r.get("role") == "icon")
        if not per:
            continue
        supplied = len(group["origins"]) if content is None else len(
            (content.get(group["content"]) or [])[:len(group["origins"])]
        )
        n += per * supplied
    return n


def _table_from_text(text: str) -> dict:
    """`a|b|c;d|e|f` as {headers, rows}.

    The shape a caller reaches for when told a slot holds "a table" and
    nothing more. Rows split on `;`, columns on `|`, first row is the
    header. Newlines work as row separators too, because that is the
    other obvious guess.
    """
    rows = [r.strip() for r in re.split(r"[;\n]", text) if r.strip()]
    cells = [[c.strip() for c in row.split("|")] for row in rows]
    cells = [row[1:] if row and row[0] == "" else row for row in cells]
    if not cells:
        return {"headers": [], "rows": []}
    width = max(len(row) for row in cells)
    cells = [row + [""] * (width - len(row)) for row in cells]
    # the header row of a comparison table has an empty first cell -- it
    # sits over the row labels -- so it comes back one short
    return {"headers": [""] * (width - len(cells[0])) + cells[0], "rows": cells[1:]}


def _coerce_content(spec: dict, content: dict) -> dict:
    """Put the caller's content into the shape the engine reads.

    Two failures this exists for, both found by handing the renderer
    copy written against the contract rather than against the code:

    A list in a text region came out as `repr` -- an external cover
    shipped with `['Maintenance and modernisation', 'Prepared for the
    property team', ...]` printed under its headline, brackets, quotes
    and all. The contract types that slot `bullets`; the region is a
    plain text box. Joining is right and failing the build is not: the
    words were correct and only the shape was wrong.

    A table written as text crashed the build outright, because the
    engine calls `.get` on it. `comparison_table` advertises `table` and
    says nothing about {headers, rows}, so a string is the reasonable
    guess and it took the whole deck down with an AttributeError.
    """
    if not content:
        return content
    kinds: dict = {}
    for region in spec.get("regions") or []:
        if region.get("content"):
            kinds[region["content"]] = (
                (region.get("dg") or {}).get("kind") or region.get("role") or "text")
    grouped = {g.get("content") for g in spec.get("groups") or []}

    out = dict(content)
    for key, value in content.items():
        kind = kinds.get(key)
        if kind == "table" and isinstance(value, str) and value.strip():
            out[key] = _table_from_text(value)
        elif (kind not in (None, "bullets", "table") and key not in grouped
              and isinstance(value, list)):
            out[key] = "\n".join(str(v) for v in value if str(v).strip())
    return out


def render(slide, name: str, content: dict, archetypes_module,
           audience: str = "") -> None:
    """Draw an archetype onto `slide`.

    Colour panels are painted first, then the engine draws the text and
    pictures on top. The order is the whole point: an earlier version
    of the gallery port delegated first and painted its background
    afterwards, which turned a white content field into a solid blue
    slide.

    Archetypes carrying gallery `chrome` (the cut covers, the dividers,
    the outro) go through `archetypes.render`, because `gallery.install`
    wraps that function to draw the chrome and calling the engine
    directly would silently drop the cut effect and the divider
    artwork.

    Everything else goes straight to `render_archetype`, for two
    reasons that only show up on real decks. `archetypes.render`
    force-feeds its own sample chart art into `chart`/`diagram` slots
    even when the caller supplied content -- which is how a deck about
    elevators acquired a pie chart nobody asked for -- so the figures
    are applied with `setdefault` here instead. And it passes the
    engine's placeholder icon chips; real pictograms are passed
    instead, cycled to cover archetypes that want more than the three
    that exist.
    """
    from pptx.enum.shapes import MSO_SHAPE

    engine = archetypes_module.E
    spec = archetypes_module.ARCHETYPES.get(name, {})
    content = _coerce_content(spec, content)

    if spec.get("chrome"):
        archetypes_module.render(slide, name, content)
        # The gallery's ports were extracted from a reference deck that
        # carried no scrims, and the full-bleed cover is the one layout
        # the brand explicitly says needs one -- its headline landed
        # white on a sunlit treeline. `_seat_above_picture` puts the
        # gradient under the type the engine has already drawn.
        for scrim in spec.get("scrims") or _implied_scrims(spec.get("regions", [])):
            if scrim.get("content") in (None, "") or content.get(scrim.get("content")):
                _draw_scrim(slide, engine, scrim["box"], scrim.get("opacity", 78),
                            _reversed_bands(spec, content, scrim["box"]))
        # A cut cover reached here draws its own masks through the
        # gallery's chrome renderer, so it never passed `draw_cut_cover`
        # and never got the KONE marks. Which branch a cover takes
        # depends on whether the gallery is installed, so the cover was
        # shipping with a logo or without one depending on load order.
        cut = _CUT_COVERS.get((name or "").lower())
        if cut:
            _cover_marks(slide, engine, cut)
        return

    for panel in spec.get("panels", []):
        x, y, w, h = panel["box"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, engine.X(x), engine.X(y), engine.X(w), engine.X(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = engine._hex(panel["fill"])
        shape.line.fill.background()
        shape.shadow.inherit = False

    # Hairlines. Thin enough that a stroked line renders differently
    # across viewers, so they are 1px filled rectangles like every other
    # rule the master draws.
    for rule in spec.get("rules", []):
        x, y, w, h = rule["box"]
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, engine.X(x), engine.X(y), engine.X(w), engine.X(max(h, 1))
        )
        line.name = "Hairline"
        line.fill.solid()
        line.fill.fore_color.rgb = engine._hex(rule.get("fill", "D0D0D0"))
        line.line.fill.background()
        line.shadow.inherit = False

    filled = dict(content)

    # Regions the reference describes more precisely than ROLE_STYLE can
    # express are drawn here and withheld from the engine.
    engine_spec = {k: v for k, v in spec.items() if k not in ("regions", "groups")}
    engine_spec["regions"] = [
        r for r in spec.get("regions", [])
        if "dg" not in r and not _is_unsupplied_figure(r, filled, name, archetypes_module)
    ]
    engine_spec["groups"] = [
        {**g, "regions": [r for r in g["regions"] if "dg" not in r]}
        for g in spec.get("groups", [])
    ]

    # Icons are drawn here as native shapes, so they are withheld from
    # the engine along with the `dg` regions -- otherwise it draws its
    # blue placeholder chip in the same box.
    from deckguard import icons as icon_mod

    # The engine renders a plain `bullets` region with `_dash_bullets`,
    # which literally emits "—  " as the marker. BRAND_MODE section 6
    # calls a dash standing in for a bullet a brand violation and names
    # the archetypes doing it. Withhold those regions and draw real disc
    # markers here, the same way the `dg` path already does.
    plain_bullets = [(r["box"], filled[r["content"]])
                     for r in spec.get("regions", [])
                     if r.get("role") == "bullets" and filled.get(r.get("content"))]
    for group in spec.get("groups", []):
        items = filled.get(group.get("content")) or []
        for (ox, oy), item in zip(group.get("origins", []), items):
            if not isinstance(item, dict):
                continue
            for region in group.get("regions", []):
                value = item.get(region.get("content"))
                if region.get("role") == "bullets" and value:
                    rx, ry, rw, rh = region["box"]
                    plain_bullets.append(([ox + rx, oy + ry, rw, rh], value))
    if plain_bullets:
        engine_spec["regions"] = [r for r in engine_spec["regions"]
                                  if r.get("role") != "bullets"]
        engine_spec["groups"] = [
            {**g, "regions": [r for r in g["regions"] if r.get("role") != "bullets"]}
            for g in engine_spec["groups"]
        ]

    slots = _icon_slots(spec, filled)
    native = bool(icon_mod.load_icons()) and slots
    if native:
        engine_spec["regions"] = [r for r in engine_spec["regions"] if r.get("role") != "icon"]
        engine_spec["groups"] = [
            {**g, "regions": [r for r in g["regions"] if r.get("role") != "icon"]}
            for g in engine_spec["groups"]
        ]
        icons = None
    else:
        marks = pictograms()
        icons = [marks[i % len(marks)] for i in range(slots)] if marks and slots else None

    engine.render_archetype(
        slide, engine_spec, filled, icons=icons,
        bg=getattr(archetypes_module, "BG", {}).get(name),
    )

    for box, value in plain_bullets:
        _draw_bullets(slide, engine, box, value,
                      {"kind": "bullets", "px": 19, "nested_px": 17, "lead": 0.45}, None)

    # The banner is one photograph; the cut comes from masking it. Drawn
    # after the engine so the masks sit ON the picture rather than under
    # it, which is the whole mechanism.
    bg_name = getattr(archetypes_module, "BG", {}).get(name) or spec.get("background")
    draw_cut_cover(slide, engine, name, _bg_hex(bg_name))

    if native:
        chosen = _icon_names_for(spec, filled, slots)
        index = 0
        for region in spec.get("regions", []):
            if region.get("role") == "icon":
                icon_mod.add_icon(slide, chosen[index], region["box"])
                index += 1
        for group in spec.get("groups", []):
            items = filled.get(group["content"]) or []
            for (ox, oy), _item in zip(group["origins"], items):
                for region in group["regions"]:
                    if region.get("role") != "icon":
                        continue
                    rx, ry, rw, rh = region["box"]
                    icon_mod.add_icon(slide, chosen[index], [ox + rx, oy + ry, rw, rh])
                    index += 1

    # AFTER the engine, never before. `render_archetype` starts by
    # painting a full-slide background rectangle, so anything drawn
    # first is buried by it -- both numbered dividers came out as an
    # empty sand rectangle with the number, title and label underneath.
    field = _field_for(spec, filled, name, audience)
    if field:
        _paint_field(slide, engine, field[0])
    else:
        _paint_layout_background(slide, engine, spec)
    # Derived here as well as at build time: the gallery's own ports win
    # over the bound layouts for several archetypes, and they were
    # extracted from a reference deck that carried no scrims at all.
    for scrim in spec.get("scrims") or _implied_scrims(spec.get("regions", [])):
        if scrim.get("content") in (None, "") or filled.get(scrim.get("content")):
            _draw_scrim(slide, engine, scrim["box"], scrim.get("opacity", 78),
                        _reversed_bands(spec, filled, scrim["box"]))
    ink = field[1] if field else None
    for region in spec.get("regions", []):
        if "dg" in region:
            _draw(slide, engine, region, filled.get(region.get("content")), ink)
    for group in spec.get("groups", []):
        items = filled.get(group["content"]) or []
        for (ox, oy), item in zip(group["origins"], items):
            for region in group["regions"]:
                if "dg" not in region:
                    continue
                rx, ry, rw, rh = region["box"]
                shifted = {**region, "box": [ox + rx, oy + ry, rw, rh]}
                _draw(slide, engine, shifted, (item or {}).get(region.get("content")), ink)
    # Cards last: they carry a shadow, so anything drawn after them
    # would sit on top of it.
    if spec.get("cards"):
        _draw_cards(slide, engine, spec["cards"], filled)


# The brand's secondary palette, with the ink each field takes. The
# rule is the brand's own: a blue field takes white type, a secondary
# field takes black.
BRAND_FIELDS = {
    "blue": ("1450F5", "FFFFFF"),
    "light-blue": ("D2F5FF", "141414"),
    "pink": ("FFCDD7", "141414"),
    "yellow": ("FFE141", "141414"),
    "mint": ("AAE1C8", "141414"),
    "sand": ("F3EEEA", "141414"),
}
DEFAULT_FIELD = "blue"


def _field_for(spec: dict, content: dict, name: str = "", audience: str = ""):
    """(fill, ink) when this archetype paints a whole-slide colour field.

    Two sources, in this order.

    An archetype that declares `field` -- the dividers -- takes the
    colour the caller names per slide. A divider is the one place the
    secondary palette is meant to carry a whole slide, and leaving them
    all the same makes a deck monotonous.

    Otherwise the SET decides. `slide-sets.json` gives every slide a
    field per audience and nothing was reading it, so `hero_stat` came
    out white where the internal set says light-blue, `divider_title_only`
    sand where it says light-blue, and -- the one that matters -- the
    external `kone_numbers` carried a sand band on a customer deck whose
    policy is white and blue only.
    """
    if spec.get("field"):
        colour = str(content.get("colour") or content.get("color")
                     or DEFAULT_FIELD).lower()
        return BRAND_FIELDS.get(colour, BRAND_FIELDS[DEFAULT_FIELD])
    # A field the SET declares carries no ink override. Only a divider,
    # which declares its own field, reverses the type on it -- there the
    # caller picked a colour and every role has to follow. Here the
    # archetype was drawn for this field already, and forcing its ink
    # black turned `kone_numbers`' blue figures and blue eyebrow black
    # the moment the white field started being painted.
    declared = _set_field(name, audience)
    if declared == "white":
        return ("FFFFFF", None)
    return (BRAND_FIELDS[declared][0], None) if declared in BRAND_FIELDS else None


def _set_field(name: str, audience: str) -> str:
    """What `slide-sets.json` says this slide's field is, in this set.

    `photo` and `secondary` return empty: one is a photograph the
    archetype places itself, the other is the archetype's own panel, and
    neither is a flood the renderer should paint.
    """
    if not name or not audience:
        return ""
    # Locally, because `brandmode` is not imported at module level here.
    # Written as a bare `bm` first, which raised NameError straight into
    # the `except` below and returned "" for every slide in both sets --
    # the fix looked applied, the renders did not change, and nothing
    # said why. Hence KeyError only: an unknown set name is the one
    # thing this is allowed to shrug off.
    from deckguard import brandmode as bm

    try:
        slides = bm.slides_in(audience)
    except KeyError:
        return ""
    for slide in slides:
        if slide["archetype"] == name:
            field = str(slide.get("field") or "").lower()
            return field if field in {*BRAND_FIELDS, "white"} else ""
    return ""


def _is_unsupplied_figure(region: dict, content: dict, name: str, archetypes_module) -> bool:
    """A chart or diagram slot the author gave nothing for.

    The engine keeps a `FIGURES` map of sample artwork and stamps it onto
    every slide of certain archetypes whether or not anyone asked --
    `segment_breakdown` gets a donut reading 53% against satisfaction
    bands of 9-10, 7-8, 5-6 and "Don't know". On a deck built from a
    brief that is not decoration, it is INVENTED DATA wearing the
    company's own chart styling, and it went out in a half-year business
    review.

    A slot with nothing in it is dropped instead. An empty column is
    obvious and harmless; a fabricated statistic is neither.
    """
    figures = getattr(archetypes_module, "FIGURES", {}).get(name) or {}
    key = region.get("content")
    if key not in figures:
        return False
    return not content.get(key)


def _paint_layout_background(slide, engine, spec: dict) -> None:
    """Make the layout's background explicit on the slide.

    Four KONE layouts set their background by THEME REFERENCE rather
    than by colour -- `<p:bgRef><a:schemeClr val="bg2"/>` -- and
    renderers disagree about what that resolves to. PowerPoint follows
    the master's clrMap to sand; LibreOffice produces KONE Blue. Same
    file, two different decks.

    Writing the resolved colour onto the slide removes the question. It
    is a no-op where the archetype paints its own background, and where
    the layout states a literal colour it simply restates it.
    """
    if spec.get("background") or spec.get("panels"):
        return
    from pptx.enum.shapes import MSO_SHAPE

    from deckguard.logo import _page_background_hex

    layout = slide.slide_layout
    if _page_background_hex(layout.element) is not None:
        return                      # already a literal colour, unambiguous
    resolved = _page_background_hex(layout.element, layout)
    if not resolved:
        return
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, engine.X(0), engine.X(0),
                                   engine.X(1280), engine.X(720))
    shape.name = "Layout background"
    shape.fill.solid()
    shape.fill.fore_color.rgb = engine._hex(resolved)
    shape.line.fill.background()
    shape.shadow.inherit = False
    style = shape._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        shape._element.remove(style)
    # behind everything the engine just drew
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def _reversed_bands(spec: dict, content: dict, scrim_box) -> list:
    """The (top, bottom) of every white text block sitting on the scrim.

    Only blocks the slide actually fills count -- an archetype that
    declares an eyebrow the author left empty must not have the picture
    darkened where nothing is written.
    """
    _, sy, _, sh = scrim_box
    bands = []
    for region in spec.get("regions", []):
        style = region.get("dg")
        if style is not None:
            if str(style.get("color", "")).upper() not in ("FFFFFF", "FFF"):
                continue
        elif not _is_light_role(region.get("role")):
            # drawn by the engine; only its light roles are reversed out
            continue
        value = content.get(region.get("content"))
        if not value:
            continue
        x, y, w, h = region["box"]
        if y + h <= sy or y >= sy + sh:
            continue
        bands.append((y, y + _typeset_height(region, value)))
    return bands


def _typeset_height(region: dict, value) -> float:
    """Roughly how far down its box a block of text actually reaches.

    Protecting the whole BOX is what turned the full-bleed cover into a
    uniformly dark photograph: the cover's title frame is 448px tall for
    three lines of type, so the scrim covered the middle of the picture
    end to end. Three lines is what needs protecting, not the frame they
    were given.
    """
    _, _, width, height = region["box"]
    style = region.get("dg") or {}
    px = style.get("px")
    if not px:
        match = re.search(r"_i(\d+)_", str(region.get("role") or ""))
        px = int(match.group(1)) if match else 0
    if not px:
        return height

    text = " ".join(str(v) for v in value) if isinstance(value, (list, tuple)) else str(value)
    per_line = max(1, int(width / (px * 0.52)))
    lines = max(1, -(-len(text) // per_line))
    return min(height, lines * px * 1.3)


def _scrim_ramp(box, protect, opacity: int, feather: float = 0.10):
    """Gradient stops that are dark exactly where white type sits.

    A fixed dark-at-both-edges ramp is a guess, and on TEXT_PICTURE_G it
    guessed wrong: the headline starts 60% of the way down its banner,
    where a ramp that clears at the midpoint has recovered barely a
    sixth of its opacity, and the first line read white-on-sunlit-
    pavement. The bands that need protecting are not a mystery -- they
    are the archetype's own reversed-out text boxes -- so the ramp is
    built from them instead of guessed at.

    Returns `(position, alpha)` pairs sampled across the box, each band
    at full opacity with a soft edge so the picture is never cut by a
    visible line.
    """
    _, y, _, h = box
    if not h:
        return [(0, opacity), (100000, opacity)]

    bands = []
    for bx in protect:
        top = max(0.0, min(1.0, (bx[0] - y) / h))
        bottom = max(0.0, min(1.0, (bx[1] - y) / h))
        if bottom > top:
            bands.append((top, bottom))
    if not bands:
        # nothing reversed out over this picture: the brand's plain
        # edge-darkening, which still seats a photo on a white slide
        bands = [(0.0, 0.06), (0.94, 1.0)]

    def cover(position: float) -> float:
        best = 0.0
        for top, bottom in bands:
            if top <= position <= bottom:
                best = 1.0
            elif position < top:
                best = max(best, 1.0 - min(1.0, (top - position) / feather))
            else:
                best = max(best, 1.0 - min(1.0, (position - bottom) / feather))
        return best

    steps = 20
    return [
        (int(round(i / steps * 100000)), int(round(opacity * cover(i / steps))))
        for i in range(steps + 1)
    ]


def _draw_scrim(slide, engine, box, opacity: int = 78, protect=()) -> None:
    """A gradient over a photo so white type stays readable on it.

    The brand specifies this for `COVER_F_FULLBLEED` and never says how,
    and every archetype that reverses type out of a photograph needs the
    same thing: a banner title came out white-on-pale-escalator and was
    close to invisible.

    Dark behind the reversed-out type and clear elsewhere, so it protects
    the eyebrow and the headline without greying out the subject of the
    picture -- which a flat overlay does, and which is why the spec calls
    for a gradient rather than a tint.
    """
    from lxml import etree
    from pptx.enum.shapes import MSO_SHAPE

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    x, y, w, h = box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, engine.X(x), engine.X(y), engine.X(w), engine.X(h)
    )
    shape.name = "Photo protection"
    shape.line.fill.background()
    shape.shadow.inherit = False
    style = shape._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        shape._element.remove(style)

    spPr = shape._element.spPr
    for tag in ("solidFill", "noFill", "gradFill"):
        existing = spPr.find(f"{{{A}}}{tag}")
        if existing is not None:
            spPr.remove(existing)

    grad = etree.SubElement(spPr, f"{{{A}}}gradFill")
    grad.set("rotWithShape", "1")
    stops = etree.SubElement(grad, f"{{{A}}}gsLst")
    for position, alpha in _scrim_ramp(box, protect, opacity):
        stop = etree.SubElement(stops, f"{{{A}}}gs")
        stop.set("pos", str(position))
        colour = etree.SubElement(stop, f"{{{A}}}srgbClr")
        colour.set("val", "141414")
        etree.SubElement(colour, f"{{{A}}}alpha").set("val", str(alpha * 1000))
    lin = etree.SubElement(grad, f"{{{A}}}lin")
    lin.set("ang", "5400000")   # top to bottom
    lin.set("scaled", "0")

    # The gradient must sit directly on the photo and UNDER the type.
    # Appending is only correct when the caller draws the text
    # afterwards; the gallery's ports have the engine draw theirs inside
    # `render_archetype`, so an appended scrim greyed out the headline it
    # was meant to protect. Seating it immediately above the picture is
    # right for both paths.
    _seat_above_picture(shape, box, slide)
    return shape


def _seat_above_picture(shape, box, slide=None) -> None:
    """Move `shape` to just after the picture it is protecting.

    The FIRST substantial one, not the last: a cover carries the logo
    and the tagline as pictures too, and they are added after the type.
    Seating the scrim above those put it over the headline and dimmed
    the very words it exists to make readable.
    """
    from lxml import etree

    tree = shape._element.getparent()
    if tree is None or slide is None:
        return
    left, top, width, height = box
    area = max(1.0, width * height)

    for candidate in slide.shapes:
        if etree.QName(candidate._element).localname != "pic":
            continue
        # Resolved geometry, not the raw `xfrm`: a picture PLACEHOLDER
        # inherits its position from the layout and carries no `xfrm` of
        # its own, so reading the element directly found nothing and the
        # scrim stayed appended -- on top of the headline.
        try:
            px, py = candidate.left / 9525.0, candidate.top / 9525.0
            pw, ph = candidate.width / 9525.0, candidate.height / 9525.0
        except Exception:  # noqa: BLE001 -- a picture without geometry
            continue
        if pw * ph < area * 0.5:      # chrome, not the photograph
            continue
        if px < left + width and px + pw > left and py < top + height and py + ph > top:
            # lxml's own index, never a dict keyed on `id()`: lxml builds
            # a fresh Python proxy on each element access and lets the
            # old one be collected, so those ids are neither stable nor
            # unique. A recycled id matched the wrong element and seated
            # the scrim one shape too late -- over the type.
            try:
                index = tree.index(candidate._element)
            except ValueError:  # pragma: no cover -- not in this tree
                continue
            tree.remove(shape._element)
            tree.insert(index + 1, shape._element)
            return


def _paint_field(slide, engine, fill: str) -> None:
    """Flood the slide with one of the brand's secondary colours.

    The field has to REPLACE the engine's own background rectangle, not
    sit under it. `render_archetype` opens by painting the archetype's
    default fill across the slide, so a field inserted at the bottom of
    the shape tree is covered by sand and only the ink colour survives --
    which is how a blue divider came out as white type on sand, all but
    invisible.
    """
    from pptx.enum.shapes import MSO_SHAPE

    for existing in _full_slide_fills(slide, engine):
        existing.getparent().remove(existing)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, engine.X(0), engine.X(0),
                                   engine.X(1280), engine.X(720))
    shape.name = "Colour field"
    shape.fill.solid()
    shape.fill.fore_color.rgb = engine._hex(fill)
    shape.line.fill.background()
    shape.shadow.inherit = False
    style = shape._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        shape._element.remove(style)
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)
    _relay_layout_logo(slide)


def _relay_layout_logo(slide) -> bool:
    """Put the layout's logo back on top of a painted field.

    The logo lives on the LAYOUT, and every shape on a slide draws above
    every shape on its layout -- so flooding the slide with a colour
    buries it. The dividers had been shipping without a logo for exactly
    this reason and it read as the layout's own design; making the set's
    declared field authoritative would have spread it to the thirty-nine
    slides that declare white.

    Placed from the vendored asset at the layout logo's own geometry,
    NOT by copying the layout's element: a picture's `r:embed` points at
    a relationship owned by the layout part, so a deep copy onto the
    slide lands as an empty frame -- the same failure as the 54 blank
    logo frames in `All_Slides.pptx`. White mark on a dark field.
    """
    from deckguard.logo import _brand_asset, _is_dark

    if any((sh.name or "") == "Logo" for sh in slide.shapes):
        return False
    source = next((sh for sh in slide.slide_layout.shapes
                   if (sh.name or "") == "Logo"), None)
    if source is None:
        return False
    field = next((sh for sh in slide.shapes if (sh.name or "") == "Colour field"), None)
    light = False
    try:
        if field is not None and field.fill.type is not None:
            light = _is_dark(str(field.fill.fore_color.rgb))
    except Exception:  # noqa: BLE001 -- an unreadable fill just means dark type
        light = False
    path = _brand_asset("logo", light)
    if not path:
        return False
    picture = slide.shapes.add_picture(
        path, source.left, source.top, source.width, source.height)
    picture.name = "Logo"
    return True


def _full_slide_fills(slide, engine):
    """The plain filled rectangles that cover the whole slide.

    Matched on geometry and fill rather than on name: the engine names
    its background `Rectangle 1`, but that name is not reserved and a
    port could reasonably produce another.
    """
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    width, height = engine.X(1280), engine.X(720)
    found = []
    for shape in slide.shapes:
        if not shape.has_text_frame or shape.text_frame.text.strip():
            continue
        geom = shape._element.find(f".//{{{A}}}prstGeom")
        if geom is None or geom.get("prst") != "rect":
            continue
        if (shape.left, shape.top) != (0, 0):
            continue
        if abs(shape.width - width) > 1000 or abs(shape.height - height) > 1000:
            continue
        if shape._element.find(f".//{{{A}}}solidFill") is None:
            continue
        found.append(shape._element)
    return found


def _draw(slide, engine, region: dict, value, ink: Optional[str] = None) -> None:
    """Draw one reference-specified block."""
    if value in (None, "", []):
        return
    style = region["dg"]
    if ink:
        style = {**style, "color": ink}
    if style["kind"] == "bullets":
        _draw_bullets(slide, engine, region["box"], value, style, ink)
    elif style["kind"] == "ruled":
        _draw_ruled(slide, engine, region["box"], value, style)
    elif style["kind"] == "tick":
        _draw_tick(slide, engine, region["box"], value, style)
    else:
        if style.get("zero_is_black") and _reads_as_zero(value):
            # The brand's own instruction: the number that is deliberately
            # zero -- no disruption, no downtime, no escalations -- is the
            # strongest claim on the slide, and reads as a different KIND
            # of claim when it is not in the same blue as the rest.
            style = {**style, "color": "141414"}
        _draw_text(slide, engine, region["box"], value, style)


def _reads_as_zero(value) -> bool:
    return str(value).strip().strip("+%") in ("0", "0.0", "zero", "Zero", "none", "None")


def _draw_ruled(slide, engine, box, value, style) -> None:
    """Text under its own hairline.

    The rule belongs to the cell rather than to the slide because the
    grid is content-length: eight names fill two rows of four, and a
    third row of rules drawn statically would hang under them as a line
    ruling nothing.
    """
    from pptx.enum.shapes import MSO_SHAPE

    x, y, w, h = box
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, engine.X(x), engine.X(y), engine.X(w), engine.X(1))
    rule.name = "Hairline"
    rule.fill.solid()
    rule.fill.fore_color.rgb = engine._hex(style.get("rule", "D0D0D0"))
    rule.line.fill.background()
    rule.shadow.inherit = False

    pad = style.get("pad", 26)
    _draw_text(slide, engine, [x, y + pad, w, h - pad], value, {**style, "kind": "text"})


def _draw_tick(slide, engine, box, value, style) -> None:
    """A blue disc with a white check, and its label beside it.

    The only symbol the milestone slide allows, and a glyph rather than
    an emoji -- the brand forbids emoji outright and a coloured emoji
    tick would also ignore the palette.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    x, y, w, h = box
    size = style.get("badge", 20)
    disc = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, engine.X(x), engine.X(y), engine.X(size), engine.X(size)
    )
    disc.name = "Tick"
    disc.fill.solid()
    disc.fill.fore_color.rgb = engine._hex(style.get("badge_fill", "1450F5"))
    disc.line.fill.background()
    disc.shadow.inherit = False

    frame = disc.text_frame
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(frame, margin, 0)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    engine._run(paragraph, "\u2713", "Inter", size * 0.6, engine._hex("FFFFFF"))

    gap = style.get("gap", 12)
    _draw_text(slide, engine, [x + size + gap, y, w - size - gap, h], value,
               {**style, "kind": "text"})


def _draw_text(slide, engine, box, value, style) -> None:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    frame = engine._tf(slide, box, MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    engine._run(
        paragraph, value, style.get("font", "Inter"), style["px"],
        engine._hex(style.get("color", "141414")), caps=style.get("caps", False),
    )
    if style.get("align") == "r":
        paragraph.alignment = PP_ALIGN.RIGHT


def _draw_bullets(slide, engine, box, items, style, ink: Optional[str] = None) -> None:
    """A real list, not a typed dash.

    The brand rule is explicit and the engine breaks it: `_dash_bullets`
    writes an em dash followed by two spaces, which is a character in a
    paragraph, not a list marker. It does not indent, does not hang, does
    not survive editing in PowerPoint's outline view, and cannot nest.

    This writes `buChar` paragraph formatting instead -- marker in KONE
    Blue, text in black, one nested level as the spec allows. A nested
    item is any list entry given as `{"text": ..., "sub": [...]}`.
    """
    from lxml import etree
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Pt

    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def mark(paragraph, level: int, px: float) -> None:
        pPr = paragraph._pPr if paragraph._pPr is not None else paragraph._p.get_or_add_pPr()
        pPr.set("lvl", str(level))
        pPr.set("indent", str(-int(px * 0.75 * 12700)))
        pPr.set("marL", str(int((level + 1) * px * 0.9 * 0.75 * 12700)))
        colour = etree.SubElement(pPr, f"{{{A_NS}}}buClr")
        srgb = etree.SubElement(colour, f"{{{A_NS}}}srgbClr")
        # a blue marker vanishes on a blue field
        srgb.set("val", ink or "1450F5")
        etree.SubElement(pPr, f"{{{A_NS}}}buSzPct").set("val", "100000")
        etree.SubElement(pPr, f"{{{A_NS}}}buFont").set("typeface", "Arial")
        etree.SubElement(pPr, f"{{{A_NS}}}buChar").set("char", "•" if level == 0 else "◦")

    frame = engine._tf(slide, box, MSO_ANCHOR.TOP)
    first = True
    for item in items if isinstance(items, (list, tuple)) else [items]:
        text = item.get("text") if isinstance(item, dict) else item
        children = item.get("sub", []) if isinstance(item, dict) else []
        for level, line, px in [(0, text, style["px"])] + [
            (1, child, style["nested_px"]) for child in children
        ]:
            paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            paragraph.space_after = Pt(px * style.get("lead", 0.35))
            engine._run(paragraph, line, "Inter", px, engine._hex(ink or "141414"))
            mark(paragraph, level, px)


# --------------------------------------------------------------------------
# building and installing
# --------------------------------------------------------------------------


# --------------------------------------------------------------------------
# refinements from the rendered reference
# --------------------------------------------------------------------------
#
# `LAYOUTS.md` is a placeholder map, not a design. It carries four
# content roles, and 109 of its 321 boxes are the single label
# `Text/body` -- so binding it alone produces a headline over an
# undifferentiated paragraph where the design has an eyebrow, a real
# bulleted list with KONE Blue markers, and a nested level. Measured:
# the reference slides average 18.6 elements, generated archetypes 3.6.
#
# The design lives in the `.dc.html` files. These entries transcribe it
# for the archetypes that carry the most decks, and overlay the
# generated geometry. Each `dg` block is drawn by `render` rather than
# by the engine, whose ROLE_STYLE has no 34px title and no real bullet
# support. Everything here is read off the reference markup -- box,
# size and colour -- not invented.
#
# This is deliberately partial. Extending it is how the rest of the
# vocabulary gets its fidelity back; a fidelity harness should drive
# which archetype is worth doing next.

def _text(x, y, w, h, key, px, *, font="Inter", color="141414", caps=False, align="l"):
    return {"role": "dg_text", "box": [x, y, w, h], "content": key,
            "dg": {"kind": "text", "px": px, "font": font, "color": color,
                   "caps": caps, "align": align}}


def _bullets(x, y, w, h, key, px, nested_px=None, lead=0.35):
    return {"role": "dg_bullets", "box": [x, y, w, h], "content": key,
            "dg": {"kind": "bullets", "px": px, "nested_px": nested_px or px - 2,
                   "lead": lead}}


_REFINEMENTS: dict[str, dict] = {
    # KoneDeck slide 4 -- eyebrow, 34px title, disc bullets with a
    # nested circle level. The workhorse of every deck.
    "title_content": {"regions": [
        _text(45, 47, 917, 22, "eyebrow", 12, font="KONE Information", color="1450F5", caps=True),
        _text(45, 91, 917, 104, "title", 34),
        _bullets(45, 227, 917, 402, "bullets", 19, nested_px=17),
    ]},
    # ArchetypeLibraryB slide 07
    "title_subtitle_content_a": {"regions": [
        _text(45, 91, 917, 68, "title", 30),
        _text(45, 159, 917, 36, "subtitle", 17),
        _bullets(45, 227, 917, 402, "bullets", 18),
    ]},
    # KoneDeck slide 6 -- two bulleted columns, each with its own label
    "two_content": {
        "regions": [_text(45, 91, 1189, 104, "title", 34)],
        "groups": [{"content": "items", "origins": [[45, 227], [657, 227]], "regions": [
            _text(0, 0, 578, 22, "label", 12, font="KONE Information", color="1450F5", caps=True),
            _bullets(0, 30, 578, 372, "bullets", 17),
        ]}],
    },
    # ArchetypeLibraryB slide 09 -- the title is in the DESIGN but not
    # in the master layout, so binding LAYOUTS.md alone loses it and the
    # slide comes out headless.
    "three_content": {
        "regions": [_text(45, 91, 1189, 45, "title", 32)],
        "groups": [{"content": "items", "origins": [[45, 227], [453, 227], [861, 227]],
                    "regions": [
                        _text(0, 0, 374, 30, "heading", 20),
                        _text(0, 40, 374, 363, "text", 16),
                    ]}],
    },
    # ArchetypeLibraryB slide 08 -- narrow title column, wide bullets
    "two_content_narrow_title": {"regions": [
        _text(45, 91, 374, 104, "title", 28),
        _text(45, 227, 374, 402, "body", 16),
        _bullets(453, 227, 781, 402, "bullets", 18),
    ]},
    # KoneDeck slide 3. The gallery port of this bound the title into
    # the RIGHT-hand box and emitted no number at all -- which is the
    # one thing a numbered divider is for. The reference puts the
    # section label and title left, and a 420px numeral right.
    # Agenda A. The gallery port emitted four separate one-line text
    # boxes, so an agenda could only ever have exactly four items and
    # they were not a list -- no markers, no hanging indent, nothing
    # PowerPoint's outline view could edit. One bulleted box instead,
    # on the master's own 578x448 body area, with the full-height photo
    # the layout already declares.
    "agenda_a_table": {
        "regions": [
            {"role": "picture", "box": [759, 0, 521, 720], "content": "photo"},
            _text(45, 39, 578, 36, "eyebrow", 14,
                  font="KONE Information", color="1450F5", caps=True),
            _text(45, 91, 578, 60, "title", 40),
            # an agenda is four or five short lines in a 448px column;
            # at body size they huddle in the top corner and leave the
            # rest of the slide empty. Bigger type, and air between the
            # items, so the list occupies the space it is given.
            _bullets(45, 195, 578, 434, "items", 24, lead=1.15),
        ],
    },

    # Corrected against a real deck, which uses this layout five times.
    # The gallery port and the HTML reference both put the title left
    # and the numeral right; the master's own boxes -- and every real
    # slide -- do the opposite. `Text/body` at 45 carries the number,
    # `Title` at 453 carries the words.
    #
    # Colour is chosen per slide rather than fixed. The layout's own
    # background resolves to sand, which reads as washed out against
    # the rest of the deck; a divider is the one place the brand's
    # secondary palette is meant to carry a whole slide. `colour` in
    # the content picks one, and the ink follows the brand rule --
    # blue field takes white type, a secondary field takes black.
    "divider_numbering": {"field": True, "regions": [
        _text(45, 91, 374, 300, "number", 190),
        _text(453, 91, 578, 120, "eyebrow", 13),
        _text(453, 130, 578, 400, "title", 46),
    ]},

    # The next two come from REAL decks rather than the HTML reference,
    # because real decks disagree with it about what these layouts are
    # for. Across two on-brand KONE decks, `Text and picture A` is used
    # 18 times and `Text and picture G` 8 -- more than every other
    # layout combined -- and both carry a pattern the archetypes had no
    # form for: a grid of icon-plus-short-text cells. Slides using them
    # average 7-10 text blocks against the 3 regions bound from the
    # placeholder map.
    #
    # Both keep their plain form too. The engine skips a region whose
    # content key is missing, so one spec serves both: supply `body` and
    # get the paragraph version, supply `items` and get the grid.

    # Measured off slides 6 (plain) and 22 (2x2 grid) of the escalator
    # portfolio deck.
    "text_picture_a": {
        "regions": [
            {"role": "picture", "box": [759, 0, 521, 720], "content": "image"},
            _text(45, 39, 917, 36, "eyebrow", 14,
                  font="KONE Information", color="1450F5", caps=True),
            _text(45, 91, 577, 104, "title", 34),
            _text(45, 227, 577, 402, "body", 16),
        ],
        "groups": [{
            "content": "items",
            "origins": [[45, 181], [349, 181], [45, 405], [349, 405]],
            "regions": [
                {"role": "icon", "box": [0, 0, 60, 60]},
                _text(0, 75, 274, 149, "text", 15),
            ],
        }],
    },

    # Measured off slides 4 (numbered) and 7 (icons) of the same deck:
    # a banner photo across the top with the eyebrow reversed out of it,
    # then six cells along the bottom.
    "text_picture_g": {
        # only when there IS a photo -- a scrim over white is a grey band
        "scrims": [{"box": [0, 0, 1280, 440], "content": "image"}],
        "regions": [
            {"role": "picture", "box": [0, 0, 1280, 440], "content": "image"},
            _text(45, 39, 917, 36, "eyebrow", 14,
                  font="KONE Information", color="FFFFFF", caps=True),
            _text(45, 262, 697, 125, "title", 40, color="FFFFFF"),
            _text(45, 470, 1190, 40, "body", 16),
        ],
        "groups": [{
            "content": "items",
            "origins": [[45, 476], [249, 476], [453, 476], [657, 476], [861, 476], [1065, 476]],
            "regions": [
                {"role": "icon", "box": [0, 0, 60, 60]},
                _text(0, 71, 170, 83, "text", 15),
            ],
        }],
    },
}

# Refinements that must REPLACE an existing registration rather than
# defer to it. `install` normally leaves a hand-built archetype alone,
# but these are transcribed from the current rendered reference while
# the incumbent came from the superseded gallery markup -- so here the
# newer source wins.
_OVERRIDE = frozenset({"divider_numbering", "agenda_a_table"})


# Archetypes `ARCHETYPES.md` marks `no master`. All but one are already
# hand-built in the engine; TIMELINE is Grade A -- one of the twelve
# that carry most real decks -- so it is transcribed here from its
# reference rendering (`KoneDeck.dc.html` slide 10) rather than left as
# the single hole in the vocabulary. The rail is a 2px hairline behind
# four evenly spaced nodes on the 1190px content width.
_NO_MASTER: dict[str, dict] = {
    "timeline": {
        "regions": [
            {"role": "title", "box": [45, 91, 917, 45], "content": "title"},
            {"role": "axis", "box": [45, 251, 1190, 0]},
        ],
        "groups": [{
            "content": "items",
            "origins": [[45, 240], [342, 240], [639, 240], [936, 240]],
            "regions": [
                {"role": "stat_label", "box": [0, 44, 269, 40], "content": "period"},
                {"role": "body", "box": [0, 92, 269, 120], "content": "text"},
            ],
        }],
    },
}


def build_archetypes(
    grades: Iterable[str] = ("A", "B", "C", "D"),
) -> dict[str, dict]:
    """Every canonical archetype that names a master layout, as engine
    dicts keyed by its `lower_snake` engine key.

    Twins are included, resolved to their parent's geometry. The spec's
    advice ("use the parent") is about which one a designer should
    reach for; a user who writes `TITLE_CONTENT_B` in a brief still
    needs it to render, and by the spec's own definition a twin is
    geometrically identical to its parent.
    """
    layouts, archetypes = load_spec()
    wanted = set(grades)
    out: dict[str, dict] = {}

    for arch in archetypes.values():
        if arch.grade not in wanted:
            continue
        source = arch
        if arch.twin_of:
            source = archetypes.get(arch.twin_of, arch)
        if not source.master:
            if arch.engine_key in _NO_MASTER:
                out[arch.engine_key] = _NO_MASTER[arch.engine_key]
            continue
        layout = layouts.get(source.master)
        if layout is None:
            continue
        spec = _bind_roles(layout)
        # a twin inherits its parent's refinement along with its geometry
        refinement = _REFINEMENTS.get(arch.engine_key) or (
            _REFINEMENTS.get(source.engine_key) if arch.twin_of else None
        )
        if refinement:
            # the reference wins on everything it describes; anything it
            # is silent about (a panel, a picture) keeps its bound form
            spec = {**spec, **refinement}
        if not spec["regions"] and not spec.get("groups") and not arch.is_built:
            # A layout with no bindable content is a blank slide -- real
            # for `BLANK` (logo and nothing else, and the designated
            # replacement for the cut `END_LOGO`), meaningless for a
            # layout that was never built.
            continue
        out[arch.engine_key] = spec
    return out


def _correct_grey_ink(archetypes_module) -> list[str]:
    """Take the engine's caption grey off the brand's type roles.

    The brand allows exactly three inks for type -- black, white and (for
    KONE Information only) blue -- and the engine sets `caption`,
    `body_muted` and `attribution` in a `#727272` grey that appears
    nowhere in the palette. It is visible on any deck carrying a quote or
    a captioned statistic, which is most of them.

    Corrected here rather than in the engine because the engine is
    vendored from the skill and is replaced wholesale when the skill
    updates; a patch applied at install time survives that.
    """
    engine = getattr(archetypes_module, "E", None)
    styles = getattr(engine, "ROLE_STYLE", None)
    if not styles:
        return []
    grey = getattr(engine, "GREY", None)
    black = getattr(engine, "BLACK", None)
    if grey is None or black is None:
        return []

    corrected = []
    for role, style in list(styles.items()):
        if len(style) < 3 or style[2] != grey:
            continue
        styles[role] = tuple(black if i == 2 else v for i, v in enumerate(style))
        corrected.append(role)
    return sorted(corrected)


# Corrections to shipped archetypes, taken from the 25+25 handoff. Each
# names the document that settles it, so re-speccing an archetype can
# absorb the fix rather than rediscover it.
_SPEC_FIXES: dict = {
    # INTERNAL_25 slide 24: "Contact line at top:566". It shipped at 654,
    # which put it exactly on top of the footer date at 658 -- and the
    # role was `body_muted`, the grey BRAND_MODE bans.
    "resource_links": {"contact": {"y": 566, "role": "body"}},
    # A section break is a pause, and a pause reads as centred. Both the
    # numeral and the title sat hard against the top band, which made a
    # divider look like a content slide someone forgot to fill in. Only
    # the vertical changes: the 300px numeral is centred on y=360 and the
    # label and title are centred as a pair beside it. The x positions
    # are left where they are, because those were measured off a real
    # deck that uses this layout five times, and the spec's own x:620
    # disagrees with all five.
    "divider_numbering": {
        "number":  {"y": 210},
        "eyebrow": {"y": 276},
        "title":   {"y": 304, "h": 150},
    },
    "divider_title_only": {"title": {"y": 290, "h": 150}},
    # A quote slide with no quote type on it. Ported from the master's
    # boxes, every region on this one came through as 16px `body`, so
    # the quotation sat in a 349px pink panel at the size of a footnote
    # and the attribution below it looked identical to it. The keys are
    # renamed too: `body2`/`body3` tell neither the planner nor the
    # person editing the slide what belongs there.
    "quote_b": {
        "body":  {"role": "quote", "content": "quote"},
        "body3": {"role": "attribution", "content": "attribution"},
        "body2": {"content": "context"},
    },
    # The standfirst under a cover's title. `body` told the planner
    # nothing, so covers came back either empty under the headline or
    # carrying a paragraph where one line belongs. Both external
    # contracts call it `context`; so does the internal prose.
    # `quote_a` is `quote_b` in blue, with the same defect and the same
    # fix -- except the type has to be white, because black on KONE Blue
    # is unreadable.
    "quote_a": {
        "body":  {"role": "quote_light", "content": "quote"},
        "body3": {"role": "attribution_light", "content": "attribution"},
        "body2": {"content": "context"},
    },
    # Measured off the reference: the cover title runs 653px wide, not
    # the 578 the content column gives it, because a cover headline is
    # allowed to run past the column into the open right-hand side.
    # The reference sets a cover headline at 76px and its standfirst at
    # subtitle size. Bound from the master the title resolved to 30px and
    # the standfirst to 12 -- a cover that read like a content slide.
    "cover_a_cut4":      {"body": {"content": "context", "role": "cover_context",
                                   "y": 588, "h": 40},
                          "title": {"role": "cover_title", "w": 653, "y": 427}},
    "cover_b_cut3":      {"body": {"content": "context", "role": "cover_context"},
                          "title": {"role": "cover_title"}},
    "cover_f_fullbleed": {"body": {"content": "context",
                                        "role": "cover_context_light"},
                          "title": {"role": "cover_title_light"}},
    # "two 19px white lines" -- named as the handoff names them.
    "outro": {"body": {"content": "text1"}, "body2": {"content": "text2"}},
}


# Archetypes the contract says take structured content and the registry
# only ever offered a paragraph. These are not tweaks to a bound layout,
# they are the slide redrawn from the as-built prose in `INTERNAL_25.md`,
# so they replace the incumbent's geometry outright rather than patching
# a box at a time. The incumbent's background is kept: it is the one
# thing the prose does not restate.
#
# Why this matters more than it looks: an archetype whose only slot is
# `body` cannot be chosen well. A planner reaching for the timeline
# found it could put two paragraphs there, and went back to the handful
# of layouts with real slots. Every entry here is a layout returning to
# the menu.
# Slide 8 of "Life, upgraded in ONE week", which is the densest slide in
# that deck and the one it holds together best: a title at y=22 and a
# grid of white cards on sand, each with a coloured rule under a caps
# label. Registered as its own archetype because nothing in the two sets
# does this, and it is the shape a positioning or framing slide wants.
_CARD_GRID = {
    "background": "sand",
    "regions": [
        _text(45, 22, 1050, 82, "title", 32),
    ],
    "cards": None,          # filled by `card_grid()` at install time
}


_RESPEC: dict[str, dict] = {
    # 03 "Mint column, 420px, full height ... 44px title, 16px lead.
    #     Right column at x:510: five rows, 14px gap, each a sand block
    #     with 20px 24px padding, 44px blue numeral chip and 24px label."
    # The fifth row inverting to blue is not expressible per item, so
    # every row is sand; the mint column and the numerals carry it.
    "agenda_c_split": {
        "panels": [{"box": [0, 0, 420, 720], "fill": "AAE1C8"}]
                  + [{"box": [510, y, 725, 96], "fill": "F3EEEA"}
                     for y in (91, 201, 311, 421, 531)],
        "regions": [
            _text(45, 91, 330, 150, "title", 40),
            _text(45, 260, 330, 160, "lead", 16),
        ],
        "groups": [{
            "content": "items",
            "origins": [[510, 91], [510, 201], [510, 311], [510, 421], [510, 531]],
            "regions": [
                _text(20, 26, 48, 44, "number", 24, color="1450F5"),
                _text(84, 30, 610, 44, "label", 22),
            ],
        }],
    },
    # 12 "40px title in a 340px column. Pink panel at top:250, 24px
    #     padding, lead plus three bullets. Right: 770x6 blue axis at
    #     left:465 top:300, four stems below -- 6x40 blue drop, 40px blue
    #     chip, KONE Information period label, 16px text in a 165px
    #     column."
    # The drops are fixed panels rather than per-item shapes, which is
    # why they are listed out: four events is what the axis was drawn
    # for, and a fifth has nowhere to hang.
    "timeline_quarter_axis": {
        "panels": [
            {"box": [45, 250, 374, 379], "fill": "FFCDD7"},
            {"box": [465, 300, 770, 6], "fill": "1450F5"},
        ] + [{"box": [x, 306, 6, 40], "fill": "1450F5"}
             for x in (465, 663, 861, 1059)],
        "regions": [
            _text(45, 91, 374, 104, "title", 40),
            _text(69, 274, 326, 60, "lead", 17),
            _bullets(69, 350, 326, 255, "bullets", 16),
        ],
        "groups": [{
            "content": "events",
            "origins": [[465, 306], [663, 306], [861, 306], [1059, 306]],
            "regions": [
                _text(0, 56, 165, 20, "period", 12,
                      font="KONE Information", color="1450F5", caps=True),
                _text(0, 82, 165, 240, "text", 16),
            ],
        }],
    },
    # "The running order as a numbered list" -- built as two parallel
    # text boxes, an 83px column of numbers beside a 476px column of
    # words. Nothing about that is a list: the rows cannot be edited in
    # PowerPoint's outline view, and a planner had to write the numbers
    # itself and hope the line breaks lined up. Five real rows instead,
    # on the same 578px column the layout already declares.
    "agenda_b_numbered": {
        "regions": [_text(45, 91, 578, 53, "title", 32)],
        "groups": [{
            "content": "items",
            "origins": [[45, y] for y in (181, 270, 359, 448, 537)],
            "regions": [
                _text(0, 0, 83, 44, "number", 24, color="1450F5"),
                _text(102, 4, 476, 80, "label", 20),
            ],
        }],
    },
    # 02 "Full-height photo right, 574px wide from x:706. Left: blue
    #     eyebrow at top:130, 48px statement at top:166, then three rows
    #     at top:392 -- 44px blue chip with white pictogram, 19px black
    #     text, 20px gap."
    # What was registered instead was a banner photo across the top with
    # a title and one line under it -- a different, weaker slide, and
    # the reason the "why we're here" opener never carried its three
    # reasons.
    "picture_intro": {
        "regions": [
            {"role": "picture", "box": [706, 0, 574, 720], "content": "photo"},
            _text(45, 130, 610, 24, "eyebrow", 12,
                  font="KONE Information", color="1450F5", caps=True),
            _text(45, 166, 610, 190, "title", 44),
        ],
        "groups": [{
            "content": "points",
            "origins": [[45, 392], [45, 456], [45, 520]],
            "regions": [
                {"role": "icon", "box": [0, 0, 44, 44]},
                _text(60, 8, 550, 44, "text", 19),
            ],
        }],
    },
}


# Slots that draw the archetype's OWN footer line, at y:664-680. They
# predate the chrome layer: `stamp_chrome` now puts the date at 45,658
# and the page number at 1167,658 on every body slide, so these are a
# second footer sitting under the first, below the floor, in a slot the
# planner is invited to fill. Dropped rather than moved -- there is
# nowhere above the floor for a duplicate to go.
_DROP_SLOTS: dict[str, tuple] = {
    "kone_numbers": ("footer",),
    "credits": ("footer",),
    "milestone_slide": ("classification",),
}


def _drop_slots(registry) -> None:
    for name, keys in _DROP_SLOTS.items():
        spec = registry.get(name)
        if not isinstance(spec, dict) or not spec.get("regions"):
            continue
        spec["regions"] = [r for r in spec["regions"] if r.get("content") not in keys]


def hold_to_the_floor(spec: dict, floor: float = 629.0, least: float = 16.0) -> int:
    """Bring anything hanging below the content floor back above it.

    BRAND_MODE §7 is unconditional -- every region's bottom is <= 629 --
    and thirteen of the built archetypes broke it, one of them by 83px.
    They were invisible because nothing built the whole library at once
    and read the preflight back; the contract work did, and there they
    were.

    Trim first, because a box that is merely too tall loses nothing by
    being shortened. Lift only when trimming would leave less than a
    line of type. Pictures and full-bleed panels are exempt: reaching
    the bottom edge is what they are for.

    Returns how many boxes moved.
    """
    moved = 0

    def _fix(box) -> bool:
        top, height = box[1], box[3]
        if top + height <= floor or _is_full_bleed(box):
            return False
        if floor - top >= least:
            box[3] = floor - top
        else:
            box[1] = max(0, floor - height)
        return True

    for region in spec.get("regions") or []:
        if region.get("role") in _PICTURE_REGION_ROLES or region.get("role") == "image":
            continue
        if _fix(region["box"]):
            moved += 1

    for group in spec.get("groups") or []:
        origins = group.get("origins") or []
        if not origins:
            continue
        lowest = max(o[1] for o in origins)
        for region in group.get("regions") or []:
            if region.get("role") in _PICTURE_REGION_ROLES:
                continue
            box = region["box"]
            bottom = lowest + box[1] + box[3]
            if bottom <= floor:
                continue
            # A group's boxes are shared by every origin, so the trim is
            # computed against the LOWEST row and applied to all of them.
            room = floor - lowest - box[1]
            if room >= least:
                box[3] = room
            else:
                box[1] = max(0, box[1] - (bottom - floor))
            moved += 1
    return moved


def _register_brand_roles(archetypes_module) -> list:
    """Teach the engine every role BRAND_MODE defines.

    The engine ships 28 roles; `TYPE_SCALE` defines 42. The fourteen it
    had never heard of included `cover_title` -- so the moment a cover
    was told to set its headline at 76px the way the reference does, the
    build died on a KeyError mid-draw. A role that exists in the brand
    and not in the renderer is a size nobody can ask for.

    Registered rather than mapped: the engine's own entries win where
    both define a role, because those were tuned against real renders.
    """
    from pptx.dml.color import RGBColor

    from deckguard import brandmode as bm

    engine = getattr(archetypes_module, "E", None)
    styles = getattr(engine, "ROLE_STYLE", None)
    if styles is None:
        return []
    added = []
    for role, entry in bm.TYPE_SCALE.items():
        if role in styles:
            continue
        font, px, _weight, _lead, _track, colour, caps = entry
        styles[role] = (font, px, RGBColor.from_string(colour), caps, False, False)
        added.append(role)
    return sorted(added)


def _add_light_roles(archetypes_module) -> None:
    """A quote set white, for the quote slide that sits on a blue panel.

    `quote_a` is `quote_b` in blue, and it had the same defect: every
    region ported as `body`, so the quotation was 16px in a 349px panel
    and the attribution under it looked identical. `quote_b` could take
    the black `quote` and `attribution` roles; this one cannot, because
    black on KONE Blue is unreadable -- so the two roles exist here in
    white, and nowhere else.
    """
    engine = getattr(archetypes_module, "E", None)
    styles = getattr(engine, "ROLE_STYLE", None)
    if not styles:
        return
    for role, source, white in (("quote_light", "quote", True),
                                ("attribution_light", "attribution", True)):
        if role in styles or source not in styles:
            continue
        base = list(styles[source])
        if len(base) >= 3:
            from pptx.dml.color import RGBColor

            base[2] = RGBColor(0xFF, 0xFF, 0xFF) if white else base[2]
        styles[role] = tuple(base)


def tighten_band(spec: dict, name: str = "") -> int:
    """Move a content slide's title to the reference's tighter band.

    The KONE 25 prose puts a title at y=91 and content at y=227. The real
    deck puts the title at y=22 and its first content row at y=118, and
    that 109px is what lets a twelve-card grid or a three-photo row
    breathe instead of crowding the floor.

    Only slides that actually sit on the standard band move, and only
    down-shifted by the same delta, so the internal spacing every layout
    was measured with is preserved exactly -- the block moves, its parts
    do not move relative to each other. Covers, dividers and the outro
    are left alone: their type is placed against the photograph or the
    field, not against a title band.

    Returns how many boxes moved.
    """
    from deckguard import brandmode as bm

    if bm.slide_kind(name) in ("cover", "divider", "outro", "fullslide_picture",
                               "blank"):
        return 0
    title = next((r for r in spec.get("regions") or []
                  if r.get("content") == "title"), None)
    if title is None or not (80 <= title["box"][1] <= 100):
        return 0
    # An eyebrow sits ABOVE the title, so shifting the block by the
    # title's delta drove it off the top edge -- "SERVICE PLAN" came out
    # half-clipped at y=0. The reference agrees: its dense card slide has
    # no eyebrow and sets the title at y=22, while the slide that DOES
    # carry an eyebrow puts its title at y=95, which is where these
    # already are. So a slide with an eyebrow is already on its band.
    above = [r for r in spec.get("regions") or []
             if r.get("content") and r["box"][1] < title["box"][1] and r["box"][1] > 10]
    if above:
        return 0
    delta = title["box"][1] - bm.TIGHT_TITLE_Y

    moved = 0
    for region in spec.get("regions") or []:
        box = region["box"]
        if _is_full_bleed(box) or box[1] < 20:
            continue
        box[1] = max(0, box[1] - delta)
        moved += 1
    for group in spec.get("groups") or []:
        for origin in group.get("origins") or []:
            origin[1] = max(0, origin[1] - delta)
            moved += 1
    for panel in (spec.get("panels") or []) + (spec.get("rules") or []):
        box = panel["box"]
        if _is_full_bleed(box) or box[1] < 20:
            continue
        box[1] = max(0, box[1] - delta)
        moved += 1
    return moved


def _apply_respecs(registry) -> None:
    for name, spec in _RESPEC.items():
        incumbent = registry.get(name)
        if not isinstance(incumbent, dict):
            continue
        keep = {k: v for k, v in incumbent.items() if k == "background"}
        registry[name] = {**keep, **copy.deepcopy(spec)}


def _apply_spec_fixes(registry) -> None:
    for name, fixes in _SPEC_FIXES.items():
        spec = registry.get(name)
        if not spec:
            continue
        for region in spec.get("regions") or []:
            fix = fixes.get(region.get("content"))
            if not fix:
                continue
            x, y, w, h = region["box"]
            region["box"] = [fix.get("x", x), fix.get("y", y),
                             fix.get("w", w), fix.get("h", h)]
            if "role" in fix:
                region["role"] = fix["role"]
            if "content" in fix:
                region["content"] = fix["content"]


def install(archetypes_module, grades: Iterable[str] = ("A", "B", "C", "D")) -> list[str]:
    """Register the generated archetypes into the engine's registry.

    Anything already implemented -- under its canonical name or under
    one of its `ARCHETYPES.md` aliases -- is left alone. Those were
    built against a real rendering; this module works from a coarser
    description and must not clobber them.
    """
    registry = archetypes_module.ARCHETYPES
    _, meta = load_spec()
    by_key = {a.engine_key: a for a in meta.values()}
    _correct_grey_ink(archetypes_module)
    _register_brand_roles(archetypes_module)
    _add_light_roles(archetypes_module)
    for existing in registry.values():
        if isinstance(existing, dict):
            lift_low_rows(existing)


    added: list[str] = []
    # Archetypes with no master layout behind them, so `build_archetypes`
    # cannot derive them. Registered like any other, and equally never
    # allowed to overwrite an incumbent.
    samples = getattr(archetypes_module, "SAMPLES", None)
    for key, spec in _EXTRAS.items():
        if key not in registry:
            registry[key] = copy.deepcopy(spec)
            added.append(key)
        meta = _EXTRA_META.get(key, {})
        if isinstance(samples, dict) and meta.get("sample") and key not in samples:
            samples[key] = copy.deepcopy(meta["sample"])

    for key, spec in build_archetypes(grades).items():
        arch = by_key.get(key)
        if key in _OVERRIDE and key in registry:
            # keep the incumbent's background, replace its geometry
            registry[key] = {
                **({"background": registry[key]["background"]}
                   if "background" in registry[key] else {}),
                **spec,
            }
            added.append(key)
            continue
        if key in registry:
            continue
        if arch and any(alias in registry for alias in arch.aliases):
            continue
        registry[key] = spec
        added.append(key)

    if added:
        try:  # the matcher caches signatures; new names must invalidate them
            from deckguard.registry import invalidate_archetype_caches

            invalidate_archetype_caches()
        except Exception:
            pass
    # Last of all, so it reaches the archetypes generated just above as
    # well as the incumbents. Run earlier it corrected only the
    # incumbents: a generated archetype was created afterwards and kept
    # the very geometry the fix existed to replace.
    if "card_grid" not in registry:
        registry["card_grid"] = {**copy.deepcopy(_CARD_GRID), **card_grid(8)}
        added.append("card_grid")
    _apply_respecs(registry)
    _apply_spec_fixes(registry)
    _drop_slots(registry)
    # Before the floor pass, so anything the shift frees at the bottom is
    # available to the layouts rather than trimmed away.
    for archetype, spec in registry.items():
        if isinstance(spec, dict):
            tighten_band(spec, archetype)
    # Last of all and over everything, including the archetypes the two
    # passes above just rewrote: the floor is the one rule with no
    # exceptions, so it is enforced after every other hand has been on
    # the geometry rather than trusted to each of them.
    for spec in registry.values():
        if isinstance(spec, dict):
            hold_to_the_floor(spec)
    return sorted(added)


def coverage(archetypes_module) -> dict[str, tuple[int, int]]:
    """{grade: (implemented, total)} for the canonical vocabulary --
    what the tool can actually render, by grade."""
    _, meta = load_spec()
    registry = archetypes_module.ARCHETYPES
    tally: dict[str, list[int]] = {}
    for arch in meta.values():
        got = arch.engine_key in registry or any(a in registry for a in arch.aliases)
        slot = tally.setdefault(arch.grade, [0, 0])
        slot[0] += bool(got)
        slot[1] += 1
    return {g: (v[0], v[1]) for g, v in sorted(tally.items())}


# --------------------------------------------------------------------------
# whole-deck assembly
# --------------------------------------------------------------------------


def protect_photo_cover(slide, engine=None) -> bool:
    """Scrim the master's own cover, which arrives without one.

    The retained Cover F is a full-bleed photograph with the title
    reversed out of it, and the brand asks for a gradient there. The
    master ships none -- the layout assumes a designer will choose a
    photograph with a quiet corner. A deck built from a brief picks its
    photograph automatically, so the assumption does not hold, and a
    half-year review came back with a white headline lost in a sunlit
    atrium.

    Protects the type that is actually written, sized to the type, so
    the photograph keeps the rest of the frame.
    """
    import importlib

    if engine is None:
        engine = importlib.import_module("kone_engine")

    P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    picture = None
    for shape in slide.shapes:
        if shape._element.tag != f"{P}pic":
            continue
        if shape.width >= engine.X(1200) and shape.height >= engine.X(680):
            picture = shape
            break
    if picture is None:
        return False

    bands = []
    for shape in slide.shapes:
        if shape is picture or not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        top = shape.top / 9525.0
        width = shape.width / 9525.0
        height = shape.height / 9525.0
        px = _run_px(shape) or (64.0 if height > 200 else 16.0)
        region = {"box": [0, top, width, height], "dg": {"px": px}}
        bands.append((top, top + _typeset_height(region, text)))
    if not bands:
        return False

    _draw_scrim(slide, engine, [0, 0, 1280, 720], protect=bands)
    return True


def _run_px(shape) -> Optional[float]:
    """The first run's size in px, when the shape states one."""
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    return run.font.size.pt / 0.75
    except Exception:  # noqa: BLE001 -- inherited sizing is normal
        return None
    return None


def _layout_for(archetype, by_partname, blank):
    """The master layout an archetype belongs on, or blank."""
    if archetype is None or not archetype.master:
        return blank
    return by_partname.get(archetype.master) or blank


def _strip_empty_placeholders(slide) -> None:
    """Take the layout's prompt text off a slide we are drawing over.

    Cloned placeholders arrive carrying "Click to add title". The
    archetype draws its own text boxes on top, so the prompts sit
    underneath as live text -- invisible in most renders, and there in
    the outline view and in search.
    """
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)


def build_deck(spec: dict, out_path, archetypes_module=None, report=None) -> str:
    """Assemble a deck from a spec, each archetype on its OWN layout.

    The skill's own `kone_deck_creator.build_deck` puts every archetype
    on the BLANK layout and calls `archetypes.render` directly. That
    bypasses everything this module adds, and it showed on the first
    deck a user built through the web tool: icons came out as the
    engine's rasterised PNGs and its blue placeholder chips instead of
    the real KONE pictograms as editable shapes, covers had no scrim so
    white headlines sat on sunlit photographs, dividers came out sand,
    photographs were stretched rather than cropped, and a slide acquired
    a donut chart nobody asked for.

    Same shape as the skill's builder -- master cover, generated body,
    master "Thank you" -- but each body slide is added on the layout its
    archetype actually belongs to, its prompt placeholders are stripped,
    and it renders through `render` above.
    """
    import posixpath

    from pptx import Presentation
    from pptx.oxml.ns import qn

    if archetypes_module is None:
        # Through the loader, never a bare `import archetypes`: the
        # registry is only correct once the gallery's ports and then
        # THIS module's refinements have been merged, in that order.
        # `gallery.install` overwrites what it finds, so installing it
        # second silently reverted the refined agenda to a port with no
        # bullets and the divider to one with no number and no colour.
        from deckguard.registry import _load_archetypes

        archetypes_module = _load_archetypes()

    creator = _load_creator()
    prs = Presentation(creator.MASTER)
    slide_ids = prs.slides._sldIdLst
    originals = list(slide_ids)
    intro, outro = originals[creator.INTRO_IDX], originals[creator.OUTRO_IDX]
    if spec.get("title"):
        creator._set_title(list(prs.slides)[creator.INTRO_IDX], spec["title"])

    _, meta = load_spec()
    by_engine_key: dict[str, object] = {}
    for archetype in meta.values():
        by_engine_key.setdefault(archetype.engine_key, archetype)
        for alias in archetype.aliases:
            by_engine_key.setdefault(alias, archetype)

    by_partname = {
        posixpath.basename(layout.part.partname).replace(".xml", ""): layout
        for layout in prs.slide_layouts
    }
    blank = next(l for l in prs.slide_layouts if l.name.strip().lower() == "blank")

    slides = _with_cut_cover(spec)
    for position, entry in enumerate(slides, start=1):
        name = entry.get("archetype")
        content = {k: v for k, v in entry.items() if k != "archetype"}
        archetype = by_engine_key.get(str(name).lower())
        slide = prs.slides.add_slide(_layout_for(archetype, by_partname, blank))
        _strip_empty_placeholders(slide)
        if report is not None:
            dropped = unread_keys(archetypes_module.ARCHETYPES.get(name) or {}, content)
            if dropped:
                report.setdefault("dropped", {})[position] = (name, dropped)
        render(slide, name, content, archetypes_module,
               audience=str(spec.get("audience") or ""))
        stamp_chrome(
            slide, archetypes_module.E, str(name),
            page=position,              # the generated cut cover is page 1
            date=str(spec.get("date") or content.get("footer") or _deck_date()),
            classification=str(content.get("classification") or spec.get("classification") or ""),
        )

    body = [el for el in list(slide_ids) if el not in originals]
    # The master's own cover is dropped: `_with_cut_cover` has put a
    # generated four-pane cut at the front of `body`, and keeping both
    # would open every deck on two covers.
    keep = {id(outro), *(id(b) for b in body)}
    for element in originals:
        if id(element) not in keep:
            prs.part.drop_rel(element.get(qn("r:id")))
            slide_ids.remove(element)
    for element in list(slide_ids):
        slide_ids.remove(element)
    for element in body:
        slide_ids.append(element)
    slide_ids.append(outro)

    protect_photo_cover(list(prs.slides)[0])
    strip_master_classification(prs)

    prs.save(str(out_path))
    return str(out_path)


def _deck_date() -> str:
    """The master's own footer format: `12 MARCH 2026`."""
    from datetime import date

    return date.today().strftime("%d %B %Y").lstrip("0")


def _load_creator():
    from deckguard.registry import _load_creator as load

    return load()


# --------------------------------------------------------------------------
# what an archetype actually reads
# --------------------------------------------------------------------------


def content_keys(spec: dict) -> list[str]:
    """The content keys a spec reads, annotated with their shape.

    The planning prompt described archetypes from `catalog.json` slots
    and the skill's `SAMPLES`, and 39 of the 80 registered archetypes
    have neither -- the model was told a name and left to guess. Worse,
    a stale sample is an active lie: `agenda_a_table`'s said
    `text1..text4` while the renderer had been rebuilt to read `items`,
    so a planner did as it was told, emitted four keys nothing reads,
    and the agenda came out as a title on an empty slide.

    Derived from the live registry, so it cannot drift.
    """
    out: list[str] = []
    for region in spec.get("regions", []):
        key = region.get("content")
        if not key:
            continue
        out.append(f"{key} ({_shape_of(region)})")
    for group in spec.get("groups", []):
        key = group.get("content")
        if not key:
            continue
        fields = [
            # an icon region carries no content key of its own, but the
            # caller names its pictogram per item -- omitting it from the
            # guide is why decks came back with default icons
            "icon" if r.get("role") == "icon" else r.get("content")
            for r in group.get("regions", [])
            if r.get("content") or r.get("role") == "icon"
        ]
        n = len(group.get("origins") or [])
        if fields:
            shape = "{" + ", ".join(fields) + "}"
            out.append(f"{key} (list of up to {n} × {shape})")
        else:
            out.append(f"{key} (list of up to {n})")
    return out


def _shape_of(region: dict) -> str:
    role = str(region.get("role") or "")
    if role in _PICTURE_REGION_ROLES or role == "image":
        return "filled automatically -- do not supply"
    if role == "icon":
        return "icon name"
    style = region.get("dg") or {}
    # `role == "bullets"` matters as much as the dg kind: a region bound
    # from the master carries the role and no dg style, so
    # `org_functions.functions` -- a real bulleted panel -- was
    # advertised as `(text)` and got a paragraph.
    if style.get("kind") == "bullets" or role == "bullets":
        return "list of strings, or {text, sub:[...]}"
    if "stat" in role or "value" in role:
        return "short figure, e.g. 70%"
    if role == "table":
        # It was annotated `(text)`, so the one archetype built to hold a
        # real table was advertised as taking a paragraph -- and got one.
        return '{headers: [...], rows: [[...], ...]}'
    return "text"


def unread_keys(spec: dict, content: dict) -> list[str]:
    """Content the archetype has nowhere to put.

    Silent loss is the failure mode this catches: a planner emitted
    `text1..text4` for an agenda whose renderer reads `items`, and the
    slide came out as a title on an empty half. Nothing said so -- not
    the build, not the review page, not the audit. The deck simply had
    less in it than the brief did.
    """
    known = {k.split(" (")[0] for k in content_keys(spec)}
    ignore = {"archetype", "colour", "color", "notes"}
    return sorted(
        key for key, value in (content or {}).items()
        if key not in known and key not in ignore and value not in (None, "", [], {})
    )


# --------------------------------------------------------------------------
# closing dead bands
# --------------------------------------------------------------------------

# How much empty vertical space between a title and the row beneath it
# counts as a hole rather than as breathing room, and the gap left after
# closing one. 69px is the master's own step: a 104px title starting at
# 91 ends at 195, and the grid's content start is 264.
_DEAD_BAND = 200.0
_TITLE_GAP = 69.0


def lift_low_rows(spec: dict) -> int:
    """Pull an item row up under its title when nothing fills the middle.

    Four archetypes place their row of items in the bottom third with
    only a title above -- geometry taken from real KONE slides where a
    paragraph of body copy filled the band between. These archetypes have
    no such paragraph, so a Q2 review came back with 248px of blank sand
    between "Plan of action" and the six things it listed.

    Only fires where the band is demonstrably empty: no region and no
    other group occupies it. Rows only ever move UP, so nothing below can
    be collided into. Returns how many groups moved.
    """
    groups = spec.get("groups") or []
    if not groups:
        return 0

    solid = [
        r for r in spec.get("regions", [])
        if not _is_full_bleed(r["box"])
    ]
    moved = 0
    for group in sorted(groups, key=lambda g: min((o[1] for o in g["origins"]), default=0)):
        origins = group.get("origins") or []
        if not origins:
            continue
        top = min(o[1] for o in origins)
        bottoms = [r["box"][1] + r["box"][3] for r in solid if r["box"][1] + r["box"][3] <= top]
        bottoms += [
            max(o[1] for o in other["origins"]) + _group_height(other)
            for other in groups
            if other is not group and other.get("origins")
            and max(o[1] for o in other["origins"]) + _group_height(other) <= top
        ]
        if not bottoms:
            continue
        floor = max(bottoms)

        # anything straddling the band means the space is spoken for
        straddles = any(
            r["box"][1] < top and r["box"][1] + r["box"][3] > floor for r in solid
        ) or any(
            min(o[1] for o in other["origins"]) < top
            and max(o[1] for o in other["origins"]) + _group_height(other) > floor
            for other in groups if other is not group and other.get("origins")
        )
        if straddles:
            continue

        delta = (top - floor) - _TITLE_GAP
        if top - floor <= _DEAD_BAND or delta <= 0:
            continue
        group["origins"] = [[x, y - delta] for x, y in origins]
        moved += 1
    return moved


def _group_height(group: dict) -> float:
    return max((r["box"][1] + r["box"][3] for r in group.get("regions", [])), default=0.0)


def _is_full_bleed(box) -> bool:
    """A background photograph is not a block the row has to clear."""
    _, _, w, h = box
    return w >= 1200 and h >= 680


# --------------------------------------------------------------------------
# archetypes with no master layout behind them
# --------------------------------------------------------------------------

# The stat band is five cells across 1190 with a 40px gap, so each cell
# is 206 wide and they start every 246px. The row is 152 tall and
# vertically centred on content 90 tall, which puts the number at 307
# and its label at 381.
_MILESTONE_STAT_X = [45 + i * 246 for i in range(5)]

# Routing and a worked example for the extras. `catalog.json` and
# `SAMPLES` describe the archetypes the skill shipped with; an extra
# registered from here is invisible to both, so a planner would see the
# key list and no reason to ever choose it.
_EXTRA_META: dict[str, dict] = {
    "milestone_slide": {
        "purpose": "A finished thing, its proof, and who did it -- one shareable "
                   "slide built from an announcement. Not for a proposal, a "
                   "decision request, or anything needing an argument across "
                   "several beats; those are decks.",
        "keywords": ["milestone", "announcement", "launch", "now live",
                     "migration complete", "programme win", "recognition",
                     "thank you", "transformation", "results", "one slide"],
        # The brand's rules for this slide, in the terms a planner has to
        # act on. They are in the skill's own SKILL.md, which the planner
        # never sees -- it sees this guide.
        "notes": [
            "`title` must be under about 55 characters -- it is set at 42px and "
            "longer wraps onto a second line, closing the gap above `lede`.",
            "Each `stats` value is a BARE number: {\"value\": \"6\", \"label\": "
            "\"Weeks end to end\"}, never {\"value\": \"6 weeks\", \"label\": "
            "\"End-to-end migration\"} -- the unit belongs in the label, and a "
            "value carrying its own unit is both wider and says it twice.",
            "Include one stat that is deliberately zero -- no disruption, no "
            "downtime, no escalations -- written as the digit \"0\". It renders "
            "black against the others' blue because it is a different kind of "
            "claim. Four strong numbers beat five where one is filler.",
            "`done` text must fit one line (about 45 characters). Three is the "
            "ceiling and the column is sized for exactly three single-line "
            "items.",
            "`scope` must name something the numbers cannot. If it restates a "
            "count already in `stats` -- the 12 frontlines, the 100+ users -- "
            "it is saying nothing twice; cut it or name the units instead.",
            "A completion fact (\"100% transitioned\") is a `done` tick, not a "
            "stat. Do not put it in both.",
        ],
        "sample": {
            "eyebrow": "Marketing Hub · Request Management",
            "title": "From Monday.com to ServiceNow in six weeks",
            "lede": "Delivered by KBS, Global Marketing, Frontlines and Business "
                    "Partner — with full data continuity.",
            "done": [{"text": "MVP pilot-tested and now live"},
                     {"text": "100% transition completed"}],
            "stats": [{"value": "6", "label": "Weeks end to end"},
                      {"value": "100+", "label": "Users migrated"},
                      {"value": "12", "label": "Frontlines"},
                      {"value": "3+3", "label": "Regions + global teams"},
                      {"value": "0", "label": "Business disruption"}],
            "scope_label": "The frontlines",
            "scope": "KSEA · KMTA · KANZ · KEI · EEM · DACH · GIN · Nordics",
            "next_label": "What's next",
            "next": ["Hypercare and small fixes ongoing",
                     "Power BI reporting integration underway, targeted for Q2"],
            "credits_label": "Thank you",
            "credits": "Arvind, Suresh Kumar, Rupesh and the Hub specialists",
            "classification": "KONE Internal",
        },
    },
}

# One slide, from one announcement email. It exists because the master
# has no recognition slide: a finished thing, its proof, and who did it.
# Transcribed from kone-milestone-slide's published geometry rather than
# eyeballed, so it stays comparable with the reference.
# The band as a slide of its own -- same geometry, because a number
# that moved between the two forms would read as a different number.
_MILESTONE_BAND: dict = {
    "panels": [{"box": [0, 276, 1280, 196], "fill": "F3EEEA"}],
    "rules": [{"box": [45, 424, 1190, 1], "fill": "141414"}],
}

# Four columns across 1190, so 297 wide on a 297.5 pitch. Rows are
# 26px padding + a 24px name + 26px, which is 81 -- three rows of four
# reach 483 and still clear the closing line at 512.
_CREDIT_X = [45, 342, 640, 937]
_CREDIT_Y = [240, 321, 402]

_EXTRAS: dict[str, dict] = {
    "kone_numbers": {
        "background": "FFFFFF",
        "panels": list(_MILESTONE_BAND["panels"]),
        "rules": list(_MILESTONE_BAND["rules"]),
        "regions": [
            _text(45, 47, 790, 20, "eyebrow", 12,
                  font="KONE Information", color="1450F5", caps=True),
            _text(45, 91, 1189, 44, "title", 34),
            _text(45, 438, 150, 18, "scope_label", 11,
                  font="KONE Information", color="1450F5", caps=True),
            _text(205, 438, 1030, 18, "scope", 13, font="KONE Information"),
            _text(45, 664, 500, 16, "footer", 11,
                  font="KONE Information", caps=True),
        ],
        "groups": [
            {
                "content": "stats",
                "origins": [[x, 307] for x in _MILESTONE_STAT_X],
                "regions": [
                    {"role": "dg_text", "box": [0, 0, 206, 66], "content": "value",
                     "dg": {"kind": "text", "px": 62, "font": "KONE Information",
                            "color": "1450F5", "caps": False, "align": "l",
                            "zero_is_black": True}},
                    {"role": "dg_text", "box": [0, 74, 206, 32], "content": "label",
                     "dg": {"kind": "text", "px": 12, "font": "KONE Information",
                            "color": "141414", "caps": True, "align": "l"}},
                ],
            },
        ],
    },
    "credits": {
        "background": "FFFFFF",
        "regions": [
            _text(45, 47, 790, 20, "eyebrow", 12,
                  font="KONE Information", color="1450F5", caps=True),
            _text(45, 91, 1189, 44, "title", 34),
            {"role": "dg_text", "box": [45, 512, 1190, 60], "content": "note",
             "dg": {"kind": "ruled", "px": 20, "font": "Inter", "color": "141414",
                    "caps": False, "align": "l", "pad": 24, "rule": "D0D0D0"}},
            _text(45, 664, 500, 16, "footer", 11,
                  font="KONE Information", caps=True),
        ],
        "groups": [
            {
                "content": "names",
                "origins": [[x, y] for y in _CREDIT_Y for x in _CREDIT_X],
                "regions": [
                    {"role": "dg_text", "box": [0, 0, 297, 81], "content": "name",
                     "dg": {"kind": "ruled", "px": 24, "font": "Inter",
                            "color": "141414", "caps": False, "align": "l",
                            "pad": 26, "rule": "D0D0D0"}},
                ],
            },
        ],
    },
    "milestone_slide": {
        "background": "FFFFFF",
        "panels": [{"box": [0, 276, 1280, 196], "fill": "F3EEEA"}],
        "rules": [
            {"box": [45, 424, 1190, 1], "fill": "141414"},     # over the scope strip
            {"box": [45, 544, 700, 1], "fill": "D0D0D0"},      # under "What's next"
            {"box": [880, 544, 355, 1], "fill": "D0D0D0"},     # under "Thank you"
        ],
        "regions": [
            _text(45, 47, 790, 20, "eyebrow", 12,
                  font="KONE Information", color="1450F5", caps=True),
            _text(45, 82, 790, 104, "title", 42),
            _text(45, 186, 700, 76, "lede", 17),

            _text(45, 438, 150, 18, "scope_label", 11,
                  font="KONE Information", color="1450F5", caps=True),
            _text(205, 438, 1030, 18, "scope", 13, font="KONE Information"),

            _text(45, 520, 700, 18, "next_label", 12,
                  font="KONE Information", color="1450F5", caps=True),
            _bullets(45, 560, 700, 110, "next", 16, lead=0.5),

            _text(880, 520, 355, 18, "credits_label", 12,
                  font="KONE Information", color="1450F5", caps=True),
            _text(880, 560, 355, 100, "credits", 16),

            _text(45, 664, 400, 16, "classification", 11,
                  font="KONE Information", caps=True),
        ],
        "groups": [
            {
                "content": "stats",
                "origins": [[x, 307] for x in _MILESTONE_STAT_X],
                "regions": [
                    {"role": "dg_text", "box": [0, 0, 206, 66], "content": "value",
                     "dg": {"kind": "text", "px": 62, "font": "KONE Information",
                            "color": "1450F5", "caps": False, "align": "l",
                            "zero_is_black": True}},
                    {"role": "dg_text", "box": [0, 74, 206, 32], "content": "label",
                     "dg": {"kind": "text", "px": 12, "font": "KONE Information",
                            "color": "141414", "caps": True, "align": "l"}},
                ],
            },
            {
                # Three is the published ceiling, and the column has
                # 186->276 to do it in before the sand band starts: 90px,
                # which is exactly three 20px badges on a 14px gap. A
                # roomier pitch fits two and puts the third INSIDE the
                # band, which is what shipped before this was measured.
                "content": "done",
                "origins": [[880, 186], [880, 220], [880, 254]],
                "regions": [
                    {"role": "dg_tick", "box": [0, 0, 355, 20], "content": "text",
                     "dg": {"kind": "tick", "px": 16, "font": "Inter",
                            "color": "141414", "caps": False, "align": "l"}},
                ],
            },
        ],
    },
}


# --------------------------------------------------------------------------
# the recognition deck
# --------------------------------------------------------------------------

# `kone-recognition-deck`'s arc: the same announcement as the milestone
# slide, paced across four sections so it can be presented rather than
# forwarded. Built here rather than left to the planner because the
# skill is explicit about the order, and because a deterministic arc can
# be tested -- and built at all on a server with no API key.
#
# Two substitutions from the published arc, both because the registry's
# archetype of that name cannot hold the content: AGENDA_B_NUMBERED
# reads title/body/body2 and cannot carry four numbered sections, so the
# agenda goes to `agenda_contents`, which does. Everything else is the
# archetype the skill names.
# The cover and the outro are the deck's retained master slides, not
# ours to emit: `build_deck` keeps Cover F and the Thank you already.
# Emitting our own put two covers and three closing slides in the first
# build of this arc.
RECOGNITION_ARC = [
    ("agenda_contents", "agenda"),
    ("divider_numbering", "section 1"),
    ("text_picture_a", "what changed"),
    ("kone_numbers", "the numbers"),
    ("two_content", "delivered / continuity"),
    ("divider_numbering", "section 2"),
    ("three_content", "the scope groups"),
    ("divider_numbering", "section 3"),
    ("quote_b", "the benchmark line"),
    ("timeline", "what's next"),
    ("divider_numbering", "section 4"),
    ("credits", "thank you"),
]

_SECTIONS = ["What changed", "How it was delivered", "What's next", "Thank you"]


def recognition_deck(content: dict) -> dict:
    """Expand one announcement into the recognition deck's arc.

    Takes the milestone slide's own content plus what the longer form
    allows -- the context paragraph, the quotable line, the timeline --
    and drops any slide whose material is missing rather than padding
    it. The skill is explicit about that: nine slides that all carry
    weight beat twelve with three filler ones.
    """
    footer = content.get("footer") or content.get("eyebrow") or ""
    sections = content.get("sections") or _SECTIONS
    slides: list[dict] = []

    agenda = [{"number": f"{i:02d}", "item": name} for i, name in enumerate(sections, 1)]
    slides.append({"archetype": "agenda_contents", "title": "What we'll cover",
                   "items": agenda})

    def divider(index: int) -> dict:
        return {"archetype": "divider_numbering", "number": f"{index:02d}",
                "eyebrow": f"Section {index:02d}", "title": sections[index - 1]}

    # 01 -- what changed
    slides.append(divider(1))
    if content.get("context"):
        slides.append({"archetype": "text_picture_a", "eyebrow": sections[0],
                       "title": content.get("context_title") or "What moved",
                       "body": content["context"]})
    if content.get("stats"):
        slides.append({"archetype": "kone_numbers", "eyebrow": sections[0],
                       "title": content.get("stats_title") or "The migration in numbers",
                       "stats": content["stats"],
                       "scope_label": content.get("scope_label", ""),
                       "scope": content.get("scope", ""), "footer": footer})
    if content.get("done"):
        slides.append({"archetype": "two_content", "title": "What was delivered",
                       "items": _two_columns(content["done"])})

    # 02 -- how it was delivered
    if content.get("groups"):
        slides.append(divider(2))
        slides.append({"archetype": "three_content", "title": "Who delivered it",
                       "items": content["groups"][:3]})

    # 03 -- what's next
    third = []
    if content.get("quote"):
        third.append({"archetype": "quote_b", "title": "In their words",
                      "quote": f"“{content['quote']}”",
                      "attribution": content.get("quote_attribution", "")})
    if content.get("next"):
        third.append({"archetype": "timeline", "title": sections[2],
                      "items": _as_timeline(content["next"])})
    if third:
        slides.append(divider(3))
        slides.extend(third)

    # 04 -- thank you
    if content.get("credit_names"):
        slides.append(divider(4))
        slides.append({"archetype": "credits", "eyebrow": "Section 04",
                       "title": content.get("credits_title") or "Thank you to everyone involved",
                       "names": [{"name": n} for n in content["credit_names"][:12]],
                       "note": content.get("credits_note", ""), "footer": footer})

    spec = {"title": content.get("title", "Recognition"), "slides": slides}
    # `text_picture_a` carries a picture slot the arc never supplies by
    # hand; unfilled it renders as a white half-slide.
    try:
        from deckguard.registry import fill_empty_photo_slots

        fill_empty_photo_slots(spec)
    except Exception:  # noqa: BLE001 -- no photo library is not an error
        pass
    return spec


def _two_columns(done: list) -> list:
    """The completion states as two bulleted columns rather than ticks --
    the deck has room to let them breathe, and `two_content` reads
    {label, bullets}."""
    items = [d.get("text", d) if isinstance(d, dict) else d for d in done]
    half = (len(items) + 1) // 2
    columns = [items[:half], items[half:]]
    labels = ["Delivered", "Data continuity"]
    return [{"label": label, "bullets": bullets}
            for label, bullets in zip(labels, columns) if bullets]


def _as_timeline(items: list) -> list:
    """`next` is written as sentences; the timeline wants a period and a
    line. Where the sentence names its own period -- a quarter, a month,
    a date -- lift it into the period column."""
    out = []
    for i, item in enumerate(items[:4]):
        text = item.get("text", "") if isinstance(item, dict) else str(item)
        period = _period_in(text) or ("Now" if i == 0 else f"Then {i}")
        out.append({"period": period, "text": text})
    return out


_PERIOD = re.compile(
    r"\b(Q[1-4](?:\s+\d{4})?|H[12]|"
    r"(?:January|February|March|April|May|June|July|August|September|October|November|December)"
    r"(?:\s+\d{4})?)\b", re.I)


def _period_in(text: str) -> Optional[str]:
    found = _PERIOD.search(text)
    return found.group(1) if found else None


# --------------------------------------------------------------------------
# deck shapes, offered in the UI
# --------------------------------------------------------------------------


def _arc_notes() -> str:
    """The recognition arc written out for the planner, built from
    `RECOGNITION_ARC` so the instruction cannot drift from the recipe."""
    steps = "\n".join(f"  {i:02d}. {name} -- {why}"
                      for i, (name, why) in enumerate(RECOGNITION_ARC, 1))
    return (
        "Build the KONE recognition deck: one announcement paced across four "
        "sections so it can be presented rather than forwarded. Use exactly "
        "this arc, in this order, for the BODY slides (the cover and the "
        "closing Thank you are retained from the master -- do not emit "
        "them):\n" + steps + "\n"
        "Sections are: " + ", ".join(_SECTIONS) + ". Every fact from the source "
        "lands in exactly one section. The numbers appear ONCE, on the "
        "kone_numbers slide; if a number repeats on a later slide that slide "
        "is restating rather than adding, so cut it. Drop any slide whose "
        "material is missing rather than padding it -- nine slides that carry "
        "weight beat twelve with three filler ones. Keep every divider on the "
        "same treatment."
    )


# What the tool offers as a starting shape. `notes` and `target` are
# passed straight to the planner; `auto` leaves it free, which is what
# every brief did before this existed.
DECK_SHAPES: dict[str, dict] = {
    "auto": {
        "label": "Let the planner choose",
        "hint": "Picks archetypes to fit the brief. The default, and right for most decks.",
        "notes": None,
        "target": None,
    },
    "recognition": {
        "label": "Recognition deck — a milestone across ~12 slides",
        "hint": "For an announcement that needs a walkthrough on a call or at a town hall: "
                "four sections, a divider each, the numbers once, credits in a ruled grid.",
        "notes": _arc_notes,
        "target": 12,
    },
    "milestone": {
        "label": "Milestone — one shareable slide",
        "hint": "For an announcement that gets posted and read rather than presented: "
                "the claim, its proof numbers, what's next and who did it, on one slide.",
        "notes": "Build exactly ONE body slide, using the `milestone_slide` archetype. "
                 "Do not add any other slide. Follow that archetype's rules exactly.",
        "target": 1,
    },
}


def shape_notes(shape: str) -> tuple[Optional[str], Optional[int]]:
    """Planner guidance and slide target for a chosen shape, or
    (None, None) for anything unrecognised -- an unknown shape must
    leave the brief exactly as it would have been."""
    entry = DECK_SHAPES.get(shape or "auto")
    if not entry:
        return None, None
    notes = entry["notes"]
    return (notes() if callable(notes) else notes), entry["target"]


# --------------------------------------------------------------------------
# chrome, owned by the layout
# --------------------------------------------------------------------------


def strip_master_classification(prs) -> int:
    """Take the master's own `KONE Internal` stamp off every slide.

    It is a plain Arial 8pt text box at x:1204 y:7 sitting on the slide
    MASTER, so it is inherited by all fifty layouts and appears on every
    slide of every deck -- including customer-facing ones, which is the
    part that matters: every external render carried `KONE Internal` in
    its top-right corner.

    It is off-brand three ways over. Arial is not an approved face;
    BRAND_MODE §3 puts classification bottom-left at `45,640` in KONE
    Information; and an external deck carries no classification at all.
    `stamp_chrome` already draws it correctly when a deck declares one,
    so this only has to remove the thing that was never chrome.

    Returns how many shapes were removed.
    """
    removed = 0
    for master in prs.slide_masters:
        for shape in list(master.shapes):
            if not getattr(shape, "has_text_frame", False) or shape.is_placeholder:
                continue
            text = shape.text_frame.text.strip().lower()
            if text in ("kone internal", "kone confidential", "internal"):
                shape._element.getparent().remove(shape._element)
                removed += 1
    return removed


def stamp_chrome(slide, engine, archetype: str, *, page: int, date: str,
                 classification: str = "") -> None:
    """Date, page number and classification, placed by the layout.

    BRAND_MODE section 3 says chrome belongs to the layout and an
    archetype should draw none of it. Every archetype duly drew none --
    and so did the layout, so every generated body slide came out with
    nothing at all below y=629. The master's own retained cover and
    Thank you carried theirs, which is why it read as a formatting
    quirk rather than as chrome being absent.

    Covers, dividers, the outro, a full-bleed picture and a blank take
    none of this by design.
    """
    from deckguard import brandmode as bm

    if not bm.wants_footer(archetype):
        return

    dark = _slide_is_dark(slide, engine)
    ink = bm.WHITE if dark else bm.BLACK
    if date:
        _chrome_text(slide, engine, [bm.FOOTER_DATE_X, bm.FOOTER_Y, 500, 16],
                     date, "footer", ink)
    if page:
        _chrome_text(slide, engine, [bm.FOOTER_PAGE_X, bm.FOOTER_Y, 68, 16],
                     f"{page:02d}", "footer", ink)
    if classification:
        _chrome_text(slide, engine, [bm.FOOTER_DATE_X, bm.CLASSIFICATION_Y, 500, 14],
                     classification, "classification", ink)


def _chrome_text(slide, engine, box, value: str, role: str, ink: str) -> None:
    from deckguard import brandmode as bm

    before = len(slide.shapes._spTree)   # shapes are not hashable
    style = bm.resolve(role) or {}
    _draw_text(slide, engine, box, value, {
        "kind": "text", "px": style.get("px", 11),
        "font": style.get("font", bm.KONE_INFO),
        "color": ink, "caps": True, "align": "l",
    })
    # Named so preflight can hold content to the real floor (629) while
    # letting chrome sit at 658 where it belongs. Without the name the
    # check had to be loosened to 680, which let a region overlapping
    # the footer by three pixels through.
    for element in list(slide.shapes._spTree)[before:]:
        shape = next((s for s in slide.shapes if s._element is element), None)
        if shape is not None:
            shape.name = f"Chrome {role}"


def _slide_is_dark(slide, engine) -> bool:
    """Is the bottom-left of this slide a dark ground?

    The handoff is explicit that chrome colour is set against the
    image, not against the archetype -- external slide 12 keeps a black
    page number because the photograph is light in that corner while
    slide 14's is dark. Approximated here by the largest shape covering
    the footer band.
    """
    from pptx.util import Emu

    px = None
    try:
        px = slide.part.package.presentation_part.presentation.slide_width / 1280
    except Exception:  # noqa: BLE001
        return False
    band_top = 620 * px
    best, best_area = None, 0
    for shape in slide.shapes:
        try:
            top, height = shape.top, shape.height
            left, width = shape.left, shape.width
        except Exception:  # noqa: BLE001
            continue
        if None in (top, height, left, width):
            continue
        if top + height < band_top or left > 700 * px:
            continue
        area = width * height
        if area > best_area:
            best, best_area = shape, area
    if best is None:
        return False
    try:
        if best.shape_type is not None and best.shape_type.name == "PICTURE":
            return True   # a photograph under the footer: white reads on both
        fill = best.fill
        if fill.type is not None and str(fill.fore_color.rgb) in ("1450F5", "141414"):
            return True
    except Exception:  # noqa: BLE001
        return False
    return False


# --------------------------------------------------------------------------
# the cut cover
# --------------------------------------------------------------------------

# One photograph behind background-coloured masks -- never a photo
# pre-sliced into panes, so dropping in a new image reproduces the
# chopped effect and PowerPoint's Change Picture still works.
#
# COVER_B_CUT3's panes are published in INTERNAL_25.md slide 01.
# COVER_A_CUT4's are not published anywhere in the handoff: it is
# described only as "four photo panes cut across the top". The four
# below keep the internal cover's rhythm -- top-anchored, staggered
# heights, 10px gutters -- across the content column, and are the one
# piece of this geometry that is derived rather than measured.
_CUT_COVERS: dict[str, dict] = {
    "cover_b_cut3": {
        "band": [267, 0, 968, 430],
        "panes": [[267, 300, 340], [577, 300, 430], [887, 348, 380]],
    },
    # No longer derived. Measured off "Life, upgraded in ONE week", where
    # the cut is baked into a transparent PNG: reading its alpha channel
    # back gives four panes 289px wide on a 330px pitch, 41px gutters,
    # FULL BLEED from x=0 to x=1280 rather than inside the content
    # column, all four flush to the top edge with staggered depths.
    #
    # The stagger is the thing. Four equal panes read as a filmstrip;
    # 249/322/421/372 reads as a cut.
    "cover_a_cut4": {
        "band": [0, 0, 1280, 422],
        "panes": [[0, 289, 249], [330, 289, 322], [660, 289, 421], [990, 290, 372]],
    },
}


def _cover_band(slide, engine, band) -> bool:
    """Grow the banner photograph so it spans the whole cut band."""
    bx, by, bw, bh = band
    left, top = engine.X(bx), engine.X(by)
    right, bottom = engine.X(bx + bw), engine.X(by + bh)
    banner, area = None, 0
    for shape in slide.shapes:
        kind = getattr(shape, "shape_type", None)
        if kind is None or kind.name != "PICTURE":
            continue
        size = (shape.width or 0) * (shape.height or 0)
        if size > area:
            banner, area = shape, size
    if banner is None:
        return False
    # Set to the band, not unioned with it. A union left the banner
    # wider than the cut on COVER_A_CUT4, so photograph bled past the
    # leftmost pane into the margin with no mask over it.
    banner.left, banner.top = left, top
    banner.width, banner.height = right - left, bottom - top
    return True


def _bg_hex(bg) -> str:
    """A background name or hex to the hex the masks paint in.

    A mask has to match the field exactly or the cut shows as a seam.
    """
    from deckguard import brandmode as bm

    if not bg:
        return bm.WHITE
    text = str(bg).strip().lstrip("#")
    if len(text) == 6 and all(c in "0123456789abcdefABCDEF" for c in text):
        return text.upper()
    return {
        "white": bm.WHITE, "sand": bm.SAND, "light_blue": bm.LIGHT_BLUE,
        "lightblue": bm.LIGHT_BLUE, "pink": bm.PINK, "mint": bm.MINT,
        "yellow": bm.YELLOW, "blue": bm.BLUE, "black": bm.BLACK,
    }.get(text.lower(), bm.WHITE)


def _cut_masks(band, panes) -> list:
    """The rectangles that hide everything the panes do not show.

    Three kinds: the gutters between panes, the strip under a pane
    shorter than the band, and any run of band before the first pane or
    after the last. Computed rather than listed so a pane change cannot
    leave a stale mask behind.
    """
    bx, by, bw, bh = band
    masks, cursor = [], bx
    for x, w, h in sorted(panes):
        if x > cursor:
            masks.append([cursor, by, x - cursor, bh])
        if h < bh:
            masks.append([x, by + h, w, bh - h])
        cursor = x + w
    if cursor < bx + bw:
        masks.append([cursor, by, bx + bw - cursor, bh])
    return masks


def draw_cut_cover(slide, engine, name: str, background: str = "FFFFFF") -> int:
    """Mask a cut cover's banner into panes. Returns how many it drew.

    The cut covers rendered their photograph full-width with no masks
    at all, so the chopped effect the brand asks for simply never
    appeared. The test that would have caught it was skipped rather
    than passing.
    """
    from pptx.enum.shapes import MSO_SHAPE

    cut = _CUT_COVERS.get((name or "").lower())
    if not cut:
        return 0
    # The banner has to cover the band it is cut from, or a pane shows
    # background instead of photograph. cover_b_cut3's picture starts at
    # x:330 while its published panes start at 267, so the leftmost pane
    # would have come out empty.
    _cover_band(slide, engine, cut["band"])

    drawn = 0
    for x, y, w, h in _cut_masks(cut["band"], cut["panes"]):
        if w <= 0 or h <= 0:
            continue
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, engine.X(x), engine.X(y), engine.X(w), engine.X(h)
        )
        shape.name = "Cut mask"
        shape.fill.solid()
        shape.fill.fore_color.rgb = engine._hex(background)
        shape.line.fill.background()
        shape.shadow.inherit = False
        drawn += 1
    _cover_marks(slide, engine, cut)
    return drawn


_SHADOW_XML = (
    '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:outerShdw blurRad="{blur}" dist="{dist}" dir="{dir}" rotWithShape="0">'
    '<a:srgbClr val="000000"><a:alpha val="{alpha}"/></a:srgbClr>'
    "</a:outerShdw></a:effectLst>"
)


def _soft_shadow(shape) -> None:
    """The card shadow, which python-pptx cannot add on its own.

    `shape.shadow` can only turn inheritance off; there is no API for
    setting one. Written as raw DrawingML because the alternative is a
    flat card, and the shadow is half of what makes the reference's grid
    read as cards rather than as a table.
    """
    from deckguard import brandmode as bm

    spec = bm.CARD_SHADOW
    shape.shadow.inherit = False
    spPr = shape._element.spPr
    for existing in spPr.findall(f"{{{_A}}}effectLst"):
        spPr.remove(existing)
    # 9525 EMU per px at the 96dpi the whole geometry is written in.
    xml = _SHADOW_XML.format(
        blur=int(spec["blur"] * 9525),
        dist=int(spec["distance"] * 9525),
        dir=int(spec["direction"] * 60000),
        alpha=int(spec["alpha"] * 100000),
    )
    from lxml import etree

    spPr.append(etree.fromstring(xml))


_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _draw_cards(slide, engine, cards: list, content: dict) -> int:
    """Draw a grid of reference-style cards. Returns how many landed."""
    from pptx.enum.shapes import MSO_SHAPE

    from deckguard import brandmode as bm

    items = content.get("cards") or []
    drawn = 0
    for spot, item in zip(cards, items):
        if not isinstance(item, dict):
            continue
        x, y, w, h = spot["box"]
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, engine.X(x), engine.X(y),
            engine.X(w), engine.X(h))
        card.name = "Card"
        card.fill.solid()
        card.fill.fore_color.rgb = engine._hex(bm.CARD_FILL)
        card.line.fill.background()
        try:
            card.adjustments[0] = bm.CARD_RADIUS_PX / min(w, h)
        except (IndexError, ValueError):
            pass
        _soft_shadow(card)

        rx, ry, rw, rh = spot["rule"]
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, engine.X(rx), engine.X(ry),
            engine.X(rw), engine.X(max(rh, 1)))
        rule.name = "Card rule"
        rule.fill.solid()
        rule.fill.fore_color.rgb = engine._hex(spot["accent"])
        rule.line.fill.background()
        rule.shadow.inherit = False

        _card_arrow(slide, engine, spot)
        if item.get("label"):
            _draw_text(slide, engine, spot["label"], item["label"],
                       {"kind": "text", "px": 12, "font": bm.KONE_INFO,
                        "color": bm.BLUE, "caps": True, "align": "l"})
        if item.get("text"):
            _draw_text(slide, engine, spot["body"], item["text"],
                       {"kind": "text", "px": 15, "font": bm.INTER,
                        "color": bm.BLACK, "caps": False, "align": "l"})
        drawn += 1
    return drawn


def _card_arrow(slide, engine, spot: dict) -> bool:
    """The small arrow the reference puts at the head of every card.

    It recurs across the deck -- a down-right arrow in the card headers,
    a right arrow before a highlighted line -- and it is the one motif
    that makes a run of these slides read as one deck. Drawn in the
    card's own accent colour.
    """
    from deckguard import brandmode as bm

    box = spot.get("arrow")
    if not box:
        return False
    # The glyph, not a pictogram: at 17px a sprite icon is mush, and the
    # reference's own mark is a plain arrow.
    #
    # Blue, not the card's accent. Colouring it mint or pale pink is
    # exactly the mistake the labels were moved off -- it is TYPE, and
    # BRAND_MODE allows three inks. Preflight caught it here, which is
    # the check doing its job on its own author.
    _draw_text(slide, engine, box, "\u2198", {
        "kind": "text", "px": 15, "font": bm.INTER,
        "color": bm.BLUE, "caps": False, "align": "l",
    })
    return True


def _card(x, y, w, h, accent: str) -> dict:
    """One card of the grid, as the reference draws it.

    White rounded rectangle on sand with a soft shadow, a caps label
    inside the top, a full-width rule under the label in the card's
    accent colour, and body copy beneath.

    The label is NOT set in the accent. The reference does that -- mint
    on white at 12px, pale blue on white -- and those two cards are the
    only ones on the slide you cannot read. The rule carries the colour
    coding; the label stays blue.
    """
    from deckguard import brandmode as bm

    return {
        "box": [x, y, w, h],
        "accent": accent,
        "label": [x + 40, y + bm.CARD_LABEL_Y, w - 54, 22],
        "rule": [x, y + bm.CARD_RULE_Y, w, bm.CARD_RULE_H],
        "body": [x + 14, y + bm.CARD_BODY_Y, w - 28, h - bm.CARD_BODY_Y - 14],
        "arrow": [x + 14, y + bm.CARD_LABEL_Y - 1, 18, 17],
    }


def card_grid(cells: int = 8) -> dict:
    """The reference's twelve-cell grid, as an archetype spec.

    Four columns on a 300px pitch, three rows on a 193px pitch, 288x176
    cards. Measured off "Recognize the need and importance of
    modernizing elevators", which is the densest slide in the reference
    deck and the one it holds together best.
    """
    from deckguard import brandmode as bm

    columns = len(bm.CARD_COL_X)
    rows = max(1, -(-cells // columns))
    # Centre the rows actually used in the content band rather than
    # hanging them off the top. Eight cards in a three-row grid left the
    # bottom third of the slide empty and the block read as unfinished.
    gap = bm.CARD_ROW_Y[1] - bm.CARD_ROW_Y[0] - bm.CARD_H
    band = bm.FLOOR - bm.TIGHT_CONTENT_Y
    # The reference runs its third row of cards to y=681, past the floor.
    # A twelve-cell grid therefore shortens its cards to fit rather than
    # hanging them over the footer -- the floor is the one rule with no
    # exceptions, and 159px still holds three lines of 15px body.
    height = min(bm.CARD_H, (band - (rows - 1) * gap) // rows)
    used = rows * height + (rows - 1) * gap
    top = bm.TIGHT_CONTENT_Y + max(0, (band - used) // 2)

    cards = []
    for index in range(cells):
        x = bm.CARD_COL_X[index % columns]
        y = top + (index // columns) * (height + gap)
        cards.append(_card(x, y, bm.CARD_W, height,
                           bm.CARD_ACCENTS[index % len(bm.CARD_ACCENTS)]))
    return {"cards": cards}


DEFAULT_COVER = "cover_a_cut4"


def _with_cut_cover(spec: dict) -> list:
    """Put the four-pane cut cover at the front of every deck.

    What the master retains as slide 1 is a full-bleed photograph with a
    white title on it, and on a sunlit frame the title all but vanishes
    -- the tool's own decks were opening with a single legible letter.
    The reference deck opens on the cut instead, and that is the cover
    the brand is known for, so it is the default rather than something
    to be picked.

    A deck that already opens on a cover keeps the one it has.
    """
    slides = [dict(s) for s in (spec.get("slides") or [])]
    first = str((slides[0] if slides else {}).get("archetype") or "").lower()
    if first.startswith("cover"):
        return slides
    cover = {"archetype": DEFAULT_COVER, "title": spec.get("title") or ""}
    context = spec.get("context") or spec.get("subtitle")
    if context:
        cover["context"] = context
    # Photo slots are filled in `assemble.build`, which runs before this
    # cover exists, so it fills its own or opens on a sand rectangle.
    try:
        from deckguard.registry import fill_empty_photo_slots

        fill_empty_photo_slots({"slides": [cover]})
    except Exception:  # noqa: BLE001 -- no library is not a build failure
        pass
    return [cover] + slides


def _cover_marks(slide, engine, cut: dict) -> None:
    """The logo and the tagline, drawn on the cut cover itself.

    Both live on the master's Cover A layout and neither reaches the
    page: the logo is an EMPTY picture placeholder -- the same defect as
    the 54 blank logo frames -- and the tagline, though it carries a real
    black PNG at 1102,633, does not render through layout inheritance.
    The reference cover has both, and a cover with no KONE mark on it is
    not a KONE cover, so they are placed here rather than inherited.

    White logo, because it sits at 45,45 on the first pane and the first
    pane is a photograph.
    """
    from deckguard.logo import _brand_asset

    if any((sh.name or "") in ("Logo", "Tagline") for sh in slide.shapes):
        return
    logo = _brand_asset("logo", light=True)
    if logo:
        mark = slide.shapes.add_picture(
            logo, engine.X(45), engine.X(45), engine.X(81), engine.X(31))
        mark.name = "Logo"
    tagline = _brand_asset("tagline", light=False)
    if tagline:
        # Bottom-right, clear of the cut band and of the footer line.
        mark = slide.shapes.add_picture(
            tagline, engine.X(1102), engine.X(633), engine.X(133), engine.X(45))
        mark.name = "Tagline"
