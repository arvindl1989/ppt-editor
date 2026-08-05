"""Deterministic fix engine.

Runs in two stages:

1. Deck-wide: theme color remap, theme font remap, and explicit font
   overrides on master/layout placeholders. These correct every
   downstream shape that inherits from them in one operation.
2. Per-violation, to a fixpoint: re-audits the (now theme-corrected) deck
   and applies every `auto_fixable` violation directly via the live
   object each `Violation.target` record points at, then repeats until a
   pass makes no further changes. Looping matters because fixing one
   thing on a run (e.g. its font) can newly unlock another check on that
   same run (e.g. that font's allowed text colors) — a single pass would
   otherwise leave mechanically-resolvable cases in "manual review".

Never touches the input file — callers pass a `Presentation` already
loaded from a copy, and `--dry-run` simply skips the final `prs.save()`.
"""

from __future__ import annotations

from collections import Counter
from dataclasses import dataclass, field
from functools import lru_cache
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deckguard import colors as colors_mod
from deckguard import effects as effects_mod
from deckguard import fonts as fonts_mod
from deckguard import logo as logo_mod
from deckguard.fonts import FontTables, normalize_key, remap_theme_fonts
from deckguard.inventory import ALIGN_BY_NAME, _has_explicit_run_color, build_inventory
from deckguard.rules_engine import Violation, audit_deck, sort_violations, summarize

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _p(tag: str) -> str:
    return f"{{{P_NS}}}{tag}"


@dataclass
class Change:
    scope: str  # "theme" | "master" | "layout" | "slide"
    rule: str
    field: str
    old: object
    new: object
    slide_index: int = 0
    shape_id: Optional[int] = None
    shape_name: Optional[str] = None
    location: Optional[str] = None

    def to_dict(self) -> dict:
        return {
            "scope": self.scope,
            "rule": self.rule,
            "field": self.field,
            "old": self.old,
            "new": self.new,
            "slide_index": self.slide_index,
            "shape_id": self.shape_id,
            "shape_name": self.shape_name,
            "location": self.location,
        }


@dataclass
class FixReport:
    source_path: str
    output_path: Optional[str]
    dry_run: bool
    changes: list = field(default_factory=list)
    manual_review: list = field(default_factory=list)  # list[Violation], auto_fixable=False

    @property
    def summary(self) -> dict:
        s = summarize(self.manual_review) if self.manual_review else {"critical": 0, "major": 0, "minor": 0}
        return {
            "changes_applied": len(self.changes),
            "manual_review_required": len(self.manual_review),
            "manual_review_by_severity": {
                k: s.get(k, 0) for k in ("critical", "major", "minor")
            },
        }


def _remap_explicit_fonts_in_masters_and_layouts(prs, font_remap: dict[str, str]) -> list[Change]:
    """Fix explicit (non-inherited) font overrides on master/layout placeholders.

    Theme fontScheme remap (see fonts.remap_theme_fonts) already fixes
    everything that inherits the default; this catches shapes on a
    master/layout that set a legacy font name directly.
    """
    changes: list[Change] = []
    remap_by_key = {normalize_key(k): v for k, v in font_remap.items()}
    if not remap_by_key:
        return changes

    seen_masters = set()
    for master in colors_mod.iter_slide_masters(prs):
        if id(master.part) not in seen_masters:
            seen_masters.add(id(master.part))
            changes += _remap_shapes_fonts(master.shapes, remap_by_key, "master", master.name)
        for layout in master.slide_layouts:
            changes += _remap_shapes_fonts(layout.shapes, remap_by_key, "layout", layout.name)
    return changes


def _remap_large_panel_fills_in_masters_and_layouts(
    prs, layout_panel_remap: dict[str, str], min_area_sq_in: float
) -> list[Change]:
    """Fix explicit solid fills on large background panels defined on a
    slide layout/master rather than the slide itself.

    Slide-level color remap only ever sees shapes actually placed on a
    slide. A shape's fill can instead live entirely on the layout it
    inherits from -- e.g. a full-panel background rectangle behind a
    picture/content placeholder -- and would otherwise never be touched.
    Scoped to `layout_panel_remap` (a separate table from `colors.remap`)
    and to shapes at least `min_area_sq_in` in area, because the same
    literal color is often reused on a layout for small accent/placeholder
    rectangles (e.g. unselected tab backgrounds) whose color is
    intentionally static regardless of brand and must not be touched.
    """
    changes: list[Change] = []
    if not layout_panel_remap:
        return changes
    remap = {colors_mod.normalize_hex(k): v for k, v in layout_panel_remap.items()}
    min_area_emu2 = min_area_sq_in * 914400 * 914400

    seen_masters = set()
    for master in colors_mod.iter_slide_masters(prs):
        if id(master.part) not in seen_masters:
            seen_masters.add(id(master.part))
            changes += _remap_shapes_fills(master.shapes, remap, min_area_emu2, "master", master.name)
        for layout in master.slide_layouts:
            changes += _remap_shapes_fills(layout.shapes, remap, min_area_emu2, "layout", layout.name)
    return changes


def _replace_old_logo_everywhere(
    prs, old_hashes: list[str], new_logo_path: Optional[str], threshold: int
) -> list[Change]:
    """Fix old-logo instances the per-slide/per-shape audit can never see:

    - A logo placed directly on a slide MASTER or LAYOUT rather than any
      individual slide -- a common, in fact standard, way to put a logo
      on every slide at once by inheritance. Ordinary slide-level
      scanning (`inventory.slides[].shapes`) never visits master/layout
      shapes, so a logo living only there is invisible to it.
    - A logo baked into a page-level background-FILL image
      (`<p:cSld><p:bg>`) on a slide, layout, or master, rather than
      being a picture *shape* at all -- outside the shape tree entirely,
      so even a full recursive shape scan can't see it either. This is
      "hard-coded" in the sense that it isn't an element deckguard's
      shape-based model can address.
    """
    changes: list[Change] = []
    if not old_hashes or not new_logo_path or not Path(new_logo_path).exists():
        return changes

    def _pic_changes(shapes, scope: str, location: str) -> None:
        for match in logo_mod.find_old_logo_matches(shapes, old_hashes, threshold):
            logo_mod.replace_logo_image(match.shape, new_logo_path)
            changes.append(
                Change(scope=scope, rule="old_logo", field="image",
                       old=match.matched_hash, new=new_logo_path, location=location)
            )

    def _bg_change(part, container_element, scope: str, location: str, slide_index: int = 0) -> None:
        match = logo_mod.find_old_logo_background_match(part, container_element, old_hashes, threshold)
        if match is None:
            return
        logo_mod.replace_background_image(part, match.shape, new_logo_path)
        changes.append(
            Change(scope=scope, rule="old_logo", field="background image",
                   old=match.matched_hash, new=new_logo_path, location=location, slide_index=slide_index)
        )

    seen_masters = set()
    for master in colors_mod.iter_slide_masters(prs):
        if id(master.part) not in seen_masters:
            seen_masters.add(id(master.part))
            _pic_changes(master.shapes, "master", master.name)
            _bg_change(master.part, master._element, "master", f"{master.name} [background]")
        for layout in master.slide_layouts:
            _pic_changes(layout.shapes, "layout", layout.name)
            _bg_change(layout.part, layout._element, "layout", f"{layout.name} [background]")

    for i, slide in enumerate(prs.slides, start=1):
        _bg_change(slide.part, slide._element, "slide", "(slide background)", slide_index=i)

    return changes


def _replace_old_logo_region_everywhere(prs, region_in, new_logo_path: Optional[str]) -> list[Change]:
    """Fix an old logo `old_logo_hashes` can never match at all: one not
    built from a raster image, so there's no picture to perceptual-hash
    against -- e.g. a wordmark drawn as vector freeform shapes directly
    on a slide master (confirmed against a real legacy deck: a `<p:grpSp>`
    of freeform paths, not a `<p:pic>`, sitting in the master's own top
    corner). `region_in` (`[left, top, width, height]` in inches, from
    `logo.old_logo_region_in`) identifies WHERE the mark lives instead of
    WHAT it looks like -- every top-level shape on a slide MASTER fully
    inside that region is removed and replaced with the current logo
    image, sized to fit. Opt-in and unset by default, same as
    `old_logo_hashes`: this deletes shapes based on position alone, which
    is only safe once a human has confirmed the region against their own
    deck (see README for how).
    """
    changes: list[Change] = []
    if not region_in or not new_logo_path or not Path(new_logo_path).exists():
        return changes
    if len(region_in) != 4:
        return changes

    from pptx.util import Inches

    from deckguard.slide_import import default_template_path

    region_emu = tuple(Inches(v) for v in region_in)
    # The org template's own actual logo size/position -- NOT the (deliberately
    # generous) search region -- is what the replacement should be sized to.
    # See reference_logo_geometry's own docstring for why.
    target_emu = logo_mod.reference_logo_geometry(default_template_path())

    seen_masters = set()
    for master in colors_mod.iter_slide_masters(prs):
        if id(master.part) in seen_masters:
            continue
        seen_masters.add(id(master.part))
        matches = logo_mod.find_shapes_in_region(master.shapes, region_emu)
        if not matches:
            continue
        shape_names = [s.name for s in matches]
        logo_mod.replace_shapes_in_region_with_logo(master.shapes, matches, new_logo_path, region_emu, target_emu)
        changes.append(
            Change(
                scope="master", rule="old_logo_region", field="image",
                old=f"{len(matches)} shape(s): {', '.join(shape_names)}", new=new_logo_path, location=master.name,
            )
        )

    return changes


@lru_cache(maxsize=32)
def _reference_placeholder_geometry(template_path, ph_type_name: str) -> Optional[tuple]:
    """Majority (left, top, width, height) EMU for every layout's own
    `ph_type_name`-type placeholder (e.g. "DATE", "SLIDE_NUMBER") in the
    org template -- same "the current template defines the ground
    truth" principle as `logo.reference_logo_geometry`, generalized to
    any placeholder type. Returns None if the template can't be read or
    has no such placeholder at all. Cached for the same reason
    `reference_logo_geometry` is -- see its docstring.
    """
    try:
        prs = Presentation(str(template_path))
    except Exception:  # noqa: BLE001 -- missing/corrupt template, not fatal to the caller
        return None

    sizes: Counter = Counter()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type is not None and ph_type.name == ph_type_name and None not in (ph.left, ph.top, ph.width, ph.height):
                    sizes[(ph.left, ph.top, ph.width, ph.height)] += 1
    if not sizes:
        return None
    return sizes.most_common(1)[0][0]


def _normalize_footer_chrome_position(prs) -> list[Change]:
    """Reposition DATE/SLIDE_NUMBER placeholders on every master/layout
    in the deck onto the org template's own position -- date on the
    left, slide number on the right, confirmed by direct inspection of
    the bundled template. Real, confirmed defect this fixes: an old
    deck's own layout had these swapped (slide number at the LEFT edge,
    its confidentiality/footer text occupying the RIGHT). Only ever
    repositions -- never touches text/content -- and skips a
    placeholder already at the reference position, so an already-
    correct deck is left byte-for-byte untouched on this axis.
    """
    from deckguard.slide_import import default_template_path

    changes: list[Change] = []
    reference = {}
    for ph_type_name in ("DATE", "SLIDE_NUMBER"):
        geom = _reference_placeholder_geometry(default_template_path(), ph_type_name)
        if geom:
            reference[ph_type_name] = geom
    if not reference:
        return changes

    containers = []
    seen_masters = set()
    for master in colors_mod.iter_slide_masters(prs):
        if id(master.part) not in seen_masters:
            seen_masters.add(id(master.part))
            containers.append(("master", master.name, master.shapes))
        for layout in master.slide_layouts:
            containers.append(("layout", layout.name, layout.shapes))

    for scope, location, shapes in containers:
        for shp in shapes:
            if not getattr(shp, "is_placeholder", False):
                continue
            ph_type = shp.placeholder_format.type
            if ph_type is None or ph_type.name not in reference:
                continue
            if None in (shp.left, shp.top, shp.width, shp.height):
                continue
            target = reference[ph_type.name]
            if (shp.left, shp.top, shp.width, shp.height) == target:
                continue
            shp.left, shp.top, shp.width, shp.height = target
            changes.append(
                Change(
                    scope=scope, rule="footer_chrome_position", field=ph_type.name.lower(),
                    old=None, new=None, location=location,
                )
            )
    return changes


CONFIDENTIALITY_FOOTER_MARKER = "confidential"


def _remove_confidentiality_footer_text(prs) -> list[Change]:
    """Clear a FOOTER placeholder's text on every slide where it reads
    as a confidentiality/copyright boilerplate line (matched
    case-insensitively on the word "confidential", e.g. "Confidential |
    (c) KONE Corporation") -- removed outright on explicit direction,
    not just recolored/refonted like ordinary footer chrome. Only ever
    touches a FOOTER placeholder specifically, and only when it
    contains that marker word -- never guesses at any other text.
    """
    changes: list[Change] = []
    for i, slide in enumerate(prs.slides, start=1):
        for shp in slide.placeholders:
            ph_type = shp.placeholder_format.type
            if ph_type is None or ph_type.name != "FOOTER" or not shp.has_text_frame:
                continue
            text = shp.text_frame.text
            if CONFIDENTIALITY_FOOTER_MARKER not in text.lower():
                continue
            shp.text_frame.clear()
            changes.append(
                Change(
                    scope="slide", rule="confidentiality_footer_removed", field="text",
                    old=text.strip(), new="", slide_index=i, shape_id=shp.shape_id, shape_name=shp.name,
                )
            )
    return changes


def _dedupe_shape_ids_in_part(root_element) -> int:
    """Renumber any `<p:cNvPr id=...>` value used more than once within
    this single part's XML (a slide, a layout, or a master), keeping the
    first occurrence and reassigning every later duplicate to the next
    free id in that same part. Shape ids only need to be unique WITHIN a
    part, not across the whole file, so this is scoped per-part on
    purpose.

    Exists to heal a real, if rare, defect found in a genuinely old deck:
    an embedded object duplicated at some point in the deck's edit
    history that kept the SAME id on both copies. PowerPoint's own
    repair-on-open logic is often silently tolerant of this in a file it
    opens directly -- but this project's own output shouldn't propagate
    a source file's pre-existing XML defect unchanged, and a file
    re-saved by python-pptx doesn't reliably get the same silent
    auto-heal treatment (confirmed: PowerPoint refused to open a
    redesign/rebrand output built from a deck with this exact defect).
    Returns how many ids were renumbered.
    """
    cNvPr_els = root_element.findall(f".//{_p('cNvPr')}")
    seen: set[int] = set()
    next_id = 1
    renumbered = 0
    for el in cNvPr_els:
        raw = el.get("id")
        if raw is None or not raw.isdigit():
            continue
        id_ = int(raw)
        if id_ not in seen:
            seen.add(id_)
            continue
        while next_id in seen:
            next_id += 1
        el.set("id", str(next_id))
        seen.add(next_id)
        renumbered += 1
        next_id += 1
    return renumbered


def _dedupe_shape_ids(prs) -> list[Change]:
    changes: list[Change] = []
    seen_masters = set()
    for master in colors_mod.iter_slide_masters(prs):
        if id(master.part) not in seen_masters:
            seen_masters.add(id(master.part))
            n = _dedupe_shape_ids_in_part(master._element)
            if n:
                changes.append(Change(scope="master", rule="duplicate_shape_id", field="id", old=None, new=None, location=f"{master.name} ({n} renumbered)"))
        for layout in master.slide_layouts:
            n = _dedupe_shape_ids_in_part(layout._element)
            if n:
                changes.append(Change(scope="layout", rule="duplicate_shape_id", field="id", old=None, new=None, location=f"{layout.name} ({n} renumbered)"))
    for i, slide in enumerate(prs.slides, start=1):
        n = _dedupe_shape_ids_in_part(slide._element)
        if n:
            changes.append(Change(scope="slide", rule="duplicate_shape_id", field="id", old=None, new=None, slide_index=i, location=f"({n} renumbered)"))
    return changes


FOOTER_CHROME_PLACEHOLDER_TYPES = {"DATE", "FOOTER", "SLD_NUM", "SLIDE_NUMBER"}


def _normalize_footer_chrome_text(prs, config: dict) -> list[Change]:
    """Force date/footer/slide-number placeholder text onto brand values
    when it has no explicit font/color of its own.

    These three placeholder types are structurally different from
    ordinary body/title text: PowerPoint auto-fields (date, slide
    number) and short literal footer text almost never carry an
    explicit run color, and almost never sit on a fill the `contrast`
    rule can compute against (they sit directly on the page canvas) --
    so `inventory.py`'s color resolver (deliberately conservative: an
    unresolvable inherited color is recorded as `ColorInfo(kind="none")`
    rather than guessed at) can never flag them as a violation in the
    first place, and they silently keep whatever the OLD deck's own
    master defined. Confirmed on a real reported deck: Arial, no color
    override at all -- exactly what a brand-compliance pass exists to
    catch, just invisible to the per-violation audit/fix loop below.

    Rather than build a general inherited-color resolver just for this,
    force these three placeholder types specifically onto the current
    org template's own default for them (its master's own
    `<p:otherStyle>`, which in the bundled template resolves to Inter /
    `#141414` -- matching `typography_rules.text_colors`' own "black
    first" fallback rule for text with no resolvable background) --
    but only for a run that has no explicit font/color of its own. A
    run that already sets one is a deliberate choice by whoever built
    the deck, not this project's to override.
    """
    fonts_cfg = config.get("fonts", {}) or {}
    typo = config.get("typography_rules", {}) or {}
    font_name = (fonts_cfg.get("approved") or ["Inter"])[0]
    color_hex = None
    for candidate in typo.get("text_colors", {}).get(font_name, []) or []:
        color_hex = colors_mod.normalize_hex(candidate)
        break
    if color_hex is None:
        color_hex = "141414"

    changes: list[Change] = []
    for i, slide in enumerate(prs.slides, start=1):
        for shp in slide.placeholders:
            ph_type = shp.placeholder_format.type
            if ph_type is None or ph_type.name not in FOOTER_CHROME_PLACEHOLDER_TYPES or not shp.has_text_frame:
                continue
            changed = False
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if not run.font.name:
                        run.font.name = font_name
                        changed = True
                    if not _has_explicit_run_color(run):
                        run.font.color.rgb = RGBColor.from_string(color_hex)
                        changed = True
            # Date and slide-number placeholders are PowerPoint auto-fields
            # (<a:fld>), not <a:r> runs at all -- python-pptx's own
            # paragraph.runs never returns them, so they need their own
            # raw-XML pass, via the exact same "only touch what's unset" rule.
            for fld in shp.text_frame._txBody.iter(effects_mod.a_qn("fld")):
                rPr = fld.find(effects_mod.a_qn("rPr"))
                if rPr is None:
                    rPr = etree.SubElement(fld, effects_mod.a_qn("rPr"))
                    fld.insert(0, rPr)  # rPr must precede pPr/t per schema order
                if _set_fld_rPr_defaults(rPr, font_name, color_hex):
                    changed = True
            if changed:
                changes.append(
                    Change(
                        scope="slide", rule="footer_chrome_default", field="font/color",
                        old=None, new=f"{font_name} / #{color_hex}", slide_index=i,
                        shape_id=shp.shape_id, shape_name=shp.name,
                    )
                )
    return changes


def _set_fld_rPr_defaults(rPr, font_name: str, color_hex: str) -> bool:
    """Add explicit color/font to an `<a:fld>`'s `<a:rPr>` if it doesn't
    already have one of its own -- same "only touch what's unset" rule
    `_normalize_footer_chrome_text` applies to ordinary runs, just via
    raw XML since python-pptx doesn't wrap `<a:fld>` as a Run object at
    all. `<a:rPr>`'s child order is schema-fixed (fill before latin), so
    a new `solidFill` is inserted before any existing `latin`/`ea`/`cs`,
    not just appended.
    """
    changed = False
    has_fill = (
        rPr.find(effects_mod.a_qn("solidFill")) is not None
        or rPr.find(effects_mod.a_qn("noFill")) is not None
    )
    has_latin = rPr.find(effects_mod.a_qn("latin")) is not None

    if not has_fill:
        solid_fill = etree.Element(effects_mod.a_qn("solidFill"))
        srgb_clr = etree.SubElement(solid_fill, effects_mod.a_qn("srgbClr"))
        srgb_clr.set("val", color_hex)
        ref = rPr.find(effects_mod.a_qn("latin"))
        if ref is not None:
            ref.addprevious(solid_fill)
        else:
            rPr.append(solid_fill)
        changed = True

    if not has_latin:
        latin = etree.SubElement(rPr, effects_mod.a_qn("latin"))
        latin.set("typeface", font_name)
        changed = True

    return changed


def _remap_shapes_fills(shapes, remap: dict[str, str], min_area_emu2: float, scope: str, location: str) -> list[Change]:
    changes = []
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            changes += _remap_shapes_fills(shape.shapes, remap, min_area_emu2, scope, location)
            continue
        width, height = getattr(shape, "width", None), getattr(shape, "height", None)
        if not width or not height or (width * height) < min_area_emu2:
            continue
        spPr = effects_mod.get_spPr(shape)
        if spPr is None:
            continue
        solid_fill = spPr.find(effects_mod.a_qn("solidFill"))
        if solid_fill is None:
            continue
        srgb = solid_fill.find(effects_mod.a_qn("srgbClr"))
        if srgb is None:
            continue
        old_hex = colors_mod.normalize_hex(srgb.get("val"))
        if old_hex not in remap:
            continue
        new_hex = colors_mod.normalize_hex(remap[old_hex])
        if new_hex == old_hex:
            continue
        srgb.set("val", new_hex)
        changes.append(
            Change(
                scope=scope,
                rule="legacy_color",
                field="fill color",
                old=f"#{old_hex}",
                new=f"#{new_hex}",
                shape_id=shape.shape_id,
                shape_name=shape.name,
                location=location,
            )
        )
    return changes


def _remap_shapes_fonts(shapes, remap_by_key: dict[str, str], scope: str, location: str) -> list[Change]:
    changes = []
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            changes += _remap_shapes_fonts(shape.shapes, remap_by_key, scope, location)
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                current = run.font.name
                target = remap_by_key.get(normalize_key(current))
                if target and target != current:
                    run.font.name = target
                    changes.append(
                        Change(
                            scope=scope,
                            rule="legacy_font",
                            field="font",
                            old=current,
                            new=target,
                            shape_id=shape.shape_id,
                            shape_name=shape.name,
                            location=location,
                        )
                    )
    return changes


def _apply_violation(v: Violation, config: dict) -> Optional[Change]:
    target = v.target
    if target is None:
        return None

    if v.rule in ("legacy_color", "near_miss_color"):
        new_hex = v.details["target"].lstrip("#")
        old_hex = v.details["current"]
        if v.element == "fill":
            color_fmt = target.obj.fill.fore_color
        elif v.element == "line":
            color_fmt = target.obj.line.color
        elif v.element == "gradient stop":
            idx = v.details.get("gradient_index", 0)
            color_fmt = target.obj.fill.gradient_stops[idx].color
        elif v.element == "text":
            color_fmt = target.obj.font.color
        else:
            return None
        # Scheme-typed colors are normally left for manual review rather
        # than silently converted to a hardcoded literal (e.g. a
        # different theme not covered, or an out-of-range tint) --
        # except for legacy_color specifically. Unlike near_miss_color (a
        # fuzzy tolerance match), legacy_color is always an exact,
        # fully-confident target -- either a colors.remap hit or the
        # deterministic unlisted_panel_fallback -- so converting a
        # schemeClr+lumMod/lumOff tint (e.g. bg1 at 85% luminance, which
        # is how a "grey panel" is often actually built) to the literal
        # target is correct, not a guess.
        if color_fmt.type != MSO_COLOR_TYPE.RGB and v.rule != "legacy_color":
            return None
        color_fmt.rgb = RGBColor.from_string(new_hex)
        return Change(
            scope="slide",
            rule=v.rule,
            field=f"{v.element} color",
            old=old_hex,
            new=v.details["target"],
            slide_index=v.slide_index,
            shape_id=v.shape_id,
            shape_name=v.shape_name,
        )

    if v.rule == "unapproved_font":
        target_font = v.details.get("target")
        if not target_font:
            return None
        old_font = target.font_raw
        target.obj.font.name = target_font
        return Change(
            scope="slide",
            rule=v.rule,
            field="font",
            old=old_font,
            new=target_font,
            slide_index=v.slide_index,
            shape_id=v.shape_id,
            shape_name=v.shape_name,
        )

    if v.rule == "font_weight":
        target.obj.font.bold = False
        return Change(
            scope="slide",
            rule=v.rule,
            field="font weight",
            old="bold",
            new="regular",
            slide_index=v.slide_index,
            shape_id=v.shape_id,
            shape_name=v.shape_name,
        )

    if v.rule in ("text_color", "text_contrast"):
        new_hex = v.details["target"].lstrip("#")
        target.obj.font.color.rgb = RGBColor.from_string(new_hex)
        return Change(
            scope="slide",
            rule=v.rule,
            field="text color",
            old=v.details["current"],
            new=v.details["target"],
            slide_index=v.slide_index,
            shape_id=v.shape_id,
            shape_name=v.shape_name,
        )

    if v.rule == "text_effect":
        if v.element == "run":
            removed = effects_mod.remove_run_effects(target.obj)
        else:
            removed = effects_mod.remove_shape_effects(target.obj)
        if not removed:
            return None
        return Change(
            scope="slide",
            rule=v.rule,
            field="effects",
            old=sorted(removed),
            new=[],
            slide_index=v.slide_index,
            shape_id=v.shape_id,
            shape_name=v.shape_name,
        )

    if v.rule == "alignment":
        new_align = v.details["target"]
        target.obj.alignment = ALIGN_BY_NAME.get(new_align)
        return Change(
            scope="slide",
            rule=v.rule,
            field="alignment",
            old=v.details["current"],
            new=new_align,
            slide_index=v.slide_index,
            shape_id=v.shape_id,
            shape_name=v.shape_name,
        )

    if v.rule == "old_logo":
        logo_cfg = config.get("logo", {}) or {}
        new_logo_path = logo_cfg.get("new_logo_path")
        if not new_logo_path or not Path(new_logo_path).exists():
            return None
        logo_mod.replace_logo_image(target.obj, new_logo_path)
        return Change(
            scope="slide",
            rule=v.rule,
            field="image",
            old=v.details.get("matched"),
            new=new_logo_path,
            slide_index=v.slide_index,
            shape_id=v.shape_id,
            shape_name=v.shape_name,
        )

    return None


def fix_deck(prs, config: dict, source_path: str, output_path: Optional[str], dry_run: bool) -> FixReport:
    colors_cfg = config.get("colors", {}) or {}
    remap = {colors_mod.normalize_hex(k): v for k, v in (colors_cfg.get("remap", {}) or {}).items()}
    fonts_cfg = config.get("fonts", {}) or {}
    font_remap = fonts_cfg.get("remap", {}) or {}

    changes: list[Change] = []

    theme_color_changes = colors_mod.remap_theme_colors(prs, remap)
    for c in theme_color_changes:
        changes.append(
            Change(
                scope="theme",
                rule="legacy_color",
                field="theme color",
                old=c["old"],
                new=c["new"],
                location=f"{c['theme_part']} [{c['slot']}]",
            )
        )

    theme_font_changes = remap_theme_fonts(prs, font_remap)
    for c in theme_font_changes:
        changes.append(
            Change(
                scope="theme",
                rule="legacy_font",
                field="theme font",
                old=c["old"],
                new=c["new"],
                location=f"{c['theme_part']} [{c['slot']}]",
            )
        )

    changes += _remap_explicit_fonts_in_masters_and_layouts(prs, font_remap)

    txstyles_font_changes = fonts_mod.remap_literal_fonts_in_master_txstyles(prs, font_remap)
    for c in txstyles_font_changes:
        changes.append(
            Change(
                scope="master",
                rule="legacy_font",
                field="master default font",
                old=c["old"],
                new=c["new"],
                location=f"{c['master']} [{c['style']}]",
            )
        )

    layout_panel_remap = colors_cfg.get("layout_panel_remap", {}) or {}
    layout_panel_min_area_sq_in = colors_cfg.get("layout_panel_min_area_sq_in", 8.0)
    changes += _remap_large_panel_fills_in_masters_and_layouts(
        prs, layout_panel_remap, layout_panel_min_area_sq_in
    )

    logo_cfg = config.get("logo", {}) or {}
    changes += _replace_old_logo_everywhere(
        prs,
        logo_cfg.get("old_logo_hashes", []) or [],
        logo_cfg.get("new_logo_path"),
        logo_cfg.get("match_threshold", logo_mod.DEFAULT_MATCH_THRESHOLD),
    )
    changes += _replace_old_logo_region_everywhere(
        prs,
        logo_cfg.get("old_logo_region_in"),
        logo_cfg.get("new_logo_path"),
    )

    # Apply auto-fixable violations to a fixpoint. Fixing one violation can
    # mechanically unlock another check on the very same run — e.g. renaming
    # an off-brand font from Arial to Inter means Inter's text-color rule
    # now applies to that run for the first time, and its existing color may
    # no longer be allowed. That's not ambiguous, just sequential, so a
    # single pass would wrongly leave it for "manual review" when it's
    # actually fully resolvable. Loop until a pass makes no further changes;
    # capped defensively since each pass only ever unlocks checks that were
    # previously gated on a now-fixed field, so this converges in a couple
    # of iterations for any real rule set.
    for _ in range(5):
        inventory = build_inventory(prs)
        violations = sort_violations(audit_deck(inventory, config))
        pass_changes = []
        for v in violations:
            if not v.auto_fixable:
                continue
            change = _apply_violation(v, config)
            if change:
                pass_changes.append(change)
        changes += pass_changes
        if not pass_changes:
            break

    # Anything still detected here needs a human either way: genuinely
    # unfixable violations, or (defensively) anything a fix attempt didn't
    # actually clear.
    manual_review = sort_violations(audit_deck(build_inventory(prs), config))

    # Last, structural rather than brand-rule-driven: heal a duplicate
    # shape id inherited from the source file (or introduced by any
    # shape surgery earlier in this same run) before saving -- see
    # _dedupe_shape_ids_in_part's own docstring for why this matters.
    changes += _dedupe_shape_ids(prs)
    changes += _remove_confidentiality_footer_text(prs)
    changes += _normalize_footer_chrome_text(prs, config)
    changes += _normalize_footer_chrome_position(prs)

    if not dry_run and output_path:
        prs.save(output_path)

    return FixReport(
        source_path=source_path,
        output_path=None if dry_run else output_path,
        dry_run=dry_run,
        changes=changes,
        manual_review=manual_review,
    )
