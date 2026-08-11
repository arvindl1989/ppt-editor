"""Cross-presentation slide import: copy a cover/outro slide (and its
layout, master, theme, and referenced media) from the org template into
a target deck, then swap it in for whatever cover/outro that deck
already has.

python-pptx has no built-in support for copying slides across
presentations -- a slide's layout, that layout's master, the master's
theme, and any media they reference are all separate OPC parts tied
together by relationship IDs that are only unique *within* the file
they came from. This module does that copy at the raw zip/XML level
(the only way to do it at all), then hands off to python-pptx's normal
object model for everything else (reordering slides, setting title
text) once the imported parts are legitimate parts of the target
package.

Kept deliberately narrow: importing exactly one master with exactly the
two layouts actually needed (not all 62), rather than a general
"merge two arbitrary presentations" tool -- that's a much bigger, much
riskier problem this doesn't need to solve.
"""

from __future__ import annotations

import re
import zipfile
from pathlib import Path
from typing import Optional

from pptx import Presentation

HEADING_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}
SUBTITLE_LIKE_PLACEHOLDER_TYPES = {"SUBTITLE", "BODY"}

DEFAULT_COVER_LAYOUTS = {"Cover A", "Cover B", "Cover C", "Cover D", "Cover E", "Cover F"}
DEFAULT_OUTRO_LAYOUTS = {"Outro", "End"}


def default_template_path() -> Path:
    return (Path(__file__).parent.parent / "assets") / "kone_master_template.pptx"


def _read_zip(path) -> dict[str, bytes]:
    with zipfile.ZipFile(path) as z:
        return {n: z.read(n) for n in z.namelist()}


def _write_zip(parts: dict[str, bytes], path) -> None:
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for name, data in parts.items():
            z.writestr(name, data)


def _next_part_number(parts: dict[str, bytes], dirname: str, stem: str) -> int:
    pattern = re.compile(rf"^ppt/{dirname}/{stem}(\d+)\.")
    nums = [int(m.group(1)) for n in parts if (m := pattern.match(n))]
    return (max(nums) if nums else 0) + 1


def _next_master_or_layout_id(target: dict[str, bytes]) -> int:
    """sldMasterId and sldLayoutId share one ID namespace (a layout ID
    reused as a master ID elsewhere in the same file is a hard
    uniqueness violation -- confirmed by direct validation against the
    real OOXML schema, ISO/IEC 29500, in production), and per that same
    schema both types are constrained to >= 2147483648 (ST_SlideMasterId/
    ST_SlideLayoutId's minInclusive) -- a disjoint range from sldId
    (ST_SlideId: >= 256 and < 2147483648), so the two never need
    cross-checking against each other. A same-looking-safe hardcoded
    constant isn't actually safe -- it collided with a real deck's own
    pre-existing ID in production. Scans every sldMasterId in
    presentation.xml plus every sldLayoutId in every existing
    slideMasterN.xml, and returns one past the max found.
    """
    max_id = 2147483648 - 1
    pres_xml = target.get("ppt/presentation.xml", b"").decode()
    for m in re.finditer(r'<p:sldMasterId id="(\d+)"', pres_xml):
        max_id = max(max_id, int(m.group(1)))
    for name, data in target.items():
        if re.match(r"^ppt/slideMasters/slideMaster\d+\.xml$", name):
            for m in re.finditer(r'<p:sldLayoutId id="(\d+)"', data.decode()):
                max_id = max(max_id, int(m.group(1)))
    return max_id + 1


def _next_sld_id(target: dict[str, bytes]) -> int:
    """sldId's own, separate ID pool -- see _next_master_or_layout_id for
    why these two ranges are disjoint by schema design."""
    max_id = 255
    pres_xml = target.get("ppt/presentation.xml", b"").decode()
    for m in re.finditer(r'<p:sldId id="(\d+)"', pres_xml):
        max_id = max(max_id, int(m.group(1)))
    return max_id + 1


def _rels_path(part_path: str) -> str:
    d = part_path.rsplit("/", 1)
    return f"{d[0]}/_rels/{d[1]}.rels"


def _parse_rels(rels_bytes: Optional[bytes]) -> dict[str, tuple[str, str]]:
    """{rId: (type, target)}"""
    if rels_bytes is None:
        return {}
    text = rels_bytes.decode()
    return {
        m.group(1): (m.group(2), m.group(3))
        for m in re.finditer(r'<Relationship Id="(rId\d+)" Type="([^"]+)" Target="([^"]+)"[^/]*/>', text)
    }


def _resolve_relative(part_path: str, rel_target: str) -> str:
    parts = part_path.split("/")[:-1]
    for seg in rel_target.split("/"):
        if seg == "..":
            parts.pop()
        else:
            parts.append(seg)
    return "/".join(parts)


def import_cover_and_outro_parts(
    target_path,
    template_path,
    out_path,
    cover_slide_num: int,
    outro_slide_num: int,
) -> tuple[str, str]:
    """Copy `template_path`'s slide `cover_slide_num` and `outro_slide_num`
    -- plus a single shared master trimmed to just their two layouts, that
    master's theme, and every media file they reference -- into a copy of
    `target_path`, written to `out_path`.

    The new slides are appended to the end of the deck's slide list, not
    yet positioned or populated with content -- see `replace_intro_outro`
    for that. Returns (new_cover_slide_partname, new_outro_slide_partname).
    """
    tmpl = _read_zip(template_path)
    target = dict(_read_zip(target_path))

    next_media_num = _next_part_number(target, "media", "image")
    next_master_num = _next_part_number(target, "slideMasters", "slideMaster")
    next_layout_num = _next_part_number(target, "slideLayouts", "slideLayout")
    next_theme_num = _next_part_number(target, "theme", "theme")
    next_slide_num = _next_part_number(target, "slides", "slide")

    next_master_layout_id = _next_master_or_layout_id(target)
    new_master_id = next_master_layout_id
    new_cover_layout_id = next_master_layout_id + 1
    new_outro_layout_id = next_master_layout_id + 2
    next_sld_id = _next_sld_id(target)
    new_cover_sld_id = next_sld_id
    new_outro_sld_id = next_sld_id + 1

    pres_rels_text = target["ppt/_rels/presentation.xml.rels"].decode()
    next_rid = max((int(m) for m in re.findall(r'Id="rId(\d+)"', pres_rels_text)), default=0) + 1

    media_map: dict[str, str] = {}

    def copy_media(tmpl_media_path: str) -> str:
        nonlocal next_media_num
        if tmpl_media_path in media_map:
            return media_map[tmpl_media_path]
        ext = tmpl_media_path.rsplit(".", 1)[-1]
        new_path = f"ppt/media/image{next_media_num}.{ext}"
        target[new_path] = tmpl[tmpl_media_path]
        media_map[tmpl_media_path] = new_path
        next_media_num += 1
        return new_path

    tmpl_prs = Presentation(template_path)
    cover_slide = tmpl_prs.slides[cover_slide_num - 1]
    outro_slide = tmpl_prs.slides[outro_slide_num - 1]
    cover_layout_src = str(cover_slide.slide_layout.part.partname).lstrip("/")
    outro_layout_src = str(outro_slide.slide_layout.part.partname).lstrip("/")
    cover_slide_src = str(cover_slide.part.partname).lstrip("/")
    outro_slide_src = str(outro_slide.part.partname).lstrip("/")

    # --- theme ---
    new_theme_path = f"ppt/theme/theme{next_theme_num}.xml"
    target[new_theme_path] = tmpl["ppt/theme/theme1.xml"]

    # --- layouts ---
    layout_src_to_new: dict[str, str] = {}
    for src_layout in (cover_layout_src, outro_layout_src):
        new_layout_path = f"ppt/slideLayouts/slideLayout{next_layout_num}.xml"
        target[new_layout_path] = tmpl[src_layout]

        rels_bytes = tmpl.get(_rels_path(src_layout))
        rels_text = rels_bytes.decode() if rels_bytes else ""
        for _rid, (rtype, rtarget) in _parse_rels(rels_bytes).items():
            if rtype.endswith("/image"):
                resolved = _resolve_relative(src_layout, rtarget)
                new_media_path = copy_media(resolved)
                rels_text = rels_text.replace(
                    f'Target="{rtarget}"', f'Target="../media/{new_media_path.rsplit("/", 1)[-1]}"'
                )
            elif rtype.endswith("/slideMaster"):
                rels_text = rels_text.replace(
                    f'Target="{rtarget}"', f'Target="../slideMasters/slideMaster{next_master_num}.xml"'
                )
        target[_rels_path(new_layout_path)] = rels_text.encode()
        layout_src_to_new[src_layout] = new_layout_path
        next_layout_num += 1

    # --- master (trimmed to just the two imported layouts) ---
    new_master_path = f"ppt/slideMasters/slideMaster{next_master_num}.xml"
    master_xml = tmpl["ppt/slideMasters/slideMaster1.xml"].decode()
    master_xml = re.sub(
        r"<p:sldLayoutIdLst>.*?</p:sldLayoutIdLst>",
        "<p:sldLayoutIdLst>"
        f'<p:sldLayoutId id="{new_cover_layout_id}" r:id="rId101"/>'
        f'<p:sldLayoutId id="{new_outro_layout_id}" r:id="rId102"/>'
        "</p:sldLayoutIdLst>",
        master_xml,
        flags=re.S,
    )
    target[new_master_path] = master_xml.encode()

    cover_layout_new = layout_src_to_new[cover_layout_src]
    outro_layout_new = layout_src_to_new[outro_layout_src]
    target[_rels_path(new_master_path)] = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        f'<Relationship Id="rId100" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" '
        f'Target="../theme/theme{next_theme_num}.xml"/>'
        f'<Relationship Id="rId101" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" '
        f'Target="../slideLayouts/{cover_layout_new.rsplit("/", 1)[-1]}"/>'
        f'<Relationship Id="rId102" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" '
        f'Target="../slideLayouts/{outro_layout_new.rsplit("/", 1)[-1]}"/>'
        "</Relationships>"
    ).encode()

    # --- slides ---
    slide_src_to_new: dict[str, str] = {}
    for src_slide, layout_new in ((cover_slide_src, cover_layout_new), (outro_slide_src, outro_layout_new)):
        new_slide_path = f"ppt/slides/slide{next_slide_num}.xml"
        target[new_slide_path] = tmpl[src_slide]

        rels_bytes = tmpl.get(_rels_path(src_slide))
        rels_text = rels_bytes.decode() if rels_bytes else (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"></Relationships>'
        )
        for _rid, (rtype, rtarget) in _parse_rels(rels_bytes).items():
            if rtype.endswith("/image"):
                resolved = _resolve_relative(src_slide, rtarget)
                new_media_path = copy_media(resolved)
                rels_text = rels_text.replace(
                    f'Target="{rtarget}"', f'Target="../media/{new_media_path.rsplit("/", 1)[-1]}"'
                )
            elif rtype.endswith("/slideLayout"):
                rels_text = rels_text.replace(
                    f'Target="{rtarget}"', f'Target="../slideLayouts/{layout_new.rsplit("/", 1)[-1]}"'
                )
        target[_rels_path(new_slide_path)] = rels_text.encode()
        slide_src_to_new[src_slide] = new_slide_path
        next_slide_num += 1

    cover_new = slide_src_to_new[cover_slide_src]
    outro_new = slide_src_to_new[outro_slide_src]

    # --- [Content_Types].xml ---
    ct = target["[Content_Types].xml"].decode()
    inserts = [
        f'<Override PartName="/{new_master_path}" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>',
        *(
            f'<Override PartName="/{p}" '
            'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>'
            for p in layout_src_to_new.values()
        ),
        *(
            f'<Override PartName="/{p}" '
            'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
            for p in (cover_new, outro_new)
        ),
        f'<Override PartName="/{new_theme_path}" '
        'ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>',
    ]
    ct = ct.replace("</Types>", "".join(inserts) + "</Types>")
    # Every media extension actually being introduced needs a Default
    # entry -- a target with no pre-existing images of a given type (e.g.
    # a brand-new deck with no images at all) won't have declared one,
    # and a part with no content-type mapping fails to load at all, not
    # just render wrong.
    media_content_types = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "svg": "image/svg+xml",
        "gif": "image/gif",
        "bmp": "image/bmp",
        "tiff": "image/tiff",
        "wmf": "image/x-wmf",
        "emf": "image/x-emf",
    }
    for new_media_path in media_map.values():
        ext = new_media_path.rsplit(".", 1)[-1].lower()
        content_type = media_content_types.get(ext)
        if content_type and f'Extension="{ext}"' not in ct:
            ct = ct.replace(
                "<Default Extension=", f'<Default Extension="{ext}" ContentType="{content_type}"/><Default Extension=', 1
            )
    target["[Content_Types].xml"] = ct.encode()

    # --- presentation.xml: add the master; append both slides at the end ---
    pres_xml = target["ppt/presentation.xml"].decode()
    master_rid = next_rid
    next_rid += 1
    pres_xml = pres_xml.replace(
        "</p:sldMasterIdLst>",
        f'<p:sldMasterId id="{new_master_id}" r:id="rId{master_rid}"/></p:sldMasterIdLst>',
    )

    cover_rid = next_rid
    next_rid += 1
    outro_rid = next_rid
    next_rid += 1

    pres_xml = pres_xml.replace(
        "</p:sldIdLst>",
        f'<p:sldId id="{new_cover_sld_id}" r:id="rId{cover_rid}"/>'
        f'<p:sldId id="{new_outro_sld_id}" r:id="rId{outro_rid}"/></p:sldIdLst>',
    )
    target["ppt/presentation.xml"] = pres_xml.encode()

    pres_rels_text = pres_rels_text.replace(
        "</Relationships>",
        f'<Relationship Id="rId{master_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" '
        f'Target="slideMasters/slideMaster{next_master_num}.xml"/>'
        f'<Relationship Id="rId{cover_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
        f'Target="slides/{cover_new.rsplit("/", 1)[-1]}"/>'
        f'<Relationship Id="rId{outro_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
        f'Target="slides/{outro_new.rsplit("/", 1)[-1]}"/>'
        "</Relationships>",
    )
    target["ppt/_rels/presentation.xml.rels"] = pres_rels_text.encode()

    _write_zip(target, out_path)
    return cover_new, outro_new


def import_layouts(target_path, template_path, out_path, layout_names: list[str]) -> dict[str, str]:
    """Copy the named layouts (by exact name) from `template_path` -- plus
    a single shared master trimmed to just those layouts, that master's
    theme, and every media file they reference -- into a copy of
    `target_path`, written to `out_path`.

    Unlike `import_cover_and_outro_parts`, this copies LAYOUT parts only,
    no slide parts -- callers create slides on the imported layouts with
    plain `prs.slides.add_slide(layout)` once the output is reopened,
    rather than this module hand-building slide XML. That keeps this
    generic (any number of layouts, any names) without duplicating
    python-pptx's own slide-creation logic.

    Returns {layout_name: new_layout_partname}. Raises KeyError if a
    requested name isn't found among the template's layouts.
    """
    tmpl = _read_zip(template_path)
    target = dict(_read_zip(target_path))

    next_media_num = _next_part_number(target, "media", "image")
    next_master_num = _next_part_number(target, "slideMasters", "slideMaster")
    next_layout_num = _next_part_number(target, "slideLayouts", "slideLayout")
    next_theme_num = _next_part_number(target, "theme", "theme")

    next_id = _next_master_or_layout_id(target)
    new_master_id = next_id

    pres_rels_text = target["ppt/_rels/presentation.xml.rels"].decode()
    next_rid = max((int(m) for m in re.findall(r'Id="rId(\d+)"', pres_rels_text)), default=0) + 1

    tmpl_prs = Presentation(template_path)
    layout_by_name = {layout.name: layout for master in tmpl_prs.slide_masters for layout in master.slide_layouts}
    src_paths = []
    for name in layout_names:
        if name not in layout_by_name:
            raise KeyError(f"template has no layout named {name!r}")
        src_paths.append(str(layout_by_name[name].part.partname).lstrip("/"))

    # All requested layouts share one imported master (this function's own
    # "single shared master trimmed to just those layouts" design) -- that
    # master's own TOP-LEVEL shapes (e.g. a logo picture placed directly
    # on the master itself, not any one layout -- a real case: a deck
    # whose branding lives on the master, not the layout) can reference
    # media too, so its own rels need to survive the trim below, not just
    # each layout's.
    src_master_path = str(layout_by_name[layout_names[0]].slide_master.part.partname).lstrip("/")

    media_map: dict[str, str] = {}

    def copy_media(tmpl_media_path: str) -> str:
        nonlocal next_media_num
        if tmpl_media_path in media_map:
            return media_map[tmpl_media_path]
        ext = tmpl_media_path.rsplit(".", 1)[-1]
        new_path = f"ppt/media/image{next_media_num}.{ext}"
        target[new_path] = tmpl[tmpl_media_path]
        media_map[tmpl_media_path] = new_path
        next_media_num += 1
        return new_path

    # --- theme ---
    new_theme_path = f"ppt/theme/theme{next_theme_num}.xml"
    target[new_theme_path] = tmpl["ppt/theme/theme1.xml"]

    # --- layouts ---
    layout_ids: list[int] = []
    layout_new_paths: list[str] = []
    for src_path in src_paths:
        new_layout_path = f"ppt/slideLayouts/slideLayout{next_layout_num}.xml"
        target[new_layout_path] = tmpl[src_path]

        rels_bytes = tmpl.get(_rels_path(src_path))
        rels_text = rels_bytes.decode() if rels_bytes else ""
        for _rid, (rtype, rtarget) in _parse_rels(rels_bytes).items():
            if rtype.endswith("/image"):
                resolved = _resolve_relative(src_path, rtarget)
                new_media_path = copy_media(resolved)
                rels_text = rels_text.replace(
                    f'Target="{rtarget}"', f'Target="../media/{new_media_path.rsplit("/", 1)[-1]}"'
                )
            elif rtype.endswith("/slideMaster"):
                rels_text = rels_text.replace(
                    f'Target="{rtarget}"', f'Target="../slideMasters/slideMaster{next_master_num}.xml"'
                )
        target[_rels_path(new_layout_path)] = rels_text.encode()
        layout_ids.append(next_id + 1 + len(layout_new_paths))
        layout_new_paths.append(new_layout_path)
        next_layout_num += 1

    # --- master (trimmed to just the imported layouts) ---
    new_master_path = f"ppt/slideMasters/slideMaster{next_master_num}.xml"
    master_xml = tmpl[src_master_path].decode()

    # The trimmed master_xml above still has whatever r:embed="rIdN" its
    # own top-level shapes originally used (unchanged -- only the
    # <p:sldLayoutIdLst> is rewritten below), so any of the SOURCE
    # master's own media relationships must survive at their ORIGINAL
    # rIds, pointing at the now-locally-copied media file. Its other
    # original relationships (its own full ~60-layout list, its own old
    # theme) aren't needed -- this master keeps only the requested
    # layouts and gets a fresh theme relationship, built below at rIds
    # guaranteed not to collide with whatever the source master already
    # used for its own media.
    src_master_rels = _parse_rels(tmpl.get(_rels_path(src_master_path)))
    max_src_rid = 0
    preserved_media_rels = []
    for rid, (rtype, rtarget) in src_master_rels.items():
        max_src_rid = max(max_src_rid, int(rid[len("rId"):]))
        if rtype.endswith("/image"):
            resolved = _resolve_relative(src_master_path, rtarget)
            new_media_path = copy_media(resolved)
            preserved_media_rels.append(
                f'<Relationship Id="{rid}" Type="{rtype}" Target="../media/{new_media_path.rsplit("/", 1)[-1]}"/>'
            )

    theme_rid_num = max_src_rid + 1
    layout_rid_start = theme_rid_num + 1
    layout_entries = "".join(
        f'<p:sldLayoutId id="{layout_id}" r:id="rId{layout_rid_start + i}"/>' for i, layout_id in enumerate(layout_ids)
    )
    master_xml = re.sub(
        r"<p:sldLayoutIdLst>.*?</p:sldLayoutIdLst>",
        f"<p:sldLayoutIdLst>{layout_entries}</p:sldLayoutIdLst>",
        master_xml,
        flags=re.S,
    )
    target[new_master_path] = master_xml.encode()

    master_rels = list(preserved_media_rels)
    master_rels.append(
        f'<Relationship Id="rId{theme_rid_num}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" '
        f'Target="../theme/theme{next_theme_num}.xml"/>'
    )
    for i, new_layout_path in enumerate(layout_new_paths):
        master_rels.append(
            f'<Relationship Id="rId{layout_rid_start + i}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" '
            f'Target="../slideLayouts/{new_layout_path.rsplit("/", 1)[-1]}"/>'
        )
    target[_rels_path(new_master_path)] = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        + "".join(master_rels) + "</Relationships>"
    ).encode()

    # --- [Content_Types].xml ---
    ct = target["[Content_Types].xml"].decode()
    inserts = [
        f'<Override PartName="/{new_master_path}" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>',
        *(
            f'<Override PartName="/{p}" '
            'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>'
            for p in layout_new_paths
        ),
        f'<Override PartName="/{new_theme_path}" '
        'ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>',
    ]
    ct = ct.replace("</Types>", "".join(inserts) + "</Types>")
    media_content_types = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "svg": "image/svg+xml",
        "gif": "image/gif",
        "bmp": "image/bmp",
        "tiff": "image/tiff",
        "wmf": "image/x-wmf",
        "emf": "image/x-emf",
    }
    for new_media_path in media_map.values():
        ext = new_media_path.rsplit(".", 1)[-1].lower()
        content_type = media_content_types.get(ext)
        if content_type and f'Extension="{ext}"' not in ct:
            ct = ct.replace(
                "<Default Extension=", f'<Default Extension="{ext}" ContentType="{content_type}"/><Default Extension=', 1
            )
    target["[Content_Types].xml"] = ct.encode()

    # --- presentation.xml: add the master only (no slides yet) ---
    pres_xml = target["ppt/presentation.xml"].decode()
    master_rid = next_rid
    pres_xml = pres_xml.replace(
        "</p:sldMasterIdLst>",
        f'<p:sldMasterId id="{new_master_id}" r:id="rId{master_rid}"/></p:sldMasterIdLst>',
    )
    target["ppt/presentation.xml"] = pres_xml.encode()

    pres_rels_text = pres_rels_text.replace(
        "</Relationships>",
        f'<Relationship Id="rId{master_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" '
        f'Target="slideMasters/slideMaster{next_master_num}.xml"/></Relationships>',
    )
    target["ppt/_rels/presentation.xml.rels"] = pres_rels_text.encode()

    _write_zip(target, out_path)
    return dict(zip(layout_names, layout_new_paths))


def _extract_lead_text(slide, placeholder_types: set[str]) -> Optional[str]:
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph_type = shape.placeholder_format.type
        if ph_type is not None and ph_type.name in placeholder_types:
            if shape.has_text_frame and shape.text_frame.text.strip():
                return shape.text_frame.text
    return None


def _set_placeholder_text(slide, placeholder_types: set[str], text: str) -> bool:
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph_type = shape.placeholder_format.type
        if ph_type is not None and ph_type.name in placeholder_types:
            shape.text_frame.text = text
            return True
    return False


def _move_slide(prs, old_index: int, new_index: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def _delete_slide(prs, index: int) -> None:
    """Remove both the slide's <p:sldId> entry AND its relationship in
    presentation.xml.rels.

    Dropping only the sldId entry leaves a dangling slide relationship
    that nothing references anymore -- structurally harmless (a
    conformant reader ignores it), but it's still reachable via
    `package.iter_parts()`'s relationship walk independently of
    `prs.slides`. In practice this makes python-pptx instantiate a
    second, uncached Part object for that partname the next time
    anything lazily touches an rPr elsewhere in the deck (confirmed by
    direct repro), and it writes both on save -- a real
    'Duplicate name' zip corruption, not just clutter. Dropping the
    relationship here removes the only path to the orphaned part, so
    there's nothing left to double-instantiate.
    """
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    sld_id_el = slides[index]
    rId = sld_id_el.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    xml_slides.remove(sld_id_el)
    if rId:
        prs.part.drop_rel(rId)


def replace_intro_outro(
    deck_path,
    out_path,
    template_path=None,
    cover_slide_num: int = 3,
    outro_slide_num: int = 11,
    cover_layouts: Optional[set[str]] = None,
    outro_layouts: Optional[set[str]] = None,
) -> dict:
    """If the deck's first slide isn't built on one of the org template's
    Cover layouts, replace it with the template's `cover_slide_num` slide
    (carrying over the old title/subtitle text into the new placeholders).
    Same for the last slide vs. `outro_slide_num`. Writes the result to
    `out_path` regardless of whether either replacement happened.
    """
    template_path = template_path or default_template_path()
    cover_layouts = cover_layouts or DEFAULT_COVER_LAYOUTS
    outro_layouts = outro_layouts or DEFAULT_OUTRO_LAYOUTS

    prs = Presentation(deck_path)
    result = {"cover_replaced": False, "outro_replaced": False}

    first_layout = prs.slides[0].slide_layout.name
    last_layout = prs.slides[-1].slide_layout.name
    need_cover = first_layout not in cover_layouts
    need_outro = last_layout not in outro_layouts

    if not need_cover and not need_outro:
        import shutil

        shutil.copy(deck_path, out_path)
        return result

    old_cover_title = _extract_lead_text(prs.slides[0], HEADING_PLACEHOLDER_TYPES) if need_cover else None
    old_cover_subtitle = _extract_lead_text(prs.slides[0], SUBTITLE_LIKE_PLACEHOLDER_TYPES) if need_cover else None
    old_outro_title = _extract_lead_text(prs.slides[-1], HEADING_PLACEHOLDER_TYPES) if need_outro else None

    cover_part, outro_part = import_cover_and_outro_parts(
        deck_path, template_path, out_path, cover_slide_num, outro_slide_num
    )

    prs2 = Presentation(out_path)
    part_to_index = {str(s.part.partname).lstrip("/"): i for i, s in enumerate(prs2.slides)}
    new_cover_index = part_to_index[cover_part]
    new_outro_index = part_to_index[outro_part]

    if need_cover:
        cover_slide = prs2.slides[new_cover_index]
        if old_cover_title:
            _set_placeholder_text(cover_slide, HEADING_PLACEHOLDER_TYPES, old_cover_title)
        if old_cover_subtitle:
            _set_placeholder_text(cover_slide, SUBTITLE_LIKE_PLACEHOLDER_TYPES, old_cover_subtitle)
        result["cover_replaced"] = True
        result["cover_title"] = old_cover_title

    if need_outro:
        outro_slide = prs2.slides[new_outro_index]
        if old_outro_title:
            _set_placeholder_text(outro_slide, HEADING_PLACEHOLDER_TYPES, old_outro_title)
        result["outro_replaced"] = True
        result["outro_title"] = old_outro_title

    # Move the new slides into position, then drop whichever originals
    # they're replacing. Indices shift after every move/delete, so this
    # re-resolves them each step rather than computing them all up front.
    if need_cover:
        idx = [i for i, s in enumerate(prs2.slides) if str(s.part.partname).lstrip("/") == cover_part][0]
        _move_slide(prs2, idx, 0)
        old_cover_idx = 1  # the deck's original first slide, now pushed to index 1
        _delete_slide(prs2, old_cover_idx)

    if need_outro:
        idx = [i for i, s in enumerate(prs2.slides) if str(s.part.partname).lstrip("/") == outro_part][0]
        last_index = len(prs2.slides) - 1
        _move_slide(prs2, idx, last_index)
        old_outro_idx = last_index - 1  # the deck's original last slide, now just before it
        _delete_slide(prs2, old_outro_idx)

    prs2.save(out_path)
    return result
