"""Bridge to the `kone-deck-generator` Claude Code skill, for deckguard's
"build a brand-new deck from a brief, no source deck at all" capability
(`redesign_deck(deck_path=None, brief=...)`).

That capability used to go through `compose.py`'s own `build_deck`: a
content/stat/timeline slide gets FIT onto whichever of deckguard's ~25
org-template layouts matches its shape (`match_layout`), placeholder
text is left theme-inherited, then one `fix_deck` pass resolves that
inheritance to an explicit brand color. That machinery exists to serve
THREE different callers of `redesign_deck` at once (redesign real
content, fill gaps in an existing deck, or author one from nothing),
and for the "real content, possibly from an existing deck" cases it's
still the right tool -- there's no source content here to hand to a
purpose-built, hand-tuned renderer instead.

But for the pure "nothing but a brief" case, the `kone-deck-generator`
skill (an independently maintained sibling of the `kone-design` skill,
normally installed under `~/.claude/skills/`) is a BETTER renderer for
exactly this one job: a library of 23 real slide ARCHETYPES (data, not
code -- `kone_engine.py`'s declarative region/group renderer plus
`archetypes_batch1/2/3.py`'s definitions), each a near-literal
transcription of a real KONE slide with real icons/charts/backgrounds,
already-compliant KONE brand colors and shrink-to-fit text -- compliant
by construction, not by a subsequent inherited-color-resolution pass.
This module imports that skill's own `kone_deck_creator.build_deck`
(never copies it, so improving an archetype there benefits both the
standalone skill and deckguard at once) and does the planning step
(brief -> spec JSON) the same way `redesign.py`'s own
`call_claude_for_outline` does -- same client/model/Usage conventions
-- but targeting the skill's own archetype-keyed spec schema instead
of compose.py's, with per-archetype guidance built at call time from
the skill's own `catalog.json` (purpose/keywords/slots) and
`archetypes.SAMPLES` (a worked content example per archetype), not
hand-duplicated here. Content shape varies too much archetype-to-
archetype (a 5-icon row vs. a 2x2 matrix vs. a comparison table) for a
single rigid structured-output schema the way the old 6-layout system
had one, so this asks for plain JSON in the response and validates/
parses it here instead -- the same approach the skill's own (now
retired) `kone_planner.call_claude` used before deckguard had any
tighter integration at all.

For the other two `redesign_deck` starting points (an existing deck's
real content, with or without a brief to fill gaps), the skill's own
whole-deck `build_deck` still doesn't apply directly -- it always
starts fresh from its OWN bundled master and has no notion of
appending onto compose.py's outline-driven build. But the underlying
per-slide primitive, `archetypes.render(slide, name, content)`, is
just a renderer for one already-added slide -- not tied to building a
whole deck from nothing -- so `select_archetype_overrides` +
`build_deck_with_archetypes` below let an archetype coexist
SLIDE-BY-SLIDE with compose.py's own org-template layouts for THOSE
two starting points too: `call_claude_for_outline`'s already-planned
outline (verbatim content, unchanged) is offered a second, additive,
fail-closed pass asking whether any "content"/"stat"/"timeline" slide
would read meaningfully better as one of the 23 archetypes than its
assigned generic layout, and if so that one slide renders through the
archetype engine while every other slide keeps using compose.py's
`_select_layout`/`_populate` exactly as it does today. Archetype slides
are rendered AFTER `fix_deck`'s theme/inherited-color pass runs over
everything else, then spliced back into the outline's own order -- see
`build_deck_with_archetypes`'s own docstring for why (their hardcoded
role colors, e.g. `kone_engine.py`'s muted caption grey `#727272`,
aren't all in `brand_rules.yaml`'s approved-colors list, so running
`fix_deck` over them would flag genuinely-fine archetype styling as
non-compliant).

Where the skill is loaded from (`_skill_dir`, in order):
1. `KONE_DECK_GENERATOR_DIR`, if set -- an explicit override, e.g. for
   a developer actively co-editing the live skill.
2. `~/.claude/skills/kone-deck-generator`, if present -- the standard
   Claude Code interactive-session location, so nothing changes for
   anyone developing deckguard and the skill together right here.
3. deckguard's OWN bundled copy under `deckguard/assets/kone_deck_generator/`
   -- a plain file copy (not a symlink, not re-fetched at build time),
   vendored so the capability actually works wherever deckguard itself
   is deployed (a Railway build has no `~/.claude/skills/` at all).
   `kone_deck_creator.py`/`kone_planner.py` are copied byte-for-byte,
   never modified -- keeping them re-syncable from the skill by just
   copying the files again, not a fork to keep merging. Its `MASTER`
   pptx path resolves as a sibling `kone-design/uploads/...` directory
   (the skill's own default, unmodified), which is why that file lives
   at `deckguard/assets/kone-design/uploads/...` here too.
"""

from __future__ import annotations

# The registry lives in `deckguard.registry` now. It used to live here,
# and two copies meant two caches: a test that invalidated one left the
# other stale, so what the loader returned depended on import order.
from deckguard.registry import (  # noqa: F401  -- re-exported for callers
    _archetype_image_slots,
    _derived_content_keys,
    _ensure_skill_on_path,
    _kone_catalog,
    _load_archetypes,
    _load_creator,
    _photo_library,
    _sample_agrees,
    _sample_without_image_paths,
    _skill_dir,
    archetype_image_capacity,
    fill_empty_photo_slots,
    invalidate_archetype_caches,
)


import importlib
import json
import os
import re
import sys
import tempfile
from pathlib import Path
from typing import Optional

from pptx import Presentation

import anthropic

from deckguard.legacy.compose import ComposeResult
from deckguard.legacy.redesign import (
    DEFAULT_MODEL,
    RedesignError,
    RedesignResult,
    Usage,
    _stream_final_message,
)




# --------------------------------------------------------------------------
# Planning: brief -> the skill's own spec JSON (title + slides, each a
# {"archetype": <name>, ...that archetype's own content fields...}).
# Content shape varies too much archetype-to-archetype (a 5-icon row vs.
# a 2x2 matrix vs. a comparison table) for one rigid structured-output
# schema, so this builds a plain-text prompt (rules + the full archetype
# guide, generated from the skill's own catalog.json + archetypes.SAMPLES)
# and parses the model's JSON response directly, same approach the
# skill's own (now retired) kone_planner.call_claude used.
# --------------------------------------------------------------------------

_KONE_SYSTEM_RULES = """\
You are the deck planner for KONE's deck generator. You do NOT design visuals -- a
separate engine renders each archetype to exact KONE brand spec (geometry, fonts,
colors, icons, charts, backgrounds). Your only job: turn a brief into a slide-by-slide
plan as strict JSON.

RULES
- The cover (intro) and closing (outro) are added automatically from the KONE master.
  Do NOT plan a cover or a thank-you slide. Plan the BODY only.
- For each idea in the brief, pick the archetype whose PURPOSE and shape fit best --
  use the "Use when"/keywords guidance below, not variety for its own sake.
- Fill only that archetype's own content fields, matching the shape of its worked
  example below (same keys; a "groups"-style example is a list of dicts, keep it a
  list of dicts with the same per-item keys).
- Never invent facts, numbers, or names beyond what the brief supports -- if the
  brief doesn't give you enough for a specific claim, write a more general but still
  concrete statement rather than a fabricated number.
- KONE voice: sentence case, plain, confident, no marketing fluff, no emoji.
- Vary archetypes across the deck; don't reuse one archetype for everything.
- Output ONLY a JSON object, no prose, no markdown fences, of this exact shape:
  {"title": "<deck title, fills the retained cover>",
   "slides": [{"archetype": "<name>", ...that archetype's own content fields...}, ...]}
"""


def _kone_archetype_guide() -> str:
    """One entry per known archetype: purpose, routing keywords, slots
    (all from catalog.json when present) and a worked content example
    (from archetypes.SAMPLES) -- built at call time from the skill's own
    data so it can never drift out of sync with whatever archetypes the
    skill actually ships."""
    archetypes = _load_archetypes()
    catalog = dict(_kone_catalog())
    # Extras carry their own routing: they are not in `catalog.json`,
    # and an archetype the planner has no reason to choose is one that
    # never gets chosen.
    from deckguard.layouts import _EXTRA_META

    for key, meta in _EXTRA_META.items():
        if key in archetypes.ARCHETYPES:
            entry = dict(catalog.get(key) or {})
            entry.setdefault("purpose", meta.get("purpose", ""))
            entry.setdefault("keywords", meta.get("keywords", []))
            entry.setdefault("notes", meta.get("notes", []))
            catalog[key] = entry
    parts = []
    for name in sorted(archetypes.ARCHETYPES.keys()):
        info = catalog.get(name, {})
        lines = [f"### {name}"]
        if info.get("purpose"):
            lines.append(f"Purpose: {info['purpose']}")
        if info.get("keywords"):
            lines.append(f"Use when: {', '.join(info['keywords'])}")
        if info.get("slots"):
            lines.append(f"Slots: {info['slots']}")
        # The authority, derived from the registry the renderer will
        # actually use. `catalog.json` covers 22 of 80 archetypes and
        # `SAMPLES` 41, so most were described by name alone -- and a
        # stale sample is worse than none: `agenda_a_table`'s advertised
        # `text1..text4` long after the renderer was rebuilt to read
        # `items`, so a planner emitted four keys nothing reads and the
        # agenda came out as a title on an empty slide.
        keys = _derived_content_keys(name)
        if keys:
            lines.append(
                "Content keys (authoritative -- anything else is discarded): "
                + "; ".join(keys)
            )
        slot = _archetype_image_slots().get(name)
        if slot is not None:
            lines.append(
                "Pictures: this archetype has picture slot(s), filled automatically from the slide's own "
                "images -- do NOT emit an image path or filename yourself."
            )
        for note in (catalog.get(name, {}) or {}).get("notes", []):
            lines.append(f"Rule: {note}")
        sample = archetypes.SAMPLES.get(name)
        if sample is not None and _sample_agrees(name, sample):
            lines.append(f"Example content: {json.dumps(_sample_without_image_paths(name, sample), ensure_ascii=False)}")
        parts.append("\n".join(lines))
    return "\n\n".join(parts)


def _kone_system_prompt() -> str:
    return _KONE_SYSTEM_RULES + "\n\nAvailable archetypes:\n\n" + _kone_archetype_guide()


def _build_kone_messages(brief: str, target_slides: Optional[int], notes: Optional[str]) -> list:
    sections = [f"Brief describing the deck to build:\n{brief}"]
    if target_slides is not None:
        sections.append(f"Aim for approximately {target_slides} total BODY slides (excluding the retained cover/thank-you).")
    if notes:
        sections.append(f"Additional guidance from the operator running this tool:\n{notes}")
    return [{"role": "user", "content": "\n\n".join(sections)}]


def call_claude_for_kone_spec(
    brief: str,
    target_slides: Optional[int] = None,
    model: str = DEFAULT_MODEL,
    effort: str = "high",
    notes: Optional[str] = None,
    api_key: Optional[str] = None,
    client=None,
) -> tuple[dict, Usage]:
    """Same call shape as `redesign.call_claude_for_outline`, targeting
    the skill's archetype-keyed spec instead of compose.py's outline
    schema. `effort` is accepted for interface parity with the rest of
    `redesign_deck`'s callers but currently unused: this call doesn't use
    structured-output's `output_config` (see module docstring for why),
    and that's the only place this codebase's `effort` knob attaches.
    `client` is the same test-injection point (any `.messages.stream(...)`
    Anthropic-SDK-shaped object) the rest of the redesign test suite uses.
    """
    del effort
    if client is None:
        key = api_key or os.environ.get("ANTHROPIC_API_KEY")
        if not key:
            raise RedesignError("ANTHROPIC_API_KEY is not set -- redesign needs an Anthropic API key")
        client = anthropic.Anthropic(api_key=key)

    system = _kone_system_prompt()
    messages = _build_kone_messages(brief, target_slides, notes)
    response = _stream_final_message(
        client,
        model=model,
        max_tokens=16000,
        thinking={"type": "adaptive"},
        system=system,
        messages=messages,
    )

    if response.stop_reason == "refusal":
        raise RedesignError("Claude declined the deck-planning request (safety refusal) -- try again or adjust the brief")

    text = next((b.text for b in response.content if b.type == "text"), None)
    if not text:
        raise RedesignError(f"no text content in the model response (stop_reason={response.stop_reason!r})")
    text = text.strip()
    if text.startswith("```"):
        text = text.strip("`")
        text = text[text.find("\n") + 1:] if "\n" in text else text
    try:
        spec = json.loads(text)
    except json.JSONDecodeError as exc:
        raise RedesignError(f"model response was not valid JSON: {exc}") from exc

    usage = Usage(input_tokens=response.usage.input_tokens, output_tokens=response.usage.output_tokens, model=model)
    return spec, usage


# --------------------------------------------------------------------------
# Validation: STRUCTURAL only (known archetype name, a slide dict that
# actually names one) -- content shape is deliberately NOT re-validated
# per-archetype here: `kone_engine.render_archetype` is already
# defensive about a missing/short content list (a region whose content
# key is absent is just skipped; a `groups` list shorter than its origin
# slots just leaves the extra slots empty -- see its own `zip` over
# `origins`/`items`), so there's no missing-content failure mode this
# needs to catch before rendering the way the old fixed 6-layout system
# did (exactly-N-columns, exactly-N-stats). Character/length limits
# were dropped from validation entirely back when the skill's own
# shrink-to-fit renderer made them advisory (see the commit that did
# that for the previous archetype set) -- still true here.
# --------------------------------------------------------------------------


def _validate_kone_spec(spec: dict, known_archetypes: set) -> None:
    errors: list = []
    if not spec.get("title"):
        errors.append("spec: missing deck 'title'")
    slides = spec.get("slides") or []
    if not slides:
        errors.append("spec: zero slides")

    for i, s in enumerate(slides, start=1):
        if not isinstance(s, dict):
            errors.append(f"slide {i}: not a JSON object")
            continue
        archetype = s.get("archetype")
        if archetype not in known_archetypes:
            errors.append(f"slide {i}: unknown archetype {archetype!r} -- not one of {sorted(known_archetypes)}")

    if errors:
        raise RedesignError(
            "the model's deck spec doesn't fit the kone-deck-generator skill's archetypes:\n  - "
            + "\n  - ".join(errors)
        )


def build_deck_via_skill(
    brief: str,
    out_path,
    target_slides: Optional[int] = None,
    model: str = DEFAULT_MODEL,
    effort: str = "high",
    notes: Optional[str] = None,
    api_key: Optional[str] = None,
    client=None,
) -> tuple[ComposeResult, RedesignResult]:
    """Plan a spec from `brief` and render it with the skill's own
    `build_deck`, wrapped into the same `(ComposeResult, RedesignResult)`
    shape every other `redesign_deck` path returns, so callers (the CLI,
    tests) don't need to know which renderer produced a given result.
    """
    creator = _load_creator()
    archetypes = _load_archetypes()
    spec, usage = call_claude_for_kone_spec(
        brief, target_slides=target_slides, model=model, effort=effort, notes=notes,
        api_key=api_key, client=client,
    )
    _validate_kone_spec(spec, set(archetypes.ARCHETYPES.keys()))

    # Nothing carries images into a from-a-brief build, so every photo
    # slot would otherwise render as a blank sand block.
    fill_empty_photo_slots(spec)
    creator.build_deck(spec, str(out_path))
    from deckguard.logo import repair_empty_logo_frames

    repair_empty_logo_frames(out_path)

    layouts_used = ["Cover F"] + [s["archetype"] for s in spec["slides"]] + ["Outro"]
    compose_result = ComposeResult(slide_count=len(layouts_used), layouts_used=layouts_used, manual_review=[])
    # No compose.py `Outline` exists for this path -- nothing downstream
    # reads `RedesignResult.outline`, so this is left None rather than
    # forcing a fake Outline into being just to satisfy the type.
    redesign_result = RedesignResult(outline=None, skipped=[], usage=usage)
    return compose_result, redesign_result


# --------------------------------------------------------------------------
# Archetype coexistence for redesign_deck's other two starting points (an
# existing deck's real content, with or without a brief to fill gaps): a
# second, additive pass over an outline `call_claude_for_outline` already
# planned onto compose.py's org-template layouts, offering an archetype
# for any slide where one's shape is a clearly better fit. Deliberately
# fail-closed throughout -- this is a quality enhancement on an already-
# working plan, never a required step, so nothing in here may turn a
# working redesign into a failed one.
# --------------------------------------------------------------------------

# Only these compose.py kinds have flexible-enough shape for an archetype
# to plausibly present the SAME content better -- cover/agenda/section/
# quote/statement/end already have one dedicated, well-fitted org-template
# layout each, so there's nothing for a second opinion to improve there.
_ARCHETYPE_CANDIDATE_KINDS = {"content", "stat", "timeline"}

_ARCHETYPE_OVERRIDE_RULES = """\
You are reviewing a deck outline that has already been planned onto generic org-template
layouts. Your ONLY job: for any slide whose CONTENT would read meaningfully better on one
of KONE's own slide archetypes below (a 2x2 matrix, an org chart, a lifecycle, a comparison
table, an icon row, etc.) than its current generic layout, say so -- and place that SAME
content, verbatim, into the archetype's slots (never invent new wording; only reorganize
what's already there). Most slides are already well served by a generic layout and should
NOT be overridden -- flag only a genuine, clearly-better fit, never for variety's own sake.

Each candidate reports an `image_count`: how many pictures that slide already carries. Those
images are attached automatically afterward -- never put an image path or filename in your
output, and never invent one. Just bear the count in mind when choosing: a slide WITH images
is a good fit for a picture-carrying archetype (its images fill those slots in order), and a
slide with image_count 0 should NOT be put on an archetype whose whole point is its pictures.

Output ONLY a JSON object, no prose, no markdown fences, of this shape:
{"overrides": [{"outline_index": <int>, "archetype": "<name>", ...that archetype's own content fields...}, ...]}
Omit any slide you are not overriding -- an empty "overrides" list is a completely valid,
expected answer when nothing here is a strong fit.
"""

# Which content key(s) each archetype's picture slots read, so a source
# deck's own images can be carried into an archetype-rendered slide (see
# `_inject_source_images`). DERIVED AT RUNTIME from each archetype's own
# `picture`/`image_band` regions in the skill's ARCHETYPES data -- so a
# skill update that adds/renames/removes archetypes (or their picture
# slots) flows through with no code change here, per the "update the
# skill and the tool just works" requirement. (An earlier version kept
# this as a hand-maintained map; it drifted the moment the skill grew.)
#
# ("single", key)            -> one image, at content[key]
# ("group", group_key, item) -> N images, one per content[group_key][i][item]
#
# `figure`-role regions are deliberately NOT picture slots (the role
# filter below excludes them): `archetypes.render` overwrites a figure
# key (chart/diagram) with the skill's own bundled art every time, so
# anything injected there would be silently discarded -- and a source
# photo is not a substitute for a real chart anyway. The FIGURES-key
# exclusion is belt-and-braces on top of the role filter.


def _inject_source_images(archetype_name: str, content: dict, image_blobs: list, tmpdir: str) -> dict:
    """Return a copy of `content` with the source deck's own images
    written to `tmpdir` and their PATHS placed in the archetype's picture
    slots -- `kone_engine._image()` takes a path, while everything
    upstream of here (`SlideProfile.images`, `_attach_source_images`)
    carries raw bytes, so this is the adapter between the two.

    A no-op (returns `content` unchanged) when the archetype has no
    picture slots or there are no images -- so it's always safe to call.
    The caller owns `tmpdir`'s lifetime; python-pptx reads each file into
    the package during `render`, so it only needs to outlive that call.
    """
    slot = _archetype_image_slots().get(archetype_name)
    if slot is None or not image_blobs:
        return content

    def _write(index: int, blob) -> Optional[str]:
        # Already a path (compose.py's own image entries can be either) --
        # pass it straight through rather than round-tripping through bytes.
        if not isinstance(blob, (bytes, bytearray)):
            return str(blob)
        path = os.path.join(tmpdir, f"img_{index}.img")
        try:
            with open(path, "wb") as fh:
                fh.write(blob)
        except Exception:
            return None
        return path

    out = dict(content)
    if slot[0] == "single":
        path = _write(0, image_blobs[0])
        if path:
            out[slot[1]] = path
        return out

    _, group_key, item_key = slot
    items = out.get(group_key)
    if not isinstance(items, list) or not items:
        return out  # the model gave this archetype no group items to hang pictures on
    new_items = []
    for i, item in enumerate(items):
        item = dict(item) if isinstance(item, dict) else {}
        if i < len(image_blobs):
            path = _write(i, image_blobs[i])
            if path:
                item[item_key] = path
        new_items.append(item)
    out[group_key] = new_items
    return out


def _run_archetype_override_call(
    candidates: list[dict], model: str, effort: str, api_key: Optional[str], client,
) -> list:
    """Shared plumbing for both override-selection entry points below:
    given already-prepared candidate dicts (each carrying its own
    caller-defined `outline_index` plus whatever verbatim content
    fields), calls the model and returns the raw, unvalidated
    "overrides" list from its response -- or `[]` on ANY failure (the
    skill not installed, an API error, a refusal, malformed JSON).
    del `effort`: accepted for interface parity, unused -- see
    `call_claude_for_kone_spec`'s own docstring for why (no
    `output_config` on this call).
    """
    del effort
    if not candidates:
        return []
    try:
        if client is None:
            key = api_key or os.environ.get("ANTHROPIC_API_KEY")
            if not key:
                return []
            local_client = anthropic.Anthropic(api_key=key)
        else:
            local_client = client

        system = _ARCHETYPE_OVERRIDE_RULES + "\n\nAvailable archetypes:\n\n" + _kone_archetype_guide()
        content = (
            "Already-planned slides to consider (their own outline_index and full "
            f"already-decided content):\n\n{json.dumps(candidates, indent=2, ensure_ascii=False)}"
        )
        messages = [{"role": "user", "content": content}]
        response = _stream_final_message(
            local_client, model=model, max_tokens=16000, thinking={"type": "adaptive"},
            system=system, messages=messages,
        )
        if response.stop_reason == "refusal":
            return []
        text = next((b.text for b in response.content if b.type == "text"), None)
        if not text:
            return []
        text = text.strip()
        if text.startswith("```"):
            text = text.strip("`")
            text = text[text.find("\n") + 1:] if "\n" in text else text
        return json.loads(text).get("overrides") or []
    except Exception:
        return []


def _validate_overrides(raw_overrides: list, candidate_indices: set) -> dict:
    """Keeps only entries naming a real archetype and an index the
    caller actually offered up -- anything else (a hallucinated
    archetype name, an index outside the candidate set) is silently
    dropped, not raised; see this module's fail-closed contract."""
    try:
        known = set(_load_archetypes().ARCHETYPES.keys())
    except RedesignError:
        return {}

    overrides: dict = {}
    for item in raw_overrides:
        if not isinstance(item, dict):
            continue
        idx = item.get("outline_index")
        archetype = item.get("archetype")
        if idx not in candidate_indices or archetype not in known:
            continue
        content = {k: v for k, v in item.items() if k not in ("outline_index", "archetype")}
        overrides[idx] = {"archetype": archetype, **content}
    return overrides


# --------------------------------------------------------------------------
# Structural archetype matching -- mapping an old slide onto a new format
# WITHOUT a model. See `match_archetypes` for why this exists.
# --------------------------------------------------------------------------

# Roles whose content this can't synthesise from an old slide's text.
# A table needs real headers/rows; a photo can be filled from the KONE
# library (see `fill_empty_photo_slots`) and a figure gets the skill's
# own bundled chart art, so neither disqualifies an archetype.
# Roles whose content this can't synthesise from an old slide's text.
# `figure` is here because the skill FORCES its own bundled chart art
# into that slot (`archetypes.FIGURES`), so structurally matching a
# figure-bearing archetype puts a chart on a slide that never had one --
# a campaign slide of screenshots came back with a pie chart on it. A
# model asked to choose deliberately may still pick these; a shape-fit
# match may not.
_UNFILLABLE_ROLES = {"table", "figure"}
# Structural roles the engine handles before it ever consults
# ROLE_STYLE. Everything else that resolves in ROLE_STYLE is, by
# definition, a text role -- which is what lets archetypes parsed from
# the HTML gallery or mined from a reference deck (whose roles are
# derived at runtime, e.g. `ref_i18_141414`) be understood here. A
# hardcoded list of the engine's own 23 role names reported every mined
# group as holding no text, so mined designs scored as capacity-0 and
# never surfaced.
_STRUCTURAL_ROLES = {
    "picture", "image_band", "figure", "panel", "panel_sand", "table", "axis", "icon",
}
_VALUE_ROLES = {"stat_value", "stat_value_md", "hero_value"}
_LIST_ROLES = {"bullets"}


def _is_text_role(role: Optional[str]) -> bool:
    if not role or role in _STRUCTURAL_ROLES or role in _LIST_ROLES:
        return False
    if role in _VALUE_ROLES:
        return True
    try:
        import importlib

        return role in importlib.import_module("kone_engine").ROLE_STYLE
    except Exception:  # noqa: BLE001
        return False

_signature_cache: Optional[list] = None


# --------------------------------------------------------------------------
# Two archetype vocabularies, and the gap between them
# --------------------------------------------------------------------------
#
# `kone-design`'s deck template ships an ARCHETYPES.md gallery naming 56
# slide archetypes in CAPS. `kone-deck-generator` -- the engine that
# actually renders .pptx -- defines 23. Only 17 names appear in both.
# So someone reads the gallery, asks for COVER_A_CUT4 / DIVIDER_D /
# END_LOGO / TITLE_TEXT_SPLIT, and the generator quietly builds
# something else, because those four are among the 39 gallery-only
# names. Reported by a user who got exactly that and had no way to see
# why.
#
# The engine can't grow 39 archetypes here. What it can do is stop
# substituting in silence: resolve the names that genuinely correspond,
# and say plainly what happened to the ones that don't.

# Gallery name -> engine archetype, only where the correspondence is
# real. Anything not here resolves to None and is reported, never
# guessed at.
_GALLERY_ALIASES = {
    "text_stats_picture": "text_stats_picture_right",
    "four_picture_cards": "four_point_value",
    "quote_panel": "quote_context",
    "quote_plain": "quote_context",
    "statement_full": "statement_links",
    "statement_two_col": "statement_links",
    "statement_three_col": "statement_links",
    "statement_on_picture": "statement_links",
    "statement_picture_note": "statement_links",
    "agenda_a_text": "agenda_contents",
    "agenda_a_bullets": "agenda_contents",
    "agenda_a_table": "agenda_contents",
    "divider_a": "image_section_divider",
    "divider_b": "image_section_divider",
    "divider_c": "image_section_divider",
    "divider_d": "image_section_divider",
    "divider_numbering": "image_section_divider",
    "picture_intro": "image_section_divider",
}

# Names that describe the deck's retained master slides rather than any
# archetype: `kone_deck_creator` always keeps the master's own cover and
# "Thank you", and never picks among their variants.
_MASTER_SLIDE_NAMES = {
    "cover_a_cut4", "cover_b_cut3", "cover_c_cut4_wide", "cover_d_cut3_wide",
    "cover_e_side", "cover_f_fullbleed", "end_logo",
}

_ARCHETYPE_TOKEN = re.compile(r"\b[A-Za-z][A-Za-z0-9]*(?:_[A-Za-z0-9]+)+\b")


def resolve_archetype_name(name: str) -> tuple:
    """Map a name someone typed onto an engine archetype.

    Returns `(archetype_or_None, status)` where status is one of
    "exact", "alias", "master_slide" or "unknown".
    """
    key = (name or "").strip().lower()
    if not key:
        return None, "unknown"
    try:
        known = {n.lower(): n for n in _load_archetypes().ARCHETYPES}
    except RedesignError:
        known = {}
    if key in known:
        return known[key], "exact"
    if key in _MASTER_SLIDE_NAMES:
        return None, "master_slide"
    alias = _GALLERY_ALIASES.get(key)
    if alias and alias in known.values():
        return alias, "alias"
    return None, "unknown"


def check_brief_archetypes(brief: str) -> dict:
    """Find archetype names in a brief and say what the engine can do
    with each -- so a request for a gallery-only archetype comes back as
    a stated substitution rather than a silent one.

    Only considers snake/CAPS tokens that look like archetype names AND
    resolve to something the vocabularies know about, so ordinary prose
    never trips it.
    """
    requested, seen = [], set()
    for token in _ARCHETYPE_TOKEN.findall(brief or ""):
        key = token.lower()
        if key in seen:
            continue
        seen.add(key)
        archetype, status = resolve_archetype_name(token)
        if status == "unknown" and key not in _GALLERY_NAMES:
            continue  # ordinary snake_case prose, not an archetype request
        requested.append({"requested": token, "archetype": archetype, "status": status})
    return {
        "exact": [r for r in requested if r["status"] == "exact"],
        "alias": [r for r in requested if r["status"] == "alias"],
        "master_slide": [r for r in requested if r["status"] == "master_slide"],
        "unknown": [r for r in requested if r["status"] == "unknown"],
    }


def _gallery_names() -> set:
    """Every archetype name the kone-design gallery documents, read from
    its own ARCHETYPES.md so this tracks the design system rather than a
    list copied here."""
    for base in (
        os.environ.get("KONE_DESIGN_DIR") and Path(os.environ["KONE_DESIGN_DIR"]),
        Path.home() / ".claude" / "skills" / "kone-design",
        Path(__file__).parent.parent / "assets" / "kone-design",
    ):
        if not base:
            continue
        doc = Path(base) / "templates" / "kone-deck" / "ARCHETYPES.md"
        try:
            text = doc.read_text(encoding="utf-8")
        except OSError:
            continue
        return {n.lower() for n in re.findall(r"\b[A-Z][A-Z0-9]*(?:_[A-Z0-9]+)+\b", text)}
    # No gallery reachable: fall back to the names this module knows are
    # gallery-only, so the check still catches the reported case.
    return set(_GALLERY_ALIASES) | _MASTER_SLIDE_NAMES


_GALLERY_NAMES = _gallery_names()


def archetype_signatures() -> list:
    """One machine-readable shape per archetype, derived at runtime from
    the skill's own region/group data -- so a skill update changes what
    can be matched with no code change here, the same contract the
    renderer and the previews already follow.

    Each signature says what the archetype can HOLD: how many repeating
    items, which content key each slot reads, and whether it needs a
    value, a picture, a chart or a table.
    """
    global _signature_cache
    if _signature_cache is not None:
        return _signature_cache
    try:
        mod = _load_archetypes()
    except RedesignError:
        return []

    catalog = _kone_catalog() or {}
    signatures = []
    for name, arch in mod.ARCHETYPES.items():
        regions = arch.get("regions", []) or []
        groups = arch.get("groups", []) or []
        roles = {r.get("role") for r in regions}
        for grp in groups:
            roles |= {r.get("role") for r in grp.get("regions", [])}

        group = None
        for grp in groups:
            item_roles = [(r.get("role"), r.get("content")) for r in grp.get("regions", [])]
            if any(_is_text_role(role) or role in _LIST_ROLES for role, _key in item_roles):
                group = {
                    "key": grp["content"],
                    "capacity": len(grp["origins"]),
                    "slots": [(role, key) for role, key in item_roles if key],
                }
                break

        signatures.append({
            "name": name,
            "regions": [(r.get("role"), r.get("content")) for r in regions if r.get("content")],
            "group": group,
            "capacity": group["capacity"] if group else 0,
            "needs_value": bool(roles & _VALUE_ROLES),
            "unfillable": bool(roles & _UNFILLABLE_ROLES),
            "keywords": [k.lower() for k in (catalog.get(name, {}).get("keywords") or [])],
            "purpose": (catalog.get(name, {}).get("purpose") or ""),
        })
    _signature_cache = signatures
    return signatures


def _looks_like_a_value(text: str) -> bool:
    """A short token carrying a number -- "91.2%", "739", "2x", "€1.4M"."""
    stripped = text.strip()
    return bool(stripped) and len(stripped) <= 12 and any(c.isdigit() for c in stripped)


def _chunk_slide(title: Optional[str], text_blocks: list) -> tuple:
    """Turn an old slide's extracted text into (title, chunks, values).

    A chunk is one repeatable unit -- a heading plus whatever supports
    it -- which is exactly the shape every archetype's group takes. A
    text block whose first line is short and whose remainder is longer
    reads as heading + body; a block of similar-length lines reads as a
    list of separate points.
    """
    blocks = [[line for line in block if str(line).strip()] for block in text_blocks or []]
    blocks = [b for b in blocks if b]

    resolved_title = title
    if not resolved_title and blocks:
        resolved_title = blocks[0][0]
        blocks[0] = blocks[0][1:]
        blocks = [b for b in blocks if b]

    chunks, values = [], []
    for block in blocks:
        head = str(block[0]).strip()
        rest = [str(line).strip() for line in block[1:]]
        if _looks_like_a_value(head):
            values.append(head)
        if len(block) == 1:
            chunks.append({"heading": head, "body": [], "text": head})
        elif len(head) <= 60:
            chunks.append({"heading": head, "body": rest, "text": " ".join([head] + rest)})
        else:
            for line in block:
                chunks.append({"heading": str(line).strip(), "body": [], "text": str(line).strip()})
    return resolved_title, chunks, values


def _fit_content(signature: dict, title: Optional[str], chunks: list, values: list, image_count: int) -> dict:
    """Pour the slide's own words into this archetype's slots."""
    content: dict = {"archetype": signature["name"]}
    spare = [c for c in chunks]

    for role, key in signature["regions"]:
        if role in ("title", "title_light", "statement", "quote"):
            content[key] = title or (spare[0]["heading"] if spare else "")
        elif role in _VALUE_ROLES:
            content[key] = values[0] if values else (spare[0]["heading"] if spare else "")
        elif role in ("eyebrow", "eyebrow_light", "label"):
            content[key] = (title or "").split(":")[0][:40] if title else ""
        elif _is_text_role(role):
            source = spare.pop() if spare else None
            content[key] = source["text"] if source else ""

    group = signature["group"]
    if group:
        capacity = group["capacity"]
        items = []
        for i, chunk in enumerate(spare[:capacity], start=1):
            item = {}
            for role, key in group["slots"]:
                if role == "number":
                    item[key] = f"{i:02d}"
                elif role in _LIST_ROLES:
                    item[key] = chunk["body"] or [chunk["heading"]]
                elif role in _VALUE_ROLES:
                    item[key] = values[i - 1] if i - 1 < len(values) else chunk["heading"]
                elif _is_text_role(role):
                    item[key] = chunk["heading"] if key not in item else " ".join(chunk["body"])
            items.append(item)
        content[group["key"]] = items
    del image_count
    return content


def match_archetypes(
    title: Optional[str], text_blocks: list, image_count: int = 0, limit: int = 3,
    prefer: Optional[set] = None,
) -> list:
    """Rank archetypes that could hold this slide, WITHOUT a model.

    Choosing a new format for an old slide reads like a job only a model
    can do, and the part that actually needs judgement is narrower than
    that: which of several plausible formats reads best, and how the copy
    should be reworded. Deciding what CAN hold the content is a fitting
    problem -- the archetype library states its own capacity (a title
    plus three [label, value, caption] groups; a title plus five
    [number, heading] rows), and an old slide's content shape is
    extractable. So this matches shape to shape.

    It exists because the model was the ONLY route to an archetype: on a
    server with no API key every dense slide was offered nothing but
    "keep", which is the tool switching off rather than degrading. With
    this, a slide always has somewhere to go, and the model -- when
    present -- upgrades the choice and the wording instead of gating it.

    Returns up to `limit` candidates, best first, each:
        {"archetype", "content", "score", "capacity", "dropped"}
    where `dropped` is how many content chunks the archetype cannot
    hold -- the honest cost of the mapping, for the review page to show.
    """
    resolved_title, chunks, values = _chunk_slide(title, text_blocks)
    haystack = " ".join([resolved_title or ""] + [c["text"] for c in chunks]).lower()

    scored = []
    for signature in archetype_signatures():
        if signature["unfillable"]:
            continue
        if signature["needs_value"] and not values:
            continue

        capacity = signature["capacity"]
        if capacity:
            # closest capacity wins; overfilling costs more than a spare slot
            if len(chunks) >= capacity:
                score = 100.0 - (len(chunks) - capacity) * 6.0
            else:
                score = 100.0 - (capacity - len(chunks)) * 10.0
        else:
            # no repeating group: only a good home for a small slide
            score = 100.0 - abs(len(chunks) - 2) * 18.0

        score += 8.0 * sum(1 for kw in signature["keywords"] if kw in haystack)
        if values and signature["needs_value"]:
            score += 12.0
        # Designs mined from the user's own reference deck outrank the
        # built-in library when one was supplied: "make it look like
        # that deck" means that deck's layouts, not a lookalike. But the
        # bonus is earned, not flat -- a reference design that would
        # drop half the slide's content must not beat a built-in one
        # that holds all of it.
        if prefer and signature["name"] in prefer:
            fits = capacity and len(chunks) <= capacity
            score += 40.0 if fits else 12.0

        dropped = max(0, len(chunks) - capacity) if capacity else max(0, len(chunks) - 2)
        scored.append({
            "archetype": signature["name"],
            "content": _fit_content(signature, resolved_title, chunks, values, image_count),
            "score": round(score, 1),
            "capacity": capacity,
            "dropped": dropped,
        })

    scored.sort(key=lambda c: (-c["score"], c["dropped"], c["archetype"]))
    return scored[:limit]


def archetype_suggestions_available(api_key: Optional[str] = None, client=None) -> bool:
    """Whether an archetype-suggestion call can be made at all.

    The selection helpers below fail CLOSED -- no key, an API error and
    a model that legitimately suggests nothing all return `{}`. That is
    the right contract for them and the wrong one for the review page,
    which told a user "archetype suggestions ran" on a server with no
    `ANTHROPIC_API_KEY` and then offered nothing but "keep" on ten of
    twelve slides, with no hint why. This separates "couldn't ask" from
    "asked and got nothing", so the page can say which happened.
    """
    if client is not None:
        return True
    return bool(api_key or os.environ.get("ANTHROPIC_API_KEY"))


def select_archetype_overrides(
    outline_items: list[dict],
    model: str = DEFAULT_MODEL,
    effort: str = "high",
    api_key: Optional[str] = None,
    client=None,
) -> dict:
    """For `redesign_deck`'s real-content path: `outline_items` is
    `call_claude_for_outline`'s own already-planned compose.py outline
    (list of dicts with a `kind` key). Returns
    `{outline_index: {"archetype": <name>, ...content...}}` for just
    the slides worth overriding. ANY failure here -- the skill not
    installed, an API error, malformed JSON, an invalid archetype name
    -- just means no overrides, never a raised error: see this module's
    own docstring for why that's the right contract for a pass that's
    purely additive on top of an already-good plan.
    """
    candidate_indices = [
        i for i, item in enumerate(outline_items)
        if item.get("kind") in _ARCHETYPE_CANDIDATE_KINDS
    ]
    candidates = []
    for i in candidate_indices:
        item = outline_items[i]
        # `images` carries raw bytes by this point (`_attach_source_images`
        # runs before this call so the count is accurate) -- not
        # JSON-serializable, and not something the model should see or
        # echo anyway. Report only how many there are; the blobs
        # themselves are placed later by `_inject_source_images`.
        content = {k: v for k, v in item.items() if k != "images"}
        candidates.append({
            "outline_index": i,
            "kind": item.get("kind"),
            "image_count": len(item.get("images") or []),
            "content": content,
        })
    raw_overrides = _run_archetype_override_call(candidates, model, effort, api_key, client)
    return _validate_overrides(raw_overrides, set(candidate_indices))


def select_archetype_overrides_for_rebrand(
    slide_profiles_by_index: dict,
    model: str = DEFAULT_MODEL,
    effort: str = "high",
    api_key: Optional[str] = None,
    client=None,
) -> dict:
    """For `apply_rebrand`'s (mode='brand', --review only) path:
    `slide_profiles_by_index` is `{1-based slide index: SlideProfile}`
    for slides `apply_rebrand` already accepted and rebuilt onto an
    ordinary org-template layout (cover/end swap positions and
    reference-layout carryovers excluded by the caller -- see
    `_rebrand_deck`). Their title/text_blocks/images should be
    re-extracted from the FINISHED, already-brand-fixed output deck, so
    this sees exactly the verbatim content a human reviewing the result
    would see. Same `{outline_index: {"archetype": ..., ...}}` return
    shape and fail-closed contract as `select_archetype_overrides`.
    """
    candidates = []
    for idx, profile in slide_profiles_by_index.items():
        candidates.append({
            "outline_index": idx,
            "title": profile.title,
            "text_blocks": [[text for _level, text in block] for block in profile.text_blocks],
            "image_count": len(profile.images),
        })
    raw_overrides = _run_archetype_override_call(candidates, model, effort, api_key, client)
    return _validate_overrides(raw_overrides, set(slide_profiles_by_index.keys()))


def apply_archetype_overrides_to_deck(deck_path, overrides: dict, images_by_index: Optional[dict] = None) -> dict:
    """For `apply_rebrand`'s (mode='brand', --review only) path: swaps
    each slide named in `overrides` (`{1-based slide index:
    {"archetype": <name>, ...content...}}`) for an archetype-rendered
    one, IN PLACE at that same position, in the already-saved deck at
    `deck_path` (overwritten with the result). Every other slide is
    untouched.

    `images_by_index` (`{1-based slide index: list[bytes]}`, optional):
    that slide's OWN pictures, carried into the archetype's picture
    slots (see `_inject_source_images`). Without it a picture-carrying
    archetype still renders -- its picture slots just fall back to
    `kone_engine`'s sand placeholder -- so this stays optional.

    Mirrors `retemplate._rebuild_accepted_slides`'s own two-phase
    add-then-move-then-delete technique (by OPC partname, not raw list
    index, so multiple simultaneous swaps can't collide as earlier ones
    shift slide positions underneath later ones) rather than modifying
    that function directly -- this operates on an ALREADY-COMPLETE,
    already-brand-fixed deck as an isolated final step, so there's no
    need to thread a second rendering engine through
    `_rebuild_accepted_slides`'s own delicate bookkeeping.

    Returns `{slide index: archetype name}` for what was actually
    swapped, so the caller can fold it into its own layout-used
    reporting.
    """
    if not overrides:
        return {}

    from deckguard.legacy.slide_import import _delete_slide, _move_slide

    archetypes = _load_archetypes()

    prs = Presentation(str(deck_path))
    blank_layout = next(l for l in prs.slide_layouts if l.name.strip().lower() == "blank")
    original_partnames = [str(s.part.partname).lstrip("/") for s in prs.slides]

    images_by_index = images_by_index or {}
    new_partname_by_old: dict[str, str] = {}
    layout_by_index: dict = {}
    with tempfile.TemporaryDirectory() as tmpdir:
        for idx, override in overrides.items():
            old_partname = original_partnames[idx - 1]
            archetype_name = override["archetype"]
            content = {k: v for k, v in override.items() if k != "archetype"}
            content = _inject_source_images(
                archetype_name, content, list(images_by_index.get(idx) or []), tmpdir,
            )
            new_slide = prs.slides.add_slide(blank_layout)
            archetypes.render(new_slide, archetype_name, content)
            new_partname_by_old[old_partname] = str(new_slide.part.partname).lstrip("/")
            layout_by_index[idx] = archetype_name

    for idx in overrides:
        old_partname = original_partnames[idx - 1]
        new_partname = new_partname_by_old[old_partname]
        old_pos = next(i for i, s in enumerate(prs.slides) if str(s.part.partname).lstrip("/") == old_partname)
        new_pos = next(i for i, s in enumerate(prs.slides) if str(s.part.partname).lstrip("/") == new_partname)
        _move_slide(prs, new_pos, old_pos)
        _delete_slide(prs, old_pos + 1)

    prs.save(deck_path)
    return layout_by_index


def build_deck_with_archetypes(
    outline, out_path, template_path=None, rules_config: Optional[dict] = None,
    overrides: Optional[dict] = None,
) -> ComposeResult:
    """Same contract and (when `overrides` is empty) same behavior as
    `compose.build_deck`'s fresh-start path -- `redesign_deck` never
    passes `existing_deck_path`, so that's the only branch this needs
    to cover. When `overrides` names slides (from
    `select_archetype_overrides`), those render through the archetype
    engine instead of an org-template layout; every other slide still
    goes through compose.py's own `_select_layout`/`_populate`,
    unchanged.

    Ordering matters here: `fix_deck` runs FIRST, over only the
    org-template slides, exactly as `compose.build_deck` already does --
    archetype slides don't exist in the file yet at that point, so
    there's nothing for it to (wrongly) "fix" on them. They're rendered
    and spliced into the deck's slide order AFTER, via a raw
    `_sldIdLst` reorder (the same technique the skill's own
    `kone_deck_creator.build_deck`/`archetypes.build_gallery` already
    use to reassemble intro/body/outro). This is deliberate: the
    archetype engine's hardcoded role colors (e.g. `kone_engine.py`'s
    muted caption grey `#727272`) aren't all in `brand_rules.yaml`'s
    approved-colors list -- they're genuinely on-brand (KONE's own
    `--ink-muted` token), just not ones this project's rules currently
    list, so a `fix_deck` pass over them would flag correct styling as
    a violation instead of leaving it alone.
    """
    overrides = overrides or {}
    if not overrides:
        from deckguard.legacy.compose import build_deck as compose_build_deck

        return compose_build_deck(outline, out_path, template_path=template_path, rules_config=rules_config)

    from deckguard.legacy.compose import _populate, _select_layout
    from deckguard.config import default_config_path, load_config
    from deckguard.legacy.fixer import fix_deck
    from deckguard.legacy.slide_import import _delete_slide
    from deckguard.legacy.slide_import import default_template_path as _default_template_path

    creator = _load_creator()
    archetypes = _load_archetypes()

    template_path = Path(template_path) if template_path else _default_template_path()
    template_prs = Presentation(str(template_path))
    layouts_by_name = {layout.name: layout for master in template_prs.slide_masters for layout in master.slide_layouts}
    layout_profile_cache: dict = {}
    config = rules_config if rules_config is not None else load_config(default_config_path())

    prs = Presentation(str(template_path))
    for i in range(len(prs.slides) - 1, -1, -1):
        _delete_slide(prs, i)
    own_layouts_by_name = {layout.name: layout for master in prs.slide_masters for layout in master.slide_layouts}
    lst = prs.slides._sldIdLst

    layout_by_index: dict[int, str] = {}
    sldid_by_index: dict[int, object] = {}

    non_archetype = [(i, spec) for i, spec in enumerate(outline.slides) if i not in overrides]
    for i, spec in non_archetype:
        layout_name = _select_layout(spec, layouts_by_name, layout_profile_cache)
        new_slide = prs.slides.add_slide(own_layouts_by_name[layout_name])
        _populate(new_slide, own_layouts_by_name[layout_name], spec)
        layout_by_index[i] = layout_name
        sldid_by_index[i] = list(lst)[-1]

    # Applies every inherited-color/theme fix in memory but doesn't save --
    # archetype slides are added below, after this, then the whole deck is
    # saved once at the end.
    fix_report = fix_deck(prs, config, source_path=str(out_path), output_path=None, dry_run=True)
    manual_review = fix_report.manual_review

    blank_layout = next(l for l in prs.slide_layouts if l.name.strip().lower() == "blank")
    # Temp dir outlives every render() below -- kone_engine._image() reads
    # each path with PIL and python-pptx copies the bytes into the package
    # during the call, so nothing needs these files afterward.
    with tempfile.TemporaryDirectory() as tmpdir:
        for i, override in overrides.items():
            archetype_name = override["archetype"]
            content = {k: v for k, v in override.items() if k != "archetype"}
            content = _inject_source_images(
                archetype_name, content, list(outline.slides[i].images or []), tmpdir,
            )
            new_slide = prs.slides.add_slide(blank_layout)
            archetypes.render(new_slide, archetype_name, content)
            layout_by_index[i] = archetype_name
            sldid_by_index[i] = list(lst)[-1]

    for el in list(lst):
        lst.remove(el)
    for i in range(len(outline.slides)):
        lst.append(sldid_by_index[i])

    prs.save(out_path)

    layouts_used = [layout_by_index[i] for i in range(len(outline.slides))]
    return ComposeResult(slide_count=len(outline.slides), layouts_used=layouts_used, manual_review=manual_review)
