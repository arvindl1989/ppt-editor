"""Mine a reference .pptx into archetypes the engine can render.

Everything deckguard did with a reference deck until now was a form of
PATCHING: learn its colours, carry its layouts across, transplant its
per-shape styles onto the old deck's shapes. All of that nudges old
geometry toward new-looking values, and a recoloured 2011 slide is still
a 2011 slide. The pair of decks that started this work proves the point
-- the human who made the "after" deck did not recolour the "before"
one, they took what each slide SAID and re-laid it out in the new
system.

So: read the reference's own designs out as archetypes, register them
beside the built-in ones, and let old content be RE-RENDERED through
them. "Make this look like that deck" stops meaning "approximate its
colours" and starts meaning "pour this content into its actual layouts".

This is the same trick `gallery.py` plays on kone-design's HTML, and a
.pptx is the easier source: exact boxes in EMU, explicit fonts and
fills, no CSS to interpret. A slide is already
`{regions: [{role, content, box}], groups: [...]}` -- it just needs
reading out.

Three things make the result usable rather than a pile of one-off
slides:

- REPEATING GROUPS. Seven identically-sized pills along one row are not
  seven regions, they are a group with seven origins. Detecting that is
  what makes a mined archetype able to hold a different number of items
  than the slide it came from.
- DEDUPING. A deck uses each design several times. Slides are collapsed
  by structural signature so a 40-slide reference yields a handful of
  archetypes, not 40.
- REAL SAMPLE CONTENT. The reference's own words become the archetype's
  `SAMPLES` entry, so the planning prompt shows the model a worked
  example in the house style.

Skipped rather than approximated: tables, charts, embedded objects and
media (no way to synthesise their content), and anything whose geometry
is missing.
"""

from __future__ import annotations

import re
from collections import defaultdict
from typing import Optional

EMU_PER_PX = 9525.0
PT_TO_PX = 4.0 / 3.0  # a 13.333in slide is 960pt and 1280px wide

# A shape must be at least this big to be worth mining: below it the
# shape is a bullet glyph or a rule, not a content slot.
MIN_REGION_PX = 12.0

# Two shapes belong to the same repeating group when their sizes match
# this closely (px) and they share a row or column.
GROUP_SIZE_TOL = 4.0
GROUP_ALIGN_TOL = 6.0

# Shape types whose content cannot be synthesised, so an archetype
# containing one could never be filled from other content.
_UNMINEABLE = {"TABLE", "CHART", "EMBEDDED_OLE_OBJECT", "OLE_OBJECT", "MEDIA"}

# Above this many free-standing regions with no repeating group, a slide
# is a one-off composition rather than a design worth reusing.
MAX_SINGLETON_REGIONS = 12

_TITLE_PH = {"TITLE", "CENTER_TITLE"}


def _px(emu) -> Optional[float]:
    return None if emu is None else emu / EMU_PER_PX


def _shape_box(shape) -> Optional[list]:
    box = [_px(shape.left), _px(shape.top), _px(shape.width), _px(shape.height)]
    if any(v is None for v in box):
        return None
    if box[2] < MIN_REGION_PX or box[3] < MIN_REGION_PX:
        return None
    return box


def _fill_hex(shape) -> Optional[str]:
    try:
        if shape.fill.type is None or shape.fill.type != 1:  # 1 == solid
            return None
        return str(shape.fill.fore_color.rgb).upper()
    except Exception:  # noqa: BLE001 -- theme/gradient/picture fills
        return None


def _placeholder_name(shape) -> Optional[str]:
    try:
        if not shape.is_placeholder or shape.placeholder_format.type is None:
            return None
        return shape.placeholder_format.type.name
    except Exception:  # noqa: BLE001
        return None


def _first_run(shape):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip():
                return run
    return None


def _inherited_size_px(shape) -> float:
    """A run with no explicit size inherits from the layout/master, which
    python-pptx will not resolve. Placeholder type is a good enough
    stand-in, and only affects the derived text STYLE, never geometry."""
    name = _placeholder_name(shape)
    if name in _TITLE_PH:
        return 40 * PT_TO_PX
    if name == "SUBTITLE":
        return 24 * PT_TO_PX
    return 18 * PT_TO_PX


def _text_style(shape) -> Optional[tuple]:
    """(role_key, (font, px, hex, caps, bold, fit)) for a text shape."""
    run = _first_run(shape)
    if run is None:
        return None
    size_px = run.font.size.pt * PT_TO_PX if run.font.size else _inherited_size_px(shape)
    font = run.font.name or "Inter"
    bold = bool(run.font.bold)
    try:
        colour = str(run.font.color.rgb).upper() if run.font.color.type is not None else "141414"
    except Exception:  # noqa: BLE001 -- theme colour with no resolvable rgb
        colour = "141414"
    text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs)
    caps = bool(text.strip()) and text.upper() == text and any(c.isalpha() for c in text)
    if "KONE Information" in font:
        family = "KONE Information"
    elif bold:
        family = "Inter SemiBold"
    else:
        family = "Inter"
    role = f"ref_{'k' if family.startswith('KONE') else 'i'}{int(size_px)}_{colour}{'_c' if caps else ''}"
    return role, (family, round(size_px, 1), colour, caps, bold, False)


def _slide_text(shape) -> str:
    return "\n".join(
        p.text.strip() for p in shape.text_frame.paragraphs if p.text and p.text.strip()
    )


def _candidates(slide) -> tuple:
    """(items, unmineable) for one slide. An item is a dict describing
    one top-level shape in mineable terms."""
    items, unmineable = [], False
    for shape in slide.shapes:
        type_name = shape.shape_type.name if shape.shape_type is not None else ""
        if type_name in _UNMINEABLE:
            unmineable = True
            continue
        box = _shape_box(shape)
        if box is None:
            continue
        if type_name == "PICTURE" or getattr(shape, "image", None) is not None:
            items.append({"kind": "picture", "box": box})
            continue
        if type_name == "GROUP":
            # A group is one visual unit; its internals are that unit's
            # business. Mined as a picture-shaped slot only when it
            # carries no text worth keeping.
            items.append({"kind": "group", "box": box})
            continue
        style = _text_style(shape) if getattr(shape, "has_text_frame", False) else None
        fill = _fill_hex(shape)
        if style is None:
            if fill:
                items.append({"kind": "fill", "box": box, "hex": fill})
            continue
        items.append({
            "kind": "text", "box": box, "role": style[0], "style": style[1],
            "text": _slide_text(shape), "fill": fill,
        })
    return items, unmineable


def _cluster_repeats(items: list) -> tuple:
    """Split items into (singles, groups).

    A group is three or more items of the same kind and near-identical
    size sharing a row or a column. That is what lets a mined archetype
    hold a different number of items than the slide it came from -- the
    difference between a reusable design and a screenshot.
    """
    buckets = defaultdict(list)
    for item in items:
        key = (
            item["kind"],
            item.get("role"),
            round(item["box"][2] / GROUP_SIZE_TOL),
            round(item["box"][3] / GROUP_SIZE_TOL),
        )
        buckets[key].append(item)

    grouped_ids, groups = set(), []
    for members in buckets.values():
        if len(members) < 3:
            continue
        rows = {round(m["box"][1] / GROUP_ALIGN_TOL) for m in members}
        cols = {round(m["box"][0] / GROUP_ALIGN_TOL) for m in members}
        if len(rows) > 1 and len(cols) > 1 and len(members) < 4:
            continue  # scattered, not a row/column/grid
        ordered = sorted(members, key=lambda m: (m["box"][1], m["box"][0]))
        first = ordered[0]
        groups.append({
            "origins": [[m["box"][0], m["box"][1]] for m in ordered],
            "prototype": first,
            "members": ordered,
            "count": len(ordered),
        })
        grouped_ids.update(id(m) for m in ordered)

    singles = [i for i in items if id(i) not in grouped_ids]
    return singles, groups


def _slot_name(item: dict, index: int, used: set) -> str:
    """Name a content slot from what the shape evidently is."""
    if item["kind"] in ("picture", "group"):
        base = "image"
    else:
        size = item["style"][1]
        text = (item.get("text") or "").strip()
        if size >= 32:
            base = "title"
        elif item["style"][3]:  # ALL CAPS
            base = "eyebrow"
        elif len(text) > 120 or "\n" in text:
            base = "body"
        else:
            base = f"text{index}"
    name = base
    n = 2
    while name in used:
        name = f"{base}{n}"
        n += 1
    used.add(name)
    return name


def _archetype_from_slide(slide, background: str = "FFFFFF") -> Optional[dict]:
    items, unmineable = _candidates(slide)
    if unmineable or not items:
        return None

    singles, groups = _cluster_repeats(items)
    regions, role_styles, sample, chrome = [], {}, {}, []
    used: set = set()

    for index, item in enumerate(sorted(singles, key=lambda i: (i["box"][1], i["box"][0]))):
        if item["kind"] == "fill":
            chrome.append({"kind": "fill", "box": item["box"], "hex": item["hex"]})
            continue
        slot = _slot_name(item, index, used)
        if item["kind"] in ("picture", "group"):
            regions.append({"role": "picture", "content": slot, "box": item["box"]})
            sample[slot] = ""
            continue
        role_styles[item["role"]] = item["style"]
        regions.append({"role": item["role"], "content": slot, "box": item["box"]})
        sample[slot] = item["text"]

    group_specs = []
    for gi, group in enumerate(groups):
        proto = group["prototype"]
        key = "items" if gi == 0 else f"items{gi + 1}"
        box = proto["box"]
        if proto["kind"] == "fill":
            # A repeating band of unlabelled colour blocks is decoration,
            # not a content slot -- drawn, never filled in. Each block
            # keeps its OWN colour: the reference's category band runs
            # blue/amber/pink/yellow/green/orange across seven cells,
            # and reusing the first cell's fill for all seven turned it
            # into a solid blue bar.
            chrome.extend(
                {"kind": "fill", "box": [m["box"][0], m["box"][1], m["box"][2], m["box"][3]],
                 "hex": m["hex"]}
                for m in group["members"]
            )
            continue
        if proto["kind"] in ("picture", "group"):
            member = {"role": "picture", "content": "image", "box": [0, 0, box[2], box[3]]}
            items_sample = [{"image": ""} for _ in range(group["count"])]
        else:
            role_styles[proto["role"]] = proto["style"]
            member = {"role": proto["role"], "content": "text", "box": [0, 0, box[2], box[3]]}
            items_sample = [{"text": (m.get("text") or "")} for m in group["members"]]
        # Text items can carry their own fill too (a coloured pill with a
        # label on it), and it is per-item for the same reason.
        chrome.extend(
            {"kind": "fill", "box": list(m["box"]), "hex": m["fill"]}
            for m in group["members"] if m.get("fill")
        )
        group_specs.append({"content": key, "origins": group["origins"], "regions": [member]})
        sample[key] = items_sample

    if not regions and not group_specs:
        return None
    # A slide of twenty loose text boxes with no repeating structure is
    # a one-off composition, not a reusable design. Mining it produces a
    # twenty-slot archetype nothing else will ever fit.
    if not group_specs and len(regions) > MAX_SINGLETON_REGIONS:
        return None
    return {
        "archetype": {"background": background, "regions": regions,
                      "groups": group_specs, "chrome": chrome},
        "role_styles": role_styles,
        "sample": sample,
    }


def _signature(archetype: dict) -> tuple:
    """What makes two slides "the same design": the arrangement, not the
    words. Boxes are rounded to a 12px tolerance so hand-nudged copies
    of one layout collapse together."""
    # Height is content-driven -- the same design holding a longer
    # paragraph grows its text box -- so matching on it split four
    # copies of one layout into four archetypes. Position and width
    # are what the design actually fixes.
    def anchor(box):
        return round(box[0] / 24), round(box[1] / 24), round(box[2] / 24)

    regions = tuple(sorted((r["role"], anchor(r["box"])) for r in archetype["regions"]))
    groups = tuple(sorted((len(g["origins"]), anchor(g["regions"][0]["box"]))
                          for g in archetype["groups"]))
    return regions, groups


def _slug(text: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "_", (text or "").lower()).strip("_")
    return slug or "slide"


def mine_reference(deck_path, prefix: str = "ref") -> dict:
    """Read a reference deck's designs out as archetypes.

    Returns `{archetypes, samples, role_styles, sources}` where `sources`
    maps each archetype name to the reference slide numbers it was mined
    from -- so a review page can say where a design came from.
    """
    from pptx import Presentation

    prs = Presentation(str(deck_path))
    result = {"archetypes": {}, "samples": {}, "role_styles": {}, "sources": {}}
    by_signature: dict = {}

    for number, slide in enumerate(prs.slides, start=1):
        try:
            mined = _archetype_from_slide(slide)
        except Exception:  # noqa: BLE001 -- one odd slide must not stop the mine
            continue
        if mined is None:
            continue
        signature = _signature(mined["archetype"])
        if signature in by_signature:
            result["sources"][by_signature[signature]].append(number)
            continue

        layout = _slug(getattr(slide.slide_layout, "name", ""))
        shape_count = len(mined["archetype"]["regions"])
        group_count = sum(len(g["origins"]) for g in mined["archetype"]["groups"])
        name = f"{prefix}_{layout}"
        if group_count:
            name += f"_{group_count}up"
        elif shape_count:
            name += f"_{shape_count}"
        base, n = name, 2
        while name in result["archetypes"]:
            name = f"{base}_{n}"
            n += 1

        by_signature[signature] = name
        result["archetypes"][name] = mined["archetype"]
        result["samples"][name] = mined["sample"]
        result["role_styles"].update(mined["role_styles"])
        result["sources"][name] = [number]
    return result


def install_reference(archetypes_module, deck_path, prefix: str = "ref") -> dict:
    """Mine `deck_path` and merge the result into the loaded archetype
    registry, exactly as the gallery archetypes are merged -- so the
    matcher, previews, picture-slot detection and the planning prompt
    all gain the reference's own designs with no further plumbing.

    Returns the mine result. Never raises: a reference that yields
    nothing simply adds nothing.
    """
    import importlib

    try:
        mined = mine_reference(deck_path, prefix=prefix)
    except Exception:  # noqa: BLE001 -- additive, never load-bearing
        return {"archetypes": {}, "samples": {}, "role_styles": {}, "sources": {}}
    if not mined["archetypes"]:
        return mined

    engine = importlib.import_module("kone_engine")
    engine.ROLE_STYLE.update({
        role: (font, px, engine._hex(colour), caps, bold, fit)
        for role, (font, px, colour, caps, bold, fit) in mined["role_styles"].items()
    })
    archetypes_module.ARCHETYPES.update(mined["archetypes"])
    from deckguard.skill_bridge import invalidate_archetype_caches

    invalidate_archetype_caches()  # the matcher's derived views predate these
    if hasattr(archetypes_module, "SAMPLES"):
        archetypes_module.SAMPLES.update(mined["samples"])
    if hasattr(archetypes_module, "BG"):
        for name, arch in mined["archetypes"].items():
            archetypes_module.BG[name] = arch.get("background")
    return mined


__all__ = ["install_reference", "mine_reference"]
