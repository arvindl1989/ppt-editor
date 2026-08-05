"""The unified Transform pipeline: plan -> human review -> execute -> audit.

This is the consolidation seam over deckguard's three primitives --
analyze (`propose_rebrand`), rebuild (org layout via `apply_rebrand` /
archetype via `apply_archetype_overrides_to_deck`), and patch
(`fix_deck`, already inside `apply_rebrand`) -- exposed as ONE flow with
a human decision point in the middle, instead of four separate blind
upload->download tabs. Nothing here re-implements those primitives: the
plan IS `propose_rebrand`'s output (so plan and execution can never
disagree), the AI archetype suggestions ARE the same fail-closed
`select_archetype_overrides_for_rebrand` pass brand-mode review uses,
and execution IS `apply_rebrand(accepted_indexes=...)` -- the
guided-review hook that existed in the engine all along but never had a
caller.

Two starting points, one result shape:

- `plan_transform(deck_path, ...)` -> per-slide plan for an EXISTING
  deck (with optional reference deck driving layout carryover exactly
  as Learn/brand-mode always did). Each slide gets a default action a
  human can override: keep / rebuild on a named org layout (or
  reference layout) / rebuild as a suggested archetype.
- `plan_transform_from_brief(brief, ...)` -> per-slide plan for a NEW
  deck (the old Create/Redesign brief-only path): the archetype spec
  the model planned, shown slide by slide for approval before anything
  renders.

`execute_transform` / `execute_transform_from_brief` then run exactly
the approved subset, and `audit_transform_result` reports the result's
brand-rule audit -- with archetype-rendered slides excluded from the
violation list, because they are compliant by construction and
`rules_engine`'s generic text rules would false-positive on their
deliberate styling (the known `#727272` caption case) -- plus, when a
reference deck was given, a similarity report against it.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation

from deckguard.config import default_config_path, load_config
from deckguard.inventory import build_inventory, iter_shapes_recursive
from deckguard.redesign import DEFAULT_MODEL
from deckguard.retemplate import (
    SlideProfile,
    _extract_slide_content,
    apply_rebrand,
    propose_rebrand,
)
from deckguard.rules_engine import audit_deck, sort_violations, summarize


@dataclass
class SlidePlan:
    index: int  # 1-based
    default_action: str  # "keep" | "rebuild" | "reference_layout" | "archetype" | "new"
    layout_name: Optional[str] = None  # org/reference layout the rebuild would use
    archetype: Optional[dict] = None  # {"archetype": name, ...content...} when the AI suggested one
    reason: Optional[str] = None  # why a slide is keep-only (ineligible)
    title_preview: Optional[str] = None
    body_preview: Optional[str] = None
    image_count: int = 0
    # Verbatim content, for rendering a faithful proposed-slide preview
    # (org-layout previews place these into the layout's real
    # placeholder boxes). None for slides that were never extracted
    # (ineligible/keep-only ones).
    title: Optional[str] = None
    text_blocks: list = field(default_factory=list)  # list[list[str]]
    # Where the archetype came from ("model" | "structural") and what it
    # costs: how many content chunks it can't hold, out of how many.
    archetype_source: Optional[str] = None
    archetype_dropped: int = 0
    content_chunks: int = 0


@dataclass
class TransformPlan:
    slides: list  # list[SlidePlan]
    ai_suggestions_ran: bool = False
    deck_title: Optional[str] = None  # brief-only plans: the planned cover title
    # What became of any archetype the brief asked for BY NAME -- see
    # skill_bridge.check_brief_archetypes. Empty for existing-deck plans.
    archetype_requests: dict = field(default_factory=dict)
    # Archetypes mined from the uploaded reference deck -> the reference
    # slide numbers each design was read from.
    reference_designs: dict = field(default_factory=dict)


def _readable_text(slide) -> tuple:
    """Every non-empty line on a slide, groups included, as
    (title, text_blocks) -- enough to offer an archetype for a slide no
    org-template layout can accept. Deliberately dumber than
    `_extract_slide_content`: no placeholder-role reasoning, no
    eligibility judgement, just the words."""
    blocks: list = []
    title = None
    for shape in slide.shapes:
        for candidate in _iter_shapes(shape):
            if not getattr(candidate, "has_text_frame", False):
                continue
            lines = [
                para.text.strip()
                for para in candidate.text_frame.paragraphs
                if para.text and para.text.strip()
            ]
            if not lines:
                continue
            if title is None and candidate == slide.shapes.title:
                title = " ".join(lines)
                continue
            blocks.append(lines)
    if title is None and blocks:
        title = blocks[0][0]
    return title, blocks


def _iter_shapes(shape):
    """Depth-first over a live python-pptx shape and any group children."""
    yield shape
    try:
        if shape.shape_type is not None and shape.shape_type.name == "GROUP":
            for child in shape.shapes:
                yield from _iter_shapes(child)
    except (AttributeError, ValueError):
        return


_CHROME_PLACEHOLDERS = ("DATE", "FOOTER", "SLIDE_NUMBER")


def restore_footer_chrome(out_path, source_path) -> int:
    """Put the date, footer and page number back on rebuilt slides.

    python-pptx deliberately does NOT clone a layout's date, footer and
    slide-number placeholders when a slide is added -- OOXML calls them
    "latent" and PowerPoint materialises them itself. The result here is
    that every slide the transform REBUILT lost its footer entirely,
    while kept slides held on to theirs: a deck came back with page
    numbers on three slides out of thirteen.

    The layout already says where each one goes; this clones them in and
    fills them with the source deck's own footer text and the slide's
    new position.
    """
    import copy

    from pptx import Presentation as _Presentation

    def _chrome_text(deck):
        """The date and footer strings the source deck actually used."""
        date_text = footer_text = None
        for slide in deck.slides:
            for shape in slide.placeholders:
                try:
                    name = shape.placeholder_format.type.name
                except Exception:  # noqa: BLE001
                    continue
                text = " ".join(shape.text_frame.text.split()) if shape.has_text_frame else ""
                if not text:
                    continue
                if name == "DATE" and date_text is None:
                    date_text = text
                elif name == "FOOTER" and footer_text is None:
                    footer_text = text
            if date_text and footer_text:
                break
        return date_text, footer_text

    try:
        source = _Presentation(str(source_path))
        out = _Presentation(str(out_path))
    except Exception:  # noqa: BLE001
        return 0
    date_text, footer_text = _chrome_text(source)

    restored = 0
    for number, slide in enumerate(out.slides, start=1):
        present = set()
        for shape in slide.placeholders:
            try:
                present.add(shape.placeholder_format.type.name)
            except Exception:  # noqa: BLE001
                continue
        for layout_ph in slide.slide_layout.placeholders:
            try:
                name = layout_ph.placeholder_format.type.name
            except Exception:  # noqa: BLE001
                continue
            if name not in _CHROME_PLACEHOLDERS or name in present:
                continue
            element = copy.deepcopy(layout_ph._element)
            slide.shapes._spTree.append(element)
            added = slide.shapes[-1]
            value = {"DATE": date_text, "FOOTER": footer_text,
                     "SLIDE_NUMBER": str(number)}[name]
            if value and added.has_text_frame:
                added.text_frame.text = value
                restored += 1
    if restored:
        out.save(str(out_path))
    return restored


def plan_transform(
    deck_path,
    reference_path: Optional[str] = None,
    template_path=None,
    rules_config: Optional[dict] = None,
    model: str = DEFAULT_MODEL,
    api_key: Optional[str] = None,
    client=None,
    suggest_archetypes: bool = True,
) -> TransformPlan:
    """Per-slide plan for an existing deck. Deterministic core
    (`propose_rebrand`), plus -- when an API key/client is available and
    `suggest_archetypes` is on -- one fail-closed model call offering an
    archetype for any eligible ordinary-content slide; a failure or
    missing key just means no suggestions, never a failed plan."""
    del rules_config  # planning needs no rules; execution loads them itself
    plan = propose_rebrand(deck_path, template_path=template_path, reference_path=reference_path)
    reference_indices = set(plan.reference_layout_by_index)

    # Verbatim content for every slide whose shapes can be read at all --
    # NOT just the ones an org-template layout can hold. That distinction
    # matters enormously on real decks: `classify_slide` refuses any
    # slide with more than MAX_TEXT_BLOCKS (3) text boxes, because that's
    # the most placeholders any org layout has, and on a real 168-slide
    # catalog deck that single cap refused 129 slides -- 77% of the deck,
    # every one of them shown to a human as "can't happen" with nothing
    # offered. But those slides have perfectly good extractable content;
    # they're only too big for the ORG TEMPLATE. Archetypes are exactly
    # the thing that holds them (a comparison table, a 6-item numbered
    # row, a 4-column quarterly plan), so they're offered here as the
    # alternative rather than the slide being written off.
    prs = Presentation(str(deck_path))
    slide_height_in = prs.slide_height / 914400 if prs.slide_height else None
    content_by_index: dict = {}
    readable_by_index: dict = {}  # {index: (title, text_blocks)} for reason-bearing slides
    for p in plan.proposals:
        # Reference-carryover slides are extracted TOO. Carryover
        # reparents the old slide onto a layout borrowed from the
        # reference -- the old shapes, restyled. Now that the
        # reference's own designs are available as archetypes, the
        # better answer for these slides is usually to re-render their
        # content through one, so they need content extracted to be
        # offered the choice at all.
        title, text_blocks, images, reason = _extract_slide_content(
            prs.slides[p.slide_index - 1], slide_height_in
        )
        # A `reason` here is a genuinely unreadable slide (table, chart,
        # embedded object, media, group) -- nothing can reflow that, so
        # it stays keep-only and says so.
        if not reason:
            content_by_index[p.slide_index] = (title, text_blocks, images)
        else:
            # A reason means no ORG-TEMPLATE rebuild is safe -- a group
            # or an embedded object can't be reflowed into placeholders.
            # It does not mean the slide has nothing to say. Reading its
            # words straight off the shape tree (groups included) is
            # enough to offer an archetype, which is the only route
            # these slides ever had; leaving them with none is what made
            # a third of a real deck un-actionable.
            readable_by_index[p.slide_index] = _readable_text(prs.slides[p.slide_index - 1])

    # A reference deck's own designs, read out as archetypes and
    # registered alongside the built-in ones. This is what turns "make
    # it look like that deck" from approximating its colours into
    # re-rendering this deck's content through its actual layouts.
    reference_archetypes: set = set()
    reference_sources: dict = {}
    if reference_path and suggest_archetypes:
        try:
            from deckguard.mine import install_reference
            from deckguard.skill_bridge import _ensure_skill_on_path, _load_archetypes

            _ensure_skill_on_path()
            mined = install_reference(_load_archetypes(), reference_path)
            reference_archetypes = set(mined["archetypes"])
            reference_sources = mined["sources"]
        except Exception:  # noqa: BLE001 -- additive; the built-in library still applies
            reference_archetypes = set()

    suggestions: dict = {}
    ai_ran = False
    if suggest_archetypes:
        from deckguard.skill_bridge import (
            archetype_suggestions_available,
            match_archetypes,
            select_archetype_overrides_for_rebrand,
        )

        cover_end = {plan.cover_index, plan.end_index}
        profile_by_index = {
            idx: SlideProfile(title=t, text_blocks=blocks, images=images, eligible=True)
            for idx, (t, blocks, images) in content_by_index.items()
            if idx not in cover_end
        }
        if profile_by_index:
            # `ai_ran` says whether the model could be ASKED, not whether
            # it happened to suggest anything -- a keyless server used to
            # report suggestions as having run and then offer nothing but
            # "keep", which reads as the tool failing rather than the
            # tool being switched off.
            ai_ran = archetype_suggestions_available(api_key=api_key, client=client)
            if ai_ran:
                suggestions = select_archetype_overrides_for_rebrand(
                    profile_by_index, model=model, api_key=api_key, client=client,
                )

    # Every slide the tool can READ now has an archetype it could become,
    # model or no model. Choosing a new format for an old slide sounds
    # like a job only a model can do; the part that genuinely needs
    # judgement is narrower -- which of several plausible formats reads
    # best, and how the copy should be reworded. What CAN hold the
    # content is a fitting problem, and the archetype library states its
    # own capacity. Before this, the model was the ONLY route, so a
    # keyless server offered nothing but "keep" on every dense slide.
    structural: dict = {}
    if suggest_archetypes:
        from deckguard.skill_bridge import match_archetypes as _match

        cover_end = {plan.cover_index, plan.end_index}
        pool = {
            idx: ([[text for _level, text in block] for block in blocks], len(images))
            for idx, (title, blocks, images) in content_by_index.items()
        }
        titles = {idx: t for idx, (t, _b, _i) in content_by_index.items()}
        for idx, (title, blocks) in readable_by_index.items():
            pool[idx] = (blocks, 0)
            titles[idx] = title
        # Picking each slide's single best fit independently gave a deck
        # where six of eight slides became the same archetype -- every
        # one individually defensible, the deck as a whole monotonous.
        # A short memory of what the previous slides used breaks the tie
        # in favour of variety, and only ever among candidates that
        # already fit.
        recent: list = []
        for idx in sorted(pool):
            if idx in suggestions or idx in cover_end:
                continue
            flat, image_count = pool[idx]
            # A slide with nothing to say gets nothing proposed. Found on
            # a real deck: an empty "Slogan" slide carrying only its
            # date, footer and page number was handed an archetype,
            # which then rendered an empty slide. Keeping it is the
            # honest answer.
            if not any(line.strip() for block in flat for line in block) and not image_count:
                continue
            candidates = _match(titles.get(idx), flat, image_count=image_count, limit=4,
                                prefer=reference_archetypes)
            if not candidates:
                continue
            fresh = [c for c in candidates if c["archetype"] not in recent]
            chosen = fresh[0] if fresh else candidates[0]
            structural[idx] = chosen
            recent.append(chosen["archetype"])
            del recent[:-2]  # remember only the last two

    slides = []
    for p in plan.proposals:
        if p.slide_index in suggestions:
            default_action = "archetype"
        elif p.slide_index in reference_indices:
            default_action = "reference_layout"
        elif p.eligible:
            default_action = "rebuild"
        elif p.slide_index in content_by_index:
            # Readable, just too big for any org layout -- keep by
            # default (nothing here changes its structure), but the
            # review page can still offer an archetype for it if one was
            # suggested, and says WHY the org route is unavailable.
            default_action = "keep"
        else:
            default_action = "keep"
        match = structural.get(p.slide_index)
        archetype = suggestions.get(p.slide_index)
        source = "model" if archetype else ("structural" if match else None)
        if archetype is None and match:
            archetype = match["content"]
            from_reference = match["archetype"] in reference_archetypes
            if from_reference:
                source = "reference"
            # A structural match that loses nothing is strictly better
            # than leaving an old slide as it was; one that would drop
            # content is offered, never imposed -- that call is the
            # reviewer's, and the card says what it costs.
            #
            # A REFERENCE design that loses nothing also beats layout
            # carryover: carryover reparents the old shapes onto a
            # borrowed layout, where this re-renders the content through
            # the design the reference actually uses. That is what "make
            # it look like that deck" was always asking for.
            if match["dropped"] == 0 and (
                default_action == "keep"
                or (from_reference and default_action == "reference_layout")
            ):
                default_action = "archetype"
        slides.append(
            SlidePlan(
                index=p.slide_index,
                default_action=default_action,
                layout_name=p.layout_name,
                archetype=archetype,
                archetype_source=source,
                archetype_dropped=match["dropped"] if match else 0,
                content_chunks=(match["dropped"] + match["capacity"]) if match else 0,
                reason=p.reason,
                title_preview=p.title_preview,
                body_preview=p.body_preview,
                image_count=p.image_count,
                title=content_by_index.get(p.slide_index, (None, [], []))[0],
                text_blocks=[
                    [text for _level, text in block]
                    for block in content_by_index.get(p.slide_index, (None, [], []))[1]
                ],
            )
        )
    return TransformPlan(slides=slides, ai_suggestions_ran=ai_ran,
                         reference_designs=reference_sources)


def plan_transform_from_brief(
    brief: str,
    target_slides: Optional[int] = None,
    model: str = DEFAULT_MODEL,
    notes: Optional[str] = None,
    api_key: Optional[str] = None,
    client=None,
) -> TransformPlan:
    """Per-slide plan for a NEW deck from a brief: the archetype spec
    the model planned (same call the brief-only redesign path uses),
    reshaped into reviewable SlidePlan entries -- one per BODY slide;
    the retained master cover/outro aren't choices, so they don't
    appear. Raises RedesignError on a failed planning call, same as the
    brief-only path always has (unlike existing-deck planning there is
    no deterministic fallback to degrade to)."""
    from deckguard.skill_bridge import (
        _load_archetypes,
        _validate_kone_spec,
        call_claude_for_kone_spec,
        check_brief_archetypes,
    )

    spec, _usage = call_claude_for_kone_spec(
        brief, target_slides=target_slides, model=model, notes=notes, api_key=api_key, client=client,
    )
    _validate_kone_spec(spec, set(_load_archetypes().ARCHETYPES.keys()))

    slides = []
    for i, s in enumerate(spec["slides"], start=1):
        content = {k: v for k, v in s.items() if k != "archetype"}
        slides.append(
            SlidePlan(
                index=i,
                default_action="new",
                archetype={"archetype": s["archetype"], **content},
                title_preview=str(content.get("title") or content.get("statement") or content.get("quote") or ""),
            )
        )
    # A brief naming COVER_A_CUT4 / DIVIDER_D / END_LOGO used to get
    # something else built with no explanation: those names come from
    # kone-design's 56-archetype gallery, while only 23 exist in the
    # engine that renders .pptx, and just 17 names appear in both. The
    # substitution is unavoidable; doing it in silence is not.
    return TransformPlan(
        slides=slides,
        ai_suggestions_ran=True,
        deck_title=spec.get("title"),
        archetype_requests=check_brief_archetypes(brief),
    )


@dataclass
class TransformOutcome:
    out_path: str
    rebuilt: list = field(default_factory=list)  # org-layout rebuilds actually executed
    reference_carryover: list = field(default_factory=list)
    archetype_swapped: list = field(default_factory=list)
    kept: list = field(default_factory=list)
    layouts_used: dict = field(default_factory=dict)  # {index: layout-or-archetype name}
    learned_colors: int = 0  # color remappings learned FROM the uploaded reference deck
    learned_fonts: int = 0
    transplanted_shapes: int = 0  # per-shape styles copied verbatim from the reference
    duplicate_logos_removed: int = 0  # see _dedupe_reference_master_logos
    needs_manual_redraw: list = field(default_factory=list)  # 1-based indices; see below
    photos_added: int = 0  # picture slots filled from the KONE photo library


def _dedupe_reference_master_logos(out_path, reference_path) -> int:
    """Remove the duplicate KONE logo that appears on every
    reference-carried slide.

    When a slide is re-parented onto a layout imported FROM the
    reference deck, that reference's master comes with it -- including
    its own logo, which is exactly what we want: the reference IS the
    brand authority for this run. But `fix_deck`'s logo pass then also
    stamps deckguard's OWN bundled logo onto that same master, because
    the reference's mark sits just outside `logo.old_logo_region_in`
    (measured on a real pair: reference logo at 11.76in, region starts
    at 11.9in), so the region scan doesn't recognize it as the existing
    logo and adds a second one beside it. Both render, overlapping, on
    every carried-over slide.

    Fix: on any master carrying a picture byte-identical to one on the
    REFERENCE deck's own master, drop any OTHER picture whose box
    overlaps it. The reference's mark always wins -- deckguard's bundled
    default has no business overriding the deck the user supplied as the
    target look. Masters with no reference-sourced logo (the org
    template's own) are untouched, so a no-reference run is unaffected.

    Returns how many duplicates were removed.
    """
    import hashlib

    def _pics(master):
        out = []
        for shape in master.shapes:
            try:
                if shape.shape_type is None or shape.shape_type.name != "PICTURE":
                    continue
                if None in (shape.left, shape.top, shape.width, shape.height):
                    continue
                out.append((shape, hashlib.sha1(shape.image.blob).hexdigest()))
            except Exception:  # noqa: BLE001 -- unreadable picture part
                continue
        return out

    def _overlaps(a, b):
        return not (
            a.left + a.width <= b.left or b.left + b.width <= a.left
            or a.top + a.height <= b.top or b.top + b.height <= a.top
        )

    try:
        ref_prs = Presentation(str(reference_path))
        ref_hashes = {h for m in ref_prs.slide_masters for _s, h in _pics(m)}
        if not ref_hashes:
            return 0

        prs = Presentation(str(out_path))
        removed = 0
        for master in prs.slide_masters:
            pics = _pics(master)
            authoritative = [s for s, h in pics if h in ref_hashes]
            if not authoritative:
                continue  # not a reference-sourced master -- leave it alone
            for shape, h in pics:
                if h in ref_hashes:
                    continue
                if any(_overlaps(shape, keeper) for keeper in authoritative):
                    shape._element.getparent().remove(shape._element)
                    removed += 1
        if removed:
            prs.save(str(out_path))
        return removed
    except Exception:  # noqa: BLE001 -- cosmetic cleanup, never load-bearing
        return 0


def _config_learned_from_reference(deck_path, reference_path, base_config: dict) -> tuple[dict, int, int]:
    """Diff the uploaded deck against the uploaded REFERENCE deck and
    fold what that pair actually demonstrates (this old color became
    that new one; this old font became that new one) into a per-run copy
    of the brand config.

    This is what makes a transform driven by YOUR decks rather than only
    by the bundled brand_rules.yaml: without it the reference contributes
    nothing but layout names, and every color/font decision comes from
    the backend's own canned tables. Never writes to brand_rules.yaml --
    the learned config is scoped to this one run.
    """
    from deckguard import learn as learn_mod

    try:
        result = learn_mod.learn(Presentation(str(deck_path)), Presentation(str(reference_path)), base_config)
        learned = learn_mod.apply_learned(base_config, result, min_confidence="high")
        n_colors = sum(1 for p in result.color_proposals if p.confidence == "high")
        n_fonts = sum(1 for p in result.font_proposals if p.confidence == "high")
        return learned, n_colors, n_fonts
    except Exception:  # noqa: BLE001 -- learning is additive; fall back to the base config
        return base_config, 0, 0


def execute_transform(
    deck_path,
    out_path,
    plan: TransformPlan,
    actions: Optional[dict] = None,
    reference_path: Optional[str] = None,
    template_path=None,
    rules_config: Optional[dict] = None,
) -> TransformOutcome:
    """Run exactly the approved subset of `plan`. `actions` maps
    1-based slide index -> "keep" | "rebuild" | "archetype"; a slide
    absent from `actions` follows its plan default. "rebuild" covers
    both ordinary org-layout rebuilds and reference-layout carryovers
    (`apply_rebrand` distinguishes them itself); "archetype" is only
    honored where the plan actually carries a suggestion -- a stray
    "archetype" choice with no suggestion degrades to "rebuild" rather
    than failing or inventing content."""
    actions = actions or {}
    plan_by_index = {s.index: s for s in plan.slides}

    chosen: dict = {}
    for s in plan.slides:
        choice = actions.get(s.index, s.default_action)
        # An archetype choice needs a suggestion to render; without one
        # it degrades rather than failing or inventing content.
        if choice == "archetype" and s.archetype is None:
            choice = "rebuild"
        if choice in ("rebuild", "reference_layout"):
            # An org-layout rebuild is only possible where the planner
            # found a layout that fits -- "keep" defaults never can.
            choice = "rebuild" if s.default_action != "keep" else "keep"
        elif choice not in ("keep", "archetype"):
            choice = "keep"
        # NOTE: an "archetype" choice is honored even on a keep-default
        # slide. Those are slides no ORG layout can hold (too many text
        # blocks) but an archetype can -- offering the archetype there is
        # the whole point; only the org-rebuild route is unavailable.
        chosen[s.index] = choice

    rebuild_indices = {i for i, c in chosen.items() if c == "rebuild"}
    archetype_indices = {i for i, c in chosen.items() if c == "archetype"}

    base_config = rules_config if rules_config is not None else load_config(default_config_path())
    effective_config, learned_colors, learned_fonts = base_config, 0, 0
    if reference_path:
        effective_config, learned_colors, learned_fonts = _config_learned_from_reference(
            deck_path, reference_path, base_config
        )

    rebrand_result = apply_rebrand(
        str(deck_path), str(out_path), template_path=template_path, rules_config=effective_config,
        reference_path=reference_path, accepted_indexes=rebuild_indices,
    )

    # "keep" means the slide's STRUCTURE is untouched -- deck-wide brand
    # patches (colors/fonts/effects) still apply everywhere, Transform's
    # baseline promise. apply_rebrand runs fix_deck itself whenever it
    # rebuilds anything, but its nothing-to-rebuild early path returns a
    # plain copy with no fix pass at all -- cover that case here so an
    # all-keep transform is still a brand fix, not a no-op.
    if not rebuild_indices and not rebrand_result.reference_layout_indices:
        from deckguard.fixer import fix_deck

        prs_out = Presentation(str(out_path))
        fix_deck(prs_out, effective_config, source_path=str(out_path), output_path=str(out_path), dry_run=False)

    layouts_used = {
        p.slide_index: p.layout_name
        for p in rebrand_result.proposals
        if p.slide_index in rebuild_indices and p.layout_name
    }

    if archetype_indices:
        from deckguard.skill_bridge import apply_archetype_overrides_to_deck

        prs = Presentation(str(deck_path))
        slide_height_in = prs.slide_height / 914400 if prs.slide_height else None
        images_by_index: dict = {}
        for idx in archetype_indices:
            _title, _blocks, images, reason = _extract_slide_content(prs.slides[idx - 1], slide_height_in)
            if not reason:
                images_by_index[idx] = images
        overrides = {i: plan_by_index[i].archetype for i in archetype_indices}
        swapped = apply_archetype_overrides_to_deck(str(out_path), overrides, images_by_index=images_by_index)
        layouts_used.update(swapped)

    # Per-shape styling copied straight off the reference deck, wherever
    # a shape's identity survives between the two (same slide index +
    # name, or same position). This is the other half of "match the
    # reference": `learn` above generalizes the pair into deck-wide
    # color/font rules, while this copies the reference's own answer for
    # a specific shape -- provably better where one source color has to
    # become two different target colors depending on the shape (see
    # exact_transplant's own docstring). Skips archetype slides: their
    # shapes come from the archetype engine, not the old deck, so there
    # is no surviving identity to match and their styling is already
    # correct by construction.
    #
    # The transplant also reports which slides it could NOT match --
    # slides where too few shapes line up with the reference for any
    # per-shape copy to be meaningful. In practice that means the human
    # who made the reference deck REDREW the slide (a flow diagram
    # rebuilt as grouped shapes and straight connectors where the old
    # deck used elbow connectors and loose text boxes, say). No amount
    # of restyling turns one into the other, so the honest thing is to
    # name those slides instead of shipping a quietly worse version of
    # them. Reported as `needs_manual_redraw`.
    #
    # Only slides whose ORIGINAL structure survived can be reported this
    # way. An archetype swap or an org-layout rebuild replaces the
    # slide's shapes by design, so of course they don't line up with the
    # reference's -- flagging those would be crying wolf on exactly the
    # slides the tool just handled well. That leaves "keep" and
    # reference-carryover slides, which still carry the old deck's own
    # shapes and therefore genuinely should have matched.
    transplanted = 0
    needs_manual_redraw: list = []
    structure_preserved = {i for i, c in chosen.items() if c == "keep"} | set(
        rebrand_result.reference_layout_indices
    )
    if reference_path:
        from deckguard.exact_transplant import transplant_exact_treatment

        try:
            out_prs = Presentation(str(out_path))
            ref_prs = Presentation(str(reference_path))
            result = transplant_exact_treatment(out_prs, ref_prs, rules_config=effective_config)
            changes = [c for c in result.changes if (c.slide_index + 1) not in archetype_indices]
            if changes:
                out_prs.save(str(out_path))
                transplanted = len(changes)
            # exact_transplant indexes slides from 0.
            needs_manual_redraw = [
                i + 1 for i in result.flagged_slides if (i + 1) in structure_preserved
            ]
        except Exception:  # noqa: BLE001 -- additive polish, never load-bearing
            transplanted = 0

    # Rebuilt slides lose their date/footer/page number, because
    # python-pptx does not clone those "latent" placeholders.
    restore_footer_chrome(out_path, deck_path)
    # ...and the KONE master's own logo frames ship with no image in
    # them, which PowerPoint draws as dotted boxes.
    from deckguard.logo import repair_empty_logo_frames

    repair_empty_logo_frames(out_path)

    duplicate_logos_removed = 0
    if reference_path:
        duplicate_logos_removed = _dedupe_reference_master_logos(out_path, reference_path)

    return TransformOutcome(
        out_path=str(out_path),
        rebuilt=sorted(rebuild_indices - set(rebrand_result.reference_layout_indices)),
        reference_carryover=sorted(rebrand_result.reference_layout_indices),
        archetype_swapped=sorted(archetype_indices),
        kept=sorted(i for i, c in chosen.items() if c == "keep"),
        layouts_used=layouts_used,
        learned_colors=learned_colors,
        learned_fonts=learned_fonts,
        transplanted_shapes=transplanted,
        duplicate_logos_removed=duplicate_logos_removed,
        needs_manual_redraw=needs_manual_redraw,
    )


def execute_transform_from_brief(out_path, plan: TransformPlan, approved_indices: Optional[set] = None) -> TransformOutcome:
    """Render the approved subset of a brief-only plan through the
    skill's own whole-deck builder (retained master cover + outro, same
    as the brief-only path always produced)."""
    from deckguard.skill_bridge import _load_creator, fill_empty_photo_slots

    creator = _load_creator()
    approved = approved_indices if approved_indices is not None else {s.index for s in plan.slides}
    spec_slides = [dict(s.archetype) for s in plan.slides if s.index in approved and s.archetype]
    spec = {"title": plan.deck_title or "Untitled deck", "slides": spec_slides}
    # The review previews draw picture slots as PHOTO; without this the
    # built deck came back with blank sand blocks where they were.
    photos_added = fill_empty_photo_slots(spec)
    creator.build_deck(spec, str(out_path))
    from deckguard.logo import repair_empty_logo_frames, restore_logo_chrome

    repair_empty_logo_frames(out_path)  # the master's empty logo frames
    # ...and the 15 layouts whose logo is a placeholder python-pptx never
    # clones, so the slide inherits nothing to repair. Covers, mostly.
    restore_logo_chrome(out_path)  # exactly one mark per slot
    # An author who asked for a cover/closer archetype gets theirs, not
    # theirs plus the master's.
    from deckguard.gallery import drop_redundant_master_slides, stamp_footers

    drop_redundant_master_slides(out_path, spec)
    # Page numbers depend on final slide order, so this runs last.
    stamp_footers(out_path, spec)
    return TransformOutcome(
        out_path=str(out_path),
        archetype_swapped=sorted(s.index for s in plan.slides if s.index in approved),
        kept=sorted(s.index for s in plan.slides if s.index not in approved),
        layouts_used={s.index: s.archetype["archetype"] for s in plan.slides if s.index in approved and s.archetype},
        photos_added=photos_added,
    )


def audit_transform_result(out_path, archetype_indices: Optional[set] = None, rules_config: Optional[dict] = None) -> dict:
    """Brand-rule audit of a transform result. Findings on
    archetype-rendered slides are EXCLUDED from the reported list (and
    counted separately): those slides are compliant by construction --
    the archetype engine writes only its own approved treatment -- and
    `rules_engine`'s generic text rules provably false-positive on their
    deliberate styling (e.g. flagging kone_engine's muted caption grey
    as a text_color violation). Suppressing a false alarm and SAYING SO
    beats reporting a defect that isn't one."""
    archetype_indices = archetype_indices or set()
    config = rules_config if rules_config is not None else load_config(default_config_path())
    prs = Presentation(str(out_path))
    violations = sort_violations(audit_deck(build_inventory(prs), config))
    reported = [v for v in violations if v.slide_index not in archetype_indices]
    suppressed = len(violations) - len(reported)
    return {
        "summary": summarize(reported),
        "violations": reported,
        "suppressed_archetype_findings": suppressed,
    }


def reference_similarity(out_path, reference_path) -> dict:
    """How close the transformed deck landed to the reference deck --
    deterministic, structural signals only (no rendering, no AI):
    per-index layout-name matches, plus colors and fonts present in the
    result that the reference itself never uses (each one a spot the
    result still visibly diverges)."""

    def _population(prs):
        inv = build_inventory(prs)
        colors: set = set()
        fonts: set = set()
        layouts: list = []
        for slide in inv.slides:
            layouts.append(slide.layout_name)
            for shape in iter_shapes_recursive(slide.shapes):
                if shape.fill:
                    for c in shape.fill.colors:
                        if c.hex:
                            colors.add(c.hex.upper())
                for para in shape.paragraphs:
                    for run in para.runs:
                        if run.color and run.color.hex:
                            colors.add(run.color.hex.upper())
                        if run.font_effective:
                            fonts.add(run.font_effective)
        return layouts, colors, fonts

    out_layouts, out_colors, out_fonts = _population(Presentation(str(out_path)))
    ref_layouts, ref_colors, ref_fonts = _population(Presentation(str(reference_path)))

    n = min(len(out_layouts), len(ref_layouts))
    layout_matches = sum(1 for i in range(n) if out_layouts[i] == ref_layouts[i])
    return {
        "slides_compared": n,
        "layout_matches": layout_matches,
        "colors_not_in_reference": sorted(out_colors - ref_colors),
        "fonts_not_in_reference": sorted(out_fonts - ref_fonts),
    }
