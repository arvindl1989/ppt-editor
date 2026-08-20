"""Package the current state of the tool for Claude Design.

Produces a folder (and a zip) holding everything needed to answer one
question -- what placeholder content should each of the fifty archetypes
carry -- without having to read this repository:

    README.md            what the tool is and where it got to
    INSTRUCTIONS.md      the ask, and exactly how to return it
    CONTRACTS.md         every archetype's slots, in set order
    contracts.json       the same, machine-readable, with fit budgets
    placeholders.*.json  what is used today, and an empty skeleton
    previews/            a real render of every slide, per audience
    brand/               the handoff this was all built from
    assets/              what can be referenced by name

Run:  python scripts/design_handoff.py [outdir]

Rendering needs LibreOffice and poppler. Without them the previews are
skipped and everything else is still written.
"""

from __future__ import annotations

import json
import shutil
import subprocess
import sys
import tempfile
from datetime import date
from pathlib import Path

from deckguard import brandmode as bm
from deckguard import contracts as C

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "src" / "deckguard" / "assets" / "kone-design"

# The same estimate `preview._fit_scale` uses: a sans glyph advances
# ~0.52em and a line occupies 1.25x its size. Rough, but it is the
# difference between copy that fits the box it was written for and copy
# that autofits down to something nobody chose.
GLYPH_EM = 0.52
LINE_HEIGHT = 1.25


def _built() -> set:
    from deckguard.registry import _load_archetypes

    return set(_load_archetypes().ARCHETYPES)


def _boxes(archetype: str) -> dict:
    """{slot key: (width_px, height_px, type_px)} from the live registry."""
    from deckguard.registry import _load_archetypes

    spec = _load_archetypes().ARCHETYPES.get(archetype) or {}
    out: dict = {}

    def _size(region) -> float:
        style = region.get("dg") or {}
        if style.get("px"):
            return float(style["px"])
        # `resolve`, not the raw scale: it is what the renderer calls,
        # and it is where a title in a narrow column drops to 28. Taking
        # the scale value straight told Design a cover headline had room
        # for 374 characters, because it costed a 76px line at 32px.
        role = str(region.get("role") or "")
        width = region["box"][2]
        try:
            resolved = bm.resolve(role, width=width)
            if resolved and resolved.get("px"):
                return float(resolved["px"])
        except Exception:  # noqa: BLE001
            pass
        try:
            return float(bm.TYPE_SCALE[role][1])
        except (KeyError, IndexError, TypeError):
            return 16.0

    for region in spec.get("regions") or []:
        key = region.get("content")
        if key:
            _x, _y, w, h = region["box"]
            out[key] = (w, h, _size(region))
    for group in spec.get("groups") or []:
        for region in group.get("regions") or []:
            key = region.get("content")
            if key:
                _x, _y, w, h = region["box"]
                out[f"{group['content']}.{key}"] = (w, h, _size(region))
    return out


def _budget(width: float, height: float, px: float) -> int:
    """Roughly how many characters fit in a box at that size."""
    if not (width and height and px):
        return 0
    per_line = max(int(width / (px * GLYPH_EM)), 1)
    lines = max(int(height / (px * LINE_HEIGHT)), 1)
    return per_line * lines


def measure(deck_path: Path, archetypes: list) -> dict:
    """Read the built deck back and cost each slot at the size it
    actually rendered at.

    Estimating from the registry got covers badly wrong -- the box says
    `title`, the renderer resolves it to a 76px cover headline, and the
    estimate cheerfully reported room for 374 characters. Measuring
    removes the guess: the sample content is known, so a shape can be
    matched back to the slot whose text it holds, and the budget comes
    from that shape's real box and real point size.
    """
    from pptx import Presentation

    from deckguard.preview import sample_content

    prs = Presentation(str(deck_path))
    px = prs.slide_width / 1280
    out: dict = {}
    slides = list(prs.slides)[1:-1]        # between the retained cover and outro
    for archetype, slide in zip(archetypes, slides):
        wanted: dict = {}
        for key, value in sample_content(archetype).items():
            if isinstance(value, str) and value.strip():
                wanted.setdefault(value.strip(), key)
            elif isinstance(value, list):
                for item in value:
                    if isinstance(item, dict):
                        for field_key, field_value in item.items():
                            if isinstance(field_value, str) and field_value.strip():
                                wanted.setdefault(field_value.strip(),
                                                  f"{key}.{field_key}")
                    elif isinstance(item, str) and item.strip():
                        wanted.setdefault(item.strip(), key)

        found: dict = {}
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            key = wanted.get(text) or next(
                (k for t, k in wanted.items() if t and t in text), None)
            if not key or key in found:
                continue
            sizes = [r.font.size.pt for para in shape.text_frame.paragraphs
                     for r in para.runs if r.font.size]
            if not sizes:
                continue
            # points at 13.333in wide == 960pt, so 1pt is 1280/960 px
            size_px = max(sizes) * (1280.0 / 960.0)
            found[key] = _budget(shape.width / px, shape.height / px, size_px)
        if found:
            out[archetype] = found
    return out


def contracts_json(measured: dict | None = None) -> dict:
    """Every archetype, per audience, with what it needs and how much of
    it fits."""
    built = _built()
    measured = measured or {}
    out: dict = {"generated": date.today().isoformat(), "sets": {}}
    for audience in bm.set_names():
        entries = []
        for slide in bm.slides_in(audience):
            name = slide["archetype"]
            contract = C.for_archetype(name, audience)
            # measured where the deck was actually rendered, estimated
            # from the registry where it was not
            boxes = _boxes(name) if name in built else {}
            real = measured.get(name, {})

            def _fits(key, _real=real, _boxes=boxes):
                if key in _real:
                    return _real[key]
                return _budget(*_boxes[key]) if key in _boxes else None
            slots = []
            for slot in (contract.slots if contract else ()):
                item = {
                    "key": slot.key,
                    "role": slot.role,
                    "required": not slot.optional,
                    "picture": slot.is_picture,
                }
                if slot.is_list:
                    item["count"] = {"min": slot.minimum, "max": slot.maximum}
                    item["fields"] = [
                        {"key": f.key, "role": f.role,
                         "fits_chars": _fits(f"{slot.key}.{f.key}")}
                        for f in slot.fields
                    ]
                else:
                    fits = _fits(slot.key)
                    if fits:
                        item["fits_chars"] = fits
                slots.append(item)
            entries.append({
                "n": slide["n"],
                "archetype": name,
                "group": slide["group"],
                "field": slide["field"],
                "built": name in built,
                "job": (contract.job if contract else "") or bm.job_for(name, audience),
                "slots": slots,
            })
        out["sets"][audience] = entries
    return out


def contracts_md(data: dict) -> str:
    lines = [
        "# What each slide needs",
        "",
        "Generated from the live registry -- this is what the renderer",
        "actually reads, not a description of it. `?` marks an optional",
        "slot. `fits` is roughly how many characters the box holds at the",
        "size it is set in; going over does not break the build, it",
        "autofits the type down, which is the thing to avoid.",
        "",
    ]
    for audience, entries in data["sets"].items():
        lines += [f"## {audience.title()} 25", ""]
        for entry in entries:
            head = f"### {entry['n']:02d} · `{entry['archetype']}`"
            if not entry["built"]:
                head += "  — NOT BUILT YET, no placeholder needed"
            lines += [head, "", f"{entry['group']} · field: {entry['field']}", ""]
            if entry["job"]:
                lines += [entry["job"].strip(), ""]
            for slot in entry["slots"]:
                if slot["picture"]:
                    lines.append(f"- `{slot['key']}` — photograph, filled "
                                 "automatically. Do not supply.")
                    continue
                mark = "" if slot["required"] else "?"
                if "count" in slot:
                    count = slot["count"]
                    span = (str(count["max"]) if count["min"] == count["max"]
                            else f"{count['min']}-{count['max']}")
                    inner = ", ".join(
                        f"{f['key']}"
                        + (f" (fits ~{f['fits_chars']})" if f.get("fits_chars") else "")
                        for f in slot["fields"])
                    lines.append(f"- `{slot['key']}`{mark} — a list of {span}, "
                                 f"each {{{inner}}}")
                else:
                    fits = (f", fits ~{slot['fits_chars']} chars"
                            if slot.get("fits_chars") else "")
                    lines.append(f"- `{slot['key']}`{mark} — {slot['role'] or 'text'}{fits}")
            lines.append("")
    return "\n".join(lines)


def placeholders_current() -> dict:
    """What the previews use today. Mine, and it shows."""
    from deckguard.preview import sample_content

    built = _built()
    out: dict = {}
    for audience in bm.set_names():
        out[audience] = {
            slide["archetype"]: sample_content(slide["archetype"])
            for slide in bm.slides_in(audience) if slide["archetype"] in built
        }
    return out


def placeholders_template(data: dict) -> dict:
    """The exact shape to fill in, with every slot present and empty."""
    out: dict = {}
    for audience, entries in data["sets"].items():
        per: dict = {}
        for entry in entries:
            if not entry["built"]:
                continue
            slide: dict = {}
            for slot in entry["slots"]:
                if slot["picture"]:
                    continue
                if "count" in slot:
                    # A field typed `bullets` inside a list is still a
                    # list. Emitting "" for it asked for one clause where
                    # three lines belong, and six slots came back written
                    # that way -- `two_content.items[].bullets` and its
                    # five siblings. The renderer had always accepted an
                    # array there; only the template said otherwise.
                    slide[slot["key"]] = [
                        {f["key"]: (["", "", ""] if f.get("role") == "bullets" else "")
                         for f in slot["fields"]}
                        for _ in range(slot["count"]["max"])
                    ]
                elif slot["role"] == "bullets":
                    slide[slot["key"]] = ["", "", ""]
                else:
                    slide[slot["key"]] = ""
            per[entry["archetype"]] = slide
        out[audience] = per
    return out


def render_previews(out_dir: Path) -> int:
    """One PNG per slide, per audience, numbered in set order."""
    for tool in ("soffice", "pdftoppm"):
        if shutil.which(tool) is None:
            print(f"  (skipping previews: {tool} not on PATH)")
            return 0, {}

    from deckguard import assemble
    from deckguard.preview import sample_content

    built = _built()
    written = 0
    measured: dict = {}
    for audience in bm.set_names():
        slides = [s for s in bm.slides_in(audience) if s["archetype"] in built]
        target = out_dir / audience
        target.mkdir(parents=True, exist_ok=True)
        spec = {
            "title": f"KONE {audience} 25",
            "date": assemble.bm_date(),
            "slides": [{"archetype": s["archetype"], **sample_content(s["archetype"])}
                       for s in slides],
        }
        with tempfile.TemporaryDirectory() as tmp:
            work = Path(tmp)
            deck = work / f"{audience}.pptx"
            assemble.build(spec, str(deck))
            subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                            "--outdir", str(work), str(deck)],
                           check=True, capture_output=True, timeout=900)
            # JPEG, not PNG: these are photographs as often as they are
            # type, and 45 full-slide PNGs came to 16MB, which is a
            # package nobody wants to move around. At 1280px and q90 the
            # type is still sharp enough to judge.
            subprocess.run(["pdftoppm", "-jpeg", "-jpegopt", "quality=90",
                            "-r", "96",
                            str(work / f"{audience}.pdf"), str(work / "page")],
                           check=True, capture_output=True, timeout=900)
            pages = sorted(work.glob("page-*.jpg"))
            # page 1 is the master's retained cover, so body slide i is
            # page i + 1
            for index, slide in enumerate(slides):
                if index + 1 >= len(pages):
                    break
                name = f"{slide['n']:02d}_{slide['archetype']}.jpg"
                shutil.copyfile(pages[index + 1], target / name)
                written += 1
            # the deck itself, so the geometry can be inspected directly
            shutil.copyfile(deck, out_dir / f"{audience}-25.pptx")
            measured.update(measure(deck, [s["archetype"] for s in slides]))
        contact_sheet(target, out_dir / f"{audience}-contact-sheet.jpg")
    return written, measured


def contact_sheet(shots_dir: Path, out_path: Path, columns: int = 5) -> bool:
    """All 25 on one page, numbered, so a set can be read as a set.

    Opening twenty-five files one at a time is the wrong way to judge
    whether a deck holds together, which is exactly the question being
    asked here.
    """
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return False
    shots = sorted(shots_dir.glob("*.jpg"))
    if not shots:
        return False

    cell_w, gap, label_h, pad = 420, 18, 26, 28
    cell_h = round(cell_w * 720 / 1280)
    rows = -(-len(shots) // columns)
    sheet = Image.new(
        "RGB",
        (pad * 2 + columns * cell_w + (columns - 1) * gap,
         pad * 2 + rows * (cell_h + label_h) + (rows - 1) * gap),
        "white")
    draw = ImageDraw.Draw(sheet)
    for index, shot in enumerate(shots):
        col, row = index % columns, index // columns
        x = pad + col * (cell_w + gap)
        y = pad + row * (cell_h + label_h + gap)
        with Image.open(shot) as art:
            sheet.paste(art.convert("RGB").resize((cell_w, cell_h)), (x, y))
        draw.rectangle([x, y, x + cell_w - 1, y + cell_h - 1], outline="#D8D8D8")
        draw.text((x, y + cell_h + 7), shot.stem.replace("_", " · ", 1), fill="#141414")
    sheet.save(out_path, quality=92)
    return True


def asset_inventory(out_dir: Path) -> dict:
    out_dir.mkdir(parents=True, exist_ok=True)
    from deckguard import icons

    names = sorted(icons.load_icons())
    (out_dir / "icon-names.txt").write_text("\n".join(names), encoding="utf-8")

    illus = sorted(p.name for p in (ASSETS / "illustrations").glob("*.svg"))
    (out_dir / "illustration-names.txt").write_text("\n".join(illus), encoding="utf-8")

    photos = sorted(p.name for p in (ASSETS / "photos").iterdir() if p.is_file())
    (out_dir / "photo-names.txt").write_text("\n".join(photos), encoding="utf-8")

    return {"icons": len(names), "illustrations": len(illus), "photos": len(photos)}


README = """# deckguard — where it got to, and what I need from you

## What the tool is

One page, one button. You give it any combination of a brief (often a
raw announcement email), a set of slides picked by hand, and a .pptx of
your own; it returns a finished KONE deck plus an editable list of what
it built, on the same page. No approval step — every decision is
reversible after the fact rather than demanded before.

Slides are drawn with python-pptx onto KONE's own master, from the
Internal 25 and External 25 you specified. The master's cover and
"Thank you" are retained; everything between them is generated.

## What changed since the last package

- **Every archetype now declares a contract.** The `contract` column in
  your `EXTERNAL_25.md` was being parsed and discarded; it is now the
  authority for the external set, and the internal set has the same
  thing transcribed from the as-built prose in `INTERNAL_25.md`. See
  `CONTRACTS.md`.
- **Five layouts were redrawn** because they could not carry structured
  content at all — they offered a title and a paragraph. `agenda_c_split`
  and `agenda_b_numbered` are real numbered lists; `timeline_quarter_axis`
  has four dated stems on its axis; `picture_intro` is the photo-right
  slide the prose describes rather than the banner-photo variant that had
  been registered; `quote_a` and `quote_b` set their quotations in quote
  type instead of 16px body.
- **The floor is enforced.** Thirteen archetypes drew below y=629, one by
  83px. Nothing had built the whole library at once and read the preflight
  back until now.
- **Previews are real renders**, not wireframes — the same PNGs as in
  `previews/`.
- Three archetypes carried their own footer line under the one the layout
  already stamps; those are gone.

## What is in here

    INSTRUCTIONS.md            the ask
    CONTRACTS.md               every slide's slots, in set order
    contracts.json             the same, machine-readable
    placeholders.current.json  the copy in use today — mine, and it shows
    placeholders.template.json the exact shape to return, empty
    OPEN_QUESTIONS.md          decisions I could not make on my own
    previews/internal/*.jpg    a real render of each of the 25
    previews/external/*.jpg    likewise
    previews/*-contact-sheet.jpg  all 25 on one page, to read as a set
    previews/*-25.pptx         the decks those renders came from
    brand/                     the handoff this was all built from
    assets/                    every icon, illustration and photo, by name

## Where this sits

Phase 1 (contracts) is done. Phase 2 is an extraction pass that reads a
brief into typed material — figures, quotes, dated events, owners — with
no layouts mentioned. Phase 3 is a matcher that picks layouts in code
from that material and the contracts, so the model only writes copy.
Better placeholders make phase 2 and 3 easier to judge, which is why
they come first.
"""


INSTRUCTIONS = """# What I need from you

## The main ask: placeholder content for all 45 built slides

The previews people pick from currently carry copy I wrote, and it
shows: "What changed", "Marketing Hub", "Frontlines" repeated down a
row. It is filler standing where a designed example should be. Someone
choosing between fifty layouts is reading these previews to decide, so
the copy is doing real work and should be yours.

Return **one file, `placeholders.json`**, in exactly the shape of
`placeholders.template.json`:

```json
{
  "internal": {
    "agenda_c_split": {
      "title": "...",
      "lead": "...",
      "items": [{"number": "01", "label": "..."}, ...]
    },
    ...
  },
  "external": { ... }
}
```

### Rules that are not negotiable, because the renderer enforces them

1. **Use exactly the keys in the template.** Anything else is silently
   discarded — that is not a policy, it is what the renderer does. The
   template is generated from the live registry, so it cannot be stale.
2. **Fill every list to its `max`.** A five-row agenda with three rows
   in it previews with two empty blocks. `CONTRACTS.md` gives the count.
3. **Bullets are arrays of strings**, never one string with line breaks
   or dashes in it. BRAND_MODE §6: a dash standing in for a bullet is a
   violation the preflight reports.
4. **Do not supply image paths.** Picture slots are filled from the
   photo library automatically; they are marked in `CONTRACTS.md`.
5. **Leave `icon` empty unless you specifically want one.** An empty
   icon slot runs the pictogram rotation, which is what a real slide
   does. If you do want a particular one, use a name from
   `assets/icon-names.txt` — there are 609 and they are checked. A name
   that is not in that list is not an error and does not leave a hole:
   the engine falls back to the rotation, so a typo fails quietly and
   you get a pictogram nobody chose.
6. **Respect `fits`.** It is measured off the built deck, at the size the
   text actually rendered at. Going over does not break the build — the
   type autofits down, which is worse, because the slide then carries a
   size nobody chose. Treat it as a ceiling, not a target: the role says
   the intent, and an `eyebrow` with room for 595 characters still wants
   two to five words.

### Tone

- **Internal** is a regional programme talking to itself: specific,
  operational, slightly dry. Owners, quarters, what changed, what is
  next.
- **External** is customer-facing: a pitch, a proposal, a QBR, a launch.
  Blue, white, black and photography — no secondary colour as a field.

Both should read as plausible and obviously placeholder. Please do not
write invented KONE figures that could be mistaken for real ones — round
numbers and evidently illustrative claims are ideal. A shared archetype
appears under both audiences and should get different copy in each.

## Second ask: look at the previews

`previews/internal/` and `previews/external/` are real renders, in set
order, numbered. Five slides in particular were redrawn by me from your
prose rather than from a design, and I would like them checked:

    internal 02  picture_intro
    internal 03  agenda_c_split
    internal 12  timeline_quarter_axis
    external 03  agenda_b_numbered
    external 20  quote_a  (and internal 20 quote_b)

Anything else in those folders that is off-brand, tell me — the geometry
is data and cheap to change.

## Third ask: the five that are not built

`text_picture_b`, `two_pictures_text_b`, `statement_b`,
`value_prop_four_point`, `quote_e`. They are in the sets and show as
"not built yet" in the picker. Are they worth building, and if so is the
prose in `brand/` enough to build from? `value_prop_four_point` in
particular has four picture panes and four numbered features and I would
rather not derive that geometry by guesswork.

## How to return it

    placeholders.json    the main deliverable
    NOTES.md             answers to OPEN_QUESTIONS.md and anything you
                         changed your mind about

If a slide's contract is wrong — the slots do not match what the slide
should hold — say so in NOTES.md rather than bending the copy to fit. The
contract is a claim about the design, and where the two disagree the
design wins.
"""


OPEN_QUESTIONS = """# Decisions I could not make on my own

Each of these is a place where two sources disagree, or where I built
something from prose and want it checked.

## 1. Where the numbered divider's parts sit

`INTERNAL_25.md` says "300px blue numeral at left:38 top:150. Section
label and 56px title at x:620." A real KONE deck that uses this layout
five times puts the numeral at x:45 and the title at x:453. I kept the
measured positions and only centred the pair vertically, which is what
was asked for. Which is right?

## 2. How big a cover title actually is

`INTERNAL_25.md` gives `COVER_B_CUT3` a 76px title. Measured off the
built deck it renders at 40px, and the `fits` budgets in `CONTRACTS.md`
are computed from the measurement. If 76px is correct the budget drops
to roughly a quarter of what is stated.

## 3. `picture_intro` — I changed the slide

What was registered was a banner photograph across the top with a title
and one line under it. Your prose describes a full-height photo on the
right with an eyebrow, a 48px statement and three icon rows. I built
your version, because the registered one could not carry the three
reasons the slide exists to give. Confirm.

## 4. The fifth agenda row

"The final row inverts to blue with a white chip." Per-item colour is
not expressible in the current group model — every row in a repeat gets
the same treatment. Worth adding, or let all five be sand?

## 5. How tall the pink panel on `TIMELINE_QUARTER_AXIS` is

The prose gives `top:250` and no bottom. I ran it to the floor at 629,
which leaves the lower half of the panel empty under three bullets. A
shorter panel that ends under the last bullet may be what was meant.

## 6. Two charts do not draw

`CHART_COMMENTARY` and `SEGMENT_BREAKDOWN` both declare a chart region
and neither renders anything into it — you can see the hole in the
previews. The prose describes them (eight stacked pairs; five bars with
a blue-tint ramp) but not as data a caller supplies. Should the chart be
a fixed illustration, or should the contract take series data?

## 7. `stat_value` on things that are not statistics

The external contracts set several cover and outro titles in
`stat_value`. Read literally that tells a planner a cover headline
should be "a short figure, e.g. 91%". I am suppressing the hint for
those keys, but the underlying overload is worth a second look — is
there a display role these should be using instead?
"""


def main(argv: list) -> int:
    out = Path(argv[1]) if len(argv) > 1 else (
        ROOT / f"design-handoff-{date.today().isoformat()}")
    if out.exists():
        shutil.rmtree(out)
    out.mkdir(parents=True)

    counts = asset_inventory(out / "assets")
    shots, measured = render_previews(out / "previews")

    data = contracts_json(measured)
    (out / "contracts.json").write_text(json.dumps(data, indent=2, ensure_ascii=False))
    (out / "CONTRACTS.md").write_text(contracts_md(data), encoding="utf-8")
    (out / "placeholders.current.json").write_text(
        json.dumps(placeholders_current(), indent=2, ensure_ascii=False))
    (out / "placeholders.template.json").write_text(
        json.dumps(placeholders_template(data), indent=2, ensure_ascii=False))

    brand = out / "brand"
    brand.mkdir()
    for name in ("BRAND_MODE.md", "INTERNAL_25.md", "EXTERNAL_25.md",
                 "README.md", "CLAUDE_CODE.md", "slide-sets.json"):
        source = ASSETS / "handoff-25" / name
        if source.is_file():
            shutil.copyfile(source, brand / name)

    (out / "README.md").write_text(README, encoding="utf-8")
    (out / "INSTRUCTIONS.md").write_text(INSTRUCTIONS, encoding="utf-8")
    (out / "OPEN_QUESTIONS.md").write_text(OPEN_QUESTIONS, encoding="utf-8")

    gaps = C.gaps()
    summary = {
        "generated": date.today().isoformat(),
        "archetypes_built": len([e for entries in data["sets"].values()
                                 for e in entries if e["built"]]),
        "archetypes_total": sum(len(e) for e in data["sets"].values()),
        "previews_rendered": shots,
        "contract_gaps": gaps,
        **counts,
    }
    (out / "state.json").write_text(json.dumps(summary, indent=2))
    print(json.dumps(summary, indent=2))

    archive = shutil.make_archive(str(out), "zip", root_dir=out.parent,
                                  base_dir=out.name)
    print(f"\n{out}\n{archive}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
